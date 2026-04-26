"""
「祟り神の章」ゲーム性提案 PowerPoint ジェネレーター
出力: proposals/atarigami_proposal_v1.pptx
"""
import io
import os
import sys
import requests
from PIL import Image as PILImage, ImageEnhance, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__), "proposals", "atarigami_proposal_v1.pptx")

# ── Colors ──────────────────────────────────────────────
C_BG      = RGBColor(0x08, 0x08, 0x18)   # 深夜の紺黒
C_CARD    = RGBColor(0x14, 0x14, 0x2C)   # カード背景
C_CARD2   = RGBColor(0x1A, 0x0A, 0x2A)   # 紫寄りカード
C_GOLD    = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_RED     = RGBColor(0xCC, 0x22, 0x22)   # 怨念の赤
C_BLUE    = RGBColor(0x33, 0x55, 0xCC)   # 成仏の青
C_PURPLE  = RGBColor(0x88, 0x22, 0xAA)   # 紫
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LTGRAY  = RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY    = RGBColor(0x88, 0x88, 0x88)
C_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
C_GREEN   = RGBColor(0x22, 0xCC, 0x66)
C_YELLOW  = RGBColor(0xFF, 0xEE, 0x44)
C_DARKRED = RGBColor(0x44, 0x00, 0x00)
C_DARKBLUE= RGBColor(0x00, 0x00, 0x44)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── Image download & process ────────────────────────────
IMAGE_URLS = [
    "https://upload.wikimedia.org/wikipedia/commons/thumb/8/8f/Shunkosai_Hokuei_Obake.jpg/480px-Shunkosai_Hokuei_Obake.jpg",
    "https://upload.wikimedia.org/wikipedia/commons/thumb/3/3e/Katsushika_Hokusai_-_Hyaku_monogatari_-_Sarayashiki.jpg/320px-Katsushika_Hokusai_-_Hyaku_monogatari_-_Sarayashiki.jpg",
    "https://upload.wikimedia.org/wikipedia/commons/thumb/4/44/Toriyama_Sekiyen_-_%E7%99%BE%E5%99%A8%E5%BE%92%E7%84%B6%E5%A4%9C%E8%A1%8C_-_Hyakki_Yagyo_Emaki.jpg/480px-Toriyama_Sekiyen_-_%E7%99%BE%E5%99%A8%E5%BE%92%E7%84%B6%E5%A4%9C%E8%A1%8C_-_Hyakki_Yagyo_Emaki.jpg",
]

def download_image(url, timeout=8):
    try:
        resp = requests.get(url, timeout=timeout,
                            headers={"User-Agent": "Mozilla/5.0"})
        if resp.status_code == 200:
            return PILImage.open(io.BytesIO(resp.content)).convert("RGBA")
    except Exception as e:
        print(f"  画像取得失敗: {url} → {e}")
    return None


def generate_ghost_art(width=960, height=540):
    """ネット画像が取得できない場合のPillow製幽霊画像"""
    import random
    from PIL import ImageDraw
    random.seed(42)
    img = PILImage.new("RGBA", (width, height), (4, 4, 18, 255))
    draw = ImageDraw.Draw(img, "RGBA")

    # 靄(もや)のレイヤー
    for _ in range(120):
        x = random.randint(0, width)
        y = random.randint(0, height)
        r = random.randint(30, 180)
        alpha = random.randint(8, 35)
        col = random.choice([
            (80, 20, 120, alpha),   # 紫
            (20, 20, 80, alpha),    # 青
            (120, 20, 20, alpha),   # 赤
            (40, 40, 80, alpha),    # 薄青
        ])
        draw.ellipse([x - r, y - r, x + r, y + r], fill=col)

    # 縦の光条（怪光）
    for _ in range(18):
        x = random.randint(0, width)
        w = random.randint(1, 4)
        alpha = random.randint(15, 55)
        col = random.choice([
            (180, 160, 60, alpha),  # 金
            (100, 20, 160, alpha),  # 紫
            (20, 40, 180, alpha),   # 青
        ])
        draw.rectangle([x, 0, x + w, height], fill=col)

    # 中央に人影シルエット（抽象的）
    cx, cy = width // 2 + 80, height // 2 + 20
    # 頭
    draw.ellipse([cx - 28, cy - 160, cx + 28, cy - 100], fill=(15, 5, 30, 200))
    # 体
    draw.polygon([
        (cx - 22, cy - 105), (cx + 22, cy - 105),
        (cx + 45, cy + 80), (cx - 45, cy + 80),
    ], fill=(10, 5, 25, 190))
    # 着物の裾（広がる）
    draw.polygon([
        (cx - 45, cy + 70), (cx + 45, cy + 70),
        (cx + 90, cy + 200), (cx - 90, cy + 200),
    ], fill=(8, 4, 22, 170))
    # 腕
    draw.polygon([
        (cx - 22, cy - 90), (cx - 80, cy + 30),
        (cx - 90, cy + 40), (cx - 30, cy - 80),
    ], fill=(12, 5, 28, 180))
    draw.polygon([
        (cx + 22, cy - 90), (cx + 80, cy + 30),
        (cx + 90, cy + 40), (cx + 30, cy - 80),
    ], fill=(12, 5, 28, 180))

    # 足（ない → 幽霊感）
    # 目の光
    draw.ellipse([cx - 12, cy - 140, cx - 2, cy - 128], fill=(180, 60, 60, 200))
    draw.ellipse([cx + 2,  cy - 140, cx + 12, cy - 128], fill=(180, 60, 60, 200))

    # 周囲ビネット
    for r in range(0, min(width, height) // 2, 12):
        alpha = int(60 * (1 - r / (min(width, height) / 2)))
        cx2, cy2 = width // 2, height // 2
        draw.ellipse([cx2 - r - 12, cy2 - r - 12, cx2 + r + 12, cy2 + r + 12],
                     outline=(0, 0, 0, alpha), width=12)

    return img.convert("RGB")

def process_ghost_image(pil_img, width=480, height=320,
                         darken=0.70, purple_alpha=60, blur=1.0):
    """ダーク＆パープルティント処理（すでに暗い生成画像にも対応）"""
    img = pil_img.resize((width, height), PILImage.LANCZOS).convert("RGBA")
    # 暗くする（生成画像は既に暗いので控えめに）
    img = ImageEnhance.Brightness(img).enhance(darken)
    # パープルオーバーレイ
    overlay = PILImage.new("RGBA", img.size, (80, 10, 120, purple_alpha))
    img = PILImage.alpha_composite(img, overlay)
    # ソフトブラー
    img = img.filter(ImageFilter.GaussianBlur(blur))
    return img.convert("RGB")

def pil_to_stream(pil_img):
    buf = io.BytesIO()
    pil_img.save(buf, format="PNG")
    buf.seek(0)
    return buf

# ── pptx helpers ────────────────────────────────────────
def new_slide(prs, bg_color=C_BG):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_image_to_slide(slide, pil_img, left, top, width, height):
    stream = pil_to_stream(pil_img)
    slide.shapes.add_picture(stream, left, top, width, height)

def rect(slide, left, top, width, height, color, alpha=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    return shp

def rect_border(slide, left, top, width, height, fill_color, border_color, border_pt=1.5):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = border_color
    shp.line.width = Pt(border_pt)
    return shp

def tb(slide, left, top, width, height, text, fontsize=11,
       bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "メイリオ"
    return txb

def label(slide, left, top, text, color=C_GOLD, size=10, bold=True):
    return tb(slide, left, top, Inches(9), Emu(360000), text, size, bold=bold, color=color)

def card(slide, left, top, width, height, color=C_CARD):
    return rect(slide, left, top, width, height, color)

def arrow_right(slide, cx, cy, size=Emu(150000), color=C_GRAY):
    shp = slide.shapes.add_shape(13, cx - size//2, cy - size//4, size, size//2)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def arrow_down(slide, cx, cy, w=Emu(120000), h=Emu(200000), color=C_GRAY):
    shp = slide.shapes.add_shape(13, cx - w//2, cy, w, h)
    # rotate 90
    shp.rotation = 90
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def divider(slide, y, color=C_GRAY):
    rect(slide, Inches(0.2), y, Inches(9.6), Emu(20000), color)

# ── Slide builders ──────────────────────────────────────

def slide_title(prs, ghost_img=None):
    s = new_slide(prs)

    # 背景に霊の画像（あれば）
    if ghost_img:
        processed = process_ghost_image(ghost_img, width=960, height=540,
                                         darken=0.25, purple_alpha=130)
        add_image_to_slide(s, processed,
                           Inches(0), Inches(0), SLIDE_W, SLIDE_H)

    # 左側グラデーション風オーバーレイ
    rect(s, Inches(0), Inches(0), Inches(5.5), SLIDE_H, RGBColor(0x04, 0x04, 0x10))

    # 赤い縦線アクセント
    rect(s, Inches(0.35), Inches(0.5), Emu(30000), Inches(2.2), C_RED)

    # タイトル
    tb(s, Inches(0.5), Inches(0.55), Inches(5), Inches(0.7),
       "新機種ゲーム性提案書", 14, bold=True, color=C_GOLD)
    tb(s, Inches(0.5), Inches(1.15), Inches(5.2), Inches(1.2),
       "祟り神の章", 46, bold=True, color=C_WHITE)
    tb(s, Inches(0.5), Inches(2.3), Inches(5.0), Inches(0.5),
       "〜物語と出玉が交差するスロット〜", 13, italic=True, color=C_LTGRAY)

    # コンセプトキャッチ
    rect(s, Inches(0.5), Inches(3.1), Inches(4.8), Emu(580000), RGBColor(0x22, 0x00, 0x00))
    rect(s, Inches(0.5), Inches(3.1), Emu(60000), Emu(580000), C_RED)
    tb(s, Inches(0.65), Inches(3.18), Inches(4.5), Emu(520000),
       "「倒した相手を、後から理解する。」\n\n普通ATで祟り神を「戦って倒す」台。\n特別ATでその祟り神の「なぜそうなったか」が分かる台。",
       11, color=C_WHITE)

    # 右下タグ
    tb(s, Inches(6.5), Inches(4.8), Inches(3.3), Emu(300000),
       "スマスロ（L型）／ 実稼働データ453機種173週 知見ベース",
       8, color=C_GRAY, align=PP_ALIGN.RIGHT)

    return s


def slide_concept(prs):
    s = new_slide(prs)

    # ヘッダー
    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "CONCEPT  ──  なぜこれが新しいのか？", 14, bold=True, color=C_GOLD)

    # 中央：感情逆転の図
    CX = Inches(5)

    # 左ブロック：普通AT
    rect_border(s, Inches(0.3), Inches(1.0), Inches(3.8), Inches(2.4),
                RGBColor(0x22, 0x00, 0x00), C_RED, 2)
    tb(s, Inches(0.4), Inches(1.05), Inches(3.5), Emu(380000),
       "普通AT（バトル型）", 12, bold=True, color=C_RED)
    tb(s, Inches(0.4), Inches(1.5), Inches(3.6), Inches(1.6),
       "祟り神を「敵」として戦い、倒す\n\n⚔ ザコ戦 → 中ボス → ボス\n⚔ 勝利でセット継続・SP蓄積\n⚔ 敗北で悔しさ→リベンジ欲求",
       10, color=C_LTGRAY)

    # 中央：矢印＋感情逆転ラベル
    rect(s, Inches(4.3), Inches(1.9), Inches(1.4), Emu(300000), RGBColor(0x10, 0x10, 0x10))
    tb(s, Inches(4.1), Inches(1.95), Inches(1.8), Emu(280000),
       "感情が\n逆転する", 10, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    tb(s, Inches(4.05), Inches(2.5), Inches(1.9), Emu(300000),
       "敵 → 共感", 13, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

    # 右ブロック：特別AT
    rect_border(s, Inches(5.9), Inches(1.0), Inches(3.8), Inches(2.4),
                RGBColor(0x00, 0x00, 0x33), C_BLUE, 2)
    tb(s, Inches(6.0), Inches(1.05), Inches(3.5), Emu(380000),
       "特別AT（章型）", 12, bold=True, color=C_BLUE)
    tb(s, Inches(6.0), Inches(1.5), Inches(3.6), Inches(1.6),
       "祟り神の「なぜそうなったか」を知る\n\n📖 村人の視点・僧侶の視点・家族の視点\n📖 断片が集まって真実が見える\n📖 100〜1500枚獲得チャンス",
       10, color=C_LTGRAY)

    # 下段：業界比較テーブル
    rect(s, Inches(0.3), Inches(3.55), Inches(9.4), Emu(50000), C_GOLD)
    headers = ["比較軸", "従来機", "本提案（祟り神の章）"]
    hx = [Inches(0.3), Inches(2.1), Inches(5.5)]
    hw = [Inches(1.7), Inches(3.3), Inches(4.1)]
    for i, (h, x, w) in enumerate(zip(headers, hx, hw)):
        tb(s, x, Inches(3.6), w, Emu(320000), h, 9, bold=True,
           color=C_BG if i < 2 else C_WHITE)

    rows = [
        ("来店動機", "設定・期待値", "物語の続きが見たい"),
        ("リピート設計", "1回来店で完結", "複数来店で完成する"),
        ("プレイ体験", "全員が同じ", "章のランダム順で変わる"),
    ]
    ry = Inches(3.95)
    for j, (axis, old, new) in enumerate(rows):
        bg = RGBColor(0x10, 0x10, 0x22) if j % 2 == 0 else RGBColor(0x14, 0x14, 0x28)
        rect(s, Inches(0.3), ry, Inches(9.4), Emu(340000), bg)
        tb(s, Inches(0.35), ry + Emu(30000), Inches(1.65), Emu(290000),
           axis, 9, color=C_LTGRAY)
        tb(s, Inches(2.15), ry + Emu(30000), Inches(3.2), Emu(290000),
           old, 9, color=C_GRAY)
        tb(s, Inches(5.55), ry + Emu(30000), Inches(3.9), Emu(290000),
           new, 9, bold=True, color=C_GOLD)
        ry += Emu(350000)

    return s


def slide_emotion_map(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "PLAYER EMOTION MAP  ──  どこで楽しませ、どこで揺さぶるか", 14, bold=True, color=C_GOLD)

    stages = [
        ("①着席", "自分の進捗\nが迎えてくれる", C_GREEN),
        ("②通常時", "レア役で\nSP積み上げ", C_LTGRAY),
        ("③普通AT", "ザコ→中ボス\n→ボス 一喜一憂", C_RED),
        ("④特別AT\n突入", "どの章が来る?\nワクワク", C_BLUE),
        ("⑤章の中", "別視点で\n同じ事件を見る", C_PURPLE),
        ("⑥加護\n獲得", "コレクション\n達成感", C_GOLD),
        ("⑦第5章AT", "同時多発\n「超楽しい」", C_ORANGE),
        ("⑧第0章", "1500枚\n感動のラスト", C_YELLOW),
    ]

    bw = Inches(1.12)
    bh = Inches(2.5)
    by = Inches(0.95)
    for i, (name, desc, col) in enumerate(stages):
        bx = Inches(0.15) + i * (bw + Emu(30000))
        # 下から伸びる感情バー
        bar_h_ratio = 0.3 + (i / 7) * 0.6
        bar_h = int(bh * bar_h_ratio)
        bar_y = by + bh - bar_h

        rect(s, bx, by, bw, bh, RGBColor(0x12, 0x12, 0x28))
        rect(s, bx + Emu(100000), bar_y, bw - Emu(200000), bar_h, col)
        tb(s, bx, by + bh + Emu(60000), bw, Emu(350000),
           name, 8, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, bx, bar_y + Emu(60000), bw, Emu(400000),
           desc, 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)

        # 矢印
        if i < len(stages) - 1:
            ax = bx + bw + Emu(0)
            ay = by + bh // 2
            rect(s, ax, ay - Emu(30000), Emu(30000), Emu(60000), C_GRAY)

    # 下段：特記事項
    rect(s, Inches(0.15), Inches(4.8), Inches(9.7), Emu(330000),
         RGBColor(0x1A, 0x10, 0x00))
    tb(s, Inches(0.3), Inches(4.86), Inches(9.3), Emu(280000),
       "★ 核心：普通ATで「倒した敵」が、特別ATで「共感できる存在」になる。  この感情の逆転が来店継続の最大動機。",
       10, bold=True, color=C_GOLD)

    return s


def slide_chapter_system(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "CHAPTER SYSTEM  ──  DQ4オムニバス型・ランダム章解放設計", 14, bold=True, color=C_GOLD)

    chapters = [
        ("第1章", "村人の視点", "なぜ祟られたのか\n分からない", C_LTGRAY, "リプレイの力"),
        ("第2章", "僧侶の視点", "封じるしかなかった\nあの日の決断", C_LTGRAY, "ベルの力"),
        ("第3章", "家族の視点", "あの人がなぜ\n祟り神に…", C_LTGRAY, "選択の力"),
        ("第4章", "祟り神の記憶", "守りたかった\nだけなのに", C_RED, "魂の力"),
        ("第5章", "全視点収束", "加護4個で\n超楽しいAT", C_GOLD, "同時多発解放"),
        ("第0章", "真実＆成仏", "1500枚一か八か\n感動のラスト", C_BLUE, "全章クリア限定"),
    ]

    cw = Inches(1.5)
    ch = Inches(3.0)
    cy = Inches(0.95)
    for i, (num, viewpoint, story, col, kago) in enumerate(chapters):
        cx = Inches(0.2) + i * (cw + Emu(100000))

        # 特別カラー
        if num == "第5章":
            bg = RGBColor(0x28, 0x20, 0x00)
            border_col = C_GOLD
        elif num == "第0章":
            bg = RGBColor(0x00, 0x00, 0x35)
            border_col = C_BLUE
        else:
            bg = RGBColor(0x14, 0x14, 0x2A)
            border_col = C_PURPLE

        rect_border(s, cx, cy, cw, ch, bg, border_col, 1.5)
        tb(s, cx + Emu(50000), cy + Emu(60000), cw - Emu(100000), Emu(350000),
           num, 11, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, cx + Emu(50000), cy + Emu(380000), cw - Emu(100000), Emu(350000),
           viewpoint, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(s, cx + Emu(50000), cy + Emu(700000), cw - Emu(100000), Emu(500000),
           story, 8.5, color=C_LTGRAY, align=PP_ALIGN.CENTER)

        # 加護
        if num not in ("第5章", "第0章"):
            rect(s, cx + Emu(50000), cy + ch - Emu(500000), cw - Emu(100000), Emu(420000),
                 RGBColor(0x20, 0x10, 0x30))
            tb(s, cx + Emu(80000), cy + ch - Emu(460000), cw - Emu(160000), Emu(380000),
               f"加護：\n{kago}", 7.5, color=C_PURPLE, align=PP_ALIGN.CENTER)
        elif num == "第5章":
            rect(s, cx + Emu(50000), cy + ch - Emu(500000), cw - Emu(100000), Emu(420000),
                 RGBColor(0x28, 0x18, 0x00))
            tb(s, cx + Emu(80000), cy + ch - Emu(460000), cw - Emu(160000), Emu(380000),
               "加護4個\n全部発動！", 7.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        else:
            rect(s, cx + Emu(50000), cy + ch - Emu(500000), cw - Emu(100000), Emu(420000),
                 RGBColor(0x00, 0x00, 0x30))
            tb(s, cx + Emu(80000), cy + ch - Emu(460000), cw - Emu(160000), Emu(380000),
               "成仏\nチャレンジ", 7.5, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # ランダム順の説明
    rect(s, Inches(0.2), Inches(4.1), Inches(4.6), Emu(550000), RGBColor(0x10, 0x18, 0x10))
    tb(s, Inches(0.35), Inches(4.15), Inches(4.3), Emu(500000),
       "📌 章はランダム順で解放。「俺は3章から入った、お前は？」\n　→ ホールで会話が生まれる、体験がプレイヤーごとに変わる",
       9, color=C_GREEN)

    rect(s, Inches(5.0), Inches(4.1), Inches(4.8), Emu(550000), RGBColor(0x18, 0x10, 0x00))
    tb(s, Inches(5.15), Inches(4.15), Inches(4.5), Emu(500000),
       "📌 章クリアで加護を獲得→第5章のATが強化される\n　→ 来るたびに体験が豊かになる「来店価値の積み上げ」",
       9, color=C_GOLD)

    return s


def slide_kago_flow(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "加護システム → 第5章「超楽しい」の設計", 14, bold=True, color=C_GOLD)

    # 加護4つ
    kagos = [
        ("第1章", "リプレイの力", "リプレイ成立で\nAT抽選が走る", C_LTGRAY),
        ("第2章", "ベルの力", "打順ベル成立で\n上乗せ抽選", C_LTGRAY),
        ("第3章", "選択の力", "毎セット\n2択チャレンジ発生", C_LTGRAY),
        ("第4章", "魂の力", "セット継続率が\n底上げされる", C_RED),
    ]

    kw = Inches(1.8)
    kh = Inches(1.7)
    ky = Inches(0.95)
    for i, (chap, name, effect, col) in enumerate(kagos):
        kx = Inches(0.25) + i * (kw + Emu(120000))
        rect_border(s, kx, ky, kw, kh, RGBColor(0x18, 0x10, 0x28), C_PURPLE, 1.5)
        tb(s, kx, ky + Emu(60000), kw, Emu(300000),
           chap, 9, color=C_PURPLE, align=PP_ALIGN.CENTER)
        tb(s, kx, ky + Emu(340000), kw, Emu(330000),
           name, 10, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, kx, ky + Emu(650000), kw, Emu(420000),
           effect, 8.5, color=C_LTGRAY, align=PP_ALIGN.CENTER)

    # 矢印 DOWN
    tb(s, Inches(4.5), Inches(2.7), Inches(1), Inches(0.4),
       "⬇ 全部集まると", 10, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 第5章ボックス
    rect_border(s, Inches(2.0), Inches(3.1), Inches(6.0), Inches(1.5),
                RGBColor(0x28, 0x20, 0x00), C_GOLD, 2.5)
    tb(s, Inches(2.0), Inches(3.15), Inches(6.0), Emu(380000),
       "第5章AT ── 同時多発で何かが起き続ける「超楽しい」",
       13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    events = [
        "リプレイ → AT抽選！",
        "ベル → 上乗せ！",
        "セット開始 → 2択！",
        "継続率底上げ → なかなか終わらない",
    ]
    ex = Inches(2.15)
    for ev in events:
        tb(s, ex, Inches(3.6), Inches(5.7), Emu(290000),
           f"✔ {ev}", 10, color=C_WHITE)
        ex  # same x (intentional: stack vertically via y)
    # Actually render them stacked
    ey = Inches(3.6)
    for ev in events:
        tb(s, Inches(2.15), ey, Inches(5.7), Emu(290000),
           f"✔ {ev}", 10, color=C_WHITE)
        ey += Emu(290000)

    # 注記
    rect(s, Inches(0.25), Inches(4.7), Inches(9.5), Emu(360000), RGBColor(0x10, 0x10, 0x10))
    tb(s, Inches(0.4), Inches(4.75), Inches(9.1), Emu(310000),
       "加護が少ないまま第5章に来ると体験が薄い → 「また章を集めに来る」動機が生まれる設計",
       9.5, bold=True, color=C_ORANGE)

    return s


def slide_chapter0(prs):
    s = new_slide(prs)

    # 暗い青黒背景
    bg = slide = s
    rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, RGBColor(0x00, 0x00, 0x18))

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x00, 0x00, 0x30))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "第0章  ──  全章コンプリート者だけが辿り着く「真実と成仏」", 14, bold=True, color=C_BLUE)

    # 左：ストーリー収束
    rect_border(s, Inches(0.2), Inches(0.95), Inches(4.5), Inches(3.5),
                RGBColor(0x00, 0x00, 0x28), C_BLUE, 2)
    tb(s, Inches(0.3), Inches(1.0), Inches(4.2), Emu(380000),
       "「真の章が解禁されました」", 11, bold=True, color=C_BLUE)
    tb(s, Inches(0.3), Inches(1.55), Inches(4.2), Inches(2.5),
       "全視点が交差して、真実が一枚の絵として完成する\n\n"
       "・なぜ村人は祟られたのか\n"
       "・なぜ僧侶は封じるしかなかったのか\n"
       "・なぜ家族は気づけなかったのか\n"
       "・なぜ祟り神はそうなったのか\n\n"
       "──  全部が繋がる",
       10, color=C_LTGRAY)

    # 右：成仏チャレンジ
    rect_border(s, Inches(5.0), Inches(0.95), Inches(4.8), Inches(3.5),
                RGBColor(0x10, 0x00, 0x00), C_RED, 2)
    tb(s, Inches(5.1), Inches(1.0), Inches(4.5), Emu(380000),
       "成仏チャレンジ  ──  一か八か", 11, bold=True, color=C_RED)

    # 成否分岐
    rect(s, Inches(5.1), Inches(1.6), Inches(2.1), Inches(1.0), RGBColor(0x00, 0x22, 0x00))
    tb(s, Inches(5.1), Inches(1.65), Inches(2.0), Emu(350000),
       "成功", 16, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tb(s, Inches(5.1), Inches(2.05), Inches(2.0), Emu(380000),
       "1500枚\n＋「救えた」感動", 10, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    rect(s, Inches(7.5), Inches(1.6), Inches(2.1), Inches(1.0), RGBColor(0x22, 0x00, 0x00))
    tb(s, Inches(7.5), Inches(1.65), Inches(2.0), Emu(350000),
       "失敗", 16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    tb(s, Inches(7.5), Inches(2.05), Inches(2.0), Emu(380000),
       "100枚\n＋「また来て救おう」", 10, color=C_RED, align=PP_ALIGN.CENTER)

    tb(s, Inches(5.1), Inches(3.0), Inches(4.5), Emu(400000),
       "失敗しても「終わり」ではなく\n「また来たい理由」に変わる。\n物語はまだ終わっていないから。",
       9, italic=True, color=C_LTGRAY)

    # 下段
    rect(s, Inches(0.2), Inches(4.6), Inches(9.6), Emu(420000), RGBColor(0x00, 0x00, 0x30))
    tb(s, Inches(0.35), Inches(4.65), Inches(9.2), Emu(360000),
       "★ 第0章は「ゴール」ではなく「また来たい気持ち」の設計。次の祟り神ストーリーへ続く可能性も。",
       10, bold=True, color=C_BLUE)

    return s


def slide_spec(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "SPEC & BENCHMARK  ──  スペックと近似機種比較", 14, bold=True, color=C_GOLD)

    # 左：スペック
    rect_border(s, Inches(0.2), Inches(0.95), Inches(4.0), Inches(3.3),
                C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.35), Inches(1.0), Inches(3.6), Emu(360000),
       "推奨スペック（目安）", 11, bold=True, color=C_GOLD)

    specs = [
        ("タイプ", "スマスロ（L型）"),
        ("天井", "600G（SP MAX時 300G）"),
        ("純増", "3.8枚/G"),
        ("機械割", "設定1: 97.5%  /  設定6: 109%"),
        ("特別AT出玉", "100〜1500枚（加護数・章で変動）"),
        ("設計CV目標", "0.20〜0.25（設定非依存）"),
    ]
    sy = Inches(1.4)
    for j, (k, v) in enumerate(specs):
        bg = RGBColor(0x12, 0x12, 0x26) if j % 2 == 0 else RGBColor(0x16, 0x16, 0x2C)
        rect(s, Inches(0.25), sy, Inches(3.85), Emu(340000), bg)
        tb(s, Inches(0.3), sy + Emu(30000), Inches(1.2), Emu(290000),
           k, 8.5, color=C_GRAY)
        tb(s, Inches(1.55), sy + Emu(30000), Inches(2.4), Emu(290000),
           v, 9, bold=True, color=C_WHITE)
        sy += Emu(350000)

    # 右：ベンチマーク表
    bx = Inches(4.4)
    rect(s, bx, Inches(0.95), Inches(5.45), Emu(380000), RGBColor(0x22, 0x18, 0x00))
    bench_headers = ["比較軸", "東京喰種", "モンキーV", "番長4", "本提案"]
    bhx = [bx + Emu(50000), bx + Emu(950000), bx + Emu(1750000),
           bx + Emu(2550000), bx + Emu(3350000)]
    bhw = [Emu(870000)] * 5
    for h, hx, hw in zip(bench_headers, bhx, bhw):
        tb(s, hx, Inches(0.98), hw, Emu(340000), h, 9, bold=True,
           color=C_GOLD, align=PP_ALIGN.CENTER)

    bench_rows = [
        ("CV値", "0.18", "0.23", "0.68", "0.20〜0.25"),
        ("来店動機", "世界観", "SP引継", "期待値", "物語続き"),
        ("一か八か", "○", "△", "◎", "◎（最大\n1500枚）"),
        ("20代訴求", "○", "△", "△", "◎"),
        ("後半維持率", "73.7%", "68.2%", "29.6%", "目標65%+"),
    ]
    rby = Inches(1.42)
    for j, row in enumerate(bench_rows):
        rbg = RGBColor(0x10, 0x10, 0x22) if j % 2 == 0 else RGBColor(0x14, 0x14, 0x28)
        rect(s, bx, rby, Inches(5.45), Emu(350000), rbg)
        for k, (val, hx, hw) in enumerate(zip(row, bhx, bhw)):
            col = C_GOLD if k == 4 else (C_RED if val in ("0.68", "29.6%") else C_LTGRAY)
            tb(s, hx, rby + Emu(30000), hw, Emu(295000), val, 8.5,
               color=col, align=PP_ALIGN.CENTER)
        rby += Emu(355000)

    # 下コメント
    rect(s, Inches(0.2), Inches(4.55), Inches(9.6), Emu(480000), RGBColor(0x12, 0x18, 0x12))
    tb(s, Inches(0.35), Inches(4.6), Inches(9.2), Emu(430000),
       "CV0.20〜0.25（東京喰種・モンキーターンV水準）を目標とし、設定1でも「体験の質が変わらない」台を設計。\n"
       "後半維持率65%以上、3ヶ月後も稼働している台を目指す。",
       9, color=C_GREEN)

    return s


def slide_differentiators(prs, ghost_img=None):
    s = new_slide(prs)

    if ghost_img:
        processed = process_ghost_image(ghost_img, width=480, height=540,
                                         darken=0.18, purple_alpha=150)
        add_image_to_slide(s, processed,
                           Inches(6.5), Inches(0), Inches(3.5), SLIDE_H)

    rect(s, Inches(0), Inches(0), SLIDE_W, Emu(440000), RGBColor(0x10, 0x00, 0x20))
    tb(s, Inches(0.3), Emu(50000), Inches(9), Emu(360000),
       "DIFFERENTIATORS  ──  差別化の核心3点", 14, bold=True, color=C_GOLD)

    points = [
        (C_RED, "①「倒した敵を後から理解する」感情逆転設計",
         "普通ATで倒した祟り神が、特別ATで共感できる存在になる。\n"
         "「敵 → 理解 → 救いたい」という感情の逆転は業界に前例がない。"),
        (C_BLUE, "②DQ4型オムニバス章システム（業界初）",
         "ランダム順に届く5つの視点。全部集まって初めて真実が見える。\n"
         "「俺は3章から入った」という会話がホールで生まれる体験設計。"),
        (C_PURPLE, "③加護の持ち越しで「第5章が超楽しくなる」設計",
         "章をクリアするほど第5章のATが強化される。\n"
         "来るたびに体験が豊かになる「来店価値の積み上げ」構造。"),
    ]

    py = Inches(0.95)
    for col, title, desc in points:
        rect(s, Inches(0.2), py, Emu(60000), Inches(1.2), col)
        rect(s, Inches(0.35), py, Inches(6.0), Inches(1.2), RGBColor(0x14, 0x14, 0x2A))
        tb(s, Inches(0.45), py + Emu(60000), Inches(5.7), Emu(340000),
           title, 11, bold=True, color=col)
        tb(s, Inches(0.45), py + Emu(380000), Inches(5.7), Emu(480000),
           desc, 9.5, color=C_LTGRAY)
        py += Inches(1.35)

    # 右上：議論したい論点
    rect_border(s, Inches(6.6), Inches(0.95), Inches(3.2), Inches(3.5),
                RGBColor(0x10, 0x10, 0x10), C_GRAY, 1)
    tb(s, Inches(6.7), Inches(1.0), Inches(2.9), Emu(360000),
       "議論したい論点", 10, bold=True, color=C_GRAY)
    qs = [
        "Q1: 第0章後のリセット設計\n　ループ か 次の祟り神IP へ？",
        "Q2: 加護は全台共通 or\n　スマスロ個人管理？",
        "Q3: IPはオリジナル推奨\n　「この台でしか見られない\n　　ストーリー」の価値最大化",
    ]
    qy = Inches(1.45)
    for q in qs:
        tb(s, Inches(6.7), qy, Inches(3.0), Emu(480000), q, 8.5, color=C_GRAY)
        qy += Emu(490000)

    # フッター
    rect(s, Inches(0.2), Inches(4.6), Inches(9.6), Emu(390000), RGBColor(0x08, 0x08, 0x18))
    tb(s, Inches(0.4), Inches(4.65), Inches(9.0), Emu(340000),
       "分析ベース：実稼働データ 453機種・173週（2022〜2026）/ 設計CV目標 0.20〜0.25",
       8.5, color=C_GRAY)

    return s


# ── Main ────────────────────────────────────────────────
def main():
    print("=" * 50)
    print("  祟り神の章 ゲーム性提案 PowerPoint ジェネレーター")
    print("=" * 50)

    # 画像取得
    print("\n🖼  画像をダウンロード中...")
    ghost_img = None
    for url in IMAGE_URLS:
        print(f"  試行: {url[:60]}...")
        ghost_img = download_image(url)
        if ghost_img:
            print(f"  ✅ 取得成功！")
            break
    if not ghost_img:
        print("  ⚠ 画像取得失敗。Pillowで幽霊画像を生成します...")
        ghost_img = generate_ghost_art(960, 540)
        print("  ✅ 生成完了")

    print("\n📊 スライド生成中...")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, ghost_img)
    print("  1/7 タイトル")
    slide_concept(prs)
    print("  2/7 コンセプト")
    slide_emotion_map(prs)
    print("  3/7 感情マップ")
    slide_chapter_system(prs)
    print("  4/7 章システム")
    slide_kago_flow(prs)
    print("  5/7 加護フロー")
    slide_chapter0(prs)
    print("  6/7 第0章")
    slide_spec(prs)
    print("  7/7 スペック・ベンチマーク")
    slide_differentiators(prs, ghost_img)
    print("  8/8 差別化3点")

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n✅ 保存完了: {OUT_PATH}")

if __name__ == "__main__":
    main()
