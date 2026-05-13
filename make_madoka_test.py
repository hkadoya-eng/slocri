"""
まどか外伝 カラースキーム試作（スライド1・2のみ）
明るい背景 × 枠色に合わせた薄塗りボックス
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFilter
import io

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
OUT_PATH = r"C:\Users\h.kadoya\Desktop\madoka_test_color.pptx"

# ── 新カラーパレット ──────────────────────────────────────────
# テキスト
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)  # 濃紺（本文）
C_MID_GRAY  = RGBColor(0x55, 0x55, 0x70)  # 中間グレー（補足）
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # 白（ヘッダーバー内）

# アクセント枠・ラベル（やや濃いめ）
C_PINK  = RGBColor(0xCC, 0x22, 0xAA)
C_LPINK = RGBColor(0xDD, 0x55, 0xBB)
C_CYAN  = RGBColor(0x11, 0x99, 0xBB)
C_GOLD  = RGBColor(0xBB, 0x88, 0x00)
C_GREEN = RGBColor(0x22, 0xAA, 0x66)
C_RED   = RGBColor(0xCC, 0x33, 0x33)
C_GRAY  = RGBColor(0x88, 0x88, 0x99)

# ボックス塗り（枠色の薄い類似色）
C_FILL_PINK    = RGBColor(0xFF, 0xEE, 0xF8)
C_FILL_CYAN    = RGBColor(0xEE, 0xF6, 0xFF)
C_FILL_GOLD    = RGBColor(0xFF, 0xFB, 0xEE)
C_FILL_GREEN   = RGBColor(0xEE, 0xFF, 0xF5)
C_FILL_RED     = RGBColor(0xFF, 0xEE, 0xEE)
C_FILL_DEFAULT = RGBColor(0xF5, 0xF0, 0xFA)  # 薄ラベンダー
C_FILL_GRAY    = RGBColor(0xF2, 0xF2, 0xF5)

FONT_H = "游明朝"
FONT_B = "メイリオ"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ── ヘルパー ─────────────────────────────────────────────────

def make_bg(slide, glow_x=0.5, glow_y=0.4):
    """薄ラベンダー×薄ピンクグロー（明るい背景）"""
    W, H = 1280, 720
    img = Image.new("RGB", (W, H), (0xED, 0xE8, 0xF5))
    draw = ImageDraw.Draw(img)
    cx, cy = int(W * glow_x), int(H * glow_y)
    for r in range(300, 0, -10):
        a = int(12 * (1 - r / 300))
        col = (min(0xED + a * 3, 255), min(0xD8 + a, 0xE8), min(0xF5, 255))
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=col)
    img = img.filter(ImageFilter.GaussianBlur(radius=3))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(0), Inches(0), SLIDE_W, SLIDE_H)


def rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def rect_b(slide, left, top, width, height, fill_color, border_color=None, border_pt=1.5):
    shape = rect(slide, left, top, width, height, fill_color)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def tb(slide, text, left, top, width, height,
       font_name=None, font_size=14, color=None,
       bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name or FONT_B
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color or C_DARK_TEXT
    return txb


def hdr(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x55, 0x00, 0x44)
    bar.line.fill.background()
    tb(slide, title, 0.15, 0.08, 8.5, 0.55,
       font_name=FONT_H, font_size=16, color=C_WHITE, bold=True)
    if subtitle:
        tb(slide, subtitle, 6.0, 0.08, 3.8, 0.52,
           font_name=FONT_B, font_size=10, color=RGBColor(0xFF, 0xCC, 0xEE),
           bold=False, align=PP_ALIGN.RIGHT)


def net_note(slide, text, design_comment):
    bar = rect(slide, 0, 5.15, 10, 0.475, RGBColor(0x22, 0x10, 0x30))
    tb(slide, f"設計：{design_comment}", 0.15, 5.17, 5.5, 0.4,
       font_name=FONT_B, font_size=9, color=C_GOLD, bold=True)
    tb(slide, text, 5.5, 5.17, 4.35, 0.4,
       font_name=FONT_B, font_size=8, color=C_GRAY,
       align=PP_ALIGN.RIGHT)


def badge(slide, text, left, top, w=1.4, h=0.32, bg=None, fc=None):
    rect_b(slide, left, top, w, h, bg or C_PINK,
           border_color=C_GOLD, border_pt=1)
    tb(slide, text, left + 0.05, top + 0.02, w - 0.1, h - 0.04,
       font_size=9, color=fc or C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def arrow_r(slide, left, top, length=0.5, label="", color=None):
    c = color or C_PINK
    body = slide.shapes.add_shape(1,
        Inches(left), Inches(top + 0.08), Inches(length * 0.75), Inches(0.09))
    body.fill.solid(); body.fill.fore_color.rgb = c; body.line.fill.background()
    tip = slide.shapes.add_shape(1,
        Inches(left + length * 0.72), Inches(top), Inches(length * 0.28), Inches(0.25))
    tip.fill.solid(); tip.fill.fore_color.rgb = c; tip.line.fill.background()


# ── スライド1：タイトル ───────────────────────────────────────

def slide1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.3, glow_y=0.5)
    hdr(slide, "スマスロ マギアレコード 魔法少女まどか☆マギカ外伝", "Part A: 機種解説")

    # 基本スペック（薄ピンク塗り × ピンク枠）
    rect_b(slide, 0.2, 0.78, 4.7, 2.30, C_FILL_PINK, border_color=C_PINK, border_pt=1.8)
    tb(slide, "基本スペック", 0.35, 0.82, 2.0, 0.35,
       font_name=FONT_H, font_size=12, color=C_PINK, bold=True)

    specs = [
        ("メーカー", "ミズホ（ユニバーサルエンタ）"),
        ("導入日",   "2025年4月7日"),
        ("純増",     "約2.6枚/G（AT中）"),
        ("天井",     "ボーナス間950pt+α"),
        ("機械割",   "97.6〜114.9%（設定1〜6）"),
        ("AT確率",   "1/654.6〜1/416.7"),
    ]
    y = 1.18
    for k, v in specs:
        tb(slide, f"◆ {k}：", 0.35, y, 1.55, 0.28,
           font_size=9, color=C_PINK, bold=True)
        tb(slide, v, 1.88, y, 2.85, 0.28,
           font_size=9, color=C_DARK_TEXT)
        y += 0.28

    # この台の3ポイント（薄シアン塗り × シアン枠）
    rect_b(slide, 5.15, 0.78, 4.65, 2.30, C_FILL_CYAN, border_color=C_CYAN, border_pt=1.8)
    tb(slide, "この台の3ポイント", 5.3, 0.82, 3.5, 0.35,
       font_name=FONT_H, font_size=12, color=C_CYAN, bold=True)

    pts = [
        ("①", "ストーリー×コンプリート体験",
         "8種ストーリーを集めてエンディングへ到達する物語体験型AT"),
        ("②", "穢れシステム「ドッペルモード」",
         "AT中に穢れが溜まると期待枚数3000枚超の爆発モードへ"),
        ("③", "まどマギ世界観を完全再現",
         "初代ファン向けサウンド・演出・キャラが深く作り込まれた完成度"),
    ]
    y = 1.18
    for num, title_pt, desc in pts:
        tb(slide, f"{num} {title_pt}", 5.3, y, 4.3, 0.28,
           font_size=10, color=C_GOLD, bold=True)
        tb(slide, desc, 5.3, y + 0.28, 4.3, 0.26,
           font_size=8, color=C_MID_GRAY)
        y += 0.54

    # アワードバッジ（濃い金バー × 金テキスト → アクセントで目立たせる）
    badge(slide, "パチスロアワード 2025 SILVER", 5.25, 2.84, w=4.45, h=0.22,
          bg=RGBColor(0x44, 0x22, 0x00), fc=C_GOLD)

    # キャッチコピー（薄ピンク塗り）
    rect_b(slide, 0.2, 3.12, 9.6, 0.52, C_FILL_PINK,
           border_color=C_PINK, border_pt=1)
    tb(slide,
       "初代『SLOT魔法少女まどか☆マギカ』のDNAを受け継ぎ、スマスロで進化──"
       "ストーリー攻略＋穢れ爆発の二重構造が今の等価市場で輝く",
       0.35, 3.16, 9.3, 0.44,
       font_size=9.5, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    # ボーナス4種（薄ラベンダー × シアン枠）
    rect_b(slide, 0.2, 3.73, 9.6, 1.35, C_FILL_DEFAULT,
           border_color=C_CYAN, border_pt=1.2)
    tb(slide, "ボーナス4種類（AT突入への足掛かり）",
       0.35, 3.77, 5.0, 0.3, font_size=10, color=C_CYAN, bold=True)

    bonuses = [
        ("エピソードボーナス", "50G継続",    C_PINK,  C_FILL_PINK),
        ("ビッグボーナス",     "30G+α",      C_GOLD,  C_FILL_GOLD),
        ("みたまボーナス",     "ベルナビ8回", C_CYAN,  C_FILL_CYAN),
        ("アリナミュージアム", "30G継続",    C_LPINK, C_FILL_PINK),
    ]
    bx = 0.3
    for bname, bdesc, bc, bf in bonuses:
        rect_b(slide, bx, 4.10, 2.28, 0.82, bf, border_color=bc, border_pt=1.5)
        tb(slide, bname, bx + 0.08, 4.14, 2.1, 0.32,
           font_size=9, color=bc, bold=True, align=PP_ALIGN.CENTER)
        tb(slide, bdesc, bx + 0.08, 4.46, 2.1, 0.28,
           font_size=8.5, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)
        bx += 2.38

    net_note(slide,
             "情報源: 一撃・ちょんぼりすた・パチセブン・P-WORLD",
             "スマスロ×物語体験×穢れ爆発の三位一体設計")


# ── スライド2：ゲームフロー ───────────────────────────────────

def slide2_gameflow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.5, glow_y=0.3)
    hdr(slide, "ゲームフロー全体図", "全ルートを蛇行2段で可視化")

    _bw  = 1.50
    _bh  = 0.80
    _bh2 = 1.00
    _gap = 0.13
    _mx  = 0.12
    _top1_y = 0.98
    _top2_y = 2.65
    _ann_y  = 3.90

    # 上段ラベル
    tb(slide, "上段ルート（メインフロー）", _mx, 0.72, 5.0, 0.26,
       font_size=9, color=C_PINK, bold=True)

    boxes_top = [
        ("通常時\n（ポイント蓄積）",  C_FILL_DEFAULT, C_PINK),
        ("マギア\nチャレンジCZ",      C_FILL_PINK,    C_PINK),
        ("ボーナス\n4種",             C_FILL_CYAN,    C_CYAN),
        ("マギア\nラッシュAT",        C_FILL_PINK,    C_PINK),
        ("ストーリー\nコンプリート",  C_FILL_GREEN,   C_GREEN),
        ("エンディング\n到達",        C_FILL_GOLD,    C_GOLD),
    ]
    for i, (label, fill, bc) in enumerate(boxes_top):
        lx = _mx + i * (_bw + _gap)
        rect_b(slide, lx, _top1_y, _bw, _bh, fill, border_color=bc, border_pt=1.5)
        tb(slide, label, lx + 0.06, _top1_y + 0.08, _bw - 0.12, _bh - 0.14,
           font_size=9.5, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    # 天井アノテーション
    tb(slide, "▼ 天井 950pt+α", _mx, _top1_y + _bh + 0.04, _bw, 0.22,
       font_size=8, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    # 上段矢印
    for i in range(len(boxes_top) - 1):
        ax = _mx + i * (_bw + _gap) + _bw + 0.01
        arrow_r(slide, ax, _top1_y + _bh / 2 - 0.12, length=_gap - 0.02)

    # 折り返し矢印
    r_edge = _mx + 5 * (_bw + _gap) + _bw

    fold_v = slide.shapes.add_shape(1,
        Inches(r_edge), Inches(_top1_y + _bh * 0.5),
        Inches(0.12), Inches(0.70))
    fold_v.fill.solid(); fold_v.fill.fore_color.rgb = C_PINK; fold_v.line.fill.background()

    fold_h_y = _top1_y + _bh * 0.5 + 0.70
    fold_h = slide.shapes.add_shape(1,
        Inches(_mx), Inches(fold_h_y),
        Inches(r_edge - _mx + 0.12), Inches(0.10))
    fold_h.fill.solid(); fold_h.fill.fore_color.rgb = C_PINK; fold_h.line.fill.background()

    tb(slide, "↩ 折り返し（エンディング → 上位AT突入）",
       3.3, fold_h_y + 0.02, 3.5, 0.22,
       font_size=8, color=C_PINK, align=PP_ALIGN.CENTER)

    # 下段ラベル
    tb(slide, "下段ルート（上位AT・爆発ルート）", _mx, _top2_y - 0.26, 5.5, 0.26,
       font_size=9, color=C_GOLD, bold=True)

    boxes_bot = [
        ("エンブリオ\nイブ覚醒\n(上位AT)",    C_FILL_GOLD,    C_GOLD),
        ("エンブリオ\nイブアタック\n(特化)",   C_FILL_GOLD,    C_GOLD),
        ("決戦\n神浜聖女\n(ST特化)",           C_FILL_CYAN,    C_CYAN),
        ("ドッペル\nモード\n(穢れ爆発)",       C_FILL_RED,     C_RED),
        ("マギウス\nバトル",                   C_FILL_PINK,    C_PINK),
        ("AT終了\n→有利区間\nリセット",        C_FILL_GRAY,    C_GRAY),
    ]
    for i, (label, fill, bc) in enumerate(boxes_bot):
        lx = _mx + i * (_bw + _gap)
        bpt = 2.5 if bc == C_RED else 1.5
        rect_b(slide, lx, _top2_y, _bw, _bh2, fill, border_color=bc, border_pt=bpt)
        tb(slide, label, lx + 0.06, _top2_y + 0.10, _bw - 0.12, _bh2 - 0.18,
           font_size=9.0, color=C_DARK_TEXT, align=PP_ALIGN.CENTER)

    # 下段矢印
    for i in range(len(boxes_bot) - 1):
        ax = _mx + i * (_bw + _gap) + _bw + 0.01
        arrow_r(slide, ax, _top2_y + _bh2 / 2 - 0.12, length=_gap - 0.02)

    # キーポイント 3パネル
    n_panels = 3
    gap_p = 0.14
    panel_w = (10.0 - _mx * 2 - gap_p * (n_panels - 1)) / n_panels
    ann_items = [
        ("天井 950pt+α",
         "ポイント蓄積上限で強制ボーナス当選。スルー回数が増えるほど初当たりAT期待値が上昇",
         C_GOLD, C_FILL_GOLD),
        ("ドッペル爆発 3000枚超",
         "穢れ解放で突入。AT終了まで上乗せ倍増。設定6以外でも3000枚超は十分現実的",
         C_RED, C_FILL_RED),
        ("コンプ体験",
         "8ストーリーをATで収集してエンディングへ。明確なゴールがAT継続モチベを維持",
         C_CYAN, C_FILL_CYAN),
    ]
    for j, (ak, av, ac, af) in enumerate(ann_items):
        lx = _mx + j * (panel_w + gap_p)
        rect_b(slide, lx, _ann_y, panel_w, 0.72, af, border_color=ac, border_pt=1)
        tb(slide, ak, lx + 0.09, _ann_y + 0.05, panel_w - 0.18, 0.26,
           font_size=9.5, color=ac, bold=True)
        tb(slide, av, lx + 0.09, _ann_y + 0.30, panel_w - 0.18, 0.38,
           font_size=7.5, color=C_MID_GRAY)

    net_note(slide,
             "ゲームフロー：一撃/flick7/ちょんぼりすた 各解析より構成",
             "通常→CZ→ボーナス→AT→コンプ→上位ATの二段構造が特徴")


# ── main ─────────────────────────────────────────────────────

def main():
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    slide1_title(prs)
    slide2_gameflow(prs)
    prs.save(OUT_PATH)
    print(f"保存完了: {OUT_PATH}")

if __name__ == "__main__":
    main()
