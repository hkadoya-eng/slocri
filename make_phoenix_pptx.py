"""
スマスロ新機種企画提案 「PHOENIX LEGACY」 PowerPointジェネレーター
出力: proposals/phoenix_legacy_v2.pptx
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__), "proposals", "phoenix_legacy_v2.pptx")

# ── カラーパレット（フェニックス：橙・金・赤テーマ）────────────────
C_BG      = RGBColor(0x08, 0x04, 0x00)
C_CARD    = RGBColor(0x14, 0x08, 0x00)
C_GOLD    = RGBColor(0xC9, 0x9A, 0x1E)
C_GOLD2   = RGBColor(0xFF, 0xD7, 0x00)
C_PHOENIX = RGBColor(0xFF, 0x6B, 0x1A)
C_EMBER   = RGBColor(0xFF, 0x30, 0x00)
C_FEATHER = RGBColor(0xFF, 0xB0, 0x40)
C_CRIMSON = RGBColor(0xA8, 0x1C, 0x1C)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_TEAL    = RGBColor(0x20, 0xA0, 0x90)
C_LTBLUE  = RGBColor(0x93, 0xC5, 0xFD)
C_CREAM   = RGBColor(0xF5, 0xE6, 0xC8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LTGRAY  = RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY    = RGBColor(0x88, 0x88, 0x88)
C_YELLOW  = RGBColor(0xFF, 0xE0, 0x60)
C_ORANGE  = RGBColor(0xFF, 0xA0, 0x30)
C_GREEN   = RGBColor(0x40, 0xB0, 0x40)

FONT_H = "游明朝"
FONT_B = "メイリオ"


# ── 背景生成（フェニックスの炎イメージ）──────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (8, 4, 0))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 45):
        draw.line([(i, 0), (0, i)], fill=(25, 12, 0), width=1)
    for y in range(h - 180, h):
        t = (y - (h - 180)) / 180
        r = int(60 * t)
        g = int(20 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, 0))
    for y in range(0, 80):
        t = (80 - y) / 80 * 0.3
        r = int(40 * t)
        draw.line([(0, y), (w, y)], fill=(r, int(r * 0.4), 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー関数 ──────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, Inches(10), Inches(5.625))
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s

def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT_B
    if color:
        run.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp

def arrow_r(slide, x, y, w, color):
    h = Emu(150000)
    shp = slide.shapes.add_shape(13, x, y - h // 2, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def hdr(slide, text):
    rect(slide, Inches(0.15), Inches(0.08), Inches(9.7), Emu(420000),
         RGBColor(0x18, 0x06, 0x00))
    rect(slide, Inches(0.15), Inches(0.08), Emu(60000), Emu(420000), C_PHOENIX)
    tb(slide, Inches(0.4), Inches(0.1), Inches(9.2), Emu(380000),
       text, 12, bold=True, color=C_GOLD, font=FONT_H)

def net_note(slide, text="※ネットより"):
    tb(slide, Inches(8.5), Inches(5.38), Inches(1.4), Emu(200000),
       text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)
    rect(s, Inches(0), Inches(1.85), Inches(10), Emu(6000), C_PHOENIX)
    rect(s, Inches(0), Inches(1.85), Inches(10), Emu(3000), C_GOLD2)
    rect(s, Inches(0), Inches(3.65), Inches(10), Emu(6000), C_PHOENIX)
    rect(s, Inches(0), Inches(3.65), Inches(10), Emu(3000), C_GOLD2)

    tb(s, Inches(0.5), Inches(0.28), Inches(9), Emu(360000),
       "新機種企画提案", 14, color=C_FEATHER, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(1000000),
       "PHOENIX LEGACY", 54, bold=True, color=C_GOLD2, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.3), Inches(1.60), Inches(9.4), Emu(280000),
       "― 受 け 継 が れ る 炎 ―", 17, bold=True, color=C_PHOENIX, font=FONT_H, align=PP_ALIGN.CENTER)

    rect(s, Inches(1.2), Inches(1.98), Inches(7.6), Emu(570000), RGBColor(0x1C, 0x08, 0x00))
    tb(s, Inches(1.3), Inches(2.03), Inches(7.4), Emu(520000),
       "「 炎は消えない。世代を超えて、燃え続ける。 」",
       16, bold=True, color=C_CREAM, font=FONT_H, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.5), Inches(3.80), Inches(4.3), Emu(360000),
       "タイプ：スマスロ（L型）\n純増：4.0枚/G（FLAME）/ 8.0枚/G（ETERNAL FLAME）",
       10, color=C_LTGRAY)
    tb(s, Inches(5.2), Inches(3.80), Inches(4.5), Emu(360000),
       "ターゲット：30〜50代 / RPG・ロールプレイング世代\nMY目標：約3,500枚（ミドル）",
       10, color=C_LTGRAY)

    tb(s, Inches(7.8), Inches(5.2), Inches(2.0), Emu(300000),
       "v2.0  2026.05", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: セールスポイント（旧スライド2+3を統合）
# ══════════════════════════════════════════════════════════════
def s_sellingpoints(prs):
    s = new_slide(prs)
    hdr(s, "3つの設計軸  ──  PHOENIX LEGACYが目指す体験")

    pillars = [
        (Inches(0.2),
         "①  PHOENIX CYCLE",
         "1回のAT = 1世代\nATが終わるたびに\n「REBIRTH演出」が流れ\n次の炎が点火する\n\n「また負けた」ではなく\n「炎を受け継いだ」\nになる感情設計\n\n負けをポジティブに\n転換する構造",
         RGBColor(0x1E, 0x08, 0x00), C_PHOENIX),
        (Inches(3.55),
         "②  IGNITION RITE",
         "AT消化中に赤7揃いで\n「点火方法」を選ぶ\n\nSkill（技術）\nBond（絆）\nPower（力）\nの3択\n\n選んだ内容で\n次ATの性能が変わる\n\nプレイヤーが\n物語を作る台",
         RGBColor(0x14, 0x10, 0x00), C_GOLD),
        (Inches(6.9),
         "③  NEMESIS GAUGE",
         "ATを重ねるほど\n宿敵が弱体化する\n1セッション完結設計\n\nゲージMAX\n＋PHOENIX FLAME\n＋IGNITION完全制覇\nでETERNAL FLAME解放\n\n「今日が決戦」という\n感動を1日で完結",
         RGBColor(0x1E, 0x04, 0x04), C_CRIMSON),
    ]
    for x, title, desc, fill, bdr in pillars:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.55), fill, bdr, 2.0)
        tb(s, x + Emu(100000), Inches(0.92), Inches(2.8), Emu(360000),
           title, 11, bold=True, color=bdr, font=FONT_H)
        tb(s, x + Emu(100000), Inches(1.40), Inches(2.8), Inches(2.80),
           desc, 9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.52), Inches(9.6), Emu(40000), C_GOLD)
    rect(s, Inches(0.2), Inches(4.60), Inches(9.6), Emu(560000), RGBColor(0x16, 0x06, 0x00))
    tb(s, Inches(0.35), Inches(4.65), Inches(9.2), Emu(240000),
       "祟り神の章（感情逆転）→ ミリゴ・吉宗（実機メカニクス）→ PHOENIX LEGACY（継承×炎×選択）",
       9, bold=True, color=C_GOLD)
    tb(s, Inches(0.35), Inches(4.92), Inches(9.2), Emu(230000),
       "タイムリープ型（Re:ゼロ等）とは異なり、「炎の継承」という積み上げ型の感情設計でユーザーの来店動機を強化する。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: ゲームフロー全体図（overflow修正）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図  ──  1セッションで炎が積み重なる")

    # ボックスサイズを縮小して右余白を確保（⊓コネクター用）
    BW  = Inches(2.75)
    BH  = Inches(1.1)
    GAP = Inches(0.22)
    R1Y = Inches(0.60)
    R2Y = Inches(2.18)
    X1  = Inches(0.18)
    X2  = X1 + BW + GAP   # ≈ 3.15"
    X3  = X2 + BW + GAP   # ≈ 6.12"
    # X3 + BW ≈ 8.87" → 右余白 1.13" （⊓コネクター用に確保）

    row1 = [
        (X1, "SLUMBER（スランバー）",
         "フェニックスが眠る時代\n4ステージ / 周期100GでAWAKENING",
         RGBColor(0x10, 0x06, 0x00), C_FEATHER),
        (X2, "AWAKENING（アウェイクニング）",
         "目覚めの試練 / 3Gバトル\nCZ成功でFLAME突入！",
         RGBColor(0x1A, 0x06, 0x00), C_PHOENIX),
        (X3, "FLAME（フレイム）",
         "1世代50G / 純増4.0枚/G\n赤7揃い → IGNITION発動！",
         RGBColor(0x20, 0x0C, 0x00), C_GOLD),
    ]
    for x, title, desc, fill, bdr in row1:
        rect_b(s, x, R1Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R1Y + Emu(55000), BW - Emu(160000), Emu(320000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R1Y + Emu(390000), BW - Emu(160000), BH - Emu(450000),
           desc, 8.5, color=C_CREAM)
    for x_l in [X1, X2]:
        arrow_r(s, x_l + BW + Emu(40000), R1Y + BH // 2, GAP - Emu(80000), C_GOLD)

    # 右端の⊓コネクター（Row1→Row2の折り返し）
    CON_X = X3 + BW + Emu(80000)   # ≈ 8.96"
    CON_R = CON_X + Emu(550000)    # ≈ 9.57"
    LW    = Emu(55000)
    MID_Y = (R1Y + BH + R2Y) // 2  # Row1底 〜 Row2天の中間
    # 縦線（左）
    rect(s, CON_X, R1Y + BH // 2, LW, MID_Y - (R1Y + BH // 2) + Emu(28000), C_GOLD)
    # 水平バー
    rect(s, CON_X, MID_Y, CON_R - CON_X + LW, LW, C_GOLD)
    # 縦線（右）→ Row2のX3右端
    rect(s, CON_R, R2Y, LW, MID_Y - R2Y + LW, C_GOLD)
    # ↓先端矢印っぽいテキスト
    tb(s, CON_X - Emu(60000), MID_Y - Emu(350000), Emu(800000), Emu(330000),
       "↺", 16, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)

    row2 = [
        (X3, "フェニックスループ",
         "炎が重なるほど強化\nストック残 → 次の炎へ！",
         RGBColor(0x20, 0x0C, 0x00), C_GOLD),
        (X2, "BLAZING ALLIANCE",
         "3世代以上で炎が集結\n継続率75% / 純増6.5枚/G",
         RGBColor(0x12, 0x0A, 0x00), C_FEATHER),
        (X1, "EMBER（エンバー）",
         "AT終了後5Gの残り火\n奇数揃いでFLAME再突入！",
         RGBColor(0x18, 0x06, 0x00), C_PHOENIX),
    ]
    for x, title, desc, fill, bdr in row2:
        rect_b(s, x, R2Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R2Y + Emu(55000), BW - Emu(160000), Emu(320000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R2Y + Emu(390000), BW - Emu(160000), BH - Emu(450000),
           desc, 8.5, color=C_CREAM)
    for x_r in [X3, X2]:
        _w = GAP - Emu(80000)
        _h = Emu(150000)
        shp = s.shapes.add_shape(13, x_r - GAP + Emu(40000),
                                  R2Y + BH // 2 - _h // 2, _w, _h)
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = C_PHOENIX
        shp.line.fill.background()

    # 下部: ETERNAL FLAME
    BOT_Y = Inches(3.48)
    rect_b(s, X1, BOT_Y, Inches(2.75), Emu(840000),
           RGBColor(0x22, 0x04, 0x04), C_CRIMSON, 2.0)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(55000), Inches(2.5), Emu(330000),
       "ETERNAL FLAME", 10.5, bold=True, color=C_CRIMSON, font=FONT_H)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(420000), Inches(2.5), Emu(380000),
       "NEMESISゲージMAX\n＋PHOENIX FLAME\n＋IGNITION完全制覇\n純増8.0枚/G", 8, color=C_CREAM)

    rect(s, Inches(3.1), BOT_Y, Inches(6.72), Emu(840000), RGBColor(0x18, 0x06, 0x00))
    tb(s, Inches(3.25), BOT_Y + Emu(80000), Inches(6.3), Emu(330000),
       "★ 前後のFLAME積み上げ込みでMY約3,500枚。",
       9, bold=True, color=C_FEATHER)
    tb(s, Inches(3.25), BOT_Y + Emu(440000), Inches(6.3), Emu(360000),
       "「自分が今日育てた炎が、ここで完全に燃え上がる」\n1セッションの物語が完結する瞬間。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 通常時の仕組み
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の仕組み  ──  SLUMBER × AWAKENING × レガシーポイント")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.1), Inches(2.45),
           C_CARD, C_FEATHER, 1.5)
    tb(s, Inches(0.3), Inches(0.90), Inches(2.9), Emu(320000),
       "① ステージ4種（遊技数連動）", 10, bold=True, color=C_FEATHER, font=FONT_H)
    stages = [
        ("灰の荒野",   "SLUMBER通常",          C_LTGRAY),
        ("燃える草原", "高確示唆",             C_YELLOW),
        ("炎の王都",   "前兆・AWAKENING高確率", C_ORANGE),
        ("天空の炎柱", "FLAME超高確率状態",    C_GOLD2),
    ]
    sy = Inches(1.33)
    for sname, sdesc, scol in stages:
        tb(s, Inches(0.3),  sy, Inches(1.5),  Emu(268000), sname, 9, bold=True, color=scol, wrap=False)
        tb(s, Inches(1.85), sy, Inches(1.35), Emu(268000), sdesc, 9, color=C_CREAM, wrap=False)
        sy += Emu(270000)

    rect_b(s, Inches(3.45), Inches(0.85), Inches(3.1), Inches(2.45),
           C_CARD, C_PHOENIX, 1.5)
    tb(s, Inches(3.55), Inches(0.90), Inches(2.9), Emu(320000),
       "② AWAKENING（アウェイクニング）", 10, bold=True, color=C_PHOENIX, font=FONT_H)
    czs = [
        ("周期",    "100G毎に規則的に到来",      C_CREAM),
        ("バトル",  "3G間の試練演出",            C_LTGRAY),
        ("成功",    "FLAME突入",                 C_GOLD2),
        ("失敗",    "レガシーポイント＋10獲得",  C_FEATHER),
    ]
    cy = Inches(1.33)
    for label, desc, col in czs:
        tb(s, Inches(3.55), cy, Inches(1.1),  Emu(268000), label, 9, bold=True, color=col,   wrap=False)
        tb(s, Inches(4.70), cy, Inches(1.75), Emu(268000), desc,  9, color=C_CREAM, wrap=False)
        cy += Emu(270000)

    rect_b(s, Inches(6.7), Inches(0.85), Inches(3.1), Inches(2.45),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(6.8), Inches(0.90), Inches(2.9), Emu(320000),
       "③ レガシーポイント（Legacy Point）", 10, bold=True, color=C_GOLD, font=FONT_H)
    pts = [
        ("獲得方法",   "AWAKENING失敗・リプレイ等",      C_CREAM),
        ("50pt到達",  "天井100G短縮（600→500G）",      C_YELLOW),
        ("100pt到達", "天井200G短縮（400G）",           C_ORANGE),
        ("200pt到達", "AWAKENING突破率大幅UP",          C_GOLD2),
    ]
    py = Inches(1.33)
    for label, desc, col in pts:
        tb(s, Inches(6.8),  py, Inches(1.3),  Emu(268000), label, 9, bold=True, color=col,   wrap=False)
        tb(s, Inches(8.15), py, Inches(1.55), Emu(268000), desc,  9, color=C_CREAM, wrap=False)
        py += Emu(270000)

    rect(s, Inches(0.2), Inches(3.40), Inches(9.6), Emu(40000), C_PHOENIX)
    rect(s, Inches(0.2), Inches(3.48), Inches(9.6), Emu(620000), RGBColor(0x14, 0x06, 0x00))
    tb(s, Inches(0.35), Inches(3.53), Inches(9.2), Emu(260000),
       "通常時の読み方", 10, bold=True, color=C_PHOENIX, font=FONT_H)
    tb(s, Inches(0.35), Inches(3.86), Inches(9.2), Emu(260000),
       "「天空の炎柱」以上になったら好機。レガシーポイントを積みながらAWAKENINGを迎える。失敗しても「育てている」感覚が大事。",
       9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.18), Inches(9.6), Emu(380000), RGBColor(0x14, 0x06, 0x00))
    tb(s, Inches(0.35), Inches(4.22), Inches(9.2), Emu(330000),
       "★ 周期設計は吉宗の規則的CZ設計を参考。「いつAWAKENINGが来るか分かる」ストレスの少ない通常時。",
       9, color=C_GOLD)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: FLAME・IGNITION
# ══════════════════════════════════════════════════════════════
def s_flame_at(prs):
    s = new_slide(prs)
    hdr(s, "FLAME・IGNITION  ──  選ぶことで炎の色が変わる")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.0), Inches(3.45),
           RGBColor(0x20, 0x0A, 0x00), C_GOLD, 2.0)
    tb(s, Inches(0.3), Inches(0.90), Inches(2.8), Emu(320000),
       "FLAME（フレイム）", 13, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(0.3), Inches(1.28), Inches(2.8), Inches(2.7),
       "1世代50G / 純増4.0枚/G\nストックを1個消費して消化\n\n"
       "【フレイムストック4種】\n"
       "  First Flame：ループ率20%\n"
       "  Second Flame：ループ率40%\n"
       "  Third Flame：ループ率60%\n"
       "  PHOENIX FLAME：80%（レア）\n\n"
       "ミリゴA〜Dストックを\nフェニックスで昇華",
       9, color=C_CREAM)

    arrow_r(s, Inches(3.3), Inches(2.52), Emu(270000), C_PHOENIX)

    rect_b(s, Inches(3.75), Inches(0.85), Inches(3.1), Inches(3.45),
           RGBColor(0x20, 0x10, 0x00), C_PHOENIX, 2.5)
    tb(s, Inches(3.85), Inches(0.90), Inches(2.9), Emu(320000),
       "IGNITION（イグニション）", 12, bold=True, color=C_PHOENIX, font=FONT_H)
    tb(s, Inches(3.85), Inches(1.28), Inches(2.9), Inches(2.7),
       "赤7揃いで発動\n（ミリゴZ-ZONEを昇華）\n\n"
       "【3択の点火方法】\n\n"
       "  Skill（スキル）\n  → 次FLAMEの特殊演出解放\n\n"
       "  Bond（ボンド）\n  → 次FLAME継続率+10%\n\n"
       "  Power（パワー）\n  → NEMESISゲージ大幅UP",
       9, color=C_CREAM)

    arrow_r(s, Inches(6.95), Inches(2.52), Emu(270000), C_GOLD)

    rect_b(s, Inches(7.4), Inches(0.85), Inches(2.4), Inches(3.45),
           RGBColor(0x1E, 0x12, 0x00), C_GOLD2, 2.5)
    tb(s, Inches(7.5), Inches(0.90), Inches(2.2), Emu(320000),
       "炎の積み重ね", 13, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(7.5), Inches(1.28), Inches(2.2), Inches(2.7),
       "1st FLAME\n  ↓ IGNITION\n2nd FLAME\n  ↓ IGNITION\n3rd FLAME\n  ↓\nBLAZING\nALLIANCEへ！\n\n★ 炎が重なるほど\n   演出が豊かに",
       9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(700000),
         RGBColor(0x18, 0x06, 0x00))
    rect(s, Inches(0.2), Inches(4.40), Emu(60000), Emu(700000), C_PHOENIX)
    tb(s, Inches(0.45), Inches(4.45), Inches(9.2), Emu(270000),
       "IGNITIONは「炎に何を注ぐか」を選ぶ唯一の瞬間", 10, bold=True, color=C_PHOENIX, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.78), Inches(9.2), Emu(270000),
       "Skill・Bond・Powerの選択が次のFLAMEの質を変える。打ち手ごとに異なる炎の物語が展開する。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: BLAZING ALLIANCE・ETERNAL FLAME
# ══════════════════════════════════════════════════════════════
def s_climax(prs):
    s = new_slide(prs)
    hdr(s, "BLAZING ALLIANCE・ETERNAL FLAME  ──  炎が頂点に達する")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.5), Inches(3.45),
           RGBColor(0x1A, 0x0A, 0x00), C_FEATHER, 2.0)
    tb(s, Inches(0.3), Inches(0.90), Inches(4.3), Emu(320000),
       "BLAZING ALLIANCE", 13, bold=True, color=C_FEATHER, font=FONT_H)
    tb(s, Inches(0.3), Inches(1.28), Inches(4.3), Inches(2.7),
       "発動条件：FLAME 3回以上ループ\n　　　　　＋ 特定IGNITION組み合わせ\n\n"
       "複数世代のフェニックスが「連合」する演出\n\n"
       "継続率：75%\n純増：6.5枚/G\n期待枚数：1,000〜1,500枚\n\n"
       "ALLIANCE終了時「EMBER（残り火）」5G\n  → 吉宗1G連の感動をここで再現",
       9, color=C_CREAM)

    rect(s, Inches(4.85), Inches(1.6), Emu(220000), Inches(2.0),
         RGBColor(0x20, 0x08, 0x00))
    tb(s, Inches(4.85), Inches(2.35), Emu(220000), Emu(400000),
       "→", 22, bold=True, color=C_EMBER, align=PP_ALIGN.CENTER)

    rect_b(s, Inches(5.25), Inches(0.85), Inches(4.55), Inches(3.45),
           RGBColor(0x24, 0x04, 0x04), C_EMBER, 2.5)
    tb(s, Inches(5.35), Inches(0.90), Inches(4.35), Emu(320000),
       "ETERNAL FLAME", 15, bold=True, color=C_EMBER, font=FONT_H)
    tb(s, Inches(5.35), Inches(1.28), Inches(4.35), Inches(2.7),
       "発動条件：\n"
       "  NEMESISゲージ MAX\n"
       "  ＋ PHOENIX FLAMEストック\n"
       "  ＋ IGNITION 3種コンプリート\n\n"
       "純増：8.0枚/G\n歴代全フェニックスが集結する演出\n\n"
       "★ 前後のFLAME積み上げ込みで\n   MY 約3,500枚（ミドル）",
       9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(700000),
         RGBColor(0x1E, 0x06, 0x00))
    rect(s, Inches(0.2), Inches(4.40), Emu(60000), Emu(700000), C_EMBER)
    tb(s, Inches(0.45), Inches(4.45), Inches(9.2), Emu(270000),
       "MY設計：FLAME(800枚) ＋ BLAZING ALLIANCE(1,200枚) ＋ ETERNAL FLAME(1,500枚) ＝ 約3,500枚",
       10, bold=True, color=C_FEATHER, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.78), Inches(9.2), Emu(270000),
       "通常遊技でもBLAZING ALLIANCEまで楽しめる設計。ETERNAL FLAMEは「積み上げれば届く」距離感。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 基本スペック
# ══════════════════════════════════════════════════════════════
def s_spec(prs):
    s = new_slide(prs)
    hdr(s, "基本スペック  ──  ミドルMY・ストレスフリー設計")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.7), Inches(3.55),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.3), Inches(0.90), Inches(4.5), Emu(320000),
       "基本スペック", 11, bold=True, color=C_GOLD, font=FONT_H)
    specs = [
        ("タイプ",      "スマスロ（L型）"),
        ("天井",        "600G（Legacy Pointで最短400G）"),
        ("純増",        "4.0枚/G（FLAME）/ 6.5枚/G（ALLIANCE）/ 8.0枚/G（ETERNAL）"),
        ("機械割",      "設定1：97.5%  /  設定6：109%"),
        ("MY目標",      "約3,500枚（ミドル）"),
        ("コイン単価",  "20円（3枚ベット / 1G＝60円）"),
        ("CV",          "0.22〜0.27"),
    ]
    sy = Inches(1.36)
    for label, val in specs:
        rect(s, Inches(0.3), sy, Inches(1.6), Emu(255000), RGBColor(0x20, 0x0C, 0x00))
        tb(s, Inches(0.35), sy + Emu(28000), Inches(1.5), Emu(200000),
           label, 8.5, bold=True, color=C_GOLD2, wrap=False)
        tb(s, Inches(1.95), sy + Emu(28000), Inches(2.85), Emu(200000),
           val, 8.5, color=C_CREAM, wrap=False)
        sy += Emu(265000)

    rect_b(s, Inches(5.1), Inches(0.85), Inches(4.7), Inches(3.55),
           C_CARD, C_PHOENIX, 1.5)
    tb(s, Inches(5.2), Inches(0.90), Inches(4.5), Emu(320000),
       "出玉期待値", 11, bold=True, color=C_PHOENIX, font=FONT_H)
    outs = [
        ("FLAME",           "300〜400枚",    "1〜2世代の日"),
        ("BLAZING ALLIANCE","1,000〜1,500枚","3世代ループの日"),
        ("ETERNAL FLAME",   "1,500枚以上",   "全条件成立の日"),
        ("MY（目安）",       "約3,500枚",    "Eternal+前後込みの1日"),
    ]
    oy = Inches(1.36)
    for label, amount, note in outs:
        rect(s, Inches(5.2), oy, Inches(1.6), Emu(280000), RGBColor(0x24, 0x08, 0x00))
        tb(s, Inches(5.25), oy + Emu(38000), Inches(1.5), Emu(200000),
           label, 8, bold=True, color=C_FEATHER, wrap=False)
        tb(s, Inches(6.85), oy + Emu(38000), Inches(1.45), Emu(200000),
           amount, 9, bold=True, color=C_GOLD2, wrap=False)
        tb(s, Inches(8.35), oy + Emu(38000), Inches(1.35), Emu(200000),
           note, 7.5, color=C_LTGRAY, wrap=False)
        oy += Emu(295000)

    rect(s, Inches(0.2), Inches(4.50), Inches(9.6), Emu(560000),
         RGBColor(0x14, 0x06, 0x00))
    rect(s, Inches(0.2), Inches(4.50), Emu(60000), Emu(560000), C_GOLD)
    tb(s, Inches(0.45), Inches(4.55), Inches(9.2), Emu(250000),
       "時間コスト目安：1G=60円 × 設定1(97.5%) → 期待損失1.5円/G。400G/時で約600円/時。",
       8.5, bold=True, color=C_GOLD)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.2), Emu(250000),
       "BLAZING ALLIANCE狙いなら3FLAME（約150G）が一区切り。ミドル層が3,000〜5,000円で十分楽しめる時間設計。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: ターゲット・市場考察
# ══════════════════════════════════════════════════════════════
def s_market(prs):
    s = new_slide(prs)
    hdr(s, "ターゲット・市場考察  ──  RPG世代×フェニックスの必然性")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.65), Inches(3.45),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.3), Inches(0.90), Inches(4.45), Emu(320000),
       "メインターゲット", 11, bold=True, color=C_GOLD, font=FONT_H)
    targets = [
        ("年齢層",         "30〜50代男性",                         C_GOLD2),
        ("世代",           "ロマサガ2・FF・DQ体験世代",            C_CREAM),
        ("プレイスタイル", "週1〜2来店のリピーター",               C_CREAM),
        ("感情ニーズ",     "「物語に没入」「積み上げを感じたい」", C_FEATHER),
        ("許容投資",       "3,000〜8,000円/日",                    C_YELLOW),
    ]
    ty = Inches(1.35)
    for label, val, col in targets:
        tb(s, Inches(0.3),  ty, Inches(1.55), Emu(268000), label, 8.5, bold=True, color=C_GRAY,  wrap=False)
        tb(s, Inches(1.90), ty, Inches(2.85), Emu(268000), val,   8.5, color=col, wrap=False)
        ty += Emu(270000)

    rect_b(s, Inches(5.05), Inches(0.85), Inches(4.75), Inches(3.45),
           C_CARD, C_PHOENIX, 1.5)
    tb(s, Inches(5.15), Inches(0.90), Inches(4.55), Emu(320000),
       "競合との差別化", 11, bold=True, color=C_PHOENIX, font=FONT_H)
    comps = [
        ("Re:ゼロ",    "タイムリープ・失敗前提",    "炎の継承＝積み上げ（前向き）"),
        ("ミリゴ",     "PGG 1/16,384の夢",          "ETERNAL FLAME＝積み上げれば届く"),
        ("吉宗",       "1G連の瞬間爆発",            "フェニックスループで連続興奮"),
        ("祟り神の章", "感情逆転・他者理解",         "炎の積み上げ・自己の物語"),
    ]
    cy = Inches(1.35)
    for title, their, ours in comps:
        tb(s, Inches(5.15), cy, Inches(1.25), Emu(265000), title, 8, bold=True, color=C_LTGRAY, wrap=False)
        tb(s, Inches(6.45), cy, Inches(1.55), Emu(265000), their, 7.5, color=C_LTGRAY, wrap=False)
        tb(s, Inches(8.05), cy, Inches(1.65), Emu(265000), f"→ {ours}", 7.5, color=C_FEATHER, wrap=False)
        cy += Emu(268000)

    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(700000),
         RGBColor(0x14, 0x06, 0x00))
    rect(s, Inches(0.2), Inches(4.40), Emu(60000), Emu(700000), C_PHOENIX)
    tb(s, Inches(0.45), Inches(4.45), Inches(9.2), Emu(260000),
       "「なぜ今フェニックスなのか」", 10, bold=True, color=C_PHOENIX, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.78), Inches(9.2), Emu(280000),
       "RPG世代（FF・ロマサガ）にはPhoenixが刺さる。「炎は消えない・再誕」というモチーフは継承テーマと完全一致。"
       "タイムリープ型との差別化軸として有効と考える。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ  ──  PHOENIX LEGACYが生み出す体験")

    cols_data = [
        (Inches(0.2),
         "1セッションで感じる",
         "「今日の遊技が\n1つの炎の物語」\n\nFLAMEが重なるたびに\n演出が豊かになり\nAT終了時も\n「炎を受け継いだ」\nというポジティブな\n余韻が残る",
         C_PHOENIX, RGBColor(0x1E, 0x08, 0x00)),
        (Inches(3.55),
         "選ぶことで感じる",
         "「この台は\n自分の炎だ」\n\nIGNITIONで選んだ\n点火方法が\n炎の色を変える\n\n同じ台でも\n2人と同じ炎がない\n設計",
         C_GOLD, RGBColor(0x1E, 0x10, 0x00)),
        (Inches(6.9),
         "決着で感じる",
         "「積み上げた\n炎が報われた」\n\nETERNAL FLAMEは\n「たまたま当たった」\nではなく\n「自分が育てた」\n感覚\n\nMY3,500枚の\n完全燃焼",
         C_EMBER, RGBColor(0x22, 0x04, 0x00)),
    ]
    for x, title, desc, col, fill in cols_data:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.3), fill, col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(330000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.38), Inches(2.8), Inches(2.5),
           desc, 9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.25), Inches(9.6), Emu(40000), C_GOLD)
    rect(s, Inches(0.2), Inches(4.33), Inches(9.6), Emu(760000), RGBColor(0x1C, 0x08, 0x00))
    tb(s, Inches(0.35), Inches(4.38), Inches(9.2), Emu(270000),
       "3つの設計軸：① PHOENIX CYCLE（炎の継承）  ② IGNITION RITE（AT中の3択）  ③ NEMESIS GAUGE（1セッション完結）",
       9, bold=True, color=C_GOLD2)
    tb(s, Inches(0.35), Inches(4.70), Inches(9.2), Emu(350000),
       "ミリゴZ-ZONE × 吉宗1G連 × 祟り神の感情設計を統合。\n"
       "「炎は消えない。世代を超えて、燃え続ける。」——PHOENIX LEGACYの核心です。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = [
        ("タイトル",                  s_title),
        ("3つの設計軸",               s_sellingpoints),
        ("ゲームフロー全体図",         s_flow),
        ("通常時の仕組み",             s_normal),
        ("FLAME・IGNITION",           s_flame_at),
        ("BLAZING ALLIANCE・ETERNAL", s_climax),
        ("基本スペック",               s_spec),
        ("ターゲット・市場考察",       s_market),
        ("まとめ",                     s_matome),
    ]

    print("=" * 55)
    print("  PHOENIX LEGACY v2 企画提案書ジェネレーター")
    print("=" * 55)
    print("\n🔥 スライド生成中...")
    for i, (name, func) in enumerate(slides, 1):
        print(f"   {i:2d}/{len(slides)} {name}")
        func(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n✅ 保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
