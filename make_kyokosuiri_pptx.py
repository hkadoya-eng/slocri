"""
L虚構推理 機種説明＋分析 統合資料 v2  （ディライト（D-LIGHT）・2026年4月6日導入）
出力: proposals/機種分析/虚構推理/kyokosuiri_guide_v2.pptx
テーマ: 深紺 × ミステリー紫(C_PUR) × シアン(C_CYAN) × 赤(C_RED)
構成: Part A（プレイヤー視点説明）6枚 + Part B（分析）3枚 = 計9枚
v2変更: ARROW告知色を正確な当選契機告知に修正、虚構連モードループ率を実数値に更新
        エピソード成功期待度（EP1:25%〜EP5:75%）を明記
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "虚構推理", "kyokosuiri_guide_v2.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深紺×ミステリー紫×シアン×赤）──────────────────────────────
C_BG    = RGBColor(0x04, 0x04, 0x18)   # 深紺
C_CARD  = RGBColor(0x08, 0x08, 0x24)
C_CARD2 = RGBColor(0x10, 0x10, 0x2C)
C_ROW   = RGBColor(0x0C, 0x0C, 0x28)
C_PUR   = RGBColor(0x66, 0x22, 0xAA)   # ミステリー紫
C_PUR2  = RGBColor(0x99, 0x44, 0xFF)   # 明るい紫
C_CYAN  = RGBColor(0x22, 0xAA, 0xCC)   # シアン（ARROW告知色）
C_CYAN2 = RGBColor(0x44, 0xDD, 0xFF)
C_RED   = RGBColor(0xCC, 0x11, 0x11)   # 赤（確定告知）
C_CRIM  = RGBColor(0xFF, 0x33, 0x33)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_BLUE  = RGBColor(0x22, 0x66, 0xCC)
C_YEL   = RGBColor(0xDD, 0xCC, 0x22)
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景PNG生成（深紺＋謎の斜めライン）──────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 4, 24))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(8, 8, 32), width=1)
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        draw.line([(0, y), (w, y)], fill=(int(10 * t), 0, int(20 * t)))
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


# ── テキストボックス ─────────────────────────────────────────────────────────
def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


# ── 矩形ヘルパー ─────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


# ── ヘッダー（スライド上部バー）─────────────────────────────────────────────
def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_PUR)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_CYAN, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_PUR)


# ── ネット解析注記 ────────────────────────────────────────────────────────────
def net_note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ── 右矢印 ────────────────────────────────────────────────────────────────────
def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


# ── フッター（設計コメント＋補足説明）──────────────────────────────────────
def footer(slide, bold_text, sub_text=""):
    """各スライド下部に設計コメント（太字）＋補足説明を入れる"""
    fy = Inches(5.10)
    fh = Emu(380000)
    rect(slide, 0, fy, SLIDE_W, fh, RGBColor(0x06, 0x04, 0x1C))
    rect(slide, 0, fy, Emu(20000), fh, C_PUR)
    tb(slide, Inches(0.22), fy + Emu(40000), Inches(6.0), Emu(160000),
       bold_text, 7.5, bold=True, color=C_PUR2)
    if sub_text:
        tb(slide, Inches(0.22), fy + Emu(200000), Inches(9.2), Emu(160000),
           sub_text, 7, color=C_GRAY)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント  [Part A - 1/6]
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル背景
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x02, 0x02, 0x10))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_PUR)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, RGBColor(0x44, 0x11, 0x77))

    tb(s, Inches(0.22), Inches(0.42), Inches(5.0), Emu(290000),
       "機種説明＋分析資料　Part A: プレイヤー視点", 10, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.88), Inches(5.1), Emu(700000),
       "L虚構推理", 38, bold=True, color=C_PUR2, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.65), Inches(5.1), Emu(290000),
       "── エピソード突破型CZ × ARROW告知 × 虚構連モード", 10, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.20), Inches(4.9), Emu(210000),
       "メーカー：ディライト（D-LIGHT）　　導入：2026年4月6日", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.46), Inches(4.9), Emu(210000),
       "設定：1〜6段階　　天井：虚構真偽間1000G（リセット後700G）", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.72), Inches(4.9), Emu(210000),
       "CZ：鋼人七瀬攻略議会（6G・5エピソード突破型）", 8.5, color=C_GRAY)

    # 右：この台の3ポイント
    kws = [
        (C_PUR,   "① エピソード突破CZ",
         "失敗しても次回CZに持越し（キャリーオーバー）\n「諦めなければ必ず成功する」積み上げ設計"),
        (C_CYAN,  "② ARROW告知デバイス",
         "筐体上部の専用デバイスが当選契機を色で告知\n白=一撃成功・赤点滅=強チェリー確定の興奮体験"),
        (C_GREEN, "③ 虚構連モード",
         "Short(56%) / Middle(81%) / Long(95%) 3種\n六花ステージ=Long到達で9400枚超の爆発力"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.28) + i * Emu(1530000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Emu(1430000), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Emu(1430000), ac)
        tb(s, Inches(5.85), y0 + Emu(60000), Inches(3.8), Emu(300000),
           kw, 11, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(360000), Inches(3.8), Emu(900000),
           desc, 8.5, color=C_WHITE)

    footer(s,
           "設計コメント：エピソードキャリーオーバー＋ARROW告知の組み合わせが本機の最大の個性",
           "補足：原作「虚構推理」（城平京・白浜鴎）のミステリー世界観をゲーム性に落とし込んだ意欲作")

    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図  [Part A - 2/6]
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→CZ→ボーナス→虚構連モードへの全ルート", "2/9")

    # ── 上段フロー（4ボックス）──────────────────────────────────────────
    tb(s, Inches(0.28), Inches(0.66), Inches(4.0), Emu(250000),
       "▶ 通常ルート（基本フロー）", 8.5, bold=True, color=C_CYAN)

    boxes_top = [
        (C_LTGRY, C_WHITE,  "通常時",
         "エピソード進行\n規定G数でCZ突入\n天井:1000G"),
        (C_PUR,   C_PUR2,   "CZ\n攻略議会",
         "6G / 初回7G\n5エピソード突破\nARROW告知"),
        (C_CARD2, C_CYAN,   "ボーナス",
         "各種ボーナス\n直撃当選もあり"),
        (RGBColor(0x04, 0x14, 0x0C), C_GREEN, "虚構連モード",
         "Short/Middle/Long\n高確率ループ"),
    ]
    bw, bh = Inches(1.85), Emu(1550000)
    gap = Inches(0.23)
    total = 4 * bw + 3 * gap
    sx = (SLIDE_W - total) / 2
    ty = Inches(0.90)
    cy_top = ty + bh // 2

    for i, (fill, bc, lbl, sub) in enumerate(boxes_top):
        bx0 = sx + i * (bw + gap)
        rect_b(s, bx0, ty, bw, bh, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), ty + Emu(80000),
           bw - Emu(80000), Emu(400000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), ty + Emu(490000),
           bw - Emu(60000), Emu(920000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw + Emu(10000), cy_top, col=C_PUR)

    # ── 下段フロー（サブルート説明）────────────────────────────────────────
    tb(s, Inches(0.28), Inches(2.90), Inches(4.0), Emu(250000),
       "▶ CZキャリーオーバー（積み上げ設計）", 8.5, bold=True, color=C_PUR2)

    by2 = Inches(3.12)
    bh2 = Emu(1580000)

    # キャリーオーバー説明ブロック（左2/3）
    co_bw = Inches(6.0)
    rect_b(s, Inches(0.28), by2, co_bw, bh2, C_CARD, C_PUR, 1.5)
    rect(s, Inches(0.28), by2, Emu(45000), bh2, C_PUR)

    # キャリーオーバーの視覚フロー
    co_steps = [
        ("1回目CZ", "3エピソード突破\n2失敗→持越し", C_PUR2),
        ("2回目CZ", "2エピソード突破\n0失敗→成功!", C_GREEN),
    ]
    csw = Inches(2.2)
    csx = Inches(0.5)
    for j, (ct, cs, cc) in enumerate(co_steps):
        cx0 = csx + j * (csw + Inches(0.45))
        rect_b(s, cx0, by2 + Emu(120000), csw, bh2 - Emu(240000),
               RGBColor(0x0C, 0x08, 0x24), cc, 1.2)
        tb(s, cx0 + Emu(50000), by2 + Emu(180000), csw - Emu(80000), Emu(290000),
           ct, 9, bold=True, color=cc, align=PP_ALIGN.CENTER)
        tb(s, cx0 + Emu(40000), by2 + Emu(470000), csw - Emu(70000), Emu(780000),
           cs, 8, color=C_WHITE, align=PP_ALIGN.CENTER)
        if j == 0:
            arrow_r(s, cx0 + csw + Emu(30000), by2 + bh2 // 2, col=C_PUR2)

    tb(s, Inches(0.32), by2 + Emu(60000), co_bw - Emu(80000), Emu(250000),
       "失敗エピソードは次回CZへ持越し ── 「積み上げ」が成功への道に変わる", 8, bold=True,
       color=C_PUR2)

    # 右：虚構連モード3種説明
    mx = Inches(6.45)
    mw = Inches(3.3)
    rect_b(s, mx, by2, mw, bh2, C_CARD, C_GREEN, 1.5)
    rect(s, mx, by2, Emu(45000), bh2, C_GREEN)
    tb(s, mx + Emu(70000), by2 + Emu(60000), mw - Emu(90000), Emu(260000),
       "虚構連モード3種", 9.5, bold=True, color=C_GREEN, font=FONT_H)
    modes3 = [
        ("Short",  "ループ率約56%\n短期集中型"),
        ("Middle", "ループ率約81%\n中期継続型"),
        ("Long",   "ループ率約95%\n爆発力最大"),
    ]
    mentry_h = (bh2 - Emu(320000)) // 3
    for k, (mn, md) in enumerate(modes3):
        mey = by2 + Emu(320000) + k * mentry_h
        ac2 = [C_CYAN, C_PUR2, C_GOLD][k]
        rect(s, mx + Emu(50000), mey + Emu(30000),
             Emu(30000), mentry_h - Emu(60000), ac2)
        tb(s, mx + Emu(110000), mey + Emu(60000),
           mw - Emu(140000), Emu(250000), mn, 9, bold=True, color=ac2)
        tb(s, mx + Emu(110000), mey + Emu(310000),
           mw - Emu(140000), Emu(280000), md, 7.5, color=C_GRAY)

    footer(s,
           "設計コメント：失敗が必ず次回CZに繋がるキャリーオーバー設計がこの台最大の差別化要素",
           "補足：蛇行フロー上段=基本ルート、下段=失敗時の積み上げルート。どちらもボーナス・虚構連に繋がる")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方  [Part A - 3/6]
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 規定G数・虚構真偽ゾーン・エピソード進行", "3/9")

    # 3カラム構成
    cols = [
        (C_PUR,  "規定G数管理",
         "通常時はモードによって\n天井G数が変わる。\n\n"
         "【天井】\n虚構真偽間 1000G\nリセット後は 700G\n\n"
         "規定G数到達でCZ\n「鋼人七瀬攻略議会」へ\n強制突入する。\n\n"
         "モード移行によって\nCZ突入G数は短縮される。",
         "G数管理が\n投資の目安"),
        (C_CYAN, "虚構真偽ゾーン",
         "高確率でCZが\n連続発生するゾーン。\n\n"
         "CZ入口となる特殊状態で\n一度入れば怒濤の\nCZ連続発生が期待できる。\n\n"
         "設定が高いほど\n虚構真偽ゾーン突入率も\n高くなる（設定差あり）。\n\n"
         "天井到達でも\n虚構真偽ゾーン経由が\n多い。",
         "高確率CZ発生\nゾーン"),
        (C_GOLD, "エピソード進行",
         "通常時は背景で\nアニメのエピソードが\n進行していく。\n\n"
         "5エピソード構成で\n進行度合いがCZの\n積み上げ状況を示す。\n\n"
         "「何話まで進んでいるか」\nがCZ攻略進捗の\n目安になる。\n\n"
         "エピソード完結で\n特別演出が発生する\nこともある。",
         "物語進行が\n攻略の指標"),
    ]
    col_w = Inches(2.90)
    col_gap = Inches(0.20)
    col_y = Inches(0.72)
    col_h = Emu(3800000)

    for i, (ac, ch, cb, badge) in enumerate(cols):
        cx0 = Inches(0.28) + i * (col_w + col_gap)
        rect_b(s, cx0, col_y, col_w, col_h, C_CARD, ac, 1.8)
        rect(s, cx0, col_y, Emu(45000), col_h, ac)
        # 見出し
        rect(s, cx0 + Emu(45000), col_y, col_w - Emu(45000), Emu(350000),
             RGBColor(0x0C, 0x08, 0x28))
        tb(s, cx0 + Emu(75000), col_y + Emu(60000),
           col_w - Emu(100000), Emu(270000),
           ch, 11, bold=True, color=ac, font=FONT_H)
        # 本文
        tb(s, cx0 + Emu(75000), col_y + Emu(380000),
           col_w - Emu(100000), col_h - Emu(620000),
           cb, 8, color=C_WHITE)
        # バッジ
        rect_b(s, cx0 + col_w - Emu(820000), col_y + Emu(60000),
               Emu(790000), Emu(230000), C_CARD2, ac, 1.0)
        tb(s, cx0 + col_w - Emu(810000), col_y + Emu(70000),
           Emu(780000), Emu(210000), badge, 7, bold=True,
           color=ac, align=PP_ALIGN.CENTER)

    # 右端：ポイントまとめ
    rx = Inches(0.28) + 3 * (col_w + col_gap)
    rw = SLIDE_W - rx - Emu(200000)
    rect_b(s, rx, col_y, rw, col_h, RGBColor(0x06, 0x04, 0x1C), C_PUR2, 1.5)
    tb(s, rx + Emu(60000), col_y + Emu(60000), rw - Emu(80000), Emu(260000),
       "攻略TIP", 10, bold=True, color=C_PUR2, font=FONT_H)
    tips = [
        (C_PUR2,  "リセット後は700G\n天井が短縮される"),
        (C_CYAN,  "虚構真偽ゾーン中は\nCZが連続しやすい"),
        (C_GOLD,  "エピソード進行が\n進んでいるほど\nCZ積み上げ済み"),
    ]
    tip_h = (col_h - Emu(320000)) // 3
    for j, (tc, tt) in enumerate(tips):
        ty0 = col_y + Emu(320000) + j * tip_h
        rect(s, rx + Emu(30000), ty0 + Emu(30000), Emu(15000),
             tip_h - Emu(60000), tc)
        tb(s, rx + Emu(80000), ty0 + Emu(60000),
           rw - Emu(110000), tip_h - Emu(80000), tt, 7.5, color=C_WHITE)

    footer(s,
           "設計コメント：通常時の3要素（G数/ゾーン/エピソード）が全てCZ攻略と連動した一貫設計",
           "補足：リセット狙い（700G天井）や高確狙い（虚構真偽ゾーン中）が有効な立ち回りになる")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ「鋼人七瀬攻略議会」の攻略  [Part A - 4/6]
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ「鋼人七瀬攻略議会」の攻略 ── エピソード突破・キャリーオーバー・ARROW告知", "4/9")

    # ── 左半分: エピソード突破＋キャリーオーバー ────────────────────────────
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    # CZ基本情報バー
    rect(s, lx, ly, lw, Emu(330000), RGBColor(0x44, 0x11, 0x77))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(240000),
       "CZ基本仕様: 6G消化（初回7G）・5エピソード突破型", 9, bold=True, color=C_PUR2)

    # エピソード5段ゲージ
    ep_y = ly + Emu(380000)
    ep_h = Emu(320000)
    ep_gap = Emu(20000)
    ep_colors = [C_PUR, C_PUR, C_PUR, C_GREEN, C_GREEN]  # 3失敗→2成功の例
    ep_texts = ["EP1\n突破", "EP2\n突破", "EP3\n失敗\n→持越", "EP4\n突破", "EP5\n突破\n成功!"]
    ep_states = [True, True, False, True, True]  # False=持越し
    epw = (lw - Emu(4 * 40000)) / 5
    for k in range(5):
        ekx = lx + k * (epw + Emu(40000))
        ec = C_GREEN if ep_states[k] else C_PUR
        border_c = C_GREEN if ep_states[k] else C_CRIM
        rect_b(s, ekx, ep_y, epw, ep_h, C_CARD, border_c, 1.5)
        rect(s, ekx, ep_y, Emu(15000), ep_h, ec)
        tb(s, ekx + Emu(30000), ep_y + Emu(30000), epw - Emu(40000), ep_h - Emu(50000),
           ep_texts[k], 7, bold=(k == 2), color=C_CRIM if not ep_states[k] else C_GREEN,
           align=PP_ALIGN.CENTER)

    # 凡例
    tb(s, lx, ep_y + ep_h + Emu(40000), lw, Emu(170000),
       "▲ EP3失敗→次回CZに持越し（キャリーオーバー）　各EP突破率≒75%　EP1開始時の全突破期待度≒25%", 7,
       color=C_PUR2, bold=True)

    # キャリーオーバー詳細
    ko_y = ep_y + ep_h + Emu(280000)
    ko_h = Emu(1750000)
    rect_b(s, lx, ko_y, lw, ko_h, C_CARD, C_PUR, 1.5)
    rect(s, lx, ko_y, Emu(45000), ko_h, C_PUR)
    tb(s, lx + Emu(75000), ko_y + Emu(50000), lw - Emu(100000), Emu(260000),
       "キャリーオーバー設計", 10, bold=True, color=C_PUR, font=FONT_H)
    tb(s, lx + Emu(75000), ko_y + Emu(310000), lw - Emu(100000), ko_h - Emu(370000),
       "・各EP突破率≒75%　一撃全突破（EP1から）の期待度≒25%\n"
       "・EP2開始:約33% / EP3開始:約44% / EP4開始:約57% / EP5:約75%\n"
       "・失敗エピソードは消えず「次回CZ」にそのまま持越される\n"
       "・「今回EP3まで突破→次回はEP4・5の2個のみ」\n"
       "・失敗が「無」ではなく「前進」になる唯一無二の設計",
       8, color=C_WHITE)

    # ── 右半分: ARROW告知 ───────────────────────────────────────────────────
    rx, ry = Inches(5.05), Inches(0.72)
    rw = Inches(4.65)

    # ARROW告知基本説明
    rect_b(s, rx, ry, rw, Emu(680000), C_CARD, C_CYAN, 1.5)
    rect(s, rx, ry, Emu(45000), Emu(680000), C_CYAN)
    tb(s, rx + Emu(75000), ry + Emu(50000), rw - Emu(100000), Emu(250000),
       "ARROW告知デバイス", 10, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, rx + Emu(75000), ry + Emu(305000), rw - Emu(100000), Emu(340000),
       "筐体上部搭載の専用デバイス。CZ最終G（JUDGEMENT）でランプ色が変化し\n当選契機を告知。白=一撃成功・赤点滅=強チェリーが最高格。",
       8, color=C_WHITE)

    # ARROW色変化カラーボックス（当選契機告知：JUDGEMENTゲームで光る色が当選契機を示す）
    arrow_y = ry + Emu(740000)
    tb(s, rx, arrow_y - Emu(60000), rw, Emu(200000),
       "ARROWランプ色 = 当選契機告知（色が強いほど高設定期待度UP）", 7.5, bold=True, color=C_CYAN)

    arrow_colors = [
        (RGBColor(0xE0, 0xE0, 0xE0), "白",    "一撃成功\n白7BB期待"),
        (C_BLUE,                      "青",    "リプレイ\n当選"),
        (C_YEL,                       "黄",    "ベル\n当選"),
        (C_GREEN,                     "緑",    "スイカ\n当選"),
        (C_RED,                       "赤",    "弱チェリー\n当選"),
        (C_PUR2,                      "紫点滅", "チャンス目\n当選"),
        (C_CRIM,                      "赤点滅", "強チェリー\n当選・激アツ"),
    ]
    ab_w = (rw - Emu(7 * 20000)) / 7
    for m, (ac, albl, aexp) in enumerate(arrow_colors):
        ax0 = rx + m * (ab_w + Emu(20000))
        ah = Emu(800000)
        rect_b(s, ax0, arrow_y, ab_w, ah, C_CARD, ac, 1.5)
        rect(s, ax0, arrow_y, ab_w, Emu(180000), ac)
        tb(s, ax0, arrow_y + Emu(200000), ab_w, Emu(230000),
           albl, 7.5, bold=True, color=ac, align=PP_ALIGN.CENTER)
        tb(s, ax0, arrow_y + Emu(430000), ab_w, Emu(320000),
           aexp, 6.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 告知が生む体験
    exp_y = arrow_y + Emu(860000)
    exp_h = Emu(1300000)
    rect_b(s, rx, exp_y, rw, exp_h, C_CARD, C_CYAN2, 1.5)
    rect(s, rx, exp_y, Emu(45000), exp_h, C_CYAN2)
    tb(s, rx + Emu(75000), exp_y + Emu(50000), rw - Emu(100000), Emu(260000),
       "ARROWが生む「告知の瞬間」", 10, bold=True, color=C_CYAN2, font=FONT_H)
    tb(s, rx + Emu(75000), exp_y + Emu(310000), rw - Emu(100000), exp_h - Emu(370000),
       "・JUDGEMENTゲームに全員の視線がARROWに集まる「1点集中の緊張感」\n"
       "・「白が光れば一撃成功」という明快な判定基準\n"
       "・赤点滅=強チェリー当選確定という直感的な色設計\n"
       "・画面外デバイスが光る「非日常感」が体験報告として続出",
       8, color=C_WHITE)

    footer(s,
           "設計コメント：CZはエピソード突破型で失敗が積み上がり、ARROWのJUDGEMENT告知が最終Gに緊張の頂点を作る二重構造",
           "補足：ARROWランプ色は当選契機告知。白=一撃成功（白7BB期待）、赤点滅=強チェリー確定が最高格。")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: ボーナス後の遊び方  [Part A - 5/6]
# ══════════════════════════════════════════════════════════════
def s_bonus(prs):
    s = new_slide(prs)
    hdr(s, "ボーナス後の遊び方 ── 何をすれば虚構連モードへ行くのか", "5/9")

    # ── 上段: ボーナス種別 ──────────────────────────────────────────────────
    tb(s, Inches(0.28), Inches(0.68), Inches(9.4), Emu(220000),
       "▶ ボーナス種別と虚構連モード突入の仕組み", 9, bold=True, color=C_CYAN)

    bonus_types = [
        (C_PUR2,  "スペシャル\nボーナス",
         "虚構連モード\n確定突入\n最上位ボーナス"),
        (C_CYAN,  "レギュラー\nボーナス",
         "虚構連モード\n突入抽選あり\n基本ボーナス"),
        (C_GOLD,  "設定変更後\n特典",
         "CZ初成功時\n約50%で高確\nスタート"),
        (C_GREEN, "高設定時\n直撃",
         "CZを経由せず\nボーナス当選\n設定差大"),
    ]
    btw = Inches(2.0)
    btgap = Inches(0.32)
    bty = Inches(0.92)
    bth = Emu(1350000)
    btsx = (SLIDE_W - 4 * btw - 3 * btgap) / 2

    for n, (bc, btt, bts) in enumerate(bonus_types):
        btx0 = btsx + n * (btw + btgap)
        rect_b(s, btx0, bty, btw, bth, C_CARD, bc, 1.8)
        rect(s, btx0, bty, btw, Emu(200000), bc)
        tb(s, btx0 + Emu(30000), bty + Emu(30000), btw - Emu(50000), Emu(170000),
           btt, 8.5, bold=True, color=RGBColor(0x04, 0x04, 0x1C),
           align=PP_ALIGN.CENTER)
        tb(s, btx0 + Emu(30000), bty + Emu(240000), btw - Emu(50000),
           bth - Emu(280000), bts, 8.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── 下段: 虚構連モードへの道 ───────────────────────────────────────────
    tb(s, Inches(0.28), Inches(2.60), Inches(9.4), Emu(220000),
       "▶ 虚構連モード突入条件と流れ", 9, bold=True, color=C_GREEN)

    path_y = Inches(2.82)
    path_h = Emu(1800000)

    # 左: 突入条件
    pw_l = Inches(4.5)
    rect_b(s, Inches(0.28), path_y, pw_l, path_h, C_CARD, C_GREEN, 1.5)
    rect(s, Inches(0.28), path_y, Emu(45000), path_h, C_GREEN)
    tb(s, Inches(0.50), path_y + Emu(50000), pw_l - Emu(100000), Emu(260000),
       "虚構連モード突入条件", 10, bold=True, color=C_GREEN, font=FONT_H)
    tb(s, Inches(0.50), path_y + Emu(310000), pw_l - Emu(100000), path_h - Emu(370000),
       "① スペシャルボーナス → 確定突入\n\n"
       "② レギュラーボーナス → 突入抽選（設定差あり）\n\n"
       "③ 設定変更後CZ初成功時 → 約50%で高確スタート\n\n"
       "④ 高設定ほど突入率・継続率が高い（Long率が設定差の指標）",
       8.5, color=C_WHITE)

    # 矢印
    arrow_r(s, Inches(0.28) + pw_l + Emu(30000),
            path_y + path_h // 2, col=C_GREEN)

    # 右: 虚構連モードで得られるもの
    pw_r = Inches(4.65)
    rx2 = Inches(0.28) + pw_l + Emu(260000)
    rect_b(s, rx2, path_y, pw_r, path_h, C_CARD, C_GOLD, 1.5)
    rect(s, rx2, path_y, Emu(45000), path_h, C_GOLD)
    tb(s, rx2 + Emu(75000), path_y + Emu(50000), pw_r - Emu(100000), Emu(260000),
       "虚構連モードで得られるもの", 10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, rx2 + Emu(75000), path_y + Emu(310000), pw_r - Emu(100000), path_h - Emu(370000),
       "・高確率ボーナス抽選ループが継続する上位状態\n"
       "・Short / Middle / Long の3種（Longが最も継続）\n"
       "・高設定での9400枚一撃事例あり（爆発力は本物）\n"
       "・虚構連中にARROW告知で再度の興奮体験が続く",
       8.5, color=C_WHITE)

    footer(s,
           "設計コメント：スペシャルボーナス=虚構連確定という明確な頂点設計がプレイヤーの目標を単純化",
           "補足：設定変更後のCZ初成功約50%高確スタートが設定6据え置き狙いの重要な判断材料となる")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 虚構連モード詳細  [Part A - 6/6]
# ══════════════════════════════════════════════════════════════
def s_kykoren(prs):
    s = new_slide(prs)
    hdr(s, "虚構連モード ── Short / Middle / Long の違いと遊び方・ループ", "6/9")

    # 3列: Short/Middle/Long
    mode_data = [
        (C_CYAN,  "Short",
         "短期集中型",
         "ループ率：約56%\n転落率：1/20.98\n\n数回のボーナスループ\n後に転落することが多い\n\nCZ成功のたびに発生\n「次のモードへの足がかり」\n\n「すぐ終わったが枚数は\n出た」という報告も",
         "到達頻度\n高い"),
        (C_PUR2,  "Middle",
         "中期継続型",
         "ループ率：約81%\n転落率：1/53.55\n\n安定したボーナス\nループが継続する\n標準的な虚構連\n\n琴子ステージ滞在中が\n目安・設定差あり",
         "安定した\n枚数獲得"),
        (C_GOLD,  "Long",
         "長期継続型（最上位）",
         "ループ率：約95%\n転落率：1/185.74\n\n9400枚一撃事例が\nあるのはこのモード\n\n六花ステージ滞在が\n最上位の目印\n設定判別の重要指標",
         "爆発力\n最大"),
    ]
    mw = Inches(2.85)
    mgap = Inches(0.21)
    my = Inches(0.72)
    mh = Emu(3400000)

    for i, (mc, mn, msub, mbody, mbadge) in enumerate(mode_data):
        mx0 = Inches(0.28) + i * (mw + mgap)
        rect_b(s, mx0, my, mw, mh, C_CARD, mc, 2.0)
        # ヘッダーバー
        rect(s, mx0, my, mw, Emu(400000), RGBColor(0x0C, 0x08, 0x28))
        rect(s, mx0, my, Emu(45000), mh, mc)
        # ラベル
        rect_b(s, mx0 + mw - Emu(770000), my + Emu(80000),
               Emu(740000), Emu(230000), mc, C_WHITE, 0.8)
        tb(s, mx0 + mw - Emu(760000), my + Emu(90000),
           Emu(730000), Emu(210000), mbadge, 7, bold=True,
           color=RGBColor(0x04, 0x04, 0x1C), align=PP_ALIGN.CENTER)
        # タイトル
        tb(s, mx0 + Emu(65000), my + Emu(60000), mw - Emu(820000), Emu(290000),
           mn, 14, bold=True, color=mc, font=FONT_H)
        tb(s, mx0 + Emu(65000), my + Emu(330000), mw - Emu(90000), Emu(210000),
           msub, 8, color=C_GRAY)
        # 本文
        tb(s, mx0 + Emu(65000), my + Emu(560000), mw - Emu(90000), mh - Emu(620000),
           mbody, 8.5, color=C_WHITE)

    # 右端：虚構連ループ図
    rx = Inches(0.28) + 3 * (mw + mgap)
    rw = SLIDE_W - rx - Emu(200000)
    rect_b(s, rx, my, rw, mh, RGBColor(0x04, 0x10, 0x08), C_GREEN, 1.5)
    tb(s, rx + Emu(60000), my + Emu(60000), rw - Emu(80000), Emu(270000),
       "虚構連ループ", 11, bold=True, color=C_GREEN, font=FONT_H)

    loop_items = [
        (C_GREEN, "高確率\nボーナス抽選"),
        (C_PUR2,  "ボーナス\n当選"),
        (C_CYAN,  "ARROW JUDGEMENT\n当選契機を色で告知"),
        (C_GOLD,  "Short56%/Mid81%\nLong95%でループ"),
    ]
    lih = (mh - Emu(420000)) // 4
    for j, (lc, lt) in enumerate(loop_items):
        liy = my + Emu(380000) + j * lih
        rect_b(s, rx + Emu(30000), liy + Emu(30000),
               rw - Emu(60000), lih - Emu(50000), C_CARD, lc, 1.2)
        tb(s, rx + Emu(70000), liy + Emu(65000),
           rw - Emu(110000), lih - Emu(90000), lt, 8,
           bold=True, color=lc, align=PP_ALIGN.CENTER)
        if j < 3:
            # 小さな下矢印表現
            arrow_mid_y = liy + lih - Emu(10000)
            rect(s, rx + rw // 2 - Emu(30000), arrow_mid_y,
                 Emu(60000), Emu(50000), lc)

    footer(s,
           "設計コメント：Short56%→Middle81%→Long95%のループ率段階設計。Long（転落率1/185.74）突入時は9400枚超の爆発力",
           "補足：六花ステージ=Long（約95%ループ）が高設定判別の最重要指標。ループするたびにARROW告知で緊張体験が続く")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計  [Part B - 1/3]
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "Part B 分析 ── 面白さの設計（エピソードキャリーオーバー＋ARROW告知の独自性）", "7/9")

    # 上段横断バー：設計の核心
    rect(s, Inches(0.28), Inches(0.72), Inches(9.44), Emu(280000),
         RGBColor(0x44, 0x11, 0x77))
    tb(s, Inches(0.50), Inches(0.74), Inches(9.0), Emu(250000),
       "核心：「失敗が前進になる」キャリーオーバー設計 ＋ 「外部デバイスが告知する」ARROW体験", 9,
       bold=True, color=C_PUR2)

    # 下段2カラム
    col_l_x = Inches(0.28)
    col_l_w = Inches(4.5)
    col_r_x = Inches(5.05)
    col_r_w = Inches(4.65)
    col_y = Inches(1.08)
    col_h = Emu(3200000)

    # 左: キャリーオーバー設計の分析
    rect_b(s, col_l_x, col_y, col_l_w, col_h, C_CARD, C_PUR, 1.5)
    rect(s, col_l_x, col_y, Emu(45000), col_h, C_PUR)
    tb(s, col_l_x + Emu(75000), col_y + Emu(50000), col_l_w - Emu(100000), Emu(270000),
       "エピソードキャリーオーバーの革新性", 10, bold=True, color=C_PUR, font=FONT_H)

    # 比較ボックス（一般CZ vs 虚構推理CZ）
    comp_y = col_y + Emu(330000)
    comp_h = Emu(680000)
    comp_items = [
        (C_LTGRY, "一般的なCZ",
         "失敗 → 振り出しに戻る\n積み上げが消える（ロス感）"),
        (C_PUR,   "虚構推理CZ",
         "失敗 → エピソード持越し\n積み上げが残る（前進感）"),
    ]
    for ci, (cc, ct, cx) in enumerate(comp_items):
        ciy = comp_y + ci * (comp_h + Emu(50000))
        rect_b(s, col_l_x + Emu(50000), ciy,
               col_l_w - Emu(100000), comp_h, C_CARD2, cc, 1.2)
        rect(s, col_l_x + Emu(50000), ciy, Emu(30000), comp_h, cc)
        tb(s, col_l_x + Emu(120000), ciy + Emu(60000),
           col_l_w - Emu(200000), Emu(250000), ct, 9, bold=True, color=cc)
        tb(s, col_l_x + Emu(120000), ciy + Emu(310000),
           col_l_w - Emu(200000), Emu(320000), cx, 8, color=C_WHITE)

    # 心理効果まとめ
    psy_y = comp_y + 2 * (comp_h + Emu(50000)) + Emu(80000)
    tb(s, col_l_x + Emu(60000), psy_y, col_l_w - Emu(100000), Emu(220000),
       "生まれる心理効果", 8.5, bold=True, color=C_PUR2)
    psys = [
        "① 「諦めなくていい」という心理的安堵感",
        "② 「あと○エピソード」という継続動機",
        "③ 長期投資を「前進感」で正当化できる",
        "④ 失敗がストーリーの一部になる没入感",
    ]
    for pi, ps in enumerate(psys):
        tb(s, col_l_x + Emu(60000), psy_y + Emu(240000) + pi * Emu(290000),
           col_l_w - Emu(100000), Emu(260000), ps, 8, color=C_WHITE)

    # 右: ARROW告知の独自性分析
    rect_b(s, col_r_x, col_y, col_r_w, col_h, C_CARD, C_CYAN, 1.5)
    rect(s, col_r_x, col_y, Emu(45000), col_h, C_CYAN)
    tb(s, col_r_x + Emu(75000), col_y + Emu(50000), col_r_w - Emu(100000), Emu(270000),
       "ARROW告知という外部デバイスの独自価値", 10, bold=True, color=C_CYAN, font=FONT_H)

    arrow_pts = [
        (C_CYAN,  "視覚的差別化",
         "画面内の演出と異なり「筐体が光る」体験\n視野の端で色変化を察知する非日常感"),
        (C_PUR2,  "集中力の1点集中",
         "最終Gにすべての情報がARROWに集約される\n「何色が光るか」という純粋な緊張の頂点"),
        (C_GREEN, "社会的共有性",
         "周囲からも見えるため「話題になる」設計\n「ARROWが赤く光った」を隣と共有できる"),
        (C_GOLD,  "色と契機の直感的対応",
         "白=一撃成功・赤=弱チェリー・赤点滅=強チェリーという\n「色が当選契機を直接示す」7段階の明快設計"),
    ]
    apt_h = (col_h - Emu(320000)) // 4
    for ai, (ac, at, ab) in enumerate(arrow_pts):
        aiy = col_y + Emu(320000) + ai * apt_h
        rect(s, col_r_x + Emu(50000), aiy + Emu(30000),
             Emu(25000), apt_h - Emu(60000), ac)
        tb(s, col_r_x + Emu(110000), aiy + Emu(60000),
           col_r_w - Emu(160000), Emu(250000), at, 9, bold=True, color=ac)
        tb(s, col_r_x + Emu(110000), aiy + Emu(310000),
           col_r_w - Emu(160000), apt_h - Emu(360000), ab, 8, color=C_WHITE)

    footer(s,
           "設計コメント：「失敗が前進になる」キャリーオーバー設計と「色が当選契機を示す」ARROW告知の組み合わせは他機種に前例なし",
           "補足：ARROWは単なる演出でなく「何が当たったか」を教えるデバイス。白光=一撃成功で白7BB期待という設計意図が明快")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題  [Part B - 2/3]
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 独自設計の革新性 vs 到達性の問題", "8/9")

    # 上段帯：評価の現実
    rect_b(s, Inches(0.28), Inches(0.72), Inches(9.44), Emu(540000),
           C_CARD, C_RED, 1.5)
    rect(s, Inches(0.28), Inches(0.72), Emu(45000), Emu(540000), C_RED)
    tb(s, Inches(0.50), Inches(0.76), Inches(9.0), Emu(260000),
       "評価の現実：DMMレビュー平均1.4点（134件）という低評価　vs　高設定9400枚一撃事例", 9,
       bold=True, color=C_CRIM, font=FONT_H)
    tb(s, Inches(0.50), Inches(1.06), Inches(9.0), Emu(220000),
       "独自性の高さとプレイヤーが感じる「到達できない」課題が評価を二極化させている", 8,
       color=C_WHITE)

    # 2カラム：良い点 vs 課題
    cy2 = Inches(1.42)
    ch2 = Emu(2750000)
    cl_w = Inches(4.5)
    cr_w = Inches(4.65)

    # 左: 良い点
    rect_b(s, Inches(0.28), cy2, cl_w, ch2, C_CARD, C_GREEN, 1.5)
    rect(s, Inches(0.28), cy2, Emu(45000), ch2, C_GREEN)
    tb(s, Inches(0.50), cy2 + Emu(50000), cl_w - Emu(100000), Emu(270000),
       "良い点（革新性・独自設計）", 11, bold=True, color=C_GREEN, font=FONT_H)

    pros = [
        (C_GREEN, "エピソードキャリーオーバー",
         "失敗が前進になる唯一無二の設計。\nパチスロ史上でも稀なCZ設計思想。"),
        (C_CYAN,  "ARROW告知デバイス（当選契機告知）",
         "色=当選契機対応（白=一撃/赤点滅=強チェリー確定）という\n他機種にない外部デバイス告知体験。話題性・コミュニティ形成に貢献。"),
        (C_PUR2,  "虚構連モード3種の明確化",
         "Short56%→Middle81%→Long95%というループ率段階設計が\nプレイヤーの目標を明確化。六花ステージ=Longの判別も容易。"),
        (C_GOLD,  "高設定時の爆発力",
         "9400枚一撃事例が示す爆発ポテンシャル。\nスペック上の夢は本物。"),
    ]
    pro_h = (ch2 - Emu(320000)) // 4
    for pi, (pc, pt, pb) in enumerate(pros):
        piy = cy2 + Emu(320000) + pi * pro_h
        rect_b(s, Inches(0.40), piy + Emu(20000),
               cl_w - Emu(180000), pro_h - Emu(40000), C_CARD2, pc, 1.0)
        rect(s, Inches(0.40), piy + Emu(20000), Emu(20000),
             pro_h - Emu(40000), pc)
        tb(s, Inches(0.55), piy + Emu(50000),
           cl_w - Emu(240000), Emu(240000), pt, 8.5, bold=True, color=pc)
        tb(s, Inches(0.55), piy + Emu(290000),
           cl_w - Emu(240000), pro_h - Emu(330000), pb, 7.5, color=C_WHITE)

    # 右: 課題
    rect_b(s, Inches(5.05), cy2, cr_w, ch2, C_CARD, C_RED, 1.5)
    rect(s, Inches(5.05), cy2, Emu(45000), ch2, C_RED)
    tb(s, Inches(5.27), cy2 + Emu(50000), cr_w - Emu(100000), Emu(270000),
       "課題（到達性・演出の問題）", 11, bold=True, color=C_CRIM, font=FONT_H)

    cons = [
        (C_CRIM,  "虚構連モードへの到達が渋い",
         "高設定以外ではなかなか辿り着けない\n体験設計が低設定時の離脱を早める。"),
        (C_RED,   "CZ中の演出強度が弱い",
         "「演出が弱い・引きが弱い」という\n口コミが多数。緊張感の持続が課題。"),
        (C_CRIM,  "強チェリー→CZ不発が多い",
         "引いたのに入らない体験が\nフラストレーションとなり低評価に直結。"),
        (C_GRAY,  "低設定体験の改善が必要",
         "DMMレビュー低評価の主因は低設定時。\n低設定でも楽しめる設計の補完が求められる。"),
    ]
    con_h = (ch2 - Emu(320000)) // 4
    for ci, (cc, ct, cb) in enumerate(cons):
        ciy = cy2 + Emu(320000) + ci * con_h
        rect_b(s, Inches(5.17), ciy + Emu(20000),
               cr_w - Emu(180000), con_h - Emu(40000), C_CARD2, cc, 1.0)
        rect(s, Inches(5.17), ciy + Emu(20000), Emu(20000),
             con_h - Emu(40000), cc)
        tb(s, Inches(5.32), ciy + Emu(50000),
           cr_w - Emu(250000), Emu(240000), ct, 8.5, bold=True, color=cc)
        tb(s, Inches(5.32), ciy + Emu(290000),
           cr_w - Emu(250000), con_h - Emu(330000), cb, 7.5, color=C_WHITE)

    footer(s,
           "設計コメント：独自性は高いが到達性の問題が評価を押し下げる。設定投入＋演出補強で大化けする可能性",
           "補足：低評価1.4点の多くは低設定体験。高設定稼働時のレビューは「最高の台」という評価が多数存在する")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること  [Part B - 3/3]
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── L虚構推理 設計から学べること", "9/9")

    # 左: 3つの設計的学び
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(300000), RGBColor(0x44, 0x11, 0x77))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(225000),
       "設計から学べる3つのエッセンス", 10, bold=True, color=C_CYAN, font=FONT_H)

    elems = [
        (C_PUR,   "① 失敗キャリーオーバーという設計的発明",
         "失敗を「無」にせず「前進」に変える設計は\n"
         "プレイヤーの継続動機として機能する。\n"
         "「あと○エピソード」が長期稼働を支える。"),
        (C_CYAN,  "② 外部告知デバイスARROWの体験価値",
         "JUDGEMENTゲームで「色=当選契機」を直接告知する革新。\n"
         "白=一撃成功・赤点滅=強チェリー確定という明快設計。\n"
         "「みんなで見守る」ホール内コミュニティ形成を実現。"),
        (C_GREEN, "③ 3種モードによる上位状態の可視化",
         "Short56% / Middle81% / Long95%という段階設計が\n"
         "「今自分がどのループ率にいるか」を意識させる。\n"
         "六花ステージ滞在が高設定判別の重要指標にもなる。"),
    ]
    for i, (ac, t, b) in enumerate(elems):
        ey = ly + Emu(300000) + i * Emu(1250000)
        rect_b(s, lx, ey, lw, Emu(1200000), C_CARD, ac, 1.5)
        rect(s, lx, ey, Emu(45000), Emu(1200000), ac)
        tb(s, lx + Emu(75000), ey + Emu(50000), lw - Emu(95000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), ey + Emu(305000), lw - Emu(95000), Emu(830000),
           b, 8, color=C_WHITE)

    # 右: 設計原則＋総括
    rx, ry = Inches(5.05), Inches(0.72)
    rw = Inches(4.65)

    rect(s, rx, ry, rw, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(215000),
       "設計原則（他機種への応用）", 10, bold=True, color=C_CYAN, font=FONT_H)

    principles = [
        (C_PUR,   "失敗をストックする設計は「諦めない」心理を生む（EP1開始25%→EP5開始75%の設計）"),
        (C_CYAN,  "ARROWは「色＝当選契機」対応で情報を直感的に伝える外部デバイスの革新"),
        (C_PUR2,  "Short56%→Middle81%→Long95%の3段階ループが明確な目標階層を作る"),
        (C_GOLD,  "六花ステージ＝Long到達の可視化がプレイヤーの設定判別指標にもなる"),
        (C_CRIM,  "到達可能性の改善こそ本機が化けるための最重要課題"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(280000) + i * Emu(480000)
        rect(s, rx, py0, Emu(20000), Emu(440000), ac)
        tb(s, rx + Emu(50000), py0 + Emu(55000), rw - Emu(60000), Emu(360000),
           p, 8, bold=(i == 4), color=C_CRIM if i == 4 else C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(2700000), rw, Emu(1050000),
           RGBColor(0x06, 0x04, 0x1C), C_PUR, 1.5)
    rect(s, rx, ry + Emu(2700000), Emu(45000), Emu(1050000), C_PUR)
    tb(s, rx + Emu(75000), ry + Emu(2750000), rw - Emu(95000), Emu(270000),
       "総括", 10, bold=True, color=C_PUR2, font=FONT_H)
    tb(s, rx + Emu(75000), ry + Emu(3020000), rw - Emu(95000), Emu(680000),
       "独自の設計思想（キャリーオーバー＋ARROW）を持つ意欲作。\n"
       "到達性の改善と演出強化が実現すれば\n"
       "唯一無二の体験として業界に残る台になり得る。",
       8.5, color=C_WHITE)

    footer(s,
           "設計コメント：失敗を前進に変える発想と外部デバイス告知の組み合わせは後世のパチスロ設計に影響を与えうる革新",
           "補足：高設定での9400枚一撃実績が示す通り、スペックの爆発力は本物。設定投入ありきで真価を発揮する機種")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Part A: プレイヤー視点説明 6枚
    s_title(prs)    # 1: タイトル・スペック・この台の3ポイント
    s_flow(prs)     # 2: ゲームフロー全体図
    s_normal(prs)   # 3: 通常時の遊び方
    s_cz(prs)       # 4: CZ「鋼人七瀬攻略議会」の攻略
    s_bonus(prs)    # 5: ボーナス後の遊び方
    s_kykoren(prs)  # 6: 虚構連モード詳細

    # Part B: 分析パート 3枚
    s_design(prs)       # 7: 面白さの設計
    s_pros_cons(prs)    # 8: 良い点と課題
    s_matome(prs)       # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
