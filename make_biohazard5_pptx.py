"""
スマスロ バイオハザード5 機種説明＋分析 統合版 PPTXジェネレーター v1
出力: proposals/機種分析/バイオハザード5/biohazard5_guide_v1.pptx
テーマ: 深黒×赤(#CC1111)×オレンジ(#FF6600)×白（バイオ世界観）
情報源: ちょんぼりすた・一撃・なな徹・アミュタメ・altema（2025年3月〜）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "バイオハザード5", "biohazard5_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深黒×赤×オレンジ×白 バイオ世界観）──────────────────────
C_BG    = RGBColor(0x06, 0x02, 0x02)
C_CARD  = RGBColor(0x12, 0x04, 0x04)
C_CARD2 = RGBColor(0x1A, 0x06, 0x06)
C_ROW   = RGBColor(0x16, 0x05, 0x05)
C_RED   = RGBColor(0xCC, 0x11, 0x11)   # メイン赤 #CC1111
C_RED2  = RGBColor(0xEE, 0x22, 0x22)   # 明るい赤
C_ORG   = RGBColor(0xFF, 0x66, 0x00)   # オレンジ #FF6600
C_ORG2  = RGBColor(0xFF, 0x99, 0x00)   # 明るいオレンジ
C_GOLD  = RGBColor(0xCC, 0xAA, 0x22)   # ゴールド
C_WHITE = RGBColor(0xF0, 0xEC, 0xE8)
C_CREAM = RGBColor(0xF0, 0xD8, 0xB0)
C_GRAY  = RGBColor(0x99, 0x88, 0x80)
C_LTGRY = RGBColor(0x55, 0x44, 0x40)
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_PINK  = RGBColor(0xFF, 0x44, 0x88)
C_BLOOD = RGBColor(0x88, 0x00, 0x00)   # 暗い血の赤

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（廃墟・暗黒×赤グロー）─────────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (6, 2, 2))
    draw = ImageDraw.Draw(img)
    # 廃墟感：斜めのひび割れ的ライン
    for i in range(0, w + h, 90):
        draw.line([(i, 0), (0, i)], fill=(10, 3, 3), width=1)
    for i in range(0, w + h, 200):
        draw.line([(i, 0), (0, i)], fill=(16, 4, 4), width=2)
    # 下部からの赤グロー（血溜まり・バイオ感）
    for y in range(h - 100, h):
        t = (y - (h - 100)) / 100
        r = int(80 * t)
        g = int(4 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, 0))
    # 上部の暗さ
    for y in range(0, 35):
        t = (35 - y) / 35 * 0.6
        draw.line([(0, y), (w, y)], fill=(int(12 * t), 0, 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_RED)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_RED)


def net_note(slide):
    tb(slide, Inches(7.0), Inches(5.38), Inches(2.9), Emu(180000),
       "※ネット解析情報より（ちょんぼりすた・一撃・なな徹・アミュタメ）", 6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x0E, 0x04, 0x04))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_RED)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_ORG)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_RED
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x08, 0x02, 0x02))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_RED)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_RED)

    tb(s, Inches(0.22), Inches(0.4), Inches(5.0), Emu(330000),
       "機種説明＋分析 統合ガイド  v1（パチスロアワード2025ノミネート）", 10, color=C_RED2, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.88), Inches(5.1), Emu(900000),
       "スマスロ\nバイオハザード5", 28, bold=True, color=C_RED, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.70), Inches(5.0), Emu(280000),
       "スマスロ（Lパチスロ）── 増殖分裂×ゲーム数上乗せ×ホラーサバイバル設計", 9, color=C_CREAM, font=FONT_H)

    # スペック表
    specs = [
        ("メーカー",        "エンターライズ　2025年3月3日導入"),
        ("設定",           "1〜6段階"),
        ("AT純増",         "ハザードラッシュ：約2.5枚/G"),
        ("上位AT純増",     "プレミアムハザードラッシュ：約4.0枚/G"),
        ("設定1機械割",     "97.3%（参考値）"),
        ("設定6機械割",     "110.8%（参考値）"),
        ("天井",           "通常時999G（恩恵：AT確定）"),
        ("初期AT",         "50G＋α（増殖・変異上乗せあり）"),
        ("最大上乗せ",      "810G（30G×27分裂）"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(3.00) + i * Emu(235000)
        tb(s, Inches(0.22), ry, Inches(1.7), Emu(210000),
           k, 7.5, color=C_GRAY)
        tb(s, Inches(1.92), ry, Inches(3.3), Emu(210000),
           v, 7.5, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_RED,   "① 増殖×変異の上乗せシステム（核心）",
         "上乗せゲーム数が「増殖」で2倍に分裂。\n10G→20G→40G→80Gと連鎖。\n最大30G×27分裂＝810Gの爆発力。"),
        (C_ORG,   "② 2層AT設計（通常AT＋上位AT）",
         "ハザードラッシュ（純増2.5枚）で土台を作り\nプレミアムハザードラッシュ（純増4.0枚）で\n爆発。上位ATはゲーム数減算なしの特殊仕様。"),
        (C_CYAN,  "③ ホラー×サバイバル世界観の継承",
         "5号機バイオ5の完全再現ゲーム性。\nパニックゾーン・ウェスカーゾーンの\nCZ演出で原作ホラー感を忠実に再現。"),
    ]
    bx_h = Inches(1.25)   # ボックス高さ（1.3→1.25に縮小して下端はみ出し防止）
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.40) + i * Emu(1490000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), bx_h, C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), bx_h, ac)
        tb(s, Inches(5.85), y0 + Emu(55000), Inches(3.8), Emu(290000),
           kw, 11.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(350000), Inches(3.8), Emu(430000),
           desc, 8.5, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「増殖×分裂上乗せ×2層AT」── 5号機バイオ5のゲーム性をスマスロで完全再現",
           "パチスロアワード2025ノミネート機種。純増2.5枚/Gのゲーム数上乗せ型で増殖連鎖が出玉の本体")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（全ルートを蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→CZ→AT→上位ATの全ルート", "2/9")

    # 上段：通常→CZ→AT の基本フロー
    top_y = Inches(0.75)
    top_h = Emu(1100000)

    flow1 = [
        (C_CARD2, C_GRAY,  "通常遊技",          "モード管理\nレア役/ベル3連\nでCZ抽選"),
        (C_CARD,  C_RED2,  "パニック\nゾーン",   "AT期待度\n約30〜50%"),
        (C_CARD,  C_RED,   "ウェスカー\nゾーン", "高期待度CZ\n期待度アップ"),
        (C_CARD,  C_ORG,   "ハザード\nラッシュ", "初期50G+α\n純増2.5枚/G"),
        (C_CARD,  C_GOLD,  "プレミアム\nHR",     "純増4.0枚\n減算なし"),
    ]
    bw1 = Inches(1.6)
    gap1 = Inches(0.17)
    sx1 = Inches(0.28)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), top_y + Emu(70000), bw1 - Emu(60000), Emu(360000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), top_y + Emu(480000), bw1 - Emu(45000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), cy1)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1120000), Inches(3.0), Emu(260000),
       "天井：通常時999G → AT確定", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1360000), Inches(3.5), Emu(5000), C_CYAN)
    tb(s, sx1 + Inches(3.6), top_y + Emu(1120000), Inches(4.0), Emu(260000),
       "直撃：中段チェリー1/10 or ロングフリーズ → 上位AT直行", 7.5, color=C_RED2)

    # 中段区切り線
    rect(s, 0, Inches(2.1), SLIDE_W, Emu(5000), RGBColor(0x44, 0x10, 0x10))

    # 下段：AT内フロー
    bot_y = Inches(2.18)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,  C_ORG,   "ハザード\nラッシュ",   "初期50G\n純増2.5枚/G\n上乗せ型"),
        (C_CARD,  C_RED2,  "インフェク\nション",    "初期特化ゾーン\n全役で上乗せ\n増殖期待大"),
        (C_CARD,  C_RED,   "上乗せ\nゾーン",       "特化ゾーン\n変異上乗せ\n連鎖可能"),
        (RGBColor(0x18,0x06,0x02), C_GOLD,
         "プレミアム\nHR",        "上位AT\n純増4.0枚/G\n減算なし"),
    ]
    bw2 = Inches(1.95)
    gap2 = Inches(0.22)
    sx2 = Inches(0.28)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8)
        tb(s, bx + Emu(40000), bot_y + Emu(70000), bw2 - Emu(70000), Emu(370000),
           lbl, 10, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), bot_y + Emu(490000), bw2 - Emu(55000), Emu(470000),
           sub, 8, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw2 + Emu(12000), cy2)

    # 右端：CWZへの道（右端はみ出し対策: 幅を残り幅に合わせる）
    rx_al = sx2 + 4 * (bw2 + gap2)
    rw_al = SLIDE_W - rx_al - Emu(90000)   # 右端9.9inchに収まるよう計算
    rect_b(s, rx_al, bot_y, rw_al, bot_h, C_CARD, C_CYAN, 1.5)
    rect(s, rx_al, bot_y, Emu(30000), bot_h, C_CYAN)
    tb(s, rx_al + Emu(50000), bot_y + Emu(70000), rw_al - Emu(60000), Emu(310000),
       "クライマックス\nウェスカーZ\n(CZ)", 8, bold=True, color=C_CYAN, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, rx_al + Emu(50000), bot_y + Emu(490000), rw_al - Emu(60000), Emu(460000),
       "AT終了後\n一部突入\nPHR確定", 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 下部アノテーション
    tb(s, sx2, bot_y + Emu(1100000), Inches(5.0), Emu(260000),
       "増殖：10G→20G→40G→80G（連鎖で最大810G）", 7.5, color=C_ORG)
    rect(s, sx2, bot_y + Emu(1340000), Inches(5.0), Emu(5000), C_ORG)

    net_note(s)
    footer(s, "上段=通常〜AT突入ルート（CZ2種＋直撃＋天井）、下段=AT内昇格フロー（増殖→インフェクション→プレミアムHR）",
           "プレミアムハザードラッシュ（上位AT）はAT終了後のクライマックスウェスカーゾーン成功でも到達可能")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（天井含む全ルート）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 全AT突入ルート・モード・天井管理", "3/9")

    # 左：AT突入ルート図
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x55, 0x10, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜AT突入ルート（全4系統）", 10, bold=True, color=C_ORG)

    routes = [
        (C_RED2,  "ルート①  CZ「パニックゾーン / ウェスカーゾーン」経由",
         "ベル3連orレア役でCZ突入を抽選。\nウェスカーゾーンなら期待度アップ。\nCZ成功でハザードラッシュ突入。"),
        (C_RED,   "ルート②  中段チェリー特殊ルート（直行）",
         "中段チェリー成立の約1/10で発生。\nプレミアムハザードラッシュへ直行。\n出現頻度は低いが最高ルートの一つ。"),
        (C_ORG,   "ルート③  ロングフリーズ直撃",
         "ロングフリーズ発生でプレミアムHR直行。\nフリーズ中の演出を楽しみながら待機。\n出現頻度は極めてレア。"),
        (C_CYAN,  "ルート④  天井（通常時999G）",
         "通常時999GでAT確定の天井。\n設定変更後は天井リセットの可能性あり。\n深いゲーム数での立ち回りに有効。"),
    ]
    # 4ボックスを4.9inch以内に収める: 利用可能高さ = 4.9 - 1.037 = 3.863inch
    # 各ボックス高さ880000EMU(0.962inch)、間隔 = 880000EMU → 合計 = 4 * 880000 = 3520000EMU(3.847inch) OK
    rte_h = Emu(880000)
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(290000) + i * rte_h
        rect_b(s, lx, iy, lw, rte_h - Emu(15000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), rte_h - Emu(15000), ac)
        tb(s, lx + Emu(75000), iy + Emu(40000), lw - Emu(100000), Emu(240000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(290000), lw - Emu(100000), Emu(540000),
           b, 7.5, color=C_WHITE)

    # 右：モードと内部状態
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x55, 0x10, 0x10))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "内部モード（4段階）とCZ期待度", 10, bold=True, color=C_ORG)

    modes = [
        (C_GRAY,  "LOW モード",      "最も低い",    "CZ当選期待度が最低"),
        (C_RED2,  "MID モード",      "低〜中",      "通常遊技の基本モード"),
        (C_RED,   "HI モード",       "高め",        "レア役でCZ当選しやすい"),
        (C_GOLD,  "SP モード",       "最高確",      "AT/CZへの当選最優先"),
    ]
    mode_h = Emu(690000)
    for i, (ac, mode, expect, desc) in enumerate(modes):
        cy = ry + Emu(290000) + i * mode_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, mode_h, bg)
        rect(s, rx, cy, Emu(35000), mode_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(55000), Inches(1.2), Emu(260000),
           mode, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.2), cy + Emu(60000), Inches(0.7), Emu(240000),
           expect, 8, bold=True, color=C_ORG2, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(310000), rw - Emu(70000), Emu(320000),
           desc, 7.5, color=C_WHITE)

    # 状態遷移メモ＋チャンス役確率（1ボックスに統合して下端はみ出し防止）
    # mode_y_end = ry + Emu(290000) + 4 * Emu(690000) = ry + Emu(3050000)
    memo_y = ry + Emu(3060000)
    memo_h = Inches(4.88) - memo_y   # 4.88inchまで = フッター直前
    rect_b(s, rx, memo_y, rw, memo_h, C_CARD2, C_RED, 1.5)
    tb(s, rx + Emu(60000), memo_y + Emu(50000), rw - Emu(80000), Emu(200000),
       "モード管理ルール", 8.5, bold=True, color=C_RED)
    tb(s, rx + Emu(60000), memo_y + Emu(260000), rw - Emu(80000), Emu(320000),
       "設定変更時・AT終了時にモード再抽選。レア役でモード昇格抽選。", 7.5, color=C_GRAY)
    tb(s, rx + Emu(60000), memo_y + Emu(590000), rw - Emu(80000), Emu(200000),
       "CZ抽選対象役：ベル3連・スイカ・チェリー・チャンス目", 7.5, bold=True, color=C_ORG)

    net_note(s)
    footer(s, "通常時の戦略：モード（LOW/MID/HI/SP）を意識しCZ当選を見極め、天井999Gを基準に立ち回る",
           "状態はCZ・AT当選まで転落しない安心設計。設定変更後のモード状況確認も重要")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み（サバイバル要素と絡めて）
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ/前兆の仕組み ── パニックゾーン×ウェスカーゾーン×サバイバル演出", "4/9")

    # 左：2種のCZ詳細
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x10, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "CZ 2種：パニックゾーン vs ウェスカーゾーン", 10, bold=True, color=C_ORG)

    czs = [
        (C_RED2, "パニックゾーン（通常CZ）",
         "通常時のCZ突入で発生。\n"
         "成功でハザードラッシュ突入確定。\n\n"
         "期待度：約30〜50%（モードにより変動）\n\n"
         "ゾーン中はゾンビ軍団との攻防演出。\n"
         "小役成立でサバイバー（主人公）のHP回復\n"
         "→HP維持でCZ継続→成功確定の流れ。\n"
         "ゾンビに囲まれる絶体絶命感がホラー体験を演出。"),
        (C_RED, "ウェスカーゾーン（上位CZ）",
         "パニックゾーンより高期待度のCZ。\n"
         "モードが高い状態ほど突入しやすい。\n\n"
         "期待度：パニックゾーンより明確に高め\n\n"
         "アルバート・ウェスカーが登場する特殊ゾーン。\n"
         "成功時は直接ハザードラッシュへ。\n"
         "ウェスカーの登場演出が成功の鍵となり\n"
         "原作プレイヤーには感情移入できる演出設計。"),
    ]
    for i, (ac, t, b) in enumerate(czs):
        iy = ly + Emu(260000) + i * Emu(2170000)
        rect_b(s, lx, iy, lw, Emu(2100000), C_CARD, ac, 2.0)
        rect(s, lx, iy, Emu(45000), Emu(2100000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 10, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(330000), lw - Emu(100000), Emu(1650000),
           b, 8, color=C_WHITE)

    # 右：前兆・クライマックスウェスカーゾーン・設計
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x55, 0x10, 0x10))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "前兆と特殊CZ：クライマックスウェスカーゾーン", 10, bold=True, color=C_ORG)

    specials = [
        (C_CYAN, "クライマックスウェスカーゾーン（AT後CZ）",
         "ハザードラッシュ（AT）終了後の一部で突入。\n"
         "成功でプレミアムハザードラッシュへ昇格。\n\n"
         "ウェスカーとの最終決戦演出が展開。\n"
         "\"810%到達\"の期待度表示で到達濃厚演出。\n"
         "セットストック5個以上保有時に突入しやすい。"),
        (C_GOLD, "セットストックと上位AT昇格抽選",
         "AT突入時のセットストック保有数で\n"
         "上位ATへの昇格抽選が変化。\n\n"
         "セットストック1個以上：昇格抽選対象\n"
         "セットストック5個以上：昇格チャンス大\n\n"
         "通常時のCZ・AT中の役成立で\n"
         "セットストックを獲得できる。"),
        (C_ORG,  "ホラー×サバイバルとパチスロの融合点",
         "CZ中のHP演出＝「生き残れるか？」の\n"
         "サバイバル体験をゲーム内に取り込む設計。\n\n"
         "成功時の緊張→解放の感情フローが\n"
         "通常のボーナス告知より強い達成感を生む。\n"
         "原作ゲームのジャンル感を維持した秀逸設計。"),
    ]
    for i, (ac, t, b) in enumerate(specials):
        iy = ry + Emu(260000) + i * Emu(1380000)
        rect_b(s, rx, iy, rw, Emu(1310000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(1310000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(870000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "CZ設計の核心：「パニック（通常）vs ウェスカー（上位）」の2段階構造でCZ出現時の期待感に差をつける",
           "HP制サバイバル演出はゲーム内の緊張→解放フローをスロットCZ体験に直接移植した秀逸な設計")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（何をすれば出玉が伸びる）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「ハザードラッシュ」── 出玉を伸ばす増殖×変異の仕組み", "5/9")

    # 左：AT中の上乗せシステム
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x10, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "AT中の上乗せシステム（増殖×変異）", 10, bold=True, color=C_ORG)

    # 上乗せゲーム数の種類
    rect_b(s, lx, ly + Emu(260000), lw, Emu(500000), C_CARD, C_RED2, 1.5)
    rect(s, lx, ly + Emu(260000), Emu(40000), Emu(500000), C_RED2)
    tb(s, lx + Emu(70000), ly + Emu(310000), lw - Emu(90000), Emu(200000),
       "上乗せゲーム数の種類（6パターン）", 9, bold=True, color=C_RED2)

    nums = [("10G", C_GRAY), ("20G", C_RED2), ("30G", C_RED),
            ("33G", C_ORG), ("50G", C_ORG2), ("55G", C_GOLD)]
    num_w = lw // 6
    for i, (n, nc) in enumerate(nums):
        nx = lx + i * num_w
        rect_b(s, nx + Emu(8000), ly + Emu(580000), num_w - Emu(16000), Emu(160000),
               C_CARD2, nc, 1.2)
        tb(s, nx + Emu(8000), ly + Emu(595000), num_w - Emu(16000), Emu(140000),
           n, 8, bold=True, color=nc, align=PP_ALIGN.CENTER, wrap=False)

    tb(s, lx + Emu(60000), ly + Emu(760000), lw - Emu(80000), Emu(200000),
       "33G・55G成立時は「増殖 or 変異上乗せ」濃厚！", 7.5, bold=True, color=C_ORG, wrap=False)

    # 増殖の仕組み
    rect_b(s, lx, ly + Emu(1010000), lw, Emu(1200000), C_CARD, C_ORG, 2.0)
    rect(s, lx, ly + Emu(1010000), Emu(40000), Emu(1200000), C_ORG)
    tb(s, lx + Emu(70000), ly + Emu(1060000), lw - Emu(90000), Emu(260000),
       "増殖上乗せ（分裂システム）", 10, bold=True, color=C_ORG, font=FONT_H)
    tb(s, lx + Emu(70000), ly + Emu(1330000), lw - Emu(90000), Emu(780000),
       "上乗せゲーム数が「増殖」で2倍に分裂して連鎖。\n\n"
       "10G → 20G → 40G → 80G（3連鎖）\n"
       "20G → 40G → 80G（2連鎖）\n"
       "30Gスタートの場合も連鎖で膨らむ\n\n"
       "最大：30G × 27分裂 ＝ 810G上乗せが可能！\n"
       "増殖成功ごとにゲーム数が加算される。",
       8, color=C_WHITE)

    # 変異の仕組み
    rect_b(s, lx, ly + Emu(2280000), lw, Emu(860000), C_CARD, C_RED, 1.5)
    rect(s, lx, ly + Emu(2280000), Emu(40000), Emu(860000), C_RED)
    tb(s, lx + Emu(70000), ly + Emu(2330000), lw - Emu(90000), Emu(260000),
       "変異上乗せ（特殊大量上乗せ）", 10, bold=True, color=C_RED, font=FONT_H)
    tb(s, lx + Emu(70000), ly + Emu(2600000), lw - Emu(90000), Emu(490000),
       "増殖とは別系統の大量上乗せ抽選。\n"
       "変異が発動すると一気に大量Gを上乗せ。\n"
       "33G・55Gの一部で変異に発展する。\n"
       "上乗せ状態（通常/高確/超高確）により発生率変化。",
       8, color=C_WHITE)

    # 右：上乗せ状態と消化フロー
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x55, 0x10, 0x10))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "上乗せ状態と消化フロー", 10, bold=True, color=C_ORG)

    jotai = [
        (C_GRAY,  "上乗せ状態：通常",
         "AT消化中の基本状態。\n"
         "役成立時に上乗せ抽選。\n"
         "増殖発生率は低め。"),
        (C_RED2,  "上乗せ状態：高確",
         "上乗せ当選率・増殖率アップ。\n"
         "ステージ開始時に一部移行。\n"
         "積極的に出玉を積み重ねる状態。"),
        (C_GOLD,  "上乗せ状態：超高確",
         "最高の上乗せ状態。\n"
         "増殖・変異の発生率が大幅向上。\n"
         "突入時は大量上乗せのチャンス。"),
    ]
    for i, (ac, t, b) in enumerate(jotai):
        iy = ry + Emu(260000) + i * Emu(990000)
        rect_b(s, rx, iy, rw, Emu(920000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(920000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(530000),
           b, 8, color=C_WHITE)

    # インフェクション解説
    rect_b(s, rx, ry + Emu(3230000), rw, Emu(1140000), C_CARD, C_CYAN, 2.0)
    rect(s, rx, ry + Emu(3230000), Emu(40000), Emu(1140000), C_CYAN)
    tb(s, rx + Emu(70000), ry + Emu(3280000), rw - Emu(90000), Emu(260000),
       "インフェクション（初期特化ゾーン）", 10, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, rx + Emu(70000), ry + Emu(3560000), rw - Emu(90000), Emu(780000),
       "AT突入時に一定確率で発生する初期特化ゾーン。\n"
       "全役で上乗せ抽選が行われる最高の初期状態。\n"
       "上乗せ当選時は増殖発生に大きく期待できる。\n"
       "インフェクション→増殖連鎖が出玉爆発の黄金ルート。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "出玉を伸ばす鍵：33G/55G上乗せからの増殖連鎖＋インフェクション突入＋超高確状態の重複が最大の上振れ条件",
           "増殖システムは「2倍ずつ分裂」という直感的な楽しさ。最大810G上乗せが理論上の夢上限")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位AT「プレミアムハザードラッシュ」── 到達ルートと遊び方", "6/9")

    # 上段：到達ルート
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(260000), RGBColor(0x44, 0x10, 0x10))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(210000),
       "プレミアムハザードラッシュ（PHR）到達ルート", 9, bold=True, color=C_ORG)

    route_boxes = [
        (C_RED2,  "ハザード\nラッシュ",         "通常AT\n純増2.5枚/G\n土台を作る"),
        (C_RED,   "CWZ成功",                   "AT終了後\nCZ成功で\nPHR昇格"),
        (C_ORG,   "中段チェリー\n1/10",          "通常時\n直撃で\nPHR直行"),
        (C_GOLD,  "プレミアム\nHR",             "上位AT\n純増4.0枚\n減算なし"),
    ]
    bw_r = Inches(2.1)
    gap_r = Inches(0.26)
    sx_r = Inches(0.35)
    cy_r = Inches(1.42)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x18, 0x06, 0x00), ac, 2.0 if i >= 2 else 1.5)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(55000), Emu(370000), lbl, 10, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), cy_r - bh_r // 2 + Emu(460000),
           bw_r - Emu(45000), Emu(480000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw_r + Emu(12000), cy_r)

    # 中区切り
    rect(s, 0, Inches(2.06), SLIDE_W, Emu(5000), RGBColor(0x44, 0x10, 0x10))

    # 下段左：PHRの詳細仕様
    lx2, ly2 = Inches(0.28), Inches(2.12)
    lw2 = Inches(4.55)
    lh2 = Emu(2620000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_GOLD, 2.0)
    rect(s, lx2, ly2, Emu(45000), lh2, C_GOLD)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "プレミアムHRの仕様と遊び方", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【純増約4.0枚/G】\n"
       "通常AT（2.5枚/G）から大幅に純増アップ。\n"
       "同じゲーム数でも1.6倍の出玉スピード。\n\n"
       "【ゲーム数減算なし】\n"
       "PHR消化中はゲーム数が減算されない特殊仕様。\n"
       "消化中に上乗せしたゲーム数は\n"
       "PHR終了後の下位ATで消化できる仕組み。\n\n"
       "【810%到達でPHR濃厚】\n"
       "ウェスカーゾーン中に期待度表示が810%になれば\n"
       "プレミアムHR突入が濃厚。\n\n"
       "【消化中の注意点】\n"
       "PHR中はレア役ごとに上乗せ・特化ゾーン移行を抽選。\n"
       "消化終了まで画面を見逃さないこと。",
       8, color=C_WHITE)

    # 下段右：PHR突入条件の詳細
    rx2, ry2 = Inches(5.05), Inches(2.12)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2, C_CARD, C_CYAN, 1.8)
    rect(s, rx2, ry2, Emu(45000), lh2, C_CYAN)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "PHR突入条件の詳細と期待ゲーム数", 11, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【条件①：クライマックスウェスカーゾーン成功】\n"
       "AT終了後の一部で突入するCZ。\n"
       "成功でプレミアムHR確定。\n"
       "セットストック5個以上で突入チャンス大。\n\n"
       "【条件②：初当たり時のセットストック昇格】\n"
       "AT突入時にセットストック1個以上保有→昇格抽選。\n"
       "セットストック数が多いほど昇格期待度が上がる。\n\n"
       "【条件③：中段チェリー（通常時）直撃】\n"
       "通常時の中段チェリー成立の約1/10で\n"
       "直接プレミアムHRへ突入する。\n\n"
       "【条件④：ロングフリーズ】\n"
       "発生確率は極めて低いが最も華やかな突入ルート。\n"
       "プレミアムHR直行が確定する。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "PHRの価値：純増4.0枚/G×ゲーム数減算なし×上乗せゲーム数は下位ATで消化という構造が出玉の最高峰",
           "810%到達演出はウェスカーゾーン内の特殊演出。到達時点でPHR突入が濃厚なサプライズ告知")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（ホラー×サバイバルとパチスロの融合）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── ホラー×サバイバルとパチスロが融合する理由", "7/9")

    principles = [
        (C_RED,   "① 増殖分裂という「目に見える爆発感」",
         "10G→20G→40G→80Gと分裂表示を\n"
         "リアルタイムで見せることで\n"
         "「増えていく！」の快感を直接体験させる。\n"
         "数字の可視化が興奮を加速させる設計。"),
        (C_RED2,  "② HP制サバイバル演出の緊張→解放フロー",
         "CZ中のHP管理演出がスロットに\n"
         "ゲーム的緊張感を持ち込む。\n"
         "瀕死からの逆転サバイバルで\n"
         "達成感が通常演出の数倍に跳ね上がる。"),
        (C_ORG,   "③ 5号機完全再現という「懐かしさ×新鮮さ」",
         "5号機時代の人気機種をスマスロ化。\n"
         "往年のプレイヤーへのノスタルジー。\n"
         "一方で増殖システムは当時にない新要素。\n"
         "懐かしさと驚きの二重設計。"),
        (C_CYAN,  "④ ウェスカー＝危機という直感的演出設計",
         "バイオ原作の「ウェスカー＝最強の脅威」を\n"
         "ゾーン名に転用し期待度の高さを直感で伝える。\n"
         "原作知識がなくてもウェスカーゾーンは\n"
         "見た目の演出から「強い」と認識できる。"),
        (C_GOLD,  "⑤ 2層AT構造による「もう一段上」の設計",
         "ハザードラッシュ（通常AT）でも\n"
         "プレミアムHR（上位AT）という\n"
         "\"次のステージ\"が常に存在する。\n"
         "目標の連鎖設計がプレイ継続を促す。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1200000)
    gy = Inches(0.10)

    positions = [
        (Inches(0.28),  Inches(0.72)),
        (Inches(5.17),  Inches(0.72)),
        (Inches(0.28),  Inches(0.72) + bh_p + gy),
        (Inches(5.17),  Inches(0.72) + bh_p + gy),
        (Inches(0.28),  Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1100000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(330000), pw - Emu(90000), ph - Emu(390000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "面白さの核心：「増殖分裂の可視化×HP演出の緊張感×ノスタルジー×ウェスカー期待度×2層AT目標」の5要素",
           "バイオハザードというIPの「恐怖・緊張・サバイバル」をパチスロのCZ/AT体験に直接落とし込んだ傑作")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN, "増殖分裂という直感的・視覚的上乗せ演出",
         "10→20→40→80Gの分裂表示で\n"
         "出玉が増える瞬間を視覚的に体験できる。\n"
         "数字が増えていく快感は普遍的な楽しさ。"),
        (C_GREEN, "5号機の完全再現によるIP活用の巧みさ",
         "当時の人気機種を忠実に再現しつつ\n"
         "スマスロの純増スペックを活用。\n"
         "既存ファン層と新規層の両方に訴求可能。"),
        (C_GREEN, "2層AT設計で常に上位状態への目標がある",
         "ハザードラッシュからプレミアムHRへ\n"
         "という明確な昇格目標。\n"
         "AT中に次の目標が常に見えている設計。"),
        (C_GREEN, "HP制サバイバル演出によるゲーム感の強化",
         "CZ中のサバイバル体験が通常演出と差別化。\n"
         "原作の世界観とスロット演出が自然融合。\n"
         "没入感と達成感が一般演出より高い。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1140000)
        rect_b(s, lx, iy, lw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(300000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,   "純増2.5枚/Gは現スマスロ世代には低め",
         "スマスロ全般の高純増（5〜8枚）に慣れた\n"
         "プレイヤーにはAT中の出玉速度が遅く感じる。\n"
         "「高純増に慣れた世代には刺さらない」との声も。"),
        (C_RED2,  "ユーザー評価の低さ（1.74/5.0）",
         "DMMぱちタウン等での評価が低め。\n"
         "完成度は評価されるも「時代のニーズに合わない」\n"
         "という指摘が目立つ。稼働定着が課題。"),
        (C_ORG,   "増殖の分かりやすさと複雑さのバランス",
         "増殖・変異・上乗せ状態と仕組みが多岐に渡り\n"
         "初心者には理解しにくい面もある。\n"
         "解析情報を事前に調べておくことが重要。"),
        (C_GRAY,  "セットストック概念の複雑さ",
         "上位AT昇格の鍵となるセットストックだが\n"
         "その仕組みを理解しないと立ち回りが困難。\n"
         "ヘビーユーザー向けの設計にとどまっている。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1140000)
        rect_b(s, rx, iy, rw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(300000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "総評：設計完成度は高く5号機再現として秀逸だが、純増2.5枚/Gというスペックが現代スマスロ市場での壁",
           "「完成度は高いが時代のニーズに合うかは難しい」という評価が示す通り、IP活用と時代感の両立が今後の課題")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x10, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "スマスロ バイオハザード5 ── 設計的強み総括", 10, bold=True, color=C_ORG)

    strengths = [
        (C_RED,   "増殖分裂という核心システム",
         "10→20→40→80Gという分裂連鎖が\n"
         "出玉増加を直感的に体験させる。\n"
         "最大810G上乗せの夢が常に存在する設計。"),
        (C_ORG,   "ホラー×サバイバル演出の融合",
         "HP制CZ・ウェスカーゾーンという\n"
         "原作IPを活かした世界観整合設計。\n"
         "緊張→解放フローがAT体験の質を高める。"),
        (C_CYAN,  "2層AT設計による目標の連鎖",
         "ハザードラッシュ→プレミアムHRという\n"
         "明確な昇格構造が離席を防ぎ\n"
         "「もう少し回そう」という動機を作る。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1210000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(325000), lw - Emu(90000), Emu(780000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_ORG, font=FONT_H)

    principles = [
        (C_RED,   "爆発感は「見える化」で体験させよ",
         "増殖の分裂表示は数字増加を\nリアルタイムで見せる最良の爆発感設計"),
        (C_RED2,  "IPの感情構造をゲーム体験に転用せよ",
         "バイオ5の「恐怖→サバイバル→達成」を\nCZ体験に直接落とし込んだ設計の秀逸さ"),
        (C_ORG,   "常に「上の状態」を見せて目標を作れ",
         "2層AT（通常→上位）という\n常在する昇格目標がリプレイ意欲を生む"),
        (C_CYAN,  "懐かしさと新しさを融合させよ",
         "5号機完全再現×増殖分裂（新要素）の\n組み合わせが旧来ファンと新規を同時に引く"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(760000)
        rect_b(s, rx, py0, rw, Emu(710000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(710000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(230000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(285000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3310000), rw, Emu(1060000),
           RGBColor(0x1A, 0x04, 0x04), C_RED, 2.0)
    rect(s, rx, ry + Emu(3310000), Emu(40000), Emu(1060000), C_RED)
    tb(s, rx + Emu(65000), ry + Emu(3360000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_RED)
    tb(s, rx + Emu(65000), ry + Emu(3620000), rw - Emu(85000), Emu(690000),
       "増殖×分裂上乗せ×ホラー演出×2層AT設計の統合機種。\n"
       "5号機バイオ5を忠実に再現しパチスロアワード2025に\n"
       "ノミネートされた注目作。純増2.5枚/Gという\n"
       "スペックが課題も設計の完成度は高い一台。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "設計思想：「増殖可視化・IP世界観転用・2層AT目標・懐古×新規融合」の4原則を次世代機設計に活用せよ",
           "パチスロアワード2025ノミネートが示す通り、設計的注目度は高い。純増スペックの時代適合が今後の鍵")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)     # 1: タイトル・スペック・3ポイント
    s_flow(prs)      # 2: ゲームフロー全体図
    s_normal(prs)    # 3: 通常時の遊び方
    s_cz(prs)        # 4: CZ/前兆の仕組み
    s_at(prs)        # 5: AT/出玉の伸ばし方
    s_upper(prs)     # 6: 上位ATへの道
    s_design(prs)    # 7: 面白さの設計
    s_pros_cons(prs) # 8: 良い点と課題
    s_matome(prs)    # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
