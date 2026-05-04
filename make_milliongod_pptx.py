"""
スマスロ ミリオンゴッド-神々の軌跡- 機種分析資料
出力: proposals/機種分析/ミリオンゴッド/milliongod_analysis.pptx
テーマ: 深紺 × 金 × 赤（GOD世界観）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "ミリオンゴッド", "milliongod_analysis.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深紺×金×赤）───────────────────────────────
C_BG    = RGBColor(0x04, 0x06, 0x18)   # 深紺
C_CARD  = RGBColor(0x0A, 0x0E, 0x28)
C_CARD2 = RGBColor(0x12, 0x16, 0x34)
C_ROW   = RGBColor(0x0E, 0x12, 0x2C)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金（GOD色）
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)   # 明るい金
C_RED   = RGBColor(0xCC, 0x22, 0x11)   # 赤（SGG色）
C_CRIM  = RGBColor(0xFF, 0x44, 0x22)   # 明るい赤橙
C_YEL   = RGBColor(0xFF, 0xCC, 0x00)   # 黄（Z-GAME色）
C_PUR   = RGBColor(0x88, 0x44, 0xCC)   # 紫（プレミア）
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)
C_TEAL  = RGBColor(0x22, 0xAA, 0x99)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 6, 24))
    draw = ImageDraw.Draw(img)
    # 斜めライン
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(8, 10, 32), width=1)
    # 下部の金グロー
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        r = int(40 * t)
        g = int(28 * t)
        b = int(0)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    # 上部薄暗化
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_RED)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GOLD, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_RED)


def note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x02, 0x04, 0x12))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_RED)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, RGBColor(0x80, 0x60, 0x10))

    tb(s, Inches(0.22), Inches(0.52), Inches(5.0), Emu(330000),
       "機種分析資料", 12, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(1.02), Inches(5.1), Emu(900000),
       "スマスロ\nミリオンゴッド\n-神々の軌跡-", 28, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(0.22), Inches(3.1), Inches(5.0), Emu(330000),
       "── 4号機GODの魂が7.0枚純増で甦った", 11, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.65), Inches(4.9), Emu(230000),
       "メーカー：ユニバーサルエンターテインメント　　導入：2026年4月20日", 9, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.97), Inches(4.9), Emu(230000),
       "設定：1〜6段階　　AT純増：約7.0枚/G", 9, color=C_GRAY)
    tb(s, Inches(0.22), Inches(4.29), Inches(4.9), Emu(230000),
       "AT1セット：50G　　最大ループ率：80%", 9, color=C_GRAY)

    # 右：3つのキーワード
    kws = [
        (C_GOLD, "GOD GAME",       "50G×最大80%ループ\n黄7でZ-GAME連鎖"),
        (C_RED,  "SGG（赤7揃い）",  "75%以上ループ\n10G〜100G増加区間+3G復活"),
        (C_PUR,  "GOD揃い",         "1/16384のプレミア\nGGストック4個以上確定"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.7 + i * 1.55)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.25), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.25), ac)
        tb(s, Inches(5.85), y0 + Emu(55000), Inches(3.8), Emu(320000),
           kw, 13, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(370000), Inches(3.8), Emu(420000),
           desc, 8.5, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: スペック
# ══════════════════════════════════════════════════════════════
def s_spec(prs):
    s = new_slide(prs)
    hdr(s, "スペック ── 基本数値", "2/7")

    bx, by = Inches(0.3), Inches(0.78)
    cols_w = [Emu(520000), Emu(1200000), Emu(1280000)]
    col_labels = ["設定", "機械割", "特記"]
    rows = [
        ("1", "—",       ""),
        ("2", "—",       ""),
        ("3", "—",       ""),
        ("4", "—",       ""),
        ("5", "—",       ""),
        ("6", "約111%",  "設定6のみ公表"),
    ]
    row_h = Emu(370000)
    hdr_h = Emu(360000)

    rect(s, bx, by, sum(cols_w), hdr_h, RGBColor(0x70, 0x40, 0x00))
    rx = bx
    for cw, label in zip(cols_w, col_labels):
        tb(s, rx + Emu(30000), by + Emu(45000), cw - Emu(50000), hdr_h - Emu(55000),
           label, 8.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, wrap=False)
        rx += cw

    for i, row in enumerate(rows):
        ry = by + hdr_h + i * row_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, bx, ry, sum(cols_w), row_h, bg)
        rx = bx
        hi = row[0] == "6"
        for j, (cw, val) in enumerate(zip(cols_w, row)):
            col = C_GOLD if j == 0 and hi else (C_GOLD2 if j == 1 and hi else C_WHITE)
            bold = j == 0 or (j == 1 and hi)
            tb(s, rx + Emu(30000), ry + Emu(50000), cw - Emu(50000), row_h - Emu(65000),
               val, 8.5, bold=bold, color=col, align=PP_ALIGN.CENTER, wrap=False)
            rx += cw

    # 右：KVカード6個
    rx2, ry2 = Inches(4.0), Inches(0.78)
    kv = [
        ("AT純増",   "約7.0枚/G（現行トップクラス）",      C_GOLD),
        ("ATセット", "1セット50G・最大80%ループ",           C_GOLD),
        ("SGG",      "赤7揃いで突入・75%以上ループ",        C_RED),
        ("Z-GAME",   "黄7揃いで突入・G数上乗せチェーン",    C_YEL),
        ("GOD揃い",  "1/16384・GGストック4個以上確定",      C_PUR),
        ("期待枚数", "GOD揃い時3000枚以上",                 C_WHITE),
    ]
    for i, (key, val, ac) in enumerate(kv):
        ry3 = ry2 + i * Emu(530000)
        rect_b(s, rx2, ry3, Inches(5.7), Emu(485000), C_CARD, ac, 1.2)
        rect(s, rx2, ry3, Emu(40000), Emu(485000), ac)
        tb(s, rx2 + Emu(70000), ry3 + Emu(40000), Inches(2.0), Emu(210000),
           key, 8, bold=True, color=ac)
        tb(s, rx2 + Emu(70000), ry3 + Emu(250000), Inches(5.2), Emu(210000),
           val, 9, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: ゲームフロー
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー ── モード → GOD GAME → SGG → Z-GAME → GOD揃い", "3/7")

    # 上段：通常時3モード
    rect(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(280000), C_CARD2)
    tb(s, Inches(0.45), Inches(0.74), Inches(3.0), Emu(250000),
       "通常時 ── 3つのモード", 8.5, bold=True, color=C_GOLD)
    modes = [
        ("通常モード",    "GG当選が重い\n最も長い"),
        ("チャンスモード", "GG当選率UP\nチャンス役で短縮"),
        ("天国モード",    "GG当選ほぼ確定\n最速でAT突入"),
    ]
    mw = Inches(9.4) / 3
    for i, (mt, md) in enumerate(modes):
        mx = Inches(0.3) + i * mw
        bc = C_GOLD if i == 2 else (C_CRIM if i == 1 else C_LTGRY)
        rect_b(s, mx + Emu(30000), Inches(1.04), mw - Emu(50000), Emu(700000),
               C_CARD, bc, 1.2)
        tb(s, mx + Emu(60000), Inches(1.07), mw - Emu(90000), Emu(270000),
           mt, 8.5, bold=True, color=bc if i >= 1 else C_WHITE,
           align=PP_ALIGN.CENTER, wrap=False)
        tb(s, mx + Emu(60000), Inches(1.35), mw - Emu(90000), Emu(330000),
           md, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 下段フロー4ボックス
    boxes = [
        (C_CARD2,                        C_GOLD, "GOD GAME\n(GG)",  "50G/最大80%ループ\n黄7→Z-GAME"),
        (C_CARD2,                        C_RED,  "SGG",             "赤7揃いで突入\n75%以上ループ\n増加区間10〜100G"),
        (RGBColor(0x14, 0x10, 0x02),     C_YEL,  "Z-GAME",          "黄7揃いで突入\nG数上乗せチェーン\n複数連鎖で爆発"),
        (RGBColor(0x0A, 0x04, 0x18),     C_PUR,  "GOD揃い",         "1/16384\nGGストック4個以上確定\n期待3000枚以上"),
    ]
    bw, bh = Inches(1.8), Inches(1.4)
    gap = Inches(0.28)
    total = 4 * bw + 3 * gap
    sx = (Inches(10) - total) / 2
    cy = Inches(3.85)

    for i, (fill, bc, lbl, sub) in enumerate(boxes):
        bx0 = sx + i * (bw + gap)
        rect_b(s, bx0, cy - bh / 2, bw, bh, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), cy - bh / 2 + Emu(80000),
           bw - Emu(80000), Emu(380000), lbl, 10, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), cy - bh / 2 + Emu(450000),
           bw - Emu(60000), Emu(280000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw + Emu(10000), cy, col=C_GOLD)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: GOD GAME × SGG × Z-GAME 核心
# ══════════════════════════════════════════════════════════════
def s_core(prs):
    s = new_slide(prs)
    hdr(s, "AT構成 ── GOD GAME × SGG × Z-GAME × GOD揃いの4層設計", "4/7")

    lx = Inches(0.28)
    rx = Inches(5.05)
    cw = Inches(4.5)
    ch = Emu(2100000)
    top_y = Inches(0.72)
    bot_y = top_y + ch + Emu(80000)
    ch2 = Emu(2100000)

    # 左上：GG
    rect_b(s, lx, top_y, cw, ch, C_CARD, C_GOLD, 1.5)
    rect(s, lx, top_y, Emu(45000), ch, C_GOLD)
    tb(s, lx + Emu(75000), top_y + Emu(45000), cw - Emu(100000), Emu(260000),
       "GG（GOD GAME）の仕組み", 10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx + Emu(75000), top_y + Emu(300000), cw - Emu(100000), ch - Emu(360000),
       "1セット50G、純増約7.0枚/G\n"
       "消化中の役でGGストック獲得\n"
       "最大80%ループで連続GGを目指す\n"
       "黄7揃いでZ-GAME突入（上乗せモード）",
       8.5, color=C_WHITE)

    # 左下：GOD揃い
    rect_b(s, lx, bot_y, cw, ch2, RGBColor(0x0A, 0x04, 0x18), C_PUR, 1.5)
    rect(s, lx, bot_y, Emu(45000), ch2, C_PUR)
    tb(s, lx + Emu(75000), bot_y + Emu(45000), cw - Emu(100000), Emu(260000),
       "GOD揃いの衝撃", 10, bold=True, color=C_PUR, font=FONT_H)
    tb(s, lx + Emu(75000), bot_y + Emu(300000), cw - Emu(100000), ch2 - Emu(360000),
       "確率1/16384の最高プレミア演出\n"
       "GGストック4個以上確定（継続確定）\n"
       "さらにループストックを高確率で獲得\n"
       "期待枚数3000枚以上の大爆発",
       8.5, color=C_WHITE)

    # 右上：SGG
    rect_b(s, rx, top_y, cw, ch, C_CARD, C_RED, 1.5)
    rect(s, rx, top_y, Emu(45000), ch, C_RED)
    tb(s, rx + Emu(75000), top_y + Emu(45000), cw - Emu(100000), Emu(260000),
       "SGG（スーパーGOD GAME）", 10, bold=True, color=C_RED, font=FONT_H)
    tb(s, rx + Emu(75000), top_y + Emu(300000), cw - Emu(100000), ch - Emu(360000),
       "赤7揃いで突入するセット管理型AT\n"
       "75%以上の高ループ率で継続\n"
       "出玉増加区間（10G〜100G）+ 3G復活ゾーンの2部構成\n"
       "復活ゾーンでさらなるストック獲得の可能性",
       8.5, color=C_WHITE)

    # 右下：Z-GAME
    rect_b(s, rx, bot_y, cw, ch2, RGBColor(0x14, 0x10, 0x02), C_YEL, 1.5)
    rect(s, rx, bot_y, Emu(45000), ch2, C_YEL)
    tb(s, rx + Emu(75000), bot_y + Emu(45000), cw - Emu(100000), Emu(260000),
       "Z-GAME", 10, bold=True, color=C_YEL, font=FONT_H)
    tb(s, rx + Emu(75000), bot_y + Emu(300000), cw - Emu(100000), ch2 - Emu(360000),
       "黄7揃いでGG消化中に割り込む上乗せモード\n"
       "消化中にさらなる黄7揃いでチェーン連鎖\n"
       "「黄7が止まらない」体験が最大の興奮\n"
       "GG消化をより「速く・大きく」するための増幅装置",
       8.5, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: ゲーム体験の核心
# ══════════════════════════════════════════════════════════════
def s_experience(prs):
    s = new_slide(prs)
    hdr(s, "ゲーム体験の核心 ── 積み上げる快感と瞬間の衝撃が共存する設計", "5/7")

    # ── 上段：5ステップ体験フロー ─────────────────────────────
    bw = Inches(1.60)
    gap = Inches(0.36)
    bh = Emu(1380000)
    sx0 = Inches(0.20)
    flow_y = Inches(0.72)
    cy = flow_y + bh // 2

    steps = [
        (C_CARD2,                      C_GOLD, "AT突入",
         "通常→GG入場\nモードで期待度が変わる"),
        (C_CARD2,                      C_GOLD, "GG消化",
         "黄7を待ちながら\nZ-GAME連鎖を狙う"),
        (RGBColor(0x20, 0x06, 0x04),   C_RED,  "赤7揃い",
         "SGG突入の衝撃\n増加区間の爆発が始まる"),
        (RGBColor(0x14, 0x10, 0x02),   C_YEL,  "黄7揃い",
         "Z-GAME発動\n連鎖するほど上乗せ爆発"),
        (RGBColor(0x0A, 0x04, 0x18),   C_PUR,  "GOD揃い",
         "1/16384の奇跡\n世界が変わる瞬間"),
    ]
    for i, (fill, ac, title, desc) in enumerate(steps):
        bx = sx0 + i * (bw + gap)
        rect_b(s, bx, flow_y, bw, bh, fill, ac, 1.5)
        tb(s, bx + Emu(40000), flow_y + Emu(60000), bw - Emu(60000), Emu(380000),
           title, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(35000), flow_y + Emu(460000), bw - Emu(55000), Emu(820000),
           desc, 8, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw + Emu(80000), cy, col=C_GOLD)

    # ── 下段左：自力感の設計 ──────────────────────────────────
    lx = Inches(0.28)
    ly = flow_y + bh + Emu(120000)
    lw = Inches(4.5)
    lh = Emu(2650000)

    rect_b(s, lx, ly, lw, lh, C_CARD, C_GOLD, 1.5)
    rect(s, lx, ly, Emu(45000), lh, C_GOLD)
    tb(s, lx + Emu(75000), ly + Emu(45000), lw - Emu(100000), Emu(260000),
       "自力感の設計", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx + Emu(75000), ly + Emu(300000), lw - Emu(100000), lh - Emu(360000),
       "7.0枚/Gという現行最速クラスの純増が\n「GGが続く快感」を最大化する。\n\n"
       "GGストックの積み上げ→消化→再積み上げ\nというループ構造が来店継続を促す。\n\n"
       "Z-GAMEの黄7連鎖は「当たり続けるほど加速する」\n正のフィードバックループを生む設計。\n\n"
       "SGGの3G復活ゾーンが「諦めなくてよかった」\n体験を一定頻度で提供する。",
       8, color=C_WHITE)

    # ── 下段右：GOD揃いという設計的頂点 ─────────────────────
    rx = Inches(5.0)
    rw = Inches(4.7)

    rect_b(s, rx, ly, rw, lh, RGBColor(0x0A, 0x04, 0x18), C_PUR, 1.5)
    rect(s, rx, ly, Emu(45000), lh, C_PUR)
    tb(s, rx + Emu(75000), ly + Emu(45000), rw - Emu(100000), Emu(260000),
       "GOD揃いという設計的頂点", 11, bold=True, color=C_PUR, font=FONT_H)
    tb(s, rx + Emu(75000), ly + Emu(300000), rw - Emu(100000), lh - Emu(360000),
       "1/16384という絶対的レアリティが\nGOD揃いを語り草にする。\n\n"
       "打った人間にとって「あの日GOD揃いが出た」は\n一生の話題になる体験。\n\n"
       "この確率設計は:\n"
       "① レアさが語り継がれる（UGC生産）\n"
       "② 期待値ゼロでも「ありえる」という希望\n"
       "③ 出た動画がバズる（SNS拡散）\n"
       "という3重の価値を持つ。",
       8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 設定判別
# ══════════════════════════════════════════════════════════════
def s_hanbet(prs):
    s = new_slide(prs)
    hdr(s, "設定判別 ── 実戦で使えるポイント", "6/7")

    cols_x = [Inches(0.28), Inches(3.48), Inches(6.68)]
    cols_w = [Inches(3.0), Inches(3.0), Inches(3.0)]
    col_hdrs = ["GG初当たり率", "SGG発生率", "Z-GAME発生率"]
    col_colors = [C_GOLD, C_RED, C_YEL]
    contents = [
        [
            ("高設定ほどGG初当たりが早い",
             "高設定ほどGG初当たりが早い\n通常時のゲーム数を記録して判断\n複数回測定で精度UP"),
            ("通常時G数を記録",
             "複数セッションにわたって\nGG当選G数を記録する\n平均値で高設定かを判断"),
            ("モード示唆を見逃さない",
             "チャンス役後のモード移行や\n天国滞在頻度は設定差あり\n短縮ゲーム数にも注目"),
        ],
        [
            ("設定差あり・高設定ほどSGG突入率が高い",
             "赤7揃い頻度を記録\nGG消化中の挙動に注目\n設定差あり"),
            ("赤7揃い頻度を記録",
             "GG中の赤7出現頻度を\n複数GGにわたって記録する\n高設定ほど頻度が高い傾向"),
            ("増加区間の長さに注目",
             "SGG突入後の増加区間（10〜100G）の\n長さにも設定差が出る可能性\n実戦データ蓄積で傾向把握"),
        ],
        [
            ("黄7揃い頻度に設定差",
             "黄7揃い頻度に設定差\n高設定ほどZ-GAME連鎖が発生しやすい\n実戦データ蓄積で傾向把握"),
            ("連鎖回数を記録",
             "Z-GAME発動後の連鎖回数を記録\n高設定ほど長い連鎖が\n発生しやすい可能性"),
            ("GG中の黄7出現G数",
             "GG1セット中に黄7が出るまでの\nG数も指標になりえる\n積み重ねで判断精度が上がる"),
        ],
    ]

    for ci, (col_x, col_w, col_hdr, col_col, items) in enumerate(
            zip(cols_x, cols_w, col_hdrs, col_colors, contents)):
        rect(s, col_x, Inches(0.72), col_w - Inches(0.12), Emu(360000), col_col)
        tb(s, col_x + Emu(30000), Inches(0.72) + Emu(45000),
           col_w - Inches(0.17), Emu(275000),
           col_hdr, 9.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER, wrap=False)

        for ri, (title, body) in enumerate(items):
            ry0 = Inches(0.72) + Emu(360000) + ri * Emu(1270000)
            bg = C_CARD if ri % 2 == 0 else C_ROW
            rect_b(s, col_x, ry0, col_w - Inches(0.12), Emu(1210000), bg, col_col, 0.5)
            tb(s, col_x + Emu(50000), ry0 + Emu(55000), col_w - Inches(0.2), Emu(255000),
               title, 8.5, bold=True, color=col_col)
            tb(s, col_x + Emu(50000), ry0 + Emu(305000), col_w - Inches(0.2), Emu(780000),
               body, 8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: まとめ
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "7/7")

    bx, by = Inches(0.28), Inches(0.72)
    bw3 = Inches(4.5)

    rect(s, bx, by, bw3, Emu(300000), RGBColor(0x80, 0x10, 0x05))
    tb(s, bx + Emu(60000), by + Emu(50000), bw3 - Emu(80000), Emu(230000),
       "長期稼働を支えた3要素", 11, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_GOLD, "① 純増7.0枚/GのAT純増設計",
         "現行機トップクラスの純増速度が\n「速く大きく勝つ」体験を実現。\n4号機GODの「爆発力」をスマスロで再現。"),
        (C_RED,  "② SGG×Z-GAME×GOD揃いの多層設計",
         "GG→SGG→Z-GAMEという昇格構造が\n常に「次がある」希望を生む。\nGOD揃い1/16384が神話的体験を定期的に生産。"),
        (C_PUR,  "③ IP力 × 世代記憶",
         "4号機ミリオンゴッドを知る世代が\nスマスロ版でGOD揃いを再び体験する。\n記憶との接続がリピート来店を促す。"),
    ]
    for i, (ac, t, b) in enumerate(elems):
        ey = by + Emu(300000) + i * Emu(1270000)
        rect_b(s, bx, ey, bw3, Emu(1200000), C_CARD, ac, 1.5)
        rect(s, bx, ey, Emu(45000), Emu(1200000), ac)
        tb(s, bx + Emu(75000), ey + Emu(50000), bw3 - Emu(95000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, bx + Emu(75000), ey + Emu(305000), bw3 - Emu(95000), Emu(800000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(210000),
       "設計原則", 11, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_GOLD, "純増7.0枚は「速さ」を差別化の武器にする"),
        (C_RED,  "多層昇格（GG→SGG→Z-GAME）が目標を途切れさせない"),
        (C_PUR,  "1/16384のGOD揃いが語り継がれる体験を生む"),
        (C_YEL,  "黄7連鎖というポジティブフィードバックループ"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(280000) + i * Emu(540000)
        rect(s, rx, py0, Emu(20000), Emu(490000), ac)
        tb(s, rx + Emu(50000), py0 + Emu(75000), rw - Emu(60000), Emu(380000),
           p, 8.5, bold=False, color=C_WHITE)

    rect_b(s, rx, ry + Emu(2450000), rw, Emu(800000),
           RGBColor(0x10, 0x08, 0x02), C_GOLD, 1.5)
    tb(s, rx + Emu(55000), ry + Emu(2500000), rw - Emu(75000), Emu(260000),
       "総括", 9, bold=True, color=C_GOLD)
    tb(s, rx + Emu(55000), ry + Emu(2760000), rw - Emu(75000), Emu(430000),
       "IP×純増7枚×多層昇格設計の完成形。\n"
       "63件の投稿データが示す稼働の強さが台の完成度を証明している。",
       8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)
    s_spec(prs)
    s_flow(prs)
    s_core(prs)
    s_experience(prs)
    s_hanbet(prs)
    s_matome(prs)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
