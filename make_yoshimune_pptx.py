"""
L真打吉宗 機種分析資料  （大都技研・2026年4月6日導入）
出力: proposals/機種分析/真打吉宗/yoshimune_analysis.pptx
テーマ: 和黒 × 紫 × 橙（吉宗世界観）
"""
import io
import os
import sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "真打吉宗", "yoshimune_analysis.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（和黒×紫×橙）───────────────────────────────
C_BG    = RGBColor(0x08, 0x04, 0x14)   # 和黒
C_CARD  = RGBColor(0x10, 0x08, 0x20)
C_CARD2 = RGBColor(0x18, 0x10, 0x2C)
C_ROW   = RGBColor(0x14, 0x0C, 0x24)
C_PUR   = RGBColor(0x88, 0x22, 0xCC)   # 紫（吉宗メインカラー）
C_PUR2  = RGBColor(0xAA, 0x55, 0xFF)   # 明るい紫
C_ORG   = RGBColor(0xFF, 0x77, 0x11)   # 橙（真BB色）
C_ORG2  = RGBColor(0xFF, 0xAA, 0x44)   # 明るい橙
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)
C_RED   = RGBColor(0xCC, 0x22, 0x22)
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (8, 4, 20))
    draw = ImageDraw.Draw(img)
    # 斜めライン
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(12, 8, 28), width=1)
    # 下部の紫グロー
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        draw.line([(0, y), (w, y)], fill=(int(25 * t), 0, int(35 * t)))
    # 上部薄暗化
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_PUR)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_PUR)


def note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x06, 0x02, 0x10))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_PUR)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, RGBColor(0x60, 0x10, 0x99))

    tb(s, Inches(0.22), Inches(0.52), Inches(5.0), Emu(330000),
       "機種分析資料", 12, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(1.02), Inches(5.1), Emu(900000),
       "L真打吉宗", 36, bold=True, color=C_ORG, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.9), Inches(5.0), Emu(330000),
       "── 1G連×純増9枚で4号機の魂が甦る", 11, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.5), Inches(4.9), Emu(230000),
       "メーカー: 大都技研　　導入: 2026年4月6日", 9, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.82), Inches(4.9), Emu(230000),
       "設定: 1〜6段階　　BB純増: 約9.0枚/G", 9, color=C_GRAY)
    tb(s, Inches(0.22), Inches(4.14), Inches(4.9), Emu(230000),
       "1G連ループ: 真BB消化中に成立役で抽選", 9, color=C_GRAY)

    # 右：3つのキーワード
    kws = [
        (C_PUR,  "勧善懲悪RUSH",    "CZ経由でAT突入\n真BBで出玉爆発"),
        (C_ORG,  "真BB（純増9.0枚）", "1G連ループで連続BB\n究極鷹ブレイクで5000枚超"),
        (C_GOLD, "究極鷹ブレイク",   "毎G1000枚ループ上乗せ\n実質5000枚以上確定"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.7 + i * 1.55)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.25), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.25), ac)
        tb(s, Inches(5.85), y0 + Emu(55000), Inches(3.8), Emu(320000),
           kw, 13, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(370000), Inches(3.8), Emu(420000),
           desc, 8.5, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: スペック
# ══════════════════════════════════════════════════════════════
def s_spec(prs):
    s = new_slide(prs)
    hdr(s, "スペック ── 基本数値", "2/7")

    bx, by = Inches(0.3), Inches(0.78)
    cols_w = [Emu(520000), Emu(1020000), Emu(980000), Emu(1280000)]
    col_labels = ["設定", "機械割", "BB初当り", "特記"]
    rows = [
        ("1", "—",      "—", ""),
        ("2", "—",      "—", ""),
        ("3", "—",      "—", ""),
        ("4", "—",      "—", ""),
        ("5", "—",      "—", ""),
        ("6", "約110%", "—", "設定6確定で高機械割"),
    ]
    row_h = Emu(370000)
    hdr_h = Emu(360000)

    rect(s, bx, by, sum(cols_w), hdr_h, RGBColor(0x55, 0x10, 0x88))
    rx = bx
    for cw, label in zip(cols_w, col_labels):
        tb(s, rx + Emu(30000), by + Emu(45000), cw - Emu(50000), hdr_h - Emu(55000),
           label, 8.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, wrap=False)
        rx += cw

    for i, row in enumerate(rows):
        ry = by + hdr_h + i * row_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, bx, ry, sum(cols_w), row_h, bg)
        rx = bx
        hi = row[0] == "6"
        for j, (cw, val) in enumerate(zip(cols_w, row)):
            col = C_PUR2 if j == 0 and hi else (C_GOLD if j == 1 and hi else C_WHITE)
            bold = j == 0 or (j == 1 and hi)
            tb(s, rx + Emu(30000), ry + Emu(50000), cw - Emu(50000), row_h - Emu(65000),
               val, 8.5, bold=bold, color=col, align=PP_ALIGN.CENTER, wrap=False)
            rx += cw

    # 右：KVカード
    rx2, ry2 = Inches(4.6), Inches(0.78)
    kv = [
        ("真BB純増",         "約9.0枚/G（現行最高クラス）",        C_ORG),
        ("真BB獲得枚数",     "約2000枚（1セット）",                 C_ORG2),
        ("1G連",             "真BB消化中に成立役→抽選",             C_PUR),
        ("究極鷹ブレイク",   "毎G1000枚ループ",                     C_GOLD),
        ("AT天井",           "1500G",                               C_WHITE),
        ("CZ天井",           "1000G or 6周期",                      C_GRAY),
    ]
    for i, (key, val, ac) in enumerate(kv):
        ry3 = ry2 + i * Emu(530000)
        rect_b(s, rx2, ry3, Inches(5.1), Emu(485000), C_CARD, ac, 1.2)
        rect(s, rx2, ry3, Emu(40000), Emu(485000), ac)
        tb(s, rx2 + Emu(70000), ry3 + Emu(40000), Inches(2.5), Emu(210000),
           key, 8, bold=True, color=ac)
        tb(s, rx2 + Emu(70000), ry3 + Emu(250000), Inches(4.6), Emu(210000),
           val, 9, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: ゲームフロー
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー ── CZモード → 勧善懲悪RUSH → 真BB → 1G連", "3/7")

    # 上段：通常時CZモード
    rect(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(280000), C_CARD2)
    tb(s, Inches(0.45), Inches(0.74), Inches(3.0), Emu(250000),
       "通常時CZモード", 8.5, bold=True, color=C_GOLD)

    modes = [
        ("6周期制",       "CZ発生タイミングが\nモードで管理される"),
        ("天国モード相当", "1周期目でCZ確定"),
        ("鷹CZ",          "最高格CZ"),
    ]
    mw = Inches(9.4) / 3
    for i, (mt, md) in enumerate(modes):
        mx = Inches(0.3) + i * mw
        bc = C_GOLD if i == 0 else (C_PUR if i == 1 else C_ORG)
        rect_b(s, mx + Emu(30000), Inches(1.04), mw - Emu(50000), Emu(700000),
               C_CARD, bc, 1.2)
        tb(s, mx + Emu(60000), Inches(1.07), mw - Emu(90000), Emu(270000),
           mt, 8.5, bold=True, color=bc,
           align=PP_ALIGN.CENTER, wrap=False)
        tb(s, mx + Emu(60000), Inches(1.35), mw - Emu(90000), Emu(330000),
           md, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 下段フロー4ボックス
    boxes = [
        (C_CARD2,                        C_PUR,  "CZ\n「悪人成敗チャンス」",   "CZモードから発生\nAT当選を目指す"),
        (C_CARD2,                        C_PUR2, "AT\n「勧善懲悪RUSH」",        "AT突入\n真BBを目指すメインルート"),
        (RGBColor(0x20, 0x08, 0x04),     C_ORG,  "真BB",                        "純増9.0枚/G\n2000枚獲得\n1G連抽選が走る"),
        (RGBColor(0x18, 0x10, 0x00),     C_GOLD, "1G連 / 究極鷹ブレイク",       "1G連ループで連続BB\n究極鷹ブレイクで5000枚超"),
    ]
    bw, bh = Inches(1.8), Inches(1.4)
    gap = Inches(0.28)
    total = 4 * bw + 3 * gap
    sx = (Inches(10) - total) / 2
    cy = Inches(3.85)

    for i, (fill, bc, lbl, sub) in enumerate(boxes):
        bx0 = sx + i * (bw + gap)
        rect_b(s, bx0, cy - bh / 2, bw, bh, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), cy - bh / 2 + Emu(80000),
           bw - Emu(80000), Emu(380000), lbl, 10, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), cy - bh / 2 + Emu(450000),
           bw - Emu(60000), Emu(280000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw + Emu(10000), cy)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 真BB × 1G連 × 究極鷹ブレイク
# ══════════════════════════════════════════════════════════════
def s_bb(prs):
    s = new_slide(prs)
    hdr(s, "真BB構成 ── 1G連ループ × 究極鷹ブレイクの爆裂設計", "4/7")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)
    half_h = Emu(2150000)
    gap = Emu(80000)
    by2 = ly + half_h + gap
    half_h2 = SLIDE_H - by2 - Emu(180000)

    # 左上：真BBの仕組み
    rect_b(s, lx, ly, lw, half_h, C_CARD, C_ORG, 1.5)
    rect(s, lx, ly, Emu(45000), half_h, C_ORG)
    tb(s, lx + Emu(75000), ly + Emu(45000), lw - Emu(100000), Emu(260000),
       "真BBの仕組み", 11, bold=True, color=C_ORG, font=FONT_H)
    tb(s, lx + Emu(75000), ly + Emu(305000), lw - Emu(100000), half_h - Emu(360000),
       "純増約9.0枚/G（現行機最高クラス）\n"
       "1セットで約2000枚獲得\n"
       "BB消化中、成立役に応じて1G連抽選\n"
       "1G連=次のBBがBB終了の1G後に開始",
       8, color=C_WHITE)

    # 左下：究極鷹ブレイク
    rect_b(s, lx, by2, lw, half_h2, C_CARD, C_GOLD, 1.5)
    rect(s, lx, by2, Emu(45000), half_h2, C_GOLD)
    tb(s, lx + Emu(75000), by2 + Emu(45000), lw - Emu(100000), Emu(260000),
       "究極鷹ブレイク", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx + Emu(75000), by2 + Emu(305000), lw - Emu(100000), half_h2 - Emu(360000),
       "本機最強の特化ゾーン\n"
       "毎ゲーム1000枚のループ上乗せが発生\n"
       "実質5000枚以上確定\n"
       "「天文学的な出玉」体験を提供する最高峰",
       8, color=C_WHITE)

    # 右上：1G連の衝撃体験
    rx = Inches(5.0)
    rw = Inches(4.7)
    rect_b(s, rx, ly, rw, half_h, C_CARD, C_PUR, 1.5)
    rect(s, rx, ly, Emu(45000), half_h, C_PUR)
    tb(s, rx + Emu(75000), ly + Emu(45000), rw - Emu(100000), Emu(260000),
       "1G連の衝撃体験", 11, bold=True, color=C_PUR, font=FONT_H)
    tb(s, rx + Emu(75000), ly + Emu(305000), rw - Emu(100000), half_h - Emu(360000),
       "真BB終了の1G後に即座に次のBBが始まる\n"
       "「終わった！」から「え、また!?」の衝撃体験\n"
       "4号機吉宗の代名詞を現代スマスロで再現\n"
       "1G連が連続すると指数的に興奮が高まる",
       8, color=C_WHITE)

    # 右下：コイン単価と荒波
    rect_b(s, rx, by2, rw, half_h2,
           RGBColor(0x18, 0x04, 0x04), C_RED, 1.5)
    rect(s, rx, by2, Emu(45000), half_h2, C_RED)
    tb(s, rx + Emu(75000), by2 + Emu(45000), rw - Emu(100000), Emu(260000),
       "コイン単価と荒波", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, rx + Emu(75000), by2 + Emu(305000), rw - Emu(100000), half_h2 - Emu(360000),
       "真BB純増9.0枚は速いが荒波も激しい\n"
       "真BBを引けないと消化不良になりやすい\n"
       "設定1は特に初当たりが遠い設計\n"
       "「当たれば爆発・当たらなければ深みにハマる」二極体験",
       8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: ゲーム体験の核心
# ══════════════════════════════════════════════════════════════
def s_experience(prs):
    s = new_slide(prs)
    hdr(s, "ゲーム体験の核心 ── 4号機の記憶と現行最速9枚が生む興奮体験", "5/7")

    # 上段：5ステップ体験フロー
    bw = Inches(1.60)
    gap = Inches(0.36)
    bh = Emu(1380000)
    sx0 = Inches(0.20)
    flow_y = Inches(0.72)
    cy = flow_y + bh // 2

    steps = [
        (C_CARD2,                        C_PUR,   "CZ突入",
         "悪人成敗チャンス\n鷹CZで期待度MAX"),
        (RGBColor(0x14, 0x08, 0x20),     C_PUR2,  "AT突入",
         "勧善懲悪RUSH\n真BBを目指す"),
        (RGBColor(0x20, 0x08, 0x04),     C_ORG,   "真BB発動",
         "純増9.0枚の嵐\n2000枚を爆速獲得"),
        (RGBColor(0x14, 0x0A, 0x00),     C_GOLD,  "1G連！",
         "終わったと思ったら\n次のBBが始まる衝撃"),
        (RGBColor(0x18, 0x12, 0x04),     C_GOLD2, "究極鷹ブレイク",
         "毎G1000枚ループ\n現実が変わる体験"),
    ]
    for i, (fill, ac, title, desc) in enumerate(steps):
        bx = sx0 + i * (bw + gap)
        rect_b(s, bx, flow_y, bw, bh, fill, ac, 1.5)
        tb(s, bx + Emu(40000), flow_y + Emu(60000), bw - Emu(60000), Emu(380000),
           title, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(35000), flow_y + Emu(460000), bw - Emu(55000), Emu(820000),
           desc, 8, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw + Emu(80000), cy)

    # 下段左：4号機記憶との接続
    lx = Inches(0.28)
    ly = flow_y + bh + Emu(120000)
    lw = Inches(4.5)
    lh = Emu(2650000)

    rect_b(s, lx, ly, lw, lh, C_CARD, C_PUR, 1.5)
    rect(s, lx, ly, Emu(45000), lh, C_PUR)
    tb(s, lx + Emu(75000), ly + Emu(45000), lw - Emu(100000), Emu(260000),
       "4号機記憶との接続", 11, bold=True, color=C_PUR, font=FONT_H)
    tb(s, lx + Emu(75000), ly + Emu(300000), lw - Emu(100000), lh - Emu(360000),
       "「1G連」は4号機吉宗の代名詞だった。\n\n"
       "スマスロ版で同じ体験ができるという事実が\n"
       "30〜40代の休眠層を引き戻す力を持つ。\n\n"
       "純増9.0枚という圧倒的なスピードは\n"
       "「あの頃の爆発感」を現代のスペックで再現。\n\n"
       "ノスタルジアと現行最高水準の性能が\n"
       "一台に共存している稀有な設計。",
       8, color=C_WHITE)

    # 下段右：1G連という設計的天才
    rx = Inches(5.0)
    rw = Inches(4.7)

    rect_b(s, rx, ly, rw, lh, RGBColor(0x18, 0x08, 0x00), C_ORG, 1.5)
    rect(s, rx, ly, Emu(45000), lh, C_ORG)
    tb(s, rx + Emu(75000), ly + Emu(45000), rw - Emu(100000), Emu(260000),
       "1G連という設計的天才", 11, bold=True, color=C_ORG, font=FONT_H)
    tb(s, rx + Emu(75000), ly + Emu(300000), rw - Emu(100000), lh - Emu(360000),
       "1G連は単なるボーナス連続ではなく\n"
       "「終わり → 始まり」の体験を1Gで完結させる設計。\n\n"
       "通常のAT継続と違い:\n"
       "① 「終わった」という感情が生まれる（一旦落ちる）\n"
       "② 「また始まった」という驚きが来る（爆上がる）\n"
       "この感情の波が興奮を最大化する。\n\n"
       "究極鷹ブレイクはこの1G連体験の\n"
       "究極形として機能する。",
       8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 設定判別 + 課題
# ══════════════════════════════════════════════════════════════
def s_hanbet(prs):
    s = new_slide(prs)
    hdr(s, "設定判別 ── 実戦で使えるポイント", "6/7")

    cols_x = [Inches(0.28), Inches(3.48), Inches(6.68)]
    cols_w = [Inches(3.0), Inches(3.0), Inches(3.0)]
    col_hdrs = ["CZ確率", "真BB中の1G連率", "終了画面示唆"]
    col_colors = [C_PUR, C_ORG, C_GOLD]
    contents = [
        [
            ("設定差あり",        "設定差あり・高設定ほど\nCZ発生が早い。"),
            ("周期の短さで判断",  "周期の短さで判断する。"),
            ("高設定の目安",      "CZが早く来る台は\n高設定の可能性あり。"),
        ],
        [
            ("高設定優遇",        "高設定ほど1G連抽選優遇。"),
            ("BB後すぐ確認",      "BB後すぐ次BBが始まるか確認。"),
            ("連続1G連",          "1G連が連続するほど\n高設定期待度UP。"),
        ],
        [
            ("終了画面示唆",      "BBやAT終了後の画面に\n設定示唆が出る。"),
            ("鷹の衣装・色",      "鷹の衣装・色で示唆が変化。"),
            ("複数回で精度UP",    "複数回の確認で\n総合的に判断する。"),
        ],
    ]

    for ci, (col_x, col_w, col_hdr, col_col, items) in enumerate(
            zip(cols_x, cols_w, col_hdrs, col_colors, contents)):
        rect(s, col_x, Inches(0.72), col_w - Inches(0.12), Emu(360000), col_col)
        tb(s, col_x + Emu(30000), Inches(0.72) + Emu(45000),
           col_w - Inches(0.17), Emu(275000),
           col_hdr, 9.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER, wrap=False)

        for ri, (title, body) in enumerate(items):
            ry0 = Inches(0.72) + Emu(360000) + ri * Emu(1270000)
            bg = C_CARD if ri % 2 == 0 else C_ROW
            rect_b(s, col_x, ry0, col_w - Inches(0.12), Emu(1210000), bg, col_col, 0.5)
            tb(s, col_x + Emu(50000), ry0 + Emu(55000), col_w - Inches(0.2), Emu(255000),
               title, 8.5, bold=True, color=col_col)
            tb(s, col_x + Emu(50000), ry0 + Emu(305000), col_w - Inches(0.2), Emu(780000),
               body, 8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: まとめ
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "7/7")

    bx, by = Inches(0.28), Inches(0.72)
    bw3 = Inches(4.5)

    rect(s, bx, by, bw3, Emu(300000), RGBColor(0x55, 0x10, 0x88))
    tb(s, bx + Emu(60000), by + Emu(50000), bw3 - Emu(80000), Emu(230000),
       "設計から学べること", 11, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_PUR,  "1G連という設計資産",
         "4号機から引き継いだ1G連は\n"
         "「終わりと始まり」の感情波を1Gで生む。\n"
         "設計的天才として今も機能している。"),
        (C_ORG,  "純増9.0枚の爆発力",
         "現行最高クラスの純増速度が\n"
         "「速く大きく勝つ」体験を実現。\n"
         "究極鷹ブレイクとの組み合わせが頂点体験。"),
        (C_GOLD, "荒波という諸刃の剣",
         "爆発力と引き換えに荒波も激しい。\n"
         "勝つ時の快感が大きい分、負ける時の傷も深い。\n"
         "コアユーザーへの訴求が最大の課題。"),
    ]
    for i, (ac, t, b) in enumerate(elems):
        ey = by + Emu(300000) + i * Emu(1270000)
        rect_b(s, bx, ey, bw3, Emu(1200000), C_CARD, ac, 1.5)
        rect(s, bx, ey, Emu(45000), Emu(1200000), ac)
        tb(s, bx + Emu(75000), ey + Emu(50000), bw3 - Emu(95000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, bx + Emu(75000), ey + Emu(305000), bw3 - Emu(95000), Emu(800000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(210000),
       "設計原則", 11, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_PUR,  "1G連の「終わり→始まり」体験が興奮を最大化する"),
        (C_ORG,  "純増9枚×1G連×究極鷹ブレイクの三位一体が爆発を生む"),
        (C_GOLD, "IP記憶との接続が休眠層を呼び戻す武器になる"),
        (C_GRAY, "荒波設計はコアファン向けに徹した潔い設計判断"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(280000) + i * Emu(540000)
        rect(s, rx, py0, Emu(20000), Emu(490000), ac)
        tb(s, rx + Emu(50000), py0 + Emu(75000), rw - Emu(60000), Emu(380000),
           p, 8.5, color=C_WHITE)

    rect_b(s, rx, ry + Emu(2450000), rw, Emu(800000),
           RGBColor(0x14, 0x06, 0x20), C_PUR, 1.5)
    tb(s, rx + Emu(55000), ry + Emu(2500000), rw - Emu(75000), Emu(260000),
       "総括", 9, bold=True, color=C_PUR)
    tb(s, rx + Emu(55000), ry + Emu(2760000), rw - Emu(75000), Emu(430000),
       "4号機吉宗の遺産を現行最高スペックで昇華した一台。\n"
       "1G連という普遍的な興奮体験は時代を超えて機能し続ける。",
       8, color=C_WHITE)

    note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)
    s_spec(prs)
    s_flow(prs)
    s_bb(prs)
    s_experience(prs)
    s_hanbet(prs)
    s_matome(prs)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
