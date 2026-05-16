"""
スマスロ ミリオンゴッド-神々の軌跡- 包括的ガイド＋機種分析 統合資料
出力: proposals/機種分析/ミリゴ/milliongod_guide_v1.pptx
テーマ: 深紺 × 金 × 赤 × 黄 × 紫（神々しいカラー）

WebSearch確認済み情報:
- 通常時: 6種の表モード(低確A/B・通常・天国準備・天国・超天国)+裏モードでGG当選率が変化
- 天井: AT間最大1480G、恩恵はGGストック1個+ループストック振り分け
- ガイアステージ: GG超高確率状態(9G+α)、奇数揃いでGG確定、0揃いでZ-ZONE突入
- Z-ZONE: 0揃い等で突入する5GのZ-GAME昇格チャレンジゾーン(フルナビ)
- Z-GAME: 黄7揃いでGGストック上乗せ、チェーンで爆発(1/1.4で黄7成立)
- GOD GAME(GG): 1セット50G、純増約7.0枚/G、最大80%ループ
- SGG: 赤7揃い(1/6900)で突入、1セット10〜100G増加区間+3G引き戻しゾーン、75%以上ループ
- GOD揃い: 1/16384、GGストック4セット以上+高継続ループストック、期待3000枚以上
- PREMIUM GOD GAME(PGG): GOD揃いorロングフリーズで突入
- 赤7揃い: 1/6900確定役
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
           "proposals", "機種分析", "ミリゴ", "milliongod_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深紺×金×赤×黄×紫）───────────────────────────
C_BG    = RGBColor(0x04, 0x06, 0x18)   # 深紺
C_CARD  = RGBColor(0x0A, 0x0E, 0x28)
C_CARD2 = RGBColor(0x12, 0x16, 0x34)
C_ROW   = RGBColor(0x0E, 0x12, 0x2C)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金（GOD色）
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)   # 明るい金
C_RED   = RGBColor(0xCC, 0x22, 0x11)   # 赤（SGG色）
C_CRIM  = RGBColor(0xFF, 0x44, 0x22)   # 明るい赤橙
C_YEL   = RGBColor(0xFF, 0xCC, 0x00)   # 黄（Z-GAME色）
C_PUR   = RGBColor(0x88, 0x44, 0xCC)   # 紫（プレミア）
C_PUR2  = RGBColor(0xBB, 0x77, 0xFF)   # 明るい紫
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)
C_TEAL  = RGBColor(0x22, 0xAA, 0x99)
C_DGOLD = RGBColor(0x70, 0x40, 0x00)   # 濃い金（ヘッダ帯）
C_DRED  = RGBColor(0x20, 0x06, 0x04)   # 濃い赤
C_DYEL  = RGBColor(0x14, 0x10, 0x02)   # 濃い黄
C_DPUR  = RGBColor(0x0A, 0x04, 0x18)   # 濃い紫

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景（神殿・深紺グラデーション）─────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 6, 24))
    draw = ImageDraw.Draw(img)
    # 斜めライン（神殿の柱イメージ）
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(8, 10, 32), width=1)
    # 下部の金グロー（神聖な光）
    for y in range(h - 100, h):
        t = (y - (h - 100)) / 100
        r = int(50 * t)
        g = int(35 * t)
        b = int(0)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    # 中央のうっすら光芒
    for x in range(w // 3, 2 * w // 3, 4):
        t = 1 - abs(x - w // 2) / (w // 6)
        if t > 0:
            draw.line([(x, 0), (x, h // 3)], fill=(int(5 * t), int(4 * t), int(10 * t)))
    # 上部薄暗化
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_RED)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GOLD, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_RED)


def net_note(slide):
    tb(slide, Inches(8.0), Inches(5.38), Inches(1.9), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, design_text, sub_text=""):
    """フッター: 設計コメント（太字）＋補足説明"""
    fy = Inches(5.05)
    rect(slide, 0, fy, SLIDE_W, Inches(0.575), C_CARD)
    rect(slide, 0, fy, Emu(35000), Inches(0.575), C_GOLD)
    tb(slide, Inches(0.18), fy + Emu(30000), Inches(6.5), Emu(250000),
       design_text, 8, bold=True, color=C_GOLD)
    if sub_text:
        tb(slide, Inches(0.18), fy + Emu(265000), Inches(7.8), Emu(220000),
           sub_text, 7, color=C_GRAY)
    net_note(slide)


def arrow_r(slide, x, cy, col=None, size=1.0):
    """右向き矢印"""
    ew = int(Emu(200000) * size)
    eh = int(Emu(180000) * size)
    shp = slide.shapes.add_shape(13, x, cy - eh // 2, ew, eh)
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    """下向き矢印"""
    shp = slide.shapes.add_shape(17, cx - Emu(80000), y, Emu(160000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GOLD
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル（タイトル）
    rect(s, 0, 0, Inches(5.3), SLIDE_H, RGBColor(0x02, 0x04, 0x12))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_RED)
    rect(s, Inches(5.3), 0, Emu(10000), SLIDE_H, RGBColor(0x80, 0x60, 0x10))

    tb(s, Inches(0.22), Inches(0.45), Inches(5.0), Emu(280000),
       "スロクリ機種ガイド＋分析資料", 10, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.88), Inches(5.1), Emu(900000),
       "スマスロ\nミリオンゴッド\n-神々の軌跡-", 26, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.92), Inches(5.0), Emu(320000),
       "── 4号機GODの魂が7.0枚純増で甦った", 10, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.42), Inches(4.9), Emu(220000),
       "メーカー：ユニバーサルエンターテインメント　導入：2026年4月", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.68), Inches(4.9), Emu(220000),
       "設定：1〜6段階　　AT純増：約7.0枚/G", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.94), Inches(4.9), Emu(220000),
       "1セット：50G　　最大ループ率：80%　　天井：1480G", 8.5, color=C_GRAY)

    # 右：この台の3ポイント
    tb(s, Inches(5.55), Inches(0.15), Inches(4.3), Emu(280000),
       "この台の3ポイント", 10, bold=True, color=C_GOLD, font=FONT_H)

    kws = [
        (C_GOLD,  "GOD GAME 80%ループ",
         "1セット50G×純増7.0枚/G\n最大80%ループでGGを重ねる\n黄7でZ-GAME突入が爆発のカギ"),
        (C_RED,   "SGG（赤7揃い）1/6900",
         "赤7で突入するセット管理型AT\n75%以上ループ＋引き戻しゾーン付き\nGGとは別軸の強力な出玉源"),
        (C_PUR2,  "GOD揃い 1/16384",
         "全ての上を超えるプレミアフラグ\nGGストック4個以上+大量ループストック\n期待枚数3000枚以上の神体験"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.50 + i * 1.65)
        rect_b(s, Inches(5.55), y0, Inches(4.25), Inches(1.42), C_CARD, ac, 2.0)
        rect(s, Inches(5.55), y0, Emu(55000), Inches(1.42), ac)
        tb(s, Inches(5.75), y0 + Emu(60000), Inches(3.8), Emu(290000),
           kw, 11.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.75), y0 + Emu(360000), Inches(3.9), Emu(600000),
           desc, 8, color=C_WHITE)

    footer(s,
           "設計核心: 4号機GODのIP力×現行最速クラス7.0枚純増×GG→SGG→Z-GAME多層昇格の三位一体設計",
           "2026年4月導入のスマスロ最新作。既存GOD世代の郷愁と新規プレイヤーへの爆発力を両立した完成形。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常 → GG → 上位昇格 → GOD揃いまでの全ルート", "2/10")

    # ─── 上段（通常時 → GG突入ルート）─────────────────────────
    label_y = Inches(0.68)
    rect(s, Inches(0.15), label_y, Inches(9.7), Emu(240000), C_CARD2)
    tb(s, Inches(0.28), label_y + Emu(30000), Inches(4.0), Emu(200000),
       "【通常時】 GG当選ルート", 8.5, bold=True, color=C_GOLD)

    # 通常時ボックス群（横一列）
    row1_y = Inches(0.96)
    row1_h = Emu(1120000)

    boxes_r1 = [
        (C_CARD,  C_LTGRY, "通常モード\n(低確A/B)",
         "6種の表モード\nの最下位層\nGG当選は重め"),
        (C_CARD,  C_CRIM,  "チャンスモード\n(通常・天国準備)",
         "チャンス役で\nモード昇格\nGG期待度UP"),
        (C_CARD,  C_GOLD,  "天国モード\n(天国・超天国)",
         "GG当選ほぼ確定\n最速でAT突入\n最優先滞在を狙う"),
        (C_CARD,  C_TEAL,  "ガイアステージ\n(超高確9G+α)",
         "奇数揃い→GG確定\n0揃い→Z-ZONE突入\n毎G高確率でGG抽選"),
        (C_CARD,  C_GRAY,  "天井\n(1480G)",
         "AT間1480Gで\nGGストック1個\n+ループストック獲得"),
    ]
    bw1 = Inches(1.72)
    gap1 = Inches(0.22)
    sx1 = Inches(0.22)
    for i, (fill, bc, lbl, sub) in enumerate(boxes_r1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, row1_y, bw1, row1_h, fill, bc, 1.5)
        tb(s, bx + Emu(40000), row1_y + Emu(60000), bw1 - Emu(65000), Emu(420000),
           lbl, 8.5, bold=True, color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), row1_y + Emu(490000), bw1 - Emu(50000), Emu(540000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), row1_y + row1_h // 2, col=C_GOLD, size=0.85)

    # 大矢印（下向き）: GG突入
    mid_x = sx1 + 2 * (bw1 + gap1) + bw1 // 2
    arrow_d(s, mid_x, row1_y + row1_h, col=C_GOLD)

    # ─── 下段（GG → 上位AT → GOD揃い）────────────────────────
    row2_y = row1_y + row1_h + Emu(200000)
    row2_h = Emu(1180000)
    label2_y = row2_y - Emu(240000)
    rect(s, Inches(0.15), label2_y, Inches(9.7), Emu(200000), C_CARD2)
    tb(s, Inches(0.28), label2_y + Emu(25000), Inches(5.0), Emu(180000),
       "【AT中】 GOD GAME → 上位昇格ルート", 8.5, bold=True, color=C_RED)

    boxes_r2 = [
        (C_CARD2,                       C_GOLD, "GOD GAME\n(GG)",
         "50G / 7.0枚/G\n80%ループ\nストック積み上げ"),
        (C_CARD2,                       C_TEAL, "Z-ZONE\n(チャレンジ)",
         "0揃いで突入\n5Gフルナビ\nZ-GAME昇格を狙う"),
        (RGBColor(0x14, 0x10, 0x02),    C_YEL,  "Z-GAME\n(上乗せ特化)",
         "黄7でGGストック\n1/1.4で連鎖中\n爆発上乗せゾーン"),
        (RGBColor(0x20, 0x06, 0x04),    C_RED,  "SGG\n(赤7揃い)",
         "1/6900で突入\n75%以上ループ\n10〜100G増加区間"),
        (RGBColor(0x0A, 0x04, 0x18),    C_PUR2, "GOD揃い\n(PGG)",
         "1/16384\nGGストック4個以上\n期待3000枚以上"),
    ]
    sx2 = Inches(0.22)
    bw2 = Inches(1.72)
    gap2 = Inches(0.22)
    for i, (fill, bc, lbl, sub) in enumerate(boxes_r2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, row2_y, bw2, row2_h, fill, bc, 1.8)
        tb(s, bx + Emu(40000), row2_y + Emu(65000), bw2 - Emu(65000), Emu(440000),
           lbl, 9, bold=True, color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), row2_y + Emu(510000), bw2 - Emu(50000), Emu(570000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw2 + Emu(10000), row2_y + row2_h // 2, col=C_GOLD, size=0.85)

    footer(s,
           "設計核心: 通常時はモード管理×ガイアステージで多様な当選ルートを提供。AT中は5つのゾーンが連鎖する多層昇格設計",
           "モード(6種)→ガイアステージ→GG→Z-ZONE→Z-GAME/SGG→GOD揃い(PGG)という完全な階層構造を持つ。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 2つのルートでGOD GAMEを目指す", "3/10")

    # ─── 左：モードルート ────────────────────────────────────
    lx, ly = Inches(0.22), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(310000), C_DGOLD)
    tb(s, lx + Emu(50000), ly + Emu(45000), lw - Emu(70000), Emu(240000),
       "ルート① モード昇格", 10.5, bold=True, color=C_GOLD2, font=FONT_H)

    modes = [
        (C_LTGRY, "低確A/B（通常モード）",
         "最初の状態。成立役でのGG当選は低め。\nチャンス役を引いてモード昇格を目指す。"),
        (C_CRIM,  "通常モード → 天国準備",
         "チャンス役が重なるとモードが上昇。\n小役履歴（青7・黄7の連続）がモード示唆。"),
        (C_GOLD,  "天国モード（天国・超天国）",
         "GGがほぼ確定する最上位モード。\n天国移行を確認したら打ち続けること。"),
        (C_TEAL,  "ガイアステージ（超高確）",
         "ガイアベル規定回数到達で突入(9G+α)。\n毎G超高確率でGG抽選。奇数揃いでGG確定！"),
    ]
    mh = Emu(890000)
    for i, (bc, mt, md) in enumerate(modes):
        my = ly + Emu(310000) + i * mh
        rect_b(s, lx, my, lw, mh - Emu(20000), C_CARD, bc, 1.2)
        rect(s, lx, my, Emu(38000), mh - Emu(20000), bc)
        tb(s, lx + Emu(68000), my + Emu(50000), lw - Emu(90000), Emu(260000),
           mt, 8.5, bold=True, color=bc)
        tb(s, lx + Emu(68000), my + Emu(300000), lw - Emu(90000), Emu(520000),
           md, 8, color=C_WHITE)
        if i < 3:
            arrow_d(s, lx + lw // 2, my + mh - Emu(20000), col=bc)

    # ─── 右：Z-ZONEルート ───────────────────────────────────
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.78)

    rect(s, rx, ry, rw, Emu(310000), RGBColor(0x0A, 0x18, 0x20))
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(240000),
       "ルート② Z-ZONEからの突入", 10.5, bold=True, color=C_TEAL, font=FONT_H)

    zinfo = [
        (C_TEAL, "Z-ZONEとは？",
         "ガイアステージ中の0揃い等で突入する\nGG＋Z-GAME両睨みのチャレンジゾーン。\n5G間フルナビで黄7連続成立を狙う。"),
        (C_YEL,  "Z-ZONEの抽選",
         "5G中に黄7が5連続成立→Z-GAME突入確定。\n途中でハズレ・青7が出るとチャレンジ終了。\n黄7を引き続ける確率は約1/1.4×5G分。"),
        (C_GOLD, "失敗してもGG当選あり",
         "Z-ZONEでチャレンジ失敗でもGG当選する\nケースがある。ガイアステージ中の奇数揃いは\nGG確定。Z-ZONEは「上振れ」を狙う区間。"),
    ]
    zh = Emu(1180000)
    for i, (bc, zt, zd) in enumerate(zinfo):
        zy = ry + Emu(310000) + i * zh
        rect_b(s, rx, zy, rw, zh - Emu(20000), C_CARD, bc, 1.2)
        rect(s, rx, zy, Emu(38000), zh - Emu(20000), bc)
        tb(s, rx + Emu(68000), zy + Emu(50000), rw - Emu(90000), Emu(260000),
           zt, 8.5, bold=True, color=bc)
        tb(s, rx + Emu(68000), zy + Emu(310000), rw - Emu(90000), Emu(760000),
           zd, 8, color=C_WHITE)

    footer(s,
           "設計核心: 6段階モード管理＋ガイアステージ＋Z-ZONEで通常時に「進んでいる感」を生む多軸構造",
           "小役履歴(青7・黄7の連続)がモード示唆になるため、プレイヤーは自然と次の役に期待を持てる。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: AT「GOD GAME」の遊び方
# ══════════════════════════════════════════════════════════════
def s_gg(prs):
    s = new_slide(prs)
    hdr(s, "AT「GOD GAME（GG）」の遊び方 ── 役を引いてストックを積み上げる", "4/10")

    # ─── 上段：GGの基本スペック ──────────────────────────────
    specs = [
        (C_GOLD,  "純増",        "約7.0枚/G",    "現行スマスロトップクラスの純増速度"),
        (C_GOLD2, "1セット",     "50G",           "1セット単発なら約350枚獲得"),
        (C_CRIM,  "ループ率",    "最大80%",       "GGストック次第で長期連戦が可能"),
        (C_TEAL,  "ストック方式", "セット数管理",  "GGストック消化でGGを重ねていく"),
    ]
    sw = Inches(2.25)
    sh = Emu(820000)
    sy = Inches(0.72)
    for i, (ac, key, val, desc) in enumerate(specs):
        sx = Inches(0.22) + i * (sw + Inches(0.1))
        rect_b(s, sx, sy, sw, sh, C_CARD, ac, 1.5)
        rect(s, sx, sy, Emu(35000), sh, ac)
        tb(s, sx + Emu(65000), sy + Emu(55000), sw - Emu(80000), Emu(260000),
           key, 8, bold=True, color=ac)
        tb(s, sx + Emu(65000), sy + Emu(320000), sw - Emu(80000), Emu(280000),
           val, 14, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, sx + Emu(65000), sy + Emu(600000), sw - Emu(80000), Emu(200000),
           desc, 7.5, color=C_GRAY)

    # ─── 下段左：GG消化中の行動指針 ─────────────────────────
    lx, ly = Inches(0.22), sy + sh + Emu(120000)
    lw, lh = Inches(4.5), Emu(2600000)

    rect_b(s, lx, ly, lw, lh, C_CARD, C_GOLD, 1.5)
    rect(s, lx, ly, Emu(42000), lh, C_GOLD)
    tb(s, lx + Emu(72000), ly + Emu(50000), lw - Emu(95000), Emu(270000),
       "GG中に「何をすれば出玉が伸びる」か", 9.5, bold=True, color=C_GOLD, font=FONT_H)

    actions = [
        ("①  成立役でGGストック抽選",
         "GG消化中は毎ゲーム成立役に応じてGGストック抽選が行われる。\n"
         "強チェリー・強スイカ等の強役ほど当選率が高い。\n"
         "役を引くほどストックが積み上がり次のGGへ繋がる。"),
        ("②  黄7を狙えZ-ZONE",
         "GG消化中に黄7が揃うとZ-ZONE突入のチャンス。\n"
         "Z-ZONEを経由してZ-GAMEへ昇格すれば上乗せ爆発。\n"
         "黄7連続成立が「出玉が大きく伸びる」瞬間。"),
        ("③  ループストックも確認",
         "GGには「GGストック」と「ループストック」の2種が存在。\n"
         "ループストックが多いほど継続確率が上昇する。\n"
         "GOD揃い・SGG発生時にまとめて獲得できる。"),
    ]
    for i, (at, ad) in enumerate(actions):
        ay = ly + Emu(320000) + i * Emu(750000)
        tb(s, lx + Emu(72000), ay, lw - Emu(95000), Emu(280000),
           at, 8.5, bold=True, color=C_GOLD2)
        tb(s, lx + Emu(72000), ay + Emu(265000), lw - Emu(95000), Emu(430000),
           ad, 8, color=C_WHITE)

    # ─── 下段右：黄7とZ-ZONEの連動 ──────────────────────────
    rx, ry2 = Inches(5.0), ly
    rw, rh = Inches(4.78), lh

    rect_b(s, rx, ry2, rw, rh, C_DYEL, C_YEL, 1.5)
    rect(s, rx, ry2, Emu(42000), rh, C_YEL)
    tb(s, rx + Emu(72000), ry2 + Emu(50000), rw - Emu(95000), Emu(270000),
       "黄7→Z-ZONE→Z-GAMEのメカニズム", 9.5, bold=True, color=C_YEL, font=FONT_H)
    tb(s, rx + Emu(72000), ry2 + Emu(320000), rw - Emu(95000), rh - Emu(380000),
       "GG消化中に黄7が揃う\n          ↓\nZ-ZONE突入（5Gフルナビ）\n          ↓\n黄7が5連続成立 → Z-GAME突入\n          ↓\nZ-GAME中は1/1.4で黄7が揃い続ける\n          ↓\n黄7が揃うたびGGストック上乗せ\n          ↓\nハズレ or 青7が出るまで延々と連鎖！",
       8.5, color=C_WHITE)

    footer(s,
           "設計核心: GG中は「役を引く→ストック蓄積→黄7でZ-GAME連鎖」という明確な行動目標と報酬ループが設計されている",
           "プレイヤーは常に「次の役で何かが起きる」という期待感を維持できる。強役を引いた瞬間の興奮が台離れを防ぐ。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 上位「SGG / Z-GAME」
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位AT ── SGG（赤7揃い）× Z-GAME（黄7連鎖）の突入ルートと遊び方", "5/10")

    # ─── 上半分：Z-GAME ──────────────────────────────────────
    zx, zy = Inches(0.22), Inches(0.72)
    zw, zh = Inches(4.5), Emu(2000000)

    rect_b(s, zx, zy, zw, zh, C_DYEL, C_YEL, 2.0)
    rect(s, zx, zy, Emu(50000), zh, C_YEL)
    tb(s, zx + Emu(80000), zy + Emu(50000), zw - Emu(110000), Emu(310000),
       "Z-GAME（黄7上乗せ特化）", 12, bold=True, color=C_YEL, font=FONT_H)

    z_info = (
        "【突入条件】\n"
        "Z-ZONE（5Gフルナビ）中に黄7が5連続成立\n\n"
        "【消化中の挙動】\n"
        "・1/1.4の確率で黄7が揃い続ける\n"
        "・黄7が揃うたびGGストック上乗せ\n"
        "・ハズレまたは青7が出るまで継続\n\n"
        "【チェーン連鎖】\n"
        "Z-GAME中にZ-ZONEへ再突入することもある\n"
        "→ Z-GAME→Z-ZONE→Z-GAMEの無限連鎖が爆発の正体"
    )
    tb(s, zx + Emu(80000), zy + Emu(360000), zw - Emu(110000), zh - Emu(420000),
       z_info, 8.5, color=C_WHITE)

    # ─── 右上：SGG ───────────────────────────────────────────
    sx, sy2 = Inches(5.0), Inches(0.72)
    sw, sgh = Inches(4.78), Emu(2000000)

    rect_b(s, sx, sy2, sw, sgh, C_DRED, C_RED, 2.0)
    rect(s, sx, sy2, Emu(50000), sgh, C_RED)
    tb(s, sx + Emu(80000), sy2 + Emu(50000), sw - Emu(110000), Emu(310000),
       "SGG（スーパーGOD GAME）", 12, bold=True, color=C_RED, font=FONT_H)

    sgg_info = (
        "【突入条件】\n"
        "赤7揃い（確率：約1/6900）で突入\n"
        "GGストック1個+ループストック獲得\n\n"
        "【1セットの構成】\n"
        "増加区間 10G〜100G（出玉増加パート）\n"
        "＋ 引き戻しゾーン 3G（ループ抽選）\n\n"
        "【ループ率・継続】\n"
        "ループ率75%以上で継続\n"
        "引き戻しゾーンでさらにストック獲得の可能性あり"
    )
    tb(s, sx + Emu(80000), sy2 + Emu(360000), sw - Emu(110000), sgh - Emu(420000),
       sgg_info, 8.5, color=C_WHITE)

    # ─── 下段：比較まとめ ────────────────────────────────────
    comp_y = Inches(0.72) + Emu(2000000) + Emu(100000)
    comp_h = Emu(1480000)

    rect_b(s, Inches(0.22), comp_y, Inches(9.56), comp_h, C_CARD, C_GOLD, 1.2)
    tb(s, Inches(0.42), comp_y + Emu(50000), Inches(9.0), Emu(280000),
       "SGG vs Z-GAME 比較", 9.5, bold=True, color=C_GOLD, font=FONT_H)

    comp_items = [
        ("", "Z-GAME（黄7特化）", "SGG（赤7特化）"),
        ("突入条件", "Z-ZONE中に黄7×5連続成立", "赤7揃い（1/6900）"),
        ("消化方式", "黄7が揃い続ける上乗せ型", "セット数管理型（10〜100G）"),
        ("ループ",   "1/1.4で黄7成立中は継続", "75%以上のループ率"),
        ("特徴",     "連鎖するほど加速・爆発的上乗せ", "引き戻しゾーン付きで粘れる"),
    ]
    rh2 = comp_h // len(comp_items)
    cws = [Emu(1400000), Emu(3300000), Emu(3300000)]
    col_colors = [C_GRAY, C_YEL, C_RED]
    for ri, row in enumerate(comp_items):
        ry3 = comp_y + Emu(320000) + ri * rh2
        bg = C_CARD if ri % 2 == 0 else C_ROW
        cx = Inches(0.42)
        for ci, (cw, cv, cc) in enumerate(zip(cws, row, col_colors)):
            rect(s, cx, ry3, cw, rh2 - Emu(20000), bg)
            bold = ri == 0 or ci == 0
            tb(s, cx + Emu(30000), ry3 + Emu(30000), cw - Emu(50000), rh2 - Emu(50000),
               cv, 8, bold=bold, color=cc if (ri == 0 or ci == 0) else C_WHITE)
            cx += cw

    footer(s,
           "設計核心: SGGは「確定役の衝撃＋セット継続」、Z-GAMEは「連鎖上乗せの爆発」と、全く異なる興奮を2軸で提供",
           "赤7揃いは頻度は低いが確実なボーナス。Z-GAMEは黄7連鎖という純粋なスキル感（運感）を最大化する設計。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: GOD揃い・最大恩恵の解説
# ══════════════════════════════════════════════════════════════
def s_god(prs):
    s = new_slide(prs)
    hdr(s, "GOD揃い（PGG）── 1/16384の奇跡・最高体験の設計", "6/10")

    # 中央のGOD揃いドラム演出イメージ（ボックス）
    gx, gy = Inches(0.22), Inches(0.72)
    gw, gh = Inches(9.56), Emu(1680000)

    rect_b(s, gx, gy, gw, gh, C_DPUR, C_PUR2, 2.5)
    rect(s, gx, gy, Emu(60000), gh, C_PUR2)
    rect(s, gx + gw - Emu(60000), gy, Emu(60000), gh, C_PUR2)

    tb(s, gx + Emu(90000), gy + Emu(50000), gw - Emu(160000), Emu(380000),
       "GOD揃い（PREMIUM GOD GAME / PGG）", 14, bold=True, color=C_PUR2,
       align=PP_ALIGN.CENTER, font=FONT_H)

    # 恩恵数値を横並びで
    pgd = [
        (C_PUR2,  "発生確率",  "1/16384",   "全設定共通"),
        (C_GOLD,  "GGストック", "4セット以上", "継続確定レベル"),
        (C_CRIM,  "ループストック", "大量獲得", "高継続確率を保証"),
        (C_GOLD2, "期待枚数",   "3000枚以上", "単発最大クラスの恩恵"),
    ]
    pw = (gw - Emu(100000)) // len(pgd)
    for i, (ac, key, val, sub) in enumerate(pgd):
        px = gx + Emu(50000) + i * pw
        rect_b(s, px + Emu(30000), gy + Emu(430000), pw - Emu(50000), Emu(1100000),
               RGBColor(0x10, 0x08, 0x1C), ac, 1.5)
        tb(s, px + Emu(60000), gy + Emu(490000), pw - Emu(90000), Emu(280000),
           key, 8, bold=True, color=ac, align=PP_ALIGN.CENTER)
        tb(s, px + Emu(60000), gy + Emu(770000), pw - Emu(90000), Emu(360000),
           val, 13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, px + Emu(60000), gy + Emu(1120000), pw - Emu(90000), Emu(280000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ─── 下段：3つの設計的意味 ───────────────────────────────
    by = gy + gh + Emu(80000)
    bh2 = Emu(1700000)
    bw3 = (gw - Emu(80000)) // 3

    meanings = [
        (C_PUR2, "① 語り継がれる体験",
         "1/16384という絶対的なレア度が\n「あの日GOD揃いが出た」を\n一生の話題にする。\n\n打ち手にとって一度でも経験すれば\n「この台で起きた最高の瞬間」として\n記憶に刻まれる設計。"),
        (C_GOLD2,"② SNS・UGC生産装置",
         "GOD揃い動画は高い確率でSNSでバズる。\n\n① バズる動画が広告塔になる\n② 「ありえる」という期待感を拡散\n③ 「俺も出したい」来店動機を創出\n\nこの1フラグだけで自走するマーケティングが成立する。"),
        (C_CRIM, "③ 期待値ゼロでも希望を残す",
         "設定1でも・ハマっていても\n「次でGOD揃いが出るかもしれない」\nという希望は同じ。\n\n3000枚確定という圧倒的な恩恵が\n閉店まで離席を防ぐ\n最後のストッパーとして機能する。"),
    ]
    for i, (ac, mt, md) in enumerate(meanings):
        mx = gx + Emu(10000) + i * (bw3 + Emu(30000))
        rect_b(s, mx, by, bw3, bh2, C_CARD, ac, 1.5)
        rect(s, mx, by, Emu(42000), bh2, ac)
        tb(s, mx + Emu(72000), by + Emu(55000), bw3 - Emu(95000), Emu(290000),
           mt, 9, bold=True, color=ac)
        tb(s, mx + Emu(72000), by + Emu(340000), bw3 - Emu(95000), bh2 - Emu(400000),
           md, 8, color=C_WHITE)

    footer(s,
           "設計核心: 1/16384のGOD揃いは「語り継がれる体験」「SNS拡散」「希望を残すストッパー」という3重の設計価値を持つ",
           "この確率は狙って出せないからこそ価値がある。稀少性と恩恵の大きさが神話的体験を生む。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（なぜこの台は面白いのか）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── なぜ「ミリオンゴッド-軌跡-」は面白いのか", "7/10")

    # ─── 5つの設計原則カード ─────────────────────────────────
    cards = [
        (C_GOLD,  "① 速度の快感",
         "7.0枚/Gという現行最速クラスの\n純増スピードが「GGが続く快感」を\n最大化する。\n50G×7.0枚=350枚が最小単位として\n常に「速く増える」体感を与える。"),
        (C_YEL,   "② 連鎖の爆発感",
         "Z-GAME黄7連鎖は\n「当たり続けるほど加速する」\n正のフィードバックループ。\n1/1.4を何連鎖できるかという\nシンプルなドキドキが設計の核心。"),
        (C_RED,   "③ 確定役の衝撃",
         "赤7揃い（1/6900）という\nひとつのレア役がSGGという\n別次元のATを起動する。\n「揃った瞬間の衝撃」が\n台の価値を上げる設計。"),
        (C_TEAL,  "④ 多層の目標",
         "GG→SGG→Z-GAME→GOD揃いという\n4層の目標が常に「次がある」\n希望を生む。\nどの層にいても上位への\n期待感が途切れない構造。"),
        (C_PUR2,  "⑤ 記憶との接続",
         "4号機ミリオンゴッドを知る世代が\n「GOD揃いをもう一度」という\n動機で来店する。\n記憶・郷愁がリピート来店を\n促す最強のIPカード。"),
    ]
    cw = Inches(1.78)
    ch = Emu(3800000)
    gap = Inches(0.115)
    cx0 = Inches(0.22)
    cy0 = Inches(0.72)

    for i, (ac, ct, cd) in enumerate(cards):
        cx = cx0 + i * (cw + gap)
        rect_b(s, cx, cy0, cw, ch, C_CARD, ac, 1.8)
        rect(s, cx, cy0, Emu(40000), ch, ac)
        tb(s, cx + Emu(65000), cy0 + Emu(55000), cw - Emu(85000), Emu(310000),
           ct, 9, bold=True, color=ac, font=FONT_H)
        tb(s, cx + Emu(65000), cy0 + Emu(380000), cw - Emu(85000), ch - Emu(450000),
           cd, 8, color=C_WHITE)

    # ─── 下段：設計の核心まとめ ──────────────────────────────
    by = cy0 + ch + Emu(80000)
    bh3 = Emu(830000)
    rect_b(s, Inches(0.22), by, Inches(9.56), bh3, RGBColor(0x10, 0x08, 0x02), C_GOLD, 1.5)
    tb(s, Inches(0.42), by + Emu(55000), Inches(9.0), Emu(290000),
       "設計の本質", 10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.42), by + Emu(340000), Inches(9.0), Emu(420000),
       "「速さ（7.0枚）× 連鎖（Z-GAME）× 衝撃（赤7/GOD揃い）× IP（4号機の記憶）」の4要素が\n"
       "互いを補強し合うことで、短期体験（1セット）も長期体験（一日）も成立する完全な設計。",
       8.5, color=C_WHITE)

    footer(s,
           "設計核心: 速度・連鎖・衝撃・記憶という4軸が互いを補強し合い、短期〜長期の両方の遊び方で満足感を生む",
           "単純に「出玉が多い」だけでなく、複数の快感設計が重なった結果として長期稼働が実現している。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_proscons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── プレイヤー・ホール両視点の評価", "8/10")

    lx, ly = Inches(0.22), Inches(0.72)
    lw = Inches(4.5)
    col_h = Emu(4100000)

    # 良い点
    rect(s, lx, ly, lw, Emu(340000), C_TEAL)
    tb(s, lx + Emu(60000), ly + Emu(50000), lw - Emu(80000), Emu(260000),
       "良い点（Pros）", 11, bold=True, color=C_BG, font=FONT_H)

    pros = [
        (C_TEAL,  "7.0枚/Gの純増速度",
         "出玉の増加速度が体感として速い。\n1セット50Gで約350枚という明確な数値が\n「打った感」を高める。"),
        (C_GOLD,  "GOD揃いの神話的体験",
         "1/16384のレア度が「打ち手の間の語り草」を生む。\nSNS拡散効果で新規来店促進にもなる。"),
        (C_YEL,   "Z-GAME黄7連鎖の爆発感",
         "連鎖するほど加速する上乗せが\n他機種では味わえない興奮を提供する。\n正のフィードバックループの設計が優秀。"),
        (C_RED,   "4号機GODのIP継続力",
         "既存ユーザーへの訴求力が高い。\n「またGOD揃いを体験したい」という\nリピート動機が強力。"),
    ]
    ph = (col_h - Emu(340000)) // len(pros)
    for i, (bc, pt, pd) in enumerate(pros):
        py = ly + Emu(340000) + i * ph
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect_b(s, lx, py, lw, ph - Emu(15000), bg, bc, 0.8)
        rect(s, lx, py, Emu(30000), ph - Emu(15000), bc)
        tb(s, lx + Emu(55000), py + Emu(45000), lw - Emu(70000), Emu(250000),
           pt, 8.5, bold=True, color=bc)
        tb(s, lx + Emu(55000), py + Emu(290000), lw - Emu(70000), ph - Emu(360000),
           pd, 8, color=C_WHITE)

    # 課題
    rx, ry2 = Inches(5.05), Inches(0.72)
    rw2 = Inches(4.73)
    rect(s, rx, ry2, rw2, Emu(340000), C_RED)
    tb(s, rx + Emu(60000), ry2 + Emu(50000), rw2 - Emu(80000), Emu(260000),
       "課題（Cons）", 11, bold=True, color=C_WHITE, font=FONT_H)

    cons = [
        (C_CRIM,  "ハマり時の待ち時間が長い",
         "天井1480Gは現行機の中でも深め。\n低設定での長期ハマりは\nプレイヤーに大きなストレスを与える。"),
        (C_CRIM,  "GOD揃いへの依存",
         "最高体験がGOD揃い（1/16384）に集中しすぎると\n大多数のプレイヤーは生涯体験できない可能性。\n「2番手の体験」を充実させる余地がある。"),
        (C_CRIM,  "設定差の開示が少ない",
         "設定6の機械割しか公表されていない。\nプレイヤーが設定を判別しにくく\n通い続ける理由が作りにくい。"),
        (C_CRIM,  "初心者には複雑すぎる可能性",
         "モード・ガイアステージ・Z-ZONE・Z-GAME・SGG\nと覚える要素が多い。\n初見プレイヤーが「何を狙えばいいか」\n理解するまでに時間がかかる。"),
    ]
    ch2 = (col_h - Emu(340000)) // len(cons)
    for i, (bc, ct, cd) in enumerate(cons):
        cy2 = ry2 + Emu(340000) + i * ch2
        bg = C_DRED if i % 2 == 0 else RGBColor(0x18, 0x06, 0x04)
        rect_b(s, rx, cy2, rw2, ch2 - Emu(15000), bg, bc, 0.8)
        rect(s, rx, cy2, Emu(30000), ch2 - Emu(15000), bc)
        tb(s, rx + Emu(55000), cy2 + Emu(45000), rw2 - Emu(70000), Emu(250000),
           ct, 8.5, bold=True, color=bc)
        tb(s, rx + Emu(55000), cy2 + Emu(290000), rw2 - Emu(70000), ch2 - Emu(360000),
           cd, 8, color=C_WHITE)

    footer(s,
           "分析視点: 「出玉性能×IP力」の強みは明確。課題は天井深さと初心者導線の不足で、それを補う店舗施策が重要",
           "強みを活かすには設定6を使った高時給の演出と、天井狙い客を取り込む台数・配置戦略がカギになる。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること・応用できる原則", "9/10")

    # ─── 左：長期稼働を支えた3要素 ──────────────────────────
    lx, ly = Inches(0.22), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(310000), C_DGOLD)
    tb(s, lx + Emu(60000), ly + Emu(50000), lw - Emu(80000), Emu(240000),
       "長期稼働を支えた3要素", 11, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_GOLD,  "① 純増7.0枚の「速さ」差別化",
         "現行機トップクラスの純増速度が\n「速く大きく勝つ」体験を実現。\n4号機GODの「爆発力」をスマスロで再現。\n速さそのものが競合との差別化ポイント。"),
        (C_YEL,   "② GG→SGG→Z-GAMEの多層昇格",
         "常に「次がある」希望を生む4層構造。\nGG中に赤7・黄7を狙い続ける\nという明確な行動指針が\nプレイヤーの集中を維持する。"),
        (C_PUR2,  "③ IP力×世代記憶の活用",
         "4号機ミリオンゴッドを知る世代が\nスマスロ版で「GOD揃いを再び」と\n来店する。記憶との接続が\nリピート来店を促す最強の動機。"),
    ]
    eh = Emu(1230000)
    for i, (ac, et, ed) in enumerate(elems):
        ey = ly + Emu(310000) + i * eh
        rect_b(s, lx, ey, lw, eh - Emu(20000), C_CARD, ac, 1.5)
        rect(s, lx, ey, Emu(42000), eh - Emu(20000), ac)
        tb(s, lx + Emu(72000), ey + Emu(50000), lw - Emu(95000), Emu(270000),
           et, 9, bold=True, color=ac)
        tb(s, lx + Emu(72000), ey + Emu(310000), lw - Emu(95000), eh - Emu(380000),
           ed, 8, color=C_WHITE)

    # ─── 右：設計原則 & 総括 ─────────────────────────────────
    rx, ry3 = Inches(5.0), Inches(0.72)
    rw3 = Inches(4.78)

    rect(s, rx, ry3, rw3, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry3 + Emu(45000), rw3 - Emu(70000), Emu(210000),
       "他機種設計に応用できる原則", 10, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_GOLD,  "純増速度は「体感」を最大化する武器になる"),
        (C_YEL,   "Z-GAME型「連鎖上乗せ」は最も興奮度が高い形式"),
        (C_RED,   "確定役（1/6900）の衝撃は通常ループを飽きさせない"),
        (C_PUR2,  "1/16384のGOD揃いが語り継がれる体験を生む"),
        (C_TEAL,  "モード管理＋ガイアステージで通常時を「進んでいる感」に"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry3 + Emu(280000) + i * Emu(500000)
        rect(s, rx, py0, Emu(22000), Emu(460000), ac)
        tb(s, rx + Emu(52000), py0 + Emu(70000), rw3 - Emu(65000), Emu(350000),
           p, 8.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry3 + Emu(2850000), rw3, Emu(870000),
           RGBColor(0x10, 0x08, 0x02), C_GOLD, 1.5)
    tb(s, rx + Emu(55000), ry3 + Emu(2900000), rw3 - Emu(75000), Emu(280000),
       "総括", 9, bold=True, color=C_GOLD)
    tb(s, rx + Emu(55000), ry3 + Emu(3170000), rw3 - Emu(75000), Emu(480000),
       "IP×純増7枚×多層昇格設計の完成形。\n"
       "GOD揃い1/16384が神話的体験を生産し続ける\n唯一無二のコンテンツ力を持つスマスロ機種。",
       8, color=C_WHITE)

    footer(s,
           "設計原則: 速度・連鎖・衝撃・記憶の4軸設計を他機種に応用するとき「それぞれの軸が互いを補強しているか」を問うべき",
           "ミリオンゴッド-軌跡-の成功は単一要素ではなく4つの軸が化学反応を起こした結果である。")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)    # 1: タイトル・スペック・3ポイント
    s_flow(prs)     # 2: ゲームフロー全体図
    s_normal(prs)   # 3: 通常時の遊び方
    s_gg(prs)       # 4: GOD GAMEの遊び方
    s_upper(prs)    # 5: SGG / Z-GAME
    s_god(prs)      # 6: GOD揃い・最大恩恵
    s_design(prs)   # 7: 面白さの設計
    s_proscons(prs) # 8: 良い点と課題
    s_matome(prs)   # 9: まとめ・設計原則

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
