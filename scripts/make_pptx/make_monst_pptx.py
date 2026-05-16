"""
スマスロ モンスターストライク 企画提案書 PPTX
9スライド構成 / 画像素材: C:/Users/h.kadoya/Desktop/monst_imgs/
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ---- パス設定 ----
IMG_DIR = r"C:\Users\h.kadoya\Desktop\monst_imgs"
OUT_DIR = r"C:\Users\h.kadoya\Desktop\slocri\proposals\新規提案\スマスロ_モンスト"
OUT_PATH = os.path.join(OUT_DIR, "monst_kikaku_v1.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

def img(name):
    return os.path.join(IMG_DIR, name)

# ---- カラー定義 ----
C_NAVY    = RGBColor(0x0D, 0x1B, 0x2A)
C_NAVY2   = RGBColor(0x14, 0x26, 0x3A)
C_ORANGE  = RGBColor(0xFF, 0x66, 0x00)
C_GOLD    = RGBColor(0xFF, 0xB8, 0x00)
C_BLUE    = RGBColor(0x00, 0xC8, 0xFF)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY   = RGBColor(0xCC, 0xCC, 0xCC)
C_DGRAY   = RGBColor(0x22, 0x33, 0x44)
C_RED     = RGBColor(0xFF, 0x33, 0x00)
C_GREEN   = RGBColor(0x00, 0xCC, 0x66)
C_FILL_OR = RGBColor(0xFF, 0x88, 0x00)  # オレンジ薄
C_FILL_BL = RGBColor(0x00, 0x55, 0x88)  # ブルー暗
C_FILL_GD = RGBColor(0x88, 0x66, 0x00)  # ゴールド暗

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # 完全ブランク

# ===== ユーティリティ =====
def add_rect(slide, x, y, w, h, fill=None, line=None, alpha=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    return sh

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or C_WHITE
    return tb

def add_label_box(slide, label, x, y, w, h, bg=None, text_color=None, size=14):
    """背景付きラベルボックス"""
    add_rect(slide, x, y, w, h, fill=bg or C_FILL_BL)
    add_text(slide, label, x, y, w, h,
             size=size, bold=True, color=text_color or C_WHITE,
             align=PP_ALIGN.CENTER)

def add_bg(slide, color=None):
    color = color or C_NAVY
    add_rect(slide, 0, 0, W, H, fill=color)

def add_header_bar(slide, title, subtitle=None):
    """スライド上部ヘッダーバー"""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=C_FILL_BL)
    add_rect(slide, 0, Inches(1.0), W, Pt(3), fill=C_ORANGE)
    add_text(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.8),
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.7), Inches(10), Inches(0.45),
                 size=13, color=C_LGRAY)

def add_picture_safe(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    try:
        if w and h:
            return slide.shapes.add_picture(path, x, y, w, h)
        elif w:
            return slide.shapes.add_picture(path, x, y, width=w)
        elif h:
            return slide.shapes.add_picture(path, x, y, height=h)
        else:
            return slide.shapes.add_picture(path, x, y)
    except Exception as e:
        print(f"  [WARN] {path}: {e}")
        return None

def add_flow_arrow(slide, x, y, w=Inches(0.5), h=Inches(0.35)):
    """→ 矢印"""
    sh = slide.shapes.add_shape(13, x, y, w, h)  # Right Arrow
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_ORANGE
    sh.line.fill.background()

# ===== SLIDE 1: タイトル =====
s1 = prs.slides.add_slide(blank)
add_bg(s1, C_NAVY)

# 背景グラデーション風（下部を少し明るく）
add_rect(s1, 0, Inches(5.0), W, Inches(2.5), fill=RGBColor(0x0A, 0x12, 0x1E))

# 右側に映画2ポスターをフェード配置
add_picture_safe(s1, img("movie2_poster.jpg"), Inches(8.3), Inches(0.5), h=Inches(6.5))
# 上から暗くするオーバーレイ（右側）
add_rect(s1, Inches(7.8), 0, Inches(5.53), H, fill=RGBColor(0x0D, 0x1B, 0x2A))

# ロゴ
add_picture_safe(s1, img("app_icon.png"), Inches(0.5), Inches(0.4), h=Inches(0.9))

# メインタイトル
add_text(s1, "スマスロ", Inches(1.6), Inches(0.3), Inches(8), Inches(0.85),
         size=30, bold=True, color=C_ORANGE)
add_text(s1, "モンスターストライク", Inches(1.5), Inches(1.0), Inches(9), Inches(1.4),
         size=54, bold=True, color=C_WHITE)

# オレンジライン
add_rect(s1, Inches(1.5), Inches(2.35), Inches(7.0), Pt(4), fill=C_ORANGE)

# キャッチコピー
add_text(s1, "「超絶攻略世代」が、帰ってくる。",
         Inches(1.5), Inches(2.5), Inches(8.5), Inches(0.8),
         size=26, bold=True, color=C_GOLD)

# サブコピー
add_text(s1, "元モン廃 × スマスロ × 感情ストーリー\n覇者の塔を頂点とする3層射幸性で、\n「懐かしさ」と「熱さ」を同時に届ける。",
         Inches(1.5), Inches(3.2), Inches(7.0), Inches(1.6),
         size=16, color=C_LGRAY)

# キャラクター: Ren
add_picture_safe(s1, img("ren.png"), Inches(9.3), Inches(0.4), h=Inches(6.5))

# オラゴン（マスコット）
add_picture_safe(s1, img("oragon3.png"), Inches(1.2), Inches(4.3), h=Inches(2.7))

# ボトムバー
add_rect(s1, 0, Inches(7.1), W, Inches(0.4), fill=RGBColor(0x0A, 0x12, 0x1E))
add_text(s1, "KEY-CRE  /  企画提案資料  /  CONFIDENTIAL",
         Inches(0.4), Inches(7.15), Inches(12), Inches(0.3),
         size=9, color=RGBColor(0x55, 0x66, 0x77), align=PP_ALIGN.CENTER)

# ===== SLIDE 2: モンストとは？ =====
s2 = prs.slides.add_slide(blank)
add_bg(s2, C_NAVY)
add_header_bar(s2, "モンストとは？", "Monster Strike — IPの概要（パチスロ担当向け早わかりガイド）")

# 左エリア：IP説明（4ブロック）
blocks = [
    ("🎮 オーブ（ガチャ通貨）",
     "ゲーム内通貨「オーブ」を消費してモンスターガチャ。\n5個で単発・50個で10連。レア度はN〜SS。"),
    ("👥 4人マルチプレイ",
     "友人とリアルタイム協力プレイ。\n全員でモンスターを弾いてボスを倒す。\n\"仲間と攻略する体験\"が核心。"),
    ("⚡ 超絶 / 爆絶 / 轟絶",
     "難関クエストの格付け。超絶＜爆絶＜轟絶の順で難易度UP。\nクリアに強力な4人パーティが必要で、\n廃人プレイヤーの腕の見せ所だった。"),
    ("🏯 覇者の塔",
     "毎月更新の難関タワー型コンテンツ。\n全48フロアを踏破する達成感・繰り返し性が高い。\nスロットの継続AT設計と相性◎"),
]

bx, by = Inches(0.35), Inches(1.25)
bw, bh = Inches(4.7), Inches(1.35)
gap = Inches(0.08)

for i, (ttl, body) in enumerate(blocks):
    row = i // 2
    col = i % 2
    x = bx + col * (bw + Inches(0.15))
    y = by + row * (bh + gap)
    bg_c = C_FILL_BL if i % 2 == 0 else RGBColor(0x0C, 0x2A, 0x40)
    add_rect(s2, x, y, bw, bh, fill=bg_c, line=C_ORANGE)
    add_text(s2, ttl, x + Inches(0.12), y + Inches(0.06), bw - Inches(0.2), Inches(0.38),
             size=13, bold=True, color=C_ORANGE)
    add_text(s2, body, x + Inches(0.12), y + Inches(0.42), bw - Inches(0.2), bh - Inches(0.5),
             size=10.5, color=C_LGRAY)

# 映画セクション
add_rect(s2, Inches(0.35), Inches(4.0), Inches(9.7), Inches(0.35), fill=RGBColor(0x0C, 0x2A, 0x40))
add_text(s2, "🎬 映画展開（2作品）",
         Inches(0.4), Inches(4.0), Inches(6), Inches(0.35),
         size=12, bold=True, color=C_GOLD)

# 映画ポスター
add_picture_safe(s2, img("movie1_poster.jpg"), Inches(0.35), Inches(4.4), w=Inches(4.6))
add_picture_safe(s2, img("movie2_poster.jpg"), Inches(5.1), Inches(4.4), w=Inches(4.65))

add_text(s2, "①はじまりの場所へ（2016）",
         Inches(0.35), Inches(6.35), Inches(4.6), Inches(0.4),
         size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s2, "②ソラノカナタ（2018）",
         Inches(5.1), Inches(6.35), Inches(4.65), Inches(0.4),
         size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 右エリア：キャラクター
add_picture_safe(s2, img("elle.png"), Inches(10.3), Inches(2.2), h=Inches(2.8))
add_picture_safe(s2, img("mirai.png"), Inches(11.8), Inches(2.1), h=Inches(2.9))
add_picture_safe(s2, img("oragon1.png"), Inches(10.6), Inches(4.8), h=Inches(2.3))

add_text(s2, "キャラクター\n多数登場", Inches(10.2), Inches(1.5), Inches(3.0), Inches(0.7),
         size=11, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 3: ターゲット =====
s3 = prs.slides.add_slide(blank)
add_bg(s3, C_NAVY)
add_header_bar(s3, "ターゲット", "超絶攻略世代 — パチスロに最もハマりやすいモンスト経験者層")

# 左：メインターゲット
add_rect(s3, Inches(0.35), Inches(1.3), Inches(6.2), Inches(5.8),
         fill=RGBColor(0x0C, 0x22, 0x38), line=C_ORANGE)

add_text(s3, "メインターゲット",
         Inches(0.5), Inches(1.4), Inches(3), Inches(0.4),
         size=12, color=C_ORANGE, bold=True)

add_text(s3, "超絶攻略世代",
         Inches(0.5), Inches(1.75), Inches(6.0), Inches(0.75),
         size=36, bold=True, color=C_GOLD)

add_text(s3, "男性 30〜35歳  /  パチスロ経験者",
         Inches(0.5), Inches(2.45), Inches(5.5), Inches(0.45),
         size=16, color=C_WHITE)

add_rect(s3, Inches(0.5), Inches(2.9), Inches(5.8), Pt(2), fill=C_ORANGE)

persona_items = [
    ("元モン廃", "超絶/爆絶を毎日攻略。仲間とマルチで深夜まで遊んでいた。"),
    ("スマホゲー離脱", "2019年頃に引退。\n\"あの頃は楽しかった\"という強い郷愁を持つ。"),
    ("パチスロ現役", "週2〜3回来店。AT機・スマスロに移行済み。\n好きなIPが台になると打ちたくなる。"),
    ("再点火のトリガー", "オーブ図柄、オラゴン、ガチャ演出などで\n\"懐かしさ×射幸性\"の同時体験ができる。"),
]

for i, (k, v) in enumerate(persona_items):
    y = Inches(3.0) + i * Inches(0.85)
    add_text(s3, f"▶ {k}", Inches(0.55), y, Inches(1.6), Inches(0.38),
             size=12, bold=True, color=C_BLUE)
    add_text(s3, v, Inches(2.1), y, Inches(4.2), Inches(0.78),
             size=10.5, color=C_LGRAY)

# 右：キャラクター配置
add_picture_safe(s3, img("ren.png"), Inches(6.8), Inches(0.9), h=Inches(5.8))
add_picture_safe(s3, img("char_101.png"), Inches(9.8), Inches(1.5), h=Inches(5.1))
add_picture_safe(s3, img("char_103.png"), Inches(11.3), Inches(1.8), h=Inches(4.8))

add_text(s3, "モンスト経験者ならこの顔を覚えている",
         Inches(6.5), Inches(6.7), Inches(6.5), Inches(0.5),
         size=11, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 4: ゲームコンセプト =====
s4 = prs.slides.add_slide(blank)
add_bg(s4, C_NAVY)
add_header_bar(s4, "ゲームコンセプト", "\"通常時をすべてCZ化\" × 3層射幸性で最後まで飽きさせない設計")

# コンセプト帯
add_rect(s4, Inches(0.35), Inches(1.25), W - Inches(0.7), Inches(0.6),
         fill=RGBColor(0x20, 0x10, 0x00), line=C_ORANGE)
add_text(s4, "コア設計思想：通常時のすべてがCZへの道筋になる",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         size=15, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# ゲームフロー図
flow = [
    ("通常時\nクエスト進行", C_FILL_BL,    "毎ゲーム\nダメージ蓄積"),
    ("オーブ獲得\nガチャ演出",  RGBColor(0x33,0x22,0x00), "5個→単発\n50個→10連"),
    ("仲間集め\n4キャラ揃う",  RGBColor(0x00,0x22,0x11), "ランク\nA〜S設定"),
    ("CZ突入\nマルチ体験",     RGBColor(0x22,0x00,0x22), "突破で\nAT確定"),
    ("AT\n超絶ステージ",       RGBColor(0x33,0x11,0x00), "映画1\nストーリー"),
    ("AT+\n爆絶ステージ",      RGBColor(0x00,0x11,0x33), "映画2\n解禁"),
    ("頂点\n轟絶/覇者の塔",    RGBColor(0x33,0x22,0x00), "最高射幸\n体験"),
]

fx = Inches(0.3)
fy = Inches(2.1)
fw = Inches(1.7)
fh = Inches(1.6)
gap = Inches(0.12)

for i, (label, bg, note) in enumerate(flow):
    x = fx + i * (fw + gap)
    add_rect(s4, x, fy, fw, fh, fill=bg, line=C_ORANGE)
    add_text(s4, label, x, fy + Inches(0.08), fw, Inches(0.75),
             size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, note, x, fy + Inches(0.82), fw, Inches(0.7),
             size=9, color=C_LGRAY, align=PP_ALIGN.CENTER)
    # 矢印（最後以外）
    if i < len(flow) - 1:
        ax = x + fw + Inches(0.01)
        ay = fy + fh / 2 - Inches(0.18)
        add_flow_arrow(s4, ax, ay, Inches(0.1), Inches(0.36))

# レベル帯 ラベル
for lbl, x1, x2, c in [
    ("NORMAL ZONE（通常時）", 0, 2, C_BLUE),
    ("BONUS ZONE（CZ）",      3, 3, C_GREEN),
    ("AT ZONE", 4, 6, C_ORANGE),
]:
    cx = fx + x1 * (fw + gap)
    cw = (x2 - x1 + 1) * (fw + gap) - gap
    add_rect(s4, cx, fy + fh + Inches(0.06), cw, Inches(0.3), fill=c)
    add_text(s4, lbl, cx, fy + fh + Inches(0.06), cw, Inches(0.3),
             size=9, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# 3層射幸性の説明
add_rect(s4, Inches(0.35), Inches(4.25), Inches(12.6), Inches(2.8),
         fill=RGBColor(0x08, 0x16, 0x24), line=C_BLUE)

add_text(s4, "3層射幸性の設計",
         Inches(0.5), Inches(4.3), Inches(4), Inches(0.4),
         size=14, bold=True, color=C_BLUE)

layers = [
    ("① 短期射幸", "ガチャ・仲間集め",
     "1〜50G単位の小さな期待感\nオーブ揃いで毎ゲーム緊張感"),
    ("② 中期射幸", "CZ→AT突入",
     "仲間4人集め→CZ→AT\n中程度の達成感の山を作る"),
    ("③ 長期射幸", "覇者の塔・轟絶",
     "AT中の高純増区間\n最終到達でフィーバー感"),
]
for i, (t1, t2, body) in enumerate(layers):
    x = Inches(0.5) + i * Inches(4.2)
    add_text(s4, t1, x, Inches(4.7), Inches(4.0), Inches(0.38),
             size=12, bold=True, color=C_GOLD)
    add_text(s4, t2, x, Inches(5.05), Inches(4.0), Inches(0.35),
             size=11, color=C_ORANGE)
    add_text(s4, body, x, Inches(5.38), Inches(3.9), Inches(0.9),
             size=10, color=C_LGRAY)

# キャラ
add_picture_safe(s4, img("oragon2.png"), Inches(11.5), Inches(3.5), h=Inches(3.6))

# ===== SLIDE 5: 通常時〜ガチャ演出 =====
s5 = prs.slides.add_slide(blank)
add_bg(s5, C_NAVY)
add_header_bar(s5, "通常時 〜 ガチャ演出", "オーブ図柄でゲームループを作る / 毎ゲームがガチャへの積み上げ")

# 左：通常時の流れ
add_text(s5, "通常時の仕組み",
         Inches(0.35), Inches(1.25), Inches(6), Inches(0.4),
         size=15, bold=True, color=C_ORANGE)

normal_steps = [
    ("STEP 1", "オーブ図柄が揃う", "青・赤・虹の3種類。揃うたびにオーブ+1〜3。"),
    ("STEP 2", "オーブ数が蓄積",   "ボーナスゲームのオーブは\n5個、10個、50個の節目で演出変化。"),
    ("STEP 3", "ガチャ発動",       "5個 → 単発ガチャ（レア演出）\n50個 → 10連ガチャ（強演出）"),
    ("STEP 4", "モンスター排出",   "キャラのランク（C〜S）が\n仲間の強さに直結。\nSランク4体で最強パーティ。"),
]

for i, (step, ttl, body) in enumerate(normal_steps):
    y = Inches(1.75) + i * Inches(1.35)
    add_rect(s5, Inches(0.35), y, Inches(6.3), Inches(1.28),
             fill=RGBColor(0x0C, 0x22, 0x38), line=C_BLUE)
    add_text(s5, step, Inches(0.4), y + Inches(0.08), Inches(1.0), Inches(0.35),
             size=10, bold=True, color=C_BLUE)
    add_text(s5, ttl, Inches(1.3), y + Inches(0.05), Inches(5.0), Inches(0.45),
             size=13, bold=True, color=C_WHITE)
    add_text(s5, body, Inches(1.3), y + Inches(0.45), Inches(5.0), Inches(0.78),
             size=10, color=C_LGRAY)

# 右上：オーブ種別
add_rect(s5, Inches(6.8), Inches(1.25), Inches(6.15), Inches(2.2),
         fill=RGBColor(0x08, 0x1A, 0x2C), line=C_GOLD)
add_text(s5, "オーブ図柄 — 種別と価値",
         Inches(6.9), Inches(1.3), Inches(5.5), Inches(0.4),
         size=13, bold=True, color=C_GOLD)

orbs = [
    ("青オーブ", "+1個", "C_BLUE", "通常オーブ。毎ゲーム狙える普通の当たり"),
    ("属性オーブ", "+2個", "C_GREEN", "属性一致時に光る。パーティ強化のヒント"),
    ("虹オーブ", "+3個 ★", "C_GOLD", "最上位。10連ガチャに一気に近づく大チャンス"),
]
c_map = {"C_BLUE": C_BLUE, "C_GREEN": C_GREEN, "C_GOLD": C_GOLD}
for i, (name, val, col_name, desc) in enumerate(orbs):
    y = Inches(1.75) + i * Inches(0.55)
    col = c_map[col_name]
    add_rect(s5, Inches(6.9), y, Inches(1.1), Inches(0.42), fill=col)
    add_text(s5, name, Inches(6.9), y, Inches(1.1), Inches(0.42),
             size=10, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    add_text(s5, val, Inches(8.1), y, Inches(0.8), Inches(0.42),
             size=11, bold=True, color=col)
    add_text(s5, desc, Inches(9.0), y, Inches(3.8), Inches(0.42),
             size=9, color=C_LGRAY)

# 右下：ガチャ演出の比較
add_rect(s5, Inches(6.8), Inches(3.55), Inches(6.15), Inches(3.5),
         fill=RGBColor(0x12, 0x08, 0x00), line=C_ORANGE)
add_text(s5, "ガチャ演出の2ルート",
         Inches(6.9), Inches(3.6), Inches(5.5), Inches(0.4),
         size=13, bold=True, color=C_ORANGE)

routes = [
    ("単発ガチャ（5個）",
     "→ オーブ5個到達で発動\n→ 1体排出・演出短め\n→ ハズレC/Dランクでも積み重ね確認できる\n→ 期待感は小〜中"),
    ("10連ガチャ（50個）",
     "→ 50個一気に溜まると自動発動\n→ 10体一覧演出（大型演出）\n→ 上位4体がパーティ候補としてハイライト\n→ 映像・BGMで大盛り上がり"),
]
for i, (t, b) in enumerate(routes):
    y = Inches(4.05) + i * Inches(1.45)
    add_text(s5, t, Inches(6.9), y, Inches(5.8), Inches(0.4),
             size=12, bold=True, color=C_GOLD)
    add_text(s5, b, Inches(6.9), y + Inches(0.38), Inches(5.8), Inches(1.0),
             size=10, color=C_LGRAY)

# キャラ
add_picture_safe(s5, img("elle.png"), Inches(11.5), Inches(3.3), h=Inches(3.7))

# ===== SLIDE 6: 仲間集め → CZ突入 =====
s6 = prs.slides.add_slide(blank)
add_bg(s6, C_NAVY)
add_header_bar(s6, "仲間集め → CZ突入", "ガチャで4キャラが揃ったとき、ストーリーが動き出す")

# 中央：4キャラ集合図
add_text(s6, "パーティ4人 揃え完了！", Inches(0.3), Inches(1.25), Inches(12.7), Inches(0.6),
         size=20, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

char_imgs = [
    ("ren.png",      "焔 レン\nリーダー",  C_RED),
    ("char_101.png", "アオイ\nスピード型", C_BLUE),
    ("char_103.png", "ミナミ\nサポート型", C_GREEN),
    ("char_104.png", "マナ\nパワー型",     C_GOLD),
]
char_w = Inches(2.2)
char_x_start = Inches(0.4)
char_gap = Inches(0.2)

for i, (fn, lbl, c) in enumerate(char_imgs):
    x = char_x_start + i * (char_w + char_gap)
    add_picture_safe(s6, img(fn), x, Inches(1.75), h=Inches(3.5))
    add_rect(s6, x, Inches(5.25), char_w, Inches(0.5), fill=c)
    add_text(s6, lbl, x, Inches(5.25), char_w, Inches(0.5),
             size=10, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

# 矢印 → CZ
ax_x = Inches(10.2)
add_rect(s6, ax_x - Inches(0.05), Inches(2.8), Pt(3), Inches(1.2), fill=C_ORANGE)
sh = s6.shapes.add_shape(13, ax_x, Inches(2.6), Inches(0.9), Inches(0.45))
sh.fill.solid(); sh.fill.fore_color.rgb = C_ORANGE; sh.line.fill.background()

# CZボックス
add_rect(s6, Inches(11.2), Inches(1.8), Inches(1.8), Inches(4.5),
         fill=RGBColor(0x33, 0x00, 0x44), line=C_ORANGE)
add_text(s6, "CZ\n突入！", Inches(11.2), Inches(1.9), Inches(1.8), Inches(1.0),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
add_text(s6,
         "▶ マルチクエスト体験\n▶ 4人でボス撃破\n▶ AT確定ルート\n▶ 突破率：30〜70%\n　（パーティランクで変動）",
         Inches(11.25), Inches(2.9), Inches(1.7), Inches(3.2),
         size=9, color=C_LGRAY)

# 下部：CZの2パターン
add_rect(s6, Inches(0.35), Inches(6.1), Inches(10.4), Inches(1.1),
         fill=RGBColor(0x08, 0x18, 0x28), line=C_BLUE)
add_text(s6, "CZパターン", Inches(0.5), Inches(6.1), Inches(2.5), Inches(0.4),
         size=12, bold=True, color=C_BLUE)
cz_pts = [
    "■ 通常CZ（単発ガチャ4回分）：ランクCDEF4体。突破率低め。強化感薄い。",
    "■ 強CZ（10連ガチャ経由）　：ランク上位4体。突破率高い。フルムービー演出。",
]
for i, t in enumerate(cz_pts):
    add_text(s6, t, Inches(0.5), Inches(6.5) + i * Inches(0.32),
             Inches(10.0), Inches(0.32),
             size=10, color=C_LGRAY)

# ===== SLIDE 7: AT設計 映画×覇者の塔 =====
s7 = prs.slides.add_slide(blank)
add_bg(s7, C_NAVY)
add_header_bar(s7, "AT設計 — 映画 × 覇者の塔", "2本の映画ストーリーを駆け上がり、最頂点「轟絶」で覇者の塔へ")

# 映画1 エリア
add_rect(s7, Inches(0.35), Inches(1.25), Inches(5.8), Inches(5.8),
         fill=RGBColor(0x08, 0x18, 0x10), line=C_GREEN)
add_text(s7, "AT前半 — 映画①「はじまりの場所へ」",
         Inches(0.4), Inches(1.28), Inches(5.6), Inches(0.4),
         size=11, bold=True, color=C_GREEN)

movie1_scenes = [
    ("1F", "OP", "レンとオラゴンの出会い"),
    ("2F", "試練1", "初めてのモンスト対戦"),
    ("3F", "試練2", "仲間との絆・作戦会議"),
    ("4F", "ピンチ", "ゲノム軍の侵攻・窮地"),
    ("5F", "逆転", "4人の力を合わせて反撃"),
    ("6F", "クライマックス", "ゲノムとの最終決戦"),
    ("7F", "エンディング", "オラゴンとの別れ・感動ED"),
]
for i, (floor, tag, desc) in enumerate(movie1_scenes):
    y = Inches(1.75) + i * Inches(0.68)
    is_peak = (i == 3 or i == 6)
    bg = RGBColor(0x33, 0x00, 0x00) if is_peak else RGBColor(0x0C, 0x22, 0x18)
    add_rect(s7, Inches(0.4), y, Inches(5.6), Inches(0.63), fill=bg)
    add_text(s7, floor, Inches(0.42), y + Inches(0.08), Inches(0.5), Inches(0.4),
             size=10, bold=True, color=C_LGRAY)
    add_text(s7, tag, Inches(0.9), y + Inches(0.08), Inches(1.3), Inches(0.4),
             size=10, bold=True, color=C_GREEN if not is_peak else C_RED)
    add_text(s7, desc, Inches(2.2), y + Inches(0.1), Inches(3.7), Inches(0.42),
             size=9.5, color=C_LGRAY)

# 映画2 エリア
add_rect(s7, Inches(6.3), Inches(1.25), Inches(5.8), Inches(5.8),
         fill=RGBColor(0x08, 0x10, 0x28), line=C_BLUE)
add_text(s7, "AT後半 — 映画②「ソラノカナタ」",
         Inches(6.35), Inches(1.28), Inches(5.6), Inches(0.4),
         size=11, bold=True, color=C_BLUE)

movie2_scenes = [
    ("8F",  "OP",      "カナタと謎の少女ソラの出会い"),
    ("9F",  "謎",      "センジュの力・世界の異変"),
    ("10F", "仲間",    "マナたちと合流・再結成"),
    ("11F", "激闘",    "爆絶クラスのモンスター登場"),
    ("12F", "ピンチ",  "センジュ覚醒・全滅寸前"),
    ("13F", "逆転",    "オーブの力で仲間を取り戻す"),
    ("14F", "轟絶",    "覇者の塔 最終決戦\n★最高射幸体験★"),
]
for i, (floor, tag, desc) in enumerate(movie2_scenes):
    y = Inches(1.75) + i * Inches(0.68)
    is_peak = (i >= 5)
    bg = RGBColor(0x22, 0x00, 0x44) if is_peak else RGBColor(0x0C, 0x18, 0x2C)
    add_rect(s7, Inches(6.35), y, Inches(5.6), Inches(0.63), fill=bg)
    add_text(s7, floor, Inches(6.38), y + Inches(0.08), Inches(0.7), Inches(0.4),
             size=10, bold=True, color=C_LGRAY)
    add_text(s7, tag, Inches(7.05), y + Inches(0.08), Inches(1.3), Inches(0.4),
             size=10, bold=True, color=C_BLUE if not is_peak else C_GOLD)
    add_text(s7, desc, Inches(8.35), y + Inches(0.08), Inches(3.4), Inches(0.52),
             size=9.5, color=C_LGRAY)

# 中央分岐ライン
add_rect(s7, Inches(6.1), Inches(1.25), Pt(3), Inches(5.8), fill=C_ORANGE)
add_text(s7, "→", Inches(5.85), Inches(3.8), Inches(0.5), Inches(0.5),
         size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# 感情グラフ説明
add_rect(s7, Inches(0.35), Inches(7.05), Inches(11.75), Inches(0.35),
         fill=RGBColor(0x18, 0x0A, 0x00), line=C_ORANGE)
add_text(s7,
         "感情設計：1F〜7F（希望→絶望→逆転）→ 8F〜14F（謎→恐怖→覚醒→感動クライマックス）",
         Inches(0.5), Inches(7.07), Inches(11.5), Inches(0.3),
         size=9.5, color=C_GOLD)

# ===== SLIDE 8: スペック概要 =====
s8 = prs.slides.add_slide(blank)
add_bg(s8, C_NAVY)
add_header_bar(s8, "スペック概要", "数値イメージ / ベンチマーク機種との比較")

# スペックテーブル
add_text(s8, "想定スペック（イメージ）",
         Inches(0.35), Inches(1.25), Inches(7), Inches(0.4),
         size=14, bold=True, color=C_ORANGE)

specs = [
    ("機種タイプ",    "スマスロ（5.9号機相当）"),
    ("純増",          "約3.0〜4.5枚/G（通常AT）\n　8〜12枚/G（覇者の塔区間）"),
    ("AT初当たり",    "1/250〜1/350（設定1〜6）"),
    ("AT期待枚数",    "700〜1,200枚（超絶AT基本）"),
    ("CZ突破率",      "30%〜70%（パーティランクに連動）"),
    ("天井",          "999G（AT非当選時）"),
    ("設定差ポイント","CZ突破率・オーブ図柄当選率・\n覇者の塔突入率"),
    ("コイン持ち",    "50〜60G/50枚（通常時）"),
]

header_cols = ["項目", "スペック値"]
col_w = [Inches(3.0), Inches(4.5)]
hx = [Inches(0.35), Inches(3.4)]

for j, (hdr, cw) in enumerate(zip(header_cols, col_w)):
    add_rect(s8, hx[j], Inches(1.7), cw, Inches(0.4), fill=C_FILL_BL)
    add_text(s8, hdr, hx[j], Inches(1.7), cw, Inches(0.4),
             size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for i, (k, v) in enumerate(specs):
    y = Inches(2.15) + i * Inches(0.55)
    bg = RGBColor(0x0A, 0x1C, 0x2C) if i % 2 == 0 else RGBColor(0x0E, 0x22, 0x34)
    add_rect(s8, hx[0], y, col_w[0], Inches(0.52), fill=bg)
    add_rect(s8, hx[1], y, col_w[1], Inches(0.52), fill=bg)
    add_text(s8, k, hx[0] + Inches(0.08), y + Inches(0.05), col_w[0] - Inches(0.1), Inches(0.45),
             size=11, bold=True, color=C_LGRAY)
    add_text(s8, v, hx[1] + Inches(0.08), y + Inches(0.04), col_w[1] - Inches(0.1), Inches(0.5),
             size=10.5, color=C_WHITE)

# ベンチマーク
add_rect(s8, Inches(8.1), Inches(1.25), Inches(4.9), Inches(5.8),
         fill=RGBColor(0x08, 0x18, 0x28), line=C_BLUE)
add_text(s8, "ベンチマーク機種",
         Inches(8.2), Inches(1.3), Inches(4.5), Inches(0.4),
         size=13, bold=True, color=C_BLUE)

benchmarks = [
    ("スマスロ からくりサーカス",
     "純増3.8枚・AT出玉の山谷設計\nIP×感動ストーリーの成功例\n→ 感情設計の参考にする"),
    ("スマスロ 北斗の拳",
     "タワー型AT×高純増区間\n長尺ATの緊張感構造\n→ 覇者の塔AT設計の参考"),
    ("スマスロ モンスターハンターライズ",
     "仲間集め→大型クエスト型\nアクションゲームIP × マルチ感\n→ CZ設計・通常時演出の参考"),
]
for i, (title, desc) in enumerate(benchmarks):
    y = Inches(1.8) + i * Inches(1.7)
    add_rect(s8, Inches(8.2), y, Inches(4.7), Inches(1.6),
             fill=RGBColor(0x0C, 0x22, 0x38))
    add_text(s8, f"【{i+1}】{title}", Inches(8.3), y + Inches(0.08),
             Inches(4.5), Inches(0.42), size=11, bold=True, color=C_GOLD)
    add_text(s8, desc, Inches(8.3), y + Inches(0.48), Inches(4.5), Inches(1.0),
             size=9.5, color=C_LGRAY)

# キャラ（右下隅に小さく配置）
add_picture_safe(s8, img("char_105.png"), Inches(12.0), Inches(4.5), h=Inches(2.8))

# ===== SLIDE 9: まとめ / 強み =====
s9 = prs.slides.add_slide(blank)
add_bg(s9, C_NAVY)
add_header_bar(s9, "まとめ — 3つの強み", "スマスロ モンスターストライク が成立する理由")

# 強み3点
strengths = [
    (C_ORANGE, "① IPの記憶力",
     "モンストは2013〜2018年に最大ユーザー数5,000万人を記録。\n今の30代男性の多くが強烈な遊び体験を持つ。\n「懐かしさ」は最強の導線。"),
    (C_BLUE,   "② 設計の親和性",
     "オーブ（ガチャ通貨）＝コインという概念の置き換えが自然。\nマルチ体験→CZ突入・覇者の塔→タワーATと\nゲームIPの構造がスロットと素直に対応する。"),
    (C_GREEN,  "③ 感情の山・谷",
     "映画2作品にはドラマ的な起承転結がある。\nAT前半（映画1）→後半（映画2）→覇者の塔（頂点）で\n\"絶望→逆転→達成感\"の感情曲線を描ける。"),
]

for i, (color, title, body) in enumerate(strengths):
    y = Inches(1.35) + i * Inches(1.85)
    add_rect(s9, Inches(0.35), y, Inches(8.3), Inches(1.78),
             fill=RGBColor(0x0A, 0x1A, 0x28), line=color)
    add_rect(s9, Inches(0.35), y, Inches(0.12), Inches(1.78), fill=color)
    add_text(s9, title, Inches(0.6), y + Inches(0.1), Inches(7.8), Inches(0.5),
             size=18, bold=True, color=color)
    add_text(s9, body, Inches(0.6), y + Inches(0.58), Inches(7.8), Inches(1.1),
             size=11, color=C_LGRAY)

# ボトムまとめ帯
add_rect(s9, Inches(0.35), Inches(6.85), Inches(8.3), Inches(0.4),
         fill=RGBColor(0x22, 0x11, 0x00), line=C_GOLD)
add_text(s9, "\"懐かしさ\" × \"射幸性\" × \"感情ストーリー\" の三位一体が、スマスロ モンスターストライクの核心",
         Inches(0.45), Inches(6.86), Inches(8.1), Inches(0.38),
         size=10, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# 右：全キャラ集合
add_picture_safe(s9, img("ren.png"),      Inches(8.6),  Inches(0.6), h=Inches(5.5))
add_picture_safe(s9, img("char_101.png"), Inches(10.3), Inches(1.1), h=Inches(4.8))
add_picture_safe(s9, img("char_103.png"), Inches(11.6), Inches(1.4), h=Inches(4.5))
add_picture_safe(s9, img("oragon3.png"),  Inches(8.9),  Inches(4.8), h=Inches(2.4))

add_text(s9, "提案：スマスロ モンスターストライク", Inches(8.5), Inches(6.9), Inches(4.5), Inches(0.4),
         size=10, italic=True, color=C_LGRAY, align=PP_ALIGN.CENTER)

# ===== 保存 =====
prs.save(OUT_PATH)
print(f"保存完了: {OUT_PATH}")
