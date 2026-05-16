"""
スマスロ モンキーターンV 機種分析資料 PowerPointジェネレーター
出力: proposals/機種分析/モンキーターンV/monkeyturn_guide_v1.pptx
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(ROOT_DIR,
           "proposals", "機種分析", "モンキーターンV", "monkeyturn_guide_v2.pptx")

# ── カラーパレット（海×速さ：深海紺×ティール×オレンジ）──────────
C_BG    = RGBColor(0x04, 0x0A, 0x1C)
C_CARD  = RGBColor(0x0C, 0x14, 0x2C)
C_WATER = RGBColor(0x14, 0x88, 0x7C)   # ティール（水面）
C_CYAN  = RGBColor(0x00, 0xC0, 0xC8)   # 明るいシアン
C_SPEED = RGBColor(0xFF, 0xA0, 0x20)   # オレンジ（スピード）
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM = RGBColor(0xE0, 0xE8, 0xF0)   # 水色寄りクリーム
C_LTGRY = RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY  = RGBColor(0x88, 0x88, 0x88)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_RED   = RGBColor(0xCC, 0x22, 0x22)
C_PUR   = RGBColor(0x88, 0x44, 0xCC)
C_NAVY  = RGBColor(0x04, 0x0A, 0x1C)
C_STEEL = RGBColor(0x40, 0x60, 0x9A)   # 通常時：スチール青（最冷）
C_LIME  = RGBColor(0x22, 0xAA, 0x44)   # CZ：チャンス緑
C_FIRE  = RGBColor(0xFF, 0x44, 0x00)   # エキシビション：炎赤橙
C_ELEC  = RGBColor(0x00, 0x7A, 0xFF)   # 青島SG：電気青（最熱）

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（深海紺×スピードライン×底部ティールグロー）──────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 10, 28))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 50):
        draw.line([(i, 0), (0, i)], fill=(8, 18, 45), width=1)
    for y in range(h - 110, h):
        t = (y - (h - 110)) / 110
        draw.line([(0, y), (w, y)], fill=(0, int(40 * t), int(65 * t)))
    for y in range(0, 55):
        t = (55 - y) / 55 * 0.4
        draw.line([(0, y), (w, y)], fill=(0, int(25 * t), int(45 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー ──────────────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s

def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT_B
    if color:
        run.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp

def arrow_r(slide, x, y, w, color):
    h = Emu(150000)
    shp = slide.shapes.add_shape(13, x, y - h // 2, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def hdr(slide, text):
    rect(slide, Inches(0.15), Inches(0.08), Inches(9.7), Emu(420000),
         RGBColor(0x04, 0x10, 0x28))
    rect(slide, Inches(0.15), Inches(0.08), Emu(60000), Emu(420000), C_WATER)
    tb(slide, Inches(0.4), Inches(0.1), Inches(9.2), Emu(380000),
       text, 12, bold=True, color=C_GOLD, font=FONT_H)

def net_note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(200000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)

def kv_row(slide, x, y, w, key, val, key_w=Inches(1.4), col_key=None, col_val=None, bg=None):
    if bg:
        rect(slide, x, y, w, Emu(295000), bg)
    tb(slide, x + Emu(50000), y + Emu(28000), key_w - Emu(80000), Emu(245000),
       key, 8.5, bold=True, color=col_key or C_WATER, wrap=False)
    tb(slide, x + key_w, y + Emu(28000), w - key_w - Emu(50000), Emu(245000),
       val, 8.5, color=col_val or C_CREAM, wrap=False)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), Inches(5.6), SLIDE_H, RGBColor(0x02, 0x08, 0x18))
    rect(s, Inches(0.35), Inches(0.5), Emu(30000), Inches(2.0), C_WATER)

    tb(s, Inches(0.5), Inches(0.52), Inches(5), Emu(360000),
       "機種分析資料", 14, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.5), Inches(1.05), Inches(5.2), Emu(900000),
       "モンキーターンV", 38, bold=True, color=C_WHITE, font=FONT_H)
    tb(s, Inches(0.5), Inches(2.22), Inches(5.0), Emu(300000),
       "スマスロ（山佐）── 複層設計と継続率管理の完成形", 10, color=C_LTGRY, font=FONT_H)

    rect(s, Inches(0.5), Inches(2.78), Inches(4.8), Emu(580000), RGBColor(0x04, 0x16, 0x24))
    rect(s, Inches(0.5), Inches(2.78), Emu(55000), Emu(580000), C_WATER)
    tb(s, Inches(0.65), Inches(2.86), Inches(4.5), Emu(510000),
       "「遊びやすさとパンチの両立」\n\n複層モード × シナリオ14種 × 上位AT青島SG\nゲーム性の完成度が高いと評価される一台", 10, color=C_WHITE)

    rect_b(s, Inches(5.8), Inches(0.75), Inches(3.9), Inches(2.7), C_CARD, C_WATER, 1.5)
    tb(s, Inches(5.95), Inches(0.82), Inches(3.6), Emu(320000),
       "基本スペック", 10, bold=True, color=C_WATER, font=FONT_H)
    specs = [
        ("タイプ",   "スマスロ 6.5号機 AT機"),
        ("メーカー", "山佐"),
        ("純増",     "約2.5枚/G（AT）/ 4.0枚/G（青島SG）"),
        ("天井",     "795G+α（設定変更後 495G+α）"),
        ("機械割",   "設定1: 97.9%  /  設定6: 114.9%"),
    ]
    qy = Inches(1.30)
    for k, v in specs:
        tb(s, Inches(5.95), qy, Inches(1.35), Emu(255000), k, 8.5, bold=True, color=C_GRAY, wrap=False)
        tb(s, Inches(7.35), qy, Inches(2.25), Emu(255000), v, 8.5, color=C_CREAM, wrap=False)
        qy += Emu(263000)

    rect_b(s, Inches(5.8), Inches(3.6), Inches(3.9), Inches(1.55),
           RGBColor(0x08, 0x18, 0x28), C_SPEED, 1.0)
    tb(s, Inches(5.95), Inches(3.68), Inches(3.6), Emu(280000),
       "3つの分析ポイント", 9, bold=True, color=C_SPEED, font=FONT_H)
    tb(s, Inches(5.95), Inches(4.02), Inches(3.6), Emu(650000),
       "I    激走ポイント × 周期 × ライバルモード\nII   14種シナリオ × 示唆システム\nIII  グランドスラム → 青島SG への昇格設計",
       8.5, color=C_CREAM)

    tb(s, Inches(7.5), Inches(5.2), Inches(2.3), Emu(300000),
       "v1.0  2026.05", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図  ──  メインゴールはグランドスラム（8セット完走）")

    BW  = Inches(2.5)
    GAP = Inches(0.22)
    FBH = Inches(1.65)   # 通常時・AT の縦幅
    HBH = Inches(0.76)   # 優出モード・CZ の縦幅
    HG  = Emu(80000)     # 優出モード⇔CZ の隙間
    R1Y = Inches(0.62)
    R2Y = Inches(2.62)
    X1  = Inches(0.18)
    X2  = X1 + BW + GAP
    X3  = X2 + BW + GAP
    CZ_TOP = R1Y + HBH + HG
    CZ_MID = CZ_TOP + HBH // 2

    # ── 通常時（full box）─────────────────────────────
    rect_b(s, X1, R1Y, BW, FBH, RGBColor(0x06, 0x10, 0x22), C_STEEL, 1.8)
    tb(s, X1 + Emu(80000), R1Y + Emu(50000), BW - Emu(160000), Emu(270000),
       "通常時", 10, bold=True, color=C_WHITE, font=FONT_H)
    tb(s, X1 + Emu(80000), R1Y + Emu(340000), BW - Emu(160000), Emu(900000),
       "激走ポイント毎G加算\n\n① 222pt到達（周期天井）\n   → 前兆「優出モード」\n\n② レア役直撃\n   → CZ「超抜チャレンジ」",
       8, color=C_CREAM)

    # ── 優出モード（upper half）─────────────────────────
    rect_b(s, X2, R1Y, BW, HBH, RGBColor(0x00, 0x12, 0x26), C_CYAN, 1.8)
    tb(s, X2 + Emu(80000), R1Y + Emu(42000), BW - Emu(160000), Emu(260000),
       "前兆「優出モード」", 9.5, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, X2 + Emu(80000), R1Y + Emu(320000), BW - Emu(160000), HBH - Emu(355000),
       "周期天井（222pt等）で突入 / AT当選確定的", 8, color=C_CREAM)

    # ── CZ「超抜チャレンジ」（lower half）────────────────
    rect_b(s, X2, CZ_TOP, BW, HBH, RGBColor(0x04, 0x18, 0x08), C_LIME, 1.8)
    tb(s, X2 + Emu(80000), CZ_TOP + Emu(42000), BW - Emu(160000), Emu(260000),
       "CZ「超抜チャレンジ」", 9.5, bold=True, color=C_LIME, font=FONT_H)
    tb(s, X2 + Emu(80000), CZ_TOP + Emu(320000), BW - Emu(160000), HBH - Emu(355000),
       "ベル/リプ/レア役カウントアップ\n9カウントで強制成功", 8, color=C_CREAM)

    # ── AT「SG RUSH」（full box）──────────────────────────
    rect_b(s, X3, R1Y, BW, FBH, RGBColor(0x0C, 0x16, 0x04), C_SPEED, 1.8)
    tb(s, X3 + Emu(80000), R1Y + Emu(50000), BW - Emu(160000), Emu(270000),
       "AT「SG RUSH」", 10, bold=True, color=C_WHITE, font=FONT_H)
    tb(s, X3 + Emu(80000), R1Y + Emu(340000), BW - Emu(160000), Emu(900000),
       "純増2.5枚/G\n周回31G + SGバトル9G\n\nメインゴール：\n8セット完走\n= グランドスラム達成！",
       8, color=C_CREAM)

    # 矢印：通常時 → 優出モード
    arrow_r(s, X1 + BW + Emu(40000), R1Y + HBH // 2, GAP - Emu(80000), C_CYAN)
    # 矢印：通常時 → CZ
    arrow_r(s, X1 + BW + Emu(40000), CZ_MID, GAP - Emu(80000), C_LIME)
    # 矢印：優出モード → AT（成功）
    arrow_r(s, X2 + BW + Emu(40000), R1Y + HBH // 2, GAP - Emu(80000), C_GREEN)
    tb(s, X2 + BW + Emu(8000), R1Y + HBH // 2 - Emu(220000), GAP + Emu(50000), Emu(195000),
       "成功", 7, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    # 矢印：CZ → AT（成功）
    arrow_r(s, X2 + BW + Emu(40000), CZ_MID, GAP - Emu(80000), C_GREEN)

    # ⊓ コネクター（AT → Row2）
    CON_X = X3 + BW + Emu(80000)
    CON_R = CON_X + Emu(550000)
    LW    = Emu(55000)
    AT_MID = R1Y + FBH // 2
    MID_Y  = (R1Y + FBH + R2Y) // 2
    rect(s, CON_X, AT_MID, LW, MID_Y - AT_MID + Emu(28000), C_GOLD)
    rect(s, CON_X, MID_Y, CON_R - CON_X + LW, LW, C_GOLD)
    rect(s, CON_R, MID_Y, LW, R2Y - MID_Y + LW, C_GOLD)
    tb(s, CON_X - Emu(55000), MID_Y - Emu(380000), Emu(790000), Emu(360000),
       "8セット\n完走", 8, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)

    # ── Row2（R→L）: グランドスラム → エキシビション → 青島SG ─────
    BH2 = Inches(0.88)
    row2 = [
        (X3, "グランドスラム達成！",
         "8完走でメインゴール到達\n→ エキシビション（ボーナス）へ",
         RGBColor(0x26, 0x16, 0x00), C_GOLD),
        (X2, "エキシビションレース",
         "グランドスラム後のボーナス区間\n継続率50/66/80/90%に固定",
         RGBColor(0x20, 0x08, 0x00), C_FIRE),
        (X1, "上位AT「青島SG」",
         "真の頂点 / 純増4.0枚/G\n継続率83% / 温泉モードあり",
         RGBColor(0x00, 0x0A, 0x22), C_ELEC),
    ]
    for x, title, desc, fill, bdr in row2:
        rect_b(s, x, R2Y, BW, BH2, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R2Y + Emu(42000), BW - Emu(160000), Emu(265000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R2Y + Emu(322000), BW - Emu(160000), BH2 - Emu(360000),
           desc, 8, color=C_CREAM)
    for x_r in [X3, X2]:
        _w = GAP - Emu(80000)
        _h = Emu(150000)
        shp = s.shapes.add_shape(13, x_r - GAP + Emu(40000), R2Y + BH2 // 2 - _h // 2, _w, _h)
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = C_FIRE
        shp.line.fill.background()

    BOT_Y = R2Y + BH2 + Emu(110000)
    rect(s, X1, BOT_Y, Inches(9.64), Emu(820000), RGBColor(0x06, 0x10, 0x20))
    rect(s, X1, BOT_Y, Emu(55000), Emu(820000), C_WATER)
    tb(s, X1 + Emu(100000), BOT_Y + Emu(60000), Inches(9.2), Emu(290000),
       "フロー設計の核心：「グランドスラム（8完走）を目指す」ことが来店動機のすべてになる",
       9, bold=True, color=C_WATER, font=FONT_H)
    tb(s, X1 + Emu(100000), BOT_Y + Emu(375000), Inches(9.2), Emu(405000),
       "優出モード（周期天井）とCZ直撃の2ルートでAT突入。AT後はシナリオ14種で継続率が決まり8完走を狙う。\n"
       "CZ失敗の1.2%でGSC（AT確定・艇王シナリオ）という逆転も存在。失敗にも次がある設計。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時設計（激走ポイント × 周期 × ライバルモード）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時設計  ──  激走ポイント × 周期 × ライバルモード 3層構造")

    # 左：激走ポイント
    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.1), Inches(3.5), C_CARD, C_CYAN, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(2.9), Emu(320000),
       "激走ポイント", 11, bold=True, color=C_CYAN, font=FONT_H)
    pts = [
        (C_CYAN,  "毎G加算",       "通常は1pt以上が毎ゲーム加算\n液晶右上のメーターで可視化"),
        (C_SPEED, "激走チャージ",  "チャンス役で5G突入\n10〜200ptの大量加算ゾーン"),
        (C_GOLD,  "222pt到達",     "優出モード突入（AT結果確定直前）\n1周期目は平均80Gで到達"),
        (C_GREEN, "チャージ天井",  "81G or 131G無チャージ\n→ 次チャンス役でCZ確定"),
    ]
    py = Inches(1.38)
    for col, key, val in pts:
        rect(s, Inches(0.28), py, Emu(50000), Emu(510000), col)
        tb(s, Inches(0.48), py + Emu(25000), Inches(2.6), Emu(240000),
           key, 9, bold=True, color=col)
        tb(s, Inches(0.48), py + Emu(265000), Inches(2.6), Emu(260000),
           val, 8, color=C_CREAM)
        py += Emu(555000)

    # 中央：内部モード（周期）
    rect_b(s, Inches(3.45), Inches(0.85), Inches(3.1), Inches(3.5), C_CARD, C_WATER, 1.5)
    tb(s, Inches(3.57), Inches(0.92), Inches(2.9), Emu(320000),
       "内部モード（周期管理）", 11, bold=True, color=C_WATER, font=FONT_H)
    modes = [
        (C_GOLD2, "天国モード",   "1周期目にAT当選濃厚\n最上位・最もレアなモード"),
        (C_WATER, "通常Aモード",  "6周期天井 / 2・5周期がチャンス\n最も標準的な推移"),
        (C_CYAN,  "通常Bモード",  "3周期天井 / 2周期がチャンス\n短めの天井が特徴"),
    ]
    my = Inches(1.38)
    for col, name, desc in modes:
        rect_b(s, Inches(3.53), my, Inches(2.88), Emu(760000), RGBColor(0x08, 0x14, 0x24), col, 1.0)
        tb(s, Inches(3.65), my + Emu(50000), Inches(2.6), Emu(280000),
           name, 9.5, bold=True, color=col)
        tb(s, Inches(3.65), my + Emu(330000), Inches(2.6), Emu(380000),
           desc, 8.5, color=C_CREAM)
        my += Emu(800000)

    # 右：ライバルモード
    rect_b(s, Inches(6.7), Inches(0.85), Inches(3.1), Inches(3.5), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(6.82), Inches(0.92), Inches(2.9), Emu(320000),
       "ライバルモード（6種）", 11, bold=True, color=C_SPEED, font=FONT_H)
    rivals = [
        (C_LTGRY, "榎木・蒲生・浜岡",
         "通常時の周期抽選に作用\n設定差あり・高設定ほど移行率UP"),
        (C_SPEED, "洞口・青島・波多野",
         "AT中シナリオ性能に作用\n全設定共通移行率"),
        (C_GOLD2, "共通の恩恵",
         "滞在中は出玉率100%超\n当選まで打ち続けたい設計"),
        (C_CYAN,  "青島モードは別格",
         "特に強いとされるモード\nAT中継続率が大幅優遇"),
    ]
    ry = Inches(1.38)
    for col, name, desc in rivals:
        rect(s, Inches(6.78), ry, Emu(50000), Emu(510000), col)
        tb(s, Inches(6.98), ry + Emu(25000), Inches(2.6), Emu(240000),
           name, 9, bold=True, color=col)
        tb(s, Inches(6.98), ry + Emu(265000), Inches(2.6), Emu(260000),
           desc, 8, color=C_CREAM)
        ry += Emu(555000)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(560000), RGBColor(0x06, 0x10, 0x20))
    rect(s, Inches(0.2), Inches(4.45), Emu(55000), Emu(560000), C_SPEED)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.1), Emu(250000),
       "3層構造の設計：「内部モードでいつ当たるか」＋「ポイントで天井を短縮」＋「ライバルモードで出玉保証」",
       9.5, bold=True, color=C_SPEED)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.1), Emu(250000),
       "どのモードに滞在しているかを意識することで、立ち回りの深みが生まれる。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ「超抜チャレンジ」
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ「超抜チャレンジ」  ──  結末は3通り")

    # ── 左：基本仕様 ─────────────────────────────────────
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.45), Inches(3.45), C_CARD, C_LIME, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.1), Emu(320000),
       "基本仕様", 11, bold=True, color=C_LIME, font=FONT_H)
    items = [
        ("突入契機",  "レア役成立時の直撃抽選",              C_CREAM),
        ("ゲーム数",  "10G（毎G成功抽選が並行）",          C_WHITE),
        ("成功率",    "約50%（強役ほど高確率）",            C_SPEED),
        ("9カウント", "ベル/リプ/レア役の合計9個→強制成功", C_GOLD),
        ("失敗",      "通常時に戻る（次の周期へ）",          C_LTGRY),
    ]
    iy = Inches(1.38)
    for j, (k, v, col) in enumerate(items):
        bg = RGBColor(0x0A, 0x14, 0x28) if j % 2 == 0 else RGBColor(0x0E, 0x18, 0x2E)
        rect(s, Inches(0.28), iy, Inches(4.27), Emu(315000), bg)
        tb(s, Inches(0.38), iy + Emu(28000), Inches(1.35), Emu(265000),
           k, 8.5, bold=True, color=C_LIME, wrap=False)
        tb(s, Inches(1.78), iy + Emu(28000), Inches(2.6), Emu(265000),
           v, 8.5, color=col, wrap=False)
        iy += Emu(325000)

    # ── 右：3結末フロー ──────────────────────────────────
    RX = Inches(4.85)
    RW = Inches(4.95)

    # CZ概要ヘッダー
    rect_b(s, RX, Inches(0.88), RW, Emu(340000),
           RGBColor(0x04, 0x18, 0x08), C_LIME, 2.0)
    tb(s, RX + Emu(80000), Inches(0.94), RW - Emu(160000), Emu(280000),
       "超抜チャレンジ（10G）  ──  ベル9カウント or 毎G抽選で突破", 10, bold=True, color=C_LIME, font=FONT_H)

    FORK_Y = Inches(0.88) + Emu(340000) + Emu(100000)
    S_W    = Inches(2.2)
    F_X    = RX + S_W + Emu(60000)
    F_W    = RW - S_W - Emu(60000)
    TOTAL_H = Emu(1860000)

    # 成功（左）
    rect_b(s, RX, FORK_Y, S_W, TOTAL_H, RGBColor(0x00, 0x18, 0x06), C_GREEN, 2.0)
    tb(s, RX + Emu(80000), FORK_Y + Emu(55000), S_W - Emu(140000), Emu(300000),
       "成功（約50%）", 11, bold=True, color=C_GREEN, font=FONT_H)
    tb(s, RX + Emu(80000), FORK_Y + Emu(375000), S_W - Emu(140000), Emu(1280000),
       "AT「SG RUSH」突入\n\n第3停止を離す瞬間に\n14種シナリオが確定\n\n純増2.5枚/G\n8完走を目指す", 9, color=C_CREAM)

    # 失敗（通常）
    N_H = Emu(840000)
    rect_b(s, F_X, FORK_Y, F_W, N_H, RGBColor(0x14, 0x10, 0x10), C_GRAY, 1.2)
    tb(s, F_X + Emu(70000), FORK_Y + Emu(50000), F_W - Emu(120000), Emu(285000),
       "失敗（通常）約49%", 10, bold=True, color=C_LTGRY, font=FONT_H)
    tb(s, F_X + Emu(70000), FORK_Y + Emu(350000), F_W - Emu(120000), Emu(460000),
       "通常時に戻る\n次の周期でまた抽選", 9, color=C_GRAY)

    # グランドスラムCH
    GSC_Y = FORK_Y + N_H + Emu(80000)
    GSC_H = TOTAL_H - N_H - Emu(80000)
    rect_b(s, F_X, GSC_Y, F_W, GSC_H, RGBColor(0x1E, 0x14, 0x00), C_GOLD, 2.0)
    tb(s, F_X + Emu(70000), GSC_Y + Emu(50000), F_W - Emu(120000), Emu(285000),
       "失敗→GSC（約1.2%）", 10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, F_X + Emu(70000), GSC_Y + Emu(350000), F_W - Emu(120000), Emu(580000),
       "AT確定\n波多野フリーズ発生\n艇王シナリオ確定\n→ 8完走ほぼ確定", 9, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.42), Inches(9.6), Emu(580000), RGBColor(0x06, 0x10, 0x20))
    rect(s, Inches(0.2), Inches(4.42), Emu(55000), Emu(580000), C_WATER)
    tb(s, Inches(0.45), Inches(4.47), Inches(9.1), Emu(250000),
       "CZ設計のポイント：50%で失敗しても通常に戻るだけ、1.2%で逆転AT確定という絶妙なバランス",
       9.5, bold=True, color=C_WATER)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.1), Emu(250000),
       "「外れたと思ったら最高の展開」体験がCZへの参加し続けるモチベーションを維持させる。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT「SG RUSH」の構造
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「SG RUSH」  ──  2部構造 × グランドスラム達成設計")

    # 左：2部構造
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.5), Inches(3.45), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(320000),
       "1セットの2部構造", 11, bold=True, color=C_SPEED, font=FONT_H)

    # 周回パート
    rect(s, Inches(0.28), Inches(1.40), Inches(4.32), Emu(360000), RGBColor(0x10, 0x18, 0x06))
    rect(s, Inches(0.28), Inches(1.40), Emu(55000), Emu(360000), C_SPEED)
    tb(s, Inches(0.53), Inches(1.46), Inches(3.8), Emu(280000),
       "周回パート（31G）", 10, bold=True, color=C_SPEED)
    tb(s, Inches(0.53), Inches(1.78), Inches(3.8), Emu(300000),
       "チェリー → ぶっちぎりバトルのメイン契機\nバトル高確中は全役で突入抽選",
       8.5, color=C_CREAM)

    # SGバトルパート
    rect(s, Inches(0.28), Inches(2.30), Inches(4.32), Emu(360000), RGBColor(0x00, 0x10, 0x28))
    rect(s, Inches(0.28), Inches(2.30), Emu(55000), Emu(360000), C_CYAN)
    tb(s, Inches(0.53), Inches(2.36), Inches(3.8), Emu(280000),
       "SGバトルパート（9G）── 青島VS波多野", 10, bold=True, color=C_CYAN)
    tb(s, Inches(0.53), Inches(2.68), Inches(3.8), Emu(260000),
       "青島が勝利 → 次セット継続\n継続率はシナリオで管理（14種）",
       8.5, color=C_CREAM)

    # グランドスラム
    rect_b(s, Inches(0.28), Inches(2.98), Inches(4.32), Emu(680000),
           RGBColor(0x20, 0x16, 0x00), C_GOLD, 2.0)
    tb(s, Inches(0.40), Inches(3.04), Inches(4.0), Emu(290000),
       "8セット完走 = グランドスラム達成", 10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.40), Inches(3.38), Inches(4.0), Emu(320000),
       "セット継続のたびに積み上がる「8回勝つ」目標\n→ エキシビションレースへ自動移行",
       8.5, color=C_CREAM)

    # 右：継続率シナリオ概要
    rect_b(s, Inches(4.9), Inches(0.85), Inches(4.9), Inches(3.45), C_CARD, C_CYAN, 1.5)
    tb(s, Inches(5.02), Inches(0.92), Inches(4.6), Emu(320000),
       "シナリオ選択（14種）と継続率管理", 11, bold=True, color=C_CYAN, font=FONT_H)

    tb(s, Inches(5.02), Inches(1.38), Inches(4.6), Emu(280000),
       "AT突入時（第3停止を離す瞬間）に14種のシナリオが確定",
       9, color=C_CREAM)

    scenarios = [
        ("継続率段階", "シナリオ名（例）",          "特徴"),
        ("低（〜25%）",  "関東ガマシ 等",           "早期終了が多い"),
        ("中（50%）",    "ギャンブラー 等",          "標準的な継続"),
        ("高（66〜80%）","艇界のヒロイン 等",        "粘れるシナリオ"),
        ("最高（艇王）", "艇王（GS濃厚）",            "8完走ほぼ確定"),
    ]
    sy = Inches(1.78)
    for j, (a, b, c) in enumerate(scenarios):
        bg = RGBColor(0x10, 0x18, 0x28) if j == 0 else (
             RGBColor(0x0A, 0x14, 0x24) if j % 2 == 1 else RGBColor(0x0E, 0x18, 0x2C))
        rect(s, Inches(5.02), sy, Inches(4.7), Emu(315000), bg)
        col_a = C_CYAN if j == 0 else C_SPEED if "艇王" in a else C_LTGRY
        tb(s, Inches(5.1), sy + Emu(28000), Inches(1.5), Emu(265000), a, 8, bold=(j==0), color=col_a, wrap=False)
        tb(s, Inches(6.65), sy + Emu(28000), Inches(1.35), Emu(265000), b, 7.5, color=C_CREAM, wrap=False)
        tb(s, Inches(8.05), sy + Emu(28000), Inches(1.45), Emu(265000), c, 7.5, color=C_GRAY, wrap=False)
        sy += Emu(325000)

    tb(s, Inches(5.02), Inches(3.50), Inches(4.6), Emu(240000),
       "ハマりゲーム数・設定の影響なし。突入時に公平に抽選される。",
       8, color=C_GRAY)

    # フッター
    rect(s, Inches(0.2), Inches(4.42), Inches(9.6), Emu(580000), RGBColor(0x06, 0x10, 0x20))
    rect(s, Inches(0.2), Inches(4.42), Emu(55000), Emu(580000), C_SPEED)
    tb(s, Inches(0.45), Inches(4.47), Inches(9.1), Emu(250000),
       "AT設計の核心：「8回勝つ」という明確な目標がセット継続のドラマを作る",
       9.5, bold=True, color=C_SPEED)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.1), Emu(250000),
       "艇王シナリオ引きで「今日は行ける」という確信がプレイヤーに芽生える。シナリオ推測が遊技の深みになる。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: ぶっちぎりバトル × 連戦
# ══════════════════════════════════════════════════════════════
def s_battle(prs):
    s = new_slide(prs)
    hdr(s, "ぶっちぎりバトル × 連戦  ──  SGバトルを超える「横道」の自力感")

    # ── 左：ぶっちぎりバトル基本仕様 ─────────────────────────
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.45), Inches(3.45), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.1), Emu(320000),
       "ぶっちぎりバトル", 11, bold=True, color=C_SPEED, font=FONT_H)
    items = [
        ("突入契機",  "周回パートのチェリー主契機（高確=全役）", C_SPEED),
        ("内容",      "通常SGバトルと独立した自力バトル",        C_WHITE),
        ("勝率",      "約50%〜（強役ほど高確率）",              C_GOLD),
        ("恩恵A〜E",  "勝利後1Gで恩恵レベルを告知",             C_GREEN),
        ("Eが頂点",   "究極Vフリーズ（100G以上確定）\nレア役で300G以上に昇格", C_GOLD),
        ("連戦突入",  "赤=通常連戦 / 紫=裏連戦（80%ループ）",  C_CYAN),
    ]
    iy = Inches(1.38)
    for j, (k, v, col) in enumerate(items):
        bg = RGBColor(0x0A, 0x14, 0x28) if j % 2 == 0 else RGBColor(0x0E, 0x18, 0x2E)
        rect(s, Inches(0.28), iy, Inches(4.27), Emu(315000), bg)
        tb(s, Inches(0.38), iy + Emu(28000), Inches(1.5), Emu(265000),
           k, 8.5, bold=True, color=C_WATER, wrap=False)
        tb(s, Inches(1.93), iy + Emu(28000), Inches(2.4), Emu(265000),
           v, 8.5, color=col, wrap=False)
        iy += Emu(325000)

    # ── 右：連戦システム ─────────────────────────────────────
    RX = Inches(4.85)
    RW = Inches(4.95)

    rect_b(s, RX, Inches(0.88), RW, Emu(340000),
           RGBColor(0x14, 0x08, 0x00), C_SPEED, 2.0)
    tb(s, RX + Emu(80000), Inches(0.94), RW - Emu(160000), Emu(280000),
       "連戦システム  ──  エフェクト色が強さを決める", 10, bold=True, color=C_SPEED, font=FONT_H)

    FORK_Y = Inches(0.88) + Emu(340000) + Emu(100000)
    HALF_H = Emu(900000)

    # 赤エフェクト（通常連戦）
    rect_b(s, RX, FORK_Y, RW, HALF_H, RGBColor(0x1A, 0x04, 0x04), C_RED, 2.0)
    tb(s, RX + Emu(80000), FORK_Y + Emu(55000), RW - Emu(140000), Emu(290000),
       "赤エフェクト → 通常連戦", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, RX + Emu(80000), FORK_Y + Emu(350000), RW - Emu(140000), HALF_H - Emu(400000),
       "敗北するまでぶっちぎりバトルが継続。\n"
       "1バトルごとに追加報酬のチャンス。",
       9, color=C_CREAM)

    # 紫エフェクト（裏連戦）
    URA_Y = FORK_Y + HALF_H + Emu(80000)
    rect_b(s, RX, URA_Y, RW, HALF_H, RGBColor(0x12, 0x04, 0x1A), C_PUR, 2.0)
    tb(s, RX + Emu(80000), URA_Y + Emu(55000), RW - Emu(140000), Emu(290000),
       "紫エフェクト → 裏連戦", 11, bold=True, color=C_PUR, font=FONT_H)
    tb(s, RX + Emu(80000), URA_Y + Emu(350000), RW - Emu(140000), HALF_H - Emu(400000),
       "初回勝利がほぼ確定。\n"
       "以降は約80%ループで継続。\n"
       "「最強の横道」で大量報酬に期待。",
       9, color=C_CREAM)

    # 設計コメント
    CMT_Y = URA_Y + HALF_H + Emu(80000)
    rect_b(s, RX, CMT_Y, RW, Emu(620000), RGBColor(0x06, 0x10, 0x20), C_CYAN, 1.2)
    tb(s, RX + Emu(70000), CMT_Y + Emu(50000), RW - Emu(100000), Emu(270000),
       "設計上の役割", 9.5, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, RX + Emu(70000), CMT_Y + Emu(310000), RW - Emu(100000), Emu(270000),
       "SGバトル（継続判定）とは独立した「横道」として\n"
       "AT中に予期せず割り込む。驚きと興奮が打感を豊かにする。",
       8, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.42), Inches(9.6), Emu(580000), RGBColor(0x06, 0x10, 0x20))
    rect(s, Inches(0.2), Inches(4.42), Emu(55000), Emu(580000), C_SPEED)
    tb(s, Inches(0.45), Inches(4.47), Inches(9.1), Emu(250000),
       "ぶっちぎりバトルの価値：SGバトルとは別軸の自力感で、AT消化が「ただ待つ」にならない",
       9.5, bold=True, color=C_SPEED)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.1), Emu(250000),
       "裏連戦（紫）を引いた瞬間「今日は何かある」という直感が生まれる。シナリオ推測とは別軸の興奮。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 14種シナリオ × 示唆システム
# ══════════════════════════════════════════════════════════════
def s_scenario(prs):
    s = new_slide(prs)
    hdr(s, "シナリオ × 示唆システム  ──  読み解く楽しさが遊技の深みになる")

    # 左：ランプ示唆
    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.0), Inches(3.5), C_CARD, C_CYAN, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(2.75), Emu(320000),
       "フレームランプ示唆（6色）", 10, bold=True, color=C_CYAN, font=FONT_H)
    lamps = [
        (RGBColor(0xFF, 0xFF, 0xFF), "白",  "継続率：低"),
        (RGBColor(0x44, 0x88, 0xFF), "青",  "継続率：やや低"),
        (RGBColor(0xFF, 0xEE, 0x00), "黄",  "継続率：中"),
        (RGBColor(0x22, 0xCC, 0x44), "緑",  "継続率：やや高"),
        (RGBColor(0xCC, 0x22, 0x22), "赤",  "継続率：高"),
        (RGBColor(0xFF, 0x80, 0xFF), "虹",  "最高（艇王濃厚）"),
    ]
    ly = Inches(1.40)
    for lcol, lname, ldesc in lamps:
        rect(s, Inches(0.35), ly, Emu(340000), Emu(290000), lcol)
        tb(s, Inches(0.85), ly + Emu(30000), Inches(1.95), Emu(240000),
           f"{lname}  ──  {ldesc}", 8.5, color=C_CREAM)
        ly += Emu(320000)

    # 中央：セリフ示唆
    rect_b(s, Inches(3.4), Inches(0.85), Inches(3.1), Inches(3.5), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(3.52), Inches(0.92), Inches(2.85), Emu(320000),
       "サブ液晶・セリフ示唆", 10, bold=True, color=C_SPEED, font=FONT_H)
    clues = [
        ("登場キャラ",  "キャラによってモード・シナリオ段階を示唆"),
        ("セリフ内容",  "ポジティブ/ネガティブで継続率の高低を示唆"),
        ("ランプ×セリフ", "組み合わせで絞り込み精度が上がる"),
        ("ラウンド開始画面", "シナリオ段階と設定を直接示唆する要素"),
        ("第3停止タイミング", "第3停止を離す瞬間のカットインで確定示唆"),
    ]
    cy2 = Inches(1.40)
    for j, (k, v) in enumerate(clues):
        bg = RGBColor(0x10, 0x12, 0x06) if j % 2 == 0 else RGBColor(0x14, 0x16, 0x08)
        rect(s, Inches(3.48), cy2, Inches(2.94), Emu(320000), bg)
        tb(s, Inches(3.58), cy2 + Emu(25000), Inches(1.25), Emu(275000),
           k, 8, bold=True, color=C_SPEED, wrap=False)
        tb(s, Inches(4.88), cy2 + Emu(25000), Inches(1.45), Emu(275000),
           v, 7.5, color=C_CREAM)
        cy2 += Emu(332000)

    # 右：シナリオ推測の楽しさ
    rect_b(s, Inches(6.7), Inches(0.85), Inches(3.1), Inches(3.5), C_CARD, C_GOLD, 1.5)
    tb(s, Inches(6.82), Inches(0.92), Inches(2.85), Emu(320000),
       "シナリオ推測の楽しさ", 10, bold=True, color=C_GOLD, font=FONT_H)
    fun_points = [
        (C_GOLD,  "情報収集の喜び",
                  "ランプ・セリフ・画面を組み合わせて\n「このシナリオかも」と絞り込む楽しさ"),
        (C_SPEED, "期待感の変化",
                  "序盤に赤ランプが出た瞬間\n「今日は行ける」という確信が芽生える"),
        (C_CYAN,  "仲間との共有",
                  "「艇王引いたから8セット行くわ」という\nホールでの会話・コミュニティ形成"),
        (C_GREEN, "公平な抽選設計",
                  "ハマり度・設定に依存しないため\n「今回のAT」だけに集中できる"),
    ]
    fy = Inches(1.40)
    for col, title, desc in fun_points:
        rect(s, Inches(6.78), fy, Emu(50000), Emu(510000), col)
        tb(s, Inches(6.98), fy + Emu(25000), Inches(2.65), Emu(240000),
           title, 9, bold=True, color=col)
        tb(s, Inches(6.98), fy + Emu(265000), Inches(2.65), Emu(260000),
           desc, 8, color=C_CREAM)
        fy += Emu(555000)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(560000), RGBColor(0x06, 0x10, 0x20))
    rect(s, Inches(0.2), Inches(4.45), Emu(55000), Emu(560000), C_CYAN)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.1), Emu(250000),
       "示唆設計の本質：「完全に分からないが、ある程度読める」絶妙なバランス",
       9.5, bold=True, color=C_CYAN)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.1), Emu(250000),
       "全部わかると緊張感がなくなり、全く分からないと虚無感が残る。その中間に「楽しい」がある。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 上位AT「青島SG」
# ══════════════════════════════════════════════════════════════
def s_aoshima(prs):
    s = new_slide(prs)
    hdr(s, "上位AT「青島SG」  ──  純増4.0枚 × 継続率83% の頂点体験")

    # 左：スペックと構造
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.5), Inches(3.45), C_CARD, C_ELEC, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(320000),
       "青島SG スペック", 11, bold=True, color=C_ELEC, font=FONT_H)
    specs = [
        ("純増",     "約4.0枚/G（SG RUSHの1.6倍）",  C_GOLD2),
        ("継続率",   "約75%（Vストック込みで約83%）",  C_ELEC),
        ("突入経路", "グランドスラム後の青島VS波多野勝利", C_CREAM),
        ("構成",     "青島周回パート＋青島VS波多野の2部", C_CREAM),
        ("上乗せ",   "V揃い単：次セット継続濃厚\n  V揃いダブル：継続＋100〜300G上乗せ", C_SPEED),
    ]
    sy = Inches(1.38)
    for j, (k, v, col) in enumerate(specs):
        bg = RGBColor(0x06, 0x14, 0x28) if j % 2 == 0 else RGBColor(0x0A, 0x18, 0x2E)
        rect(s, Inches(0.28), sy, Inches(4.32), Emu(330000) if "\n" not in v else Emu(620000), bg)
        tb(s, Inches(0.38), sy + Emu(28000), Inches(1.3), Emu(260000),
           k, 8.5, bold=True, color=C_WATER, wrap=False)
        lines = v.count("\n") + 1
        tb(s, Inches(1.73), sy + Emu(28000), Inches(2.75), Emu(280000) * lines,
           v, 8.5, color=col)
        sy += Emu(340000) if "\n" not in v else Emu(640000)

    # 右：温泉モード＋突入の価値
    rect_b(s, Inches(4.9), Inches(0.85), Inches(4.9), Inches(3.45), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(5.02), Inches(0.92), Inches(4.6), Emu(320000),
       "特化ゾーン「温泉モード」＆突入の価値", 11, bold=True, color=C_SPEED, font=FONT_H)

    rect_b(s, Inches(5.08), Inches(1.38), Inches(4.64), Emu(840000),
           RGBColor(0x10, 0x18, 0x06), C_SPEED, 1.0)
    tb(s, Inches(5.2), Inches(1.44), Inches(4.3), Emu(280000),
       "温泉モード（30G）", 10, bold=True, color=C_SPEED)
    tb(s, Inches(5.2), Inches(1.78), Inches(4.3), Emu(400000),
       "青島VS波多野レース中のレア役成立で突入\n消化中はゲーム数上乗せ抽選が強化\n弱チェリー・強チャンス役で上乗せ確定",
       8.5, color=C_CREAM)

    rect_b(s, Inches(5.08), Inches(2.38), Inches(4.64), Emu(1700000),
           RGBColor(0x00, 0x0A, 0x22), C_ELEC, 1.0)
    tb(s, Inches(5.2), Inches(2.44), Inches(4.3), Emu(280000),
       "青島SGに突入する意味", 10, bold=True, color=C_ELEC)
    tb(s, Inches(5.2), Inches(2.78), Inches(4.3), Emu(1300000),
       "SG RUSH（純増2.5枚）との差は1.5枚/G\n100G消化で+150枚の差が生まれる\n\n"
       "継続率83%で長く続く×純増4.0枚\n= 1日の大台を作る唯一のルート\n\n"
       "「グランドスラム達成 → 青島SG」が\n来店時の最大目標として機能する",
       8.5, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.42), Inches(9.6), Emu(580000), RGBColor(0x00, 0x10, 0x24))
    rect(s, Inches(0.2), Inches(4.42), Emu(55000), Emu(580000), C_ELEC)
    tb(s, Inches(0.45), Inches(4.47), Inches(9.1), Emu(250000),
       "青島SGの設計意図：「到達した人だけが味わえる体験」として来店継続の最大動機を作る",
       9.5, bold=True, color=C_ELEC)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.1), Emu(250000),
       "SG RUSHだけでも楽しめるが、青島SGを知ると「もう一度グランドスラムを取りたい」が生まれる。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: グランドスラム × エキシビションレース
# ══════════════════════════════════════════════════════════════
def s_grandslam(prs):
    s = new_slide(prs)
    hdr(s, "グランドスラム × エキシビションレース  ──  8セット完走の先にあるもの")

    # 左：グランドスラム達成の仕組み
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.55), Inches(3.5), C_CARD, C_GOLD, 2.0)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(320000),
       "グランドスラム（8セット完走）", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.32), Inches(1.35), Inches(4.2), Emu(280000),
       "SGバトル（青島VS波多野）を8回連続で勝利した状態", 8.5, color=C_CREAM)

    gs_items = [
        (C_SPEED, "積み上げの快感",
                  "1セット目から8セット目まで\n「まだ続くかも」を8回繰り返す体験"),
        (C_GOLD,  "シナリオとの連動",
                  "艇王シナリオなら達成がほぼ確実\n低シナリオでの達成が「奇跡」体験"),
        (C_CYAN,  "達成後の特典",
                  "エキシビションレースへ自動移行\n青島SGへの入口が開く"),
    ]
    gi = Inches(1.78)
    for col, title, desc in gs_items:
        rect(s, Inches(0.28), gi, Emu(55000), Emu(560000), col)
        tb(s, Inches(0.52), gi + Emu(30000), Inches(4.0), Emu(260000),
           title, 9.5, bold=True, color=col)
        tb(s, Inches(0.52), gi + Emu(295000), Inches(4.0), Emu(280000),
           desc, 8.5, color=C_CREAM)
        gi += Emu(600000)

    # 右：エキシビションレース
    rect_b(s, Inches(5.0), Inches(0.85), Inches(4.75), Inches(3.5), C_CARD, C_SPEED, 1.5)
    tb(s, Inches(5.12), Inches(0.92), Inches(4.5), Emu(320000),
       "エキシビションレース", 11, bold=True, color=C_SPEED, font=FONT_H)
    tb(s, Inches(5.12), Inches(1.33), Inches(4.5), Emu(250000),
       "グランドスラム後の継続率固定AT区間", 8.5, color=C_CREAM)

    ex_rates = [
        ("A", "50.0%",  "低め",    C_LTGRY),
        ("B", "66.4%",  "中",      C_SPEED),
        ("C", "80.1%",  "高め",    C_CYAN),
        ("D", "90.2%",  "最高",    C_GOLD2),
    ]
    ey = Inches(1.68)
    for pat, rate, level, col in ex_rates:
        bg = RGBColor(0x0A, 0x14, 0x24) if ex_rates.index((pat, rate, level, col)) % 2 == 0 \
             else RGBColor(0x0E, 0x18, 0x2C)
        rect(s, Inches(5.1), ey, Inches(4.55), Emu(360000), bg)
        tb(s, Inches(5.2), ey + Emu(30000), Emu(380000), Emu(300000),
           f"型{pat}", 9, bold=True, color=col, wrap=False)
        tb(s, Inches(5.85), ey + Emu(30000), Emu(820000), Emu(300000),
           rate, 13, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, Inches(7.05), ey + Emu(30000), Inches(2.45), Emu(300000),
           f"継続率{level}", 8.5, color=C_CREAM)
        ey += Emu(375000)

    rect_b(s, Inches(5.1), Inches(3.42), Inches(4.55), Emu(820000),
           RGBColor(0x00, 0x0A, 0x22), C_ELEC, 1.5)
    tb(s, Inches(5.22), Inches(3.48), Inches(4.2), Emu(280000),
       "継続率に漏れた場合", 9.5, bold=True, color=C_ELEC, font=FONT_H)
    tb(s, Inches(5.22), Inches(3.82), Inches(4.2), Emu(420000),
       "青島VS波多野へ移行\n→ 勝利で上位AT「青島SG」突入\n（エキシビション失敗が青島SG入口になる）",
       8.5, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(560000), RGBColor(0x10, 0x10, 0x06))
    rect(s, Inches(0.2), Inches(4.45), Emu(55000), Emu(560000), C_GOLD)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.1), Emu(250000),
       "エキシビションの巧みな設計：「継続に漏れる = 青島SGへの入口」という逆転構造",
       9.5, bold=True, color=C_GOLD)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.1), Emu(250000),
       "エキシビションで多く続けば出玉増、漏れても青島SGチャンス。どちらに転んでも「良いことがある」設計。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: 設計哲学：なぜこの台は面白いのか
# ══════════════════════════════════════════════════════════════
def s_philosophy(prs):
    s = new_slide(prs)
    hdr(s, "設計哲学  ──  なぜモンキーターンVは面白いのか")

    pillars = [
        (Inches(0.2), C_WATER,
         "複層的な目標設計",
         "1G先：次の激走ポイント加算\n1AT先：次のセット継続\n1日先：グランドスラム達成\n将来：上位AT青島SG体験\n\n「今ここ」と「将来」を\n同時に意識させる設計",
         RGBColor(0x04, 0x14, 0x20)),
        (Inches(3.55), C_SPEED,
         "失敗を無駄にしない",
         "CZ失敗\n→ グランドスラムCH（救済）\nエキシビション漏れ\n→ 青島SG入口\n激走チャージ天井\n→ CZ確定\n\nどの失敗にも\n「次の期待」が残る",
         RGBColor(0x18, 0x10, 0x00)),
        (Inches(6.9), C_CYAN,
         "「読める」楽しさ",
         "シナリオはランプで示唆\n周期はモードで予測可能\nライバルモードは滞在確定\n\n「完全にわかる」ではなく\n「ある程度読める」バランス\n\n推測が当たる快感が\n遊技を豊かにする",
         RGBColor(0x00, 0x10, 0x20)),
    ]
    for x, col, title, desc, fill in pillars:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.45), fill, col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(340000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.40), Inches(2.8), Inches(2.6),
           desc, 9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(40000), C_WATER)
    rect(s, Inches(0.2), Inches(4.48), Inches(9.6), Emu(640000), RGBColor(0x06, 0x12, 0x20))
    tb(s, Inches(0.35), Inches(4.53), Inches(9.2), Emu(270000),
       "総評：「遊びやすさ × パンチ × 深み」が高水準で共存している数少ない台",
       9.5, bold=True, color=C_GOLD)
    tb(s, Inches(0.35), Inches(4.85), Inches(9.2), Emu(260000),
       "機械割設定6で114.9%・シナリオ14種・上位AT純増4.0枚。スペック・ゲーム性・演出のどこを切っても水準以上。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10: ベンチマーク × まとめ
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ  ──  設計から学べること")

    # 上段：他機種との比較
    rect(s, Inches(0.2), Inches(0.85), Inches(9.6), Emu(360000),
         RGBColor(0x08, 0x14, 0x28))
    bench_h = ["比較軸", "東京喰種", "番長4", "Re:ゼロ", "モンキーターンV"]
    bx = Inches(0.2)
    bcolw = [Emu(1006000), Emu(1924000), Emu(1924000), Emu(1924000), Emu(1924000)]
    bhx = []
    cx = bx + Emu(30000)
    for cw in bcolw:
        bhx.append(cx)
        cx += cw
    for h, hx_val, cw in zip(bench_h, bhx, bcolw):
        col = C_CYAN if h == "モンキーターンV" else C_GOLD
        tb(s, hx_val, Inches(0.91), cw - Emu(30000), Emu(285000),
           h, 8.5, bold=True, color=col, align=PP_ALIGN.CENTER)

    rows = [
        ("CV値",      "0.18", "0.68", "0.22", "0.23"),
        ("機械割(設6)","108%", "112%", "110%", "114.9%"),
        ("継続設計",  "シナリオ型", "設定依存", "有利区間型", "シナリオ14種"),
        ("来店継続",  "世界観",    "期待値",   "コンプガチャ型", "グランドスラム目標"),
        ("後半維持率","73.7%",  "29.6%",  "70%前後",  "68.2%"),
    ]
    rby = Inches(1.28)
    for j, row in enumerate(rows):
        rbg = RGBColor(0x0A, 0x14, 0x24) if j % 2 == 0 else RGBColor(0x0E, 0x18, 0x2C)
        rect(s, bx, rby, Inches(9.6), Emu(320000), rbg)
        for k, (val, hx_val, cw) in enumerate(zip(row, bhx, bcolw)):
            col = C_CYAN if k == 4 else (C_RED if val in ("0.68", "29.6%") else C_LTGRY)
            tb(s, hx_val, rby + Emu(25000), cw - Emu(30000), Emu(275000),
               val, 8, color=col, align=PP_ALIGN.CENTER)
        rby += Emu(330000)

    # 下段：学べること
    rect(s, Inches(0.2), Inches(3.95), Inches(9.6), Emu(760000),
         RGBColor(0x06, 0x12, 0x1E))
    rect(s, Inches(0.2), Inches(3.95), Emu(55000), Emu(760000), C_WATER)
    tb(s, Inches(0.45), Inches(4.00), Inches(9.1), Emu(260000),
       "モンキーターンVから学べる設計原則", 10, bold=True, color=C_WATER, font=FONT_H)

    learnings = [
        (C_SPEED, "複数の目標を同時に持たせる",  "短期（次セット）× 中期（グランドスラム）× 長期（青島SG）"),
        (C_CYAN,  "失敗に必ず「次の期待」を残す",  "どの失敗にも救済ルートがある = 諦めない設計"),
        (C_GOLD,  "「読める」示唆が遊技を豊かにする", "シナリオ推測という知的楽しさが長期稼働を支える"),
    ]
    ly = Inches(4.35)
    for col, title, desc in learnings:
        tb(s, Inches(0.52), ly, Emu(55000), Emu(270000), "│", 14, bold=True, color=col)
        tb(s, Inches(0.72), ly, Inches(3.5), Emu(270000), title, 9, bold=True, color=col, wrap=False)
        tb(s, Inches(4.3), ly, Inches(5.2), Emu(270000), desc, 8.5, color=C_CREAM, wrap=False)
        ly += Emu(285000)

    net_note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = [
        ("タイトル",                    s_title),
        ("ゲームフロー全体図",           s_flow),
        ("通常時設計",                   s_normal),
        ("CZ「超抜チャレンジ」",         s_cz),
        ("AT「SG RUSH」",               s_at),
        ("ぶっちぎりバトル × 連戦",      s_battle),
        ("シナリオ × 示唆システム",      s_scenario),
        ("上位AT「青島SG」",             s_aoshima),
        ("グランドスラム × エキシビション", s_grandslam),
        ("設計哲学",                     s_philosophy),
        ("まとめ",                       s_matome),
    ]

    print("=" * 55)
    print("  スマスロ モンキーターンV 機種分析資料ジェネレーター")
    print("=" * 55)
    print()
    for i, (name, func) in enumerate(slides, 1):
        print(f"  {i:2d}/{len(slides)} {name}")
        func(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
