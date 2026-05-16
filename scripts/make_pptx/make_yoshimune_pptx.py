"""
スマスロ真打吉宗 機種説明＋分析 統合版資料 v12  （大都技研・2026年4月6日導入）
出力: proposals/機種分析/吉宗/yoshimune_guide_v12.pptx
テーマ: 和黒 × 紫 × 橙 × 金（吉宗世界観）

Part A: 説明パート（プレイヤー視点）  スライド1〜6
Part B: 分析パート                    スライド7〜9
"""
import io
import os
import sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(ROOT_DIR,
           "proposals", "機種分析", "吉宗", "yoshimune_guide_v12.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（和黒×紫×橙×金）───────────────────────────
C_BG    = RGBColor(0x08, 0x04, 0x14)   # 和黒（深夜の空）
C_CARD  = RGBColor(0x10, 0x08, 0x20)   # カード背景
C_CARD2 = RGBColor(0x18, 0x10, 0x2C)   # カード背景2
C_ROW   = RGBColor(0x14, 0x0C, 0x24)   # テーブル奇数行
C_PUR   = RGBColor(0x88, 0x22, 0xCC)   # 紫（吉宗メインカラー）
C_PUR2  = RGBColor(0xAA, 0x55, 0xFF)   # 明るい紫
C_ORG   = RGBColor(0xFF, 0x77, 0x11)   # 橙（真BB色）
C_ORG2  = RGBColor(0xFF, 0xAA, 0x44)   # 明るい橙
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)   # 輝く金
C_RED   = RGBColor(0xCC, 0x22, 0x22)   # 赤（警告）
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)   # オフホワイト
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)   # クリーム
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)   # グレー
C_LTGRY = RGBColor(0x44, 0x44, 0x66)   # ライトグレー
C_GREEN = RGBColor(0x22, 0xCC, 0x66)   # 緑

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

TOTAL_SLIDES = 9


# ── 背景生成（和の深黒・墨色）─────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (8, 4, 20))
    draw = ImageDraw.Draw(img)
    # 斜めライン（格子模様・和紙風）
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(12, 8, 28), width=1)
    # 下部の紫グロー
    for y in range(h - 100, h):
        t = (y - (h - 100)) / 100
        draw.line([(0, y), (w, y)], fill=(int(30 * t), 0, int(45 * t)))
    # 上部薄暗化
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    # 右端のアクセントライン
    for x in range(w - 6, w):
        draw.line([(x, 0), (x, h)], fill=(0x60, 0x10, 0x99))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


# ── ヘルパー関数 ───────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    """塗りつぶし矩形（ボーダーなし）"""
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    """塗りつぶし矩形（ボーダーあり）"""
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    """テキストボックス"""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def hdr(slide, title_text, pg=""):
    """スライドヘッダー（紫ライン＋タイトル）"""
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_PUR)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_PUR)


def net_note(slide):
    """右下の※ネット解析情報より"""
    tb(slide, Inches(8.0), Inches(5.35), Inches(1.85), Emu(200000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, design_comment, sub_text=""):
    """フッター（設計コメント＋補足）"""
    fy = Inches(5.05)
    rect(slide, 0, fy, SLIDE_W, Inches(0.55), RGBColor(0x08, 0x04, 0x18))
    tb(slide, Inches(0.2), fy + Emu(30000), Inches(7.0), Emu(380000),
       "【設計】" + design_comment, 7.5, bold=True, color=C_GOLD)
    if sub_text:
        tb(slide, Inches(0.2), fy + Emu(230000), Inches(7.5), Emu(200000),
           sub_text, 6.5, color=C_GRAY)
    net_note(slide)


def arrow_r(slide, x, cy, col=None):
    """右向き矢印"""
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    """下向き矢印"""
    shp = slide.shapes.add_shape(17, cx - Emu(90000), y, Emu(180000), Emu(200000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    """Part A – スライド1: タイトル・スペック・3ポイント"""
    s = new_slide(prs)

    # 左パネル（タイトル領域）
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x06, 0x02, 0x10))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_PUR)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, RGBColor(0x60, 0x10, 0x99))

    # PartAバッジ
    rect(s, Inches(0.22), Inches(0.2), Inches(1.4), Emu(260000), C_PUR)
    tb(s, Inches(0.22), Inches(0.2), Inches(1.4), Emu(260000),
       "Part A 説明編", 7.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.22), Inches(0.6), Inches(5.0), Emu(300000),
       "機種説明＋分析 統合版資料 v12", 10, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(1.05), Inches(5.1), Emu(900000),
       "スマスロ\n真打吉宗", 30, bold=True, color=C_ORG, font=FONT_H)
    tb(s, Inches(0.22), Inches(3.05), Inches(5.0), Emu(300000),
       "── 1G連×純増9枚で4号機の魂が甦る", 10, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.6), Inches(4.9), Emu(230000),
       "メーカー: 大都技研　　導入: 2026年4月6日", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.9), Inches(4.9), Emu(230000),
       "設定: 1〜6段階　　真BB純増: 約9.0枚/G", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(4.2), Inches(4.9), Emu(230000),
       "AT: 勧善懲悪RUSH（純増約2.7枚/G・差枚数管理）", 8.5, color=C_GRAY)

    # 右：この台の3ポイント
    kws = [
        (C_PUR,  "勧善懲悪RUSH（AT）",  "周期CZからAT突入\n御白州チャンスで差枚数上乗せ\n真高確率で真BBを目指す"),
        (C_ORG,  "真BB（純増9.0枚）",   "月下ノ花道経由で突入\n1セット2000枚爆獲得\nBB中に1G連抽選が走る"),
        (C_GOLD, "究極鷹ブレイク",      "毎G1000枚ループ上乗せ\n終了後は月下ノ花道で真BB確定\n実質5000枚以上確定の夢体験"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.3 + i * 1.65)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.4), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.4), ac)
        tb(s, Inches(5.85), y0 + Emu(60000), Inches(3.8), Emu(300000),
           kw, 12, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(360000), Inches(3.8), Emu(500000),
           desc, 8, color=C_WHITE)

    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    """Part A – スライド2: ゲームフロー全体図"""
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常→CZ→AT→真BB→1G連→究極鷹ブレイク", f"2/{TOTAL_SLIDES}")

    # ── 上段（左→右）: 通常時 → CZ → AT → 真高確率
    row1_y = Inches(0.78)
    row1_h = Emu(1500000)
    boxes1 = [
        (C_CARD2,                    C_PUR,   "通常時",
         "周期ポイント蓄積\n(6周期制)\nレア役で加算"),
        (C_CARD2,                    C_PUR,   "CZ\n悪人成敗チャンス",
         "鷹CZが最高格\nAT当選を目指す"),
        (C_CARD2,                    C_PUR2,  "AT\n勧善懲悪RUSH",
         "純増2.7枚/G\n差枚数管理型\n40G周期で内部CZ"),
        (RGBColor(0x14, 0x06, 0x20), C_PUR2,  "真高確率\n（ジャッジ）",
         "AT中 勧善懲悪\nチャンス勝利で突入\n真BB当選を目指す"),
    ]
    bw1 = Inches(2.0)
    gap1 = Inches(0.28)
    sx1 = Inches(0.3)
    for i, (fill, bc, lbl, sub) in enumerate(boxes1):
        bx0 = sx1 + i * (bw1 + gap1)
        rect_b(s, bx0, row1_y, bw1, row1_h, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row1_y + Emu(80000),
           bw1 - Emu(80000), Emu(420000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row1_y + Emu(510000),
           bw1 - Emu(60000), Emu(850000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw1 + Emu(10000), row1_y + row1_h // 2)

    # 右端ラベル「真BB当選！」
    tb(s, Inches(9.0), row1_y + Emu(600000), Inches(0.9), Emu(600000),
       "真BB\n当選！", 9, bold=True, color=C_ORG, align=PP_ALIGN.CENTER, font=FONT_H)

    # ↓矢印（右端から下段右端へ）
    arrow_d(s, Inches(9.45), row1_y + row1_h + Emu(20000), col=C_ORG)

    # ── 下段（右→左）: 月下ノ花道 → 真BB → 1G連 → 究極鷹ブレイク
    row2_y = row1_y + row1_h + Emu(350000)
    row2_h = Emu(1550000)
    boxes2 = [
        (RGBColor(0x20, 0x10, 0x00), C_GOLD, "究極鷹ブレイク",
         "毎G1000枚ループ上乗せ\n終了後→月下ノ花道へ\n5000枚以上確定！"),
        (RGBColor(0x18, 0x06, 0x00), C_GOLD2, "1G連！",
         "真BB終了の1G後\n即次のBBが始まる\n連続するほど興奮↑"),
        (RGBColor(0x22, 0x08, 0x04), C_ORG,  "真BB\n(純増9.0枚)",
         "2000枚爆獲得\nBB中に1G連抽選\n成立役で当選率変化"),
        (RGBColor(0x14, 0x06, 0x20), C_PUR,  "月下ノ花道",
         "真BB突入前の\n演出ゾーン\n真BB確定状態"),
    ]
    # 右→左の順で配置
    bw2 = Inches(2.0)
    gap2 = Inches(0.28)
    sx2 = Inches(0.3)
    n2 = len(boxes2)
    for i, (fill, bc, lbl, sub) in enumerate(reversed(boxes2)):
        idx = n2 - 1 - i
        bx0 = sx2 + idx * (bw2 + gap2)
        rect_b(s, bx0, row2_y, bw2, row2_h, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row2_y + Emu(80000),
           bw2 - Emu(80000), Emu(420000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row2_y + Emu(510000),
           bw2 - Emu(60000), Emu(900000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if idx > 0:
            arrow_r(s, bx0 - gap2 - Emu(10000), row2_y + row2_h // 2, col=C_ORG)

    # 左端ラベル（1G連ループ矢印）
    rect_b(s, Inches(0.05), row2_y + Emu(400000), Emu(190000), Emu(750000),
           RGBColor(0x20, 0x10, 0x00), C_GOLD, 1.5)
    tb(s, Inches(0.05), row2_y + Emu(440000), Emu(190000), Emu(680000),
       "↑\nLoop", 7, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    footer(s,
           "「通常→CZ→AT→真BB→1G連→究極鷹ブレイク」の全ルートを可視化。蛇行2段で流れを追える。",
           "補足: 究極鷹ブレイク終了後は月下ノ花道へ移行（真BB確定）。1G連ループで左端に戻るイメージ。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    """Part A – スライド3: 通常時の遊び方"""
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 周期CZとレア役の2ルートでATを目指す", f"3/{TOTAL_SLIDES}")

    # 左カラム: 周期システム
    lx, ly = Inches(0.3), Inches(0.72)
    lw = Inches(4.5)
    card_h = Emu(3900000)
    rect_b(s, lx, ly, lw, card_h, C_CARD, C_PUR, 1.5)
    rect(s, lx, ly, Emu(45000), card_h, C_PUR)
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(260000),
       "ルート① 周期CZシステム", 11, bold=True, color=C_PUR, font=FONT_H)

    items_l = [
        ("6周期制",           "通常時はゲーム数消化でポイントを蓄積。\n規定ポイント到達で「CZ（悪人成敗チャンス）」が発生。"),
        ("4種のCZモード",     "滞在モードによって天井周期が異なる。\n「天国モード相当」では1周期目でCZ確定。"),
        ("鷹CZ（最高格）",    "最も期待度が高いCZ。\n突入したら大チャンス。"),
        ("CZ天井",            "最大6周期（約1000G）でCZ確定。\nCZ後・AT後にモードが再抽選される。"),
    ]
    for i, (title, body) in enumerate(items_l):
        iy = ly + Emu(320000) + i * Emu(860000)
        rect_b(s, lx + Emu(60000), iy, lw - Emu(75000), Emu(810000), C_CARD2, C_PUR, 0.6)
        tb(s, lx + Emu(90000), iy + Emu(50000), lw - Emu(130000), Emu(250000),
           title, 9, bold=True, color=C_PUR2)
        tb(s, lx + Emu(90000), iy + Emu(300000), lw - Emu(130000), Emu(440000),
           body, 7.5, color=C_WHITE)

    # 右カラム: レア役ルート
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)
    rect_b(s, rx, ry, rw, card_h, C_CARD, C_ORG, 1.5)
    rect(s, rx, ry, Emu(45000), card_h, C_ORG)
    tb(s, rx + Emu(75000), ry + Emu(50000), rw - Emu(100000), Emu(260000),
       "ルート② レア役からの自力当選", 11, bold=True, color=C_ORG, font=FONT_H)

    items_r = [
        ("チェリー・スイカ",  "レア小役成立でポイントを大量加算。\n周期到達を早める効果がある。"),
        ("チャンス目",        "特殊な出目が絡む役。\n直接CZ当選や高確率状態への抽選あり。"),
        ("AT天井",           "最大約1500GでAT確定。\n周期天井と合わせて管理するとよい。"),
        ("打ち方のポイント",  "全役対応の左リール枠内にBAR付きチェリーを狙う。\nスイカはV字押し（右・左・中）で取得。"),
    ]
    for i, (title, body) in enumerate(items_r):
        iy = ry + Emu(320000) + i * Emu(860000)
        rect_b(s, rx + Emu(60000), iy, rw - Emu(75000), Emu(810000), C_CARD2, C_ORG, 0.6)
        tb(s, rx + Emu(90000), iy + Emu(50000), rw - Emu(130000), Emu(250000),
           title, 9, bold=True, color=C_ORG2)
        tb(s, rx + Emu(90000), iy + Emu(300000), rw - Emu(130000), Emu(440000),
           body, 7.5, color=C_WHITE)

    footer(s,
           "通常時は「待つ設計」ではなく「レア役で能動的に縮められる」二層構造。周期ゲーとレア役ゲーを両立。",
           "補足: 周期モードはCZ後・AT後に抽選し直されるため、連続CZも起こりうる。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: AT/ボーナスの遊び方
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    """Part A – スライド4: AT/ボーナスの遊び方"""
    s = new_slide(prs)
    hdr(s, "AT/ボーナスの遊び方 ── 勧善懲悪RUSH → 真高確率 → 真BB", f"4/{TOTAL_SLIDES}")

    # AT説明（上段2カラム）
    lx, ly = Inches(0.3), Inches(0.72)
    lw, lh = Inches(4.5), Emu(1900000)
    rx, rw = Inches(5.0), Inches(4.7)

    # 左: AT概要
    rect_b(s, lx, ly, lw, lh, C_CARD, C_PUR, 1.5)
    rect(s, lx, ly, Emu(45000), lh, C_PUR)
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(260000),
       "AT「勧善懲悪RUSH」概要", 10.5, bold=True, color=C_PUR, font=FONT_H)
    tb(s, lx + Emu(75000), ly + Emu(320000), lw - Emu(100000), lh - Emu(380000),
       "■ 純増: 約2.7枚/G（差枚数管理型）\n"
       "■ 初期枚数: 150枚+α\n"
       "■ 40Gごとに内部周期CZ「勧善懲悪チャンス」抽選\n"
       "■ レア役 → 直乗せ（20〜500枚）or\n"
       "   御白州チャンス高確 or 御白州チャンス突入",
       8, color=C_WHITE)

    # 右: 御白州チャンス
    rect_b(s, rx, ly, rw, lh, C_CARD, C_GOLD, 1.5)
    rect(s, rx, ly, Emu(45000), lh, C_GOLD)
    tb(s, rx + Emu(75000), ly + Emu(50000), rw - Emu(100000), Emu(260000),
       "御白州チャンス（上乗せゾーン）", 10.5, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, rx + Emu(75000), ly + Emu(320000), rw - Emu(100000), lh - Emu(380000),
       "■ 全4種類の報酬から抽選で決定\n"
       "■ 最低50枚以上の差枚数上乗せ確定\n"
       "■ 消化中レア役成立で特殊上乗せ以上確定\n"
       "■ 最高報酬は「真BB確定」級の大当たり",
       8, color=C_WHITE)

    # 真高確率の流れ（下段）
    flow_y = ly + lh + Emu(120000)
    flow_h = Emu(1050000)
    flow_boxes = [
        (C_CARD2, C_PUR,   "勧善懲悪チャンス\n（3G特化ゾーン）",
         "AT中40G周期で発生\n勝利で真高確率へ"),
        (RGBColor(0x14, 0x06, 0x20), C_PUR2, "真高確率\nジャッジ（1G）",
         "真BB当選を判定\n当選で月下ノ花道"),
        (RGBColor(0x18, 0x06, 0x00), C_ORG,  "月下ノ花道",
         "真BB確定の演出ゾーン\n真BBへ突入！"),
        (RGBColor(0x22, 0x08, 0x04), C_ORG,  "真BB\n（純増9.0枚）",
         "2000枚爆獲得\n1G連抽選スタート！"),
    ]
    bfw = Inches(2.15)
    fgap = Inches(0.25)
    fsx = (SLIDE_W - 4 * bfw - 3 * fgap) / 2
    for i, (fill, bc, lbl, sub) in enumerate(flow_boxes):
        bx0 = fsx + i * (bfw + fgap)
        rect_b(s, bx0, flow_y, bfw, flow_h, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), flow_y + Emu(70000),
           bfw - Emu(60000), Emu(380000), lbl, 9, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), flow_y + Emu(460000),
           bfw - Emu(50000), Emu(490000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bfw + Emu(10000), flow_y + flow_h // 2, col=bc)

    footer(s,
           "ATは「差枚数管理＋御白州チャンス」で枚数を伸ばし、勧善懲悪チャンスで真BBを狙う二段構え。",
           "補足: AT純増2.7枚/Gはあくまで通常消化時。真BB突入後は9.0枚/Gに跳ね上がる。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 1G連の仕組み
# ══════════════════════════════════════════════════════════════
def s_1g(prs):
    """Part A – スライド5: 1G連の仕組みと体験"""
    s = new_slide(prs)
    hdr(s, "1G連の仕組み ── 「終わり→始まり」を1Gで完結させる感情設計", f"5/{TOTAL_SLIDES}")

    # 上段: 仕組み説明（3カード）
    card_y = Inches(0.75)
    card_h = Emu(1800000)
    cw = Inches(2.9)
    cgap = Inches(0.28)
    csx = Inches(0.3)
    cards = [
        (C_PUR,  "① BB中の1G連抽選",
         "真BB消化中に成立役の種類に応じて\n1G連獲得の抽選が走る。\n\n高設定ほど1G連当選率が優遇。\nレア役成立が当選チャンス。"),
        (C_ORG,  "② 1G後に次のBBが始まる",
         "真BB終了の直後（1G）に\n再び月下ノ花道が始まり\n次の真BBに突入する。\n\n「終わった！」が「え、また!?」に変わる。"),
        (C_GOLD, "③ 連続するほど加速する",
         "1G連が連続するたびに\n獲得枚数は2000枚ずつ積み重なる。\n\n興奮も線形ではなく指数的に上昇。\n究極鷹ブレイクへの布石になることも。"),
    ]
    for i, (ac, title, body) in enumerate(cards):
        bx0 = csx + i * (cw + cgap)
        rect_b(s, bx0, card_y, cw, card_h, C_CARD, ac, 1.8)
        rect(s, bx0, card_y, Emu(45000), card_h, ac)
        tb(s, bx0 + Emu(75000), card_y + Emu(55000), cw - Emu(95000), Emu(270000),
           title, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, bx0 + Emu(75000), card_y + Emu(330000), cw - Emu(95000), card_h - Emu(390000),
           body, 8, color=C_WHITE)
        if i < 2:
            arrow_r(s, bx0 + cw + Emu(15000), card_y + card_h // 2, col=ac)

    # 下段: 感情の波グラフ（テキストで可視化）
    gx, gy = Inches(0.3), Inches(0.75) + card_h + Emu(120000)
    gw, gh = Inches(9.4), Emu(1500000)
    rect_b(s, gx, gy, gw, gh, RGBColor(0x0C, 0x06, 0x20), C_PUR2, 1.2)
    tb(s, gx + Emu(60000), gy + Emu(50000), gw - Emu(80000), Emu(260000),
       "1G連の感情曲線 ── 興奮は「谷と山」の繰り返しで最大化される", 9, bold=True, color=C_PUR2, font=FONT_H)

    phases = [
        ("真BB消化\n(興奮MAX)", C_ORG),
        ("BB終了\n(落胆)", C_GRAY),
        ("1G連！\n(衝撃MAX)", C_GOLD2),
        ("真BB消化\n(再加速)", C_ORG),
        ("究極鷹ブレイク\n(現実が変わる)", C_GOLD),
    ]
    pw = (gw - Emu(120000)) // len(phases)
    for i, (label, ac) in enumerate(phases):
        px = gx + Emu(60000) + i * pw
        ph = Emu(850000) if i % 2 == 0 else Emu(400000)
        py = gy + gh - ph - Emu(50000)
        rect(s, px + Emu(30000), py, pw - Emu(60000), ph, ac)
        tb(s, px, gy + Emu(320000), pw, Emu(450000),
           label, 7.5, bold=(i % 2 == 0), color=ac, align=PP_ALIGN.CENTER)

    footer(s,
           "1G連は「終わり→始まり」を1Gで完結させる感情設計。この「谷と山」が興奮を最大化する。",
           "補足: 4号機吉宗の代名詞だった1G連を現行スマスロスペックで完全再現。30〜40代の記憶に直撃する。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 究極鷹ブレイク
# ══════════════════════════════════════════════════════════════
def s_hawk(prs):
    """Part A – スライド6: 究極鷹ブレイクの到達ルートと遊び方"""
    s = new_slide(prs)
    hdr(s, "究極鷹ブレイク ── 毎G1000枚ループ・5000枚以上確定の最高体験", f"6/{TOTAL_SLIDES}")

    # 左: 到達ルート
    lx, ly = Inches(0.3), Inches(0.72)
    lw, lh = Inches(4.4), Emu(3900000)
    rect_b(s, lx, ly, lw, lh, C_CARD, C_GOLD, 1.8)
    rect(s, lx, ly, Emu(45000), lh, C_GOLD)
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(260000),
       "到達ルート", 11, bold=True, color=C_GOLD, font=FONT_H)

    routes = [
        ("真BB中に1G連を複数獲得",
         "真BB消化中に1G連抽選を突破。\n1G連が連続するほど究極鷹ブレイクの\n"
         "抽選チャンスが増える。"),
        ("BB中の特殊演出を経て発動",
         "真BBの高確率ゾーン中などに\n特定演出（鷹の乱舞等）が発生すると\n究極鷹ブレイク確定の報知あり。"),
        ("究極鷹ブレイク終了後も継続",
         "1000枚ループ失敗で終了だが\n終了後は月下ノ花道へ移行（真BB確定）。\n1G連状態で再スタートする構造。"),
    ]
    for i, (title, body) in enumerate(routes):
        iy = ly + Emu(320000) + i * Emu(1160000)
        rect_b(s, lx + Emu(60000), iy, lw - Emu(75000), Emu(1100000),
               RGBColor(0x16, 0x10, 0x04), C_GOLD, 0.8)
        rect(s, lx + Emu(60000), iy, Emu(30000), Emu(1100000), C_GOLD)
        tb(s, lx + Emu(110000), iy + Emu(50000), lw - Emu(140000), Emu(260000),
           f"STEP {i+1}  {title}", 8.5, bold=True, color=C_GOLD2)
        tb(s, lx + Emu(110000), iy + Emu(310000), lw - Emu(140000), Emu(720000),
           body, 7.5, color=C_WHITE)

    # 右: 遊び方・インパクト
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    # 上カード: 仕組み
    top_h = Emu(1900000)
    rect_b(s, rx, ry, rw, top_h, RGBColor(0x1A, 0x0C, 0x02), C_ORG, 1.5)
    rect(s, rx, ry, Emu(45000), top_h, C_ORG)
    tb(s, rx + Emu(75000), ry + Emu(50000), rw - Emu(100000), Emu(260000),
       "仕組みと遊び方", 10.5, bold=True, color=C_ORG, font=FONT_H)
    tb(s, rx + Emu(75000), ry + Emu(320000), rw - Emu(100000), top_h - Emu(380000),
       "■ 毎ゲーム1000枚のループ上乗せが発生\n"
       "■ ループ失敗（低確率）で終了\n"
       "■ 上乗せ成功中は演出を楽しみながら消化\n"
       "■ 終了後は月下ノ花道（真BB確定）に移行\n"
       "■ つまり究極鷹ブレイク後も真BBで1G連継続！",
       8, color=C_WHITE)

    # 下カード: 出玉インパクト
    bot_y = ry + top_h + Emu(100000)
    bot_h = lh - top_h - Emu(100000)
    rect_b(s, rx, bot_y, rw, bot_h, RGBColor(0x16, 0x10, 0x00), C_GOLD, 2.0)
    rect(s, rx, bot_y, Emu(45000), bot_h, C_GOLD)
    tb(s, rx + Emu(75000), bot_y + Emu(50000), rw - Emu(100000), Emu(260000),
       "出玉インパクト", 10.5, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, rx + Emu(75000), bot_y + Emu(320000), rw - Emu(100000), bot_h - Emu(380000),
       "究極鷹ブレイク中: 1000枚 × ループ数\n"
       "終了後の真BB:     +2000枚\n"
       "1G連複数:         +2000枚 × 連数\n"
       "─────────────────────\n"
       "合計 5000枚以上確定！\n"
       "（ループが続けば10000枚超も現実的）\n\n"
       "「現実が変わる体験」を提供する本機の頂点",
       8.5, color=C_GOLD2, font=FONT_H)

    footer(s,
           "究極鷹ブレイクは「毎G1000枚ループ＋終了後も真BB確定」の二重安全網。終わらない夢体験を設計している。",
           "補足: 発生確認時点で最低でも5000枚以上が期待できる。スマスロ最高峰の出玉インパクトの一つ。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    """Part B – スライド7: なぜこの台は面白いのか"""
    s = new_slide(prs)
    hdr(s, "【分析】面白さの設計 ── なぜ真打吉宗は刺さるのか", f"7/{TOTAL_SLIDES}")

    # PartBバッジ
    rect(s, Inches(9.2), Inches(0.02), Inches(0.75), Emu(240000), C_GOLD)
    tb(s, Inches(9.2), Inches(0.02), Inches(0.75), Emu(240000),
       "Part B", 7, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # 上段3カード
    card_y = Inches(0.72)
    card_h = Emu(1600000)
    cw = Inches(2.9)
    cgap = Inches(0.28)
    csx = Inches(0.3)
    top3 = [
        (C_PUR,  "IP記憶との接続",
         "4号機吉宗の1G連は「コイン100枚が一瞬で消えた頃の記憶」。"
         "\n30〜40代の休眠層が「あの台が帰ってきた」と感じる\n"
         "強力な引力を持つ。ノスタルジアは最強のマーケティング。"),
        (C_ORG,  "感情の波を設計する",
         "通常の継続型ATは「緩やかな興奮」しか生まない。\n"
         "1G連は「落胆→衝撃」を1Gで繰り返す設計。\n"
         "この感情の谷と山が記憶に残る体験を作る。"),
        (C_GOLD, "純増9枚という体験速度",
         "2.7枚/Gと9.0枚/Gでは「感じるスピード」が3倍以上違う。\n"
         "真BBの速さは「自分が主役になった感」を生む。\n"
         "スピードそのものが演出になっている。"),
    ]
    for i, (ac, title, body) in enumerate(top3):
        bx0 = csx + i * (cw + cgap)
        rect_b(s, bx0, card_y, cw, card_h, C_CARD, ac, 1.8)
        rect(s, bx0, card_y, Emu(45000), card_h, ac)
        tb(s, bx0 + Emu(75000), card_y + Emu(55000), cw - Emu(95000), Emu(270000),
           title, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, bx0 + Emu(75000), card_y + Emu(330000), cw - Emu(95000), card_h - Emu(390000),
           body, 8, color=C_WHITE)

    # 下段: 設計原則表
    tx, ty = Inches(0.3), Inches(0.72) + card_h + Emu(120000)
    tw, th = Inches(9.4), Emu(1700000)
    rect_b(s, tx, ty, tw, th, C_CARD, C_LTGRY, 1.0)
    tb(s, tx + Emu(60000), ty + Emu(50000), tw - Emu(80000), Emu(260000),
       "1G連システムの設計原則（他機種との比較）", 9, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        ("通常AT継続",    "連続する興奮",   "単調になりやすい・慣れで麻痺",     C_GRAY),
        ("ゲーム数上乗せ", "数字が増える喜び", "抽象的・実感が遅い",              C_GRAY),
        ("1G連（吉宗）",  "落胆→衝撃の波",  "感情的記憶に残る・話題性が高い",   C_GOLD2),
    ]
    row_h_p = (th - Emu(360000)) // len(principles)
    cols = [Inches(2.5), Inches(2.5), Inches(3.0), Inches(1.4)]
    col_labels = ["方式", "体験タイプ", "プレイヤーへの効果", "評価"]
    rect(s, tx, ty + Emu(310000), tw, Emu(280000), RGBColor(0x44, 0x10, 0x70))
    cx = tx + Emu(60000)
    for j, (label, cw_) in enumerate(zip(col_labels, cols)):
        tb(s, cx, ty + Emu(330000), cw_, Emu(250000),
           label, 8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        cx += cw_

    for i, (p1, p2, p3, ac) in enumerate(principles):
        row_y = ty + Emu(590000) + i * row_h_p
        bg = C_CARD if i % 2 == 0 else C_CARD2
        rect(s, tx, row_y, tw, row_h_p - Emu(20000), bg)
        cx = tx + Emu(60000)
        for j, (val, cw_) in enumerate(zip([p1, p2, p3, ("◎" if ac == C_GOLD2 else "△")], cols)):
            col = ac if j == 0 and ac == C_GOLD2 else (C_GOLD if j == 3 and ac == C_GOLD2 else C_WHITE)
            tb(s, cx, row_y + Emu(80000), cw_, row_h_p - Emu(100000),
               val, 8, bold=(j == 0), color=col, align=PP_ALIGN.CENTER)
            cx += cw_

    footer(s,
           "1G連は「感情の波」を設計した革新。通常のAT継続では生めない「衝撃」を体験させる。",
           "補足: IP記憶×感情波×速度感の三位一体が「刺さる台」の正体。それぞれ単独では到達できない高みを生む。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    """Part B – スライド8: 良い点と課題"""
    s = new_slide(prs)
    hdr(s, "【分析】良い点と課題 ── 爆発力と荒波の二面性", f"8/{TOTAL_SLIDES}")

    # 良い点（左）
    lx, ly = Inches(0.3), Inches(0.72)
    lw = Inches(4.5)
    rect(s, lx, ly, lw, Emu(330000), RGBColor(0x10, 0x44, 0x20))
    tb(s, lx + Emu(60000), ly + Emu(50000), lw - Emu(80000), Emu(260000),
       "良い点（PROS）", 11, bold=True, color=C_GREEN, font=FONT_H)

    pros = [
        ("1G連体験の再現性",
         "4号機吉宗のコア体験を完全再現。\n"
         "「また!?」という衝撃は現代スペックで何倍も強烈。"),
        ("純増9.0枚の速度感",
         "現行最高クラスの純増速度。\n"
         "真BBが始まると体感は「秒で終わる」感覚になる。"),
        ("究極鷹ブレイクの夢",
         "5000枚以上確定という数字が「夢」として機能。\n"
         "台の前に座る理由を提供する最大の動機。"),
        ("設計のシンプルさ",
         "「1G連とればいい」「鷹ブレイクに行けばいい」\n"
         "という分かりやすいゴールがある。"),
    ]
    for i, (title, body) in enumerate(pros):
        iy = ly + Emu(330000) + i * Emu(950000)
        rect_b(s, lx, iy, lw, Emu(910000), C_CARD, C_GREEN, 0.8)
        rect(s, lx, iy, Emu(40000), Emu(910000), RGBColor(0x10, 0x88, 0x44))
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           title, 9, bold=True, color=C_GREEN)
        tb(s, lx + Emu(70000), iy + Emu(310000), lw - Emu(90000), Emu(540000),
           body, 8, color=C_WHITE)

    # 課題（右）
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)
    rect(s, rx, ry, rw, Emu(330000), RGBColor(0x44, 0x10, 0x10))
    tb(s, rx + Emu(60000), ry + Emu(50000), rw - Emu(80000), Emu(260000),
       "課題（CONS）", 11, bold=True, color=C_RED, font=FONT_H)

    cons = [
        ("荒波による消化不良",
         "真BBを引けないと枚数が伸びない設計。\n"
         "「当たらないAT」を長時間消化するストレス。"),
        ("設定1の重さ",
         "設定差が大きく、設定1は初当たりが遠い。\n"
         "低設定での長時間遊戯はコスパが厳しい。"),
        ("ライト層への敷居",
         "「1G連」の前提知識がない世代には\n"
         "驚きが半減する。IP記憶は刺さる人限定の武器。"),
        ("究極鷹ブレイクの偶然性",
         "5000枚以上という数字は夢だが\n"
         "実際の到達率は公開されておらず不透明感がある。"),
    ]
    for i, (title, body) in enumerate(cons):
        iy = ry + Emu(330000) + i * Emu(950000)
        rect_b(s, rx, iy, rw, Emu(910000), C_CARD, C_RED, 0.8)
        rect(s, rx, iy, Emu(40000), Emu(910000), RGBColor(0x88, 0x18, 0x18))
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           title, 9, bold=True, color=C_RED)
        tb(s, rx + Emu(70000), iy + Emu(310000), rw - Emu(90000), Emu(540000),
           body, 8, color=C_WHITE)

    footer(s,
           "爆発力はトップクラスだが荒波も激しい。コアファン向けに徹した潔い設計判断が吉と出るか凶と出るか。",
           "補足: 良い点・課題いずれも「1G連という設計の二面性」から来ている。設計の核心が強みでも弱みでもある。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    """Part B – スライド9: まとめ・設計から学べること"""
    s = new_slide(prs)
    hdr(s, "【分析】まとめ ── 設計から学べること", f"9/{TOTAL_SLIDES}")

    # 左: 設計から学べること
    lx, by = Inches(0.3), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, by, lw, Emu(300000), RGBColor(0x55, 0x10, 0x88))
    tb(s, lx + Emu(60000), by + Emu(50000), lw - Emu(80000), Emu(230000),
       "設計から学べること", 11, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_PUR,  "感情の「谷」が興奮を作る",
         "1G連の「終わった→また始まった」は\n"
         "一旦感情を下げることで次の山を高くする設計。\n"
         "継続型ATでは生めない感情の振れ幅がある。"),
        (C_ORG,  "三位一体の爆発設計",
         "純増9枚（速度）×1G連（継続）×究極鷹ブレイク（頂点）\n"
         "この三つが組み合わさることで\n"
         "どの一つも単独では到達できない頂点体験を生む。"),
        (C_GOLD, "IP記憶は最強の集客資産",
         "4号機吉宗の記憶がある層にとって\n"
         "「1G連が帰ってきた」は広告費ゼロの集客力。\n"
         "ノスタルジアを機能させる設計は他機種の手本。"),
    ]
    for i, (ac, t, b) in enumerate(elems):
        ey = by + Emu(300000) + i * Emu(1230000)
        rect_b(s, lx, ey, lw, Emu(1185000), C_CARD, ac, 1.5)
        rect(s, lx, ey, Emu(45000), Emu(1185000), ac)
        tb(s, lx + Emu(75000), ey + Emu(50000), lw - Emu(95000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), ey + Emu(305000), lw - Emu(95000), Emu(800000),
           b, 8, color=C_WHITE)

    # 右: 設計原則＋総括
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(210000),
       "設計原則 チェックリスト", 10, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_PUR,  "感情の「谷」を意図的に作ると「山」が高くなる"),
        (C_ORG,  "速度感（純増枚数）は演出の代わりになりうる"),
        (C_GOLD, "IP記憶への接続は広告を超えた集客力を持つ"),
        (C_PUR2, "複数要素の組み合わせが到達不可能な頂点を生む"),
        (C_GRAY, "爆発設計は必ず荒波を伴う ── 想定ユーザーを絞れ"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(280000) + i * Emu(490000)
        rect(s, rx, py0, Emu(20000), Emu(450000), ac)
        tb(s, rx + Emu(50000), py0 + Emu(60000), rw - Emu(60000), Emu(350000),
           p, 8, color=C_WHITE)

    # 総括ボックス
    sum_y = ry + Emu(280000) + len(principles) * Emu(490000) + Emu(60000)
    sum_h = SLIDE_H - Emu(320000) - sum_y
    rect_b(s, rx, sum_y, rw, sum_h,
           RGBColor(0x14, 0x06, 0x20), C_PUR, 2.0)
    tb(s, rx + Emu(55000), sum_y + Emu(45000), rw - Emu(75000), Emu(260000),
       "総括", 10, bold=True, color=C_PUR, font=FONT_H)
    tb(s, rx + Emu(55000), sum_y + Emu(305000), rw - Emu(75000), sum_h - Emu(360000),
       "4号機吉宗の遺産を現行最高スペックで昇華した稀有な一台。\n"
       "1G連という普遍的な興奮体験は時代を超えて機能し続ける。\n"
       "「感情設計」の教科書として他の機種設計にも応用可能。",
       8, color=C_WHITE)

    footer(s,
           "「感情の谷と山」「三位一体の爆発」「IP記憶接続」── この3原則が真打吉宗の本質。",
           "補足: 荒波という諸刃の剣を承知でコアファン向けに徹したことが、この台の設計的な誠実さでもある。")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Part A: 説明パート
    s_title(prs)    # 1. タイトル・スペック・3ポイント
    s_flow(prs)     # 2. ゲームフロー全体図（蛇行2段）
    s_normal(prs)   # 3. 通常時の遊び方
    s_at(prs)       # 4. AT/ボーナスの遊び方
    s_1g(prs)       # 5. 1G連の仕組み
    s_hawk(prs)     # 6. 究極鷹ブレイク

    # Part B: 分析パート
    s_design(prs)       # 7. 面白さの設計
    s_pros_cons(prs)    # 8. 良い点と課題
    s_matome(prs)       # 9. まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
