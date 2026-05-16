"""
AIゲーム性提案シリーズ 全体まとめ PowerPoint ジェネレーター
出力: proposals/まとめ/series_overview_v1.pptx
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(ROOT_DIR, "proposals", "まとめ", "series_overview_v1.pptx")

# ── カラーパレット（統合まとめ：ネイビー×シルバー×金）──────────────
C_BG     = RGBColor(0x06, 0x08, 0x14)
C_CARD   = RGBColor(0x10, 0x14, 0x28)
C_NAVY   = RGBColor(0x0C, 0x18, 0x38)
C_TEAL   = RGBColor(0x18, 0x90, 0x80)   # 機種分析アクセント
C_GOLD   = RGBColor(0xC8, 0xA8, 0x40)   # 新規提案アクセント
C_GOLD2  = RGBColor(0xFF, 0xD7, 0x00)
C_SILVER = RGBColor(0xA8, 0xB4, 0xC8)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM  = RGBColor(0xE0, 0xD8, 0xC8)
C_LTGRY  = RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY   = RGBColor(0x88, 0x88, 0x88)
C_RED    = RGBColor(0xCC, 0x22, 0x22)
C_BLUE   = RGBColor(0x33, 0x55, 0xCC)
C_PUR    = RGBColor(0x88, 0x22, 0xAA)
C_PHOENIX= RGBColor(0xFF, 0x6B, 0x1A)
C_GREEN  = RGBColor(0x22, 0xCC, 0x66)
C_ORANGE = RGBColor(0xFF, 0x99, 0x00)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（深夜ネイビー×対角グリッド）──────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (6, 8, 20))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 70):
        draw.line([(i, 0), (0, i)], fill=(12, 16, 35), width=1)
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        draw.line([(0, y), (w, y)], fill=(int(20 * t), int(15 * t), int(40 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー ──────────────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s

def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT_B
    if color:
        run.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp

def hdr(slide, text, accent=None):
    accent = accent or C_TEAL
    rect(slide, Inches(0.15), Inches(0.08), Inches(9.7), Emu(420000), C_NAVY)
    rect(slide, Inches(0.15), Inches(0.08), Emu(60000), Emu(420000), accent)
    tb(slide, Inches(0.4), Inches(0.1), Inches(9.2), Emu(380000),
       text, 12, bold=True, color=C_GOLD, font=FONT_H)

def divider_v(slide, x, y, h, color):
    rect(slide, x, y, Emu(30000), h, color)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 横帯2本
    rect(s, Inches(0), Inches(1.9), Inches(10), Emu(5000), C_TEAL)
    rect(s, Inches(0), Inches(3.55), Inches(10), Emu(5000), C_GOLD)

    tb(s, Inches(0.5), Inches(0.3), Inches(9), Emu(360000),
       "AIゲーム性提案シリーズ", 15, color=C_SILVER, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.3), Inches(0.78), Inches(9.4), Emu(1000000),
       "全体まとめ", 52, bold=True, color=C_WHITE, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.3), Inches(1.68), Inches(9.4), Emu(280000),
       "── 機種分析 × 新規提案 ──", 14, color=C_SILVER, font=FONT_H, align=PP_ALIGN.CENTER)

    # 中段：2カテゴリ
    rect(s, Inches(1.0), Inches(2.05), Inches(3.5), Emu(580000), RGBColor(0x08, 0x20, 0x28))
    rect(s, Inches(1.0), Inches(2.05), Emu(55000), Emu(580000), C_TEAL)
    tb(s, Inches(1.15), Inches(2.12), Inches(3.2), Emu(260000),
       "機種分析", 12, bold=True, color=C_TEAL, font=FONT_H)
    tb(s, Inches(1.15), Inches(2.42), Inches(3.2), Emu(200000),
       "吉宗（真打）・ミリゴ", 9.5, color=C_CREAM)

    rect(s, Inches(5.5), Inches(2.05), Inches(3.5), Emu(580000), RGBColor(0x20, 0x18, 0x00))
    rect(s, Inches(5.5), Inches(2.05), Emu(55000), Emu(580000), C_GOLD)
    tb(s, Inches(5.65), Inches(2.12), Inches(3.2), Emu(260000),
       "新規提案", 12, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(5.65), Inches(2.42), Inches(3.2), Emu(200000),
       "祟り神の章・PHOENIX LEGACY", 9.5, color=C_CREAM)

    # 下段
    rect(s, Inches(0.5), Inches(3.72), Inches(9), Emu(820000), RGBColor(0x0C, 0x10, 0x20))
    tb(s, Inches(0.65), Inches(3.80), Inches(8.6), Emu(280000),
       "実稼働データ 453機種・173週（2022〜2026）の分析知見をベースに、", 9.5, color=C_LTGRY, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.65), Inches(4.12), Inches(8.6), Emu(280000),
       "AIと共に機種を「読む」→「設計する」まで一気通貫で取り組んだシリーズです。", 9.5, color=C_LTGRY, align=PP_ALIGN.CENTER)

    tb(s, Inches(7.8), Inches(5.2), Inches(2.0), Emu(300000),
       "2026.05  KEY企画", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: 全体マップ
# ══════════════════════════════════════════════════════════════
def s_map(prs):
    s = new_slide(prs)
    hdr(s, "全体マップ  ──  機種分析と新規提案の2軸で積み上げた知見")

    # 左ゾーン：機種分析
    rect(s, Inches(0.15), Inches(0.62), Inches(4.55), Inches(4.6),
         RGBColor(0x08, 0x16, 0x20))
    rect(s, Inches(0.15), Inches(0.62), Inches(4.55), Emu(320000), RGBColor(0x0C, 0x22, 0x30))
    tb(s, Inches(0.3), Inches(0.66), Inches(4.2), Emu(280000),
       "機種分析", 11, bold=True, color=C_TEAL, font=FONT_H, align=PP_ALIGN.CENTER)

    # 吉宗カード
    rect_b(s, Inches(0.3), Inches(1.12), Inches(4.1), Inches(1.5),
           RGBColor(0x10, 0x1A, 0x28), C_TEAL, 1.5)
    tb(s, Inches(0.42), Inches(1.18), Inches(3.8), Emu(300000),
       "真打吉宗", 12, bold=True, color=C_TEAL, font=FONT_H)
    tb(s, Inches(0.42), Inches(1.55), Inches(3.8), Emu(560000),
       "1G連・ストックAT設計の徹底解剖\n説明資料 v11 まで進化（10スライド構成）\n純増4.8枚 / 周期CZ設計 / 天昇ボーナス",
       8.5, color=C_CREAM)

    # ミリゴカード
    rect_b(s, Inches(0.3), Inches(2.78), Inches(4.1), Inches(1.5),
           RGBColor(0x10, 0x1A, 0x28), C_SILVER, 1.5)
    tb(s, Inches(0.42), Inches(2.84), Inches(3.8), Emu(300000),
       "ミリオンゴッド", 12, bold=True, color=C_SILVER, font=FONT_H)
    tb(s, Inches(0.42), Inches(3.21), Inches(3.8), Emu(560000),
       "PGG（1/16,384）・A〜Dストック設計の分析\n説明資料 v2 完成\nZ-ZONE仕様・純増約8枚/G",
       8.5, color=C_CREAM)

    # 中央矢印＋ラベル
    rect(s, Inches(4.8), Inches(1.8), Emu(380000), Emu(30000), C_SILVER)
    rect(s, Inches(4.8), Inches(2.7), Emu(380000), Emu(30000), C_SILVER)
    tb(s, Inches(4.72), Inches(2.1), Emu(600000), Emu(500000),
       "知見\n活用", 9, bold=True, color=C_SILVER, align=PP_ALIGN.CENTER)
    tb(s, Inches(4.68), Inches(2.32), Emu(720000), Emu(260000),
       "→", 22, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 右ゾーン：新規提案
    rect(s, Inches(5.3), Inches(0.62), Inches(4.55), Inches(4.6),
         RGBColor(0x1A, 0x14, 0x00))
    rect(s, Inches(5.3), Inches(0.62), Inches(4.55), Emu(320000), RGBColor(0x28, 0x1C, 0x00))
    tb(s, Inches(5.45), Inches(0.66), Inches(4.2), Emu(280000),
       "新規提案", 11, bold=True, color=C_GOLD, font=FONT_H, align=PP_ALIGN.CENTER)

    # 祟り神カード
    rect_b(s, Inches(5.45), Inches(1.12), Inches(4.1), Inches(1.5),
           RGBColor(0x20, 0x08, 0x08), C_RED, 1.5)
    tb(s, Inches(5.57), Inches(1.18), Inches(3.8), Emu(300000),
       "祟り神の章", 12, bold=True, color=C_RED, font=FONT_H)
    tb(s, Inches(5.57), Inches(1.55), Inches(3.8), Emu(560000),
       "感情逆転設計・DQ4型章システム・加護の積み上げ\nv1→v2 に進化（游明朝・蛇行フロー追加）\n3設計核を軸としたオリジナル企画",
       8.5, color=C_CREAM)

    # PHOENIXカード
    rect_b(s, Inches(5.45), Inches(2.78), Inches(4.1), Inches(1.5),
           RGBColor(0x22, 0x10, 0x00), C_PHOENIX, 1.5)
    tb(s, Inches(5.57), Inches(2.84), Inches(3.8), Emu(300000),
       "PHOENIX LEGACY", 12, bold=True, color=C_PHOENIX, font=FONT_H)
    tb(s, Inches(5.57), Inches(3.21), Inches(3.8), Emu(560000),
       "炎の継承・IGNITION RITE・NEMESIS GAUGE\nv1→v2 に進化（フロー修正・業界初表現削除）\n吉宗×ミリゴ×祟り神 の知見を統合した集大成",
       8.5, color=C_CREAM)

    # 底部ライン
    rect(s, Inches(0.15), Inches(5.28), Inches(9.7), Emu(30000), C_SILVER)
    tb(s, Inches(0.3), Inches(5.32), Inches(9.4), Emu(240000),
       "機種分析の知見 → オリジナル設計に昇華  /  AIが一貫してサポートした企画プロセス",
       8.5, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 機種分析① 吉宗
# ══════════════════════════════════════════════════════════════
def s_yoshimune(prs):
    s = new_slide(prs)
    hdr(s, "機種分析①  ──  真打吉宗：1G連設計の完全解剖", accent=C_TEAL)

    # 左：設計の核心
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.55), Inches(3.5),
           C_CARD, C_TEAL, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(320000),
       "設計の核心", 11, bold=True, color=C_TEAL, font=FONT_H)

    points = [
        (C_TEAL,   "1G連（ストックAT）",
                   "ストックを1個消費して1GのATを繰り返す。\n消費枚数3枚で純増4.8枚 → 1G消化で約1.8枚増。"),
        (C_SILVER, "周期CZ設計",
                   "100G毎に規則的にCZが到来する「予測可能」設計。\nストレスなく「次はいつ来る」が分かる通常時。"),
        (C_GOLD,   "天昇ボーナス",
                   "通常時AT間のスペシャルボーナス。\n告知演出の多彩さが「吉宗らしさ」の象徴。"),
    ]
    py = Inches(1.38)
    for col, title, desc in points:
        rect(s, Inches(0.28), py, Emu(50000), Emu(540000), col)
        tb(s, Inches(0.48), py + Emu(30000), Inches(4.1), Emu(260000),
           title, 9.5, bold=True, color=col)
        tb(s, Inches(0.48), py + Emu(290000), Inches(4.1), Emu(280000),
           desc, 8.5, color=C_CREAM)
        py += Emu(590000)

    # 右：資料の進化
    rect_b(s, Inches(5.0), Inches(0.85), Inches(4.75), Inches(3.5),
           C_CARD, C_SILVER, 1.5)
    tb(s, Inches(5.12), Inches(0.92), Inches(4.5), Emu(320000),
       "説明資料の進化（v1 → v11）", 11, bold=True, color=C_SILVER, font=FONT_H)

    versions = [
        ("v1〜v3",  "基本スペック・フロー図の整備"),
        ("v4〜v6",  "演出パターン・ゾーン解説を追加"),
        ("v7〜v9",  "スライドデザイン統一（游明朝導入）"),
        ("v10",     "net_note・ベンチマーク比較追加"),
        ("v11",     "最終版・10スライド構成で完結"),
    ]
    vy = Inches(1.38)
    for ver, desc in versions:
        bg = RGBColor(0x12, 0x12, 0x26) if versions.index((ver, desc)) % 2 == 0 \
             else RGBColor(0x16, 0x16, 0x2E)
        rect(s, Inches(5.08), vy, Inches(4.58), Emu(320000), bg)
        tb(s, Inches(5.15), vy + Emu(28000), Inches(0.85), Emu(270000),
           ver, 8.5, bold=True, color=C_TEAL, wrap=False)
        tb(s, Inches(6.05), vy + Emu(28000), Inches(3.5), Emu(270000),
           desc, 8.5, color=C_CREAM, wrap=False)
        vy += Emu(330000)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(580000), RGBColor(0x08, 0x18, 0x20))
    rect(s, Inches(0.2), Inches(4.45), Emu(55000), Emu(580000), C_TEAL)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.1), Emu(250000),
       "吉宗で学んだこと：「周期設計の読みやすさ」と「1G連の爽快感」がリピート来店の構造的理由",
       9.5, bold=True, color=C_TEAL)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.1), Emu(250000),
       "→ PHOENIX LEGACYの周期AWAKENINGとEMBER（残り火5G）に直接継承された設計知見",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 機種分析② ミリゴ
# ══════════════════════════════════════════════════════════════
def s_milliongod(prs):
    s = new_slide(prs)
    hdr(s, "機種分析②  ──  ミリオンゴッド：夢の設計とストック管理の仕組み", accent=C_SILVER)

    # 左：設計の核心
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.55), Inches(3.5),
           C_CARD, C_SILVER, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(320000),
       "設計の核心", 11, bold=True, color=C_SILVER, font=FONT_H)

    points = [
        (C_GOLD2,  "PGG（プレミアムゴッドゲーム）",
                   "確率1/16,384の超大当たり。\n「夢を買う」体験として機能するゲーム性の頂点。"),
        (C_SILVER, "A〜Dストック設計",
                   "複数のストック種別が蓄積・放出を繰り返す。\n「今日の機械の調子」をストック状況で読む文化。"),
        (C_ORANGE, "Z-ZONE（ゾーン管理）",
                   "特定G数でAT突入確率が上昇するゾーン設計。\n立ち回り指標としてユーザーに深く浸透。"),
    ]
    py = Inches(1.38)
    for col, title, desc in points:
        rect(s, Inches(0.28), py, Emu(50000), Emu(540000), col)
        tb(s, Inches(0.48), py + Emu(30000), Inches(4.1), Emu(260000),
           title, 9.5, bold=True, color=col)
        tb(s, Inches(0.48), py + Emu(290000), Inches(4.1), Emu(280000),
           desc, 8.5, color=C_CREAM)
        py += Emu(590000)

    # 右：スペック概要
    rect_b(s, Inches(5.0), Inches(0.85), Inches(4.75), Inches(3.5),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(5.12), Inches(0.92), Inches(4.5), Emu(320000),
       "スペック概要と分析ポイント", 11, bold=True, color=C_GOLD, font=FONT_H)

    items = [
        ("タイプ",     "スマスロ対応 / ユニバーサルエンターテインメント"),
        ("純増",       "約8.0枚/G（神々の系譜シリーズ最高水準）"),
        ("PGG確率",   "1/16,384（約0.006%）の超プレミアム"),
        ("ストック",   "A・B・C・D 4種類の蓄積と連動放出"),
        ("CV値",       "高分散設計（番長4に近い高CV水準）"),
        ("来店動機",   "「今日PGGが出るかも」という期待値外の夢"),
    ]
    iy = Inches(1.38)
    for j, (k, v) in enumerate(items):
        bg = RGBColor(0x12, 0x12, 0x26) if j % 2 == 0 else RGBColor(0x16, 0x16, 0x2E)
        rect(s, Inches(5.08), iy, Inches(4.58), Emu(310000), bg)
        tb(s, Inches(5.15), iy + Emu(25000), Inches(1.15), Emu(265000),
           k, 8.5, bold=True, color=C_GOLD, wrap=False)
        tb(s, Inches(6.35), iy + Emu(25000), Inches(3.2), Emu(265000),
           v, 8.5, color=C_CREAM, wrap=False)
        iy += Emu(320000)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(580000), RGBColor(0x16, 0x14, 0x08))
    rect(s, Inches(0.2), Inches(4.45), Emu(55000), Emu(580000), C_GOLD)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.1), Emu(250000),
       "ミリゴで学んだこと：「夢の一枚」と「ストック管理」がリピートの2大動機",
       9.5, bold=True, color=C_GOLD)
    tb(s, Inches(0.45), Inches(4.83), Inches(9.1), Emu(250000),
       "→ PHOENIX LEGACYのA〜DストックをFLAMEストック4種に昇華・Z-ZONEをIGNITION RITEに再解釈",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 新規提案① 祟り神の章
# ══════════════════════════════════════════════════════════════
def s_atarigami(prs):
    s = new_slide(prs)
    hdr(s, "新規提案①  ──  祟り神の章：感情逆転設計の原点", accent=C_RED)

    # 3カラム
    cols = [
        (Inches(0.2), C_RED,
         "I  感情逆転設計",
         "普通ATで「倒す」\n特別ATで「理解する」\n\n敵 → 共感 への\n感情の逆転が\n来店継続の核心\n\n「倒した相手を\n後から理解する」"),
        (Inches(3.55), C_PUR,
         "II  DQ4型章システム",
         "1〜4章をランダム順解放\n視点ごとに同じ事件を語る\n\n「俺は3章から入った、\n お前は？」\nホールで会話が生まれる\n\n章クリアで加護を獲得\n→ 来店継続設計"),
        (Inches(6.9), C_BLUE,
         "III  第0章・成仏設計",
         "全章クリア後に解放\n成仏チャレンジ\n成功 1500枚\n失敗 100枚\n\n「失敗も次の動機に」\n物語がまだ終わらない\n設計"),
    ]
    for x, col, title, desc in cols:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.45),
               RGBColor(0x0C, 0x08, 0x18), col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(340000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.40), Inches(2.8), Inches(2.6),
           desc, 9, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(40000), C_RED)
    rect(s, Inches(0.2), Inches(4.48), Inches(9.6), Emu(640000), RGBColor(0x14, 0x08, 0x18))
    tb(s, Inches(0.35), Inches(4.53), Inches(9.2), Emu(270000),
       "v1（2025）→ v2（2026.05）：游明朝・蛇行フロー・純Pillow背景を導入。スタイルが格段に進化。",
       9, bold=True, color=C_RED)
    tb(s, Inches(0.35), Inches(4.85), Inches(9.2), Emu(260000),
       "シリーズ最初のオリジナル提案。「感情設計」という軸がPHOENIX LEGACYにも受け継がれた。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 新規提案② PHOENIX LEGACY
# ══════════════════════════════════════════════════════════════
def s_phoenix(prs):
    s = new_slide(prs)
    hdr(s, "新規提案②  ──  PHOENIX LEGACY：集大成の炎", accent=C_PHOENIX)

    # 3カラム
    cols = [
        (Inches(0.2), C_PHOENIX,
         "I  PHOENIX CYCLE",
         "1AT = 1世代の「炎」\nATが終わるたびに\nREBIRTH演出が流れ\n次の炎が点火する\n\n「また負けた」→\n「炎を受け継いだ」\nという感情設計"),
        (Inches(3.55), C_GOLD,
         "II  IGNITION RITE",
         "AT中に赤7揃いで発動\nSkill / Bond / Power\n3択の点火方法を選ぶ\n\n選んだ内容で\n次ATの性能が変わる\n\nプレイヤーが\n物語を作る台"),
        (Inches(6.9), RGBColor(0xA8, 0x1C, 0x1C),
         "III  NEMESIS GAUGE",
         "ATを重ねるほど\n宿敵が弱体化する\n1セッション完結設計\n\nゲージMAX+\nPHOENIX FLAME+\nIGNITION完全制覇\nでETERNAL FLAME"),
    ]
    for x, col, title, desc in cols:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.45),
               RGBColor(0x18, 0x08, 0x00), col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(340000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.40), Inches(2.8), Inches(2.6),
           desc, 9, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(40000), C_PHOENIX)
    rect(s, Inches(0.2), Inches(4.48), Inches(9.6), Emu(640000), RGBColor(0x1C, 0x0C, 0x00))
    tb(s, Inches(0.35), Inches(4.53), Inches(9.2), Emu(270000),
       "吉宗の周期設計・ミリゴのストック設計・祟り神の感情設計 ── 3機種の知見を統合した集大成。",
       9, bold=True, color=C_PHOENIX)
    tb(s, Inches(0.35), Inches(4.85), Inches(9.2), Emu(260000),
       "v1→v2でフロー表示・業界初表現・構成順の課題を修正。「育てれば届く」MY約3500枚の設計。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 4作品を貫く設計哲学
# ══════════════════════════════════════════════════════════════
def s_philosophy(prs):
    s = new_slide(prs)
    hdr(s, "4作品を貫く設計哲学  ──  数字と感情の両立")

    # 3つの軸
    axes = [
        (Inches(0.2), C_TEAL,
         "CV 0.20〜0.25を目標とする",
         "東京喰種・モンキーターンV水準\n\n「設定1でも体験の質が変わらない台」\nCV高すぎ（番長4: 0.68）は荒すぎる\nCV低すぎはリピートが生まれない\n\n全提案でこの水準を一貫して設計目標に"),
        (Inches(3.55), C_GOLD,
         "後半維持率 65%以上",
         "3ヶ月後も稼働している台を目指す\n\n番長4は後半維持率29.6%で失速\n東京喰種73.7% / モンキーV 68.2%\nが来店継続の実績ベンチマーク\n\n提案目標：65%超・3ヶ月以上稼働"),
        (Inches(6.9), C_SILVER,
         "「また来たい」を設計に組み込む",
         "出玉だけでは来店理由にならない時代\n\n祟り神：章の続きが気になる\nPHOENIX：炎を育てたい\n吉宗：次の周期CZが来る\nミリゴ：今日PGGが出るかも\n\n感情的な「また来る理由」が設計の核心"),
    ]
    for x, col, title, desc in axes:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.6),
               C_CARD, col, 1.8)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(340000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.40), Inches(2.8), Inches(2.75),
           desc, 8.5, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.55), Inches(9.6), Emu(540000), RGBColor(0x10, 0x10, 0x10))
    tb(s, Inches(0.35), Inches(4.60), Inches(9.2), Emu(450000),
       "453機種・173週の実稼働データを分析してたどり着いた3原則。\n"
       "このフレームがあるから、提案に「なぜそのスペックなのか」の根拠が生まれる。",
       9.5, bold=True, color=C_LTGRY)
    tb(s, Inches(8.5), Inches(5.38), Inches(1.4), Emu(200000),
       "※分析データより", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: まとめ・次の展開
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ  ──  AIと積み上げたゲーム性設計の軌跡")

    # 上段：成果サマリー
    rect(s, Inches(0.2), Inches(0.85), Inches(9.6), Emu(360000), RGBColor(0x10, 0x14, 0x28))
    tb(s, Inches(0.35), Inches(0.92), Inches(9.2), Emu(300000),
       "このシリーズで生まれたもの：2機種の分析資料  ＋  2本のオリジナル企画  ＋  設計哲学3原則",
       10, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, font=FONT_H)

    # 中段：4作品横断比較
    works = [
        ("吉宗",         "機種分析", "周期CZ・1G連",   "→ PHOENIX AWAKENINGに継承", C_TEAL),
        ("ミリゴ",        "機種分析", "ストック・Z-ZONE","→ FLAME4種・IGNITIONに継承", C_SILVER),
        ("祟り神の章",    "新規提案", "感情逆転・章設計", "→ 感情積み上げの原点を確立",  C_RED),
        ("PHOENIX L.",   "新規提案", "炎の継承・3軸",   "→ 全知見を統合した集大成",    C_PHOENIX),
    ]
    wy = Inches(1.38)
    for wname, wcat, wcore, wresult, wcol in works:
        bg = RGBColor(0x12, 0x12, 0x24)
        rect(s, Inches(0.2), wy, Inches(9.6), Emu(330000), bg)
        rect(s, Inches(0.2), wy, Emu(55000), Emu(330000), wcol)
        tb(s, Inches(0.45), wy + Emu(28000), Inches(1.55), Emu(278000),
           wname, 9.5, bold=True, color=wcol, wrap=False)
        tb(s, Inches(2.05), wy + Emu(28000), Inches(1.2), Emu(278000),
           wcat, 8.5, color=C_GRAY, wrap=False)
        tb(s, Inches(3.3), wy + Emu(28000), Inches(2.8), Emu(278000),
           wcore, 9, bold=True, color=C_WHITE, wrap=False)
        tb(s, Inches(6.15), wy + Emu(28000), Inches(3.5), Emu(278000),
           wresult, 8.5, color=C_CREAM, wrap=False)
        wy += Emu(340000)

    # 下段：次の展開
    rect(s, Inches(0.2), Inches(4.42), Inches(9.6), Emu(660000), RGBColor(0x0C, 0x18, 0x28))
    rect(s, Inches(0.2), Inches(4.42), Emu(55000), Emu(660000), C_TEAL)
    tb(s, Inches(0.45), Inches(4.48), Inches(9.1), Emu(270000),
       "次の展開として考えられること", 10, bold=True, color=C_TEAL, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.1), Emu(260000),
       "祟り神×PHOENIX のクロスオーバー提案  /  新機種の追加分析  /  提案をさらに肉付けした実現性検討",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = [
        ("タイトル",           s_title),
        ("全体マップ",         s_map),
        ("機種分析① 吉宗",    s_yoshimune),
        ("機種分析② ミリゴ",  s_milliongod),
        ("新規提案① 祟り神",  s_atarigami),
        ("新規提案② PHOENIX", s_phoenix),
        ("設計哲学",           s_philosophy),
        ("まとめ",             s_matome),
    ]

    print("=" * 55)
    print("  AIゲーム性提案シリーズ 全体まとめ ジェネレーター")
    print("=" * 55)
    print()
    for i, (name, func) in enumerate(slides, 1):
        print(f"  {i:2d}/{len(slides)} {name}")
        func(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
