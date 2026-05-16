"""
スマスロ いざ！番長 機種説明＋分析 統合版 PPTXジェネレーター v1
出力: proposals/機種分析/いざ番長/bansho_guide_v1.pptx
テーマ: 深黒×青(#2244CC)×金(#C8A840)×白（番長・硬派）
情報源: ちょんぼりすた・一撃・altema・p-town DMM・なな徹（2025年〜）
パチスロアワード2025ノミネート機種
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(ROOT_DIR,
           "proposals", "機種分析", "いざ番長", "bansho_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深黒×青×金×白）──────────────────────────
C_BG    = RGBColor(0x06, 0x08, 0x10)
C_CARD  = RGBColor(0x0C, 0x10, 0x1E)
C_CARD2 = RGBColor(0x10, 0x16, 0x28)
C_ROW   = RGBColor(0x0E, 0x12, 0x22)
C_BLUE  = RGBColor(0x22, 0x44, 0xCC)   # メインブルー #2244CC
C_BLUE2 = RGBColor(0x33, 0x66, 0xFF)   # 明るいブルー
C_LBLUE = RGBColor(0x44, 0x88, 0xFF)   # ライトブルー
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金 #C8A840
C_GOLD2 = RGBColor(0xFF, 0xD0, 0x60)   # 明るい金
C_WHITE = RGBColor(0xF0, 0xEC, 0xE8)
C_CREAM = RGBColor(0xF0, 0xE0, 0xB8)
C_GRAY  = RGBColor(0x88, 0x90, 0xAA)
C_LTGRY = RGBColor(0x44, 0x48, 0x60)
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_RED   = RGBColor(0xDD, 0x22, 0x22)
C_PINK  = RGBColor(0xFF, 0x44, 0x88)
C_SBLUE = RGBColor(0x11, 0x22, 0x88)   # 暗い青（背景帯）

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景・ヘルパー群 ─────────────────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (6, 8, 16))
    draw = ImageDraw.Draw(img)
    # 斜線ライン（硬派感）
    for i in range(0, w + h, 100):
        draw.line([(i, 0), (0, i)], fill=(10, 14, 28), width=1)
    # 下端に青のグラデーション
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        b_val = int(50 * t)
        draw.line([(0, y), (w, y)], fill=(0, int(20 * t), b_val))
    # 上端に青のほんのりライン
    for y in range(0, 30):
        t = (30 - y) / 30 * 0.4
        draw.line([(0, y), (w, y)], fill=(0, 0, int(20 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_BLUE)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GOLD, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_BLUE)


def net_note(slide):
    tb(slide, Inches(7.8), Inches(5.38), Inches(2.1), Emu(180000),
       "※ネット解析情報より（ちょんぼりすた・一撃・altema・p-town）", 6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x08, 0x0C, 0x1C))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_BLUE)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_GOLD)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_BLUE
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_BLUE
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x08, 0x0A, 0x18))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_BLUE)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_BLUE)

    tb(s, Inches(0.22), Inches(0.4), Inches(5.0), Emu(330000),
       "機種説明＋分析 統合ガイド  v1　【パチスロアワード2025ノミネート】", 9.5, color=C_LBLUE, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.85), Inches(5.1), Emu(900000),
       "いざ！番長", 38, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.6), Inches(5.0), Emu(280000),
       "スマスロ（大都技研）── 差枚管理型AT×刺客ゾーン×絶頂輪廻ループの設計", 9, color=C_CREAM, font=FONT_H)

    # スペック表
    specs = [
        ("メーカー",        "大都技研（サボハニ）　2025年6月2日導入"),
        ("設定",           "1〜6段階"),
        ("AT純増（通常）",   "約2.8枚/G（頂ZBASH）"),
        ("AT純増（上位）",   "約5.0枚/G（青頂ZBASH）"),
        ("設定1機械割",     "約97%"),
        ("設定6機械割",     "約111%"),
        ("天井",           "モード別規定G数（モード管理型）"),
        ("引き戻し状態",    "AT終了後あり（前兆確認後ヤメ）"),
        ("特徴",           "絶頂輪廻ループ・刺客ゾーン(CZ)・番長BB"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(3.0) + i * Emu(200000)
        tb(s, Inches(0.22), ry, Inches(1.85), Emu(185000),
           k, 7, color=C_GRAY)
        tb(s, Inches(2.1), ry, Inches(3.1), Emu(185000),
           v, 7, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_BLUE,  "① 刺客ゾーン（CZ）自力感設計",
         "通常時に刀ポイント到達でCZ「刺客ゾーン」突入。\n10G+αの自力バトルで約40%のAT当選。\nレア役ヒットでゾーン延長・色昇格の演出。"),
        (C_GOLD,  "② 絶頂輪廻ループ（最強連鎖）",
         "特化ゾーン「絶頂決戦〜巌流島〜」(平均+700枚)\n終了後→上位AT「青頂ZBASH」(純増5.0枚)へ直行。\n2つの頂をシームレスに行き来する絶頂輪廻。"),
        (C_CYAN,  "③ 番長ボーナスによる上乗せ加速",
         "赤BB（20G+α・上乗せ高確抽選）\n青BB（20G・常時上乗せ高確）で枚数大幅加速。\n差枚管理型の上乗せが出玉体験を豊かにする。"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.55) + i * Emu(1540000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.3), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.3), ac)
        tb(s, Inches(5.85), y0 + Emu(65000), Inches(3.8), Emu(310000),
           kw, 11.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(380000), Inches(3.8), Emu(450000),
           desc, 8.5, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「刺客ゾーン自力感×絶頂輪廻ループ×番長ボーナス上乗せ」── 硬派番長シリーズ第3弾",
           "通常AT2.8枚/G・上位AT5.0枚/Gの2段階純増が出玉波の緩急を生む")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（全ルートを蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→AT→上位ATの全ルート", "2/9")

    # 上段：通常→AT突入ルート
    top_y = Inches(0.75)
    top_h = Emu(1100000)

    flow1 = [
        (C_CARD2, C_GRAY,  "通常遊技",          "刀ポイント蓄積\nモード別規定G数消化\n→CZ/前兆へ"),
        (C_CARD,  C_BLUE2, "刺客ゾーン\n(CZ)",   "10G+α\n約40%でAT当選\n自力バトル"),
        (C_CARD,  C_GOLD,  "修行\n（前兆）",     "対決バトルで\nAT成否決定\n番長演出炸裂"),
        (C_CARD,  C_GOLD2, "頂ZBASH\n(AT)",      "純増約2.8枚/G\n差枚管理型\n初期150枚+α"),
        (C_CARD,  C_LBLUE, "番長ボーナス\n(BB)", "赤/青BBで\n上乗せ加速\n20G+α"),
    ]
    bw1 = Inches(1.6)
    gap1 = Inches(0.17)
    sx1 = Inches(0.28)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), top_y + Emu(70000), bw1 - Emu(60000), Emu(360000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), top_y + Emu(480000), bw1 - Emu(45000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), cy1, C_BLUE)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1120000), Inches(5.0), Emu(260000),
       "天井：モード別規定G数（設定・モードで異なる） ⇒ 修行→対決でAT当選", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1360000), Inches(5.5), Emu(5000), C_CYAN)

    # 中段区切り線
    rect(s, 0, Inches(2.1), SLIDE_W, Emu(5000), RGBColor(0x22, 0x33, 0x66))

    # 下段：AT内ループ昇格ルート
    bot_y = Inches(2.18)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,  C_BLUE2, "頂ZBASH\n（通常AT）",      "純増2.8枚/G\n決闘で上乗せ\nBB抽選"),
        (C_CARD,  C_GOLD,  "絶頂決戦\n〜巌流島〜",     "差枚上乗せ特化\n平均+700枚\n特化ゾーン"),
        (C_CARD,  C_LBLUE, "青頂ZBASH\n（上位AT）",    "純増5.0枚/G\n上乗せ強化\n3桁上乗せ濃厚"),
        (RGBColor(0x08,0x10,0x28), C_GOLD2,
         "絶頂輪廻ループ",  "巌流島終了後\n→青頂ZBASH\n繰り返しループ"),
        (C_CARD,  C_CYAN,  "引き戻し\n状態(AT後)",    "前兆なし\n確認後\nヤメ推奨"),
    ]
    bw2 = Inches(1.65)
    gap2 = Inches(0.17)
    sx2 = Inches(0.3)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), bot_y + Emu(70000), bw2 - Emu(60000), Emu(370000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), bot_y + Emu(490000), bw2 - Emu(45000), Emu(470000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw2 + Emu(8000), cy2, C_GOLD)

    # 巌流島→ループアノテーション
    tb(s, sx2, bot_y + Emu(1100000), Inches(5.5), Emu(260000),
       "絶頂輪廻ループ：巌流島終了→青頂ZBASH→巌流島→…の繰り返しが出玉爆発の核心", 7.5, color=C_GOLD)
    rect(s, sx2, bot_y + Emu(1340000), Inches(5.5), Emu(5000), C_GOLD)

    net_note(s)
    footer(s, "上段=通常〜AT突入ルート（刺客ゾーンCZ経由・モード天井経由の2本）、下段=AT内の絶頂輪廻ループ構造",
           "刺客ゾーン：AT期待度約40%の自力CZ / 絶頂輪廻ループ：巌流島特化ゾーン→青頂ZBASH→繰り返し")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（ルート・天井含む全ルートを明確に）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 全AT突入ルート・天井管理・刀ポイント", "3/9")

    # 左：AT突入ルート図
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x15, 0x22, 0x55))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜AT突入ルート（全3系統）", 10, bold=True, color=C_GOLD)

    routes = [
        (C_BLUE2, "ルート①  刺客ゾーン（CZ）経由",
         "刀ポイントが規定量に到達するとCZ「刺客ゾーン」へ。\n10G+αの自力バトル、約40%でAT当選。\nレア役ヒットで色昇格・残G延長の恩恵あり。"),
        (C_GOLD,  "ルート②  規定G数（モード天井）到達→修行→対決",
         "モードによって規定G数が異なり、天井到達で修行（前兆）へ。\n修行ステージでの対決バトルでAT成否が決定する。\n一部は刺客ゾーンをスキップして直接修行へ。"),
        (C_CYAN,  "ルート③  レア役直撃",
         "強チェリー・チャンス目等のレア役成立時は\nAT直当り抽選または刺客ゾーン高確移行を抽選。\n通常時も常に役成立に期待を持てる設計。"),
        (C_LBLUE, "天井（モード管理）",
         "設定・モードで変化する規定G数を管理。\n通常A/B・チャンスモード・天国モード等が存在。\nAT後は引き戻し状態を必ず確認してからヤメること。"),
    ]
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(290000) + i * Emu(1100000)
        rect_b(s, lx, iy, lw, Emu(1040000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1040000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(650000),
           b, 7.5, color=C_WHITE)

    # 右：刀ポイントとチャンス役
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x15, 0x22, 0x55))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "刀ポイントとチャンス役の期待度", 10, bold=True, color=C_GOLD)

    chance = [
        (C_GRAY,  "ベル/ナビ",    "毎G",      "刀ポイント蓄積の主役（規定量で刺客ゾーンへ）"),
        (C_BLUE2, "弱チェリー",  "約1/50",   "刀ポイント大量獲得・刺客ゾーン移行抽選"),
        (C_GOLD,  "スイカ",      "約1/100",  "刀ポイント大量獲得・AT直当り抽選"),
        (C_GOLD2, "強チェリー",  "約1/200",  "AT直当り濃厚・刺客ゾーン高確確定"),
        (C_LBLUE, "チャンス目",  "約1/100",  "刀ポイント加算・刺客ゾーン抽選"),
        (C_CYAN,  "レア役合算",  "約1/30前後","AT直当り・刺客ゾーン移行の全契機役"),
    ]
    ch_h = Emu(690000)
    for i, (ac, role, prob, desc) in enumerate(chance):
        cy = ry + Emu(290000) + i * ch_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, ch_h, bg)
        rect(s, rx, cy, Emu(35000), ch_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(55000), Inches(1.0), Emu(260000),
           role, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.0), cy + Emu(60000), Inches(0.85), Emu(240000),
           prob, 8, bold=True, color=C_GOLD2, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(310000), rw - Inches(0.6), Emu(320000),
           desc, 7.5, color=C_WHITE)

    # 打ち方メモ
    rect_b(s, rx, ry + Emu(4450000), rw, Emu(490000), C_CARD2, C_BLUE, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4490000), rw - Emu(80000), Emu(200000),
       "打ち方：通常は順押し推奨（機種ナビに従う）", 8.5, bold=True, color=C_BLUE2)
    tb(s, rx + Emu(60000), ry + Emu(4690000), rw - Emu(80000), Emu(220000),
       "AT終了後は引き戻し状態あり。前兆演出なしを確認してからヤメること。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "通常時の基本戦略：刀ポイント管理で刺客ゾーンを目指し、モード天井（規定G数）を把握して立ち回る",
           "引き戻し状態のヤメ判断は最重要。AT終了後すぐにヤメると当選を取り逃がすリスクあり")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み（番長バトル要素と絡めて）
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ/前兆の仕組み ── 刺客ゾーン×修行×対決バトル詳細", "4/9")

    # 左：刺客ゾーン詳細
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(270000), RGBColor(0x15, 0x22, 0x55))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(210000),
       "刺客ゾーン（CZ）── 自力ATバトル詳細", 10, bold=True, color=C_BLUE2)

    steps = [
        (C_GRAY,   "突入条件",
         "刀ポイントが規定量到達でCZへ。\n通常時の主要な攻略ポイント。\n規定量はモード・設定で変動する。"),
        (C_BLUE2,  "ゲーム数",
         "基本10G+αで消化。\nレア役成立で残G延長（＋数G）。\nチャンス演出が頻発するほど期待大。"),
        (C_GOLD,   "AT当選率",
         "CZ成功（AT当選）率：約40%以上。\n成功時は一部番長ボーナス（BB）に当選。\n成功後の一部で絶頂決戦〜巌流島〜突入。"),
        (C_LBLUE,  "演出・色昇格",
         "刺客ゾーン中の演出色が昇格するほど期待度UP。\n白→青→赤→金の順に当選期待度が高まる。\nレア役は色昇格のトリガー。"),
    ]
    for i, (ac, t, b) in enumerate(steps):
        iy = ly + Emu(270000) + i * Emu(1115000)
        rect_b(s, lx, iy, lw, Emu(1050000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1050000), ac)
        tb(s, lx + Emu(75000), iy + Emu(48000), lw - Emu(100000), Emu(260000),
           t, 9.5, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(650000),
           b, 7.5, color=C_WHITE)

    # 右：修行（前兆）と対決バトル
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(270000), RGBColor(0x15, 0x22, 0x55))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(210000),
       "修行（前兆ステージ）と対決バトル", 10, bold=True, color=C_GOLD)

    battles = [
        (C_GOLD,   "修行ステージへの移行",
         "規定G数到達・刺客ゾーン成功の後、\n修行ステージ（前兆）へ移行する。\n番長キャラが登場し対決前の緊張感を演出。"),
        (C_GOLD2,  "対決バトルの仕組み",
         "修行からの「対決」でAT入口を確定判定。\n番長が刺客と闘い、勝利すればAT突入。\nバトル中のキャラ・BGM・ナレーションで\n期待度を演出する番長シリーズ伝統の要素。"),
        (C_BLUE2,  "修行の法則・示唆",
         "修行中の服装・背景・セリフに\nAT当否や設定の示唆が含まれる場合がある。\n「逆境からの逆転」演出は当選濃厚のサイン。\n対決での敗北＝スルーはカウント管理が重要。"),
        (C_CYAN,   "スルー後の優遇",
         "修行→対決でのAT非当選（スルー）後は\n次回の当選期待度が高まる。\n低設定でも複数スルー後の当選チャンスは\n高く設定されている可能性あり。"),
    ]
    for i, (ac, t, b) in enumerate(battles):
        iy = ry + Emu(270000) + i * Emu(1115000)
        rect_b(s, rx, iy, rw, Emu(1050000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(45000), Emu(1050000), ac)
        tb(s, rx + Emu(70000), iy + Emu(48000), rw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(650000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "CZ設計の核心：刀ポイント管理（能動性）→刺客ゾーン（自力感）→修行→対決（番長バトル伝統）の三層構造",
           "刺客ゾーン中のレア役色昇格演出が期待度の可視化装置として機能。約40%というCZ成功率は適度な緊張感を維持")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（何をすれば出玉が伸びる）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「頂ZBASH」── 出玉を伸ばすための仕組みと戦略", "5/9")

    # 左：AT基本構造
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(270000), RGBColor(0x15, 0x22, 0x55))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(210000),
       "頂ZBASH（AT）── 基本構造と出玉メカニズム", 10, bold=True, color=C_BLUE2)

    at_items = [
        (C_BLUE2, "AT基本仕様",
         "差枚数管理型AT、純増約2.8枚/G。\n初期差枚数：150枚+α（刺客ゾーン成功時）。\n成立役に応じて決闘抽選・差枚上乗せ抽選を実施。"),
        (C_GOLD,  "番長ボーナス（BB）の種類",
         "赤BB（20G+α）：消化中に上乗せ高確抽選＋差枚上乗せ抽選。\n  上乗せ枚数：一撃100〜数百枚も可能。\n青BB（20G）：常時上乗せ高確状態。差枚上乗せ期待大。\n  青BBは赤BBより上乗せ期待値が高い。"),
        (C_GOLD2, "倍斬刀チャレンジ（上乗せ倍化）",
         "3G間のCZで、レア役・斬揃い成立で上乗せ枚数が\n2〜8倍に倍化。倍化成功を重ねると最大64倍まで可能。\n獲得差枚数の爆発的増加の起点となる特化CZ。"),
        (C_CYAN,  "決闘と差枚上乗せ",
         "AT中は成立役ごとに決闘抽選を実施。\n決闘勝利で報酬獲得（差枚上乗せ/番長BB/絶頂決戦）。\n報酬の約25%が絶頂決戦〜巌流島〜突入のチャンス。"),
    ]
    for i, (ac, t, b) in enumerate(at_items):
        iy = ly + Emu(270000) + i * Emu(1115000)
        rect_b(s, lx, iy, lw, Emu(1050000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1050000), ac)
        tb(s, lx + Emu(75000), iy + Emu(48000), lw - Emu(100000), Emu(260000),
           t, 9.5, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(650000),
           b, 7.5, color=C_WHITE)

    # 右：出玉を伸ばすポイント
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(270000), RGBColor(0x15, 0x22, 0x55))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(210000),
       "出玉を伸ばす3つのポイント", 10, bold=True, color=C_GOLD)

    extend = [
        (C_GOLD,   "① 青BBを引く・引かせる",
         "青BBは常時上乗せ高確のため1BB当たりの\n期待差枚数が赤BBを大きく上回る。\nAT中に青BBを何回引けるかが出玉規模を決める。"),
        (C_BLUE2,  "② 倍斬刀チャレンジで倍化積み上げ",
         "上乗せ枚数の倍化チャンスを活かすことが\n大量出玉への最短ルート。\n2倍→4倍→8倍→16倍→32倍→64倍の段階があり\nレア役が鍵を握る。"),
        (C_GOLD2,  "③ 絶頂決戦〜巌流島〜への到達",
         "決闘勝利報酬で突入できる差枚上乗せ特化ゾーン。\n平均上乗せ約700枚で終了後は上位AT「青頂ZBASH」へ。\n絶頂輪廻ループへの入口として最重要の特化ゾーン。"),
    ]
    for i, (ac, t, b) in enumerate(extend):
        iy = ry + Emu(270000) + i * Emu(1490000)
        rect_b(s, rx, iy, rw, Emu(1420000), C_CARD, ac, 2.0)
        rect(s, rx, iy, Emu(45000), Emu(1420000), ac)
        tb(s, rx + Emu(75000), iy + Emu(50000), rw - Emu(95000), Emu(270000),
           t, 10, bold=True, color=ac, font=FONT_H)
        tb(s, rx + Emu(75000), iy + Emu(340000), rw - Emu(95000), Emu(1000000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "出玉の核心：「青BB連打×倍斬刀チャレンジの倍化×絶頂決戦〜巌流島〜突入」── この3要素が重なると大量出玉",
           "差枚管理型は消化中の上乗せが全て「残差枚数」として積み上がる設計。番長BBが出玉加速装置として機能")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方（絶頂輪廻ループ）
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位ATへの道 ── 絶頂決戦〜巌流島〜×青頂ZBASH×絶頂輪廻ループ", "6/9")

    # 上段：上位AT到達ルートフロー
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(260000), RGBColor(0x15, 0x22, 0x55))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(210000),
       "上位AT到達ルートと絶頂輪廻ループの全容", 9, bold=True, color=C_GOLD)

    route_boxes = [
        (C_BLUE2, "頂ZBASH\n（通常AT）",       "純増2.8枚/G\n決闘抽選\n起点"),
        (C_GOLD,  "決闘勝利\n→報酬獲得",       "報酬の約25%が\n巌流島突入\nチャンス"),
        (C_GOLD2, "絶頂決戦\n〜巌流島〜",      "差枚特化ゾーン\n平均+700枚\n終了後→上位AT"),
        (C_LBLUE, "青頂ZBASH\n（上位AT）",      "純増5.0枚/G\n3桁上乗せ濃厚\n超高性能AT"),
    ]
    bw_r = Inches(2.1)
    gap_r = Inches(0.26)
    sx_r = Inches(0.35)
    cy_r = Inches(1.42)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x08, 0x12, 0x28), ac, 2.0 if i >= 2 else 1.5)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(55000), Emu(370000), lbl, 10, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), cy_r - bh_r // 2 + Emu(460000),
           bw_r - Emu(45000), Emu(480000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw_r + Emu(12000), cy_r, C_GOLD)

    # ループ矢印アノテーション
    tb(s, sx_r + 2 * (bw_r + gap_r), Inches(2.0), Inches(4.5), Emu(240000),
       "← 絶頂輪廻ループ：青頂ZBASH→巌流島→青頂ZBASH→…を繰り返す", 8, bold=True, color=C_GOLD2)
    rect(s, sx_r + 2 * (bw_r + gap_r), Inches(2.26), Inches(4.5), Emu(5000), C_GOLD)

    # 中区切り
    rect(s, 0, Inches(2.36), SLIDE_W, Emu(5000), RGBColor(0x22, 0x33, 0x66))

    # 下段左：絶頂決戦〜巌流島〜の遊び方
    lx2, ly2 = Inches(0.28), Inches(2.42)
    lw2 = Inches(4.55)
    lh2 = Emu(2500000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_GOLD, 1.8)
    rect(s, lx2, ly2, Emu(45000), lh2, C_GOLD)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "絶頂決戦〜巌流島〜の全容", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【突入条件】\n"
       "① AT頂ZBASH中の決闘勝利報酬（約25%で突入）\n"
       "② 刺客ゾーン成功時の一部で直接突入\n\n"
       "【特化ゾーンの内容】\n"
       "差枚数上乗せに特化した特化ゾーン。\n"
       "消化中の成立役に応じて大量の差枚数を上乗せ。\n"
       "平均上乗せ枚数：約700枚。\n\n"
       "【終了後の恩恵】\n"
       "絶頂決戦〜巌流島〜終了後は無条件で\n"
       "上位AT「青頂ZBASH」へ突入（絶頂輪廻ループ開始）。",
       8, color=C_WHITE)

    # 下段右：青頂ZBASHの遊び方
    rx2, ry2 = Inches(5.05), Inches(2.42)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2,
           RGBColor(0x08, 0x12, 0x28), C_LBLUE, 2.0)
    rect(s, rx2, ry2, Emu(45000), lh2, C_LBLUE)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "青頂ZBASH（上位AT）の遊び方", 11, bold=True, color=C_LBLUE, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【基本スペック】\n"
       "純増：約5.0枚/G（通常ATの約1.8倍）\n"
       "差枚数：巌流島で獲得した分から継続\n\n"
       "【通常ATとの違い】\n"
       "ゲーム性は頂ZBASHと同様の決闘システム。\n"
       "ただし、レア小役成立時の上乗せ確率が大幅UP。\n"
       "上乗せ時の枚数が3桁（100枚以上）濃厚。\n\n"
       "【絶頂輪廻ループへ】\n"
       "青頂ZBASH中も決闘→報酬で巌流島突入チャンス。\n"
       "巌流島→青頂ZBASH→巌流島→…の\n"
       "繰り返しループが最大の出玉爆発ルート。\n"
       "このループに入ることが最終目標。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "上位ATの設計核心：「巌流島(+700枚)→青頂ZBASH(純増5.0枚)→巌流島→…」の絶頂輪廻ループが出玉爆発の本体",
           "青頂ZBASH中の上乗せは3桁濃厚。ループに入れば自動的に大量獲得が実現する設計")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（番長シリーズの設計哲学）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── 番長シリーズの設計哲学といざ番長の進化", "7/9")

    principles = [
        (C_GOLD,   "① 「男の自力感」── 番長シリーズのDNA",
         "押忍！番長3から受け継ぐ「自力でATを引き当てる感覚」。\n刀ポイント管理→刺客ゾーン→対決という\n能動的なゲームフローが硬派プレイヤーを引きつける。"),
        (C_BLUE2,  "② 刀ポイントという「蓄積型エンジン」",
         "毎ゲームの役成立を刀ポイントとして積み上げる設計。\nベルが刺さるたびに「あと何G」という\n目標意識が生まれ、回転率と滞在率を高める。"),
        (C_GOLD2,  "③ 絶頂輪廻ループという「夢の設計」",
         "一度ループに入ると止まらない連鎖が明確に見える。\n「巌流島→青頂ZBASH→巌流島→…」という\nエンドレスの上位ループが\"夢\"として機能する。"),
        (C_LBLUE,  "④ 差枚管理型×倍化という「出玉体験の民主化」",
         "差枚数管理型はハマりでも残差枚が積み上がる。\n倍斬刀チャレンジの最大64倍倍化が\n低投資からでも一気に逆転できる「逆転劇」を可能にする。"),
        (C_CYAN,   "⑤ 2段階純増（2.8→5.0枚）という「到達感の演出」",
         "通常AT→上位ATで純増が大幅アップする構造が\n明確な「上振れ感」を生む。\n青頂ZBASHへの突入が視覚的・体感的に\n特別なイベントとして印象付けられる。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1200000)
    gy = Inches(0.10)

    positions = [
        (Inches(0.28),  Inches(0.72)),
        (Inches(5.17),  Inches(0.72)),
        (Inches(0.28),  Inches(0.72) + bh_p + gy),
        (Inches(5.17),  Inches(0.72) + bh_p + gy),
        (Inches(0.28),  Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1100000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(330000), pw - Emu(90000), ph - Emu(390000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "設計哲学の核心：「蓄積（刀ポイント）→自力感（刺客ゾーン）→バトル（対決）→ループ（絶頂輪廻）」の硬派循環設計",
           "押忍！番長3で確立した伝統のゲームフローを差枚管理型×ループ型ATで現代化した進化作")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN, "刀ポイント蓄積による能動的な通常遊技",
         "毎G役成立を蓄積するシステムが\nプレイヤーの「関与感」を高める。\n単なる待ち時間を\"積み上げる作業\"に変える設計。"),
        (C_GREEN, "絶頂輪廻ループの爽快な出玉連鎖",
         "巌流島→青頂ZBASHのループは\n視覚・体感ともに圧倒的な出玉体験。\n「止まらない感覚」が最大の差別化要素。"),
        (C_GREEN, "差枚管理型×倍化で逆転劇が可能",
         "倍斬刀チャレンジ最大64倍倍化で\n少ない投資から大量出玉への逆転劇。\nどのタイミングでも期待感を維持。"),
        (C_GREEN, "パチスロアワード2025ノミネートの話題性",
         "業界評価の高さがホール集客力に直結。\n番長ブランドの認知度×最新設計で\nライト層からベテランまで幅広く集客。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1140000)
        rect_b(s, lx, iy, lw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(300000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,   "低設定でのゲーム数解除の薄さ",
         "低設定ではCZ成功率が下がり\nゲーム数解除もほとんど期待できない。\n「CZ頼み」という声もあり、\n設定差の体感が大きくなりがち。"),
        (C_BLUE2, "モード管理の把握が必要",
         "モード別天井の理解がないと\n立ち回りが難しく感じる。\nライトユーザーへの敷居がやや高い。\n演出示唆の読み解きにも慣れが必要。"),
        (C_GOLD,  "絶頂輪廻ループ到達率の偏り",
         "ループに入れると爆発するが\n到達率は高くない。\n巌流島突入機会が少ない台は\n通常ATの消化を繰り返すだけになりがち。"),
        (C_GRAY,  "賛否両論の評価",
         "番長ファンには高評価だが\n新規ユーザーには取っつきにくいとの声も。\nシリーズ未経験者向けの\n入口設計のさらなる改善余地あり。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1140000)
        rect_b(s, rx, iy, rw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(300000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "強みと課題の両面把握が設計学習の本質：絶頂輪廻ループの爽快感を活かしつつ、ライト層への間口拡大が次の課題",
           "パチスロアワード2025ノミネートは「業界が評価する設計」の証。その設計要素を分解して活用する")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── いざ！番長の設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x15, 0x22, 0x55))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "いざ！番長 ── 設計的強み総括", 10, bold=True, color=C_GOLD)

    strengths = [
        (C_BLUE2, "刀ポイント蓄積フローという核心設計",
         "毎G役成立を蓄積→刺客ゾーン→対決という\n能動的な3ステップが「硬派な自力感」を演出。\n番長シリーズDNAを現代スマスロに移植した設計。"),
        (C_GOLD,  "絶頂輪廻ループという「夢と現実の循環」",
         "巌流島(+700枚)→青頂ZBASH(純増5.0枚)→\n巌流島→…の繰り返しが視覚的・体感的に\n「出玉が止まらない」圧倒的な爽快感を生む。"),
        (C_GOLD2, "差枚管理型×倍化×2段階純増の複合設計",
         "差枚管理で安心感、倍斬刀チャレンジで逆転劇、\n通常→上位ATの純増アップで到達感を演出。\n3要素の組み合わせが多様な出玉体験を提供。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1210000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(325000), lw - Emu(90000), Emu(780000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_BLUE2,  "蓄積型エンジンで毎Gに意味を持たせよ",
         "刀ポイントのように役成立を蓄積する\nシステムが通常時の単調さを解消する"),
        (C_GOLD,   "「夢の連鎖」は視覚化せよ",
         "絶頂輪廻ループのように繰り返しの\n美しさを見せることで離席を防ぐ"),
        (C_GOLD2,  "純増の段差で「到達感」を演出せよ",
         "2.8→5.0枚/Gの切り替わりが\n上位到達の明確な体感報酬になる"),
        (C_CYAN,   "シリーズの文法を守りつつ進化せよ",
         "番長の「硬派バトル感」というDNAを\n維持しながら差枚管理型に昇華した手法"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(760000)
        rect_b(s, rx, py0, rw, Emu(710000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(710000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(230000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(285000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3310000), rw, Emu(1060000),
           RGBColor(0x08, 0x10, 0x28), C_BLUE, 2.0)
    rect(s, rx, ry + Emu(3310000), Emu(40000), Emu(1060000), C_BLUE)
    tb(s, rx + Emu(65000), ry + Emu(3360000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_BLUE2)
    tb(s, rx + Emu(65000), ry + Emu(3620000), rw - Emu(85000), Emu(690000),
       "刀ポイント蓄積×刺客ゾーン自力感×対決番長バトル\n×絶頂輪廻ループという4要素の統合設計。\n番長シリーズのDNAを差枚管理型スマスロに\n昇華した2025年の代表的な設計完成機。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "本機の設計思想：「蓄積・自力感・番長バトル・絶頂輪廻ループ」── 番長シリーズ第3弾が示した設計進化の方向性",
           "パチスロアワード2025ノミネートは「業界が認めた現代スマスロ設計の模範」としての評価の証")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)      # 1: タイトル・スペック・3ポイント
    s_flow(prs)       # 2: ゲームフロー全体図
    s_normal(prs)     # 3: 通常時の遊び方
    s_cz(prs)         # 4: CZ/前兆の仕組み
    s_at(prs)         # 5: AT/ボーナス（出玉を伸ばす方法）
    s_upper(prs)      # 6: 上位ATへの道と遊び方
    s_design(prs)     # 7: 面白さの設計（番長シリーズの哲学）
    s_pros_cons(prs)  # 8: 良い点と課題
    s_matome(prs)     # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
