import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PATH = r"C:\Users\h.kadoya\Desktop\slocri\proposals\game_proposal_light_v3.pptx"

C_BG     = RGBColor(0x18, 0x18, 0x2E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xD8, 0x5A, 0x30)
C_GRAY   = RGBColor(0x88, 0x88, 0x99)
C_LTGRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_YELLOW = RGBColor(0xF5, 0xC5, 0x42)
C_GREEN  = RGBColor(0x4C, 0xB0, 0x7A)
C_RED    = RGBColor(0xE0, 0x50, 0x50)
C_BLUE   = RGBColor(0x55, 0x99, 0xDD)
C_PURPLE = RGBColor(0x99, 0x55, 0xDD)

W = Inches(10)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, size, bold=False, color=None,
       align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    if color: r.font.color.rgb = color
    return box


def bg(slide):
    rect(slide, 0, 0, W, H, C_BG)


def header(slide, title, accent=None):
    accent = accent or C_ORANGE
    rect(slide, 0, 0, W, Inches(1.0), RGBColor(0x10, 0x10, 0x22))
    rect(slide, 0, Inches(1.0), W, Emu(40000), accent)
    tb(slide, Inches(0.35), Inches(0.15), Inches(9.3), Inches(0.75),
       title, 20, bold=True, color=C_WHITE)


def footer(slide):
    tb(slide, Inches(0.3), Inches(7.12), Inches(9.4), Inches(0.28),
       "スロキー ／ ゲーム性提案書", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, color=None):
    rect(slide, l, t, w, h, color or RGBColor(0x22, 0x22, 0x3C))


def label(slide, l, t, text, color=None):
    tb(slide, l, t, Inches(2.5), Inches(0.32),
       text, 9, bold=True, color=color or C_GRAY)


def flow_arrow(slide, x, y):
    tb(slide, x, y, Inches(0.5), Inches(0.35), "▼", 14, color=C_GRAY,
       align=PP_ALIGN.CENTER)


# ============================================================
# Slide 1: タイトル
# ============================================================
s1 = prs.slides.add_slide(blank)
bg(s1)
rect(s1, 0, 0, W, Emu(55000), C_ORANGE)
tb(s1, Inches(0.6), Inches(1.6), Inches(8.8), Inches(1.1),
   "ゲーム性提案　2案", 40, bold=True, color=C_WHITE)
tb(s1, Inches(0.6), Inches(2.8), Inches(8.8), Inches(0.6),
   "20代ライトユーザー向け　新機種コンセプト", 17, color=C_ORANGE)
rect(s1, Inches(0.6), Inches(3.5), Inches(5.5), Emu(25000), C_GRAY)

items_s1 = [
    ("案①", "積み上げ型", "来るたびに有利になる台", C_GREEN),
    ("案③", "リスク選択型", "自分の判断で展開が変わる台", C_YELLOW),
]
for i, (num, name, desc, col) in enumerate(items_s1):
    y = Inches(3.8) + i * Inches(0.85)
    card(s1, Inches(0.6), y, Inches(8.8), Inches(0.75), RGBColor(0x1E, 0x1E, 0x35))
    tb(s1, Inches(0.8), y + Emu(12000), Inches(0.55), Inches(0.55),
       num, 12, bold=True, color=col)
    tb(s1, Inches(1.45), y + Emu(12000), Inches(1.6), Inches(0.55),
       name, 13, bold=True, color=C_WHITE)
    tb(s1, Inches(3.2), y + Emu(14000), Inches(6.0), Inches(0.45),
       desc, 11, color=C_LTGRAY)
footer(s1)

# ============================================================
# Slide 2: 案① 積み上げ型
# ============================================================
s2 = prs.slides.add_slide(blank)
bg(s2)
header(s2, "案①　積み上げ型　―　来るたびに有利になる台", C_GREEN)

# キャッチ
card(s2, Inches(0.3), Inches(1.15), Inches(9.4), Inches(0.72), C_GREEN)
tb(s2, Inches(0.55), Inches(1.22), Inches(9.0), Inches(0.58),
   "「前回の続きから始まる」設計で、来店のたびに天井が短くなる。長く付き合うほど得をする台。",
   12, bold=True, color=RGBColor(0x10, 0x10, 0x10))

# 左カラム：操作・ルール・ターゲット
LX = Inches(0.3)
card(s2, LX, Inches(1.98), Inches(4.55), Inches(4.85), RGBColor(0x1C, 0x1C, 0x30))

label(s2, LX + Emu(30000), Inches(2.05), "▌ プレイヤーは何をするの？", C_GREEN)
lines_op = [
    "通常時：レア役でSPが自動で貯まる",
    "　　　　（特別な操作は不要）",
    "AT中　：押し順ナビに従って消化",
    "　　　　セット終了ごとに継続演出を楽しむ",
]
y_op = Inches(2.42)
for line in lines_op:
    tb(s2, LX + Emu(30000), y_op, Inches(4.3), Inches(0.32), line, 10.5, color=C_LTGRAY)
    y_op += Inches(0.32)

label(s2, LX + Emu(30000), Inches(3.65), "▌ 全体ルール・何をすれば勝ち？", C_GREEN)
rules = [
    ("SP（ストックポイント）を貯める", "→ 天井Gが短縮される"),
    ("天井に到達してATに入る",         "→ 枚数を獲得"),
    ("SP引き継ぎで次回来店も有利",      "→ 来るほど楽になる"),
]
y_r = Inches(4.05)
for rule, result in rules:
    card(s2, LX + Emu(25000), y_r, Inches(4.38), Inches(0.45),
         RGBColor(0x18, 0x28, 0x18))
    tb(s2, LX + Emu(50000), y_r + Emu(8000), Inches(2.0), Inches(0.38),
       rule, 9.5, color=C_WHITE)
    tb(s2, LX + Emu(220000), y_r + Emu(8000), Inches(2.15), Inches(0.38),
       result, 9.5, bold=True, color=C_GREEN)
    y_r += Inches(0.52)

label(s2, LX + Emu(30000), Inches(5.7), "▌ ターゲット", C_GRAY)
tb(s2, LX + Emu(30000), Inches(6.05), Inches(4.3), Inches(0.55),
   "週1〜2回来店の20〜30代\n「投資した分だけ報われたい」層", 10.5, color=C_LTGRAY)

# 右カラム：SP蓄積ゲージ図
RX = Inches(5.1)
card(s2, RX, Inches(1.98), Inches(4.6), Inches(2.95), RGBColor(0x1C, 0x1C, 0x30))
label(s2, RX + Emu(30000), Inches(2.05), "▌ SP蓄積 → 天井短縮 イメージ", C_ORANGE)

bar_lx   = RX + Emu(400000)
bar_wmax = Inches(2.3)
bar_h    = Emu(220000)
visits = [
    ("初回",    0.00, "天井 600G", C_LTGRAY),
    ("2回目",   0.33, "天井 500G", C_LTGRAY),
    ("3回目",   0.67, "天井 400G", C_YELLOW),
    ("MAX来店", 1.00, "天井 300G", C_GREEN),
]
gy = Inches(2.52)
for vis_label, ratio, g_label, g_col in visits:
    tb(s2, RX + Emu(70000), gy, Emu(330000), bar_h + Emu(60000), vis_label, 8.5, color=C_GRAY)
    rect(s2, bar_lx, gy, bar_wmax, bar_h, RGBColor(0x2A, 0x2A, 0x45))
    if ratio > 0:
        rect(s2, bar_lx, gy, int(bar_wmax * ratio), bar_h, C_GREEN)
    tb(s2, bar_lx + bar_wmax + Emu(30000), gy, Emu(500000), bar_h + Emu(60000),
       g_label, 9, bold=(ratio == 1.0), color=g_col)
    if ratio < 1.0:
        tb(s2, bar_lx + bar_wmax // 2 - Emu(80000), gy + bar_h,
           Emu(250000), Emu(260000), "↓", 9, color=C_GRAY, align=PP_ALIGN.CENTER)
    gy += Inches(0.52)
rect(s2, RX + Emu(30000), gy + Emu(20000), Inches(4.38), Emu(290000),
     RGBColor(0x10, 0x28, 0x10))
tb(s2, RX + Emu(70000), gy + Emu(60000), Inches(4.1), Emu(240000),
   "★ AT後もSPを引き継ぎ、次回来店から有利スタート", 9, bold=True, color=C_GREEN)

# スペック
card(s2, RX, Inches(5.05), Inches(4.6), Inches(2.1), RGBColor(0x1A, 0x1A, 0x2C))
label(s2, RX + Emu(30000), Inches(5.12), "▌ スペック（目安）", C_ORANGE)
specs = [
    ("天井",     "SP0: 600G → SPMAX: 300G"),
    ("AT純増",   "3.8枚 / コイン単価 3.0円前後"),
    ("機械割",   "設定1: 97.5% / 設定6: 109%"),
    ("タイプ",   "スマスロ（SP記録・引き継ぎ必須）"),
    ("参考機種", "LモンキーターンV（SP引継）/ Lまどか☆マギカ4（壁引継）/ Lゴッドイーター3（レベル蓄積）"),
]
sy = Inches(5.5)
for k, v in specs:
    tb(s2, RX + Emu(30000), sy, Inches(1.1), Inches(0.28), k, 9, color=C_GRAY)
    tb(s2, RX + Emu(150000), sy, Inches(3.8), Inches(0.28), v, 9.5,
       bold=True, color=C_WHITE)
    sy += Inches(0.3)

footer(s2)

# ============================================================
# Slide 3: 案③ リスク選択型
# ============================================================
s3 = prs.slides.add_slide(blank)
bg(s3)
header(s3, "案③　リスク選択型　―　自分の判断で展開が変わる台", C_YELLOW)

card(s3, Inches(0.3), Inches(1.15), Inches(9.4), Inches(0.72),
     RGBColor(0x2C, 0x22, 0x10))
tb(s3, Inches(0.55), Inches(1.22), Inches(9.0), Inches(0.58),
   "ATの各セットで「攻める」か「守る」かを選ぶ。結果は自分の判断次第。勝っても負けても「自分が決めた」納得感がある。",
   12, bold=True, color=C_YELLOW)

LX3 = Inches(0.3)
card(s3, LX3, Inches(1.98), Inches(4.55), Inches(4.85), RGBColor(0x1C, 0x1C, 0x30))

label(s3, LX3 + Emu(30000), Inches(2.05), "▌ プレイヤーは何をするの？", C_YELLOW)
lines_op3 = [
    "通常時：押し順ナビに従って消化",
    "　　　　（スペシャルな操作なし）",
    "AT突入：各セット開始時に択",
    "　　　　「攻め」or「守り」を選ぶ",
    "AT中　：選んだルートの演出を楽しむ",
]
y_op3 = Inches(2.42)
for line in lines_op3:
    tb(s3, LX3 + Emu(30000), y_op3, Inches(4.3), Inches(0.32),
       line, 10.5, color=C_LTGRAY)
    y_op3 += Inches(0.32)

label(s3, LX3 + Emu(30000), Inches(3.85), "▌ 全体ルール・何をすれば勝ち？", C_YELLOW)

# 攻め / 守り の説明ボックス
for ci, (mode, lines, col, bgc) in enumerate([
    ("攻め", ["上乗せG数：大（+40〜100G）", "セット終了リスク：あり", "→ ハイリスク・ハイリターン"],
     C_ORANGE, RGBColor(0x28, 0x18, 0x10)),
    ("守り", ["上乗せG数：小（+10〜30G）", "セット継続：ほぼ確定", "→ 安定して枚数を積む"],
     C_BLUE, RGBColor(0x10, 0x18, 0x28)),
]):
    bx = LX3 + Emu(25000) + ci * Inches(2.22)
    card(s3, bx, Inches(4.22), Inches(2.12), Inches(1.85), bgc)
    tb(s3, bx + Emu(20000), Inches(4.3), Inches(1.95), Inches(0.35),
       mode, 13, bold=True, color=col)
    for j, ln in enumerate(lines):
        tb(s3, bx + Emu(20000), Inches(4.68) + j * Inches(0.32),
           Inches(2.0), Inches(0.32), ln, 9, color=C_LTGRAY)

label(s3, LX3 + Emu(30000), Inches(6.0), "▌ ターゲット", C_GRAY)
tb(s3, LX3 + Emu(30000), Inches(6.35), Inches(4.3), Inches(0.55),
   "「自分で決めたい・結果に納得したい」20〜30代\nゲームの判断要素が好きな層", 10.5, color=C_LTGRAY)

# 右カラム：分岐フロー図
RX3 = Inches(5.1)
card(s3, RX3, Inches(1.98), Inches(4.6), Inches(3.15), RGBColor(0x1C, 0x1C, 0x30))
label(s3, RX3 + Emu(30000), Inches(2.05), "▌ AT内 選択フロー", C_ORANGE)

CX3 = RX3 + Inches(2.3)  # 右カラム中心X

# AT突入ボックス
atw = Inches(1.5); ath = Emu(300000)
atx = CX3 - atw // 2
rect(s3, atx, Inches(2.22), atw, ath, RGBColor(0x18, 0x38, 0x18))
tb(s3, atx, Inches(2.27), atw, ath, "AT突入", 11, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

# ↓ 矢印
tb(s3, CX3 - Emu(100000), Inches(2.58), Emu(300000), Emu(230000), "↓", 10, color=C_GRAY, align=PP_ALIGN.CENTER)

# 択ボックス（オレンジ強調）
chw = Inches(2.1); chh = Emu(320000)
chx = CX3 - chw // 2
rect(s3, chx, Inches(2.78), chw, chh, C_ORANGE)
tb(s3, chx, Inches(2.83), chw, chh, "毎セット開始「択」", 11, bold=True,
   color=RGBColor(0x10, 0x10, 0x10), align=PP_ALIGN.CENTER)

# 分岐ラベル（↙ 攻め  守り ↘）
tb(s3, RX3 + Emu(80000), Inches(3.18), Inches(2.1), Emu(300000),
   "↙  攻め", 14, bold=True, color=C_ORANGE)
tb(s3, RX3 + Inches(2.5), Inches(3.18), Inches(1.9), Emu(300000),
   "守り  ↘", 14, bold=True, color=C_BLUE, align=PP_ALIGN.RIGHT)

# 攻め結果ボックス
abw = Inches(1.88); abh = Emu(720000)
abx = RX3 + Emu(55000)
rect(s3, abx, Inches(3.52), abw, abh, RGBColor(0x38, 0x18, 0x08))
tb(s3, abx + Emu(50000), Inches(3.57), abw - Emu(80000), Emu(270000),
   "成功: +40〜100G", 9.5, bold=True, color=C_GREEN)
tb(s3, abx + Emu(50000), Inches(3.88), abw - Emu(80000), Emu(270000),
   "失敗: セット終了", 9.5, bold=True, color=C_RED)
tb(s3, abx + Emu(50000), Inches(4.19), abw - Emu(80000), Emu(240000),
   "ハイリスク・ハイリターン", 8.5, color=C_LTGRAY)

# 守り結果ボックス
bbx = RX3 + Inches(2.62)
rect(s3, bbx, Inches(3.52), abw, abh, RGBColor(0x08, 0x18, 0x38))
tb(s3, bbx + Emu(50000), Inches(3.57), abw - Emu(80000), Emu(270000),
   "+10〜30G", 9.5, bold=True, color=C_GREEN)
tb(s3, bbx + Emu(50000), Inches(3.88), abw - Emu(80000), Emu(270000),
   "継続ほぼ確定", 9.5, bold=True, color=C_BLUE)
tb(s3, bbx + Emu(50000), Inches(4.19), abw - Emu(80000), Emu(240000),
   "ローリスク・ローリターン", 8.5, color=C_LTGRAY)

# スペック
card(s3, RX3, Inches(5.25), Inches(4.6), Inches(1.9), RGBColor(0x1A, 0x1A, 0x2C))
label(s3, RX3 + Emu(30000), Inches(5.32), "▌ スペック（目安）", C_ORANGE)
specs3 = [
    ("天井",     "設定1: 500G / 設定6: 250G"),
    ("AT純増",   "3.8枚 / コイン単価 3.0円前後"),
    ("機械割",   "設定1: 97.5% / 設定6: 109%"),
    ("タイプ",   "スマスロ（L型）"),
    ("参考機種", "L乃木坂46バイナリースター（択設計）/ Lバジリスク絆2（バトル択）/ L番長ZERO（AT内択）"),
]
sy3 = Inches(5.7)
for k, v in specs3:
    tb(s3, RX3 + Emu(30000), sy3, Inches(1.1), Inches(0.28), k, 9, color=C_GRAY)
    tb(s3, RX3 + Emu(150000), sy3, Inches(3.8), Inches(0.28), v, 9.5,
       bold=True, color=C_WHITE)
    sy3 += Inches(0.28)

footer(s3)

# ============================================================
# Slide 4: ②ブレスト候補
# ============================================================
s4 = prs.slides.add_slide(blank)
bg(s4)
header(s4, "案②　検討中のアイデア候補", C_PURPLE)

tb(s4, Inches(0.35), Inches(1.15), Inches(9.3), Inches(0.38),
   "「時間帯による選択」は実現困難。以下の方向で引き続き検討中。",
   11, color=C_GRAY)

ideas = [
    (
        "A　デイリーミッション型",
        "台側が「今日のミッション」を自動設定。プレイヤーは選ばない。",
        [
            "操作契機：　通常通りプレイするだけ",
            "全体ルール：今日のミッション（例:AT2回入る・○○枚獲得）を達成する",
            "勝ちの条件：ミッション達成 → スマスロに記録・称号ゲット",
            "メリット：　選択肢がないのでルールがシンプル",
            "課題：　　　「台が決めたミッション」に乗れない人には刺さらない",
        ],
        C_BLUE,
        "参考：FGO・モンスト等のデイリークエスト（パチスロでの先行事例はほぼなし）",
    ),
    (
        "B　周期明示型",
        "「次のチャンスまであと○G」を常時表示。来店ハードルを下げる設計。",
        [
            "操作契機：　通常通りプレイするだけ",
            "全体ルール：表示されたG数まで打てばチャンスタイム確定",
            "勝ちの条件：チャンスタイムでAT突入 → 枚数獲得",
            "メリット：　投資の上限が見えて20代の安心感につながる",
            "課題：　　　「見えすぎる」ことで設定読みが単純化するリスク",
        ],
        C_PURPLE,
        "参考：S 沖ドキ！（33G周期が見える）/ Lエウレカ3（周期ゾーン）/ Lカバネリ（チャージ周期）",
    ),
    (
        "C　コレクション型",
        "演出・フラグを集める楽しさ。「今日何が出たか」がスマスロに残る。",
        [
            "操作契機：　通常通りプレイするだけ",
            "全体ルール：レア演出を引くとスマスロの図鑑に登録される",
            "勝ちの条件：枚数獲得 ＋ 新しい演出の解放",
            "メリット：　SNS世代の「記録・共有」欲求に刺さる",
            "課題：　　　コレクションだけが目的化してゲーム性が薄れるリスク",
        ],
        C_GREEN,
        "参考：L 戦国コレクション（コレクション核）/ Lバジリスク絆2（スマスロ記録・図鑑）",
    ),
]

for i, (title, summary, lines, col, benchmark) in enumerate(ideas):
    y = Inches(1.65) + i * Inches(1.82)
    card(s4, Inches(0.3), y, Inches(9.4), Inches(1.72), RGBColor(0x1C, 0x1C, 0x32))
    rect(s4, Inches(0.3), y, Emu(30000), Inches(1.72), col)
    tb(s4, Inches(0.62), y + Emu(12000), Inches(3.5), Inches(0.38),
       title, 12, bold=True, color=col)
    tb(s4, Inches(0.62), y + Emu(105000), Inches(3.5), Inches(0.28),
       summary, 9.5, color=C_LTGRAY)
    tb(s4, Inches(0.62), y + Emu(900000), Inches(3.8), Inches(0.32),
       benchmark, 8.5, color=C_ORANGE)
    tx = Inches(4.6)
    for j, line in enumerate(lines):
        ty = y + Emu(15000) + j * Inches(0.3)
        tb(s4, tx, ty, Inches(5.0), Inches(0.28), line, 9,
           color=C_YELLOW if line.startswith("課題") else C_LTGRAY)

footer(s4)

# ============================================================
# Slide 5: 問いかけ
# ============================================================
s5 = prs.slides.add_slide(blank)
bg(s5)
rect(s5, 0, 0, W, Emu(55000), C_ORANGE)

tb(s5, Inches(0.5), Inches(0.72), Inches(9.0), Inches(0.6),
   "2案の論点と問いかけ", 20, bold=True, color=C_WHITE)

qs = [
    (
        "案①　積み上げ型",
        C_GREEN,
        [
            "Q1：SPの引き継ぎ上限はどう設定するか？（無限に有利になりすぎないか）",
            "Q2：スマスロ必須設計にすることへの懸念はあるか？",
            "Q3：「来るたびに有利」という訴求が設定依存に見えないか？",
        ]
    ),
    (
        "案③　リスク選択型",
        C_YELLOW,
        [
            "Q1：「攻め/守り」択はゲームバランス上どう担保するか？",
            "Q2：択がなくても成立するゲーム性に変えた方が良いか？",
            "Q3：ターゲット層（20代ライト）に「判断する楽しさ」は刺さるか？",
        ]
    ),
]

for i, (title, col, questions) in enumerate(qs):
    y = Inches(1.45) + i * Inches(2.55)
    card(s5, Inches(0.3), y, Inches(9.4), Inches(2.4), RGBColor(0x1C, 0x1C, 0x32))
    rect(s5, Inches(0.3), y, Emu(25000), Inches(2.4), col)
    tb(s5, Inches(0.65), y + Emu(12000), Inches(3.5), Inches(0.4),
       title, 13, bold=True, color=col)
    for j, q in enumerate(questions):
        tb(s5, Inches(0.65), y + Emu(130000) + j * Inches(0.55),
           Inches(8.8), Inches(0.5), q, 11, color=C_LTGRAY)

card(s5, Inches(0.3), Inches(6.65), Inches(9.4), Inches(0.62), C_ORANGE)
tb(s5, Inches(0.5), Inches(6.72), Inches(9.0), Inches(0.48),
   "②については A・B・C 案のどれかに絞るか、全く別の方向を探すか　ご意見お聞かせください。",
   11, bold=True, color=C_WHITE)

footer(s5)

prs.save(OUT_PATH)
print("保存完了:", OUT_PATH)
print("スライド数:", len(prs.slides))
