import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT_PATH = r"C:\Users\h.kadoya\Desktop\slocri\proposals\game_design_proposal_v4.pptx"

# カラー定数
C_BG     = RGBColor(0x18, 0x18, 0x2E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_ORANGE = RGBColor(0xD8, 0x5A, 0x30)
C_GRAY   = RGBColor(0x88, 0x88, 0x99)
C_LTGRAY = RGBColor(0xBB, 0xBB, 0xCC)
C_YELLOW = RGBColor(0xF5, 0xC5, 0x42)
C_GREEN  = RGBColor(0x4C, 0xB0, 0x7A)
C_RED    = RGBColor(0xE0, 0x50, 0x50)
C_BLUE   = RGBColor(0x55, 0x99, 0xDD)

W = Inches(10)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]


def add_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_tb(slide, l, t, w, h, text, size, bold=False, color=None,
           align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    if color:
        r.font.color.rgb = color
    return tb


def bg(slide):
    add_rect(slide, 0, 0, W, H, C_BG)


def title_bar(slide, title, sub=None):
    bar_h = Inches(1.05)
    add_rect(slide, 0, 0, W, bar_h, RGBColor(0x10, 0x10, 0x22))
    add_tb(slide, Inches(0.35), Inches(0.12), Inches(9.3), Inches(0.52),
           title, 19, bold=True, color=C_WHITE)
    if sub:
        add_tb(slide, Inches(0.35), Inches(0.65), Inches(9.3), Inches(0.35),
               sub, 10, color=C_GRAY)
    add_rect(slide, 0, bar_h, W, Emu(40000), C_ORANGE)


def card(slide, l, t, w, h, color=None):
    c = color or RGBColor(0x26, 0x26, 0x42)
    add_rect(slide, l, t, w, h, c)


def footer(slide, text="スロキー ／ ゲーム性提案書 v1"):
    add_tb(slide, Inches(0.3), Inches(7.1), Inches(9.4), Inches(0.3),
           text, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ==============================
# Slide 1: タイトル
# ==============================
s1 = prs.slides.add_slide(blank_layout)
bg(s1)
add_rect(s1, 0, 0, W, Emu(60000), C_ORANGE)
add_tb(s1, Inches(0.6), Inches(1.8), Inches(8.8), Inches(1.2),
       "スロット新機種　ゲーム性提案書", 36, bold=True, color=C_WHITE)
add_tb(s1, Inches(0.6), Inches(3.0), Inches(8.8), Inches(0.8),
       "173週・453機種の実稼働データが示す\n「長期稼働する台の設計法則」から逆算したゲーム性",
       16, color=C_ORANGE)
add_rect(s1, Inches(0.6), Inches(3.95), Inches(6.0), Emu(28000), C_GRAY)
add_tb(s1, Inches(0.6), Inches(4.08), Inches(8.8), Inches(0.5),
       "このゲーム性、どう思いますか？", 13, bold=True, color=C_LTGRAY)
footer(s1)

# ==============================
# Slide 2: 問題提起
# ==============================
s2 = prs.slides.add_slide(blank_layout)
bg(s2)
title_bar(s2, "問題意識：なぜ「3ヶ月後に誰も打っていない台」が生まれるのか")

card(s2, Inches(0.3), Inches(1.2), Inches(4.5), Inches(5.4), RGBColor(0x28, 0x18, 0x18))
add_tb(s2, Inches(0.45), Inches(1.32), Inches(4.2), Inches(0.42),
       "◆ よくあるパターン", 12, bold=True, color=C_RED)
fail_lines = [
    "初週：設定6が入り、爆発的稼働",
    "2〜4週：高評価口コミで来店急増",
    "5〜8週：ホールが設定を下げ始める",
    "9週〜：「あの台、最近ぜんぜん出ない」",
    "12週〜：稼働激減、他機種に流れる",
    "",
    "→ 設定次第の台だったから",
]
tb2l = s2.shapes.add_textbox(Inches(0.45), Inches(1.82), Inches(4.2), Inches(4.5))
tf2l = tb2l.text_frame; tf2l.word_wrap = True
for i, line in enumerate(fail_lines):
    p = tf2l.add_paragraph() if i > 0 else tf2l.paragraphs[0]
    p.space_before = Pt(5)
    r = p.add_run(); r.text = line
    r.font.size = Pt(11)
    r.font.bold = line.startswith("→")
    r.font.color.rgb = C_RED if line.startswith("→") else C_LTGRAY

card(s2, Inches(5.05), Inches(1.2), Inches(4.65), Inches(5.4), RGBColor(0x18, 0x28, 0x1E))
add_tb(s2, Inches(5.18), Inches(1.32), Inches(4.4), Inches(0.42),
       "◆ 長期稼働できた機種の共通点", 12, bold=True, color=C_GREEN)
good_items = [
    ("L東京喰種", "63週、稼働率182%を維持"),
    ("LモンキーターンV", "124週、稼働率144%を維持"),
    ("LディスクアップULTRAREMIX", "65週、稼働率146%を維持"),
    ("", ""),
    ("共通する設計の特徴", ""),
    ("", "「設定が低くてもゲームが面白い」"),
    ("", "AT演出の品質が設定に依存しない"),
    ("", "→ ホールが設定を下げても去らない"),
]
tb2r = s2.shapes.add_textbox(Inches(5.18), Inches(1.82), Inches(4.4), Inches(4.5))
tf2r = tb2r.text_frame; tf2r.word_wrap = True
for i, (nm, val) in enumerate(good_items):
    p = tf2r.add_paragraph() if i > 0 else tf2r.paragraphs[0]
    p.space_before = Pt(4)
    r = p.add_run()
    r.text = (nm + "　" + val) if nm and val else (nm or val)
    r.font.size = Pt(10.5)
    if nm and not nm.startswith("共通"):
        r.font.bold = True; r.font.color.rgb = C_YELLOW
    elif nm.startswith("共通"):
        r.font.bold = True; r.font.color.rgb = C_GREEN
    else:
        r.font.color.rgb = C_LTGRAY
add_tb(s2, Inches(4.55), Inches(3.3), Inches(0.9), Inches(0.5),
       "VS", 20, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
footer(s2)

# ==============================
# Slide 3: 用語解説（稼働率・後半維持率）
# ==============================
s_term = prs.slides.add_slide(blank_layout)
bg(s_term)
title_bar(s_term, "分析の見方：この資料で使う2つの指標")

# 稼働率カード（左）
card(s_term, Inches(0.3), Inches(1.2), Inches(4.65), Inches(5.5), RGBColor(0x1A, 0x1A, 0x32))
add_tb(s_term, Inches(0.45), Inches(1.28), Inches(4.3), Inches(0.45),
       "指標①　稼働率", 15, bold=True, color=C_ORANGE)
add_rect(s_term, Inches(0.45), Inches(1.78), Inches(4.2), Emu(22000), C_ORANGE)
add_tb(s_term, Inches(0.45), Inches(1.9), Inches(4.3), Inches(0.52),
       "＝  機種のアウト ÷ 全機種平均アウト × 100%", 12, bold=True, color=C_WHITE)

add_tb(s_term, Inches(0.45), Inches(2.55), Inches(4.3), Inches(0.38),
       "平均が100%。高いほど「よく打たれている台」", 10.5, color=C_LTGRAY)

examples_left = [
    ("100%", "平均と同じ稼働",             C_GRAY),
    ("150%", "平均の1.5倍（人気台）",      C_YELLOW),
    ("200%", "平均の2倍（超人気台）",      C_GREEN),
    (" 50%", "平均の半分以下（不振台）",   C_RED),
]
for i, (val, desc, col) in enumerate(examples_left):
    y_e = Inches(3.05) + i * Inches(0.62)
    card(s_term, Inches(0.45), y_e, Inches(4.25), Inches(0.55), RGBColor(0x20,0x20,0x3C))
    add_tb(s_term, Inches(0.55), y_e + Emu(12000), Inches(0.85), Inches(0.42),
           val, 15, bold=True, color=col)
    add_tb(s_term, Inches(1.45), y_e + Emu(14000), Inches(3.1), Inches(0.38),
           desc, 10.5, color=C_LTGRAY)

add_tb(s_term, Inches(0.45), Inches(5.62), Inches(4.3), Inches(0.55),
       "例）東京喰種の平均稼働率 = 182%\n→ 常に全体平均の約1.8倍打たれていた",
       10, color=C_YELLOW)

# 後半維持率カード（右）
card(s_term, Inches(5.15), Inches(1.2), Inches(4.55), Inches(5.5), RGBColor(0x1A, 0x1A, 0x32))
add_tb(s_term, Inches(5.3), Inches(1.28), Inches(4.2), Inches(0.45),
       "指標②　後半維持率", 15, bold=True, color=C_BLUE)
add_rect(s_term, Inches(5.3), Inches(1.78), Inches(4.1), Emu(22000), C_BLUE)
add_tb(s_term, Inches(5.3), Inches(1.9), Inches(4.2), Inches(0.52),
       "＝  W13以降の平均稼働率 ÷ W1〜12の平均稼働率", 12, bold=True, color=C_WHITE)

add_tb(s_term, Inches(5.3), Inches(2.55), Inches(4.2), Inches(0.52),
       "「人気が長続きしているか」を示す維持力\nW1〜12が初期12週、W13以降が長期", 10.5, color=C_LTGRAY)

examples_right = [
    ("90%以上", "ほぼ落ちていない（超長期型）", C_GREEN),
    ("65〜90%", "緩やかに落ちている（長期安定）", C_GREEN),
    ("45〜65%", "ある程度落ちた（中期維持）",    C_YELLOW),
    ("45%未満", "急速に落ちた（消費が早い）",    C_RED),
]
for i, (val, desc, col) in enumerate(examples_right):
    y_e = Inches(3.2) + i * Inches(0.65)
    card(s_term, Inches(5.3), y_e, Inches(4.15), Inches(0.58), RGBColor(0x20,0x20,0x3C))
    add_tb(s_term, Inches(5.4), y_e + Emu(12000), Inches(1.15), Inches(0.42),
           val, 13, bold=True, color=col)
    add_tb(s_term, Inches(6.6), y_e + Emu(14000), Inches(2.85), Inches(0.38),
           desc, 10, color=C_LTGRAY)

add_tb(s_term, Inches(5.3), Inches(5.62), Inches(4.2), Inches(0.55),
       "例）東京喰種の後半維持率 = 73.7%\n→ 長期でも初期の7割以上の稼働を維持",
       10, color=C_YELLOW)

# 下部まとめ
card(s_term, Inches(0.3), Inches(6.82), Inches(9.4), Inches(0.42), RGBColor(0x20,0x20,0x38))
add_tb(s_term, Inches(0.5), Inches(6.88), Inches(9.0), Inches(0.32),
       "この2指標を組み合わせることで「高くて長持ちする台」と「最初だけ高くてすぐ落ちた台」を区別できる",
       10.5, bold=True, color=C_LTGRAY)

footer(s_term)

# ==============================
# Slide 4: CVの法則
# ==============================
s3 = prs.slides.add_slide(blank_layout)
bg(s3)
title_bar(s3, "発見した法則：変動係数（CV）と長期維持率の関係",
          "n=130機種（ジャグラー系除く・30週以上）")

card(s3, Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.72), C_ORANGE)
add_tb(s3, Inches(0.5), Inches(1.28), Inches(9.0), Inches(0.55),
       "CV（週ごとアウトの変動係数）と後半維持率の相関係数  r = -0.913", 16,
       bold=True, color=C_WHITE)

add_tb(s3, Inches(0.35), Inches(2.08), Inches(9.3), Inches(0.38),
       "変動係数（CV）＝ 週ごとアウトの標準偏差 ÷ 平均  →  「設定依存度」の代理指標",
       11, color=C_LTGRAY)

rows3 = [
    ("CV値",       "意味",                  "後半維持率の傾向",  C_GRAY,   C_GRAY,   C_GRAY),
    ("CV < 0.2",   "設定に関わらず安定稼働", "70〜90%台",       C_GREEN,  C_WHITE,  C_GREEN),
    ("CV 0.3〜0.5","設定依存が中程度",       "40〜65%台",       C_YELLOW, C_WHITE,  C_YELLOW),
    ("CV > 0.6",   "高設定のときだけ爆発",   "20〜35%台",       C_RED,    C_WHITE,  C_RED),
]
col_w3 = [Inches(2.0), Inches(3.8), Inches(3.2)]
col_x3 = [Inches(0.3), Inches(2.35), Inches(6.2)]
row_bg3 = [RGBColor(0x10,0x10,0x22), RGBColor(0x1E,0x30,0x28),
           RGBColor(0x28,0x24,0x18), RGBColor(0x30,0x18,0x18)]
row_h3 = Inches(0.55)
for ri, (row, rbg) in enumerate(zip(rows3, row_bg3)):
    y3 = Inches(2.58) + ri * row_h3
    for ci, (cell, cw3, cx3) in enumerate(zip(row[:3], col_w3, col_x3)):
        add_rect(s3, cx3, y3, cw3 - Emu(20000), row_h3 - Emu(15000), rbg)
        fc3 = row[3 + ci] if ri > 0 else C_GRAY
        fb3 = ri == 0 or ci == 0
        add_tb(s3, cx3 + Emu(40000), y3 + Emu(10000),
               cw3 - Emu(80000), row_h3 - Emu(30000),
               cell, 11 if ri == 0 else 13, bold=fb3, color=fc3)

add_tb(s3, Inches(0.35), Inches(5.0), Inches(9.3), Inches(0.4),
       "★ r=-0.913 は「ほぼ完全な負の相関」。設定依存度を下げるほど長期稼働できる、と数値が示している。",
       10.5, bold=True, color=C_ORANGE)

card(s3, Inches(0.3), Inches(5.5), Inches(9.4), Inches(1.62), RGBColor(0x20,0x20,0x38))
add_tb(s3, Inches(0.45), Inches(5.58), Inches(9.0), Inches(0.3),
       "実機データ確認（2022-2026実稼働）　※稼働率100%=全体平均", 9, color=C_GRAY)
items3 = [
    ("東京喰種",    "稼働率182%", "維持率73.7%", C_GREEN),
    ("モンキーターンV", "稼働率144%", "維持率68.2%", C_GREEN),
    ("ディスクアップ", "稼働率146%", "維持率61.9%", C_GREEN),
    ("北斗の拳",    "稼働率257%→低下", "維持率42.1%", C_YELLOW),
    ("番長４",      "稼働率187%→低下", "維持率29.6%", C_RED),
]
for i, (nm, cv, rate, col3) in enumerate(items3):
    x3i = Inches(0.35) + i * Inches(1.88)
    card(s3, x3i, Inches(5.95), Inches(1.82), Inches(1.1), RGBColor(0x1A,0x1A,0x30))
    add_tb(s3, x3i + Emu(30000), Inches(5.99), Inches(1.75), Inches(0.32),
           nm, 8, bold=True, color=col3)
    add_tb(s3, x3i + Emu(30000), Inches(6.32), Inches(1.75), Inches(0.3),
           cv, 9, color=C_LTGRAY)
    add_tb(s3, x3i + Emu(30000), Inches(6.62), Inches(1.75), Inches(0.35),
           rate, 11, bold=True, color=col3)
footer(s3)

# ==============================
# Slide 4: 打ちたい気持ち vs やめにくさ
# ==============================
s4 = prs.slides.add_slide(blank_layout)
bg(s4)
title_bar(s4, "「打ちたい気持ち」と「やめにくさ」は別物だった")

for ci4, (ttl, lines, col4, bg4) in enumerate([
    ("やめにくさ", [
        "今のセッションを続けたい",
        "",
        "・天井まであと○G",
        "・ATが終わったらやめよう",
        "・残ゲーム数が勿体ない",
        "・損失回避の心理",
        "",
        "→ ホールを出れば来週は来ない",
    ], C_RED, RGBColor(0x28,0x18,0x18)),
    ("打ちたい気持ち", [
        "また来週あの台を打ちに来たい",
        "",
        "・あの演出がもう一度見たい",
        "・今日負けた、リベンジしたい",
        "・ボス撃破の達成感が忘れられない",
        "・次はもっと上手く打てる気がする",
        "",
        "→ 来週の来店動機が設計で生まれる",
    ], C_GREEN, RGBColor(0x18,0x28,0x1E)),
]):
    x4 = Inches(0.3 + ci4 * 4.9)
    card(s4, x4, Inches(1.2), Inches(4.65), Inches(5.4), bg4)
    add_tb(s4, x4 + Emu(40000), Inches(1.3), Inches(4.4), Inches(0.45),
           ttl, 16, bold=True, color=col4)
    add_rect(s4, x4 + Emu(40000), Inches(1.8), Inches(4.0), Emu(20000), col4)
    tb4 = s4.shapes.add_textbox(x4 + Emu(40000), Inches(1.9), Inches(4.4), Inches(4.4))
    tf4 = tb4.text_frame; tf4.word_wrap = True
    for j, line in enumerate(lines):
        p4 = tf4.add_paragraph() if j > 0 else tf4.paragraphs[0]
        p4.space_before = Pt(2)
        r4 = p4.add_run(); r4.text = line
        r4.font.size = Pt(11.5 if line.startswith("→") else 11)
        r4.font.bold = line.startswith("→")
        r4.font.color.rgb = col4 if line.startswith("→") else C_LTGRAY

add_tb(s4, Inches(4.55), Inches(3.2), Inches(0.9), Inches(0.5),
       "VS", 20, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)
card(s4, Inches(0.3), Inches(6.72), Inches(9.4), Inches(0.55), C_ORANGE)
add_tb(s4, Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.42),
       "長期稼働するのは「打ちたい気持ち」を生む機種。その源泉は「設定に依存しないゲーム体験の完成度」。",
       11, bold=True, color=C_WHITE)
footer(s4)

# ==============================
# Slide 5: コンセプト
# ==============================
s5 = prs.slides.add_slide(blank_layout)
bg(s5)
title_bar(s5, "提案するゲーム性コンセプト：進化するバトル型AT")

principles5 = [
    ("原則①", "AT演出・バトルは設定非依存で面白い",
     "低設定でも体験の質が変わらない設計\n→ CV目標 0.20〜0.25"),
    ("原則②", "設定差はAT突入頻度に集中",
     "設定が高いほど当たりやすい\nでも当たったら全設定で同じ体験"),
    ("原則③", "毎セッションに「達成感」が生まれる",
     "ボス撃破・上乗せの喜び\n敗北→リベンジ欲求のサイクル"),
]
for i5, (num5, ttl5, body5) in enumerate(principles5):
    x5 = Inches(0.3 + i5 * 3.2)
    card(s5, x5, Inches(1.2), Inches(3.1), Inches(3.8), RGBColor(0x20,0x20,0x3A))
    add_tb(s5, x5+Emu(30000), Inches(1.3), Inches(2.9), Inches(0.4),
           num5, 11, bold=True, color=C_ORANGE)
    add_tb(s5, x5+Emu(30000), Inches(1.75), Inches(2.9), Inches(0.62),
           ttl5, 12, bold=True, color=C_WHITE)
    add_rect(s5, x5+Emu(30000), Inches(2.42), Inches(2.7), Emu(18000), C_ORANGE)
    add_tb(s5, x5+Emu(30000), Inches(2.52), Inches(2.9), Inches(2.4),
           body5, 10.5, color=C_LTGRAY)

add_tb(s5, Inches(0.35), Inches(5.15), Inches(9.3), Inches(0.45),
       "▶ 設計根拠：東京喰種(赫眼・喰種対決)・モンキーターンV(競艇レース映像)・ディスクアップ(技術介入)",
       10.5, bold=True, color=C_ORANGE)
add_tb(s5, Inches(0.35), Inches(5.6), Inches(9.3), Inches(0.38),
       "いずれも「設定1でも体験できる盛り上がり」が核心。これがCV低下の理由。",
       10.5, color=C_LTGRAY)

card(s5, Inches(0.3), Inches(6.12), Inches(9.4), Inches(1.02), RGBColor(0x18,0x18,0x2C))
specs5 = [
    ("タイプ", "スマスロ（L型）"),
    ("単価",   "2.8〜3.2円"),
    ("機械割", "設定1: 97.5% / 設定6: 109%"),
    ("純増",   "AT中 3.8枚/G"),
    ("天井",   "設定1: 700G / 設定6: 350G"),
]
for i5, (k5, v5) in enumerate(specs5):
    x5s = Inches(0.35 + i5 * 1.88)
    add_tb(s5, x5s, Inches(6.18), Inches(1.85), Inches(0.3), k5, 8, color=C_GRAY)
    add_tb(s5, x5s, Inches(6.5), Inches(1.85), Inches(0.55), v5, 9.5,
           bold=True, color=C_WHITE)
footer(s5)

# ==============================
# Slide 6: ゲームフロー
# ==============================
s6 = prs.slides.add_slide(blank_layout)
bg(s6)
title_bar(s6, "ゲームフロー：BURST CHAIN")

flow = (
    "【通常時：挑戦ゾーン】\n"
    "  50〜300G周期で CZ「激突バトル」（30G）へ\n"
    "  失敗しても「バトルポイント（BP）」として蓄積 → 天井CZへ持ち越し\n"
    "  ※ 天井は設定差あり（設定6：350G  /  設定1：700G）\n"
    "\n"
    "          ↓ CZ突破\n"
    "\n"
    "【AT：BURST CHAIN】  純増 3.8枚 / 1セット 40G\n"
    "  ┌─ 毎G：小役でポイント蓄積 → 50ptで「上乗せ抽選」（設定差なし・自力感）\n"
    "  ├─ 10G毎：バトルフェーズ（全3段階）\n"
    "  │    ① ザコ戦（勝率80%以上・チュートリアル感・毎回の達成感）\n"
    "  │    ② 中ボス戦（勝率60〜70%・緊張感が高まる）\n"
    "  │    ③ ボス戦（勝率40〜50%・クライマックス・完全自力）\n"
    "  │    ★ 勝敗は「小役成立パターン」で決定 → デキレ感ゼロ\n"
    "  │    ★ バトル勝率は設定に影響されない（どの設定でも同じ体験）\n"
    "  ├─ ボス撃破 → セット上乗せ（+20〜60G）＆「次も来たい」演出\n"
    "  └─ ボス敗北 → セット終了、引き戻し抽選 →「悔しい…次こそ」\n"
    "\n"
    "          ↓ ボス連続撃破 x 3回\n"
    "\n"
    "【上位AT：CHAIN BURST MAX】  純増 5.2枚・上乗せフリー状態\n"
    "  IP演出最高潮 / 全設定で同じ確率で突入可能（引き次第）"
)
card(s6, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.9), RGBColor(0x1A,0x1A,0x2E))
tb6 = s6.shapes.add_textbox(Inches(0.52), Inches(1.32), Inches(9.05), Inches(5.7))
tf6 = tb6.text_frame; tf6.word_wrap = True
p6 = tf6.paragraphs[0]
r6 = p6.add_run(); r6.text = flow
r6.font.size = Pt(10.5); r6.font.color.rgb = C_LTGRAY
r6.font.name = "MS Gothic"
footer(s6)

# ==============================
# Slide 7: 設計の核心（なぜ低CVか）
# ==============================
s7 = prs.slides.add_slide(blank_layout)
bg(s7)
title_bar(s7, "設計の核心：なぜこれが「低CV」になるのか")

rows7 = [
    ("設計要素",       "設定との関係",                  "プレイヤーが感じること"),
    ("バトル勝率",     "設定は影響しない\n（小役確率で決まる）", "「自分が引いて勝った」\n達成感が毎回残る"),
    ("AT中の演出品質", "設定差なし\n（全設定で同じ映像・音楽）", "設定1でもIP体験が\n満足できる"),
    ("上乗せ抽選",     "設定はほぼ影響しない\n（小役からの自力感）", "「今日は引きが良かった」\n個人体験として完結"),
    ("AT突入頻度",     "設定差あり\n（天井G数・CZ確率に差）", "「今日は高設定か？」\n来店・設定狙いの動機"),
]
col_w7 = [Inches(2.5), Inches(3.3), Inches(3.5)]
col_x7 = [Inches(0.3), Inches(2.85), Inches(6.2)]
row_h7 = Inches(1.05)
row_bg7 = [
    [RGBColor(0x10,0x10,0x22)] * 3,
    [RGBColor(0x22,0x22,0x3A), RGBColor(0x18,0x28,0x18), RGBColor(0x1E,0x1E,0x32)],
    [RGBColor(0x22,0x22,0x3A), RGBColor(0x18,0x28,0x18), RGBColor(0x1E,0x1E,0x32)],
    [RGBColor(0x22,0x22,0x3A), RGBColor(0x18,0x28,0x18), RGBColor(0x1E,0x1E,0x32)],
    [RGBColor(0x22,0x22,0x3A), RGBColor(0x28,0x18,0x18), RGBColor(0x1E,0x1E,0x32)],
]
row_fc7 = [
    [C_GRAY, C_GRAY, C_GRAY],
    [C_WHITE, C_GREEN, C_YELLOW],
    [C_WHITE, C_GREEN, C_YELLOW],
    [C_WHITE, C_GREEN, C_YELLOW],
    [C_WHITE, C_RED,   C_YELLOW],
]
for ri7, row7 in enumerate(rows7):
    y7 = Inches(1.25) + ri7 * row_h7
    for ci7, (cell7, cw7, cx7) in enumerate(zip(row7, col_w7, col_x7)):
        add_rect(s7, cx7, y7, cw7 - Emu(15000), row_h7 - Emu(15000),
                 row_bg7[ri7][ci7])
        add_tb(s7, cx7+Emu(30000), y7+Emu(15000),
               cw7-Emu(75000), row_h7-Emu(38000),
               cell7, 9.5 if ri7 == 0 else 10.5,
               bold=(ri7 == 0 or ci7 == 0),
               color=row_fc7[ri7][ci7])

card(s7, Inches(0.3), Inches(6.55), Inches(9.4), Inches(0.72), C_ORANGE)
add_tb(s7, Inches(0.5), Inches(6.62), Inches(9.0), Inches(0.58),
       "設定差をAT頻度のみに集中させることで、「設定6は当たりやすい」「でも当たったらどの設定も同じ体験」を両立。\n"
       "→ 打ち手全員に「打って良かった」体験が残り、来週の来店動機が生まれる。",
       9.5, bold=True, color=C_WHITE)
footer(s7)

# ==============================
# Slide 8: 近似機種比較
# ==============================
s8 = prs.slides.add_slide(blank_layout)
bg(s8)
title_bar(s8, "近似機種との設計比較")

tbl8_h = ["比較軸", "東京喰種", "モンキーターンV", "番長４", "本提案（目標）"]
tbl8_d = [
    ["設定依存度（CV）", "0.18（低）★",  "0.23（低）★", "0.68（高）",   "0.20〜0.25★"],
    ["AT中の設定差",    "ほぼなし",      "ほぼなし",     "大きい",        "ほぼなし"],
    ["自力感の強さ",    "◎（赫眼）",    "◎（V入賞）",  "△（押し順）",  "◎（バトル）"],
    ["後半維持率",      "73.7%★",       "68.2%★",      "29.6%",         "目標65%以上★"],
    ["リベンジ欲求",    "高（赫眼外れ）", "高（V外れ）", "低",            "高（ボス敗北）"],
    ["ターゲット",      "幅広",          "設定狙い〜幅広", "ベテラン寄り", "幅広"],
]
col_w8 = [Inches(1.8), Inches(1.7), Inches(1.85), Inches(1.65), Inches(2.05)]
col_x8 = [Inches(0.25), Inches(2.08), Inches(3.82), Inches(5.71), Inches(7.4)]
row_h8 = Inches(0.82)

for ri8, row8 in enumerate([tbl8_h] + tbl8_d):
    y8 = Inches(1.22) + ri8 * row_h8
    for ci8, (cell8, cw8, cx8) in enumerate(zip(row8, col_w8, col_x8)):
        is_h8 = ri8 == 0
        is_last8 = ci8 == 4
        is_star8 = "★" in cell8
        if is_h8:
            rc8 = RGBColor(0x10,0x10,0x20); fc8 = C_GRAY; fs8 = 9; fb8 = True
        elif ci8 == 0:
            rc8 = RGBColor(0x22,0x22,0x38); fc8 = C_LTGRAY; fs8 = 10; fb8 = True
        elif is_last8 and is_star8:
            rc8 = RGBColor(0x28,0x1A,0x08); fc8 = C_ORANGE; fs8 = 10; fb8 = True
        elif is_last8:
            rc8 = RGBColor(0x22,0x18,0x10); fc8 = C_YELLOW; fs8 = 10; fb8 = False
        elif is_star8:
            rc8 = RGBColor(0x18,0x28,0x18); fc8 = C_GREEN; fs8 = 10; fb8 = True
        else:
            rc8 = RGBColor(0x1E,0x1E,0x30); fc8 = C_LTGRAY; fs8 = 10; fb8 = False
        add_rect(s8, cx8, y8, cw8 - Emu(12000), row_h8 - Emu(12000), rc8)
        add_tb(s8, cx8+Emu(25000), y8+Emu(12000),
               cw8-Emu(52000), row_h8-Emu(30000),
               cell8, fs8, bold=fb8, color=fc8)
footer(s8)

# ==============================
# Slide 9: 論点
# ==============================
s9 = prs.slides.add_slide(blank_layout)
bg(s9)
title_bar(s9, "論点と議論したい点")

qs9 = [
    ("Q1", "AT設定非依存への懸念",
     "設定6の魅力をどう伝えるか？\n→ 「設定6は当たりやすい＝バトルを何回も楽しめる」という価値提案でどうか"),
    ("Q2", "ボス勝率40〜50%の適切さ",
     "勝率が高すぎると達成感が薄い。低すぎると離脱する。\n→ 段階構成（①80% / ②65% / ③45%）でどうか"),
    ("Q3", "IPとの親和性",
     "バトル系IPとの相性が良い（格闘・スポーツ・ダンジョン系）\n→ どのIPが当てはまるか"),
    ("Q4", "セット完結型 vs ループ型",
     "本提案はセット完結型（40G×複数）を採用\n→ ループ型の方が「打ちたい気持ち」に有効か？"),
]
for i9, (q9, ttl9, body9) in enumerate(qs9):
    row9 = i9 // 2; col9 = i9 % 2
    x9 = Inches(0.3 + col9 * 4.85)
    y9 = Inches(1.2 + row9 * 2.68)
    card(s9, x9, y9, Inches(4.7), Inches(2.55), RGBColor(0x1C,0x1C,0x32))
    add_tb(s9, x9+Emu(30000), y9+Emu(18000), Inches(0.65), Inches(0.38),
           q9, 13, bold=True, color=C_ORANGE)
    add_tb(s9, x9+Emu(120000), y9+Emu(18000), Inches(3.9), Inches(0.42),
           ttl9, 12, bold=True, color=C_WHITE)
    add_rect(s9, x9+Emu(30000), y9+Emu(160000), Inches(4.3), Emu(18000), C_GRAY)
    add_tb(s9, x9+Emu(30000), y9+Emu(195000), Inches(4.4), Inches(1.85),
           body9, 10.5, color=C_LTGRAY)
footer(s9)

# ==============================
# Slide 10: まとめ・問いかけ
# ==============================
s10 = prs.slides.add_slide(blank_layout)
bg(s10)
add_rect(s10, 0, 0, W, Emu(60000), C_ORANGE)
add_tb(s10, Inches(0.5), Inches(0.72), Inches(9.0), Inches(0.62),
       "まとめ：この提案で実現したいこと", 20, bold=True, color=C_WHITE)

checks = [
    "実稼働データが証明した「低CV設計」を意図的に実現する",
    "AT中の演出・バトルを設定非依存にすることで「どの設定でも面白い台」にする",
    "初週の爆発ではなく「3ヶ月後も稼働している台」を作る",
    "「また来週打ちに行こう」と思わせる「打ちたい気持ち」を設計で生み出す",
]
for i10, txt10 in enumerate(checks):
    y10 = Inches(1.5) + i10 * Inches(0.72)
    card(s10, Inches(0.35), y10, Inches(9.3), Inches(0.65), RGBColor(0x1E,0x1E,0x35))
    add_tb(s10, Inches(0.5), y10+Emu(14000), Inches(0.4), Inches(0.48),
           "✅", 14, color=C_GREEN)
    add_tb(s10, Inches(1.0), y10+Emu(14000), Inches(8.5), Inches(0.48),
           txt10, 12, color=C_WHITE)

add_rect(s10, 0, Inches(4.95), W, Emu(30000), C_ORANGE)
card(s10, Inches(0.3), Inches(5.1), Inches(9.4), Inches(2.08), C_ORANGE)
add_tb(s10, Inches(0.5), Inches(5.22), Inches(9.0), Inches(1.1),
       "このゲーム性、どう思いますか？",
       34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_tb(s10, Inches(0.5), Inches(6.38), Inches(9.0), Inches(0.52),
       "気になる点・詰めたい箇所・別案など、何でもお聞かせください。"
       "データを根拠にしているので、どの設計判断も説明できます。",
       10.5, color=RGBColor(0xFF, 0xDD, 0xCC), align=PP_ALIGN.CENTER)

add_tb(s10, Inches(0.3), Inches(7.1), Inches(9.4), Inches(0.3),
       "分析対象：実稼働データ 453機種・173週（2022〜2026）"
       "　　分析手法：変動係数（CV）× 後半維持率 相関分析 r=-0.913（n=130）",
       7.5, color=RGBColor(0xFF, 0xCC, 0xAA), align=PP_ALIGN.RIGHT)

prs.save(OUT_PATH)
print("保存完了:", OUT_PATH)
print("スライド数:", len(prs.slides))
