"""
スマスロ ヨルムンガンド 機種説明＋分析 統合版資料 v2 （山佐・2026年4月6日導入）
出力: proposals/機種分析/ヨルムンガンド/yormungand_guide_v2.pptx
テーマ: 深緑黒 × ミリタリー（説明パート＋分析パートの統合版）
※ WebSearch解析情報を反映した正確版

スライド構成:
  Part A 説明パート（プレイヤー視点）
    1. タイトル・スペック・この台の3ポイント
    2. ゲームフロー全体図（通常→CZ→AT→POへの全ルートを蛇行2段で可視化）
    3. 通常時の遊び方（CZ当選ルート3種・天井）
    4. ストーリーCZの仕組み（プロローグ10G＋ジャッジパート3G・デキレ感の構造）
    5. AT「ヨルムンガンドラッシュ」（PO到達ルート・AT中の行動指針）
    6. PO「パーフェクトオーダー」（純増5枚の体験・恥の世紀ループ詳細）
  Part B 分析パート
    7. 面白さの設計（POポテンシャルの価値）
    8. 不評の構造と教訓（通常時渋さ・有利区間・デキレ感）
    9. まとめ・設計から学べること
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "ヨルムンガンド", "yormungand_guide_v2.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深緑×ミリタリー）───────────────────────────────
C_BG    = RGBColor(0x04, 0x10, 0x08)   # 深緑黒
C_CARD  = RGBColor(0x08, 0x18, 0x0C)
C_CARD2 = RGBColor(0x0C, 0x20, 0x10)
C_ROW   = RGBColor(0x0A, 0x1C, 0x0E)
C_GREEN = RGBColor(0x22, 0xAA, 0x44)   # グリーン（メインカラー）
C_GREEN2= RGBColor(0x44, 0xDD, 0x66)   # 明るいグリーン
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)   # シアン（PO色）
C_YEL   = RGBColor(0xCC, 0xAA, 0x22)   # 黄土（恥の世紀）
C_RED   = RGBColor(0xCC, 0x22, 0x22)   # 赤（不評・問題）
C_CRIM  = RGBColor(0xFF, 0x44, 0x44)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)
C_DARK  = RGBColor(0x02, 0x08, 0x04)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（ミリタリー深緑）──────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 16, 8))
    draw = ImageDraw.Draw(img)
    # 斜めライン（ミリタリーグリッド）
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(6, 20, 10), width=1)
    # 下部グリーングロー
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        draw.line([(0, y), (w, y)], fill=(0, int(25 * t), int(10 * t)))
    # 上部薄暗化
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, int(8 * t), 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_GREEN)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GREEN2, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_GREEN)


def net_note(slide):
    """右下に「※ネット解析情報より」を表示"""
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, design_comment, sub_comment=""):
    """各スライド下部にフッター（設計コメント＋補足）を表示"""
    rect(slide, 0, Inches(5.1), SLIDE_W, Inches(0.525), C_DARK)
    rect(slide, 0, Inches(5.1), Emu(12000), Inches(0.525), C_GREEN)
    tb(slide, Inches(0.18), Inches(5.13), Inches(6.5), Emu(310000),
       design_comment, 8, bold=True, color=C_GREEN2, font=FONT_B)
    if sub_comment:
        tb(slide, Inches(6.7), Inches(5.13), Inches(3.1), Emu(310000),
           sub_comment, 7.5, color=C_GRAY, align=PP_ALIGN.RIGHT, font=FONT_B)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GREEN
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    """下向き矢印"""
    shp = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(200000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GREEN
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル（タイトル）
    rect(s, 0, 0, Inches(5.3), SLIDE_H, C_DARK)
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_GREEN)
    rect(s, Inches(5.3), 0, Emu(9000), SLIDE_H, RGBColor(0x10, 0x60, 0x20))

    tb(s, Inches(0.22), Inches(0.35), Inches(5.0), Emu(330000),
       "スマスロ 機種説明＋分析 統合版", 10, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.85), Inches(5.1), Emu(900000),
       "スマスロ\nヨルムンガンド", 30, bold=True, color=C_GREEN2, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.75), Inches(5.0), Emu(280000),
       "山佐（YAMASA）　2026年4月6日導入", 9.5, color=C_CREAM, font=FONT_H)

    # スペック一覧（WebSearch解析情報より）
    specs = [
        ("設定",     "1〜6段階"),
        ("AT名",     "ヨルムンガンドラッシュ"),
        ("AT純増",   "約2.4枚/G  初期50G"),
        ("PO名",     "パーフェクトオーダー（PO）"),
        ("PO純増",   "約5.0枚/G  初期50G"),
        ("天井",     "AT間 最大999G+α"),
    ]
    for i, (k, v) in enumerate(specs):
        sy = Inches(3.15) + i * Emu(330000)
        col = C_CYAN if "PO" in k else C_GREEN
        tb(s, Inches(0.25), sy, Inches(1.5), Emu(290000),
           k, 8, bold=True, color=col, font=FONT_B)
        tb(s, Inches(1.75), sy, Inches(3.3), Emu(290000),
           v, 8, color=C_WHITE, font=FONT_B)

    # 右：この台の3ポイント
    kws = [
        (C_GREEN,  "ポイント①\nストーリーCZ経由でATへ",
         "通常時→ストーリーCZ突破で\nヨルムンガンドラッシュに突入する台"),
        (C_CYAN,   "ポイント②\nPOが純増5.0枚の高性能AT",
         "PO「パーフェクトオーダー」は\n純増5枚・恥の世紀でループするトップスペック"),
        (C_RED,    "ポイント③\n不評の実態（客観的に）",
         "通常時の体感の渋さと\n有利区間上限が不満として噴出した機種"),
    ]
    # 右3ボックス: y0=0.42+i*1.55, 高さ1.38 → 末端=0.42+2×1.55+1.38=4.9 OK
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.42) + i * Emu(1417200)  # 1.55 inch = 1417200 EMU
        bh_kw = Inches(1.38)
        rect_b(s, Inches(5.56), y0, Inches(4.2), bh_kw, C_CARD, ac, 2.0)
        rect(s, Inches(5.56), y0, Emu(55000), bh_kw, ac)
        tb(s, Inches(5.76), y0 + Emu(65000), Inches(3.8), Emu(340000),
           kw, 11, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.76), y0 + Emu(415000), Inches(3.8), Emu(470000),
           desc, 8.5, color=C_WHITE)

    footer(s, "設計コメント：高性能POを持ちながら「届かない台」として評価が割れた機種。",
           "説明はフラット・客観的に記述")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow_overview(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時 → CZ（3種） → AT → PO の全ルート", "2/9")

    # ── 上段（左→右）: 通常時 → ストーリーCZ → AT ─────────────────────
    bw1 = Inches(2.2)
    bh1 = Emu(1250000)
    gap1 = Inches(0.38)
    row1_y = Inches(0.72)
    row1_cx = row1_y + bh1 // 2

    upper_boxes = [
        (C_LTGRY,                     C_GRAY,   "通常時",
         "レア役・CZポイントで\nCZ（3種）を目指す"),
        (C_CARD2,                     C_GREEN,  "CZ（3種類）",
         "①シューティングゾーン\n②滅びの丘\n③パーフェクトチャレンジ"),
        (RGBColor(0x08, 0x20, 0x10),  C_GREEN2, "ヨルムンガンドラッシュ",
         "基本AT\n純増2.4枚/G・初期50G"),
    ]
    box_xs_u = []
    for i, (fill, bc, lbl, sub) in enumerate(upper_boxes):
        bx0 = Inches(0.35) + i * (bw1 + gap1)
        box_xs_u.append(bx0)
        rect_b(s, bx0, row1_y, bw1, bh1, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row1_y + Emu(90000), bw1 - Emu(80000), Emu(380000),
           lbl, 10, bold=True, color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row1_y + Emu(490000), bw1 - Emu(60000), Emu(600000),
           sub, 8, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            arrow_r(s, bx0 + bw1 + Emu(20000), row1_cx)

    # AT後の右端から下へ折り返し矢印（テキスト付き）
    at_bx = box_xs_u[2]
    at_right = at_bx + bw1
    rect(s, at_right + Emu(40000), row1_y + bh1 - Emu(80000),
         Inches(1.45), Emu(140000), RGBColor(0x04, 0x18, 0x20))
    tb(s, at_right + Emu(60000), row1_y + bh1 - Emu(70000),
       Inches(1.35), Emu(110000),
       "エンディングで恥の世紀へ", 7, color=C_CYAN, align=PP_ALIGN.CENTER)
    # 右端縦線（折り返しの曲がり表現）
    rect(s, at_right + Emu(40000) + Inches(1.45) // 2 - Emu(10000),
         row1_y + bh1 - Emu(80000), Emu(20000), bh1 // 2 + Emu(200000),
         C_CYAN)

    # ── 下段（右→左）: AT → PO → 恥の世紀ループ ──────────────────────
    bw2 = Inches(2.2)
    bh2 = Emu(1250000)
    gap2 = Inches(0.38)
    row2_y = Inches(2.22)
    row2_cx = row2_y + bh2 // 2

    lower_boxes = [
        (RGBColor(0x04, 0x18, 0x20),  C_CYAN,  "PO（パーフェクトオーダー）",
         "純増5.0枚/G\n初期50G・BB以上確定"),
        (RGBColor(0x10, 0x14, 0x04),  C_YEL,   "恥の世紀",
         "期待度約50%\nPOループCZ"),
        (RGBColor(0x18, 0x04, 0x04),  C_RED,   "有利区間上限\n（エンディング）",
         "枚数上限到達で\n強制終了"),
    ]
    # 右から左に並べる（右端 → 左）
    box_xs_l = []
    for i, (fill, bc, lbl, sub) in enumerate(lower_boxes):
        bx0 = Inches(9.65) - (i + 1) * bw2 - i * gap2
        box_xs_l.append(bx0)
        rect_b(s, bx0, row2_y, bw2, bh2, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row2_y + Emu(90000), bw2 - Emu(80000), Emu(380000),
           lbl, 10, bold=True, color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row2_y + Emu(490000), bw2 - Emu(60000), Emu(600000),
           sub, 8, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            # 左向き矢印（右から左なので逆方向）
            arr_x = bx0 - gap2 - Emu(30000)
            shp = s.shapes.add_shape(12, arr_x, row2_cx - Emu(90000),
                                     Emu(200000), Emu(180000))
            shp.fill.solid()
            shp.fill.fore_color.rgb = bc
            shp.line.fill.background()

    # ── AT→PO ルートラベル ───────────────────────────────────────────
    # 上段ATボックス下部からPOへの接続ラベル（右側の縦ライン）
    tb(s, Inches(8.0), row1_y + bh1 - Emu(60000), Inches(1.7), Emu(140000),
       "BB以上+フリーズ→PO確定", 7.5, bold=True, color=C_CYAN, align=PP_ALIGN.RIGHT)

    # ── 下段CZ非突破→リトライのループラベル ─────────────────────────
    rect(s, Inches(0.35), row2_y, Inches(1.5), bh2,
         RGBColor(0x0A, 0x14, 0x0C))
    rect_b(s, Inches(0.35), row2_y, Inches(1.5), bh2,
           RGBColor(0x0A, 0x14, 0x0C), C_GRAY, 1.0)
    tb(s, Inches(0.42), row2_y + Emu(120000), Inches(1.35), Emu(900000),
       "CZ\n非突破\n↓\n通常時\n戻り", 8, color=C_GRAY, align=PP_ALIGN.CENTER)

    footer(s, "設計コメント：通常時→CZ→AT→POの到達ルートが長く、POにたどり着けない体験が不評の核心。",
           "蛇行2段フロー：上段 左→右、下段 右→左")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── CZ3種へのルートと天井", "3/9")

    # 左：通常時の基本フロー（縦ステップ）
    steps = [
        (C_GREEN,  "STEP 1",  "コインを投入して通常遊技",
         "設定1〜6。演出は弱めが多く、\nCZ示唆がほぼ来ない体感。\n「当たらない」時間が続く。"),
        (C_GREEN2, "STEP 2",  "レア役・CZポイント獲得",
         "チェリー・スイカ等のレア役で\nCZポイントが加算される。\n"
         "ただし体感の当選率は低い。"),
        (C_YEL,   "STEP 3",  "CZ（3種類）当選",
         "ポイント規定到達 or 抽選で\nCZに突入。天井はAT間\n最大999G+α（設定変更後450G+α）。"),
        (C_CYAN,  "STEP 4",  "AT「ヨルムンガンドラッシュ」突入",
         "CZ突破でAT確定。\nCZ非突破は通常時に戻る\n（ループ）。"),
    ]
    sw = Inches(4.3)
    sh = Emu(960000)
    sx0 = Inches(0.28)
    for i, (ac, step, title, body) in enumerate(steps):
        sy = Inches(0.72) + i * (sh + Emu(65000))
        rect_b(s, sx0, sy, sw, sh, C_CARD, ac, 1.5)
        rect(s, sx0, sy, Emu(45000), sh, ac)
        tb(s, sx0 + Emu(75000), sy + Emu(50000), sw - Emu(100000), Emu(240000),
           step, 8, bold=True, color=ac, font=FONT_B)
        tb(s, sx0 + Emu(75000), sy + Emu(280000), sw - Emu(100000), Emu(260000),
           title, 10, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, sx0 + Emu(75000), sy + Emu(550000), sw - Emu(100000), sh - Emu(600000),
           body, 8, color=C_GRAY)
        if i < 3:
            arrow_d(s, sx0 + sw // 2, sy + sh + Emu(10000), ac)

    # 右上：CZ当選ルートまとめ
    rx = Inches(4.85)
    rw = Inches(4.9)
    rect_b(s, rx, Inches(0.72), rw, Emu(1500000), C_CARD2, C_GREEN, 1.5)
    rect(s, rx, Inches(0.72), Emu(45000), Emu(1500000), C_GREEN)
    tb(s, rx + Emu(75000), Inches(0.75), rw - Emu(100000), Emu(270000),
       "CZ当選ルート（主な方法）", 10, bold=True, color=C_GREEN, font=FONT_H)
    routes = [
        "① シューティングゾーン（10G+α・AT期待度約43%）",
        "② 滅びの丘（10G・突入時点でAT濃厚・上位AT期待度約50%）",
        "③ パーフェクトチャレンジ（AT突入濃厚・成功で上位AT）",
        "④ 天井：AT間 最大999G+α（設定変更後・恥の世紀終了後は450G+α）",
    ]
    for i, r in enumerate(routes):
        tb(s, rx + Emu(75000), Inches(1.05) + i * Emu(290000),
           rw - Emu(100000), Emu(270000), r, 8.5, color=C_WHITE)

    # 右下：プレイヤーが感じる実態
    rect_b(s, rx, Inches(0.72) + Emu(1600000), rw,
           Emu(2800000), C_CARD, C_RED, 1.5)
    rect(s, rx, Inches(0.72) + Emu(1600000),
         Emu(45000), Emu(2800000), C_RED)
    tb(s, rx + Emu(75000), Inches(0.72) + Emu(1650000),
       rw - Emu(100000), Emu(270000),
       "通常時プレイヤーの実態（客観的記録）", 10, bold=True, color=C_RED, font=FONT_H)
    facts = [
        "・「昨今のデキレや冷遇が可愛く見える」という実戦報告が複数",
        "・弱い演出ばかりでCZ示唆がほぼ来ない",
        "・レア役を引いてもCZが入らないセッションが続く",
        "・ゾーン狙いよりも純粋な天井狙いが有効という意見多数",
        "・「投資が重く、回収できないまま終わる」という声が支配的",
    ]
    for i, f in enumerate(facts):
        tb(s, rx + Emu(75000), Inches(0.72) + Emu(1930000) + i * Emu(330000),
           rw - Emu(100000), Emu(290000), f, 8, color=C_WHITE)

    footer(s, "設計コメント：通常時の「当たらない体感」がプレイヤーを消耗させ、POへの期待を消す前に離脱を招く。",
           "天井狙い以外に有効な立ち回りが少ない")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: ストーリーCZの仕組み
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZの仕組み ── 3種CZ・ストーリーCZ構造・デキレ感の正体", "4/9")

    # 左：CZ仕組みフロー
    lx = Inches(0.28)
    lw = Inches(4.4)

    rect_b(s, lx, Inches(0.72), lw, Emu(560000), C_CARD2, C_GREEN, 1.5)
    rect(s, lx, Inches(0.72), Emu(45000), Emu(560000), C_GREEN)
    tb(s, lx + Emu(75000), Inches(0.75), lw - Emu(100000), Emu(270000),
       "ストーリーCZの構造（プロローグ＋ジャッジ）", 11, bold=True, color=C_GREEN, font=FONT_H)
    tb(s, lx + Emu(75000), Inches(1.05), lw - Emu(100000), Emu(220000),
       "前半「プロローグ10G」＋後半「ジャッジパート3G」の2パート構成", 8.5, color=C_WHITE)

    # CZフロー縦ステップ
    cz_steps = [
        (C_GREEN,  "前半：プロローグ（10G）",
         "BARを狙えカットインでBAR停止（約1/3.2）\n→対応役が点灯。レア役で全対応役点灯"),
        (C_GREEN2, "後半：ジャッジパート（3G）",
         "3Gで対応役を引ければ\nヨルムンガンドチャンス突入濃厚"),
        (C_YEL,    "3G全成功ならパーフェクトチャレンジ獲得",
         "3回全て成功→CZ「パーフェクトチャレンジ」\n獲得濃厚（AT突入濃厚・上位AT当選も）"),
        (C_CYAN,   "CZ突破（成功）",  "ヨルムンガンドラッシュへ\n（おめでとう）"),
        (C_RED,    "CZ非突破（失敗）","通常時に戻る\nリトライが必要"),
    ]
    sh2 = Emu(640000)
    for i, (ac, t, b) in enumerate(cz_steps):
        sy = Inches(0.72) + Emu(560000) + Emu(100000) + i * (sh2 + Emu(50000))
        rect_b(s, lx, sy, lw, sh2, C_CARD, ac, 1.3)
        rect(s, lx, sy, Emu(40000), sh2, ac)
        tb(s, lx + Emu(65000), sy + Emu(55000), lw - Emu(80000), Emu(250000),
           t, 9, bold=True, color=ac, font=FONT_H)
        tb(s, lx + Emu(65000), sy + Emu(300000), lw - Emu(80000), sh2 - Emu(350000),
           b, 8, color=C_GRAY)
        if i < 4:
            arrow_d(s, lx + lw // 2, sy + sh2 + Emu(5000),
                    ac if i < 3 else C_GRAY)

    # 右上：突破条件
    rx = Inches(4.95)
    rw = Inches(4.75)
    rect_b(s, rx, Inches(0.72), rw, Emu(1350000), C_CARD2, C_GREEN2, 1.5)
    rect(s, rx, Inches(0.72), Emu(45000), Emu(1350000), C_GREEN2)
    tb(s, rx + Emu(75000), Inches(0.75), rw - Emu(100000), Emu(270000),
       "突破条件（解析情報）", 10, bold=True, color=C_GREEN2, font=FONT_H)
    conds = [
        "・BAR停止確率 約1/3.2（プロローグ10G）",
        "・ジャッジパート3G：対応役を引ければ突破",
        "・レア役成立時は対応役不問で成功濃厚",
        "・3G全成功でパーフェクトチャレンジ獲得濃厚",
        "・強レア役ほど成功率・上位AT当選率が高い傾向",
    ]
    for i, c in enumerate(conds):
        tb(s, rx + Emu(75000), Inches(1.04) + i * Emu(290000),
           rw - Emu(100000), Emu(265000), c, 8.5, color=C_WHITE)

    # 右中：デキレ感の正体
    rect_b(s, rx, Inches(0.72) + Emu(1450000), rw,
           Emu(1680000), C_CARD, C_RED, 1.5)
    rect(s, rx, Inches(0.72) + Emu(1450000),
         Emu(45000), Emu(1680000), C_RED)
    tb(s, rx + Emu(75000), Inches(0.72) + Emu(1500000),
       rw - Emu(100000), Emu(270000),
       "「デキレ感」の正体と構造", 10, bold=True, color=C_RED, font=FONT_H)
    dekire = [
        ("シナリオ進行型演出",
         "当否が先に決まって演出が後付けされる設計。\n"
         "プレイヤーが「自力感」を持てない。"),
        ("弱い演出で失敗を繰り返す",
         "「どうせ負ける演出だ」という学習が起き、\n"
         "CZ自体に期待感が湧かなくなる。"),
        ("強い演出でも負けがある",
         "信頼度設計が不透明で\n「信じられない」感覚が広がる。"),
    ]
    for i, (dt, dd) in enumerate(dekire):
        dy = Inches(0.72) + Emu(1800000) + i * Emu(420000)
        rect(s, rx + Emu(75000), dy, Emu(15000), Emu(390000), C_RED)
        tb(s, rx + Emu(115000), dy + Emu(30000), rw - Emu(150000), Emu(200000),
           dt, 8.5, bold=True, color=C_CRIM)
        tb(s, rx + Emu(115000), dy + Emu(215000), rw - Emu(150000), Emu(180000),
           dd, 8, color=C_GRAY)

    # 右下：ATの平均G数
    rect_b(s, rx, Inches(0.72) + Emu(3230000), rw,
           Emu(560000), RGBColor(0x04, 0x14, 0x08), C_CYAN, 1.2)
    tb(s, rx + Emu(75000), Inches(0.72) + Emu(3280000),
       rw - Emu(100000), Emu(200000),
       "CZ突破後の報酬：AT「ヨルムンガンドラッシュ」突入\n　→ 次のスライドで詳しく解説", 8.5, color=C_CYAN)

    footer(s, "設計コメント：シナリオ型CZは「自力感の欠如」を生み、繰り返すほど打感が悪化する悪循環を生む。",
           "CZ突破率の設定差も大きな判別要素")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT「ヨルムンガンドラッシュ」の遊び方
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「ヨルムンガンドラッシュ」── 遊び方とPO到達ルート", "5/9")

    # 上段：ATの基本情報
    rect_b(s, Inches(0.28), Inches(0.72), Inches(9.44), Emu(480000),
           RGBColor(0x08, 0x20, 0x10), C_GREEN2, 1.5)
    atinfo = [
        ("AT名", "ヨルムンガンドラッシュ"),
        ("純増", "約2.4枚/G"),
        ("AT期待ゲーム数", "約90G"),
        ("継続方式", "ストーリーCZ(ループ)"),
    ]
    for i, (k, v) in enumerate(atinfo):
        bx0 = Inches(0.55) + i * Inches(2.38)
        tb(s, bx0, Inches(0.76), Inches(1.1), Emu(190000),
           k, 7.5, bold=True, color=C_GRAY, font=FONT_B)
        tb(s, bx0, Inches(0.96), Inches(2.2), Emu(220000),
           v, 10, bold=True, color=C_GREEN2, font=FONT_H)

    # 中段左：AT中のやること（縦ステップ）
    lx = Inches(0.28)
    lw = Inches(4.4)

    at_steps = [
        (C_GREEN,  "AT中 STEP 1",   "50G消化・ゲーム数上乗せを狙う",
         "AT中はレバーを叩くだけ。\n純増2.4枚/Gでコインを増やす。\nレア役で上乗せ・ストーリー・疑似ボーナスを抽選。"),
        (C_GREEN2, "AT中 STEP 2",   "AT中ボーナスを引く",
         "BB（ビッグボーナス）以上の\nボーナスがPO当選の鍵。\nストーリー成功で特化ゾーン突入のチャンスも。"),
        (C_CYAN,   "AT中 STEP 3",   "BB以上フリーズでPO確定",
         "BB以上＋フリーズ発生でPO（パーフェクトオーダー）\nへ移行確定。ATエンディング到達でも恥の世紀突入。"),
        (C_YEL,    "ATエンディング到達",  "「恥の世紀」CZへ突入",
         "ATエンディングに到達すると\n「恥の世紀」（期待度約50%）突入。\n成功でPOへ、失敗で通常時戻り。"),
    ]
    sh = Emu(870000)
    for i, (ac, step, title, body) in enumerate(at_steps):
        sy = Inches(1.32) + i * (sh + Emu(55000))
        rect_b(s, lx, sy, lw, sh, C_CARD, ac, 1.5)
        rect(s, lx, sy, Emu(40000), sh, ac)
        tb(s, lx + Emu(65000), sy + Emu(50000), lw - Emu(85000), Emu(230000),
           step, 8, bold=True, color=ac, font=FONT_B)
        tb(s, lx + Emu(65000), sy + Emu(275000), lw - Emu(85000), Emu(270000),
           title, 10, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, lx + Emu(65000), sy + Emu(545000), lw - Emu(85000), sh - Emu(590000),
           body, 8, color=C_GRAY)
        if i < 3:
            arrow_d(s, lx + lw // 2, sy + sh + Emu(8000), ac)

    # 右上：PO到達ルート図解
    rx = Inches(4.95)
    rw = Inches(4.75)
    rect_b(s, rx, Inches(1.32), rw, Emu(1500000),
           RGBColor(0x04, 0x18, 0x20), C_CYAN, 2.0)
    rect(s, rx, Inches(1.32), Emu(50000), Emu(1500000), C_CYAN)
    tb(s, rx + Emu(80000), Inches(1.35), rw - Emu(110000), Emu(280000),
       "PO（パーフェクトオーダー）到達ルート", 11, bold=True, color=C_CYAN, font=FONT_H)
    po_routes = [
        ("メインルート",  "AT中 BB以上+フリーズ → PO確定"),
        ("エンディング",  "ATエンディング到達 → 恥の世紀（約50%）→ PO突入"),
        ("通常時から",    "通常時からも直接PO突入の可能性あり"),
        ("設定差あり",    "高設定ほどPO到達率・恥の世紀成功率が高い"),
    ]
    for i, (rt, rd) in enumerate(po_routes):
        ry0 = Inches(1.68) + i * Emu(290000)
        rect(s, rx + Emu(80000), ry0, Emu(15000), Emu(255000), C_CYAN)
        tb(s, rx + Emu(120000), ry0 + Emu(30000), Inches(1.0), Emu(200000),
           rt, 8, bold=True, color=C_CYAN)
        tb(s, rx + Emu(1200000), ry0 + Emu(30000), rw - Emu(1250000), Emu(200000),
           rd, 8, color=C_WHITE)

    # 右下：AT設計の課題
    rect_b(s, rx, Inches(1.32) + Emu(1600000), rw,
           Emu(2700000), C_CARD, C_RED, 1.5)
    rect(s, rx, Inches(1.32) + Emu(1600000),
         Emu(45000), Emu(2700000), C_RED)
    tb(s, rx + Emu(75000), Inches(1.32) + Emu(1650000),
       rw - Emu(100000), Emu(270000),
       "AT設計の課題（客観的評価）", 10, bold=True, color=C_RED, font=FONT_H)
    at_issues = [
        "・AT純増2.4枚/Gはそこそこだが「POに到達しないと物足りない」",
        "・BB以上当選が必要なため、SBばかりではPOに行けない",
        "・「AT何回打ってもPOに行かない」という低設定体験が口コミ化",
        "・ATが短く（90G）、到達感が薄いまま終わるケースが多い",
        "・有利区間が残り少ない状態でATに入ると満足できない問題も",
    ]
    for i, ai in enumerate(at_issues):
        tb(s, rx + Emu(75000), Inches(1.32) + Emu(1950000) + i * Emu(300000),
           rw - Emu(100000), Emu(265000), ai, 8, color=C_WHITE)

    footer(s, "設計コメント：AT中のPO到達には「BB以上」という条件があり、低設定ではPOが幻になりやすい構造。",
           "BB以上フリーズ→PO確定が最大の盛り上がりポイント")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: PO「パーフェクトオーダー」の遊び方
# ══════════════════════════════════════════════════════════════
def s_po(prs):
    s = new_slide(prs)
    hdr(s, "PO「パーフェクトオーダー」── 到達ルート・遊び方・恥の世紀ループ", "6/9")

    # 上段：PO基本スペック強調
    rect(s, 0, Inches(0.58), SLIDE_W, Emu(450000), RGBColor(0x04, 0x18, 0x20))
    po_kv = [
        ("PO名称",      "パーフェクトオーダー（PO）"),
        ("純増",        "約5.0枚/G  ← 業界トップクラス"),
        ("初期G数",     "50G（初期保証）"),
        ("ボーナス",    "PO中は全てBB以上に当選"),
        ("ループ",      "恥の世紀（期待度約50%）でPOループ"),
    ]
    # 5カラム均等: step=1.88 → i=4: bx=0.3+7.52=7.82, 幅1.8 → 末端=9.62 OK
    for i, (k, v) in enumerate(po_kv):
        bx0 = Inches(0.3) + i * Inches(1.88)
        col = C_YEL if "恥" in k else C_CYAN
        tb(s, bx0, Inches(0.65), Inches(1.75), Emu(200000),
           k, 7.5, bold=True, color=C_GRAY)
        tb(s, bx0, Inches(0.86), Inches(1.80), Emu(230000),
           v, 9, bold=True, color=col, font=FONT_H)

    # 左：PO中の遊び方フロー
    lx = Inches(0.28)
    lw = Inches(4.35)

    rect_b(s, lx, Inches(1.22), lw, Emu(550000),
           RGBColor(0x04, 0x18, 0x20), C_CYAN, 2.0)
    rect(s, lx, Inches(1.22), Emu(45000), Emu(550000), C_CYAN)
    tb(s, lx + Emu(75000), Inches(1.25), lw - Emu(100000), Emu(270000),
       "PO中の遊び方", 11, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, lx + Emu(75000), Inches(1.55), lw - Emu(100000), Emu(200000),
       "純増5.0枚/Gで豪快にコインを増やす50Gを楽しむ", 8.5, color=C_WHITE)

    po_steps = [
        (C_CYAN,  "PO STEP 1",  "50G消化スタート",
         "初期50Gを全力で消化。\n純増5.0枚/Gで毎ゲーム増加。\nボーナスは全てBB以上に当選。"),
        (C_CYAN,  "PO STEP 2",  "PO中ボーナス全てBB以上",
         "PO中のボーナスは全てBB以上。\nBB以上当選でPOループの抽選。"),
        (C_YEL,   "PO STEP 3",  "消化後「恥の世紀」CZ突入",
         "PO消化後は必ず「恥の世紀」へ突入。\n期待度約50%でPO再突入を抽選。"),
        (C_GREEN, "PO STEP 4",  "ループ継続 or 終了",
         "ループするたびに純増5枚が継続。\n強レア役→特化ゾーンも濃厚。\n万枚も射程に入る連鎖設計。"),
    ]
    sh3 = Emu(770000)
    for i, (ac, step, title, body) in enumerate(po_steps):
        sy = Inches(1.22) + Emu(550000) + Emu(70000) + i * (sh3 + Emu(50000))
        rect_b(s, lx, sy, lw, sh3, C_CARD, ac, 1.5)
        rect(s, lx, sy, Emu(40000), sh3, ac)
        tb(s, lx + Emu(65000), sy + Emu(45000), lw - Emu(85000), Emu(210000),
           step, 8, bold=True, color=ac, font=FONT_B)
        tb(s, lx + Emu(65000), sy + Emu(255000), lw - Emu(85000), Emu(245000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, lx + Emu(65000), sy + Emu(490000), lw - Emu(85000), sh3 - Emu(530000),
           body, 8, color=C_GRAY)
        if i < 3:
            arrow_d(s, lx + lw // 2, sy + sh3 + Emu(8000), ac)

    # 右上：恥の世紀の詳細
    rx = Inches(4.9)
    rw = Inches(4.8)
    rect_b(s, rx, Inches(1.22), rw, Emu(1800000),
           RGBColor(0x14, 0x10, 0x02), C_YEL, 2.0)
    rect(s, rx, Inches(1.22), Emu(50000), Emu(1800000), C_YEL)
    tb(s, rx + Emu(80000), Inches(1.25), rw - Emu(110000), Emu(280000),
       "「恥の世紀」── POループ確定演出", 11, bold=True, color=C_YEL, font=FONT_H)
    haji = [
        ("突入契機",  "ATエンディング到達後・PO消化後に突入"),
        ("期待度",    "約50%でPO当選。毎ゲーム成立役に応じてPO抽選"),
        ("レア役",    "レア役成立で成功濃厚・強レア役でPO+特化ゾーン濃厚"),
        ("体験価値",  "「恥の世紀来た！」が最高の演出体験と評価される"),
        ("ループ",    "PO消化後も恥の世紀に突入→万枚も視野に入る"),
    ]
    for i, (hk, hv) in enumerate(haji):
        hy = Inches(1.56) + i * Emu(290000)
        rect(s, rx + Emu(80000), hy, Emu(15000), Emu(260000), C_YEL)
        tb(s, rx + Emu(120000), hy + Emu(30000), Inches(1.0), Emu(200000),
           hk, 8, bold=True, color=C_YEL)
        tb(s, rx + Emu(1230000), hy + Emu(30000), rw - Emu(1280000), Emu(200000),
           hv, 8, color=C_WHITE)

    # 右下：PO純増5.0枚の体験価値
    rect_b(s, rx, Inches(1.22) + Emu(1900000), rw,
           Emu(2400000), RGBColor(0x04, 0x18, 0x20), C_CYAN, 1.8)
    rect(s, rx, Inches(1.22) + Emu(1900000),
         Emu(50000), Emu(2400000), C_CYAN)
    tb(s, rx + Emu(80000), Inches(1.22) + Emu(1950000),
       rw - Emu(110000), Emu(270000),
       "純増5.0枚/G の体験価値", 11, bold=True, color=C_CYAN, font=FONT_H)
    po_value = [
        "50G × 5.0枚 ＝ 期待値 約250枚（1POあたり）",
        "PO消化後は必ず恥の世紀（約50%）→ ループ期待",
        "ループするほど積み上がる体験設計（万枚も視野）",
        "業界水準（3〜4枚）を大きく超える純増スピード",
        "「PO中は別ゲーム」という没入感・純増5枚の実感",
    ]
    for i, pv in enumerate(po_value):
        tb(s, rx + Emu(80000), Inches(1.22) + Emu(2250000) + i * Emu(300000),
           rw - Emu(110000), Emu(265000), pv, 8.5, color=C_WHITE)

    footer(s, "設計コメント：PO純増5.0枚は本物の体験価値を持つ。到達できたプレイヤーは高評価するが、到達できない多数が不評を形成。",
           "「恥の世紀」が来れば最高潮 ── 問題は到達確率")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: Part B 開始 ── 面白さの設計（POのポテンシャル）
# ══════════════════════════════════════════════════════════════
def s_strengths(prs):
    s = new_slide(prs)
    hdr(s, "【分析①】面白さの設計 ── PO体験が持つ本物のポテンシャル", "7/9")

    # 左パネル：PO体験の価値を3段
    lx = Inches(0.28)
    lw = Inches(4.55)

    rect(s, lx, Inches(0.72), lw, Emu(300000), C_CARD2)
    tb(s, lx + Emu(60000), Inches(0.75), lw - Emu(90000), Emu(240000),
       "PO「パーフェクトオーダー」が持つ本物の設計力", 11, bold=True,
       color=C_CYAN, font=FONT_H)

    po_vals = [
        (C_CYAN,  "純増5.0枚/G は本物のトップスペック",
         "通常AT（2.4枚）の約2倍の速度でコインが増える。\n"
         "50G × 5枚 = 約250枚の期待値は、1POで十分な満足感。\n"
         "業界基準（3〜4枚級）を大きく上回る数値。"),
        (C_YEL,  "「恥の世紀ループ」は明確で分かりやすいゴール",
         "「恥の世紀が来るかどうか」という一点に集中する緊張感。\n"
         "テキスト表示で明確に確定するため達成感が明確。\n"
         "ループのたびに得られる快感の連鎖設計は優秀。"),
        (C_GREEN, "「到達できた人」は高評価する台",
         "ネット上の高評価レビューの多くがPO到達者によるもの。\n"
         "「PO中は本当に楽しい」という声は一貫している。\n"
         "設計そのものは間違っておらず、問題は到達率にある。"),
    ]
    sh = Emu(1060000)
    for i, (ac, t, b) in enumerate(po_vals):
        sy = Inches(0.72) + Emu(300000) + Emu(60000) + i * (sh + Emu(60000))
        rect_b(s, lx, sy, lw, sh, C_CARD, ac, 1.5)
        rect(s, lx, sy, Emu(45000), sh, ac)
        tb(s, lx + Emu(75000), sy + Emu(60000), lw - Emu(100000), Emu(280000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, lx + Emu(75000), sy + Emu(345000), lw - Emu(100000), sh - Emu(390000),
           b, 8.5, color=C_WHITE)

    # 右上：設計の良い点マトリクス
    rx = Inches(5.1)
    rw = Inches(4.6)

    rect_b(s, rx, Inches(0.72), rw, Emu(2100000), C_CARD2, C_GREEN, 1.5)
    rect(s, rx, Inches(0.72), Emu(45000), Emu(2100000), C_GREEN)
    tb(s, rx + Emu(75000), Inches(0.75), rw - Emu(100000), Emu(270000),
       "設計として優れている点", 10, bold=True, color=C_GREEN, font=FONT_H)
    good_pts = [
        ("ストーリー × スロット", "原作ファン向けの演出統合は丁寧"),
        ("純増スピードの差別化", "AT→PO の純増差が体感でも明確"),
        ("恥の世紀の分かりやすさ", "ループ確定演出の明瞭さは設計の正解"),
        ("設定差の多様性", "REGキャラ・PO率・恥の世紀率の3軸判別"),
        ("天井設計",       "概算1000G天井で最終救済ルートは存在する"),
    ]
    for i, (gk, gv) in enumerate(good_pts):
        gy = Inches(1.05) + i * Emu(330000)
        rect(s, rx + Emu(75000), gy, Emu(15000), Emu(295000), C_GREEN)
        tb(s, rx + Emu(120000), gy + Emu(35000), Inches(1.5), Emu(220000),
           gk, 8.5, bold=True, color=C_GREEN2)
        tb(s, rx + Emu(1740000), gy + Emu(35000), rw - Emu(1800000), Emu(220000),
           gv, 8.5, color=C_WHITE)

    # 右下：到達者 vs 非到達者の評価対比
    rect_b(s, rx, Inches(0.72) + Emu(2200000), rw,
           Emu(2100000), C_CARD, C_CYAN, 1.2)
    rect(s, rx, Inches(0.72) + Emu(2200000),
         Emu(45000), Emu(2100000), C_CYAN)
    tb(s, rx + Emu(75000), Inches(0.72) + Emu(2250000),
       rw - Emu(100000), Emu(270000),
       "評価の二極化 ── PO到達者 vs 非到達者", 10, bold=True, color=C_CYAN, font=FONT_H)
    evals = [
        (C_CYAN, "PO到達者",    "「PO中は最高」「恥の世紀が来たら神台」\n「もう一度打ちたい」という肯定評価"),
        (C_RED,  "PO非到達者",  "「全然当たらない」「投資がかさむだけ」\n「通常時が苦行」という否定評価"),
    ]
    for i, (ac, ek, ev) in enumerate(evals):
        ey = Inches(0.72) + Emu(2570000) + i * Emu(840000)
        rect(s, rx + Emu(75000), ey, Emu(15000), Emu(780000), ac)
        tb(s, rx + Emu(120000), ey + Emu(40000), rw - Emu(150000), Emu(260000),
           ek, 9, bold=True, color=ac, font=FONT_H)
        tb(s, rx + Emu(120000), ey + Emu(310000), rw - Emu(150000), Emu(450000),
           ev, 8, color=C_WHITE)

    footer(s, "設計コメント：PO体験の質は本物。問題は「体験できる割合」が低すぎる設計にあり、到達率の引き上げが解決策。",
           "高性能POが届かない台 ── これが不評の本質的構造")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 不評の構造（通常時の渋さ・有利区間・デキレ感）
# ══════════════════════════════════════════════════════════════
def s_issues(prs):
    s = new_slide(prs)
    hdr(s, "【分析②】不評の構造 ── 3つの設計課題を解剖する", "8/9")

    # 3カラム構造
    col_w = Inches(3.05)
    col_gap = Inches(0.12)
    cols_x = [Inches(0.28), Inches(0.28) + col_w + col_gap,
              Inches(0.28) + 2 * (col_w + col_gap)]
    issues_data = [
        (C_RED,  "課題①\n通常時の渋さ",
         "問題の実態",
         "・「昨今のデキレが可愛く見える」実戦報告\n"
         "・弱い演出が続きCZ示唆がほぼ来ない\n"
         "・レア役でもCZに入らないセッションが連続\n"
         "・ゾーン狙いより天井狙いのみが有効",
         "設計的原因",
         "CZ当選率が全体的に低く設定されており、\n"
         "特にCZポイントの規定到達が遠い。\n"
         "通常時に「何かが起きる」感覚が薄い。",
         "ユーザー体験",
         "「投資が重い→回収できない→離席」\nという離脱サイクルが発生"),
        (C_RED,  "課題②\n有利区間上限問題",
         "問題の実態",
         "・大量獲得中に有利区間上限でエンディング強制\n"
         "・「PO中なのに残枚数が少なくエンディングは虚無」\n"
         "・出玉が伸びるほど終わりが近づく矛盾体験\n"
         "・天井（999G+α）まで粘ると区間がほぼ残らない",
         "設計的原因",
         "有利区間の管理が枚数上限と強く連動。\n"
         "天井狙いで投資した分だけ残区間が減り、\n"
         "AT入ってもPOを十分楽しめない問題。",
         "ユーザー体験",
         "「盛り上がった瞬間に終わる」\n最大の失望体験を生む"),
        (C_RED,  "課題③\nストーリーCZのデキレ感",
         "問題の実態",
         "・シナリオ進行型でプレイヤーの自力感がない\n"
         "・「どうせ負ける演出」という学習が起きる\n"
         "・強い演出でも負けがあり信頼度不透明\n"
         "・CZ自体への期待感が失われていく",
         "設計的原因",
         "結果が先に決まってから演出が進行する\n"
         "シナリオ型の構造的欠陥。\n"
         "プレイヤーの能動性が排除されている。",
         "ユーザー体験",
         "「やらされている感覚」が蓄積し\n「打感が悪い台」として記憶される"),
    ]

    for ci, (cx, (ac, hd, l1, c1, l2, c2, l3, c3)) in enumerate(
            zip(cols_x, issues_data)):
        # ヘッダー
        rect_b(s, cx, Inches(0.72), col_w, Emu(560000), C_CARD, ac, 2.0)
        rect(s, cx, Inches(0.72), Emu(40000), Emu(560000), ac)
        tb(s, cx + Emu(65000), Inches(0.75), col_w - Emu(85000), Emu(520000),
           hd, 11, bold=True, color=ac, font=FONT_H, align=PP_ALIGN.CENTER)

        # 3段カード
        card_data = [(l1, c1), (l2, c2), (l3, c3)]
        card_h = [Emu(1350000), Emu(1000000), Emu(820000)]
        cy0 = Inches(0.72) + Emu(560000) + Emu(80000)
        for j, ((lab, cont), ch) in enumerate(zip(card_data, card_h)):
            bc = C_GRAY if j == 2 else (C_RED if j == 0 else C_YEL)
            rect_b(s, cx, cy0, col_w, ch, C_CARD, bc, 0.8)
            tb(s, cx + Emu(50000), cy0 + Emu(50000), col_w - Emu(70000),
               Emu(230000), lab, 8, bold=True, color=bc)
            tb(s, cx + Emu(50000), cy0 + Emu(280000), col_w - Emu(70000),
               ch - Emu(330000), cont, 7.5, color=C_WHITE)
            cy0 += ch + Emu(60000)

    footer(s, "設計コメント：3つの課題が複合的に作用し「打感の悪い台」を形成。単独では軽い課題も重なると致命的になる。",
           "有利区間・到達率・自力感 ── 3点同時改善が必要")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "【まとめ】設計から学べること ── 高性能だが刺さらない台の教訓", "9/9")

    # 左：不評機種から学べること
    lx = Inches(0.28)
    lw = Inches(4.55)

    rect(s, lx, Inches(0.72), lw, Emu(295000), RGBColor(0x08, 0x28, 0x14))
    tb(s, lx + Emu(60000), Inches(0.75), lw - Emu(90000), Emu(240000),
       "ヨルムンガンドが示す「設計の教訓」4選", 11, bold=True,
       color=C_GREEN2, font=FONT_H)

    lessons = [
        (C_RED,   "教訓①：到達できないスペックは「幻」になる",
         "PO純増5.0枚が本物でも、通常時の渋さで\n"
         "到達できなければ存在しないも同然。\n"
         "性能と到達可能性のバランスが設計の核心。"),
        (C_CYAN,  "教訓②：自力感がないと打感は壊れる",
         "シナリオ型CZは「受け身体験」を生む。\n"
         "「自分が引いた」という能動体験こそが\n"
         "スロットの面白さの源泉。"),
        (C_YEL,   "教訓③：有利区間の「見えない壁」がユーザーを離す",
         "大量獲得中の強制終了は最大の失望体験。\n"
         "有利区間管理の透明性確保が\n"
         "長期稼働の前提条件になる。"),
        (C_GREEN, "教訓④：通常時の演出強弱が「当たる気がする台」を作る",
         "弱い演出が続くと希望が薄れ離席を招く。\n"
         "通常時に「何かが起きる予感」を設計できるかが\n"
         "来店継続率を左右する。"),
    ]
    sh = Emu(940000)
    for i, (ac, t, b) in enumerate(lessons):
        sy = Inches(0.72) + Emu(295000) + Emu(60000) + i * (sh + Emu(45000))
        rect_b(s, lx, sy, lw, sh, C_CARD, ac, 1.5)
        rect(s, lx, sy, Emu(45000), sh, ac)
        tb(s, lx + Emu(75000), sy + Emu(55000), lw - Emu(100000), Emu(265000),
           t, 9, bold=True, color=ac, font=FONT_H)
        tb(s, lx + Emu(75000), sy + Emu(320000), lw - Emu(100000), sh - Emu(370000),
           b, 8, color=C_WHITE)

    # 右：設計原則 + 総括
    rx = Inches(5.1)
    rw = Inches(4.6)

    rect(s, rx, Inches(0.72), rw, Emu(290000), C_CARD2)
    tb(s, rx + Emu(60000), Inches(0.75), rw - Emu(90000), Emu(235000),
       "設計原則：次の台を作るとき意識すること", 10, bold=True,
       color=C_GREEN2, font=FONT_H)

    principles = [
        (C_GREEN, "通常時の自力感がAT到達への期待感を生む"),
        (C_CYAN,  "高性能ATは到達可能性が担保されて初めて機能する"),
        (C_RED,   "有利区間管理の不透明さは不信感の温床になる"),
        (C_YEL,   "演出の強弱設計が「当たる気がする台」を作る"),
        (C_GREEN2,"CZ突破率と通常時CZ当選率のバランスが打感を決める"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = Inches(0.72) + Emu(290000) + Emu(60000) + i * Emu(490000)
        rect(s, rx, py0, Emu(20000), Emu(445000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(75000), rw - Emu(70000), Emu(360000),
           p, 8.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, Inches(0.72) + Emu(290000) + Emu(60000) + 5 * Emu(490000) + Emu(80000),
           rw, Emu(1000000),
           RGBColor(0x04, 0x14, 0x08), C_GREEN, 2.0)
    summy = Inches(0.72) + Emu(290000) + Emu(60000) + 5 * Emu(490000) + Emu(80000)
    tb(s, rx + Emu(65000), summy + Emu(60000), rw - Emu(90000), Emu(270000),
       "総括", 9.5, bold=True, color=C_GREEN2, font=FONT_H)
    tb(s, rx + Emu(65000), summy + Emu(330000), rw - Emu(90000), Emu(600000),
       "ヨルムンガンドは「高性能だが刺さらない台」の典型例。\n"
       "POの設計思想は正しい。しかし通常時の苦痛が\n"
       "PO体験の価値を「見えないもの」にしてしまった。\n"
       "不評の分析こそが、次の傑作への最短ルートになる。",
       8, color=C_WHITE)

    footer(s, "設計コメント：不評の分析は批判ではなく学習。高性能POが輝く台を作るために、通常時設計から見直すことが鍵。",
           "「高性能だが刺さらない台」── この教訓を活かせる設計者が次の傑作を生む")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Part A: 説明パート
    s_title(prs)           # 1: タイトル・スペック・3ポイント
    s_flow_overview(prs)   # 2: ゲームフロー全体図（蛇行2段）
    s_normal(prs)          # 3: 通常時の遊び方
    s_cz(prs)              # 4: ストーリーCZの仕組み
    s_at(prs)              # 5: ATの遊び方とPO到達ルート
    s_po(prs)              # 6: PO「パーフェクトオーダー」

    # Part B: 分析パート
    s_strengths(prs)       # 7: 面白さの設計（POポテンシャル）
    s_issues(prs)          # 8: 不評の構造（3つの課題）
    s_matome(prs)          # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
