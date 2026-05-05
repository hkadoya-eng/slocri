"""
スマスロ スーパーブラックジャック 機種説明＋分析 統合版資料 v2
（セブンリーグ（山佐ネクス）・2025年2月3日導入）
出力: proposals/機種分析/スーパーブラックジャック/blackjack_guide_v2.pptx
テーマ: クリーム背景 × カジノグリーン × リオピンク × ゴールド（華やか女の子台）
"""
import io
import os
import sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
           "proposals", "機種分析", "スーパーブラックジャック", "blackjack_guide_v2.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（クリーム×グリーン×ピンク×ゴールド）─────────────
C_DARK_TEXT = RGBColor(0x1A, 0x14, 0x08)   # 濃いブラウン（本文）
C_MID_GRAY  = RGBColor(0x55, 0x50, 0x40)   # 中間グレー（補足）
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 白（ヘッダー内）
C_GREEN  = RGBColor(0x0C, 0x7A, 0x3F)   # カジノグリーン
C_GREEN2 = RGBColor(0x16, 0xA0, 0x58)   # 明るいグリーン
C_GOLD   = RGBColor(0xBB, 0x90, 0x15)   # ゴールド
C_RED    = RGBColor(0xCC, 0x22, 0x22)   # カードの赤
C_PINK   = RGBColor(0xDD, 0x33, 0x88)   # リオのピンク
C_MINT   = RGBColor(0x11, 0x99, 0x66)   # ミント（ミントキャラ）
C_BLUE   = RGBColor(0x22, 0x66, 0xCC)   # ブルー（補足）
C_GRAY   = RGBColor(0x88, 0x85, 0x70)   # グレー

C_FILL_GREEN   = RGBColor(0xEC, 0xF8, 0xF2)
C_FILL_GOLD    = RGBColor(0xFE, 0xFA, 0xEA)
C_FILL_RED     = RGBColor(0xFF, 0xED, 0xED)
C_FILL_PINK    = RGBColor(0xFF, 0xEE, 0xF5)
C_FILL_MINT    = RGBColor(0xEC, 0xFF, 0xF5)
C_FILL_BLUE    = RGBColor(0xEE, 0xF4, 0xFF)
C_FILL_DEFAULT = RGBColor(0xF8, 0xF4, 0xEC)  # クリーム

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

TOTAL_SLIDES = 9


# ── 背景生成（クリーム×斜線×左端グリーンライン）────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (0xF8, 0xF4, 0xEC))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 40):
        draw.line([(i, 0), (0, i)], fill=(0xF0, 0xEB, 0xDE), width=1)
    for x in range(0, 6):
        draw.line([(x, 0), (x, h)], fill=(0x0C, 0x7A, 0x3F))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


# ── ヘルパー関数 ───────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border_color=None, border_pt=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = Pt(border_pt)
    else:
        shp.line.fill.background()
    return shp


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_DARK_TEXT
    run.font.name = font or FONT_B


def hdr(slide, title_text, sub="", pg=""):
    """ヘッダー（白背景×グリーン左帯×濃色タイトル）"""
    rect_b(slide, 0, 0, SLIDE_W, Inches(0.60), C_WHITE, C_GREEN, 1.0)
    rect(slide, 0, 0, Emu(55000), Inches(0.60), C_GREEN)
    tb(slide, Inches(0.18), Emu(18000), Inches(7.5), Emu(420000),
       title_text, 13, bold=True, color=C_GREEN, font=FONT_H)
    if sub:
        tb(slide, Inches(0.18), Emu(290000), Inches(7.0), Emu(220000),
           sub, 7.5, color=C_MID_GRAY)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.60), SLIDE_W, Emu(6000), C_GREEN)


def net_note(slide):
    tb(slide, Inches(8.0), Inches(5.35), Inches(1.85), Emu(200000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, design_comment, sub_text=""):
    fy = Inches(5.05)
    rect(slide, 0, fy, SLIDE_W, Inches(0.55), RGBColor(0x0C, 0x7A, 0x3F))
    tb(slide, Inches(0.2), fy + Emu(30000), Inches(7.0), Emu(380000),
       "【設計】" + design_comment, 7.5, bold=True, color=C_WHITE)
    if sub_text:
        tb(slide, Inches(0.2), fy + Emu(230000), Inches(7.5), Emu(200000),
           sub_text, 6.5, color=RGBColor(0xCC, 0xFF, 0xDD))
    net_note(slide)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GREEN
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp = slide.shapes.add_shape(17, cx - Emu(90000), y, Emu(180000), Emu(200000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GREEN
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・キャラクター・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル（薄グリーン地）
    rect_b(s, 0, 0, Inches(4.75), SLIDE_H, C_FILL_GREEN, C_GREEN, 1.0)
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_GREEN)

    # PartAバッジ
    rect_b(s, Inches(0.20), Inches(0.18), Inches(1.35), Emu(260000),
           C_GREEN, None, 0)
    tb(s, Inches(0.20), Inches(0.18), Inches(1.35), Emu(260000),
       "Part A 説明編", 7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.20), Inches(0.62), Inches(4.4), Emu(280000),
       "機種説明＋分析 統合版資料 v2", 9.5, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.20), Inches(1.00), Inches(4.4), Emu(500000),
       "スマスロ\nスーパーブラックジャック", 25, bold=True, color=C_GREEN, font=FONT_H)
    tb(s, Inches(0.20), Inches(2.80), Inches(4.4), Emu(260000),
       "── カードゲームと華やかキャラが彩る新感覚体験", 8.5, color=C_MID_GRAY, font=FONT_H)

    # アワードバッジ
    rect_b(s, Inches(0.20), Inches(3.26), Inches(3.0), Emu(280000),
           C_FILL_GOLD, C_GOLD, 1.5)
    tb(s, Inches(0.20), Inches(3.26), Inches(3.0), Emu(280000),
       "  パチスロアワード2025 BRONZE受賞", 8.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.20), Inches(3.72), Inches(4.4), Emu(220000),
       "メーカー: セブンリーグ（山佐ネクス）　導入: 2025年2月3日", 7.5, color=C_MID_GRAY)
    tb(s, Inches(0.20), Inches(3.96), Inches(4.4), Emu(220000),
       "機械割: 97.8〜112.7%　純増: 約5.1枚/G　設定: 1〜6", 7.5, color=C_MID_GRAY)
    tb(s, Inches(0.20), Inches(4.20), Inches(4.4), Emu(220000),
       "4号機「スーパーブラックジャック」のスマスロリメイク", 7.5, color=C_MID_GRAY)

    # 右パネル（キャラクター紹介）
    rect_b(s, Inches(4.85), 0, Inches(5.15), SLIDE_H, C_WHITE, C_PINK, 1.0)

    tb(s, Inches(5.0), Inches(0.10), Inches(4.7), Emu(270000),
       "登場キャラクター", 11, bold=True, color=C_PINK, font=FONT_H)

    chars = [
        (C_PINK,  C_FILL_PINK,  "リオ",
         "主人公キャラ。元気・明るい・パワフルな赤髪の女の子。\nSBJの顔として全編に登場し、大当りを盛り上げる。"),
        (C_MINT,  C_FILL_MINT,  "ミント",
         "クールな緑髪の女の子。リオの対になる存在。\nSST（スーパーストックタイム）の盛り上げ役。"),
        (C_GOLD,  C_FILL_GOLD,  "リナ",
         "お嬢様系の金髪キャラ。上品なBJ演出を担当。\nジョーカーモード突入演出に深く関与。"),
    ]
    for i, (col, fill, name, desc) in enumerate(chars):
        cy = Inches(0.55 + i * 1.60)
        rect_b(s, Inches(5.0), cy, Inches(4.7), Inches(1.45),
               fill, col, 1.5)
        rect(s, Inches(5.0), cy, Emu(55000), Inches(1.45), col)
        tb(s, Inches(5.12), cy + Emu(60000), Inches(1.0), Emu(270000),
           name, 12, bold=True, color=col)
        tb(s, Inches(5.12), cy + Emu(330000), Inches(4.2), Emu(720000),
           desc, 8, color=C_DARK_TEXT)

    # 右下：3ポイント
    rect_b(s, Inches(5.0), Inches(5.30 - 0.53), Inches(4.7), Emu(320000),
           C_FILL_GREEN, C_GREEN, 1.0)
    tb(s, Inches(5.1), Inches(5.30 - 0.50), Inches(4.5), Emu(290000),
       "SBJループ・リオチャンス2段構え・ジョーカーモード ── 3軸で爆発する設計",
       8, bold=True, color=C_GREEN)

    footer(s,
           "女の子が活躍する華やかな世界観 × BJ演出の緊張感。稼働とアワードが実力を証明",
           "リオ・ミント・リナの3キャラがゲームの各局面を彩り、単なるカジノ台とは一線を画す")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段・コンパクト）
# ══════════════════════════════════════════════════════════════
def s_gameflow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図", pg=f"2 / {TOTAL_SLIDES}")

    _bw = Inches(1.50)
    _bh = Inches(0.78)
    _gap = Inches(0.13)
    _mx  = Inches(0.17)

    # 上段 5ボックス
    top_nodes = [
        (C_FILL_DEFAULT, C_GRAY,   "通常時\n（低確/高確）"),
        (C_FILL_GOLD,    C_GOLD,   "BIG BONUS\n(青7/赤7)"),
        (C_FILL_GREEN,   C_GREEN,  "RC高確\n(高確率)"),
        (C_FILL_PINK,    C_PINK,   "リオチャンス\n（RC）"),
        (C_FILL_MINT,    C_MINT,   "ストックタイム\n（ST）"),
    ]
    top_y = Inches(0.78)
    for i, (fill, col, label) in enumerate(top_nodes):
        bx = _mx + i * (_bw + _gap)
        rect_b(s, bx, top_y, _bw, _bh, fill, col, 1.5)
        tb(s, bx, top_y, _bw, _bh, label, 8.5, bold=True, color=col,
           align=PP_ALIGN.CENTER)
        if i < len(top_nodes) - 1:
            arrow_r(s, bx + _bw, top_y + _bh / 2, col)

    # 折り返し矢印（右端から下へ）
    last_bx = _mx + 4 * (_bw + _gap)
    arrow_d(s, last_bx + _bw / 2, top_y + _bh, C_MINT)

    # 下段 5ボックス（右→左）
    bot_nodes = [
        (C_FILL_MINT,  C_MINT,   "SST\n（超高速AT）"),
        (C_FILL_GOLD,  C_GOLD,   "ジョーカー\nモード"),
        (C_FILL_GOLD,  C_GOLD,   "金BAR\n昇格"),
        (C_FILL_RED,   C_RED,    "REG BONUS\n（シナリオ+1）"),
        (C_FILL_GREEN, C_GREEN,  "BIG再当選\n→ループ"),
    ]
    bot_y = top_y + _bh + Inches(0.42)
    bot_xs = [_mx + (4 - i) * (_bw + _gap) for i in range(5)]
    for i, (fill, col, label) in enumerate(bot_nodes):
        bx = bot_xs[i]
        rect_b(s, bx, bot_y, _bw, _bh, fill, col, 1.5)
        tb(s, bx, bot_y, _bw, _bh, label, 8.5, bold=True, color=col,
           align=PP_ALIGN.CENTER)
        if i < len(bot_nodes) - 1:
            # 左向き矢印
            shp = s.shapes.add_shape(14, bot_xs[i + 1] + _bw, bot_y + _bh / 2 - Emu(90000),
                                     Emu(200000), Emu(180000))
            shp.fill.solid()
            shp.fill.fore_color.rgb = col
            shp.line.fill.background()

    # 3パネル要約
    summaries = [
        (C_GREEN,  C_FILL_GREEN, "SBJループ",
         "BIG→RC高確→RC→ST\n→BIG のループが出玉の核心"),
        (C_PINK,   C_FILL_PINK,  "RC 2段構え",
         "前半：カード昇格期待\n後半：報酬確定の2ドキドキ"),
        (C_GOLD,   C_FILL_GOLD,  "ジョーカーモード",
         "RC3回→金BAR→SST頻発\nエンディングから確定突入"),
    ]
    sy = bot_y + _bh + Inches(0.28)
    sw = Inches(3.05)
    for i, (col, fill, ttl, desc) in enumerate(summaries):
        sx = Inches(0.22 + i * 3.22)
        rect_b(s, sx, sy, sw, Inches(1.02), fill, col, 1.5)
        tb(s, sx + Emu(50000), sy + Emu(40000), sw - Emu(70000), Emu(260000),
           ttl, 9.5, bold=True, color=col)
        tb(s, sx + Emu(50000), sy + Emu(300000), sw - Emu(70000), Emu(520000),
           desc, 7.5, color=C_DARK_TEXT)

    footer(s,
           "BIG→ループ型連チャン。REGはシナリオカウンター（3回→金BAR昇格）",
           "差枚数2000枚完走でエンディング移行→ジョーカーモード確定。有利区間管理が鍵")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 初当りへの3ルート", pg=f"3 / {TOTAL_SLIDES}")

    routes = [
        (C_GREEN, C_FILL_GREEN, "①チャンス役ルート",
         "スイカ・チェリーなどのチャンス役で\nボーナス高確への移行を抽選\n"
         "スイカは天井カウンターにもなる",
         "最もオーソドックスな突入ルート"),
        (C_RED,   C_FILL_RED,   "②小役連続ルート",
         "通常時に小役が複数回連続すると\n高確率状態（ボーナス高確）に移行\n"
         "設定が高いほど移行率がUP",
         "設定判別の重要指標でもある"),
        (C_GOLD,  C_FILL_GOLD,  "③ゲーム数天井ルート",
         "BIG間999G消化でBIG確定\nスイカ規定回数到達でSTに直撃\n"
         "スイカ天井が狙い目の核心",
         "低設定でも拾える狙い目ルート"),
    ]
    for i, (col, fill, ttl, body, note) in enumerate(routes):
        rx = Inches(0.22 + i * 3.25)
        ry = Inches(0.78)
        rect_b(s, rx, ry, Inches(3.0), Inches(3.25), fill, col, 1.5)
        rect(s, rx, ry, Inches(3.0), Emu(58000), col)
        tb(s, rx + Emu(40000), ry + Emu(70000), Inches(2.7), Emu(270000),
           ttl, 10, bold=True, color=col)
        tb(s, rx + Emu(40000), ry + Emu(340000), Inches(2.7), Emu(1150000),
           body, 8.5, color=C_DARK_TEXT)
        rect_b(s, rx + Emu(40000), ry + Emu(1560000), Inches(2.7), Emu(330000),
               C_FILL_DEFAULT, col, 0.8)
        tb(s, rx + Emu(70000), ry + Emu(1580000), Inches(2.6), Emu(310000),
           note, 7.5, color=C_MID_GRAY)

    # BIG BONUS 2種類
    rect_b(s, Inches(0.22), Inches(4.17), Inches(9.6), Emu(680000),
           C_FILL_GREEN, C_GREEN, 1.0)
    tb(s, Inches(0.35), Inches(4.20), Inches(9.2), Emu(260000),
       "BIG BONUS 2種類", 10, bold=True, color=C_GREEN)
    tb(s, Inches(0.35), Inches(4.48), Inches(4.4), Emu(310000),
       "青7揃い BIG（300枚払い出し）→ RC高確確定・SBJループ突入", 8.5, color=C_DARK_TEXT)
    tb(s, Inches(4.9), Inches(4.48), Inches(4.9), Emu(310000),
       "赤7揃い BIG（100枚払い出し）→ RC高確移行・ループ期待度やや低", 8.5, color=C_RED)

    footer(s,
           "通常時は「いかに早くBIGを引くか」に集約。スイカ天井はAT直撃（ST）で中間狙いも有効",
           "高確移行ルートを複数設けることで単調さを回避。設定判別は小役連続率が重要指標")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: リオチャンス（RC）の仕組み ── 2段構えの体験設計
# ══════════════════════════════════════════════════════════════
def s_bjmechanism(prs):
    s = new_slide(prs)
    hdr(s, "リオチャンス（RC）の仕組み ── 2段構えの体験設計", pg=f"4 / {TOTAL_SLIDES}")

    # 前半パート（左）
    rect_b(s, Inches(0.22), Inches(0.75), Inches(4.55), Inches(2.75),
           C_FILL_PINK, C_PINK, 2.0)
    rect(s, Inches(0.22), Inches(0.75), Inches(4.55), Emu(58000), C_PINK)
    tb(s, Inches(0.30), Inches(0.82), Inches(4.3), Emu(270000),
       "前半パート ── カード昇格抽選", 10.5, bold=True, color=C_PINK)
    tb(s, Inches(0.30), Inches(1.20), Inches(4.3), Emu(1550000),
       "カードをめくる演出でブラックジャック体験を再現。\n\n"
       "● 押し順ベル約1/2で昇格チャンス発生\n"
       "● レア役（スイカ・チェリー）→ 超昇格の大チャンス\n"
       "● シナリオ14種 → 最終的な報酬グレードを決定\n\n"
       "昇格: REG → エピソードBONUS → BIG の3段階",
       8.5, color=C_DARK_TEXT)

    # 前半→後半の矢印
    arrow_r(s, Inches(4.77), Inches(2.125), C_PINK)

    # 後半パート（右）
    rect_b(s, Inches(5.1), Inches(0.75), Inches(4.65), Inches(2.75),
           C_FILL_RED, C_RED, 2.0)
    rect(s, Inches(5.1), Inches(0.75), Inches(4.65), Emu(58000), C_RED)
    tb(s, Inches(5.2), Inches(0.82), Inches(4.3), Emu(270000),
       "後半パート ── ジャッジ（報酬確定）", 10.5, bold=True, color=C_RED)
    tb(s, Inches(5.2), Inches(1.20), Inches(4.3), Emu(1550000),
       "「狙えカットイン」のタイミングで報酬が確定する。\n\n"
       "● 狙え役を外す → 1段階さらに昇格（逆転演出！）\n"
       "● 周回ごとに報酬グレードが上がる仕組み\n"
       "● 報酬種別 → ST突入・SST突入・BIG確定 に分岐\n\n"
       "「外れた＝嬉しい」という逆転心理設計が秀逸",
       8.5, color=C_DARK_TEXT)

    # 下段：ST / SST説明
    st_data = [
        (C_MINT, C_FILL_MINT, "ストックタイム（ST）",
         "50G〜777G継続　リプ・レア役でRC抽選\n平均ストック数: 約2個　50Gごとに放出"),
        (C_GOLD, C_FILL_GOLD, "スーパーストックタイム（SST）",
         "100G+α　純増約5.1枚/G　500枚超可能\nRC大量ストック→放出で爆発力抜群"),
        (C_GREEN, C_FILL_GREEN, "RC3回→金BAR昇格",
         "REGをRCで3回引くと金BARシナリオへ\n4回目のREG→ジョーカーモード大チャンス"),
    ]
    for i, (col, fill, ttl, desc) in enumerate(st_data):
        sx = Inches(0.22 + i * 3.25)
        sy = Inches(3.65)
        rect_b(s, sx, sy, Inches(3.0), Inches(1.22), fill, col, 1.5)
        tb(s, sx + Emu(50000), sy + Emu(40000), Inches(2.7), Emu(270000),
           ttl, 9, bold=True, color=col)
        tb(s, sx + Emu(50000), sy + Emu(310000), Inches(2.7), Emu(620000),
           desc, 8, color=C_DARK_TEXT)

    footer(s,
           "「カードをめくって昇格確認」→「外れたら逆転」の2段ドキドキ構造が核心",
           "前半で昇格を期待、後半で確定確認。1回のRCに2つの感情ピークを設計した稀有な構造")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス ── 出玉が伸びる仕組み
# ══════════════════════════════════════════════════════════════
def s_bonus(prs):
    s = new_slide(prs)
    hdr(s, "AT/ボーナス ── 出玉が伸びる仕組み", pg=f"5 / {TOTAL_SLIDES}")

    # SBJループ図
    flow = [
        (C_GOLD,  C_FILL_GOLD,  "BIG BONUS\n青7/赤7"),
        (C_GREEN, C_FILL_GREEN, "RC高確\n（ループ土台）"),
        (C_PINK,  C_FILL_PINK,  "リオチャンス\n（RC）"),
        (C_MINT,  C_FILL_MINT,  "ストックタイム\n（ST/SST）"),
    ]
    fxs = [Inches(0.3), Inches(2.72), Inches(5.14), Inches(7.56)]
    for i, (col, fill, label) in enumerate(flow):
        rect_b(s, fxs[i], Inches(0.78), Inches(2.2), Inches(1.45), fill, col, 2.0)
        tb(s, fxs[i], Inches(0.78), Inches(2.2), Inches(1.45),
           label, 10, bold=True, color=col, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            arrow_r(s, fxs[i] + Inches(2.2), Inches(0.78) + Inches(0.725), col)

    tb(s, Inches(0.4), Inches(2.45), Inches(9.2), Emu(220000),
       "↑ RCをストックして放出 → BIG当選 → RC高確継続 → SBJループ繰り返し ↑", 9,
       bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # 出玉ポイント3点
    pts = [
        (C_GREEN, C_FILL_GREEN, "RC高確維持が最重要",
         "BIG後は必ずRC高確移行。\nRC高確中のボーナス当選でRC突入。\nRC高確が続く限りループは継続する。"),
        (C_RED,   C_FILL_RED,   "RCでの報酬昇格を狙う",
         "RC前半のカード演出でBIG昇格が理想。\nBIG獲得→ST移行→大量RC放出で爆発。\nREGはシナリオカウンター（3回→金BAR）。"),
        (C_GOLD,  C_FILL_GOLD,  "SSTで一気に積む",
         "SSTは純増5.1枚×100G以上で500枚超。\nジョーカーモード中はSST頻発。\nエンディングまで駆け抜ければ最高。"),
    ]
    for i, (col, fill, ttl, body) in enumerate(pts):
        rx = Inches(0.22 + i * 3.25)
        rect_b(s, rx, Inches(2.80), Inches(3.0), Inches(2.0), fill, col, 1.5)
        tb(s, rx + Emu(50000), Inches(2.88), Inches(2.7), Emu(280000),
           ttl, 10, bold=True, color=col)
        tb(s, rx + Emu(50000), Inches(3.22), Inches(2.7), Emu(1250000),
           body, 8, color=C_DARK_TEXT)

    footer(s,
           "「ST中にRCをどれだけストックできるか」が1セットの出玉規模を決める核心変数",
           "純増5.1枚のSSTを何度引けるかが連チャンの深さを左右")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: ジョーカーモード ── 上位AT
# ══════════════════════════════════════════════════════════════
def s_joker(prs):
    s = new_slide(prs)
    hdr(s, "上位ATへの道 ── ジョーカーモード", pg=f"6 / {TOTAL_SLIDES}")

    rect_b(s, Inches(0.22), Inches(0.75), Inches(9.6), Inches(1.05),
           C_FILL_GOLD, C_GOLD, 2.0)
    tb(s, Inches(0.40), Inches(0.82), Inches(9.2), Emu(320000),
       "ジョーカーモード：RC高確のロングバージョン。滞在中はSST当選率が大幅UP。ジョーカーランプ（筐体上部）が点灯で滞在濃厚",
       9.5, bold=True, color=C_GOLD)

    tb(s, Inches(0.3), Inches(1.92), Inches(9.0), Emu(280000),
       "ジョーカーモード 突入条件", 11, bold=True, color=C_GOLD)

    conditions = [
        (C_GREEN, C_FILL_GREEN, "①エンディング後",
         "エンディングボーナス（差枚数2000枚完走）終了後\n→ ジョーカーモード確定移行"),
        (C_RED,   C_FILL_RED,   "②スペシャルエピソードBONUS",
         "RC経由のREGを3回引いて金BARシナリオ昇格後\n4回目のREG（スペシャルエピソード発生）"),
        (C_GOLD,  C_FILL_GOLD,  "③有利区間リセット時",
         "有利区間リセット（設定変更除く）\n→ 必ずジョーカーモード突入"),
        (C_MINT,  C_FILL_MINT,  "④その他の条件",
         "金BAR獲得 / RC高確中に差枚2000枚でRC当選\n/ 有利区間差枚1800枚でRC当選"),
    ]
    for i, (col, fill, ttl, body) in enumerate(conditions):
        rx = Inches(0.22 + (i % 2) * 4.85)
        ry = Inches(2.32 + (i // 2) * 1.35)
        rect_b(s, rx, ry, Inches(4.6), Inches(1.20), fill, col, 1.5)
        rect(s, rx, ry, Emu(55000), Inches(1.20), col)
        tb(s, rx + Emu(70000), ry + Emu(50000), Inches(4.2), Emu(260000),
           ttl, 10, bold=True, color=col)
        tb(s, rx + Emu(70000), ry + Emu(320000), Inches(4.2), Emu(640000),
           body, 8, color=C_DARK_TEXT)

    footer(s,
           "ジョーカーモードはSST頻発モード。RC3回→金BAR→スペシャルエピソードが最短ルート",
           "ランプ消灯でも滞在継続の可能性あり。ジョーカーモード中REG引くと次RCに必ずSST出現")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── カードゲーム×女の子×パチスロの融合がなぜ機能するか",
        pg=f"7 / {TOTAL_SLIDES}")

    rect_b(s, Inches(9.0), Inches(0.0), Inches(1.0), Inches(0.60),
           C_RED, None, 0)
    tb(s, Inches(9.0), Emu(20000), Inches(1.0), Emu(340000),
       "Part B", 8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    design_pts = [
        (C_PINK,  C_FILL_PINK,  "キャラクター演出設計",
         "リオ・ミント・リナが各フェーズを担当。\n"
         "「誰が出るか」でゲーム状態を直感的に判断できる。\n"
         "華やかさとゲーム性が一体化した演出設計。"),
        (C_RED,   C_FILL_RED,   "2段階ドキドキ設計",
         "前半（昇格期待） × 後半（報酬確定）の2パートで\n"
         "1回のRCに2つの感情ピークを設計。\n"
         "「外れ＝逆転昇格」の心理的反転で常に期待が続く。"),
        (C_GREEN, C_FILL_GREEN, "ループ × カウンター設計",
         "SBJループ（連鎖）× REGカウンター（蓄積）。\n"
         "失敗が次回への期待に変換される設計。\n"
         "悔しさ・損失を「次は成功する」動機に変換。"),
        (C_GOLD,  C_FILL_GOLD,  "4号機リメイクの情緒価値",
         "原作プレイヤーへのノスタルジー。\n"
         "リオチャンス・ストックタイムの名称を踏襲し親近感。\n"
         "新規層には新鮮、経験者には懐かしい二重訴求。"),
    ]
    for i, (col, fill, ttl, body) in enumerate(design_pts):
        rx = Inches(0.22 + (i % 2) * 4.85)
        ry = Inches(0.85 + (i // 2) * 1.90)
        rect_b(s, rx, ry, Inches(4.6), Inches(1.70), fill, col, 1.5)
        rect(s, rx, ry, Emu(55000), Inches(1.70), col)
        tb(s, rx + Emu(70000), ry + Emu(55000), Inches(4.2), Emu(280000),
           ttl, 10, bold=True, color=col)
        tb(s, rx + Emu(70000), ry + Emu(340000), Inches(4.2), Emu(960000),
           body, 8, color=C_DARK_TEXT)

    footer(s,
           "「カジノのダーク」ではなく「女の子が活躍する華やか台」としての世界観構築が差別化の核",
           "BJ演出の知的緊張感 × キャラの明るさ × ループの高揚感を三位一体で設計")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題", pg=f"8 / {TOTAL_SLIDES}")

    rect_b(s, Inches(0.22), Inches(0.75), Inches(4.65), Inches(4.1),
           C_FILL_GREEN, C_GREEN, 2.0)
    rect(s, Inches(0.22), Inches(0.75), Inches(4.65), Emu(65000), C_GREEN)
    tb(s, Inches(0.35), Inches(0.82), Inches(4.3), Emu(280000),
       "良い点（Pros）", 11, bold=True, color=C_GREEN)
    pros = [
        "● BJ演出＋女の子キャラで「見ていて楽しい」台",
        "● RC2段構えで1回の当選体験の密度が高い",
        "● SBJループの連鎖性が高く「次もある」期待を持続",
        "● ジョーカーモードという明確な上位状態で目標設定しやすい",
        "● 4号機リメイクによる情緒価値・二重訴求（新旧ユーザー）",
        "● 天井・スイカ天井で狙い目が明確。立ち回りしやすさ◎",
        "● スペック（機械割）は設定5で110%・高設定は十分な性能",
    ]
    for i, txt in enumerate(pros):
        tb(s, Inches(0.35), Inches(1.26) + Emu(i * 390000), Inches(4.3), Emu(360000),
           txt, 8, color=C_DARK_TEXT)

    rect_b(s, Inches(5.1), Inches(0.75), Inches(4.72), Inches(4.1),
           C_FILL_RED, C_RED, 2.0)
    rect(s, Inches(5.1), Inches(0.75), Inches(4.72), Emu(65000), C_RED)
    tb(s, Inches(5.23), Inches(0.82), Inches(4.3), Emu(280000),
       "課題（Cons）", 11, bold=True, color=C_RED)
    cons = [
        "● 低設定は初当り確率が重く「何もできない」体験になりやすい",
        "● BIGを引けないとループが始まらない完全ボーナス依存設計",
        "● RC昇格なしのREGは出玉貢献が低く、連続で引くと消化感",
        "● SBJループは理解するまでゲームフローが複雑に見える",
        "● ジョーカーモードの認識にランプ依存（演出上の分かりにくさ）",
        "● BJ要素が「演出」止まりで実際の戦略性はない",
    ]
    for i, txt in enumerate(cons):
        tb(s, Inches(5.23), Inches(1.26) + Emu(i * 430000), Inches(4.3), Emu(400000),
           txt, 8, color=C_DARK_TEXT)

    footer(s,
           "高設定での爆発力・BJ演出の新鮮さはBRONZE受賞に値するが、低設定での体験改善が普及の鍵",
           "「女の子が活躍する」世界観の認知が広まれば若年層・女性層への訴求力も上がる")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_summary(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", pg=f"9 / {TOTAL_SLIDES}")

    rect_b(s, Inches(0.22), Inches(0.75), Inches(9.6), Inches(0.90),
           C_FILL_GOLD, C_GOLD, 2.0)
    tb(s, Inches(0.40), Inches(0.82), Inches(9.2), Emu(550000),
       "スマスロ SBJは「他ジャンルのゲーム性をパチスロに翻訳した実践例」として\n"
       "設計の教科書的価値を持つ。パチスロアワード2025 BRONZE受賞もその証左。",
       9.5, bold=True, color=C_GOLD)

    learnings = [
        (C_PINK,  C_FILL_PINK,  "設計学び①\nキャラ設計",
         "リオ・ミント・リナで状態を視覚化。\n"
         "「誰が出るか＝何が起きるか」の直感設計。\n"
         "世界観とゲーム性を分離せず一体化した好例。\n"
         "「女の子台」という入口で間口を広げた戦略"),
        (C_RED,   C_FILL_RED,   "設計学び②\n感情ピーク2段化",
         "前半（昇格期待）× 後半（報酬確定）の2段階。\n"
         "「外れ＝逆転昇格」の心理反転も組み込み。\n"
         "1回の当選に複数の感情ピークを設計する手法は\n"
         "様々なゲームジャンルに応用可能"),
        (C_GREEN, C_FILL_GREEN, "設計学び③\nループ×カウンター",
         "SBJループ（連鎖）× REGカウンター（蓄積）。\n"
         "「失敗が次回への期待に変換される」設計は\n"
         "プレイヤーのモチベーション維持に直結。\n"
         "悔しさ・損失を動機に変える心理工学"),
    ]
    for i, (col, fill, ttl, body) in enumerate(learnings):
        rx = Inches(0.22 + i * 3.25)
        ry = Inches(1.75)
        rect_b(s, rx, ry, Inches(3.0), Inches(2.65), fill, col, 2.0)
        rect(s, rx, ry, Inches(3.0), Emu(55000), col)
        tb(s, rx + Emu(50000), ry + Emu(65000), Inches(2.7), Emu(330000),
           ttl, 9, bold=True, color=col)
        tb(s, rx + Emu(50000), ry + Emu(400000), Inches(2.7), Emu(1680000),
           body, 7.5, color=C_DARK_TEXT)

    rect_b(s, Inches(0.22), Inches(4.57), Inches(9.6), Emu(340000),
           C_FILL_GREEN, C_GREEN, 1.5)
    tb(s, Inches(0.40), Inches(4.61), Inches(9.2), Emu(300000),
       "総合：華やかキャラ×BJ演出×ループ爆発力の三位一体。"
       "カジノの重い印象を女の子の活躍で払拭し、幅広いプレイヤーに訴求した意欲作。",
       8.5, bold=True, color=C_GREEN)

    footer(s,
           "「キャラ設計×感情2段化×ループカウンター」の3軸設計がSBJの本質",
           "4号機リメイクという情緒的入口と、スマスロの現代設計の両立が2025年のBRONZE評価を獲得した要因")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Slide 1: タイトル・キャラクター・3ポイント")
    s_title(prs)
    print("Slide 2: ゲームフロー全体図")
    s_gameflow(prs)
    print("Slide 3: 通常時の遊び方")
    s_normal(prs)
    print("Slide 4: リオチャンス 2段構えの仕組み")
    s_bjmechanism(prs)
    print("Slide 5: AT/ボーナス 出玉の伸ばし方")
    s_bonus(prs)
    print("Slide 6: ジョーカーモード")
    s_joker(prs)
    print("Slide 7: 面白さの設計")
    s_design(prs)
    print("Slide 8: 良い点と課題")
    s_pros_cons(prs)
    print("Slide 9: まとめ")
    s_summary(prs)

    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}")


if __name__ == "__main__":
    main()
