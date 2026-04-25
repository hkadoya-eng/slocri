"""
ゲーム性提案書 MD → PowerPoint 変換
使い方: python make_pptx.py proposals/proposal_xxx.md
"""
import sys
import os
import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

LIBRARY_PATH = os.path.join(os.path.dirname(__file__), "src", "gameDesignLibrary.json")

C_ORANGE = RGBColor(0xD8, 0x5A, 0x30)
C_DARK   = RGBColor(0x22, 0x22, 0x22)
C_MID    = RGBColor(0x66, 0x66, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BG     = RGBColor(0x1A, 0x1A, 0x2E)
C_CODE   = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "游ゴシック"
FONT_MONO = "Consolas"


def sf(run, size, bold=False, color=None, mono=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT_MONO if mono else FONT
    if color:
        run.font.color.rgb = color


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_tb(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def title_slide(prs, title, catch):
    slide = blank_slide(prs)
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, 0.18, 7.5, C_ORANGE)

    tf = add_tb(slide, 0.5, 2.0, 9.0, 1.8)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    sf(r, 30, bold=True, color=C_WHITE)

    if catch:
        tf2 = add_tb(slide, 0.5, 3.9, 9.0, 0.9)
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = catch
        sf(r2, 17, color=C_ORANGE)

    tf3 = add_tb(slide, 0.5, 6.85, 9.0, 0.4)
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "スロキー / ゲーム性提案書"
    sf(r3, 11, color=C_MID)


def matrix_slide(prs, lib):
    """既存機種との設計比較マトリックス"""
    slide = blank_slide(prs)
    set_bg(slide, C_WHITE)
    add_rect(slide, 0, 0, 10, 1.05, C_ORANGE)

    tf = add_tb(slide, 0.4, 0.12, 9.2, 0.8)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "既存機種 設計パターン比較"
    sf(r, 20, bold=True, color=C_WHITE)

    # テーブルデータ構築
    patterns = lib.get("gameFlowPatterns", {})
    rows = []
    for pat_name, pat_data in patterns.items():
        for ex in pat_data.get("examples", []):
            rows.append({
                "machine": ex["machine"],
                "pattern": pat_name,
                "emotion": pat_data.get("playerEmotion", ""),
                "strength": pat_data.get("strengths", [""])[0],
                "weakness": pat_data.get("weaknesses", [""])[0],
            })
    # 名機も追加
    for m_name, m_data in lib.get("classicMachines", {}).items():
        rows.append({
            "machine": m_name,
            "pattern": m_data.get("designPattern", ""),
            "emotion": m_data.get("playerEmotion", ""),
            "strength": m_data.get("designLesson", "")[:30],
            "weakness": "",
        })

    if not rows:
        return

    headers = ["機種名", "設計パターン", "プレイヤー感情", "強み"]
    col_widths = [2.5, 2.0, 2.5, 2.5]
    row_height = 0.38
    table_top = 1.2
    max_rows = min(len(rows), 13)

    # ヘッダー行
    x = 0.3
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, table_top, w - 0.02, row_height, RGBColor(0x22, 0x22, 0x22))
        tf2 = add_tb(slide, x + 0.05, table_top + 0.05, w - 0.12, row_height - 0.1)
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = h
        sf(r2, 9, bold=True, color=C_WHITE)
        x += w

    # データ行
    for ri, row in enumerate(rows[:max_rows]):
        y = table_top + row_height * (ri + 1)
        bg_color = RGBColor(0xF8, 0xF4, 0xF1) if ri % 2 == 0 else C_WHITE
        x = 0.3
        for ci, (key, w) in enumerate(zip(["machine", "pattern", "emotion", "strength"], col_widths)):
            add_rect(slide, x, y, w - 0.02, row_height, bg_color)
            val = str(row.get(key, ""))[:28]
            tf3 = add_tb(slide, x + 0.05, y + 0.04, w - 0.12, row_height - 0.06)
            p3 = tf3.paragraphs[0]
            r3 = p3.add_run()
            r3.text = val
            sf(r3, 8, color=C_DARK if ci > 0 else C_ORANGE, bold=(ci == 0))
            x += w


def content_slide(prs, heading, lines, code_block=False):
    slide = blank_slide(prs)
    set_bg(slide, C_WHITE)
    add_rect(slide, 0, 0, 10, 1.05, C_ORANGE)

    tf = add_tb(slide, 0.4, 0.12, 9.2, 0.8)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = heading
    sf(r, 20, bold=True, color=C_WHITE)

    body_t = 1.15
    body_h = 6.1

    if code_block:
        add_rect(slide, 0.35, body_t, 9.3, body_h, C_CODE)
        tf2 = add_tb(slide, 0.5, body_t + 0.1, 9.0, body_h - 0.15)
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = "\n".join(lines)
        sf(r2, 10, mono=True, color=C_DARK)
    else:
        tf2 = add_tb(slide, 0.4, body_t, 9.2, body_h)
        first = True
        for line in lines:
            stripped = line.strip()
            if first:
                p2 = tf2.paragraphs[0]
                first = False
            else:
                p2 = tf2.add_paragraph()

            if not stripped:
                p2.space_after = Pt(3)
                continue

            p2.space_before = Pt(1)
            parts = re.split(r"(\*\*[^*]+\*\*)", stripped)
            is_bullet = stripped.startswith(("- ", "• ")) or re.match(r"^\d+\.\s", stripped)

            for j, part in enumerate(parts):
                run = p2.add_run()
                if part.startswith("**") and part.endswith("**"):
                    run.text = part[2:-2]
                    sf(run, 13, bold=True, color=C_ORANGE)
                else:
                    if j == 0 and is_bullet:
                        text = re.sub(r"^[-•]\s+", "• ", part)
                        text = re.sub(r"^\d+\.\s+", lambda m: m.group(), text)
                        run.text = text
                    else:
                        run.text = part
                    sf(run, 13, color=C_DARK)


def parse_md(path):
    with open(path, encoding="utf-8") as f:
        raw = f.read()

    raw = re.sub(r"<!--.*?-->", "", raw, flags=re.DOTALL).strip()

    title_m = re.search(r"^# (.+)$", raw, re.MULTILINE)
    title = title_m.group(1).strip() if title_m else "新機種ゲーム性提案書"

    catch_m = re.search(r"\*\*キャッチ[：:]\s*「([^」]+)」\*\*", raw)
    catch = f"「{catch_m.group(1)}」" if catch_m else ""

    sections = re.split(r"^## ", raw, flags=re.MULTILINE)[1:]
    result = []
    for sec in sections:
        lines = sec.strip().splitlines()
        heading = lines[0].strip()
        body = lines[1:]
        body = [l for l in body if l.strip() != "---"]
        result.append((heading, body))

    return title, catch, result


def make_pptx(md_path):
    title, catch, sections = parse_md(md_path)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs, title, catch)

    # 比較マトリックススライド
    try:
        with open(LIBRARY_PATH, encoding="utf-8") as f:
            lib = json.load(f)
        matrix_slide(prs, lib)
    except Exception:
        pass

    for heading, body in sections:
        in_code = False
        code_lines = []
        normal_lines = []

        for line in body:
            if line.strip().startswith("```"):
                in_code = not in_code
                continue
            if in_code:
                code_lines.append(line)
            else:
                normal_lines.append(line)

        if code_lines:
            content_slide(prs, heading, normal_lines + [""], code_block=False)
            content_slide(prs, heading + "（フロー図）", code_lines, code_block=True)
        else:
            content_slide(prs, heading, normal_lines)

    out_path = os.path.splitext(md_path)[0] + ".pptx"
    prs.save(out_path)
    print(f"[OK] 保存しました: {out_path}")
    return out_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        files = [f for f in os.listdir("proposals") if f.endswith(".md")]
        if not files:
            print("使い方: python make_pptx.py proposals/proposal_xxx.md")
            sys.exit(1)
        files.sort(reverse=True)
        path = os.path.join("proposals", files[0])
        print(f"[INFO] 最新ファイルを使用: {path}")
    else:
        path = sys.argv[1]

    make_pptx(path)
