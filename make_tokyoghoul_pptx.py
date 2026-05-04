"""
スマスロ L東京喰種 機種説明＋分析 統合版資料 v1
（スパイキー・クロスアルファ / パチスロアワード2025 GOLD受賞）
出力: proposals/機種分析/東京喰種/tokyoghoul_guide_v1.pptx
テーマ: 深黒 × 紫(#8833CC) × 赤(#CC2222) × 金（喰種世界観）

Part A: 説明パート（プレイヤー視点）  スライド1〜6
Part B: 分析パート                    スライド7〜9
"""
import io
import os
import sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "東京喰種", "tokyoghoul_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深黒×紫×赤×金 ── 喰種世界観）──────────────
C_BG    = RGBColor(0x06, 0x02, 0x10)   # 深黒（喰種の夜）
C_CARD  = RGBColor(0x0E, 0x06, 0x1C)   # カード背景
C_CARD2 = RGBColor(0x16, 0x0A, 0x28)   # カード背景2
C_ROW   = RGBColor(0x12, 0x08, 0x22)   # テーブル奇数行
C_PUR   = RGBColor(0x88, 0x33, 0xCC)   # 紫（喰種メインカラー）
C_PUR2  = RGBColor(0xAA, 0x55, 0xEE)   # 明るい紫
C_RED   = RGBColor(0xCC, 0x22, 0x22)   # 赤（赫子・赫眼）
C_RED2  = RGBColor(0xFF, 0x44, 0x44)   # 明るい赤
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金（アワード金色）
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)   # 輝く金
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)   # オフホワイト
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)   # クリーム
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)   # グレー
C_LTGRY = RGBColor(0x44, 0x44, 0x66)   # ライトグレー
C_GREEN = RGBColor(0x22, 0xCC, 0x66)   # 緑

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

TOTAL_SLIDES = 9


# ── 背景生成（暗黒×紫グロー ── 喰種世界観）───────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (6, 2, 16))
    draw = ImageDraw.Draw(img)
    # 斜めライン（格子模様・暗闇感）
    for i in range(0, w + h, 90):
        draw.line([(i, 0), (0, i)], fill=(10, 4, 22), width=1)
    # 右下の紫グロー（赫子エフェクト）
    for y in range(h - 130, h):
        t = (y - (h - 130)) / 130
        draw.line([(0, y), (w, y)], fill=(int(40 * t), 0, int(60 * t)))
    # 左上の赤グロー（赫眼エフェクト）
    for y in range(0, 50):
        t = (50 - y) / 50 * 0.3
        draw.line([(0, y), (w // 3, y)], fill=(int(30 * t), 0, int(5 * t)))
    # 左端アクセントライン（紫）
    for x in range(0, 5):
        draw.line([(x, 0), (x, h)], fill=(0x88, 0x22, 0xCC))
    # 右端アクセントライン（赤）
    for x in range(w - 4, w):
        draw.line([(x, 0), (x, h)], fill=(0x99, 0x10, 0x10))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


# ── ヘルパー関数 ───────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    """塗りつぶし矩形（ボーダーなし）"""
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    """塗りつぶし矩形（ボーダーあり）"""
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    """テキストボックス"""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def hdr(slide, title_text, pg=""):
    """スライドヘッダー（紫ライン＋タイトル）"""
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_PUR)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 13, bold=True, color=C_RED2, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_PUR)


def net_note(slide):
    """右下の※ネット解析情報より"""
    tb(slide, Inches(8.0), Inches(5.35), Inches(1.85), Emu(200000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, design_comment, sub_text=""):
    """フッター（設計コメント＋補足）"""
    fy = Inches(5.05)
    rect(slide, 0, fy, SLIDE_W, Inches(0.55), RGBColor(0x06, 0x02, 0x14))
    tb(slide, Inches(0.2), fy + Emu(30000), Inches(7.0), Emu(380000),
       "【設計】" + design_comment, 7.5, bold=True, color=C_GOLD)
    if sub_text:
        tb(slide, Inches(0.2), fy + Emu(230000), Inches(7.5), Emu(200000),
           sub_text, 6.5, color=C_GRAY)
    net_note(slide)


def arrow_r(slide, x, cy, col=None):
    """右向き矢印"""
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    """下向き矢印"""
    shp = slide.shapes.add_shape(17, cx - Emu(90000), y, Emu(180000), Emu(200000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_PUR
    shp.line.fill.background()


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    """Part A – スライド1: タイトル・スペック・3ポイント"""
    s = new_slide(prs)

    # 左パネル（タイトル領域）
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x04, 0x00, 0x0E))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_PUR)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_RED)

    # PartAバッジ
    rect(s, Inches(0.22), Inches(0.2), Inches(1.4), Emu(260000), C_PUR)
    tb(s, Inches(0.22), Inches(0.2), Inches(1.4), Emu(260000),
       "Part A 説明編", 7.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # アワードバッジ
    rect(s, Inches(0.22), Inches(0.6), Inches(2.5), Emu(260000), C_GOLD)
    tb(s, Inches(0.22), Inches(0.6), Inches(2.5), Emu(260000),
       "パチスロアワード2025 GOLD受賞", 8, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    tb(s, Inches(0.22), Inches(1.05), Inches(5.1), Emu(900000),
       "L 東京喰種", 34, bold=True, color=C_RED2, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.85), Inches(5.0), Emu(300000),
       "スマスロ / スパイキー・クロスアルファ", 9.5, color=C_CREAM, font=FONT_H)
    tb(s, Inches(0.22), Inches(3.2), Inches(5.0), Emu(260000),
       "── 差枚数管理AT×喰種対決×裏AT 三層構造の傑作", 9, color=C_GRAY, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.7), Inches(4.9), Emu(230000),
       "設定: 1〜6段階　　AT純増: 約4.0枚/G（差枚数管理）", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(4.0), Inches(4.9), Emu(230000),
       "裏AT純増: 約5.0枚/G　　期待枚数（裏AT込）: 約3430枚", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(4.3), Inches(4.9), Emu(230000),
       "機械割: 設定1=97.5%〜設定6=114.9%　　導入: 2024年", 8.5, color=C_GRAY)

    # 右：この台の3ポイント
    kws = [
        (C_PUR,  "差枚数管理AT「東京喰種咬」",
         "純増4.0枚/G・初期差枚150枚\nレア役→差枚上乗せ or 喰種対決抽選\n特化ゾーンで一気に積み増し"),
        (C_RED,  "喰種対決→BITES→特化ゾーン",
         "バトル勝利でBITES（4桁上乗せ期待）\n「百足覚醒」(期待約300枚)「隻眼の梟」(約500枚)\n上乗せ連打で爆発力が生まれる"),
        (C_GOLD, "裏AT（上位AT）で一撃3430枚",
         "有馬貴将ジャッジメント成功→裏AT\n純増5.0枚・喰種対決成功率が大幅UP\nCCG死神で300〜2000枚上乗せ確定"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.3 + i * 1.65)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.4), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.4), ac)
        tb(s, Inches(5.85), y0 + Emu(60000), Inches(3.8), Emu(300000),
           kw, 11, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(360000), Inches(3.8), Emu(500000),
           desc, 8, color=C_WHITE)

    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    """Part A – スライド2: ゲームフロー全体図"""
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常→CZ→AT→喰種対決→BITES→裏AT", f"2/{TOTAL_SLIDES}")

    # ── 上段（左→右）: 通常時 → CZ → EP BONUS → AT
    row1_y = Inches(0.78)
    row1_h = Emu(1500000)
    boxes1 = [
        (C_CARD2, C_PUR,  "通常時",
         "レア役・規定G消化\n招待状50G周期\nCZ/EP BONUS目指す"),
        (C_CARD2, C_PUR,  "CZ\n追憶/大喰種討伐",
         "2種類のCZ\n大喰種討伐が高期待\n突破でAT確定"),
        (C_CARD2, C_PUR2, "EP BONUS",
         "EP BONUS当選で\nAT直行ルート\nレア役から直撃も"),
        (C_CARD2, C_RED,  "AT\n東京喰種咬",
         "純増4.0枚/G\n差枚管理・初期150枚\n喰種対決を目指す"),
    ]
    bw1 = Inches(2.0)
    gap1 = Inches(0.28)
    sx1 = Inches(0.3)
    for i, (fill, bc, lbl, sub) in enumerate(boxes1):
        bx0 = sx1 + i * (bw1 + gap1)
        rect_b(s, bx0, row1_y, bw1, row1_h, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row1_y + Emu(80000),
           bw1 - Emu(80000), Emu(420000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row1_y + Emu(510000),
           bw1 - Emu(60000), Emu(850000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw1 + Emu(10000), row1_y + row1_h // 2)

    # 右端「喰種対決」
    tb(s, Inches(9.0), row1_y + Emu(600000), Inches(0.9), Emu(600000),
       "喰種対決\n抽選！", 8.5, bold=True, color=C_RED, align=PP_ALIGN.CENTER, font=FONT_H)

    # ↓矢印（右端から下段右端へ）
    arrow_d(s, Inches(9.45), row1_y + row1_h + Emu(20000), col=C_RED)

    # ── 下段（右→左）: BITES → 特化ゾーン → ジャッジメント → 裏AT
    row2_y = row1_y + row1_h + Emu(350000)
    row2_h = Emu(1550000)
    boxes2 = [
        (RGBColor(0x20, 0x10, 0x00), C_GOLD, "裏AT（上位AT）",
         "純増5.0枚\n喰種対決成功率UP\nCCG死神300〜2000枚"),
        (RGBColor(0x14, 0x04, 0x04), C_RED,  "有馬貴将\nジャッジメント",
         "AT終了後/BITES後\n成功率約61%\n成功→裏AT/CCG死神"),
        (RGBColor(0x12, 0x06, 0x20), C_PUR,  "特化ゾーン\n(百足覚醒/隻眼の梟)",
         "百足覚醒：期待300枚\n隻眼の梟：期待500枚\nBITES勝利後に突入"),
        (RGBColor(0x18, 0x04, 0x04), C_RED2, "BITES\n報酬ゾーン",
         "喰種対決勝利で突入\n4桁上乗せ期待\n小役でUP/ハズレで終了"),
    ]
    bw2 = Inches(2.0)
    gap2 = Inches(0.28)
    sx2 = Inches(0.3)
    n2 = len(boxes2)
    for i, (fill, bc, lbl, sub) in enumerate(reversed(boxes2)):
        idx = n2 - 1 - i
        bx0 = sx2 + idx * (bw2 + gap2)
        rect_b(s, bx0, row2_y, bw2, row2_h, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), row2_y + Emu(80000),
           bw2 - Emu(80000), Emu(420000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), row2_y + Emu(510000),
           bw2 - Emu(60000), Emu(900000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if idx > 0:
            arrow_r(s, bx0 - gap2 - Emu(10000), row2_y + row2_h // 2, col=C_RED)

    # 左端ラベル（終了→通常へ）
    rect_b(s, Inches(0.05), row2_y + Emu(400000), Emu(190000), Emu(750000),
           RGBColor(0x10, 0x04, 0x20), C_PUR, 1.5)
    tb(s, Inches(0.05), row2_y + Emu(440000), Emu(190000), Emu(680000),
       "↑\n通常", 7, bold=True, color=C_PUR, align=PP_ALIGN.CENTER)

    footer(s,
           "「通常→CZ/EP→AT→喰種対決→BITES→ジャッジメント→裏AT」の全ルートを蛇行2段で可視化。",
           "補足: 裏AT終了後は通常に戻るが、高設定ほど裏AT直行率が高く爆連しやすい構造。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    """Part A – スライド3: 通常時の遊び方"""
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 招待状50G周期・CZ2種・EP BONUS直撃の3ルート", f"3/{TOTAL_SLIDES}")

    # 左カラム: 周期システム
    lx, ly = Inches(0.3), Inches(0.72)
    lw = Inches(3.0)
    card_h = Emu(3900000)
    rect_b(s, lx, ly, lw, card_h, C_CARD, C_PUR, 1.5)
    rect(s, lx, ly, Emu(45000), card_h, C_PUR)
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(260000),
       "周期システム", 11, bold=True, color=C_PUR, font=FONT_H)

    items_l = [
        ("招待状（50G周期）",
         "通常時は50G消化ごとに左下に\n「招待状」が出現。次のCZまでの\nG数や設定示唆を行う。"),
        ("CZ ① 追憶",
         "通常格CZ。AT当選を目指す。\n突破率は設定差あり。\n失敗でも再チャレンジあり。"),
        ("CZ ② 大喰種討伐",
         "高格CZ・高期待度版。\n突破成功でAT確定。\n高設定ほど出現率UP。"),
        ("EP BONUS",
         "レア役成立で直撃当選あり。\nEP BONUS成立でAT直行ルート。\n最も速い突入経路。"),
    ]
    for i, (title, body) in enumerate(items_l):
        iy = ly + Emu(320000) + i * Emu(880000)
        rect_b(s, lx + Emu(60000), iy, lw - Emu(75000), Emu(830000), C_CARD2, C_PUR, 0.6)
        tb(s, lx + Emu(90000), iy + Emu(50000), lw - Emu(130000), Emu(250000),
           title, 9, bold=True, color=C_PUR2)
        tb(s, lx + Emu(90000), iy + Emu(290000), lw - Emu(130000), Emu(480000),
           body, 7.5, color=C_WHITE)

    # 中カラム: 赫眼状態
    mx, my = Inches(3.5), Inches(0.72)
    mw = Inches(3.0)
    rect_b(s, mx, my, mw, card_h, C_CARD, C_RED, 1.5)
    rect(s, mx, my, Emu(45000), card_h, C_RED)
    tb(s, mx + Emu(75000), my + Emu(50000), mw - Emu(100000), Emu(260000),
       "赫眼状態（高確率）", 11, bold=True, color=C_RED, font=FONT_H)

    items_m = [
        ("赫眼状態とは",
         "通常時にいつでも移行し得る\nチェリー高確率状態。\nCZ/AT当選率が大幅UP。"),
        ("チェリー連続",
         "赫眼状態中はチェリー頻出。\n連続するほどCZ当選や\nAT直撃期待度が上がる。"),
        ("打ち方ポイント",
         "左リール枠内にBARを狙う\n（チェリー・スイカ対応）。\n小役取得漏れに注意。"),
        ("やめどき",
         "AT終了後はステージ確認。\nCZ前兆中やモード示唆で\n続行判断する。"),
    ]
    for i, (title, body) in enumerate(items_m):
        iy = my + Emu(320000) + i * Emu(880000)
        rect_b(s, mx + Emu(60000), iy, mw - Emu(75000), Emu(830000), C_CARD2, C_RED, 0.6)
        tb(s, mx + Emu(90000), iy + Emu(50000), mw - Emu(130000), Emu(250000),
           title, 9, bold=True, color=C_RED2)
        tb(s, mx + Emu(90000), iy + Emu(290000), mw - Emu(130000), Emu(480000),
           body, 7.5, color=C_WHITE)

    # 右カラム: 天井・設定示唆
    rx, ry = Inches(6.7), Inches(0.72)
    rw = Inches(3.0)
    rect_b(s, rx, ry, rw, card_h, C_CARD, C_GOLD, 1.5)
    rect(s, rx, ry, Emu(45000), card_h, C_GOLD)
    tb(s, rx + Emu(75000), ry + Emu(50000), rw - Emu(100000), Emu(260000),
       "天井・設定示唆", 11, bold=True, color=C_GOLD, font=FONT_H)

    items_r = [
        ("天井G数",
         "AT非当選時の天井あり。\nリセット時はCZ天井\n短縮（200G+α）。"),
        ("CZ後の設定示唆",
         "CZ終了後ボタン押しで\nカード出現→設定示唆。\n涼宮：偶数/梟：4以上\n有馬：設定6確定。"),
        ("招待状の示唆文",
         "「たっぷり楽しんで」→設定4以上\n「特別な夜を楽しもう」→設定6\nといった示唆を含む。"),
        ("モード別天井",
         "滞在モードによりCZ発生\nG数が変化。高モードほど\n早いCZが期待できる。"),
    ]
    for i, (title, body) in enumerate(items_r):
        iy = ry + Emu(320000) + i * Emu(880000)
        rect_b(s, rx + Emu(60000), iy, rw - Emu(75000), Emu(830000), C_CARD2, C_GOLD, 0.6)
        tb(s, rx + Emu(90000), iy + Emu(50000), rw - Emu(130000), Emu(250000),
           title, 9, bold=True, color=C_GOLD2)
        tb(s, rx + Emu(90000), iy + Emu(290000), rw - Emu(130000), Emu(480000),
           body, 7.5, color=C_WHITE)

    footer(s,
           "通常時は「招待状50G周期」という視覚的な周期管理+赫眼状態という二層構造でプレイヤーに常に期待感を与える。",
           "補足: CZ成功率は設定差大。設定示唆を積極的に読み取ることが高設定狙いのカギ。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    """Part A – スライド4: CZ/前兆の仕組み"""
    s = new_slide(prs)
    hdr(s, "CZ・前兆の仕組み ── 2種のCZと有馬貴将ジャッジメントの突破で道が開く", f"4/{TOTAL_SLIDES}")

    # 上段: CZ2種
    cx = Inches(0.3)
    cy = Inches(0.75)
    cw = Inches(4.5)
    ch = Emu(1700000)

    # CZ① 追憶
    rect_b(s, cx, cy, cw, ch, C_CARD, C_PUR, 1.5)
    rect(s, cx, cy, Emu(45000), ch, C_PUR)
    tb(s, cx + Emu(75000), cy + Emu(60000), cw - Emu(100000), Emu(280000),
       "CZ ① 追憶（通常格）", 11, bold=True, color=C_PUR, font=FONT_H)
    tb(s, cx + Emu(75000), cy + Emu(360000), cw - Emu(100000), ch - Emu(420000),
       "■ 発生条件: 50G周期消化 or レア役\n"
       "■ 内容: バトル演出でAT当選を目指す\n"
       "■ 成功率: 設定差あり（高設定ほど有利）\n"
       "■ 失敗時: 再チャレンジ抽選あり\n"
       "■ 演出見方: キャラ会話・技演出に注目",
       8.5, color=C_WHITE)

    # CZ② 大喰種討伐
    cx2 = Inches(5.2)
    rect_b(s, cx2, cy, cw, ch, C_CARD, C_RED, 1.5)
    rect(s, cx2, cy, Emu(45000), ch, C_RED)
    tb(s, cx2 + Emu(75000), cy + Emu(60000), cw - Emu(100000), Emu(280000),
       "CZ ② 大喰種討伐（高格）", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, cx2 + Emu(75000), cy + Emu(360000), cw - Emu(100000), ch - Emu(420000),
       "■ 発生条件: 高モード滞在時・レア役高確\n"
       "■ 内容: 成功でAT確定（失敗なし確定型も）\n"
       "■ 期待度: ① より大幅に高い\n"
       "■ 演出: 大喰種（SS級）が出現\n"
       "■ 高設定ほど出現頻度が高い",
       8.5, color=C_WHITE)

    # 下段: 有馬貴将ジャッジメント
    jy = Inches(2.6)
    jw = Inches(9.4)
    jh = Emu(1700000)
    rect_b(s, Inches(0.3), jy, jw, jh, C_CARD, C_GOLD, 2.0)
    rect(s, Inches(0.3), jy, Emu(55000), jh, C_GOLD)
    tb(s, Inches(0.55), jy + Emu(60000), jw - Emu(100000), Emu(280000),
       "上位CZ「有馬貴将ジャッジメント」── AT終了/BITES後に突入する最重要ゾーン", 11, bold=True,
       color=C_GOLD, font=FONT_H)

    # 3カラム詳細
    cols = [
        ("突入条件",
         "AT中の喰種対決で有馬乱入\n→勝利後にジャッジメント発生\nまたはAT終了時のエンディング\n達成でも突入。"),
        ("5G間の挑戦",
         "5G継続のミニバトル形式\n成功率約61%（全設定）\n成功→CCG死神ゾーン突入\n失敗→AT終了（引き戻しなし）"),
        ("演出の見方",
         "有馬貴将の攻撃演出に注目\n白エフェクト：低期待\n赤エフェクト：高期待\n金エフェクト：成功濃厚"),
    ]
    col_w = Inches(2.9)
    for i, (title, body) in enumerate(cols):
        cx3 = Inches(0.5) + i * (col_w + Emu(100000))
        rect_b(s, cx3, jy + Emu(380000), col_w, Emu(1200000), C_CARD2, C_GOLD, 0.8)
        tb(s, cx3 + Emu(50000), jy + Emu(430000), col_w - Emu(100000), Emu(280000),
           title, 9.5, bold=True, color=C_GOLD2)
        tb(s, cx3 + Emu(50000), jy + Emu(720000), col_w - Emu(100000), Emu(800000),
           body, 7.5, color=C_WHITE)

    footer(s,
           "CZ2種の格差設計が「大喰種討伐」発生時の高揚感を演出。ジャッジメントはAT体験のクライマックスを担う。",
           "補足: ジャッジメント失敗はAT終了確定。この緊張感がゲーム性の核心。成功時の爆発力との落差が面白さを作る。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（出玉を伸ばす仕組み）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    """Part A – スライド5: AT/ボーナスの遊び方"""
    s = new_slide(prs)
    hdr(s, "AT「東京喰種咬」── 差枚管理＋喰種対決＋BITESで出玉を積み上げる", f"5/{TOTAL_SLIDES}")

    # AT概要（上段）
    lx, ly = Inches(0.3), Inches(0.72)
    lw, lh = Inches(4.5), Emu(1850000)
    rx, rw = Inches(5.1), Inches(4.6)

    # 左: AT概要
    rect_b(s, lx, ly, lw, lh, C_CARD, C_RED, 1.5)
    rect(s, lx, ly, Emu(45000), lh, C_RED)
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(260000),
       "AT「東京喰種咬」概要", 10.5, bold=True, color=C_RED, font=FONT_H)
    tb(s, lx + Emu(75000), ly + Emu(320000), lw - Emu(100000), lh - Emu(380000),
       "■ 純増: 約4.0枚/G（差枚数管理型）\n"
       "■ 初期差枚数: 150枚\n"
       "■ レア役成立 → 差枚上乗せ抽選\n"
       "  強チェリー/チャンス目(A・B)/確定役\n"
       "  → 上乗せ濃厚（リールロックで100〜500枚）\n"
       "■ 規定G消化 → 喰種対決抽選",
       8.5, color=C_WHITE)

    # 右: 喰種対決の流れ
    rect_b(s, rx, ly, rw, lh, C_CARD, C_PUR, 1.5)
    rect(s, rx, ly, Emu(45000), lh, C_PUR)
    tb(s, rx + Emu(75000), ly + Emu(50000), rw - Emu(100000), Emu(260000),
       "喰種対決の流れ", 10.5, bold=True, color=C_PUR, font=FONT_H)
    tb(s, rx + Emu(75000), ly + Emu(320000), rw - Emu(100000), lh - Emu(380000),
       "■ レア役 or 規定G数 → 喰種対決抽選\n"
       "■ バトル勝利 → BITES（報酬ゾーン）\n"
       "  小役連続でUP / ハズレで終了のドキドキ構造\n"
       "  4桁上乗せも狙える\n"
       "■ 勝利後→特化ゾーン移行抽選あり",
       8.5, color=C_WHITE)

    # 下段: 特化ゾーン2種
    zones = [
        (C_RED, "特化ゾーン ①\n百足覚醒",
         "■ 上乗せ期待枚数: 約300枚\n"
         "■ 「咬みつく」度に上乗せ枚数が増加していく\n"
         "■ 連続する咬みつきで300枚超えも\n"
         "■ 終了後はAT継続 or 喰種対決へ"),
        (C_PUR, "特化ゾーン ②\n隻眼の梟",
         "■ 上乗せ期待枚数: 約500枚\n"
         "■ 1G完結型の超高速特化ゾーン\n"
         "■ 1G勝負でドカンと積み増す爽快感\n"
         "■ 百足覚醒より高い期待値"),
        (C_GOLD, "上乗せ演出法則",
         "■ リールロック発生: 100〜500枚濃厚\n"
         "■ 赫眼連続演出: 複数回上乗せ期待\n"
         "■ 金枠フラッシュ: 大量上乗せ濃厚\n"
         "■ 法則崩れはさらなる上乗せのサイン"),
    ]
    zy = Inches(2.75)
    zw = Inches(3.0)
    zh = Emu(1900000)
    for i, (bc, title, body) in enumerate(zones):
        zx = Inches(0.3) + i * (zw + Emu(120000))
        rect_b(s, zx, zy, zw, zh, C_CARD, bc, 1.5)
        rect(s, zx, zy, Emu(45000), zh, bc)
        tb(s, zx + Emu(75000), zy + Emu(60000), zw - Emu(100000), Emu(380000),
           title, 10, bold=True, color=bc, font=FONT_H)
        tb(s, zx + Emu(75000), zy + Emu(450000), zw - Emu(100000), zh - Emu(510000),
           body, 8, color=C_WHITE)

    footer(s,
           "AT内の上乗せ設計は「差枚管理＋バトル勝利ボーナス＋特化ゾーン」の三重構造。積み上げの連鎖が快感の核。",
           "補足: 強役での上乗せ濃厚演出がプレイヤーに能動的な期待感を与え続ける設計が評価の高さにつながっている。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位AT（裏AT）への道と遊び方
# ══════════════════════════════════════════════════════════════
def s_upper_at(prs):
    """Part A – スライド6: 上位ATへの道"""
    s = new_slide(prs)
    hdr(s, "上位AT「裏AT」への道 ── ジャッジメント→CCG死神→裏AT 三段階の夢", f"6/{TOTAL_SLIDES}")

    # 三段階フロー（横並び）
    stages = [
        (C_GOLD, "有馬貴将\nジャッジメント",
         "突入条件:\n・有馬乱入バトル勝利後\n・AT中エンディング達成後\n\n内容:\n5G間バトル\n成功率: 約61%\n失敗 → AT終了"),
        (C_RED, "CCG死神\n（特殊上乗せゾーン）",
         "突入条件:\nジャッジメント成功\n\n内容:\n1G完結型ゾーン\n300 / 500 / 1000 / 2000枚\nのいずれかを上乗せ\n成立役により決定"),
        (C_PUR2, "裏AT（上位AT）",
         "突入条件:\n・ジャッジメント成功一部\n・非有馬経由AT当選時の一部\n\n内容:\n純増: 約5.0枚/G\n喰種対決成功率が大幅UP\n有馬乱入率: 全設定で高確率\n期待枚数: 約3430枚（込）"),
    ]
    sw = Inches(2.8)
    sh = Emu(3500000)
    sy = Inches(0.72)
    for i, (bc, title, body) in enumerate(stages):
        sx2 = Inches(0.3) + i * (sw + Emu(500000))
        rect_b(s, sx2, sy, sw, sh, C_CARD, bc, 2.0)
        rect(s, sx2, sy, Emu(55000), sh, bc)
        tb(s, sx2 + Emu(80000), sy + Emu(60000), sw - Emu(110000), Emu(380000),
           title, 11, bold=True, color=bc, font=FONT_H, align=PP_ALIGN.CENTER)
        rect(s, sx2, sy + Emu(450000), sw, Emu(7000), bc)
        tb(s, sx2 + Emu(80000), sy + Emu(480000), sw - Emu(110000), sh - Emu(550000),
           body, 8.5, color=C_WHITE)
        if i < 2:
            arrow_r(s, sx2 + sw + Emu(50000), sy + sh // 2, col=bc)

    # 右端: 裏AT中の喰種対決強化
    rx = Inches(9.05)
    rect_b(s, rx, sy, Inches(0.8), sh, C_CARD2, C_PUR, 1.0)
    tb(s, rx + Emu(20000), sy + Emu(100000), Inches(0.65), sh - Emu(200000),
       "裏AT\n中の\n喰種\n対決\n\n有馬\n乱入\n大幅\nUP\n\n約\n3430\n枚\n確定", 7,
       bold=True, color=C_PUR2, align=PP_ALIGN.CENTER)

    # 下部: 裏AT突入の設定差
    fy2 = Inches(4.6)
    rect_b(s, Inches(0.3), fy2, Inches(9.4), Emu(380000), C_CARD2, C_GOLD, 1.0)
    tb(s, Inches(0.5), fy2 + Emu(50000), Inches(9.0), Emu(300000),
       "設定差ポイント: 高設定ほど「非有馬経由」AT当選時の裏AT直行率が高い。高設定台の爆発力の根幹がここにある。",
       8.5, bold=True, color=C_GOLD)

    footer(s,
           "「ジャッジメント成功率61%」という絶妙な数字が、期待と緊張の均衡を保ちつつ爆発体験を担保する設計の要。",
           "補足: 裏AT中の有馬乱入は全設定で 1・3戦目18.75%、5戦目62.50%。高設定は突入率で差がつく構造。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（なぜアワードGOLDを獲れたのか）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    """Part B – スライド7: 面白さの設計分析"""
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── なぜアワードGOLDを獲れたのか", f"7/{TOTAL_SLIDES}")

    # Part Bバッジ
    rect(s, Inches(9.1), Inches(0.07), Inches(0.75), Emu(240000), C_GOLD)
    tb(s, Inches(9.1), Inches(0.07), Inches(0.75), Emu(240000),
       "Part B", 7.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    # 5つの設計要素（グリッド）
    factors = [
        (C_PUR, "三層の上乗せ設計",
         "差枚管理→喰種対決→BITES→特化ゾーン\nと段階的に上乗せが重なる構造。\n「まだ続く？」という連続期待感がクセになる。"),
        (C_RED, "ジャッジメントの緊張設計",
         "AT終了間際の5Gバトル（成功率61%）。\n失敗確定のリスクがあるからこそ\n成功時の爆発が感動的になる。"),
        (C_GOLD, "裏ATによる一撃性",
         "裏AT突入で期待枚数3430枚。\n高設定ほど突入率が高く\n高設定探しへの動機がホールに客を呼ぶ。"),
        (C_PUR2, "原作IPとの一致度",
         "「喰種対決」「有馬貴将」「CCG死神」等\n東京喰種の世界観をそのままゲーム設計に転換。\nファン心理を刺激する演出と符合した報酬構造。"),
        (C_RED2, "複数の到達ルート",
         "通常→CZ→AT だけでなく\nEP BONUS直撃や赫眼状態からの直撃など\n多様な体験を生む複線的な設計。"),
        (C_GOLD2, "設定差の可視化",
         "CZ後のカード・招待状文言など\n多数の設定示唆要素を内包。\n上手い人ほど楽しめる深みを担保。"),
    ]
    fw = Inches(3.0)
    fh = Emu(1450000)
    for i, (bc, title, body) in enumerate(factors):
        row, col = i // 3, i % 3
        fx = Inches(0.3) + col * (fw + Emu(120000))
        fy = Inches(0.75) + row * (fh + Emu(100000))
        rect_b(s, fx, fy, fw, fh, C_CARD, bc, 1.5)
        rect(s, fx, fy, Emu(45000), fh, bc)
        tb(s, fx + Emu(70000), fy + Emu(60000), fw - Emu(90000), Emu(280000),
           title, 9.5, bold=True, color=bc, font=FONT_H)
        tb(s, fx + Emu(70000), fy + Emu(360000), fw - Emu(90000), fh - Emu(420000),
           body, 8, color=C_WHITE)

    footer(s,
           "「特別なレアフラグ不要で出玉トリガーまでたどり着ける」という間口の広さと、裏ATの爆発力の両立がGOLD受賞の本質。",
           "補足: ホールでも中古機150万円前後と高値安定。オペレーター側の評価も高く稼動・粗利ともに好評価。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_eval(prs):
    """Part B – スライド8: 評価・良い点と課題"""
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── ホール・プレイヤー・メーカー三視点での評価", f"8/{TOTAL_SLIDES}")

    # 左: 良い点（PROS）
    lx, ly = Inches(0.3), Inches(0.72)
    lw = Inches(4.5)
    lh = Emu(3900000)
    rect_b(s, lx, ly, lw, lh, C_CARD, C_GREEN, 1.5)
    rect(s, lx, ly, Emu(45000), lh, RGBColor(0x22, 0xCC, 0x66))
    tb(s, lx + Emu(75000), ly + Emu(50000), lw - Emu(100000), Emu(280000),
       "良い点（PROS）", 12, bold=True, color=RGBColor(0x22, 0xCC, 0x66), font=FONT_H)

    pros = [
        "間口の広さ: 特別な強役不要で喰種対決までたどり着ける設計",
        "三層上乗せで常に「続き」への期待感が持続",
        "ジャッジメントの緊張×成功時の爆発落差が高い感情体験を生む",
        "原作IPとゲーム設計の一致度が高くファン層・一般層両方に響く",
        "高設定の裏AT突入率で設定差が明確→高設定探しモチベUP",
        "CCG死神1G完結の爽快感・BITESの小役連打の没入感が独自",
    ]
    for i, pro in enumerate(pros):
        iy = ly + Emu(350000) + i * Emu(570000)
        rect(s, lx + Emu(55000), iy, Emu(40000), Emu(380000), RGBColor(0x22, 0xCC, 0x66))
        tb(s, lx + Emu(115000), iy + Emu(30000), lw - Emu(160000), Emu(480000),
           pro, 8.5, color=C_WHITE)

    # 右: 課題（CONS）
    rx, ry = Inches(5.1), Inches(0.72)
    rw = Inches(4.6)
    rh = Emu(3900000)
    rect_b(s, rx, ry, rw, rh, C_CARD, C_RED, 1.5)
    rect(s, rx, ry, Emu(45000), rh, C_RED)
    tb(s, rx + Emu(75000), ry + Emu(50000), rw - Emu(100000), Emu(280000),
       "課題・改善点（CONS）", 12, bold=True, color=C_RED2, font=FONT_H)

    cons = [
        "駆け抜けが多い: 差枚数切れでのAT終了が頻発しやすい",
        "成功時との格差大: ジャッジメント失敗でゼロリターンのストレス",
        "荒い出玉性能: 引き次第で高設定でも噛み合わない局面が発生",
        "裏AT一極集中: 裏AT非突入時の消化が単調に感じやすい",
        "設定判別難度: 示唆要素は多いが確定系が少なく長時間必要",
        "初心者導線: ゲームフローの複雑さは初見プレイヤーに難解",
    ]
    for i, con in enumerate(cons):
        iy = ry + Emu(350000) + i * Emu(570000)
        rect(s, rx + Emu(55000), iy, Emu(40000), Emu(380000), C_RED)
        tb(s, rx + Emu(115000), iy + Emu(30000), rw - Emu(160000), Emu(480000),
           con, 8.5, color=C_WHITE)

    footer(s,
           "良い点も課題も「差枚管理＋バトル構造」の特性から生まれている。荒さとエンタメ性は表裏一体の設計判断。",
           "補足: 賛否両論の声があることがむしろ「尖った個性」の証拠。万人受けより熱狂的なファンを作る設計思想。")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_summary(prs):
    """Part B – スライド9: まとめ・設計から学べること"""
    s = new_slide(prs)
    hdr(s, "まとめ ── アワードGOLD受賞機から学べる設計の本質", f"9/{TOTAL_SLIDES}")

    # 上部: 総括カード
    rect_b(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(1100000),
           C_CARD, C_GOLD, 2.0)
    rect(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(55000), C_GOLD)
    tb(s, Inches(0.5), Inches(0.82), Inches(9.0), Emu(280000),
       "L東京喰種 総括", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.5), Inches(1.12), Inches(9.0), Emu(600000),
       "差枚管理AT×喰種対決バトル×裏ATという三層構造で、「特別な役不要の間口の広さ」と「裏ATの一撃爆発力」を両立。\n"
       "IPの世界観をゲーム設計に落とし込み、有馬貴将ジャッジメントで感情的クライマックスを演出。ホール・ユーザー双方から高評価を獲得し"
       "パチスロアワード2025 GOLDを受賞。",
       8.5, color=C_WHITE)

    # 設計から学べること（4点）
    lessons = [
        (C_PUR, "間口×爆発の両立設計",
         "強役不要で高揚感を得られる導線を引きつつ、上位ルートで圧倒的な爆発力を担保する設計が最も重要な学び。\n"
         "「誰でも楽しめる」と「熱狂できる頂点」を同時に設計することがアワード級の機種の条件。"),
        (C_RED, "緊張と爆発の落差設計",
         "ジャッジメントの失敗リスク（AT終了確定）があるからこそ成功時の爆発が輝く。\n"
         "ゼロリスクの上乗せゲームより「失うかもしれない緊張感」が感動体験を生む。"),
        (C_GOLD, "IP世界観とメカニクスの一致",
         "「有馬貴将」「CCG」「赫子」「喰種対決」という原作要素がゲームフローと対応している。\n"
         "IP活用は見た目だけでなく、強敵=強い報酬・組織=システムへの落とし込みが深い体験を生む。"),
        (C_PUR2, "設定差の可視化と探す楽しみ",
         "多層の設定示唆（カード/招待状/終了画面）で「高設定探し」という別ゲームを付与。\n"
         "プレイヤーの滞在理由を「今の楽しさ」だけでなく「将来の可能性」にも分散させる設計。"),
    ]
    lw2 = Inches(4.5)
    lh2 = Emu(1350000)
    for i, (bc, title, body) in enumerate(lessons):
        row, col = i // 2, i % 2
        lx2 = Inches(0.3) + col * (lw2 + Emu(300000))
        ly2 = Inches(2.05) + row * (lh2 + Emu(80000))
        rect_b(s, lx2, ly2, lw2, lh2, C_CARD, bc, 1.5)
        rect(s, lx2, ly2, Emu(45000), lh2, bc)
        tb(s, lx2 + Emu(70000), ly2 + Emu(60000), lw2 - Emu(90000), Emu(270000),
           title, 9.5, bold=True, color=bc, font=FONT_H)
        tb(s, lx2 + Emu(70000), ly2 + Emu(350000), lw2 - Emu(90000), lh2 - Emu(410000),
           body, 8, color=C_WHITE)

    footer(s,
           "L東京喰種の設計哲学: 「誰もが出玉体験にたどり着ける設計」と「一撃爆発の感情的クライマックス」の共存。",
           "この二軸こそが長期ヒット・アワード受賞・高中古価格の根本要因。次世代機設計のベンチマークとなる一台。")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)      # 1
    s_flow(prs)       # 2
    s_normal(prs)     # 3
    s_cz(prs)         # 4
    s_at(prs)         # 5
    s_upper_at(prs)   # 6
    s_design(prs)     # 7
    s_eval(prs)       # 8
    s_summary(prs)    # 9

    prs.save(OUT_PATH)
    print(f"保存完了: {OUT_PATH}")


if __name__ == "__main__":
    main()
