"""
スマスロ モンスターハンターライズ 機種説明＋分析 統合版 PPTXジェネレーター v1
出力: proposals/機種分析/モンスターハンターライズ/mhrise_guide_v1.pptx
テーマ: 深茶×深緑×橙×金（モンハン世界観）
情報源: ちょんぼりすた・なな徹・一撃・altema・pachiseven（2024年11月〜）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "モンスターハンターライズ", "mhrise_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深茶×深緑×橙×金）──────────────────────────
C_BG    = RGBColor(0x0C, 0x08, 0x04)   # 深茶背景
C_CARD  = RGBColor(0x18, 0x12, 0x08)   # カード背景
C_CARD2 = RGBColor(0x20, 0x18, 0x0A)   # カード背景2
C_ROW   = RGBColor(0x14, 0x10, 0x06)   # 行背景
C_GREEN = RGBColor(0x22, 0xAA, 0x44)   # メイングリーン #22AA44
C_GREEN2= RGBColor(0x18, 0x80, 0x30)   # ダークグリーン
C_ORG   = RGBColor(0xFF, 0x88, 0x00)   # オレンジ #FF8800
C_ORG2  = RGBColor(0xFF, 0xAA, 0x22)   # ライトオレンジ
C_GOLD  = RGBColor(0xCC, 0xAA, 0x22)   # 金 #CCAA22
C_GOLD2 = RGBColor(0xFF, 0xDD, 0x44)   # 明るい金
C_WHITE = RGBColor(0xF0, 0xEC, 0xE0)
C_CREAM = RGBColor(0xF0, 0xD8, 0x90)
C_GRAY  = RGBColor(0x99, 0x90, 0x80)
C_LTGRY = RGBColor(0x55, 0x4C, 0x40)
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)
C_RED   = RGBColor(0xDD, 0x22, 0x11)
C_BROWN = RGBColor(0x6C, 0x40, 0x10)   # 茶色アクセント

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景：森・自然をイメージした深緑茶 ──────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (12, 8, 4))
    draw = ImageDraw.Draw(img)
    # 木目調の縦グラデーション
    for y in range(h):
        t = y / h
        r = int(12 + 8 * t)
        g = int(8 + 16 * t)
        b = int(4 + 4 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    # 森の木々を暗示する縦筋
    for x in range(0, w, 55):
        alpha = 8
        draw.line([(x, 0), (x + 10, h)], fill=(alpha, alpha + 4, alpha), width=2)
    # 地面グラデーション（下部）
    for y in range(h - 100, h):
        t = (y - (h - 100)) / 100
        r = int(30 * t)
        g = int(20 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, 0))
    # 上部を少し暗く
    for y in range(0, 30):
        t = (30 - y) / 30 * 0.4
        draw.line([(0, y), (w, y)], fill=(int(12 * t), int(8 * t), 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_GREEN)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_GREEN)


def net_note(slide):
    tb(slide, Inches(7.3), Inches(5.38), Inches(2.6), Emu(180000),
       "※ネット解析情報より（ちょんぼりすた・なな徹・一撃・pachiseven）", 6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x10, 0x0C, 0x04))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_GREEN)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_ORG)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_GREEN
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_GREEN
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x08, 0x06, 0x02))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_GREEN)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_GREEN)

    tb(s, Inches(0.22), Inches(0.4), Inches(5.0), Emu(330000),
       "機種説明＋分析 統合ガイド  v1  ｜  パチスロアワード2025 ノミネート作品", 9, color=C_ORG2, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.82), Inches(5.1), Emu(900000),
       "スマスロ\nモンスターハンターライズ", 26, bold=True, color=C_ORG, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.68), Inches(5.0), Emu(280000),
       "スマスロ（Lパチスロ）── 狩猟ループ型AT×カムラポイント蓄積×多層上位AT搭載", 8.5, color=C_CREAM, font=FONT_H)

    # スペック表
    specs = [
        ("メーカー",        "エンターライズ  2024年11月18日導入"),
        ("設定",           "1〜6段階"),
        ("AT純増①",        "約2.7枚/G（赤7・BAR揃い狩猟ボーナス）"),
        ("AT純増②",        "約4.0枚/G（紫7揃い狩猟ボーナス・上位AT）"),
        ("設定1 AT初当り", "約1/309.5"),
        ("設定6 AT初当り", "約1/230.8"),
        ("設定1 機械割",   "97.9%"),
        ("設定6 機械割",   "114.3%"),
        ("天井",           "リプレイ規定200回（5セット）でCZ確定"),
        ("有利区間",       "差枚+2100枚到達→エンディング"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(3.0) + i * Emu(220000)
        tb(s, Inches(0.22), ry, Inches(1.8), Emu(210000),
           k, 7.5, color=C_GRAY)
        tb(s, Inches(2.02), ry, Inches(3.2), Emu(210000),
           v, 7.5, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_GREEN,  "① 狩猟ループ型AT（討伐→剥ぎ取り→次戦）",
         "モンスター討伐成功で「剥ぎ取りチャンス」に突入。\n"
         "ボーナスをストックし次の狩猟へ。\n"
         "討伐→報酬→次戦のモンハンサイクルをそのまま再現。"),
        (C_ORG,    "② カムラポイント蓄積×CZ3段階構造",
         "通常時は小役でカムラポイントを貯め\n"
         "100pt毎のクエストでAT突入を目指す。\n"
         "上位CZ「百竜夜行」成功で大連続ボーナス濃厚。"),
        (C_GOLD,   "③ 上位AT「気焔万丈」で純増4.0枚/G",
         "規定10頭討伐→エンディング→百竜ノ淵源BONUS\n"
         "経由で上位AT突入。純増が2.7→4.0枚/Gに上昇。\n"
         "業界最大18.5インチ液晶で圧倒的映像体験。"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.55) + i * Emu(1540000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.3), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.3), ac)
        tb(s, Inches(5.85), y0 + Emu(65000), Inches(3.8), Emu(310000),
           kw, 10.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(380000), Inches(3.8), Emu(450000),
           desc, 8, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「狩猟→討伐→剥ぎ取り→次戦」のモンハンサイクルをスロットゲーム性に完全昇華",
           "エンターライズ製・新筐体「イマーシブ」搭載・5Dサウンド×18.5インチ液晶で圧倒的没入感")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（全ルートを蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→CZ→AT→上位ATの全ルート", "2/9")

    # 上段：通常→初当り→狩猟ボーナス の基本フロー
    top_y = Inches(0.75)
    top_h = Emu(1100000)

    flow1 = [
        (C_CARD2, C_GRAY,   "通常遊技",          "カムラポイント蓄積\nレア役直撃抽選\nリプレイ規定でCZ"),
        (C_CARD,  C_GREEN2, "アイルーだるま\n落とし(CZ)",  "リプレイ規定回数\n（最大200回）で\n突入するメインCZ"),
        (C_CARD,  C_GREEN,  "百竜夜行\n(上位CZ)",  "12G継続CZ\n成功率約70%\n大連続ボーナス濃厚"),
        (C_CARD,  C_ORG,    "狩猟ボーナス\n(メインAT)",  "純増2.7or4.0枚/G\nモンスター討伐で\n継続するループ型"),
        (C_CARD,  C_GOLD,   "剥ぎ取り\nチャンス", "討伐後の報酬タイム\n銀/金/G/超の4種\n金以上でストック確定"),
    ]
    bw1 = Inches(1.62)
    gap1 = Inches(0.14)
    sx1 = Inches(0.28)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), top_y + Emu(70000), bw1 - Emu(60000), Emu(360000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), top_y + Emu(480000), bw1 - Emu(45000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(8000), cy1, C_GREEN)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1120000), Inches(3.5), Emu(260000),
       "天井：リプレイ規定5セット目（累計200回）でCZ確定", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1360000), Inches(3.8), Emu(5000), C_CYAN)
    tb(s, sx1 + Inches(4.0), top_y + Emu(1120000), Inches(3.5), Emu(260000),
       "有利区間差枚+2100枚到達→エンディング→上位AT抽選", 7.5, color=C_GOLD)

    # 中段区切り線
    rect(s, 0, Inches(2.12), SLIDE_W, Emu(5000), RGBColor(0x30, 0x50, 0x18))

    # 下段：AT内→上位ATへの昇格ルート
    bot_y = Inches(2.20)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,  C_ORG,   "狩猟ボーナス\n（基本AT）",   "純増約2.7枚/G\nモンスターと戦闘\n討伐率約52%"),
        (C_CARD,  C_GREEN, "百竜刀\n-千変万化MODE-", "早期討伐時突入\nボーナスストック\n特化ゾーン"),
        (C_CARD,  C_GOLD2, "エンディング\n(10頭討伐)",   "規定10頭討伐で\n突入する特殊フェーズ\n上位AT抽選の起点"),
        (RGBColor(0x14,0x10,0x02), C_GOLD,
         "気焔万丈\n（上位AT）",    "純増約4.0枚/G\n通常ATと同ゲーム性\n百竜ノ淵源経由"),
    ]
    bw2 = Inches(2.0)
    gap2 = Inches(0.22)
    sx2 = Inches(0.28)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8)
        tb(s, bx + Emu(40000), bot_y + Emu(70000), bw2 - Emu(70000), Emu(370000),
           lbl, 10, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), bot_y + Emu(490000), bw2 - Emu(55000), Emu(470000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw2 + Emu(12000), cy2, C_GOLD)

    # 右端補足ボックス
    rx_al = sx2 + 4 * (bw2 + gap2)
    rect_b(s, rx_al, bot_y, Inches(1.42), bot_h, C_CARD, C_CYAN, 1.5)
    rect(s, rx_al, bot_y, Emu(30000), bot_h, C_CYAN)
    tb(s, rx_al + Emu(50000), bot_y + Emu(70000), Inches(1.2), Emu(310000),
       "百竜ノ\n淵源\nBONUS", 9, bold=True, color=C_CYAN, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, rx_al + Emu(50000), bot_y + Emu(490000), Inches(1.2), Emu(460000),
       "討伐成功\n(約72%)で\n気焔万丈へ", 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    net_note(s)
    footer(s, "上段=通常〜狩猟ボーナス突入ルート（カムラPT・CZ・レア役直撃）、下段=AT内昇格ルート（討伐→エンディング→気焔万丈）",
           "討伐するたびに次の戦いが始まるモンハンサイクルがそのままスロットのAT継続設計に")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（モンハン要素×パチスロの融合）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── カムラポイント蓄積・CZ・レア役直撃", "3/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x20, 0x50, 0x18))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜AT突入ルート（全3系統）", 10, bold=True, color=C_GREEN)

    routes = [
        (C_GREEN,  "ルート① カムラポイント蓄積→クエスト",
         "小役入賞でカムラポイントを獲得。\n"
         "最大600pt（100pt毎にチャンス）でクエスト突入。\n"
         "クエスト成功でCZまたはAT直撃抽選が行われる。\n"
         "メインルートで最も頻出の突入経路。"),
        (C_ORG,    "ルート② レア役直撃（ライズゾーン）",
         "チェリー・スイカ・チャンス目などのレア役で\n"
         "AT直撃またはライズゾーン（カムラPT特化）抽選。\n"
         "引き強時の一発突入ルート。\n"
         "ライズゾーン突入後は一定G間カムラPT高速蓄積。"),
        (C_CYAN,   "ルート③ リプレイ規定回数到達→CZ",
         "リプレイが規定回数（1セット40回）成立するとCZ\n"
         "「アイルーだるま落とし」に突入。\n"
         "最大5セット目（累計200回）でCZ確定の天井機能。\n"
         "CZを経由して上位CZ「百竜夜行」へ昇格も。"),
        (C_GOLD,   "ルート④ 百竜夜行（上位CZ）直撃",
         "レア役や特定条件で上位CZ「百竜夜行」に直撃。\n"
         "12G継続・成功率約70%の高確CZ。\n"
         "成功時は「大連続ボーナス」濃厚で大量獲得の入口。\n"
         "突入時点で高期待度確定の激熱ルート。"),
    ]
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(290000) + i * Emu(1125000)
        rect_b(s, lx, iy, lw, Emu(1060000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1060000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(670000),
           b, 7.5, color=C_WHITE)

    # 右：チャンス役と確率・打ち方
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x20, 0x50, 0x18))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "チャンス役の種類と期待度・モンハン要素との連動", 10, bold=True, color=C_GREEN)

    chance = [
        (C_GRAY,  "リプレイ",      "約1/8",   "規定回数成立でCZ「アイルーだるま落とし」へ"),
        (C_GREEN2,"弱チェリー",    "約1/64",  "カムラPT加算・ライズゾーン抽選"),
        (C_GREEN, "スイカ",        "約1/100", "カムラPT大量加算・AT直撃抽選"),
        (C_ORG,   "チャンス目",    "約1/100", "AT直撃・百竜夜行突入の抽選対象"),
        (C_GOLD2, "強チェリー",    "約1/200", "高期待度・百竜夜行直撃の主要役"),
        (C_GOLD,  "ロングフリーズ","激レア",  "設定6濃厚・大量上乗せ確定演出"),
    ]
    ch_h = Emu(690000)
    for i, (ac, role, prob, desc) in enumerate(chance):
        cy = ry + Emu(290000) + i * ch_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, ch_h, bg)
        rect(s, rx, cy, Emu(35000), ch_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(55000), Inches(1.0), Emu(260000),
           role, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.0), cy + Emu(60000), Inches(0.8), Emu(240000),
           prob, 8, bold=True, color=C_GOLD2, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(310000), rw - Inches(0.6), Emu(320000),
           desc, 7.5, color=C_WHITE)

    # 打ち方メモ
    rect_b(s, rx, ry + Emu(4440000), rw, Emu(550000), C_CARD2, C_GREEN, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4490000), rw - Emu(80000), Emu(200000),
       "打ち方：通常時は順押し推奨。レア役を取りこぼさない打順を徹底", 8.5, bold=True, color=C_GREEN)
    tb(s, rx + Emu(60000), ry + Emu(4700000), rw - Emu(80000), Emu(280000),
       "カムラポイントのゲージはモンハンの「猟団レベル」的な位置づけ。蓄積感が継続動機になる。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "通常時の核心：カムラポイントという「蓄積感×目標感」がモンハンのクエスト受注体験をスロットで再現する",
           "リプレイ規定天井（最大5セット）は立ち回り計画の明確な指標。残セット数管理が重要")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み（狩猟要素と絡めて）
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ/前兆の仕組み ── アイルーだるま落とし×百竜夜行×狩猟要素", "4/9")

    # 左：アイルーだるま落としCZ
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x18, 0x40, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "メインCZ「アイルーだるま落とし」の仕組み", 10, bold=True, color=C_GREEN)

    cz_main = [
        (C_GREEN2, "突入条件",
         "リプレイが規定回数（1セット40回）成立で突入。\n"
         "セット数が増えるほど期待度アップ。\n"
         "5セット目（累計200回）で必ずCZ当選（天井）。"),
        (C_GREEN,  "CZ内容・ゲーム性",
         "アイルーキャラが積み上がるビジュアル演出。\n"
         "成功/失敗は演出展開で徐々に告知される。\n"
         "モンハンシリーズのマスコット・アイルーを主役にした\n"
         "ファン心理をくすぐるCZ設計。"),
        (C_ORG,    "成功時の恩恵",
         "成功→狩猟ボーナス（AT）突入確定。\n"
         "演出発展次第で大連続ボーナス（複数セット確定）も。\n"
         "百竜夜行（上位CZ）へ昇格することもある。"),
    ]
    for i, (ac, t, b) in enumerate(cz_main):
        iy = ly + Emu(260000) + i * Emu(1320000)
        rect_b(s, lx, iy, lw, Emu(1250000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1250000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(330000), lw - Emu(100000), Emu(780000),
           b, 8, color=C_WHITE)

    # 右上：百竜夜行（上位CZ）
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x50, 0x38, 0x08))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "上位CZ「百竜夜行」── 12G継続・成功率約70%", 10, bold=True, color=C_GOLD)

    hyakuryu_items = [
        (C_GOLD2,  "突入条件",
         "レア役・特定条件・アイルーCZ経由で昇格。\n"
         "百竜夜行は「百体の龍が押し寄せる」設定で\n"
         "原作のビッグイベントをCZに昇華している。"),
        (C_GOLD,   "ゲーム性・成功率",
         "12G継続のバトル型CZ。\n"
         "成功率は約70%と非常に高く、突入時点で期待度大。\n"
         "突入時に50%で成功を内部抽選→継続Gで告知。"),
        (C_ORG,    "成功時の恩恵（大連続ボーナス）",
         "成功時は「大連続ボーナス」がほぼ確定。\n"
         "複数の狩猟ボーナスが連続する大量獲得モード。\n"
         "上位AT「気焔万丈」への入口ともなる重要CZ。"),
    ]
    for i, (ac, t, b) in enumerate(hyakuryu_items):
        iy2 = ry + Emu(260000) + i * Emu(1320000)
        rect_b(s, rx, iy2, rw, Emu(1250000), C_CARD, ac, 1.5)
        rect(s, rx, iy2, Emu(45000), Emu(1250000), ac)
        tb(s, rx + Emu(75000), iy2 + Emu(50000), rw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(75000), iy2 + Emu(330000), rw - Emu(100000), Emu(780000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "CZ設計の核心：メインCZ（アイルー）と上位CZ（百竜夜行）の2段構成で期待度の緩急と昇格サプライズを演出",
           "百竜夜行の成功率70%は業界上位クラス。突入した時点で「当てた感覚」を与える心理設計")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（何をすれば出玉が伸びる）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT/ボーナス ── 狩猟ボーナスで出玉を伸ばす仕組み", "5/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x18, 0x40, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "狩猟ボーナスの種類と純増・継続の仕組み", 10, bold=True, color=C_GREEN)

    bonus_types = [
        (C_ORG,   "赤7揃い（狩猟ボーナス）",
         "純増：約2.7枚/G  保障：40Gor50G\n"
         "最もスタンダードな狩猟ボーナス。\n"
         "消化中にモンスターと戦い討伐を目指す。\n"
         "討伐成功→剥ぎ取りチャンス→次ボーナスへ。"),
        (C_GOLD2, "紫7揃い（狩猟ボーナス）",
         "純増：約4.0枚/G  保障：60Gor100G\n"
         "高純増バージョン。保障Gも長く出玉効率◎。\n"
         "上位AT「気焔万丈」と同様の純増を誇る。\n"
         "引くだけで一段上の出玉スピードに突入。"),
        (C_GRAY,  "BAR揃い（狩猟ボーナス）",
         "純増：約2.7枚/G  保障：25G\n"
         "保障が短いが討伐成功で次ボーナスへ繋がる。\n"
         "早期討伐で「百竜刀-千変万化MODE-」突入のチャンス。"),
    ]
    for i, (ac, t, b) in enumerate(bonus_types):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.8)
        rect(s, lx, iy, Emu(45000), Emu(1210000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, lx + Emu(75000), iy + Emu(330000), lw - Emu(100000), Emu(770000),
           b, 8, color=C_WHITE)

    # 右：剥ぎ取りチャンスと出玉上乗せ構造
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x18, 0x40, 0x10))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "剥ぎ取りチャンス×ストックで出玉を伸ばす", 10, bold=True, color=C_GREEN)

    hagitori_types = [
        (C_GRAY,  "剥ぎ取りチャンス 銀",
         "討伐後の基本報酬。\nボーナスストック獲得率は低い。\n次セットへの繋ぎ役。"),
        (C_GREEN, "剥ぎ取りチャンス 金",
         "ボーナスストック獲得がほぼ確定。\n「金」以上で次の狩猟が保証される。\n出玉伸長の分岐点となる重要ランク。"),
        (C_GREEN2,"剥ぎ取りチャンスG\n（継続率付き）",
         "70%ループでボーナスストックを大量獲得。\n継続するたびに次のチャンスが発生。\n一度乗ると一気に残弾数が増える。"),
        (C_GOLD,  "超剥ぎ取りチャンスG",
         "最上位報酬。大量ストック確定。\n上位AT「気焔万丈」への橋渡しになることも。\n百竜ノ淵源BONUS討伐後に発生する。"),
    ]
    hag_h = Emu(1000000)
    for i, (ac, t, b) in enumerate(hagitori_types):
        iy2 = ry + Emu(260000) + i * hag_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect_b(s, rx, iy2, rw, hag_h - Emu(15000), bg, ac, 1.2)
        rect(s, rx, iy2, Emu(35000), hag_h - Emu(15000), ac)
        tb(s, rx + Emu(60000), iy2 + Emu(45000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy2 + Emu(295000), rw - Emu(80000), Emu(580000),
           b, 7.5, color=C_WHITE)

    # 早期討伐の特典
    rect_b(s, rx, ry + Emu(4250000), rw, Emu(580000), C_CARD2, C_ORG, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4300000), rw - Emu(80000), Emu(220000),
       "早期討伐（15G以上残して討伐）→百竜刀-千変万化MODE-", 8, bold=True, color=C_ORG)
    tb(s, rx + Emu(60000), ry + Emu(4530000), rw - Emu(80000), Emu(270000),
       "ボーナスストック特化ゾーン。一気に規定討伐数エンディングへの近道となる。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "出玉伸長の核心：剥ぎ取りチャンスの「ランク」でストック数が決まる。金以上が出玉伸長の分岐点",
           "早期討伐ボーナスの設計がモンハンの「効率的な狩猟」という原作体験をスロットで完全再現")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方（気焔万丈）
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位ATへの道 ── 気焔万丈（上位AT）到達ルートと遊び方", "6/9")

    # 上段：昇格ルートフロー
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(260000), RGBColor(0x20, 0x50, 0x10))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(210000),
       "上位AT「気焔万丈」到達ルート（規定討伐数を積み上げる階段式設計）", 9, bold=True, color=C_GREEN)

    route_boxes = [
        (C_ORG,   "狩猟ボーナス\n（基本AT）",    "純増2.7枚/G\n討伐率約52%\n毎セット挑戦"),
        (C_GREEN2,"10頭討伐\n（エンディング）",  "差枚+2100枚or\n規定10頭討伐で\n発生するゴール"),
        (C_CYAN,  "百竜ノ淵源\nチャレンジ",      "約67%で成功\n→百竜ノ淵源BONUS\nへ進めるCZ"),
        (C_GOLD,  "気焔万丈\n（上位AT）",       "純増4.0枚/G\nモンスター討伐は\n通常ATと同ゲーム性"),
    ]
    bw_r = Inches(2.1)
    gap_r = Inches(0.26)
    sx_r = Inches(0.35)
    cy_r = Inches(1.42)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x16, 0x12, 0x02), ac, 2.0 if i >= 2 else 1.5)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(55000), Emu(370000), lbl, 10, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), cy_r - bh_r // 2 + Emu(460000),
           bw_r - Emu(45000), Emu(480000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw_r + Emu(12000), cy_r, C_GOLD)

    # 百竜ノ淵源BONUSの補足矢印
    tb(s, Inches(7.2), Inches(0.78), Inches(2.5), Emu(240000),
       "百竜ノ淵源BONUS\n討伐成功率：約72%", 7.5, color=C_CYAN, align=PP_ALIGN.CENTER)

    # 中区切り
    rect(s, 0, Inches(2.06), SLIDE_W, Emu(5000), RGBColor(0x30, 0x50, 0x18))

    # 下段左：気焔万丈の遊び方
    lx2, ly2 = Inches(0.28), Inches(2.12)
    lw2 = Inches(4.55)
    lh2 = Emu(2620000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_GOLD, 1.8)
    rect(s, lx2, ly2, Emu(45000), lh2, C_GOLD)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "気焔万丈（上位AT）の遊び方", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【基本ゲーム性】\n"
       "通常のAT（狩猟ボーナス）と同じゲーム性で進行。\n"
       "違いは純増が約2.7枚→約4.0枚/Gにアップする点のみ。\n"
       "モンスター討伐・剥ぎ取りチャンスの流れは同じ。\n\n"
       "【早期討伐の活用】\n"
       "保障G残15G以上残して早期討伐すると\n"
       "「百竜刀-千変万化MODE-」（ストック特化）チャンス。\n"
       "気焔万丈中でも早期討伐を狙うのが出玉最大化の鍵。\n\n"
       "【終了条件と次ラウンド】\n"
       "差枚上限（+2100枚）または規定討伐数達成でエンディング。\n"
       "終了後は百竜ノ淵源チャレンジで再度気焔万丈を狙える。",
       8, color=C_WHITE)

    # 下段右：到達難易度と期待値
    rx2, ry2 = Inches(5.05), Inches(2.12)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2, C_CARD, C_GREEN, 1.8)
    rect(s, rx2, ry2, Emu(45000), lh2, C_GREEN)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "気焔万丈到達の難易度と期待収支", 11, bold=True, color=C_GREEN, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【到達までのステップ確率】\n"
       "エンディング到達（10頭討伐）\n"
       "  → 百竜ノ淵源チャレンジ成功（約67%）\n"
       "  → 百竜ノ淵源BONUS討伐成功（約72%）\n"
       "  → 気焔万丈突入（約48%の複合確率）\n\n"
       "【純増アップの恩恵】\n"
       "2.7→4.0枚/G = 約1.48倍の出玉効率。\n"
       "同じ討伐数でも気焔万丈中の方が大幅に多く獲得可能。\n\n"
       "【設定差と気焔万丈頻度】\n"
       "設定6は初当り確率が高く（1/230.8）\n"
       "エンディング到達頻度も高くなる傾向。\n"
       "高設定ほど気焔万丈に乗れる回数が増える設計。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "上位ATの設計核心：「同ゲーム性で純増だけアップ」という明快な差別化が高達成感と継続動機を生む",
           "百竜ノ淵源チャレンジ（67%）×BONUS討伐（72%）=約48%の到達率。プレイ中の「近さ」が期待感を維持")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（IPとゲーム性の融合がなぜ機能するのか）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── モンハンIPとパチスロゲーム性が融合する理由", "7/9")

    principles = [
        (C_GREEN,  "① 「狩猟サイクル」という普遍的ゲームループ",
         "モンハンの核心は「準備→狩猟→報酬→次の狩猟」の\n"
         "繰り返し。スロットの「AT→継続抽選→継続」と\n"
         "構造が完全一致。IPのゲームループとスロットの\n"
         "ゲームループが同一のため、摩擦ゼロで融合する。"),
        (C_ORG,    "② カムラポイント蓄積という「育成感」",
         "小役入賞→PT積算→クエスト突入という\n"
         "蓄積→達成のサイクルがモンハンのクエスト受注体験。\n"
         "「ゲージを貯める」行為に目的意識が生まれ\n"
         "通常時のストレスを「積み上げ」に変換する。"),
        (C_GOLD,   "③ モンスター討伐という「自力感」の演出",
         "ボーナス中に「モンスターと戦い討伐する」演出が\n"
         "「自分がボーナスを継続させた」感覚を生む。\n"
         "結果は内部抽選だが、演出の勝利体験が\n"
         "プレイヤーの主体性・没入感を高める。"),
        (C_CYAN,   "④ 剥ぎ取りチャンスという「報酬の可視化」",
         "討伐後に報酬（銀/金/G/超）が明示される設計。\n"
         "これはモンハンの剥ぎ取りアイテム入手体験そのもの。\n"
         "「何が来るか」のワクワクと「金以上で次が来る」\n"
         "緊張感が共存する完成度の高い報酬演出。"),
        (C_GREEN2, "⑤ 18.5インチ液晶×5Dサウンドによる没入感",
         "業界最大の液晶と4スピーカー＋ウーファーが\n"
         "モンハンライズの世界をホールで再現。\n"
         "IPファンにとっては「ゲームの続き」を\n"
         "スロット台でプレイする体験になる特別な筐体設計。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1190000)
    gy = Inches(0.10)

    positions = [
        (Inches(0.28), Inches(0.72)),
        (Inches(5.17), Inches(0.72)),
        (Inches(0.28), Inches(0.72) + bh_p + gy),
        (Inches(5.17), Inches(0.72) + bh_p + gy),
        (Inches(0.28), Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1090000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(330000), pw - Emu(90000), ph - Emu(390000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "融合成功の核心：モンハンの「狩猟サイクル」がスロットのゲームループと構造的に一致しているため自然な融合が実現",
           "蓄積感（カムラPT）×自力感（討伐演出）×報酬可視化（剥ぎ取り）×没入感（筐体）の4要素が相互強化する設計")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN,  "IP×ゲームループの完全一致という稀有な融合",
         "モンハンの「狩猟サイクル」がスロットのAT継続設計\n"
         "と構造一致。違和感なく世界観に入れる\n"
         "IP活用の模範事例として業界屈指の完成度。"),
        (C_GREEN,  "カムラポイントによる通常時の目的意識",
         "小役入賞をポイント蓄積として可視化し\n"
         "通常時の「退屈な消化」を「積み上げる楽しさ」に変換。\n"
         "遊技時間全体に意味を持たせるUX設計。"),
        (C_GREEN,  "剥ぎ取りチャンスの階層的報酬設計",
         "銀→金→G→超の4段階で報酬を明示。\n"
         "「金が出れば次が来る」という明確な条件が\n"
         "プレイヤーに目標と緊張感を同時提供する。"),
        (C_GREEN,  "筐体（18.5インチ+5Dサウンド）の圧倒的没入感",
         "業界最大液晶×ウーファー搭載で\n"
         "モンハンの世界観を最大出力で再現。\n"
         "IPファンの「続きをプレイする感覚」を実現。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1140000)
        rect_b(s, lx, iy, lw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(300000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,    "ユーザー評価の低さ（平均1.4点）",
         "DMMぱちタウンのユーザー評価が1.4点と低い。\n"
         "「攻撃しない」「ポイント周期が遠い」の不満が目立つ。\n"
         "通常時の出力が体感的に弱く感じられやすい。"),
        (C_ORG,    "討伐率約52%の引き弱による単調感",
         "平均討伐期待度約52%は五分五分の継続率。\n"
         "連続失敗時の出玉停滞感が強く\n"
         "\"当たっているのに増えない\"という体験になりやすい。"),
        (C_GOLD,   "上位AT到達率の低さ（複合約48%）",
         "エンディング→チャレンジ→BONUS討伐成功の\n"
         "複合確率で約48%。半数以上が上位AT未到達。\n"
         "通常ATだけでは物足りなさが残りやすい。"),
        (C_GRAY,   "設定判別要素の少なさ・わかりにくさ",
         "ロングフリーズ（設定6濃厚）は極レア。\n"
         "キャラ紹介ムービーのエンタライオン登場が\n"
         "設定6示唆だが気付けないユーザーも多い。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1140000)
        rect_b(s, rx, iy, rw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(300000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "評価分析：IP融合の設計完成度は高いが「討伐率50%の単調さ」と「上位AT到達難易度」が稼働離脱の主因",
           "ユーザー評価1.4の背景には「通常ATの単発驚け抜け率の高さ」という出玉体験の問題が色濃く反映")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x18, 0x40, 0x10))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "スマスロ モンハンライズ ── 設計的強み総括", 10, bold=True, color=C_GREEN)

    strengths = [
        (C_GREEN,  "IP×ゲームループ一致という設計の理想形",
         "「狩猟サイクル」というモンハンの本質的ループが\n"
         "スロットのAT継続構造と完全対応。\n"
         "IPライセンスを世界観だけでなく設計原理で活用した\n"
         "パチスロIP活用の最高峰事例。"),
        (C_ORG,    "蓄積×報酬×継続の三位一体設計",
         "カムラポイント蓄積（通常時の目的）×\n"
         "剥ぎ取りチャンスの階層報酬×\n"
         "討伐ループ継続という3つの軸が一体となって\n"
         "プレイ全体に意味と楽しさを提供する。"),
        (C_GOLD,   "気焔万丈（純増4.0枚）という上位目標の存在",
         "通常AT（2.7枚）から上位AT（4.0枚）への\n"
         "明確なステップアップ設計。\n"
         "「もっと効率的に狩れる状態になる」という\n"
         "RPG的成長体験がプレイ継続動機を生む。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1210000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(325000), lw - Emu(90000), Emu(780000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_ORG, font=FONT_H)

    principles = [
        (C_GREEN,  "IPのゲームループを設計原理として使え",
         "世界観だけでなく「プレイの核心構造」を\n抽出してスロット設計に組み込む"),
        (C_ORG,    "蓄積感×報酬の可視化で通常時を楽しくせよ",
         "ゲージ・ポイントで目標を作り\n報酬ランク表示で達成感を演出する"),
        (C_GOLD,   "段階的な純増アップで上位目標を設計せよ",
         "「同じゲーム性で速くなる」という\nシンプルな上位差別化が最も伝わりやすい"),
        (C_CYAN,   "IP筐体は世界観体験の最大化装置として活用せよ",
         "液晶・サウンドへの投資が\nIPファンの「続きをプレイする感覚」を生む"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(760000)
        rect_b(s, rx, py0, rw, Emu(710000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(710000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(230000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(285000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3310000), rw, Emu(1060000),
           RGBColor(0x14, 0x10, 0x02), C_GOLD, 2.0)
    rect(s, rx, ry + Emu(3310000), Emu(40000), Emu(1060000), C_GOLD)
    tb(s, rx + Emu(65000), ry + Emu(3360000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_GOLD)
    tb(s, rx + Emu(65000), ry + Emu(3620000), rw - Emu(85000), Emu(690000),
       "「狩猟ループ×カムラPT蓄積×剥ぎ取り報酬×気焔万丈」\n"
       "という4軸がモンハンの世界観をスロット設計に昇華。\n"
       "パチスロアワード2025ノミネートはIP融合設計の\n"
       "完成度と筐体体験への評価と見られる。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "設計思想：「IPのゲームループをスロット構造に組み込む」という戦略が成功した教科書的モデルケース",
           "ユーザー評価の課題（討伐率50%・上位AT到達難）は次世代改善の指針。IP融合の設計完成度は業界最高水準")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)      # 1: タイトル・スペック・3ポイント
    s_flow(prs)       # 2: ゲームフロー全体図
    s_normal(prs)     # 3: 通常時の遊び方
    s_cz(prs)         # 4: CZ/前兆の仕組み
    s_at(prs)         # 5: AT/ボーナス（出玉を伸ばす仕組み）
    s_upper(prs)      # 6: 上位ATへの道（気焔万丈）
    s_design(prs)     # 7: 面白さの設計
    s_pros_cons(prs)  # 8: 良い点と課題
    s_matome(prs)     # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
