"""
スマスロ マギアレコード 魔法少女まどか☆マギカ外伝
統合PPTX（説明＋分析）生成スクリプト
パチスロアワード2025 SILVER受賞機種
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFilter
import io

# ─── 定数 ───────────────────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
OUT_PATH = r"C:\Users\h.kadoya\Desktop\slocri\proposals\機種分析\まどか外伝\madoka_guide_v1.pptx"

# カラーテーマ（深紺×ピンク×白×金）
C_DARK  = RGBColor(0x0A, 0x0A, 0x2A)    # 深紺（背景）
C_PINK  = RGBColor(0xCC, 0x44, 0xAA)    # マゼンタピンク
C_LPINK = RGBColor(0xEE, 0x88, 0xCC)    # ライトピンク
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)    # 白
C_GOLD  = RGBColor(0xFF, 0xD7, 0x00)    # 金
C_CYAN  = RGBColor(0x44, 0xCC, 0xEE)    # シアン（アクセント）
C_GRAY  = RGBColor(0xBB, 0xBB, 0xCC)    # グレー（補足テキスト）
C_GREEN = RGBColor(0x44, 0xEE, 0x88)    # 緑（良い点）
C_RED   = RGBColor(0xFF, 0x55, 0x55)    # 赤（課題）
C_BOX   = RGBColor(0x14, 0x14, 0x3A)    # ボックス背景（少し明るい深紺）

# フォント
FONT_H = "游明朝"
FONT_B = "メイリオ"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H


# ─── ヘルパー関数 ────────────────────────────────────

def make_bg(slide, glow_x=0.5, glow_y=0.4):
    """深紺×ピンクグロー背景をPILで生成してスライドに貼る"""
    W, H = 1280, 720
    img = Image.new("RGB", (W, H), (0x0A, 0x0A, 0x2A))
    draw = ImageDraw.Draw(img)
    # グロー円（ピンク）
    cx, cy = int(W * glow_x), int(H * glow_y)
    for r in range(320, 0, -8):
        alpha = int(30 * (1 - r / 320))
        col = (min(0xCC, alpha * 4), 0x10, min(0xAA, alpha * 3))
        draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=col)
    # 星屑ドット
    import random
    random.seed(42)
    for _ in range(120):
        x = random.randint(0, W)
        y = random.randint(0, H)
        r2 = random.randint(1, 3)
        bright = random.randint(150, 255)
        draw.ellipse([x-r2, y-r2, x+r2, y+r2], fill=(bright, bright, bright))
    img = img.filter(ImageFilter.GaussianBlur(radius=1))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(0), Inches(0), SLIDE_W, SLIDE_H)


def rect(slide, left, top, width, height, fill_color, alpha_box=False):
    """塗りつぶし矩形を追加"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def rect_b(slide, left, top, width, height, fill_color, border_color=None, border_pt=1.5):
    """ボーダー付き矩形"""
    shape = rect(slide, left, top, width, height, fill_color)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def tb(slide, text, left, top, width, height,
       font_name=None, font_size=14, color=None,
       bold=False, align=PP_ALIGN.LEFT, wrap=True):
    """テキストボックス追加"""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name or FONT_B
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color or C_WHITE
    return txb


def hdr(slide, title, subtitle=None):
    """スライドヘッダー（タイトルバー）"""
    # ピンクバー
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x55, 0x00, 0x44)
    bar.line.fill.background()
    # タイトル文字
    tb(slide, title, 0.15, 0.08, 8.5, 0.55,
       font_name=FONT_H, font_size=16, color=C_WHITE, bold=True)
    # サブタイトル（右寄せ）
    if subtitle:
        tb(slide, subtitle, 6.0, 0.08, 3.8, 0.52,
           font_name=FONT_B, font_size=10, color=C_LPINK,
           bold=False, align=PP_ALIGN.RIGHT)


def arrow_r(slide, left, top, length=0.5, label="", color=None):
    """右向き矢印とラベル"""
    c = color or C_PINK
    body = slide.shapes.add_shape(1,
        Inches(left), Inches(top + 0.08), Inches(length * 0.75), Inches(0.09))
    body.fill.solid()
    body.fill.fore_color.rgb = c
    body.line.fill.background()
    # 三角形（近似：縦長の矩形で代替）
    tip = slide.shapes.add_shape(1,
        Inches(left + length * 0.72), Inches(top), Inches(length * 0.28), Inches(0.25))
    tip.fill.solid()
    tip.fill.fore_color.rgb = c
    tip.line.fill.background()
    if label:
        tb(slide, label, left, top + 0.27, length, 0.22,
           font_size=8, color=c, align=PP_ALIGN.CENTER)


def net_note(slide, text, design_comment):
    """フッター：設計コメント（太字）＋補足説明"""
    # フッターバー
    bar = slide.shapes.add_shape(1, Inches(0), Inches(5.15), SLIDE_W, Inches(0.475))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x05, 0x05, 0x1A)
    bar.line.fill.background()
    # 設計コメント（太字・金色）
    tb(slide, f"設計：{design_comment}", 0.15, 5.17, 5.5, 0.4,
       font_name=FONT_B, font_size=9, color=C_GOLD, bold=True)
    # 補足説明（白・通常）
    tb(slide, text, 5.5, 5.17, 4.35, 0.4,
       font_name=FONT_B, font_size=8, color=C_GRAY,
       align=PP_ALIGN.RIGHT)


def badge(slide, text, left, top, w=1.4, h=0.32, bg=None, fc=None):
    """バッジ風ラベル"""
    rect_b(slide, left, top, w, h, bg or C_PINK, border_color=C_GOLD, border_pt=1)
    tb(slide, text, left + 0.05, top + 0.02, w - 0.1, h - 0.04,
       font_size=9, color=fc or C_WHITE, bold=True, align=PP_ALIGN.CENTER)


def flow_box(slide, label, left, top, w=1.55, h=0.45,
             fill=None, text_color=None, font_size=9):
    """フローチャート用ボックス"""
    fill = fill or C_BOX
    rect_b(slide, left, top, w, h, fill, border_color=C_PINK, border_pt=1.2)
    tb(slide, label, left + 0.05, top + 0.04, w - 0.1, h - 0.08,
       font_size=font_size, color=text_color or C_WHITE,
       bold=False, align=PP_ALIGN.CENTER)


# ─── スライド生成関数 ──────────────────────────────────

def slide1_title(prs):
    """スライド1：タイトル・スペック・この台の3ポイント"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.3, glow_y=0.5)
    hdr(slide,
        "スマスロ マギアレコード 魔法少女まどか☆マギカ外伝",
        "Part A: 機種解説")

    # 機種基本情報（左半分）
    rect_b(slide, 0.2, 0.78, 4.7, 2.30, C_BOX, border_color=C_PINK, border_pt=1.5)
    tb(slide, "基本スペック", 0.35, 0.82, 2.0, 0.35,
       font_name=FONT_H, font_size=12, color=C_LPINK, bold=True)

    specs = [
        ("メーカー", "ミズホ（ユニバーサルエンタ）"),
        ("導入日",   "2025年4月7日"),
        ("純増",     "約2.6枚/G（AT中）"),
        ("天井",     "ボーナス間950pt+α"),
        ("機械割",   "97.6〜114.9%（設定1〜6）"),
        ("AT確率",   "1/654.6〜1/416.7"),
    ]
    y = 1.18
    for k, v in specs:
        tb(slide, f"◆ {k}：", 0.35, y, 1.55, 0.28, font_size=9, color=C_LPINK, bold=True)
        tb(slide, v, 1.88, y, 2.85, 0.28, font_size=9, color=C_WHITE)
        y += 0.28

    # この台の3ポイント（右半分）
    rect_b(slide, 5.15, 0.78, 4.65, 2.30, C_BOX, border_color=C_CYAN, border_pt=1.5)
    tb(slide, "この台の3ポイント", 5.3, 0.82, 3.5, 0.35,
       font_name=FONT_H, font_size=12, color=C_CYAN, bold=True)

    pts = [
        ("①", "ストーリー×コンプリート体験",
         "8種ストーリーを集めてエンディングへ到達する物語体験型AT"),
        ("②", "穢れシステム「ドッペルモード」",
         "AT中に穢れが溜まると期待枚数3000枚超の爆発モードへ"),
        ("③", "まどマギ世界観を完全再現",
         "初代ファン向けサウンド・演出・キャラが深く作り込まれた完成度"),
    ]
    y = 1.18
    for num, title_pt, desc in pts:
        tb(slide, f"{num} {title_pt}", 5.3, y, 4.3, 0.28,
           font_size=10, color=C_GOLD, bold=True)
        tb(slide, desc, 5.3, y + 0.28, 4.3, 0.26,
           font_size=8, color=C_GRAY)
        y += 0.54   # 3ポイント全て＋バッジが収まるステップ

    # アワードバッジ（ボックスより後に描画→最前面。③の下に配置）
    badge(slide, "パチスロアワード 2025 SILVER", 5.25, 2.84, w=4.45, h=0.22,
          bg=RGBColor(0x44, 0x22, 0x00), fc=C_GOLD)

    # キャッチコピー
    rect_b(slide, 0.2, 3.12, 9.6, 0.52, RGBColor(0x22, 0x00, 0x22),
           border_color=C_PINK, border_pt=1)
    tb(slide,
       "初代『SLOT魔法少女まどか☆マギカ』のDNAを受け継ぎ、スマスロで進化──"
       "ストーリー攻略＋穢れ爆発の二重構造が今の等価市場で輝く",
       0.35, 3.16, 9.3, 0.44,
       font_size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ボーナス4種
    rect_b(slide, 0.2, 3.73, 9.6, 1.35, C_BOX, border_color=C_CYAN, border_pt=1.2)
    tb(slide, "ボーナス4種類（AT突入への足掛かり）",
       0.35, 3.77, 5.0, 0.3, font_size=10, color=C_CYAN, bold=True)
    bonuses = [
        ("エピソードボーナス", "50G継続", C_PINK),
        ("ビッグボーナス",     "30G+α",  C_GOLD),
        ("みたまボーナス",     "ベルナビ8回", C_CYAN),
        ("アリナミュージアム", "30G継続", C_LPINK),
    ]
    bx = 0.3
    for bname, bdesc, bc in bonuses:
        rect_b(slide, bx, 4.10, 2.28, 0.82, RGBColor(0x18, 0x08, 0x28),
               border_color=bc, border_pt=1.5)
        tb(slide, bname, bx + 0.08, 4.14, 2.1, 0.32,
           font_size=9, color=bc, bold=True, align=PP_ALIGN.CENTER)
        tb(slide, bdesc, bx + 0.08, 4.46, 2.1, 0.28,
           font_size=8.5, color=C_WHITE, align=PP_ALIGN.CENTER)
        bx += 2.38

    net_note(slide,
             "情報源: 一撃・ちょんぼりすた・パチセブン・P-WORLD",
             "スマスロ×物語体験×穢れ爆発の三位一体設計")


def slide2_gameflow(prs):
    """スライド2：ゲームフロー全体図（蛇行2段可視化）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.5, glow_y=0.3)
    hdr(slide, "ゲームフロー全体図", "全ルートを蛇行2段で可視化")

    _bw  = 1.50   # ボックス幅（前: 1.3in）
    _bh  = 0.80   # 上段高さ（前: 0.72in）
    _bh2 = 1.00   # 下段高さ（前: 0.82in）
    _gap = 0.13   # ボックス間隔
    _mx  = 0.12   # 左右マージン
    # 幅チェック: 0.12 + 6*1.50 + 5*0.13 + 末端 = 9.77in ✓

    _top1_y = 0.98
    _top2_y = 2.65
    _ann_y  = 3.90

    # ─ 上段ラベル ─
    tb(slide, "上段ルート（メインフロー）", _mx, 0.72, 5.0, 0.26,
       font_size=9, color=C_LPINK, bold=True)

    boxes_top = [
        ("通常時\n（ポイント蓄積）",  C_BOX,                    C_PINK),
        ("マギア\nチャレンジCZ",      RGBColor(0x20,0x10,0x30), C_PINK),
        ("ボーナス\n4種",             RGBColor(0x10,0x10,0x35), C_CYAN),
        ("マギア\nラッシュAT",        RGBColor(0x30,0x00,0x30), C_GOLD),
        ("ストーリー\nコンプリート",  RGBColor(0x15,0x25,0x10), C_GREEN),
        ("エンディング\n到達",        RGBColor(0x30,0x10,0x00), C_GOLD),
    ]
    for i, (label, fill, bc) in enumerate(boxes_top):
        lx = _mx + i * (_bw + _gap)
        rect_b(slide, lx, _top1_y, _bw, _bh, fill, border_color=bc, border_pt=1.5)
        tb(slide, label, lx + 0.06, _top1_y + 0.08, _bw - 0.12, _bh - 0.14,
           font_size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 天井アノテーション（通常時ボックス直下）
    tb(slide, "▼ 天井 950pt+α", _mx, _top1_y + _bh + 0.04, _bw, 0.22,
       font_size=8, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    # 上段矢印
    for i in range(len(boxes_top) - 1):
        ax = _mx + i * (_bw + _gap) + _bw + 0.01
        arrow_r(slide, ax, _top1_y + _bh / 2 - 0.12, length=_gap - 0.02)

    # 折り返し矢印（右端縦バー → 横バー）
    r_edge = _mx + 5 * (_bw + _gap) + _bw   # ≈ 9.77in

    fold_v = slide.shapes.add_shape(1,
        Inches(r_edge), Inches(_top1_y + _bh * 0.5),
        Inches(0.12), Inches(0.70))
    fold_v.fill.solid(); fold_v.fill.fore_color.rgb = C_PINK; fold_v.line.fill.background()

    fold_h_y = _top1_y + _bh * 0.5 + 0.70   # ≈ 2.08
    fold_h = slide.shapes.add_shape(1,
        Inches(_mx), Inches(fold_h_y),
        Inches(r_edge - _mx + 0.12), Inches(0.10))
    fold_h.fill.solid(); fold_h.fill.fore_color.rgb = C_PINK; fold_h.line.fill.background()

    tb(slide, "↩ 折り返し（エンディング → 上位AT突入）",
       3.3, fold_h_y + 0.02, 3.5, 0.22,
       font_size=8, color=C_PINK, align=PP_ALIGN.CENTER)

    # ─ 下段ラベル ─
    tb(slide, "下段ルート（上位AT・爆発ルート）", _mx, _top2_y - 0.26, 5.5, 0.26,
       font_size=9, color=C_GOLD, bold=True)

    boxes_bot = [
        ("エンブリオ\nイブ覚醒\n(上位AT)",    RGBColor(0x35,0x10,0x30), C_GOLD),
        ("エンブリオ\nイブアタック\n(特化)",   RGBColor(0x30,0x15,0x00), C_GOLD),
        ("決戦\n神浜聖女\n(ST特化)",           RGBColor(0x00,0x15,0x35), C_CYAN),
        ("ドッペル\nモード\n(穢れ爆発)",       RGBColor(0x35,0x00,0x08), C_RED),
        ("マギウス\nバトル",                   RGBColor(0x25,0x00,0x25), C_PINK),
        ("AT終了\n→有利区間\nリセット",        RGBColor(0x10,0x10,0x10), C_GRAY),
    ]
    for i, (label, fill, bc) in enumerate(boxes_bot):
        lx = _mx + i * (_bw + _gap)
        bpt = 2.5 if bc == C_RED else 1.5
        rect_b(slide, lx, _top2_y, _bw, _bh2, fill, border_color=bc, border_pt=bpt)
        tb(slide, label, lx + 0.06, _top2_y + 0.10, _bw - 0.12, _bh2 - 0.18,
           font_size=9.0, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 下段矢印
    for i in range(len(boxes_bot) - 1):
        ax = _mx + i * (_bw + _gap) + _bw + 0.01
        arrow_r(slide, ax, _top2_y + _bh2 / 2 - 0.12, length=_gap - 0.02)

    # ─ キーポイント 3パネル（凡例に代わる設計補足）─
    n_panels = 3
    gap_p = 0.14
    panel_w = (10.0 - _mx * 2 - gap_p * (n_panels - 1)) / n_panels   # ≈ 3.16in
    ann_items = [
        ("天井 950pt+α",
         "ポイント蓄積上限で強制ボーナス当選。スルー回数が増えるほど初当たりAT期待値が上昇",
         C_GOLD),
        ("ドッペル爆発 3000枚超",
         "穢れ解放で突入。AT終了まで上乗せ倍増。設定6以外でも3000枚超は十分現実的",
         C_RED),
        ("コンプ体験",
         "8ストーリーをATで収集してエンディングへ。明確なゴールがAT継続モチベを維持",
         C_CYAN),
    ]
    for j, (ak, av, ac) in enumerate(ann_items):
        lx = _mx + j * (panel_w + gap_p)
        rect_b(slide, lx, _ann_y, panel_w, 0.72, RGBColor(0x08, 0x08, 0x20),
               border_color=ac, border_pt=1)
        tb(slide, ak, lx + 0.09, _ann_y + 0.05, panel_w - 0.18, 0.26,
           font_size=9.5, color=ac, bold=True)
        tb(slide, av, lx + 0.09, _ann_y + 0.30, panel_w - 0.18, 0.38,
           font_size=7.5, color=C_GRAY)

    net_note(slide,
             "ゲームフロー：一撃/flick7/ちょんぼりすた 各解析より構成",
             "通常→CZ→ボーナス→AT→コンプ→上位ATの二段構造が特徴")


def slide3_normal(prs):
    """スライド3：通常時の遊び方"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.2, glow_y=0.6)
    hdr(slide, "通常時の遊び方", "全ルート・天井・魔法少女モード")

    # ─ 左：ポイント蓄積システム ─
    rect_b(slide, 0.15, 0.78, 4.6, 2.00, C_BOX, border_color=C_PINK, border_pt=1.5)
    tb(slide, "ポイント蓄積システム（メイン突入経路）",
       0.3, 0.82, 4.3, 0.32, font_size=10, color=C_PINK, bold=True)

    pt_items = [
        "毎ゲーム最低1pt以上獲得（液晶右下に表示）",
        "有利区間開始・ボーナス終了・AT終了時に規定ptを再抽選",
        "規定pt到達 → ボーナス当選（天井は950pt+α）",
        "レア役（強チェリー・強スイカ等）で直接ボーナス抽選あり",
        "魔法少女モードがAorBで抽選性能が変化",
    ]
    y = 1.18
    for item in pt_items:
        tb(slide, f"・{item}", 0.3, y, 4.3, 0.28, font_size=8.5, color=C_WHITE)
        y += 0.30

    # 天井（左列・ポイント蓄積の直下）
    rect_b(slide, 0.15, 2.85, 4.6, 0.65, RGBColor(0x20,0x00,0x20),
           border_color=C_GOLD, border_pt=1.8)
    tb(slide, "天井：ボーナス間 950pt+α 到達で強制当選",
       0.3, 2.89, 4.3, 0.28, font_size=9.5, color=C_GOLD, bold=True)
    tb(slide, "＊スルー回数が増えるほど初当たりAT期待値UP（スルー狙い有効）",
       0.3, 3.17, 4.3, 0.28, font_size=8, color=C_GRAY)

    # ─ 右：CZ種別と魔法少女モード ─
    rect_b(slide, 4.95, 0.78, 4.85, 1.90, C_BOX, border_color=C_CYAN, border_pt=1.5)
    tb(slide, "CZ種別", 5.1, 0.82, 2.0, 0.32,
       font_size=10, color=C_CYAN, bold=True)

    czs = [
        ("マギアチャレンジ（通常）", "スイカ/ポイント消化で突入", C_WHITE),
        ("黒江チャレンジ",          "突入した時点でAT濃厚！",   C_GOLD),
        ("前兆中の昇格抽選",        "マギアCZ→黒江CZへアップグレード", C_LPINK),
    ]
    y = 1.18
    for cn, cd, cc in czs:
        tb(slide, f"▶ {cn}", 5.1, y, 4.6, 0.24, font_size=9, color=cc, bold=True)
        tb(slide, f"   {cd}", 5.1, y + 0.24, 4.6, 0.22, font_size=8, color=C_GRAY)
        y += 0.46   # 0.54 → 0.46 で3項目がボックス内に収まる

    rect_b(slide, 4.95, 2.77, 4.85, 1.40, C_BOX, border_color=C_PINK, border_pt=1.5)
    tb(slide, "魔法少女モード（隠し内部モード）",
       5.1, 2.81, 4.6, 0.32, font_size=10, color=C_PINK, bold=True)

    modes = [
        ("モードA（通常）", "標準的な抽選テーブルで運用"),
        ("モードB（高確）", "CZ突入率・AT直撃率が大幅アップ"),
        ("モード移行",      "ボーナス/AT終了時・特定役で再抽選"),
    ]
    y = 3.17
    for mn, md in modes:
        tb(slide, f"・{mn}：", 5.1, y, 2.0, 0.28, font_size=9, color=C_LPINK, bold=True)
        tb(slide, md, 7.05, y, 2.7, 0.28, font_size=9, color=C_WHITE)
        y += 0.34

    # ─ 下：打ち方メモ ─
    rect_b(slide, 0.15, 4.24, 9.65, 0.88, RGBColor(0x0C,0x0C,0x28),
           border_color=C_CYAN, border_pt=1)
    tb(slide, "打ち方メモ", 0.3, 4.27, 1.5, 0.28,
       font_size=9, color=C_CYAN, bold=True)
    memo = [
        "通常時：左リール枠上〜上段にBARを目押し（チェリー・スイカ取得のため）",
        "ナビ発生時：押し順ナビに従う（ペナルティあり）",
        "朝一：設定変更後は有利区間ランプ消灯→規定ptリセット→狙い目",
    ]
    my = 4.50
    for m in memo:
        tb(slide, f"・{m}", 0.3, my, 9.3, 0.21, font_size=8, color=C_GRAY)
        my += 0.20

    net_note(slide,
             "通常時の核心：ポイント蓄積＋モード管理でCZ/ボーナスを引く",
             "黒江CZ=AT濃厚という演出設計がモチベーションを生む")


def slide4_cz(prs):
    """スライド4：CZ/前兆の仕組み"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.7, glow_y=0.4)
    hdr(slide, "CZ・前兆の仕組み", "突破条件と演出の見方")

    # ─ マギアチャレンジ概要 ─
    rect_b(slide, 0.15, 0.78, 5.95, 3.1, C_BOX, border_color=C_PINK, border_pt=1.5)
    tb(slide, "マギアチャレンジ（CZ）", 0.3, 0.82, 3.5, 0.35,
       font_name=FONT_H, font_size=12, color=C_PINK, bold=True)

    cz_info = [
        ("突入条件",
         "①規定ポイント消化  ②スイカ/強レア役  ③前兆経由",
         C_WHITE),
        ("ゲーム数",
         "最大30G程度（成功/失敗で抽選）",
         C_WHITE),
        ("突破条件",
         "CZ中のレア役成立・ゲーム数消化・特定演出成功",
         C_WHITE),
        ("成功時",
         "ボーナス当選（4種のいずれか）",
         C_GOLD),
        ("失敗時",
         "通常に戻る（ポイントは引き継ぎ）",
         C_GRAY),
    ]
    y = 1.22
    for k, v, vc in cz_info:
        tb(slide, f"◆ {k}", 0.3, y, 1.75, 0.28, font_size=9, color=C_LPINK, bold=True)
        tb(slide, v, 2.0, y, 3.95, 0.28, font_size=9, color=vc)
        y += 0.35

    # 黒江チャレンジ（AT濃厚CZ）
    rect_b(slide, 0.3, 2.45, 5.6, 0.88, RGBColor(0x28,0x18,0x00),
           border_color=C_GOLD, border_pt=2)
    tb(slide, "黒江チャレンジ ＝ AT濃厚！（最重要CZ）",
       0.45, 2.50, 5.3, 0.32, font_size=10, color=C_GOLD, bold=True)
    tb(slide, "マギアCZの本前兆中に昇格抽選あり。発生時点でAT当確レベルの強CZ",
       0.45, 2.82, 5.3, 0.28, font_size=8.5, color=C_WHITE)

    # CZ中の演出示唆
    rect_b(slide, 0.15, 3.40, 5.95, 0.72, C_BOX, border_color=C_CYAN, border_pt=1)
    tb(slide, "CZ中の演出示唆（見方）", 0.3, 3.44, 3.0, 0.28,
       font_size=9, color=C_CYAN, bold=True)
    hints = [
        "ホムラ登場 → 高期待度",
        "4人集合演出 → ボーナス濃厚",
        "テロップ色変化 → 昇格示唆",
    ]
    hx = 0.3
    for h in hints:
        tb(slide, f"・{h}", hx, 3.72, 2.7, 0.28, font_size=8.5, color=C_WHITE)
        hx += 2.0

    # ─ 右：前兆の見分け方 ─
    rect_b(slide, 6.25, 0.78, 3.6, 3.34, C_BOX, border_color=C_CYAN, border_pt=1.5)
    tb(slide, "前兆の種類と見分け方",
       6.4, 0.82, 3.3, 0.35, font_name=FONT_H, font_size=11, color=C_CYAN, bold=True)

    zencho = [
        ("本前兆（弱）", "画面色が薄く変化\n小役連続で少し期待", C_WHITE),
        ("本前兆（強）", "ステージ移行＋キャラ会話\nボーナス確率大幅UP", C_GOLD),
        ("ガセ前兆",    "通常時にも発生\n演出で見分けが難しい", C_GRAY),
        ("黒江前兆",    "固有演出で出現\nAT確定レベル",          C_GOLD),
    ]
    zy = 1.22
    for zt, zd, zc in zencho:
        rect_b(slide, 6.35, zy, 3.4, 0.72, RGBColor(0x12,0x12,0x2E),
               border_color=zc, border_pt=1)
        tb(slide, zt, 6.5, zy + 0.05, 3.1, 0.28,
           font_size=9, color=zc, bold=True)
        tb(slide, zd, 6.5, zy + 0.30, 3.1, 0.36,
           font_size=8, color=C_WHITE)
        zy += 0.78

    net_note(slide,
             "CZ設計：弱→強→黒江の段階的期待感とガセの使い方が秀逸",
             "黒江チャレンジのAT確定演出が「打ち続ける動機」を作る")


def slide5_at(prs):
    """スライド5：AT/ボーナス（何をすれば出玉が伸びる）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.5, glow_y=0.5)
    hdr(slide, "AT「マギアラッシュ」の仕組み", "何をすれば出玉が伸びる？")

    # AT概要
    rect_b(slide, 0.15, 0.78, 9.7, 0.6, RGBColor(0x20,0x00,0x28),
           border_color=C_PINK, border_pt=1.5)
    info_items = [
        ("純増", "約2.6枚/G"),
        ("タイプ", "ゲーム数管理型"),
        ("初期G数", "マギアアタックで決定"),
        ("最大目標", "ストーリー8種コンプ"),
    ]
    ix = 0.25
    for k, v in info_items:
        tb(slide, k, ix, 0.82, 1.2, 0.24, font_size=8, color=C_LPINK, bold=True)
        tb(slide, v, ix, 1.06, 1.8, 0.24, font_size=9, color=C_WHITE, bold=False)
        ix += 2.4

    # ─ 上乗せ手段（左）─
    rect_b(slide, 0.15, 1.50, 4.7, 2.72, C_BOX, border_color=C_GOLD, border_pt=1.5)
    tb(slide, "上乗せ手段（出玉の伸ばし方）",
       0.3, 1.54, 4.3, 0.32, font_size=10, color=C_GOLD, bold=True)

    jouze = [
        ("レア役直乗せ",     "5〜100G\nレア役成立時に直接G数上乗せ",        C_WHITE),
        ("ストーリー当選",   "コンプ数UP\n全8種集めるとエンディング権利",    C_CYAN),
        ("決戦神浜聖女",     "ST型特化ゾーン\n成立役で連続G数上乗せ",        C_GOLD),
        ("マギウスバトル",   "バトル勝利でAT継続/上乗せ\n敗北で転落リスク", C_PINK),
        ("ドッペルモード",   "穢れ解放で突入\n上乗せ倍増・3000枚超期待",    C_RED),
    ]
    jy = 1.90
    for jn, jd, jc in jouze:
        rect_b(slide, 0.25, jy, 4.5, 0.47, RGBColor(0x12,0x10,0x28),
               border_color=jc, border_pt=1)
        tb(slide, jn, 0.38, jy + 0.04, 1.55, 0.40,
           font_size=9, color=jc, bold=True)
        tb(slide, jd, 1.95, jy + 0.04, 2.65, 0.40,
           font_size=8, color=C_WHITE)
        jy += 0.53

    # ─ ストーリーシステム（右）─
    rect_b(slide, 5.0, 1.50, 4.85, 2.72, C_BOX, border_color=C_CYAN, border_pt=1.5)
    tb(slide, "ストーリーシステム（AT継続の核）",
       5.15, 1.54, 4.5, 0.32, font_size=10, color=C_CYAN, bold=True)

    tb(slide, "AT中に8種のストーリーを収集し、コンプリートでエンディング→上位AT権利獲得",
       5.15, 1.90, 4.6, 0.42, font_size=8.5, color=C_WHITE)

    stories = [
        "① 覚醒のマギア", "② 禁断の実験",
        "③ 闇のウワサ",   "④ 聖女の誓い",
        "⑤ 魔女の誘惑",   "⑥ 因果の轍",
        "⑦ 決意の代償",   "⑧ 真実の扉（コンプ）",
    ]
    sx = 0
    sy = 2.38
    for i, s in enumerate(stories):
        col = C_GOLD if i == 7 else C_WHITE
        bg  = RGBColor(0x30, 0x15, 0x00) if i == 7 else RGBColor(0x12, 0x12, 0x30)
        rect_b(slide, 5.1 + sx * 2.35, sy, 2.28, 0.30,
               bg, border_color=col, border_pt=0.8)
        tb(slide, s, 5.15 + sx * 2.35, sy + 0.02, 2.2, 0.28,
           font_size=8.5, color=col, align=PP_ALIGN.CENTER)
        sx = 1 - sx
        if sx == 0:
            sy += 0.36

    # マギアアタック説明
    rect_b(slide, 5.0, 4.25, 4.85, 0.67, RGBColor(0x1A,0x00,0x1A),
           border_color=C_PINK, border_pt=1)
    tb(slide, "マギアアタック（AT開幕）",
       5.15, 4.29, 4.6, 0.28, font_size=9.5, color=C_PINK, bold=True)
    tb(slide, "キャラカード×属性×成立役で初期G数を決定。対応役を引けば上乗せ優遇",
       5.15, 4.57, 4.6, 0.28, font_size=8, color=C_GRAY)

    net_note(slide,
             "出玉設計：G数直乗せ×ストーリーコンプ×ドッペル爆発の三段階",
             "「引き続ける疲労感」と「コンプ達成感」がトレードオフ")


def slide6_upper_at(prs):
    """スライド6：上位ATへの道と遊び方"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.6, glow_y=0.3)
    hdr(slide, "上位ATへの道と遊び方", "エンブリオ・イブ覚醒 ＆ ドッペルモード")

    # 左：エンブリオイブ覚醒
    rect_b(slide, 0.15, 0.78, 5.75, 3.5, C_BOX, border_color=C_GOLD, border_pt=2)
    badge(slide, "エンブリオ・イブ覚醒", 0.25, 0.82, w=2.8, h=0.38,
          bg=RGBColor(0x40, 0x18, 0x00), fc=C_GOLD)

    tb(slide, "ストーリー8種コンプリートでエンディング到達後に突入する上位AT",
       0.3, 1.28, 5.5, 0.3, font_size=9, color=C_WHITE)

    steps = [
        ("STEP 1", "ストーリー8種をAT中に収集（コンプ）", C_GOLD),
        ("STEP 2", "エンディング（コンプリート報酬）到達", C_GOLD),
        ("STEP 3", "エンブリオ・イブ覚醒準備（32G）※左第一停止厳守", C_LPINK),
        ("STEP 4", "エンブリオ・イブアタック突入（10G+α特化）", C_GOLD),
        ("STEP 5", "特化中に大量G数を上乗せ → 超長期AT突入", C_GREEN),
    ]
    sy = 1.55   # 少し上から開始して警告テキストの余白を確保
    for sn, sd, sc in steps:
        rect_b(slide, 0.25, sy, 1.0, 0.36, RGBColor(0x28,0x10,0x00),
               border_color=sc, border_pt=1)
        tb(slide, sn, 0.27, sy + 0.03, 0.96, 0.30,
           font_size=8.5, color=sc, bold=True, align=PP_ALIGN.CENTER)
        tb(slide, sd, 1.32, sy + 0.05, 4.45, 0.30, font_size=8.5, color=C_WHITE)
        if sy < 3.5:
            fold_bar = slide.shapes.add_shape(1,
                Inches(0.68), Inches(sy + 0.36), Inches(0.12), Inches(0.12))
            fold_bar.fill.solid(); fold_bar.fill.fore_color.rgb = sc
            fold_bar.line.fill.background()
        sy += 0.48

    tb(slide, "※準備区間中はペナルティあり。必ず左第一停止を守ること！",
       0.3, 4.00, 5.5, 0.28, font_size=8, color=C_RED, bold=True)

    # 右：ドッペルモード
    rect_b(slide, 6.05, 0.78, 3.72, 3.5, C_BOX, border_color=C_RED, border_pt=2)
    badge(slide, "ドッペルモード（穢れ爆発）", 6.15, 0.82, w=3.52, h=0.38,
          bg=RGBColor(0x35, 0x00, 0x10), fc=C_RED)

    doppel_info = [
        ("発動条件", "穢れが一定量蓄積→解放抽選通過"),
        ("継続期間", "AT終了まで継続（永続）"),
        ("上乗せ性能", "通常の数倍の上乗せ期待値"),
        ("期待枚数", "3000枚超（6以外でも到達可）"),
        ("穢れの蓄積", "通常時の特定演出・負けた回数で加算"),
    ]
    dy = 1.28
    for dk, dv in doppel_info:
        tb(slide, f"◆ {dk}：", 6.15, dy, 1.3, 0.3, font_size=9, color=C_LPINK, bold=True)
        tb(slide, dv, 7.45, dy, 2.3, 0.3, font_size=8.5, color=C_WHITE)
        dy += 0.37

    # ドッペルモード演出ガイド
    rect_b(slide, 6.15, 3.05, 3.52, 0.95, RGBColor(0x25, 0x00, 0x08),
           border_color=C_RED, border_pt=1)
    tb(slide, "ドッペルモード演出サイン",
       6.25, 3.09, 3.3, 0.28, font_size=9, color=C_RED, bold=True)
    tb(slide, "・画面が赤黒く染まる\n・魔女側のBGMに切り替わる\n・上乗せ数字が激増",
       6.25, 3.37, 3.4, 0.55, font_size=8, color=C_WHITE)

    net_note(slide,
             "上位AT設計：コンプリート達成感 ＋ ドッペル爆発の二軸で中長期モチベ維持",
             "穢れは「負けた積み上げが報われる」逆転設計")


def slide7_design(prs):
    """スライド7：面白さの設計（なぜアワードSILVERを獲れたのか）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.5, glow_y=0.5)
    hdr(slide, "面白さの設計分析", "なぜパチスロアワード2025 SILVERを獲れたのか")

    # アワード情報バナー
    rect_b(slide, 0.15, 0.78, 9.7, 0.46, RGBColor(0x30, 0x15, 0x00),
           border_color=C_GOLD, border_pt=1.5)
    tb(slide,
       "パチスロアワード2025 SILVER受賞 ──「遊技台としての完成度の高さ」と「現在のユーザーニーズへの合致」が評価",
       0.3, 0.83, 9.3, 0.36, font_size=9.5, color=C_GOLD, bold=True,
       align=PP_ALIGN.CENTER)

    # 5つの設計要素
    designs = [
        ("物語体験の設計",
         "8種ストーリーを1AT中に少しずつ集める「コレクション欲求」を刺激。"
         "コンプリートという明確なゴールがAT継続中ずっとプレイヤーを引きつける。",
         C_CYAN, "🎯 ゴール設計"),
        ("逆転設計（穢れシステム）",
         "通常時の負けが「穢れ」として蓄積し、ドッペルモードで爆発する設計。"
         "「負けるほど強くなる」という逆転快感がロングセッションを促進。",
         C_RED, "🔄 逆転快感"),
        ("マイルドな投資/出玉速度",
         "純増2.6枚/Gというスマスロ基準では控えめな数値が「疲れにくいペース」を実現。"
         "ゆったり打ちたい層・ライトユーザー需要に合致した遊技体験。",
         C_GREEN, "⏱ 時間設計"),
        ("IPとゲーム性の融合",
         "まどかマギカの「魔法少女の穢れと絶望」というテーマをドッペルモードに昇華。"
         "世界観とゲーム性が一致しているため演出に説得力がある。",
         C_LPINK, "🌟 IP活用"),
        ("初代ファンへの継承",
         "初代SLOT魔法少女まどか☆マギカのサウンド・キャラ・演出を踏襲しつつ"
         "スマスロで新機軸を加えた「懐かし×新鮮」のバランス設計。",
         C_GOLD, "♻ 継承設計"),
    ]

    dx = 0.15
    dy_start = 1.35
    for i, (dt, dd, dc, tag) in enumerate(designs):
        col = i % 2
        row = i // 2
        lx = 0.15 + col * 4.9
        ly = dy_start + row * 1.15
        if i == 4:  # 最後は中央
            lx = 2.55
        rect_b(slide, lx, ly, 4.7, 1.05, C_BOX, border_color=dc, border_pt=1.5)
        tb(slide, f"{tag}  {dt}", lx + 0.1, ly + 0.04, 4.5, 0.3,
           font_size=9.5, color=dc, bold=True)
        tb(slide, dd, lx + 0.1, ly + 0.32, 4.5, 0.68,
           font_size=8, color=C_WHITE)

    net_note(slide,
             "受賞の核：テーマ×逆転×コレクションの三位一体が完成度を押し上げた",
             "等価市場でも戦える機械割97.6〜114.9%が後押し")


def slide8_pros_cons(prs):
    """スライド8：良い点と課題"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.4, glow_y=0.4)
    hdr(slide, "良い点と課題", "Part B: ユーザー評価＋設計視点の分析")

    # ─ 良い点（左）─
    rect_b(slide, 0.15, 0.78, 4.7, 4.15, C_BOX, border_color=C_GREEN, border_pt=2)
    badge(slide, "良い点 PROS", 0.25, 0.82, w=2.2, h=0.38,
          bg=RGBColor(0x00, 0x2A, 0x10), fc=C_GREEN)

    pros = [
        ("世界観の完成度", "まどかマギカIPを機種性能と直結させた演出作りが秀逸"),
        ("コンプリート体験", "8ストーリー集めの達成感が他機種にない一体感を生む"),
        ("穢れ→ドッペル", "「負けが報われる」逆転設計がロングセッションを促進"),
        ("マイルドな純増", "2.6枚/GでATが長く続く安心感。疲れにくい遊技体験"),
        ("設定差の明確さ", "機械割97.6〜114.9%の幅で高設定を打つ意義が明確"),
        ("万枚実績",       "稼働初期から万枚・コンプリート報告が相次いだ"),
    ]
    py = 1.28
    for pk, pv in pros:
        rect_b(slide, 0.25, py, 4.5, 0.52, RGBColor(0x08, 0x1E, 0x10),
               border_color=C_GREEN, border_pt=0.8)
        tb(slide, f"✓ {pk}", 0.38, py + 0.04, 1.7, 0.24,
           font_size=9, color=C_GREEN, bold=True)
        tb(slide, pv, 2.05, py + 0.04, 2.65, 0.44, font_size=8, color=C_WHITE)
        py += 0.60

    # ─ 課題（右）─
    rect_b(slide, 5.1, 0.78, 4.7, 4.15, C_BOX, border_color=C_RED, border_pt=2)
    badge(slide, "課題 CONS", 5.2, 0.82, w=2.0, h=0.38,
          bg=RGBColor(0x2A, 0x00, 0x00), fc=C_RED)

    cons = [
        ("疲労感の強さ", "常に何かを引き続けないとATが伸びない構造。集中力消耗が大きい"),
        ("ライト層の離脱", "ストーリーコンプの難易度が初心者には高く感じられる"),
        ("ガセ前兆の多さ", "通常時のガセ前兆が多く「空振りが続く」ストレスが蓄積"),
        ("準備区間ペナ", "エンブリオイブ覚醒準備中のペナルティが知らないと損失大"),
        ("ドッペル依存度", "爆発力がドッペルモード頼みで非突入時の出玉がやや物足りない"),
        ("IP知識格差", "原作未視聴者には演出の背景が伝わりにくく没入感に差が出る"),
    ]
    cy2 = 1.28
    for ck, cv in cons:
        rect_b(slide, 5.2, cy2, 4.5, 0.52, RGBColor(0x1E, 0x06, 0x06),
               border_color=C_RED, border_pt=0.8)
        tb(slide, f"✗ {ck}", 5.33, cy2 + 0.04, 1.7, 0.24,
           font_size=9, color=C_RED, bold=True)
        tb(slide, cv, 7.0, cy2 + 0.04, 2.65, 0.44, font_size=8, color=C_WHITE)
        cy2 += 0.60

    net_note(slide,
             "総評：コアユーザー向け高完成度機。ライト向け間口の広げ方が次世代への課題",
             "情報源: アミュタメ/DMM PALOR/ONEHALLサービス/パーラーフルスロットル")


def slide9_summary(prs):
    """スライド9：まとめ・設計から学べること"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_bg(slide, glow_x=0.5, glow_y=0.6)
    hdr(slide, "まとめ・設計から学べること", "スロクリ設計視点からの総括")

    # ─ 3つの学び（上段）─
    tb(slide, "この機種が示す「現代スマスロ設計の成功方程式」",
       0.3, 0.78, 9.3, 0.36, font_name=FONT_H, font_size=13,
       color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)

    lessons = [
        ("学び①\nゴール設計", C_CYAN,
         "コンプリートという「明確なゴール」を設けることで、AT中のプレイヤーの"
         "動機が持続する。ゴールが見えるゲームは打ち手を能動的にする。"),
        ("学び②\n逆転設計", C_RED,
         "「穢れ」のように不利な状況を貯蓄して有利に転換する仕組みは、"
         "敗北体験を「投資感」に変換し、継続遊技へのモチベを生み出す。"),
        ("学び③\nIP×ゲーム性", C_GOLD,
         "テーマ（魔法少女の絶望と穢れ）とゲームメカニクス（ドッペルモード）を"
         "直結させることで「演出に説得力」が生まれ、世界観浸透度が向上する。"),
    ]
    lx = 0.15
    for lt, lc, ld in lessons:
        rect_b(slide, lx, 1.22, 3.15, 1.75, C_BOX, border_color=lc, border_pt=2)
        tb(slide, lt, lx + 0.1, 1.26, 2.9, 0.45,
           font_name=FONT_H, font_size=12, color=lc, bold=True,
           align=PP_ALIGN.CENTER)
        tb(slide, ld, lx + 0.1, 1.70, 2.9, 1.20, font_size=8, color=C_WHITE)
        lx += 3.28

    # ─ 総括スコア（中段）─
    rect_b(slide, 0.15, 3.10, 9.7, 1.72, RGBColor(0x08,0x08,0x20),
           border_color=C_GRAY, border_pt=0.8)
    tb(slide, "総括スコア（設計視点）",
       0.3, 3.14, 3.0, 0.3, font_size=10, color=C_GRAY, bold=True)

    scores = [
        ("ゲーム性完成度",  "★★★★★", C_GOLD),
        ("出玉設計",        "★★★★☆", C_GREEN),
        ("演出・世界観",    "★★★★★", C_PINK),
        ("ライト向け間口",  "★★★☆☆", C_CYAN),
        ("設定差",          "★★★★☆", C_LPINK),
    ]
    sx = 0.3
    for sk, sv, sc in scores:
        tb(slide, sk, sx, 3.48, 1.75, 0.28, font_size=8.5, color=C_GRAY, bold=True)
        tb(slide, sv, sx, 3.76, 1.75, 0.32, font_size=10, color=sc, bold=True)
        sx += 1.95

    # 総合評価
    rect_b(slide, 0.15, 4.44, 9.7, 0.55, RGBColor(0x30, 0x10, 0x00),
           border_color=C_GOLD, border_pt=1.5)
    tb(slide,
       "総合評価：IP活用×逆転設計×コレクション体験を高水準で統合した、"
       "現代スマスロの教科書的設計機。コア層には強く刺さり、アワードSILVER受賞は必然。",
       0.3, 4.48, 9.3, 0.46,
       font_size=9, color=C_GOLD, bold=False, align=PP_ALIGN.CENTER)

    net_note(slide,
             "設計の本質：「負け→蓄積→爆発」のサイクルと「物語の完結」が共存する稀有な設計",
             "次世代機への示唆：ライト向け短期セッションへの対応が残課題")


# ─── main ────────────────────────────────────────────

def main():
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

    slide1_title(prs)
    slide2_gameflow(prs)
    slide3_normal(prs)
    slide4_cz(prs)
    slide5_at(prs)
    slide6_upper_at(prs)
    slide7_design(prs)
    slide8_pros_cons(prs)
    slide9_summary(prs)

    prs.save(OUT_PATH)
    print(f"保存完了: {OUT_PATH}")


if __name__ == "__main__":
    main()
