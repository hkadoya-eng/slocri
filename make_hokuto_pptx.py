"""
スマスロ 北斗の拳 機種説明+分析 統合版 PowerPointジェネレーター
出力: proposals/機種分析/北斗の拳/hokuto_guide_v1.pptx
テーマ: 黒×金×赤（北斗カラー）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "北斗の拳", "hokuto_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（黒×金×赤：北斗カラー）────────────────────────
C_BG    = RGBColor(0x05, 0x08, 0x18)   # 暗黒紺
C_CARD  = RGBColor(0x0C, 0x12, 0x28)   # カード背景
C_CARD2 = RGBColor(0x14, 0x1C, 0x38)   # やや明るいカード
C_ROW   = RGBColor(0x10, 0x16, 0x30)   # テーブル偶数行
C_RED   = RGBColor(0xBB, 0x11, 0x11)   # 血赤
C_RED2  = RGBColor(0x88, 0x00, 0x00)   # 濃赤
C_CRIM  = RGBColor(0xFF, 0x33, 0x33)   # 明るい赤（強調）
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)   # 明るい金
C_WHITE = RGBColor(0xE8, 0xE8, 0xF4)   # 本文白
C_CREAM = RGBColor(0xD0, 0xC0, 0xA0)   # クリーム
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)   # グレー
C_LTGRY = RGBColor(0x44, 0x44, 0x66)   # 薄グレー（ダーク版）
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_BLUE  = RGBColor(0x22, 0x77, 0xFF)
C_TEAL  = RGBColor(0x22, 0xAA, 0x99)
C_PUR   = RGBColor(0x88, 0x44, 0xCC)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（暗黒紺×斜めライン×底部赤グロー）────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (5, 8, 24))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(8, 12, 30), width=1)
    for y in range(h - 80, h):
        t = (y - (h - 80)) / 80
        draw.line([(0, y), (w, y)], fill=(int(20 * t), 0, 0))
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(0, 0, int(8 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー ──────────────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_RED)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GOLD, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_RED)


def net_note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


def footer(slide, bold_text, sub_text, col=None):
    """各スライドのフッター（設計コメント太字＋補足説明）"""
    c = col or C_RED
    rect(slide, Inches(0.2), Inches(4.42), Inches(9.6), Emu(580000),
         RGBColor(0x06, 0x08, 0x18))
    rect(slide, Inches(0.2), Inches(4.42), Emu(55000), Emu(580000), c)
    tb(slide, Inches(0.45), Inches(4.47), Inches(9.1), Emu(250000),
       bold_text, 9.5, bold=True, color=c)
    tb(slide, Inches(0.45), Inches(4.80), Inches(9.1), Emu(250000),
       sub_text, 8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x02, 0x05, 0x14))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_RED)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_RED2)

    tb(s, Inches(0.22), Inches(0.52), Inches(5.0), Emu(330000),
       "スマスロ 機種説明＋分析", 11, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.98), Inches(5.1), Emu(900000),
       "北斗の拳", 42, bold=True, color=C_CRIM, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.82), Inches(5.0), Emu(270000),
       "── 4号機世代が帰還した伝説の一台", 10, color=C_CREAM, font=FONT_H)

    # スペック
    specs = [
        ("メーカー",     "サミー"),
        ("導入",         "2023年4月3日"),
        ("設定",         "1〜6段階"),
        ("機械割",       "設定1: 98.0%  /  設定6: 113.0%"),
        ("天井",         "最大1268G → BB確定"),
        ("BB継続率",     "66 / 79 / 84 / 89%（4段階）"),
        ("無想転生バトル", "継続率94% / 期待2500枚以上"),
    ]
    sy = Inches(3.20)
    for k, v in specs:
        tb(s, Inches(0.22), sy, Inches(1.55), Emu(230000),
           k, 8, bold=True, color=C_GRAY, wrap=False)
        tb(s, Inches(1.82), sy, Inches(3.35), Emu(230000),
           v, 8, color=C_CREAM, wrap=False)
        sy += Emu(235000)

    # 右：3つの分析ポイント
    kws = [
        (C_RED,   "バトルボーナス型AT",
                  "BB×継続率4段階×Vストックの\n三位一体で出玉が伸びる"),
        (C_GOLD,  "無想転生バトル",
                  "継続率94%・最大出玉の舞台\n来店の最大目標として機能"),
        (C_TEAL,  "IP×世代回帰",
                  "30〜40代休眠層がスマスロで\n帰ってきた希有な事例"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.7 + i * 1.55)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.28), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(55000), Inches(1.28), ac)
        tb(s, Inches(5.85), y0 + Emu(55000), Inches(3.8), Emu(310000),
           kw, 12, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(380000), Inches(3.8), Emu(420000),
           desc, 8.5, color=C_WHITE)

    tb(s, Inches(7.8), Inches(5.18), Inches(2.0), Emu(270000),
       "v1.0  2026.05", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段で全ルートを可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図  ──  通常時 → AT → 上位ATへの全ルート", "2/10")

    # ── Row1（左→右）: 通常時 → 宿命バトル → BB → Vストック ─────────
    BW  = Inches(2.2)
    GAP = Inches(0.22)
    R1Y = Inches(0.70)
    BH1 = Inches(1.68)

    boxes_r1 = [
        (Inches(0.18), RGBColor(0x06, 0x08, 0x1C), C_GOLD,
         "通常時",
         "チャンス役成立を待つ\n中段チェリー・リーチ目は\nアミババトル確定（AT直行）"),
        (Inches(0.18) + BW + GAP, RGBColor(0x1A, 0x04, 0x04), C_RED,
         "宿命バトル（8G）",
         "リプレイ→チャンス\nレア役→勝利確定\n8G目レア役→ユリア復活"),
        (Inches(0.18) + 2*(BW+GAP), RGBColor(0x0A, 0x10, 0x04), C_GREEN,
         "バトルボーナス(BB)",
         "小役パート30G\n＋バトルパート8G\n1セット≈110枚"),
        (Inches(0.18) + 3*(BW+GAP), RGBColor(0x0C, 0x14, 0x28), C_TEAL,
         "Vストック",
         "次セット継続確定\nレイ/トキ協力で\nループ率UP"),
    ]
    for bx, fill, ac, title, desc in boxes_r1:
        rect_b(s, bx, R1Y, BW, BH1, fill, ac, 1.8)
        tb(s, bx + Emu(70000), R1Y + Emu(50000), BW - Emu(130000), Emu(290000),
           title, 10, bold=True, color=ac, font=FONT_H)
        tb(s, bx + Emu(60000), R1Y + Emu(340000), BW - Emu(110000), Emu(980000),
           desc, 8, color=C_CREAM)

    # Row1矢印
    for i in range(3):
        ax = Inches(0.18) + (i+1)*(BW+GAP) - GAP + Emu(40000)
        arrow_r(s, ax, R1Y + BH1//2, C_RED)

    # 「2ルート」ラベル（通常時の下）
    tb(s, Inches(0.18), R1Y + BH1 + Emu(30000), BW, Emu(230000),
       "① チャンス役→宿命バトル\n② 中段チェリー/リーチ目→アミバ確定",
       7, color=C_GOLD)

    # ⊓ コネクター（Row1右端→Row2右端）
    CON_X = Inches(0.18) + 4*(BW+GAP) - GAP + Emu(20000)
    CON_R = CON_X + Emu(650000)
    LW    = Emu(50000)
    AT_MID = R1Y + BH1//2
    R2Y   = Inches(2.58)
    BH2   = Inches(1.12)
    MID_Y = (R1Y + BH1 + R2Y) // 2
    rect(s, CON_X, AT_MID - LW//2, LW, MID_Y - AT_MID + LW, C_GOLD)
    rect(s, CON_X, MID_Y, CON_R - CON_X + LW, LW, C_GOLD)
    rect(s, CON_R, MID_Y, LW, R2Y + BH2//2 - MID_Y + LW, C_GOLD)
    tb(s, CON_X + Emu(60000), MID_Y - Emu(310000), Emu(590000), Emu(300000),
       "継続\nループ", 8, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)

    # ── Row2（右→左）: 無想転生チャンス → 無想転生バトル → 修羅の国 ──
    boxes_r2 = [
        (Inches(0.18) + 3*(BW+GAP), RGBColor(0x14, 0x06, 0x00), C_GOLD,
         "無想転生チャンス(15G)",
         "BB消化後に約33%で突入\n突入でBB中の神演出が発生"),
        (Inches(0.18) + 2*(BW+GAP), RGBColor(0x1A, 0x02, 0x02), C_CRIM,
         "無想転生バトル",
         "継続率94%\n期待2500枚以上\n「真の頂点」体験"),
        (Inches(0.18) + 1*(BW+GAP), RGBColor(0x14, 0x08, 0x24), C_PUR,
         "修羅の国",
         "高純増・高継続の\n上位AT\nさらに血涙の章あり"),
        (Inches(0.18), RGBColor(0x04, 0x0C, 0x04), C_GREEN,
         "天井（最大1268G）",
         "到達でBB確定\n777G通過で\n北斗揃い高確率"),
    ]
    for bx, fill, ac, title, desc in boxes_r2:
        rect_b(s, bx, R2Y, BW, BH2, fill, ac, 1.8)
        tb(s, bx + Emu(70000), R2Y + Emu(38000), BW - Emu(130000), Emu(265000),
           title, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, bx + Emu(60000), R2Y + Emu(305000), BW - Emu(110000), BH2 - Emu(360000),
           desc, 8, color=C_CREAM)

    # Row2矢印（←逆方向なのでrotation=180）
    for i in range(3):
        ax = Inches(0.18) + (i+1)*(BW+GAP) - GAP + Emu(50000)
        _w = GAP - Emu(100000)
        _h = Emu(150000)
        shp = s.shapes.add_shape(13, ax, R2Y + BH2//2 - _h//2, _w, _h)
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = C_CRIM
        shp.line.fill.background()

    # フッター
    footer(s,
           "フロー設計の核心：BB継続率4段階＋Vストックが「次も続くかも」を繰り返す設計",
           "宿命バトルで自力感→BB継続で期待感→無想転生で頂点体験。3段階の感情設計が長期稼働を支える。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（2ルート明確に記載）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方  ──  2つのルートでBBを目指す", "3/10")

    # ── ルート①（左）: チャンス役 → 宿命バトル ────────────────────
    rect_b(s, Inches(0.2), Inches(0.72), Inches(4.5), Inches(3.55),
           C_CARD, C_GOLD, 2.0)
    rect(s, Inches(0.2), Inches(0.72), Emu(55000), Inches(3.55), C_GOLD)
    tb(s, Inches(0.45), Inches(0.78), Inches(4.1), Emu(310000),
       "ルート① チャンス役 → 宿命バトル", 11, bold=True, color=C_GOLD, font=FONT_H)

    r1_items = [
        (C_GOLD,  "契機となる役",
                  "スイカ・チャンス役が成立したら宿命バトル抽選\n強レア役ほど突入しやすい"),
        (C_WHITE, "宿命バトルの目的",
                  "8Gバトル中にリプレイ・レア役を引いてケンシロウを勝利させる\n→ AT（バトルボーナス）に当選！"),
        (C_CREAM, "天井まで打てる",
                  "最大1268Gでバトルボーナス確定\n777G通過で北斗揃い（高継続）の確率UP\n焦らずチャンス役を引き続ける展開"),
    ]
    y1 = Inches(1.30)
    for ac, title, body in r1_items:
        rect(s, Inches(0.28), y1, Emu(50000), Emu(990000), ac)
        tb(s, Inches(0.48), y1 + Emu(28000), Inches(4.0), Emu(270000),
           title, 9, bold=True, color=ac)
        tb(s, Inches(0.48), y1 + Emu(295000), Inches(4.1), Emu(650000),
           body, 8.5, color=C_WHITE)
        y1 += Emu(1050000)

    # ── ルート②（右）: 中段チェリー/リーチ目 → アミバ確定 ──────────
    rect_b(s, Inches(5.0), Inches(0.72), Inches(4.75), Inches(3.55),
           RGBColor(0x1A, 0x04, 0x04), C_CRIM, 2.0)
    rect(s, Inches(5.0), Inches(0.72), Emu(55000), Inches(3.55), C_CRIM)
    tb(s, Inches(5.25), Inches(0.78), Inches(4.3), Emu(310000),
       "ルート② 中段チェリー/リーチ目 → アミババトル確定", 11, bold=True, color=C_CRIM, font=FONT_H)

    r2_items = [
        (C_CRIM,  "契機となる役",
                  "中段チェリー・リーチ目リプレイが成立\nこれだけでAT当選が確定する！"),
        (C_GOLD2, "アミババトルとは",
                  "ケンシロウがアミバと戦う特別バトル\n勝利で即BBに突入するルート\n宿命バトルを経由しない最速経路"),
        (C_WHITE, "引いた瞬間の体験",
                  "「確定した」という安堵と喜びが生まれる\n設置確認できたら積極的に期待してOK\nリーチ目を覚えると楽しさが倍増"),
    ]
    y2 = Inches(1.30)
    for ac, title, body in r2_items:
        rect(s, Inches(5.08), y2, Emu(50000), Emu(990000), ac)
        tb(s, Inches(5.28), y2 + Emu(28000), Inches(4.25), Emu(270000),
           title, 9, bold=True, color=ac)
        tb(s, Inches(5.28), y2 + Emu(295000), Inches(4.35), Emu(650000),
           body, 8.5, color=C_WHITE)
        y2 += Emu(1050000)

    # 中央の区切り線
    rect(s, Inches(4.88), Inches(0.72), Emu(18000), Inches(3.55), C_LTGRY)
    tb(s, Inches(4.70), Inches(1.95), Emu(320000), Emu(350000),
       "OR", 12, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER, font=FONT_H)

    footer(s,
           "通常時は「チャンス役成立の瞬間」と「中段チェリー/リーチ目の瞬間」の2段構えで楽しめる",
           "ルート①は継続的なチャンス役を引くドキドキ感、ルート②は確定役の爽快感。どちらのルートも明確な手応えがある。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 宿命バトル（突破の仕方・役の意味）
# ══════════════════════════════════════════════════════════════
def s_battle(prs):
    s = new_slide(prs)
    hdr(s, "宿命バトル  ──  8Gで自分の役がケンシロウを勝たせる", "4/10")

    # ── 左：役ごとの意味 ───────────────────────────────────────
    rect_b(s, Inches(0.2), Inches(0.72), Inches(4.55), Inches(3.55),
           C_CARD, C_RED, 1.5)
    rect(s, Inches(0.2), Inches(0.72), Emu(55000), Inches(3.55), C_RED)
    tb(s, Inches(0.45), Inches(0.78), Inches(4.15), Emu(310000),
       "役が持つ意味 ── 何が起きたら何をすればいいか", 10, bold=True, color=C_RED, font=FONT_H)

    roles = [
        (C_LTGRY, "ハズレ・ベル",  "バトル結果に変化なし\n次のゲームへ"),
        (C_CREAM, "リプレイ",      "チャンス！\n勝利への書き換え抽選が走る"),
        (C_GOLD,  "スイカ",        "強チャンス！\n勝利書き換え率がリプよりUP"),
        (C_CRIM,  "チェリー",      "高確率で勝利確定！\n7G目まで引けば大喜び"),
        (C_GOLD2, "8G目のレア役",  "ユリア復活確定\n継続率84%以上が確定する最大の自力演出"),
        (C_PUR,   "敗北後の復活",  "演出が続けば逆転あり！\n敗北確定まで諦めない"),
    ]
    ry = Inches(1.30)
    for ac, title, body in roles:
        rect(s, Inches(0.28), ry, Emu(50000), Emu(520000), ac)
        tb(s, Inches(0.48), ry + Emu(30000), Inches(1.55), Emu(245000),
           title, 9, bold=True, color=ac, wrap=False)
        tb(s, Inches(2.08), ry + Emu(30000), Inches(2.5), Emu(460000),
           body, 8.5, color=C_WHITE)
        ry += Emu(555000)

    # ── 右上：バトルの流れ ────────────────────────────────────
    rect_b(s, Inches(5.0), Inches(0.72), Inches(4.75), Inches(1.85),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(5.15), Inches(0.78), Inches(4.4), Emu(290000),
       "宿命バトルのフロー（8G）", 10, bold=True, color=C_GOLD, font=FONT_H)

    flow_steps = [
        ("突入", "チャンス役成立で宿命バトル開始"),
        ("1〜7G", "役を引くたびに書き換え抽選 or 勝利確定"),
        ("8G目", "レア役→ユリア復活確定（継続率84%以上確定）"),
        ("勝利", "BB（バトルボーナス）当選！"),
        ("敗北", "復活演出の可能性あり / 通常時に戻る"),
    ]
    fy = Inches(1.26)
    for j, (step, desc) in enumerate(flow_steps):
        bc = C_CRIM if step == "8G目" else (C_GREEN if step == "勝利" else C_GOLD)
        bg = RGBColor(0x14, 0x04, 0x04) if step == "8G目" else C_CARD2
        rect(s, Inches(5.08), fy, Inches(4.58), Emu(230000), bg)
        tb(s, Inches(5.18), fy + Emu(28000), Emu(550000), Emu(190000),
           step, 8.5, bold=True, color=bc, wrap=False)
        tb(s, Inches(6.0), fy + Emu(28000), Inches(3.4), Emu(190000),
           desc, 8.5, color=C_WHITE, wrap=False)
        fy += Emu(240000)

    # ── 右下：戦略ポイント ───────────────────────────────────
    rect_b(s, Inches(5.0), Inches(2.72), Inches(4.75), Inches(1.55),
           RGBColor(0x14, 0x04, 0x04), C_CRIM, 2.0)
    tb(s, Inches(5.15), Inches(2.78), Inches(4.4), Emu(285000),
       "プレイヤーとしての立ち回り", 10, bold=True, color=C_CRIM, font=FONT_H)
    tb(s, Inches(5.15), Inches(3.15), Inches(4.4), Emu(1000000),
       "・レア役が出たら確定演出を見逃さない\n"
       "・8G目まで諦めない（ユリア復活チャンス）\n"
       "・敗北演出が続いても画面を離れない\n"
       "・BB当選時のオーラ色を必ず確認\n"
       "  （白〜虹で継続率を「宣言」している）",
       8.5, color=C_WHITE)

    footer(s,
           "宿命バトル設計の本質：役を引く自分の行動が「ケンシロウを勝たせた」能動体験になる",
           "プレイヤーが確率を変えられないのに「自分が引いた」と感じるのがこのバトル設計の巧みさ。8G目のユリア復活が最大の自力演出。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT「世紀末モード」（何をすれば出玉が伸びるか）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「世紀末モード（BB）」  ──  継続と上乗せの設計を理解する", "5/10")

    # ── 上段：BBの2パート構成フロー ───────────────────────────
    BW = Inches(3.8)
    GAP = Inches(0.3)
    SX = Inches(0.2)

    rect_b(s, SX, Inches(0.72), BW, Inches(1.68),
           RGBColor(0x14, 0x04, 0x04), C_RED, 1.8)
    rect(s, SX, Inches(0.72), Emu(55000), Inches(1.68), C_RED)
    tb(s, SX + Emu(80000), Inches(0.78), BW - Emu(120000), Emu(285000),
       "小役パート（30G＋α）", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, SX + Emu(80000), Inches(1.18), BW - Emu(120000), Emu(1040000),
       "チャンス役でVストック・継続率アップを抽選\n"
       "宿命バトル（小役パート内）に勝利するとVストック獲得\n"
       "Vストックがあれば次セット継続が確定する！",
       8.5, color=C_WHITE)

    arrow_r(s, SX + BW + Emu(50000), Inches(0.72) + Inches(1.68)//2, C_GOLD)

    SX2 = SX + BW + GAP
    rect_b(s, SX2, Inches(0.72), BW, Inches(1.68),
           RGBColor(0x1A, 0x02, 0x02), C_CRIM, 1.8)
    rect(s, SX2, Inches(0.72), Emu(55000), Inches(1.68), C_CRIM)
    tb(s, SX2 + Emu(80000), Inches(0.78), BW - Emu(120000), Emu(285000),
       "バトルパート（8G）", 11, bold=True, color=C_CRIM, font=FONT_H)
    tb(s, SX2 + Emu(80000), Inches(1.18), BW - Emu(120000), Emu(1040000),
       "ケンシロウがラオウとバトル（＝継続判定）\n"
       "レア役成立で勝利確定・8G目ユリア復活が最大演出\n"
       "Vストックがあれば結果に関わらず次セットへ！",
       8.5, color=C_WHITE)

    # ── 下段左：継続率4段階 ─────────────────────────────────
    rect_b(s, Inches(0.2), Inches(2.54), Inches(3.8), Inches(1.75),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.35), Inches(2.60), Inches(3.5), Emu(285000),
       "BBオーラ色＝継続率を「宣言」する演出", 10, bold=True, color=C_GOLD, font=FONT_H)

    rates = [
        (C_GRAY,  "白オーラ",  "66%",  "最低継続率。敗北リスクが高い"),
        (C_WHITE, "青オーラ",  "79%",  "中継続。2〜3連が平均的"),
        (C_GOLD,  "金オーラ",  "84%",  "高継続。長連に期待大"),
        (C_CRIM,  "虹オーラ",  "89%",  "最高継続率。最強の宣言"),
    ]
    bary = Inches(2.95)
    BAR_MAX = Inches(2.5)
    for ac, color_name, rate, desc in rates:
        rect(s, Inches(0.28), bary, Inches(3.65), Emu(245000), C_ROW)
        tb(s, Inches(0.38), bary + Emu(28000), Emu(700000), Emu(200000),
           color_name, 8.5, bold=True, color=ac, wrap=False)
        tb(s, Inches(1.28), bary + Emu(28000), Emu(560000), Emu(200000),
           rate, 10, bold=True, color=ac, align=PP_ALIGN.CENTER, wrap=False)
        pct = float(rate.replace('%', '')) / 100
        rect(s, Inches(1.95), bary + Emu(90000), int(BAR_MAX * pct), Emu(90000), ac)
        rect(s, Inches(1.95), bary + Emu(90000), BAR_MAX, Emu(8000), C_LTGRY)
        bary += Emu(258000)

    # ── 下段中：Vストック ─────────────────────────────────────
    rect_b(s, Inches(4.2), Inches(2.54), Inches(2.6), Inches(1.75),
           C_CARD, C_TEAL, 1.5)
    rect(s, Inches(4.2), Inches(2.54), Emu(50000), Inches(1.75), C_TEAL)
    tb(s, Inches(4.43), Inches(2.60), Inches(2.25), Emu(285000),
       "Vストックの仕組み", 10, bold=True, color=C_TEAL, font=FONT_H)
    tb(s, Inches(4.43), Inches(2.98), Inches(2.25), Emu(1200000),
       "宿命バトル勝利で獲得\n\n"
       "Vストック1つ = 次セット継続確定\n（バトルパートの結果に関わらない）\n\n"
       "さらに「レイ・トキ協力」約40%で発生\nループ率12.5〜66%で連鎖する",
       8.5, color=C_WHITE)

    # ── 下段右：出玉を伸ばすための行動 ──────────────────────────
    rect_b(s, Inches(7.0), Inches(2.54), Inches(2.75), Inches(1.75),
           RGBColor(0x10, 0x0A, 0x02), C_GOLD2, 2.0)
    rect(s, Inches(7.0), Inches(2.54), Emu(50000), Inches(1.75), C_GOLD2)
    tb(s, Inches(7.23), Inches(2.60), Inches(2.3), Emu(285000),
       "出玉を伸ばすには", 10, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(7.23), Inches(2.98), Inches(2.35), Emu(1200000),
       "① Vストックを積む\n   （宿命バトルに勝ち続ける）\n\n"
       "② BBのバトルパートで\n   レア役を引いて自力で勝つ\n\n"
       "③ 無想転生チャンスへ\n   繋いで頂点へ到達！",
       8.5, color=C_WHITE)

    footer(s,
           "BBで出玉が伸びる条件：バトル継続×Vストック蓄積×無想転生チャンス突入の3段階",
           "上乗せはボーナス経由のみ。「バトルに勝ち続けること」がすべての行動指針になる設計。シンプルで明快。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位モード「修羅の国」（到達ルート＋遊び方）
# ══════════════════════════════════════════════════════════════
def s_jura(prs):
    s = new_slide(prs)
    hdr(s, "上位モード「修羅の国」  ──  高純増×高継続の頂点体験", "6/10")

    # ── 到達ルートフロー ──────────────────────────────────────
    rect(s, Inches(0.2), Inches(0.72), Inches(9.6), Emu(330000), C_CARD)
    tb(s, Inches(0.35), Inches(0.76), Inches(9.0), Emu(270000),
       "修羅の国への到達ルート  ──  無想転生バトルから繋がる", 10, bold=True, color=C_PUR, font=FONT_H)

    route_boxes = [
        (Inches(0.22), C_RED,  "BB継続",
         "バトルパート勝利\nまたはVストック消化で\n次セットBBへ"),
        (Inches(2.62), C_GOLD, "無想転生チャンス",
         "BB消化後に約33%で\n突入する15Gの特別区間\n→ 突入で神演出が発生"),
        (Inches(5.02), C_CRIM, "無想転生バトル",
         "継続率94%\n期待獲得2500枚以上\nここが「頂点」"),
        (Inches(7.42), C_PUR,  "修羅の国",
         "高純増・高継続の上位AT\n血涙の章がさらに上位\n最高峰の体験ゾーン"),
    ]
    for bx, ac, title, desc in route_boxes:
        rect_b(s, bx, Inches(1.15), Inches(2.2), Inches(1.20), C_CARD, ac, 1.5)
        tb(s, bx + Emu(60000), Inches(1.21), Inches(2.0), Emu(270000),
           title, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, bx + Emu(55000), Inches(1.55), Inches(2.05), Emu(720000),
           desc, 7.5, color=C_CREAM)
        if bx != Inches(7.42):
            arrow_r(s, bx + Inches(2.2) + Emu(30000), Inches(1.15) + Inches(1.20)//2, ac)

    # ── 左：無想転生バトル詳細 ───────────────────────────────
    rect_b(s, Inches(0.2), Inches(2.5), Inches(4.55), Inches(1.8),
           RGBColor(0x18, 0x02, 0x02), C_CRIM, 2.0)
    rect(s, Inches(0.2), Inches(2.5), Emu(55000), Inches(1.8), C_CRIM)
    tb(s, Inches(0.45), Inches(2.56), Inches(4.15), Emu(285000),
       "無想転生バトル  詳細", 10, bold=True, color=C_CRIM, font=FONT_H)
    tb(s, Inches(0.45), Inches(2.95), Inches(4.25), Emu(1200000),
       "・継続率94%（最高峰）\n"
       "・期待獲得枚数：2500枚以上\n"
       "・有利区間終了後も84%以上の高継続BBが再セット\n"
       "・「終わったと思ったらまだ続く」体験が生まれる\n"
       "・突入確認できたら長丁場を覚悟して楽しむ",
       8.5, color=C_WHITE)

    # ── 右：修羅の国＋血涙の章 ──────────────────────────────
    rect_b(s, Inches(5.0), Inches(2.5), Inches(4.75), Inches(1.8),
           RGBColor(0x10, 0x04, 0x20), C_PUR, 2.0)
    rect(s, Inches(5.0), Inches(2.5), Emu(55000), Inches(1.8), C_PUR)
    tb(s, Inches(5.25), Inches(2.56), Inches(4.3), Emu(285000),
       "修羅の国 ＆ 血涙の章  詳細", 10, bold=True, color=C_PUR, font=FONT_H)

    jura_items = [
        (C_PUR,  "修羅の国",  "無想転生バトル中から突入する上位AT\n高純増かつ継続率が大幅に強化"),
        (C_CRIM, "血涙の章",  "修羅の国のさらに上位のモード\n突入した時点で別格の出玉が約束される"),
        (C_GOLD, "遊び方",    "到達できたら「残枚数」を意識せず\nとにかくバトルを見守るだけでOK"),
    ]
    jy = Inches(2.95)
    for ac, title, body in jura_items:
        rect(s, Inches(5.08), jy, Emu(50000), Emu(500000), ac)
        tb(s, Inches(5.28), jy + Emu(30000), Inches(1.2), Emu(225000),
           title, 9, bold=True, color=ac, wrap=False)
        tb(s, Inches(6.52), jy + Emu(30000), Inches(3.1), Emu(430000),
           body, 8, color=C_WHITE)
        jy += Emu(540000)

    footer(s,
           "修羅の国設計の核心：「到達した人だけが体験できる上位世界」として来店継続の最大動機になる",
           "無想転生バトル→修羅の国→血涙の章という3段階のラダー設計が、初心者から上級者まで目標を分散させる。",
           C_PUR)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（なぜこの台は面白いのか）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計  ──  なぜスマスロ北斗の拳は面白いのか", "7/10")

    pillars = [
        (Inches(0.2), C_RED,
         "自力感の設計",
         "リプレイ成立→書き換えチャンス\nレア役成立→勝利確定\n8G目レア役→ユリア復活確定\n\n「役を引く自分の行動」が\nケンシロウを勝たせた\nと感じさせる巧みな設計\n\n確率は変えられないのに\n能動体験として記憶される",
         RGBColor(0x18, 0x02, 0x02)),
        (Inches(3.55), C_GOLD,
         "継続率の4段階管理",
         "66→79→84→89%の4段階が\nBB開始時にオーラで宣言される\n\n「今回はどの段階か」という\n緊張感が生まれ、虹オーラで\n「行けるかも」という確信が芽生える\n\nVストックが「敗北の恐怖」を\n和らげる緩衝材として機能する",
         RGBColor(0x14, 0x0C, 0x02)),
        (Inches(6.9), C_TEAL,
         "IP力による世代回帰",
         "4号機「北斗の拳」を\n原体験とする30〜40代が\nスマスロ版をきっかけに帰還\n\nIPの力で「初回来店」を獲得し\n実機の完成度で「リピート」を生む\n\n「北斗あるある」が\nコミュニティと会話を生み\n社会現象化を後押しした",
         RGBColor(0x02, 0x10, 0x18)),
    ]
    for x, col, title, desc, fill in pillars:
        rect_b(s, x, Inches(0.72), Inches(3.0), Inches(3.55), fill, col, 2.0)
        rect(s, x, Inches(0.72), Emu(55000), Inches(3.55), col)
        tb(s, x + Emu(80000), Inches(0.78), Inches(2.75), Emu(320000),
           title, 11, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.22), Inches(2.75), Inches(2.75),
           desc, 8.5, color=C_CREAM)

    footer(s,
           "総評：「自力感×継続率精密管理×IP」の三位一体が89週連続稼働ランキング3位を支えた",
           "どれか一つが欠けても長期稼働は実現しなかった。IP単体でも実機完成度単体でも不十分で、両方が必要だった。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_eval(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題  ──  バランス評価", "8/10")

    # ── 良い点（左）──────────────────────────────────────────
    rect_b(s, Inches(0.2), Inches(0.72), Inches(4.55), Inches(3.55),
           C_CARD, C_GREEN, 2.0)
    rect(s, Inches(0.2), Inches(0.72), Emu(55000), Inches(3.55), C_GREEN)
    tb(s, Inches(0.45), Inches(0.78), Inches(4.15), Emu(295000),
       "良い点（強み）", 11, bold=True, color=C_GREEN, font=FONT_H)

    pros = [
        (C_GREEN, "設計がシンプルで分かりやすい",
                  "BB継続とVストックという2軸だけで\n「何をすればいいか」が常に明確"),
        (C_GOLD,  "自力感と偶然性のバランス",
                  "役を引く行動が「自分で勝った」体験になる\n同時に偶然性も保たれ理不尽感が生きる"),
        (C_TEAL,  "89週連続の長期稼働",
                  "5634店設置・稼働ランキング3位維持\nIPと実機完成度の両立が証明された"),
        (C_CRIM,  "世代を超えたIP展開",
                  "30〜40代休眠層の呼び戻しに成功\n新規層との世代交流が生まれた"),
    ]
    py = Inches(1.28)
    for ac, title, body in pros:
        rect(s, Inches(0.28), py, Emu(50000), Emu(820000), ac)
        tb(s, Inches(0.48), py + Emu(30000), Inches(4.0), Emu(275000),
           title, 9, bold=True, color=ac)
        tb(s, Inches(0.48), py + Emu(305000), Inches(4.1), Emu(490000),
           body, 8.5, color=C_WHITE)
        py += Emu(870000)

    # ── 課題（右）──────────────────────────────────────────
    rect_b(s, Inches(5.0), Inches(0.72), Inches(4.75), Inches(3.55),
           RGBColor(0x16, 0x04, 0x04), C_RED, 2.0)
    rect(s, Inches(5.0), Inches(0.72), Emu(55000), Inches(3.55), C_RED)
    tb(s, Inches(5.25), Inches(0.78), Inches(4.3), Emu(295000),
       "課題（改善余地）", 11, bold=True, color=C_RED, font=FONT_H)

    cons = [
        (C_RED,   "差枚管理の不透明さ",
                  "有利区間内の差枚管理が非公開\n「冷遇区間」「デキレ」論争が続く\n透明性の欠如が信頼侵食リスク"),
        (C_CRIM,  "天然終了か冷遇か不明",
                  "有利区間終了原因が判断できない\n不透明さが陰謀論を育てる環境になっている"),
        (C_GOLD,  "設定判別が難しい",
                  "機械割差は大きいが判別要素が少ない\n設定1と設定6の差を意識する来店が難しい"),
        (C_LTGRY, "IP依存リスク",
                  "北斗の拳が無ければ成立しない要素が多い\n次世代タイトルへの応用には注意が必要"),
    ]
    cy2 = Inches(1.28)
    for ac, title, body in cons:
        rect(s, Inches(5.08), cy2, Emu(50000), Emu(820000), ac)
        tb(s, Inches(5.28), cy2 + Emu(30000), Inches(4.25), Emu(275000),
           title, 9, bold=True, color=ac)
        tb(s, Inches(5.28), cy2 + Emu(305000), Inches(4.35), Emu(490000),
           body, 8.5, color=C_WHITE)
        cy2 += Emu(870000)

    footer(s,
           "課題の本質：差枚管理の不透明さは設計上のリスクファクター。長期的な信頼構築には透明性が必要",
           "良い点の多くがIPと実機完成度の掛け算。次世代設計では「透明性」と「IP不依存の自力感設計」が課題。",
           C_RED)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ  ──  設計から学べること", "9/10")

    # ── 左：長期稼働を支えた3要素 ────────────────────────────
    rect_b(s, Inches(0.2), Inches(0.72), Inches(4.55), Inches(3.55),
           C_CARD, C_RED2, 1.5)
    rect(s, Inches(0.2), Inches(0.72), Emu(55000), Inches(3.55), C_RED2)
    tb(s, Inches(0.45), Inches(0.78), Inches(4.15), Emu(285000),
       "89週稼働を支えた3要素", 11, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_RED,   "① IP力（知名度×ノスタルジア）",
                  "4号機北斗を原体験とする30〜40代が\n復帰する動機になった\nIP単体では不十分。実機の完成度が前提"),
        (C_GOLD,  "② BB継続率4段階×Vストック",
                  "66〜89%の4段階＋Vストックで\n「次も続くかも」を繰り返す設計\n無想転生バトル94%が来店目標として機能"),
        (C_TEAL,  "③ 5634店・89週の安心感",
                  "設置台数と稼働期間の長さが\n「まだ打てる台」という消極的安心感を醸成\n長期稼働自体が来店動機になる循環"),
    ]
    ey = Inches(1.22)
    for ac, t, b in elems:
        rect_b(s, Inches(0.28), ey, Inches(4.37), Emu(1020000), C_CARD, ac, 1.2)
        rect(s, Inches(0.28), ey, Emu(48000), Emu(1020000), ac)
        tb(s, Inches(0.50), ey + Emu(50000), Inches(3.95), Emu(265000),
           t, 9, bold=True, color=ac)
        tb(s, Inches(0.50), ey + Emu(315000), Inches(3.95), Emu(680000),
           b, 8, color=C_WHITE)
        ey += Emu(1080000)

    # ── 右：設計原則と総括 ────────────────────────────────────
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.75)

    rect(s, rx, ry, rw, Emu(290000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(45000), rw - Emu(80000), Emu(220000),
       "設計原則（次機種に活かせること）", 11, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_RED,   "強IPは「休眠層の呼び水」になる\n  ── ただし実機の完成度が前提"),
        (C_GOLD,  "BB継続率の4段階が「次への期待感」を精密制御する"),
        (C_TEAL,  "Vストックが「敗北の恐怖」を和らげる緩衝材になる"),
        (C_CRIM,  "透明性の欠如（差枚非公開）は長期的な信頼リスク"),
        (C_PUR,   "自力感×継続率×IPの三位一体が奇跡的に揃った事例"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(290000) + i * Emu(480000)
        rect(s, rx, py0, Emu(22000), Emu(450000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(65000), rw - Emu(70000), Emu(390000),
           p, 8.5, bold=(i == 3 or i == 4), color=C_CRIM if i == 3 else ac)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(2740000), rw, Emu(835000),
           RGBColor(0x18, 0x04, 0x04), C_RED, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(2790000), rw - Emu(80000), Emu(265000),
       "総括", 9, bold=True, color=C_RED, font=FONT_H)
    tb(s, rx + Emu(60000), ry + Emu(3055000), rw - Emu(80000), Emu(470000),
       "IP×BB継続率×Vストックの三位一体が奇跡的に揃った事例。\n"
       "ただし差枚管理の不透明さは次世代設計で解決すべき課題。\n"
       "「長く打てる台」の設計としてベンチマーク価値が高い。",
       8, color=C_WHITE)

    net_note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        ("タイトル・スペック・3ポイント",           s_title),
        ("ゲームフロー全体図",                       s_flow),
        ("通常時の遊び方（2ルート）",                s_normal),
        ("宿命バトル（突破の仕方・役の意味）",       s_battle),
        ("AT「世紀末モード」",                       s_at),
        ("上位モード「修羅の国」",                   s_jura),
        ("面白さの設計（なぜ面白いか）",             s_design),
        ("良い点と課題",                             s_eval),
        ("まとめ・設計から学べること",               s_matome),
    ]

    print("=" * 58)
    print("  スマスロ北斗の拳 機種説明＋分析 統合版ジェネレーター")
    print("=" * 58)
    print()
    for i, (name, func) in enumerate(slides, 1):
        print(f"  {i:2d}/{len(slides)} {name}")
        func(prs)

    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
