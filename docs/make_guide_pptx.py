from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# カラーパレット
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_BLUE   = RGBColor(0x16, 0x47, 0x9A)
C_ACCENT = RGBColor(0x0F, 0xB9, 0xB9)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xFF)
C_GRAY   = RGBColor(0x55, 0x55, 0x77)
C_GREEN  = RGBColor(0x2E, 0xCC, 0x71)
C_ORANGE = RGBColor(0xF3, 0x9C, 0x12)
C_RED    = RGBColor(0xE7, 0x4C, 0x3C)
C_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
C_TEAL   = RGBColor(0x17, 0x7A, 0x6E)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, size=16, bold=False, color=C_DARK,
        align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return txb


def header(slide, title, sub=None):
    rect(slide, 0, 0, W, H, fill=C_LIGHT)
    rect(slide, 0, 0, W, Inches(1.3), fill=C_DARK)
    rect(slide, 0, Inches(1.3), W, Inches(0.06), fill=C_ACCENT)
    txt(slide, title, Inches(0.5), Inches(0.15), Inches(11), Inches(0.8),
        size=32, bold=True, color=C_WHITE)
    if sub:
        txt(slide, sub, Inches(0.5), Inches(0.85), Inches(11), Inches(0.4),
            size=15, color=C_ACCENT)


def card(slide, x, y, w, h, color, icon, title, desc, title_size=17, desc_size=13):
    rect(slide, x, y, w, Inches(0.62), fill=color)
    txt(slide, f"{icon}  {title}", x+Inches(0.1), y+Inches(0.06),
        w-Inches(0.2), Inches(0.5), size=title_size, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER)
    rect(slide, x, y+Inches(0.62), w, h-Inches(0.62), fill=C_WHITE,
         line=color, lw=Pt(1.5))
    txt(slide, desc, x+Inches(0.18), y+Inches(0.75),
        w-Inches(0.36), h-Inches(0.85), size=desc_size, color=C_DARK)


# ─────────────────────────────────────────────────
# Slide 1: タイトル
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, fill=C_DARK)
rect(s, 0, Inches(3.1), W, Inches(0.08), fill=C_ACCENT)
rect(s, 0, Inches(4.3), W, Inches(0.08), fill=C_BLUE)
txt(s, "SLOKEY", Inches(1), Inches(1.2), Inches(11), Inches(1.2),
    size=72, bold=True, color=C_ACCENT)
txt(s, "スロキー", Inches(1), Inches(2.2), Inches(11), Inches(0.8),
    size=28, color=RGBColor(0xAA,0xCC,0xFF))
txt(s, "社内向け 使い方ガイド", Inches(1), Inches(3.3), Inches(10), Inches(0.7),
    size=22, color=C_WHITE)
txt(s, "パチスロ情報を集めて・見て・分析する社内ツール",
    Inches(1), Inches(4.5), Inches(11), Inches(0.6),
    size=18, color=RGBColor(0x99,0xBB,0xDD))
txt(s, "2026年5月", Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.5),
    size=13, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────
# Slide 2: スロキーとは
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "スロキーとは？", "社内パチスロ情報共有ツール")
cards3 = [
    ("📥", "情報を集める", "ウェブ上に散らばる機種情報・ホール情報・打ち手の声を一か所にストック。AIが毎日自動収集。"),
    ("📊", "データで把握する", "「どの機種が今熱いか」「どんな情報が不足しているか」をランキング・ギャップ表で可視化。"),
    ("💬", "チームで共有する", "気になった情報をいいねで評価し、チームで使える情報の精度を上げていく。"),
]
for i, (ico, ttl, desc) in enumerate(cards3):
    x = Inches(0.4 + i*4.3)
    card(s, x, Inches(1.7), Inches(4.0), Inches(4.9), C_BLUE, ico, ttl, desc)

txt(s, "ブラウザで開くだけ・インストール不要・スマホ／PC両対応",
    Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.5),
    size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────
# Slide 3: 画面構成（5タブ）
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "画面構成", "画面下部の5タブで切り替え")

tabs5 = [
    ("📰 投稿",  C_BLUE,   "集まった情報を読む\nいいね・NG評価"),
    ("✏️ 追加",  C_GREEN,  "新しい情報を\n手動で投稿する"),
    ("📈 まとめ", C_ORANGE, "傾向・ランキング・\nギャップ表・新台カレンダーを俯瞰"),
    ("🔍 分析",  C_PURPLE, "機種別・コラム・\n機種分析・チャット・企画提案\n（パスワード保護）"),
    ("📊 稼働",  C_TEAL,   "ホール稼働データ\n（パスワード保護）"),
]
W_tab = Inches(2.45)
for i, (name, color, desc) in enumerate(tabs5):
    x = Inches(0.35 + i*2.55)
    y = Inches(1.7)
    rect(s, x, y, W_tab, Inches(0.72), fill=color)
    txt(s, name, x+Inches(0.05), y+Inches(0.08), W_tab-Inches(0.1), Inches(0.55),
        size=17, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, y+Inches(0.72), W_tab, Inches(4.4), fill=C_WHITE,
         line=color, lw=Pt(2))
    txt(s, desc, x+Inches(0.12), y+Inches(0.85), W_tab-Inches(0.24), Inches(4.0),
        size=14, color=C_DARK)

txt(s, "※「分析」「稼働」タブはパスワード保護（管理者のみ）",
    Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35),
    size=11, color=C_GRAY)


# ─────────────────────────────────────────────────
# Slide 4: 投稿タブ
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "投稿タブ：情報を読む", "カテゴリ別に絞り込みながら一覧閲覧")

txt(s, "カテゴリ一覧", Inches(0.5), Inches(1.55), Inches(5.5), Inches(0.5),
    size=15, bold=True, color=C_BLUE)
cats = [
    ("新台",   "導入前後2週間のスペック速報・新台発表"),
    ("機種情報", "天井・設定判別・機械割など攻略系"),
    ("実戦",   "打ち手の実戦報告・評判・体験談"),
    ("業界",   "ホール・メーカー・新台動向"),
    ("名機",   "4号機・5号機の思い出・伝説"),
]
for i, (cat, desc) in enumerate(cats):
    y = Inches(2.05 + i*0.82)
    rect(s, Inches(0.5), y, Inches(2.0), Inches(0.62), fill=C_BLUE)
    txt(s, cat, Inches(0.55), y+Inches(0.06), Inches(1.9), Inches(0.5),
        size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(2.7), y+Inches(0.1), Inches(3.5), Inches(0.52),
        size=13, color=C_DARK)

rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.6), fill=C_WHITE,
     line=C_ACCENT, lw=Pt(1.5))
txt(s, "基本操作", Inches(7.0), Inches(1.65), Inches(5.5), Inches(0.5),
    size=18, bold=True, color=C_BLUE)
ops = [
    "🔽  カテゴリボタンで絞り込み",
    "♥   いいねで「役に立つ」を評価",
    "👎  NGで「このネタは不要」を評価",
    "🔗  URLタップで元記事を表示",
    "📊  品質スコア（Lv.1〜5）が表示",
    "",
    "⭐ いいね・NGがAIの次回収集精度に反映",
]
for i, op in enumerate(ops):
    txt(s, op, Inches(7.1), Inches(2.25 + i*0.62), Inches(5.5), Inches(0.55),
        size=14, color=C_DARK if op else C_GRAY)


# ─────────────────────────────────────────────────
# Slide 5: 追加タブ
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "追加タブ：情報を投稿する", "気になった情報を気軽にストック")

steps = [
    ("①", "「新規投稿」ボタンをタップ", ""),
    ("②", "カテゴリ・機種名を選択", "カテゴリは5種類から選ぶ"),
    ("③", "タイトル・本文を入力", "URLを貼ると本文を自動取得。短い一言でもOK"),
    ("④", "投稿ボタンで送信", "全メンバーから閲覧可能になる"),
]
for i, (num, ttl, hint) in enumerate(steps):
    y = Inches(1.7 + i*1.3)
    rect(s, Inches(0.5), y, Inches(0.75), Inches(0.9), fill=C_BLUE)
    txt(s, num, Inches(0.5), y, Inches(0.75), Inches(0.9),
        size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, ttl, Inches(1.45), y+Inches(0.05), Inches(7.5), Inches(0.5),
        size=18, bold=True, color=C_DARK)
    if hint:
        txt(s, hint, Inches(1.45), y+Inches(0.55), Inches(7.5), Inches(0.45),
            size=13, color=C_GRAY)

rect(s, Inches(9.5), Inches(1.7), Inches(3.5), Inches(5.1), fill=C_DARK)
txt(s, "投稿のコツ", Inches(9.65), Inches(1.85), Inches(3.2), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
tips = ["URLを貼ると\n本文が自動入力", "短い一言でもOK", "自分の投稿は\nあとで削除可能", "AIも毎日自動収集\n（編集部AI名義）"]
for i, tip in enumerate(tips):
    txt(s, f"• {tip}", Inches(9.7), Inches(2.4 + i*1.0), Inches(3.1), Inches(0.9),
        size=13, color=C_WHITE)


# ─────────────────────────────────────────────────
# Slide 6: まとめタブ（7ビュー）
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "まとめタブ：7つのビュー", "傾向・ランキング・ギャップ表を一画面で俯瞰")

views7 = [
    (C_BLUE,   "🏆", "ランキング",    "機種ごとの投稿数・いいね数・品質スコア・直近2週間数を一覧比較"),
    (C_ORANGE, "📌", "機種別",        "機種を選んで投稿を深掘り。いいね順に表示"),
    (C_GREEN,  "🥧", "カテゴリ分布",  "カテゴリ別の投稿数・割合・人気投稿を確認"),
    (C_PURPLE, "👤", "投稿者",        "誰がどんな情報を多く投稿しているか一覧"),
    (C_TEAL,   "🔽", "絞り込み",      "機種名・カテゴリで絞り込んで投稿を読む"),
    (C_RED,    "🟥", "ギャップ表",    "機種×カテゴリのマトリクス。赤＝情報なし・黄＝1件・緑＝2件以上"),
    (RGBColor(0x63,0x7A,0x2F), "📅", "新台カレンダー", "直近3ヶ月の導入台を月別スケジュール表示。純増・機械割などスペック付き"),
]
cols = 4
for i, (color, ico, ttl, desc) in enumerate(views7):
    row, col = divmod(i, cols)
    x = Inches(0.3 + col*3.25)
    y = Inches(1.7 + row*2.7)
    w = Inches(3.0)
    h = Inches(2.5)
    card(s, x, y, w, h, color, ico, ttl, desc, title_size=15, desc_size=12)


# ─────────────────────────────────────────────────
# Slide 7: ランキング＆ギャップ表の活用
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "ランキング＆ギャップ表の活用", "「今どの機種に注目すべきか」が一目でわかる")

rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.6), fill=C_WHITE,
     line=C_BLUE, lw=Pt(2))
rect(s, Inches(0.4), Inches(1.6), Inches(6.0), Inches(0.6), fill=C_BLUE)
txt(s, "🏆 ランキング", Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.5),
    size=18, bold=True, color=C_WHITE)
rank_items = [
    ("投稿数ソート",    "情報量が多い＝注目度が高い機種がわかる"),
    ("いいねソート",    "チームが「使える」と評価した情報が多い機種"),
    ("品質ソート",      "信頼度の高い情報が集まっている機種"),
    ("直近2週間ソート", "今まさに話題になっている機種"),
]
for i, (sort, desc) in enumerate(rank_items):
    y = Inches(2.35 + i*0.9)
    rect(s, Inches(0.6), y, Inches(2.1), Inches(0.65), fill=C_LIGHT,
         line=C_BLUE, lw=Pt(1))
    txt(s, sort, Inches(0.65), y+Inches(0.07), Inches(2.0), Inches(0.55),
        size=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(2.85), y+Inches(0.12), Inches(3.3), Inches(0.55),
        size=13, color=C_DARK)

rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(5.6), fill=C_WHITE,
     line=C_RED, lw=Pt(2))
rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(0.6), fill=C_RED)
txt(s, "🟥 ギャップ表", Inches(6.9), Inches(1.65), Inches(5.8), Inches(0.5),
    size=18, bold=True, color=C_WHITE)
txt(s, "機種 × カテゴリのマトリクス表\n「どの情報が足りていないか」が一目でわかる",
    Inches(7.0), Inches(2.35), Inches(5.6), Inches(1.0), size=14, color=C_DARK)
legend = [
    (C_RED,    "赤", "情報なし → 収集の最優先"),
    (C_ORANGE, "黄", "1件のみ → もう少し集めたい"),
    (C_GREEN,  "緑", "2件以上 → ある程度揃っている"),
]
for i, (color, label, desc) in enumerate(legend):
    y = Inches(3.5 + i*0.9)
    rect(s, Inches(7.0), y, Inches(0.62), Inches(0.62), fill=color)
    txt(s, label, Inches(7.0), y, Inches(0.62), Inches(0.62),
        size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(7.75), y+Inches(0.1), Inches(4.8), Inches(0.55),
        size=13, color=C_DARK)
txt(s, "活用例：「この機種の実戦情報が全然ない」→ 収集リクエストを出す",
    Inches(7.0), Inches(6.35), Inches(5.8), Inches(0.55), size=12, color=C_GRAY)


# ─────────────────────────────────────────────────
# Slide 8: 分析タブ（6モード）
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "分析タブ：6つのモード", "機種別・コラム・ゲーム性・チャット・企画提案")

modes6 = [
    (C_BLUE,   "📝", "コラム",      "編集部が執筆したパチスロ情報コラムを読む"),
    (C_GREEN,  "⭐", "機種評価",    "機種ごとの「良い点・気になる点」をユーザー評価で確認"),
    (C_ORANGE, "📋", "機種分析",    "AIが蓄積投稿を分析。良い点・気になる点を箇条書きで自動生成"),
    (RGBColor(0x1A,0x78,0x8A), "🎮", "ゲーム性分析", "CZ設計・AT設計・演出など機種のゲーム性を詳細解説"),
    (C_PURPLE, "💬", "チャット",    "質問を入力→AIがDB情報を元にバックグラウンドで返答。会話はリアルタイム更新。参加型"),
    (RGBColor(0xA0,0x3A,0x90), "✏️", "企画提案",  "IP名・ターゲット入力→AIが3問ヒアリング→提案書生成→評価・修正対応まで。参加型"),
]
cols = 3
for i, (color, ico, ttl, desc) in enumerate(modes6):
    row, col = divmod(i, cols)
    x = Inches(0.35 + col*4.3)
    y = Inches(1.7 + row*2.75)
    w = Inches(4.0)
    h = Inches(2.55)
    card(s, x, y, w, h, color, ico, ttl, desc, title_size=16, desc_size=13)

txt(s, "💬 チャット・✏️ 企画提案はユーザー参加型機能です",
    Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.35),
    size=11, color=C_PURPLE)


# ─────────────────────────────────────────────────
# Slide 9: チャット・企画提案の使い方
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "チャット・企画提案の使い方", "分析タブ内のユーザー参加型機能")

# ── 左: チャット ──────────────────────────
rect(s, Inches(0.4), Inches(1.6), Inches(6.1), Inches(5.6), fill=C_WHITE,
     line=C_PURPLE, lw=Pt(2))
rect(s, Inches(0.4), Inches(1.6), Inches(6.1), Inches(0.58), fill=C_PURPLE)
txt(s, "💬 チャット", Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.48),
    size=18, bold=True, color=C_WHITE)

chat_steps = [
    ("①", "機種名やキーワードを入力して送信"),
    ("②", "AIがスロキーのDBを参照してバックグラウンドで返答"),
    ("③", "返答が届くと画面がリアルタイムで自動更新"),
    ("④", "続けて質問して情報を深掘りできる"),
]
for i, (num, desc) in enumerate(chat_steps):
    y = Inches(2.35 + i*0.88)
    rect(s, Inches(0.6), y, Inches(0.62), Inches(0.62), fill=C_PURPLE)
    txt(s, num, Inches(0.6), y, Inches(0.62), Inches(0.62),
        size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(1.35), y+Inches(0.1), Inches(4.9), Inches(0.55),
        size=13, color=C_DARK)

rect(s, Inches(0.6), Inches(6.0), Inches(5.6), Inches(0.9), fill=C_LIGHT)
txt(s, "会話履歴はSupabaseに保存。セッションをまたいで続きから質問可能",
    Inches(0.7), Inches(6.1), Inches(5.4), Inches(0.7), size=12, color=C_GRAY)

# ── 右: 企画提案 ──────────────────────────
rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(5.6), fill=C_WHITE,
     line=RGBColor(0xA0,0x3A,0x90), lw=Pt(2))
rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(0.58), fill=RGBColor(0xA0,0x3A,0x90))
txt(s, "✏️ 企画提案", Inches(6.9), Inches(1.65), Inches(5.8), Inches(0.48),
    size=18, bold=True, color=C_WHITE)

prop_steps = [
    ("①", "IP名（版権タイトル）とターゲット像を入力"),
    ("②", "AIが3問ヒアリング（コンセプト・演出方向性・訴求ポイント）"),
    ("③", "提案書を自動生成（スペック・AT設計・演出骨子を含む）"),
    ("④", "「良い」「修正希望」「やり直し」で評価→修正対応も可"),
]
for i, (num, desc) in enumerate(prop_steps):
    y = Inches(2.35 + i*0.88)
    rect(s, Inches(7.0), y, Inches(0.62), Inches(0.62), fill=RGBColor(0xA0,0x3A,0x90))
    txt(s, num, Inches(7.0), y, Inches(0.62), Inches(0.62),
        size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(7.75), y+Inches(0.1), Inches(4.9), Inches(0.55),
        size=13, color=C_DARK)

rect(s, Inches(7.0), Inches(6.0), Inches(5.6), Inches(0.9), fill=C_LIGHT)
txt(s, "提案書はリクエスト一覧に保存。過去の提案をいつでも参照可能",
    Inches(7.1), Inches(6.1), Inches(5.4), Inches(0.7), size=12, color=C_GRAY)


# ─────────────────────────────────────────────────
# Slide 10: 機種分析
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "機種分析", "良い点・気になる点をまとめてチェック")

txt(s, "機種をドロップダウンから選んで「分析する」を押すだけ",
    Inches(0.5), Inches(1.55), Inches(12), Inches(0.5), size=16, color=C_DARK)

rect(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.8), fill=C_WHITE,
     line=C_GREEN, lw=Pt(2))
rect(s, Inches(0.5), Inches(2.1), Inches(5.8), Inches(0.55), fill=C_GREEN)
txt(s, "✅ 良い点", Inches(0.6), Inches(2.15), Inches(5.5), Inches(0.45),
    size=16, bold=True, color=C_WHITE)
for i, p in enumerate(["高設定の爆発力がわかりやすい", "設定差が演出に出るため判別しやすい", "AT中の上乗せ頻度が高く体感満足度大"]):
    txt(s, f"• {p}", Inches(0.7), Inches(2.8+i*0.7), Inches(5.4), Inches(0.6), size=14, color=C_DARK)

rect(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.8), fill=C_WHITE,
     line=C_RED, lw=Pt(2))
rect(s, Inches(7.0), Inches(2.1), Inches(5.8), Inches(0.55), fill=C_RED)
txt(s, "⚠️ 気になる点", Inches(7.1), Inches(2.15), Inches(5.5), Inches(0.45),
    size=16, bold=True, color=C_WHITE)
for i, c in enumerate(["低設定は天井まで到達しやすくストレス", "CZのデキレ感をユーザーが指摘", "導入直後は情報が少なく設定判別困難"]):
    txt(s, f"• {c}", Inches(7.1), Inches(2.8+i*0.7), Inches(5.4), Inches(0.6), size=14, color=C_DARK)

txt(s, "📊 未分析（N件）→ まだ分析データなし   ⚠️ +N件追加あり → 情報が増えて再分析が必要な機種",
    Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.35), size=11, color=C_GRAY)


# ─────────────────────────────────────────────────
# Slide 11: AIいいね・NGで精度アップ
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "AIいいね・NGで精度を上げる", "評価がフィードバックされ、次回以降の収集に反映")

txt(s, "AIが毎日自動収集した投稿（編集部AI・スロキー編集部 名義）を積極的に評価してください",
    Inches(0.5), Inches(1.6), Inches(12), Inches(0.55), size=15, color=C_DARK)

for i, (ico, label, color, title, desc) in enumerate([
    ("♥", "いいね", C_GREEN,
     "「この種類の情報は役に立つ」シグナル",
     "同じソース・機種・内容パターンが\n今後も優先的に収集される"),
    ("👎", "NG", C_RED,
     "「このネタは不要・的外れ」シグナル",
     "同種のネタが次回収集から除外される\n（badカウントが増加）"),
]):
    x = Inches(1.0 + i*6.2)
    y = Inches(2.3)
    rect(s, x, y, Inches(5.5), Inches(4.2), fill=C_WHITE, line=color, lw=Pt(2))
    rect(s, x, y, Inches(5.5), Inches(1.0), fill=color)
    txt(s, f"{ico}  {label}", x+Inches(0.2), y+Inches(0.15), Inches(5.0), Inches(0.7),
        size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+Inches(0.2), y+Inches(1.15), Inches(5.0), Inches(0.6),
        size=15, bold=True, color=color)
    txt(s, desc, x+Inches(0.2), y+Inches(1.85), Inches(5.0), Inches(2.0),
        size=14, color=C_DARK)

txt(s, "チームみんなで評価するほど、スロキーの情報がチームの好みに最適化されていきます",
    Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.5),
    size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────
# Slide 12: フィードバック＆FAQ
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "フィードバック＆よくある質問", "")

rect(s, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.5), fill=C_WHITE,
     line=C_ACCENT, lw=Pt(2))
rect(s, Inches(0.4), Inches(1.6), Inches(5.8), Inches(0.55), fill=C_ACCENT)
txt(s, "💬 フィードバックの送り方", Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.45),
    size=16, bold=True, color=C_WHITE)
txt(s,
    "画面右下の 💬 ボタンをタップ\n→ 種類を選ぶ（機能要望・バグ報告・その他）\n→ 内容を入力して送信\n\n管理者のみが確認できます",
    Inches(0.6), Inches(2.25), Inches(5.4), Inches(2.8), size=14, color=C_DARK)

rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(5.5), fill=C_WHITE,
     line=C_BLUE, lw=Pt(2))
rect(s, Inches(6.8), Inches(1.6), Inches(6.1), Inches(0.55), fill=C_BLUE)
txt(s, "❓ よくある質問", Inches(6.9), Inches(1.65), Inches(5.8), Inches(0.45),
    size=16, bold=True, color=C_WHITE)

faqs = [
    ("投稿は全員に見えますか？",      "はい。ツールを開いているメンバー全員が閲覧可能です。"),
    ("投稿を間違えた場合は？",          "自分の投稿は削除できます（一覧右上のメニューから）。"),
    ("スマホとPCで内容は同じですか？",  "同じデータです。どちらからでも最新情報が表示されます。"),
    ("品質スコアとは？",               "情報の信頼度（Lv.1〜5）。Lv.5が最も信頼度高。URLありの記事はLv.4〜5。"),
]
for i, (q, a) in enumerate(faqs):
    y = Inches(2.25 + i*1.2)
    txt(s, f"Q. {q}", Inches(7.0), y, Inches(5.7), Inches(0.45),
        size=13, bold=True, color=C_BLUE)
    txt(s, f"A. {a}", Inches(7.0), y+Inches(0.45), Inches(5.7), Inches(0.6),
        size=13, color=C_DARK)

rect(s, Inches(0.4), Inches(5.3), Inches(5.8), Inches(1.8), fill=C_DARK)
txt(s, "使いながら育てるツールです\n気になること・要望はどんどん送ってください",
    Inches(0.6), Inches(5.5), Inches(5.4), Inches(1.4), size=14, color=C_WHITE)


# ─────────────────────────────────────────────────
out = r"C:\Users\h.kadoya\Desktop\slocri\docs\slokey_guide.pptx"
prs.save(out)
print(f"Saved: {out}")
