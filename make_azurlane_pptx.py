"""
L アズールレーン THE ANIMATION（スマスロ）機種説明＋分析 統合版 PPTXジェネレーター v1
出力: proposals/機種分析/アズールレーン/azurlane_guide_v1.pptx
テーマ: 深海紺×シアン×金×白（海軍・艦隊カラー）
情報源: ちょんぼりすた・一撃・altema・アミュタメ 各解析ページ（2025年〜2026年）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "アズールレーン", "azurlane_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深海紺×シアン×金×白）────────────────────────────
C_BG    = RGBColor(0x04, 0x08, 0x18)   # 深海紺
C_CARD  = RGBColor(0x06, 0x0E, 0x22)   # カード背景
C_CARD2 = RGBColor(0x08, 0x12, 0x2C)   # カード背景2
C_ROW   = RGBColor(0x0A, 0x10, 0x20)   # 行交互
C_CYAN  = RGBColor(0x00, 0xCC, 0xDD)   # シアン #00CCDD
C_CYAN2 = RGBColor(0x00, 0xAA, 0xCC)   # シアン2
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金 #C8A840
C_GOLD2 = RGBColor(0xFF, 0xCC, 0x55)   # 明るい金
C_WHITE = RGBColor(0xEC, 0xF0, 0xF8)   # 白
C_CREAM = RGBColor(0xD8, 0xE4, 0xF0)   # クリーム
C_GRAY  = RGBColor(0x88, 0x9A, 0xAA)   # グレー
C_LTGRY = RGBColor(0x44, 0x55, 0x66)   # 濃グレー
C_NAVY  = RGBColor(0x10, 0x1E, 0x3A)   # ネイビー
C_BLUE  = RGBColor(0x22, 0x55, 0xCC)   # ブルー
C_GREEN = RGBColor(0x22, 0xCC, 0x66)   # 緑
C_RED   = RGBColor(0xDD, 0x22, 0x22)   # 赤
C_PINK  = RGBColor(0xFF, 0x66, 0xAA)   # ピンク
C_ORG   = RGBColor(0xFF, 0x88, 0x22)   # オレンジ

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景・ヘルパー群 ─────────────────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (4, 8, 24))
    draw = ImageDraw.Draw(img)
    # 波紋パターン（海面）
    for i in range(0, w + h, 90):
        draw.line([(i, 0), (0, i)], fill=(6, 12, 30), width=1)
    # 底部シアングロー
    for y in range(h - 140, h):
        t = (y - (h - 140)) / 140
        r = int(0 * t)
        g = int(40 * t)
        b = int(60 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    # 上部グラデーション
    for y in range(0, 50):
        t = (50 - y) / 50 * 0.4
        draw.line([(0, y), (w, y)], fill=(0, int(10 * t), int(30 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_NAVY)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_CYAN)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_CYAN, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_CYAN)


def net_note(slide):
    tb(slide, Inches(7.8), Inches(5.38), Inches(2.1), Emu(180000),
       "※ネット解析情報より（ちょんぼりすた・一撃・altema・アミュタメ）",
       6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x04, 0x0A, 0x1C))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_CYAN)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_CYAN)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_CYAN
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_CYAN
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x04, 0x08, 0x1C))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_CYAN)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_CYAN)

    tb(s, Inches(0.22), Inches(0.4), Inches(5.0), Emu(330000),
       "機種説明＋分析 統合ガイド  v1（解析情報版）", 10, color=C_CYAN2, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.88), Inches(5.1), Emu(900000),
       "アズールレーン", 34, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.55), Inches(5.0), Emu(280000),
       "L アズールレーン THE ANIMATION ── スマスロ（京楽）陣営×艦隊×上位AT設計", 9, color=C_CREAM, font=FONT_H)

    # スペック表
    specs = [
        ("メーカー",        "京楽産業（KYORAKU）"),
        ("機種タイプ",      "スマスロ（Lシリーズ）AT機"),
        ("AT純増",         "約2.5枚/G（通常AT） / 約5.1枚/G（上位AT）"),
        ("設定1機械割",    "非公表（推定97%台）"),
        ("設定6機械割",    "非公表（推定110%超）"),
        ("AT初当り",       "疑似ボーナス経由（海戦BONUS消化後）"),
        ("天井①",         "ボーナス間350G → ボーナス確定"),
        ("天井②（スルー）","9スルー目（10回目のボーナスでAT確定）"),
        ("リセット短縮",   "リセット時スルー天井6回（7回目でAT）"),
        ("コイン持ち",     "50枚あたり約25.8G（低持ち設計）"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(3.0) + i * Emu(225000)
        tb(s, Inches(0.22), ry, Inches(1.7), Emu(210000),
           k, 7.5, color=C_GRAY)
        tb(s, Inches(1.92), ry, Inches(3.3), Emu(210000),
           v, 7.5, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_CYAN,  "① 陣営システムでATが強化される",
         "AT中に陣営が増えるほど性能がアップ。\n4陣営集結で天城BATTLEへ突入。\n「仲間を集める」体験と出玉が直結。"),
        (C_GOLD,  "② 純増5.1枚/Gの上位AT「SS RUSH」",
         "通常ATの倍以上の純増速度で出玉爆増。\n4陣営集結→天城BATTLE→上位AT昇格の\nドラマチックな昇格ルートが魅力。"),
        (C_PINK,  "③ 明石商店とキャラIPの融合",
         "通常時・AT中に「明石商店」が登場し\nリプレイで文字点灯→明石チャンス突入。\nアズレンキャラとゲーム性の完全融合。"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.55) + i * Emu(1540000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.3), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.3), ac)
        tb(s, Inches(5.85), y0 + Emu(65000), Inches(3.8), Emu(310000),
           kw, 11.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(380000), Inches(3.8), Emu(450000),
           desc, 8.5, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「陣営システム×上位AT昇格×艦隊IPキャラ」── 仲間集めがそのまま出玉性能に繋がる独自設計",
           "パチスロアワード2025ノミネート機。京楽奇跡の1台との評価あり")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（全ルートを蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→AT→上位ATの全ルート", "2/9")

    # 上段：通常→疑似ボーナス→RUSH の基本フロー
    top_y = Inches(0.75)
    top_h = Emu(1100000)

    flow1 = [
        (C_CARD2, C_GRAY,   "通常遊技",              "規定G数/レア役\n/明石チャンスで\nボーナス抽選"),
        (C_CARD,  C_CYAN2,  "海戦BONUS\n(疑似ボーナス)", "ベルナビ10回\n+10G継続\n平均約100枚"),
        (C_CARD,  C_BLUE,   "BATTLE演出\n（AT抽選）",  "勝利でAT確定\n告知発生でも\nRUSH突入"),
        (C_CARD,  C_CYAN,   "アズールレーン\nRUSH(AT)", "純増約2.5枚/G\nゲーム数上乗せ\nメインAT"),
        (C_CARD,  C_GOLD,   "陣営集結\n→天城BATTLE", "4陣営で発動\nAT後に上位AT\nをかけた決戦"),
    ]
    bw1 = Inches(1.6)
    gap1 = Inches(0.17)
    sx1 = Inches(0.28)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), top_y + Emu(70000), bw1 - Emu(60000), Emu(360000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), top_y + Emu(480000), bw1 - Emu(45000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), cy1)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1120000), Inches(3.2), Emu(260000),
       "天井①：ボーナス間350G → ボーナス確定", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1360000), Inches(3.2), Emu(5000), C_CYAN)
    tb(s, sx1 + Inches(3.3), top_y + Emu(1120000), Inches(3.5), Emu(260000),
       "天井②：9スルーで10回目にAT確定（リセット後は6スルー）", 7.5, color=C_GOLD)

    # 中段区切り線
    rect(s, 0, Inches(2.1), SLIDE_W, Emu(5000), RGBColor(0x10, 0x2A, 0x44))

    # 下段：RUSH内→上位AT昇格ルート
    bot_y = Inches(2.18)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,  C_CYAN,   "アズールレーン\nRUSH(通常AT)",   "純増約2.5枚/G\n上乗せ型\n1陣営スタート"),
        (C_CARD,  C_CYAN2,  "陣営が増える\n（AT性能UP）",   "BAR揃い等で\n陣営参戦\n最大4陣営"),
        (C_CARD,  C_GOLD,   "4陣営集結\n→天城BATTLE",     "AT終了後発動\n上位CZ\nクリアで昇格"),
        (RGBColor(0x06,0x12,0x28), C_GOLD2,
         "異次元性能SS\nRUSH(上位AT)",              "純増約5.1枚/G\n4陣営状態\n期待約3500枚+"),
    ]
    bw2 = Inches(1.95)
    gap2 = Inches(0.22)
    sx2 = Inches(0.28)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8)
        tb(s, bx + Emu(40000), bot_y + Emu(70000), bw2 - Emu(70000), Emu(370000),
           lbl, 10, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), bot_y + Emu(490000), bw2 - Emu(55000), Emu(470000),
           sub, 8, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw2 + Emu(12000), cy2)

    # 右端：共同戦線RUSH説明
    rx_al = sx2 + 4 * (bw2 + gap2)
    rect_b(s, rx_al, bot_y, Inches(1.55), bot_h, C_CARD, C_PINK, 1.5)
    rect(s, rx_al, bot_y, Emu(30000), bot_h, C_PINK)
    tb(s, rx_al + Emu(50000), bot_y + Emu(70000), Inches(1.3), Emu(310000),
       "共同戦線\nRUSH\n(特化)", 9, bold=True, color=C_PINK, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, rx_al + Emu(50000), bot_y + Emu(490000), Inches(1.3), Emu(460000),
       "1陣営時\n高上乗せ\n特化ゾーン", 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 明石チャンス補足
    tb(s, sx2, bot_y + Emu(1100000), Inches(4.5), Emu(260000),
       "明石チャンス：通常時・AT中共通でリプレイ積算→ボーナス・上乗せ抽選", 7.5, color=C_CYAN2)
    rect(s, sx2, bot_y + Emu(1340000), Inches(4.5), Emu(5000), C_CYAN2)

    net_note(s)
    footer(s, "上段=通常時〜AT突入フロー（疑似ボーナス→BATTLE勝利）、下段=AT内陣営集結→上位AT昇格ルート",
           "天城BATTLEはAT終了後に4陣営集結で発生。上位AT（SS RUSH）は純増5.1枚/Gと大幅性能UP")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（天井含む全ルート）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── AT突入全ルート・天井管理・打ち方", "3/9")

    # 左：AT突入ルート図
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x10, 0x22, 0x44))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜AT突入ルート（全3系統）", 10, bold=True, color=C_CYAN)

    routes = [
        (C_CYAN,  "ルート①  海戦BONUS→BATTLE勝利でAT突入（最頻）",
         "ゲーム消化→高確・超高確経由→海戦BONUS当選。\nボーナス内BATTLE勝利でアズールレーンRUSH突入。\n告知発生でもRUSH確定。標準的な最多発生ルート。"),
        (C_GOLD,  "ルート②  強チェリー・チャンス目→ダイレクト突入",
         "強チェリー・チャンス目成立→ボーナスまたはAT直撃。\n高確・超高確状態なら期待度大幅アップ。\n通常時にレア役が出たら要注目。"),
        (C_CYAN2, "ルート③  天井①（ボーナス間350G）",
         "ボーナス間350GでボーナスorAT確定。\n浅めの天井設定でライトユーザーにも優しい。\n設定変更後はリセット恩恵あり（状態が変わる）。"),
        (C_GOLD2, "ルート④  スルー天井（9スルー→10回目AT確定）",
         "ボーナス9スルー後、10回目のボーナスでAT確定。\nリセット後は6スルーに短縮（7回目AT確定）。\n長期ハマリ台の最終救済システム。"),
    ]
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(290000) + i * Emu(1125000)
        rect_b(s, lx, iy, lw, Emu(1060000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1060000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(670000),
           b, 7.5, color=C_WHITE)

    # 右：明石商店とチャンス役
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x10, 0x22, 0x44))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "明石商店と通常時チャンス役", 10, bold=True, color=C_CYAN)

    # 明石商店説明ボックス
    rect_b(s, rx, ry + Emu(290000), rw, Emu(900000), C_CARD, C_GOLD, 1.8)
    rect(s, rx, ry + Emu(290000), Emu(40000), Emu(900000), C_GOLD)
    tb(s, rx + Emu(70000), ry + Emu(340000), rw - Emu(90000), Emu(260000),
       "明石商店システム（通常時・AT中共通）", 9, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, rx + Emu(70000), ry + Emu(610000), rw - Emu(90000), Emu(520000),
       "リール左側に「明・石・商・店」が表示中、\nリプレイを引くたびに1文字ずつ点灯を抽選。\n全文字点灯→明石チャンスへ突入。\n通常時はボーナス抽選、AT中は上乗せ抽選。",
       8, color=C_WHITE)

    # チャンス役一覧
    chance = [
        (C_GRAY,  "リプレイ",    "約1/8.6",   "明石商店の文字点灯抽選役"),
        (C_CYAN,  "弱チェリー",  "約1/80",    "ボーナス・AT当選の抽選対象"),
        (C_CYAN2, "スイカ",      "約1/128",   "ボーナス当選・高確移行抽選"),
        (C_GOLD,  "強チェリー",  "約1/200",   "ボーナス直撃・AT直撃の高期待役"),
        (C_GOLD2, "チャンス目",  "約1/128",   "AT当選抽選・高確移行の高期待役"),
        (C_PINK,  "レア役合算",  "約1/35",    "チェリー・スイカ・チャンス目合算"),
    ]
    ch_h = Emu(610000)
    for i, (ac, role, prob, desc) in enumerate(chance):
        cy = ry + Emu(1190000) + i * ch_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, ch_h, bg)
        rect(s, rx, cy, Emu(35000), ch_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(50000), Inches(1.0), Emu(230000),
           role, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.0), cy + Emu(55000), Inches(0.8), Emu(220000),
           prob, 8, bold=True, color=C_GOLD2, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(280000), rw - Emu(60000), Emu(280000),
           desc, 7.5, color=C_WHITE)

    # 打ち方メモ
    rect_b(s, rx, ry + Emu(4870000), rw, Emu(490000), C_CARD2, C_CYAN, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4920000), rw - Emu(80000), Emu(200000),
       "打ち方：左リールにBARを狙う（チェリー・スイカをフォロー）", 8.5, bold=True, color=C_CYAN)
    tb(s, rx + Emu(60000), ry + Emu(5130000), rw - Emu(80000), Emu(220000),
       "ベスト手順を守ることでレア役フォローと高確移行を逃さない。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "通常時戦略：天井①(350G)は浅め、スルー天井(9スルー)は長期台のサイン。明石商店の文字点灯も要チェック",
           "高確・超高確状態でのレア役は大チャンス。状態の変化を意識したゲーム消化が重要")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み（海戦・艦隊要素との連携）
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ/前兆の仕組み ── 海戦BONUS・BATTLE・艦隊要素との連携", "4/9")

    fc_x = Inches(0.28)
    fc_w = Inches(4.4)

    rect(s, fc_x, Inches(0.72), fc_w, Emu(260000), RGBColor(0x10, 0x22, 0x44))
    tb(s, fc_x + Emu(60000), Inches(0.74), fc_w - Emu(80000), Emu(240000),
       "海戦BONUSからAT突入までのフロー", 9.5, bold=True, color=C_CYAN)

    # STEP 1
    n1_y = Inches(1.06)
    n1_h = Emu(560000)
    rect_b(s, fc_x, n1_y, fc_w, n1_h, C_CARD2, C_GRAY, 1.2)
    tb(s, fc_x + Emu(40000), n1_y + Emu(45000), fc_w - Emu(60000), Emu(200000),
       "STEP 1", 7, bold=True, color=C_GRAY)
    tb(s, fc_x + Emu(40000), n1_y + Emu(240000), fc_w - Emu(60000), Emu(260000),
       "通常時 → 海戦BONUS当選", 10.5, bold=True, color=C_WHITE, font=FONT_H)

    arrow_d(s, fc_x + fc_w // 2, n1_y + n1_h + Emu(15000), C_CYAN)

    # STEP 2
    n2_y = n1_y + n1_h + Emu(215000)
    n2_h = Emu(760000)
    rect_b(s, fc_x, n2_y, fc_w, n2_h, RGBColor(0x06, 0x10, 0x28), C_CYAN, 2.0)
    rect(s, fc_x, n2_y, Emu(35000), n2_h, C_CYAN)
    tb(s, fc_x + Emu(60000), n2_y + Emu(45000), fc_w - Emu(80000), Emu(200000),
       "STEP 2", 7, bold=True, color=C_CYAN)
    tb(s, fc_x + Emu(60000), n2_y + Emu(240000), fc_w - Emu(80000), Emu(260000),
       "海戦BONUS消化（ベルナビ10回+10G）", 10.5, bold=True, color=C_CYAN2, font=FONT_H)
    tb(s, fc_x + Emu(60000), n2_y + Emu(510000), fc_w - Emu(80000), Emu(200000),
       "平均約100枚獲得。消化中にBAR揃いで陣営参戦の可能性。", 7.5, color=C_GRAY)

    arrow_d(s, fc_x + fc_w // 2, n2_y + n2_h + Emu(15000), C_CYAN)

    # STEP 3
    n3_y = n2_y + n2_h + Emu(215000)
    n3_h = Emu(760000)
    rect_b(s, fc_x, n3_y, fc_w, n3_h, RGBColor(0x08, 0x14, 0x30), C_GOLD, 2.0)
    rect(s, fc_x, n3_y, Emu(35000), n3_h, C_GOLD)
    tb(s, fc_x + Emu(60000), n3_y + Emu(45000), fc_w - Emu(80000), Emu(200000),
       "STEP 3", 7, bold=True, color=C_GOLD)
    tb(s, fc_x + Emu(60000), n3_y + Emu(240000), fc_w - Emu(80000), Emu(260000),
       "BATTLE演出でAT抽選！", 10.5, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, fc_x + Emu(60000), n3_y + Emu(510000), fc_w - Emu(80000), Emu(200000),
       "勝利 or 告知発生でアズールレーンRUSH突入確定。", 7.5, color=C_GRAY)

    arrow_d(s, fc_x + fc_w // 2, n3_y + n3_h + Emu(15000), C_GOLD)

    # STEP 4
    n4_y = n3_y + n3_h + Emu(215000)
    n4_h = Emu(560000)
    rect_b(s, fc_x, n4_y, fc_w, n4_h, RGBColor(0x04, 0x10, 0x2C), C_GOLD2, 2.0)
    rect(s, fc_x, n4_y, Emu(35000), n4_h, C_GOLD2)
    tb(s, fc_x + Emu(60000), n4_y + Emu(45000), fc_w - Emu(80000), Emu(200000),
       "STEP 4", 7, bold=True, color=C_GOLD2)
    tb(s, fc_x + Emu(60000), n4_y + Emu(240000), fc_w - Emu(80000), Emu(260000),
       "アズールレーンRUSH突入！", 10, bold=True, color=C_GOLD, font=FONT_H)

    # 右パネル
    rx, ry = Inches(4.9), Inches(0.72)
    rw = Inches(4.85)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x10, 0x22, 0x44))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "艦隊・海戦要素とゲーム性の融合", 10, bold=True, color=C_CYAN)

    items = [
        (C_CYAN,  "高確・超高確（内部状態）と海戦設計",
         "通常時は「通常・高確・超高確」の3状態が存在。\n"
         "高確・超高確ほど海戦BONUS当選率がアップ。\n"
         "状態変化はレア役や規定G数消化で発生。\n"
         "艦隊の哨戒→交戦強度上昇のイメージ。"),
        (C_GOLD,  "BAR揃い＝陣営参戦（海戦BONUS中）",
         "海戦BONUS消化中にBAR揃い発生→陣営が参戦。\n"
         "1陣営から最大4陣営まで参戦数が増えると\n"
         "AT中の上乗せ性能・特化ゾーン突入率が向上。\n"
         "陣営増加は艦隊増員のIPストーリーと一致。"),
        (C_PINK,  "共同戦線RUSH（1陣営時の特化ゾーン）",
         "AT中に1陣営のみの状態で発動する特化ゾーン。\n"
         "上乗せ性能が高く、ここで陣営を追加しながら\n"
         "ゲーム数・性能を底上げするゲームプラン。\n"
         "艦隊が少ない時こそ奮起する設計思想。"),
    ]
    for i, (ac, t, b) in enumerate(items):
        iy = ry + Emu(260000) + i * Emu(1485000)
        rect_b(s, rx, iy, rw, Emu(1415000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(1415000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(970000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "海戦BONUS→BATTLE演出のCZ構造：ボーナス消化そのものが「艦隊の交戦」を演出するIP連動設計",
           "BAR揃いによる陣営参戦が出玉に直結する点がアズレン独自ゲーム性の核心")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（出玉が伸びる仕組み）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT/ボーナス ── 出玉が伸びるメカニズム詳解", "5/9")

    # 左：ATの上乗せ仕組み
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x10, 0x22, 0x44))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "アズールレーンRUSH（通常AT）の出玉構造", 10, bold=True, color=C_CYAN)

    at_items = [
        (C_CYAN,  "AT基本仕様（純増約2.5枚/G）",
         "ゲーム数上乗せ型AT。純増約2.5枚/G。\n"
         "初期ゲーム数からスタートし、\n"
         "上乗せで継続ゲーム数を積み上げる。\n"
         "ゲーム数が0になると終了→天城BATTLE発展の可能性。"),
        (C_CYAN2, "上乗せ抽選のトリガー",
         "通常AT中の上乗せは以下から発生：\n"
         "① レア役（強チェリー・スイカ・チャンス目）\n"
         "② 明石チャンス（リプレイ積算→突入）\n"
         "③ 海戦BONUS消化中の上乗せ抽選\n"
         "④ 各陣営の特殊演出（陣営依存）"),
        (C_GOLD,  "陣営数に応じたAT性能変化",
         "陣営が増えるほどAT中の上乗せ期待値が向上。\n"
         "1陣営：基本性能\n"
         "2陣営：性能アップ（上乗せ期待値+）\n"
         "3陣営：さらに上乗せ頻度・量が増加\n"
         "4陣営：最高性能＋天城BATTLE発生権利獲得"),
        (C_PINK,  "共同戦線RUSHで一気に底上げ",
         "AT中1陣営状態で発動する特化ゾーン。\n"
         "高確率で大量上乗せが発生し、\n"
         "陣営追加のチャンスでもある。\n"
         "1陣営→4陣営達成の重要な踏み台。"),
    ]
    for i, (ac, t, b) in enumerate(at_items):
        iy = ly + Emu(290000) + i * Emu(1145000)
        rect_b(s, lx, iy, lw, Emu(1075000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1075000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(650000),
           b, 7.5, color=C_WHITE)

    # 右：ボーナス種類と出玉役割
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x10, 0x22, 0x44))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "ボーナス種類と出玉・AT性能への役割", 10, bold=True, color=C_CYAN)

    bonuses = [
        (C_CYAN,  "海戦BONUS（疑似ボーナス・メイン）",
         "通常時の主要当選役。ベルナビ10回+10G継続。\n"
         "平均約100枚獲得。BATTLE演出でAT抽選。\n"
         "BAR揃い時に陣営参戦抽選。頻繁に当選するため\n"
         "スルー回数管理が立ち回り上の重要指標。"),
        (C_GOLD2, "天城BATTLE（上位CZ・AT後発生）",
         "4陣営集結後のAT終了時に発生するCZ。\n"
         "クリアすると上位AT「異次元性能SS RUSH」確定。\n"
         "上位AT自体が4陣営状態で動作するため\n"
         "通常ATと比べてAT性能が大幅に向上する。"),
        (C_CYAN2, "AT中の継続ボーナス（上乗せ源）",
         "AT中に当選するボーナスはゲーム数上乗せ源。\n"
         "上乗せ量はボーナス種別・陣営数・内部状態で変動。\n"
         "AT中はできるだけ多くのボーナスを引いて\n"
         "ゲーム数を積み上げていくことが基本戦略。"),
        (C_GRAY,  "エンディング・差枚リセット後の恩恵",
         "エンディング到達や差枚数条件で有利区間リセット。\n"
         "リセット後は天城BATTLE（上位CZ）に突入。\n"
         "上位ATへの再挑戦機会が保証される仕組み。\n"
         "上位AT後も同様の恩恵で高ループ設計。"),
    ]
    for i, (ac, t, b) in enumerate(bonuses):
        iy = ry + Emu(290000) + i * Emu(1145000)
        rect_b(s, rx, iy, rw, Emu(1075000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(45000), Emu(1075000), ac)
        tb(s, rx + Emu(75000), iy + Emu(50000), rw - Emu(100000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(75000), iy + Emu(320000), rw - Emu(100000), Emu(650000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "AT出玉の核心：陣営数×上乗せ量の比例関係。4陣営集結を早めることが最大の出玉戦略",
           "共同戦線RUSHは1陣営状態の救済特化ゾーン。入れば逆転の大チャンス")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位ATへの道 ── 天城BATTLE→異次元性能SS RUSHの遊び方", "6/9")

    # 上段：昇格ルートフロー
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(260000), RGBColor(0x10, 0x22, 0x44))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(210000),
       "上位AT到達ルート（4陣営集結から異次元性能SS RUSHまでの昇格階層）", 9, bold=True, color=C_CYAN)

    route_boxes = [
        (C_CYAN,  "アズールレーン\nRUSH（通常AT）",   "純増約2.5枚/G\n1〜4陣営変動\nメインAT"),
        (C_GOLD,  "4陣営集結\n→天城BATTLE",         "AT終了後発動\n上位AT決定戦\nクリアで昇格"),
        (C_GOLD2, "異次元性能SS\nRUSH（上位AT）",    "純増約5.1枚/G\n4陣営MAX\n期待約3500枚+"),
        (C_PINK,  "上位AT終了後\n→再天城BATTLE",     "上位AT後も\nCZ再突入\n高ループ設計"),
    ]
    bw_r = Inches(2.1)
    gap_r = Inches(0.26)
    sx_r = Inches(0.35)
    cy_r = Inches(1.42)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x06, 0x12, 0x28), ac, 2.0 if i >= 2 else 1.5)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(55000), Emu(370000), lbl, 10, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), cy_r - bh_r // 2 + Emu(460000),
           bw_r - Emu(45000), Emu(480000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw_r + Emu(12000), cy_r)

    # 中区切り
    rect(s, 0, Inches(2.06), SLIDE_W, Emu(5000), RGBColor(0x10, 0x2A, 0x44))

    # 下段左：天城BATTLEの遊び方
    lx2, ly2 = Inches(0.28), Inches(2.12)
    lw2 = Inches(4.55)
    lh2 = Emu(2620000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_GOLD, 1.8)
    rect(s, lx2, ly2, Emu(45000), lh2, C_GOLD)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "天城BATTLEの仕組みと遊び方", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【天城BATTLEとは】\n"
       "AT「アズールレーンRUSH」で4陣営集結後、\n"
       "AT終了時に発動する上位CZ（チャレンジゾーン）。\n"
       "クリアすれば純増5.1枚/Gの上位ATへ昇格。\n\n"
       "【発生タイミング】\n"
       "① 通常ATで4陣営集結→AT終了後\n"
       "② エンディング到達・差枚リセット後\n"
       "③ 上位AT終了後（再CZ突入）\n\n"
       "【CZ中の遊び方】\n"
       "演出を最後まで確認。勝利告知の\n"
       "タイミングと演出の豪華さに注目。\n"
       "キャラ登場・BGM変化で期待度把握。",
       8, color=C_WHITE)

    # 下段右：異次元性能SS RUSHの遊び方
    rx2, ry2 = Inches(5.05), Inches(2.12)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2,
           RGBColor(0x06, 0x12, 0x28), C_GOLD2, 2.0)
    rect(s, rx2, ry2, Emu(45000), lh2, C_GOLD2)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "異次元性能SS RUSH（上位AT）の遊び方", 11, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【上位ATの基本仕様】\n"
       "純増約5.1枚/G（通常ATの約2倍）。\n"
       "常に4陣営MAXの状態で動作するため\n"
       "上乗せ期待値も通常ATより大幅に高い。\n"
       "期待獲得枚数：約3,500枚以上。\n\n"
       "【ループ設計】\n"
       "上位AT終了後も再び天城BATTLEに突入。\n"
       "CZクリアで上位ATを繰り返せる設計。\n"
       "高ループによる大量出玉が期待できる。\n\n"
       "【打ち方・注意点】\n"
       "上位AT中は演出・告知を見逃さないこと。\n"
       "終了時の画面で設定示唆が表示される。\n"
       "消化中は離席せず全力消化が基本。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "上位ATの設計思想：4陣営MAX×純増5.1枚/G×ループCZで「爆発的出玉と夢の持続」を両立",
           "天城BATTLE突入時点で大きな期待値。上位AT後も再CZでループを狙えるのが本機の最大の魅力")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（艦隊美少女IPとパチスロの融合）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── 艦隊美少女IPとパチスロの完全融合", "7/9")

    principles = [
        (C_CYAN,  "① 陣営参戦がそのままAT強化になる設計",
         "キャラクター（陣営）が増えるほど出玉性能が上がる。\n"
         "「仲間を集める」というIPの物語体験が\n"
         "ゲームメカニクスと直結している点が革新的。\n"
         "原作ファンの感情移入を最大限に活用した設計。"),
        (C_GOLD,  "② 明石商店というIP活用のUIデザイン",
         "明石（整備キャラ）の商店がリプレイ積算の\n"
         "可視化UIとして機能。IPキャラを単なる演出飾りでなく\n"
         "ゲーム性の核心に組み込んでいる。\n"
         "通常時・AT中の両場面で存在感を発揮。"),
        (C_PINK,  "③ 海戦という世界観をゲームフローに昇華",
         "海戦BONUS→BATTLE演出というCZ構造が\n"
         "「艦隊が交戦する」というIPの世界観そのもの。\n"
         "「勝利でAT突入」という演出設計が\n"
         "IPストーリーとギャンブルの緊張感を融合させた。"),
        (C_CYAN2, "④ 純増2.5枚→5.1枚という体感差の演出",
         "通常ATと上位ATの純増差が約2倍。\n"
         "上位ATになった瞬間の出玉スピードの体感変化が\n"
         "「本当に別次元のATに突入した」と感じさせる。\n"
         "名前「異次元性能SS RUSH」が体感を強化。"),
        (C_GOLD2, "⑤ 天城BATTLEという上位AT挑戦の劇的設計",
         "4陣営集結という達成感→AT終了→天城BATTLE挑戦\n"
         "という流れが「積み上げ→解放」の感情設計。\n"
         "天城（キャラ）との決戦というIP演出が\n"
         "最大の興奮場面として機能する。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1200000)
    gy = Inches(0.10)

    positions = [
        (Inches(0.28),  Inches(0.72)),
        (Inches(5.17),  Inches(0.72)),
        (Inches(0.28),  Inches(0.72) + bh_p + gy),
        (Inches(5.17),  Inches(0.72) + bh_p + gy),
        (Inches(0.28),  Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1100000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(330000), pw - Emu(90000), ph - Emu(390000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "面白さの核心：「陣営参戦×AT性能直結×IP世界観×純増倍速体感×天城BATTLE劇的展開」の5要素連鎖",
           "アズールレーンのIPを単なるビジュアル飾りでなくゲームメカニクスに昇華した点がパチスロアワード評価の核心")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN, "IPキャラがゲーム性に直結（陣営システム）",
         "美少女キャラの参戦がAT性能に直結する設計。\n"
         "「好きなキャラが来ると嬉しい＆有利」という\n"
         "感情とゲーム性を同時に満たす二重構造。"),
        (C_GREEN, "浅い天井（350G）でライトユーザーに優しい",
         "350Gという浅めの天井設定が\n"
         "ライトユーザーでも安心して立ち回れる基盤。\n"
         "スルー天井との二重セーフティ設計。"),
        (C_GREEN, "通常AT→上位ATの純増倍速体感（2.5→5.1）",
         "上位AT突入時の出玉スピード変化が大きく\n"
         "「別次元に突入した」という強烈な体感を提供。\n"
         "上位AT到達後のゲーム体験の質が高い。"),
        (C_GREEN, "評価「京楽奇跡の1台」「個人的2025年ベスト台」",
         "複数のプレイヤー評価で高評価を獲得。\n"
         "ATの「やれる感」と演出の完成度が\n"
         "スマスロの中でも特に評価されている。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1140000)
        rect_b(s, lx, iy, lw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(300000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,   "コイン持ち低め（25.8G/50枚）の通常時リスク",
         "コイン持ちの悪さにより通常時の投資が\n"
         "増えやすい。ボーナスにハマると消耗が大きく\n"
         "ライトユーザーへの敷居になりうる。"),
        (C_ORG,   "優遇・冷遇システムによる設定依存度",
         "スマスロの設定優遇・冷遇で通常時もAT中も\n"
         "設定に大きく左右される可能性がある。\n"
         "「ゲーム性最高だが設定に弄られる」という評価も。"),
        (C_CYAN2, "4陣営集結までの道のりの長さ",
         "4陣営に到達するまでのAT消化時間が\n"
         "比較的長くなる場合がある。\n"
         "陣営が増えない引きが続くと単調感が出やすい。"),
        (C_GRAY,  "上位AT無しの\"通常打ち\"の物足りなさ",
         "上位ATを経験しないと本機の真価を\n"
         "感じにくい設計とも言える。\n"
         "天城BATTLE失敗時の落胆が大きい点も課題。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1140000)
        rect_b(s, rx, iy, rw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(300000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "強みと課題の両面把握：IP連動ゲーム性設計の成功事例として評価しつつ、コイン持ち・設定依存の課題も直視",
           "「京楽奇跡の1台」との評価はIP×ゲーム性融合の完成度への賛辞。設計の参考価値は非常に高い")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x10, 0x22, 0x44))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "アズールレーン ── 設計的強み総括", 10, bold=True, color=C_CYAN)

    strengths = [
        (C_CYAN,  "IPキャラ参戦＝AT強化という革新的設計",
         "陣営（キャラ）が増えるとAT性能が上がる構造。\n"
         "「仲間を集める」というIPの物語体験が\n"
         "出玉という実利と直結した設計の革新性。"),
        (C_GOLD,  "2段階AT（2.5枚→5.1枚）の体感設計",
         "通常ATから上位ATへの昇格で純増が2倍。\n"
         "「異次元」という台名通りの出玉体験差が\n"
         "高い達成感と継続プレイ動機を生む。"),
        (C_PINK,  "浅天井（350G）×スルー天井の二重セーフティ",
         "2種の天井が異なるハマりパターンをカバー。\n"
         "ライトユーザーからヘビーユーザーまで\n"
         "各層に応じた安心感を提供する設計。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1210000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(325000), lw - Emu(90000), Emu(780000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_CYAN, font=FONT_H)

    principles = [
        (C_CYAN,  "IPキャラをゲーム性の核に組み込め",
         "キャラを演出飾りに留めず\n出玉メカニクスに直結させることで\nファンと一般層の両方を取り込む"),
        (C_GOLD,  "「仲間が増える=強くなる」という普遍的快感",
         "RPGの成長感をパチスロに移植。\n陣営参戦というUIが感情と利益を\n同時に満たす双方向設計の鉄則"),
        (C_CYAN2, "二段階出玉スピードで昇格体験を演出せよ",
         "通常モード→上位モードの\n出玉速度差を設計に組み込むことで\n「別世界に来た感」を強制生成"),
        (C_GOLD2, "天城BATTLEのような「積み上げ→解放」設計",
         "4陣営達成という積み上げが\nCZ解放という感情解放に直結。\n努力報酬型の快感設計の見本"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(760000)
        rect_b(s, rx, py0, rw, Emu(710000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(710000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(230000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(285000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3310000), rw, Emu(1060000),
           RGBColor(0x04, 0x0E, 0x26), C_CYAN, 2.0)
    rect(s, rx, ry + Emu(3310000), Emu(40000), Emu(1060000), C_CYAN)
    tb(s, rx + Emu(65000), ry + Emu(3360000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_CYAN)
    tb(s, rx + Emu(65000), ry + Emu(3620000), rw - Emu(85000), Emu(690000),
       "陣営参戦×AT性能直結×IP世界観融合×二段階純増×天城BATTLE。\n"
       "これら5要素を統合した設計がパチスロアワード2025\n"
       "ノミネートの根拠。艦隊IPとパチスロの\n"
       "ゲーム性融合の模範事例として参照価値が高い。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "本機の設計思想：「IPキャラ×出玉直結×二段階AT×積み上げ解放」──アズレン設計4原則を次世代機に活用せよ",
           "パチスロアワード2025ノミネート台。京楽の設計転換点となった歴史的1台として記録すべき事例")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)     # 1: タイトル・スペック・3ポイント
    s_flow(prs)      # 2: ゲームフロー全体図
    s_normal(prs)    # 3: 通常時の遊び方
    s_cz(prs)        # 4: CZ/前兆の仕組み
    s_at(prs)        # 5: AT/ボーナス
    s_upper(prs)     # 6: 上位ATへの道
    s_design(prs)    # 7: 面白さの設計
    s_pros_cons(prs) # 8: 良い点と課題
    s_matome(prs)    # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
