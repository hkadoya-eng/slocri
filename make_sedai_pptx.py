"""
スマスロ新機種企画提案 「世代 ―継承の炎―」 PowerPointジェネレーター v2
出力: proposals/新規提案/世代の章/sedai_proposal_v2.pptx
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(
    os.path.dirname(__file__),
    "proposals", "新規提案", "世代の章", "sedai_proposal_v2.pptx"
)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# ── カラーパレット（深黒×炎オレンジ×金×紫）────────────────────────
C_BG    = RGBColor(0x06, 0x04, 0x14)
C_CARD  = RGBColor(0x10, 0x0A, 0x22)
C_ORG   = RGBColor(0xFF, 0x88, 0x00)   # 炎オレンジ
C_ORG2  = RGBColor(0xFF, 0xAA, 0x44)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)
C_PUR   = RGBColor(0x88, 0x33, 0xCC)   # 紫（継承・英霊）
C_PUR2  = RGBColor(0xAA, 0x66, 0xFF)
C_DKPUR = RGBColor(0x20, 0x08, 0x30)
C_RED   = RGBColor(0xCC, 0x22, 0x22)
C_DKRED = RGBColor(0x28, 0x04, 0x04)
C_GREEN = RGBColor(0x40, 0xAA, 0x40)
C_STEEL = RGBColor(0x44, 0x77, 0xAA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM = RGBColor(0xF0, 0xE4, 0xCC)
C_LTGRAY= RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY  = RGBColor(0x88, 0x88, 0x88)
C_YELLOW= RGBColor(0xFF, 0xEE, 0x44)
C_ORANGE= RGBColor(0xFF, 0xA0, 0x30)

FONT_H = "游明朝"
FONT_B = "メイリオ"


# ── 背景生成（炎と闇が交差するイメージ）────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (6, 4, 20))
    draw = ImageDraw.Draw(img)
    # 斜めグリッド（暗い格子）
    for i in range(0, w + h, 50):
        draw.line([(i, 0), (0, i)], fill=(16, 10, 38), width=1)
    # 底部の炎グラデーション
    for y in range(h - 140, h):
        t = (y - (h - 140)) / 140
        r = int(45 * t)
        g = int(12 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, 0))
    # 上部の薄い紫オーラ
    for y in range(0, 60):
        t = (60 - y) / 60 * 0.4
        draw.line([(0, y), (w, y)], fill=(int(20 * t), 0, int(35 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー関数 ──────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s

def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT_B
    if color:
        run.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp

def arrow_r(slide, x, y, w, color):
    h = Emu(150000)
    shp = slide.shapes.add_shape(13, x, y - h // 2, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def hdr(slide, text):
    rect(slide, Inches(0.15), Inches(0.08), Inches(9.7), Emu(420000),
         RGBColor(0x10, 0x06, 0x24))
    rect(slide, Inches(0.15), Inches(0.08), Emu(60000), Emu(420000), C_ORG)
    tb(slide, Inches(0.4), Inches(0.1), Inches(9.2), Emu(380000),
       text, 12, bold=True, color=C_GOLD, font=FONT_H)

def net_note(slide, text="※ネットより"):
    tb(slide, Inches(8.5), Inches(5.38), Inches(1.4), Emu(200000),
       text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・キャッチコピー・3本柱
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    """タイトルスライド：タイトル・キャッチコピー・この提案の3本柱"""
    s = new_slide(prs)

    # 炎ライン2本
    rect(s, Inches(0), Inches(1.82), Inches(10), Emu(7000), C_ORG)
    rect(s, Inches(0), Inches(3.60), Inches(10), Emu(7000), C_ORG)

    # サブタイトル
    tb(s, Inches(0.5), Inches(0.26), Inches(9), Emu(360000),
       "新機種企画提案  v2.0", 13, color=C_ORG2, font=FONT_H, align=PP_ALIGN.CENTER)

    # メインタイトル
    tb(s, Inches(0.3), Inches(0.72), Inches(9.4), Emu(950000),
       "世 代", 68, bold=True, color=C_GOLD2, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, Inches(0.3), Inches(1.55), Inches(9.4), Emu(280000),
       "―  継 承 の 炎  ―", 17, bold=True, color=C_ORG, font=FONT_H, align=PP_ALIGN.CENTER)

    # キャッチコピー帯
    rect(s, Inches(1.2), Inches(1.98), Inches(7.6), Emu(550000),
         RGBColor(0x14, 0x06, 0x28))
    tb(s, Inches(1.3), Inches(2.04), Inches(7.4), Emu(490000),
       "「 1人では届かない。でも、意志を継ぐことで、必ず倒せる。 」",
       14, bold=True, color=C_CREAM, font=FONT_H, align=PP_ALIGN.CENTER)

    # 3本柱
    pillars = [
        (Inches(0.2),  "① 感情設計",   "「負けても前進感がある」\n唯一の設計（積み上げ型）", C_ORG),
        (Inches(3.55), "② 世界観",     "不死の魔王×英雄の系譜\n1セッションが1つの戦争",     C_PUR),
        (Inches(6.9),  "③ 業界初",     "継承の儀（AT中3択）×\n宿敵弱体化×世代ループ複合",   C_GOLD),
    ]
    for x, label, desc, col in pillars:
        rect_b(s, x, Inches(3.75), Inches(2.9), Emu(950000),
               RGBColor(0x12, 0x08, 0x22), col, 1.8)
        tb(s, x + Emu(120000), Inches(3.81), Inches(2.6), Emu(330000),
           label, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(120000), Inches(4.22), Inches(2.6), Emu(480000),
           desc, 9, color=C_CREAM)

    tb(s, Inches(7.7), Inches(5.2), Inches(2.1), Emu(300000),
       "2026.05  v2.0", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: コアコンセプト（世代交代ループとは何か・感情設計）
# ══════════════════════════════════════════════════════════════
def s_concept(prs):
    """コアコンセプト：Re:ゼロとの対比・感情設計の6段階山場"""
    s = new_slide(prs)
    hdr(s, "コアコンセプト  ──  世代交代ループとは何か")

    # 左：Re:ゼロ対比表
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.5), Inches(3.5),
           C_CARD, C_STEEL, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.2), Emu(310000),
       "Re:ゼロ vs 世代  ――  設計思想の違い", 10, bold=True, color=C_STEEL, font=FONT_H)

    rows_compare = [
        ("失敗の意味", "タイムリープ（リセット）",  "継承（積み上げ）",     C_LTGRAY, C_ORG),
        ("感情",      "また失敗した…次へ",        "意志を次世代に渡した", C_LTGRAY, C_GOLD),
        ("設計思想",  "失敗前提ループ",            "積み上げ前提ループ",   C_LTGRAY, C_ORG2),
        ("終了感",    "あり（リセット感）",        "なし（継続感）",       C_LTGRAY, C_PUR2),
    ]
    # ヘッダー行
    rect(s, Inches(0.22), Inches(1.32), Inches(4.46), Emu(280000), RGBColor(0x1A, 0x14, 0x30))
    for txt, cx in [("軸", Inches(0.28)), ("Re:ゼロ", Inches(1.38)), ("世代", Inches(3.02))]:
        tb(s, cx, Inches(1.35), Inches(1.55), Emu(240000), txt, 8, bold=True, color=C_GOLD)
    ry = Inches(1.65)
    for axis, rzero, sedai, c1, c2 in rows_compare:
        bg = RGBColor(0x10, 0x0C, 0x22) if rows_compare.index((axis, rzero, sedai, c1, c2)) % 2 == 0 \
             else RGBColor(0x14, 0x10, 0x28)
        rect(s, Inches(0.22), ry, Inches(4.46), Emu(290000), bg)
        tb(s, Inches(0.28), ry + Emu(28000), Inches(1.05), Emu(240000), axis,  8,   color=C_GRAY, wrap=False)
        tb(s, Inches(1.38), ry + Emu(28000), Inches(1.58), Emu(240000), rzero, 7.5, color=c1,    wrap=False)
        tb(s, Inches(3.02), ry + Emu(28000), Inches(1.58), Emu(240000), sedai, 7.5, bold=True, color=c2, wrap=False)
        ry += Emu(295000)

    # 強調テキスト
    rect(s, Inches(0.22), Inches(3.05), Inches(4.46), Emu(250000), RGBColor(0x22, 0x10, 0x00))
    tb(s, Inches(0.3), Inches(3.09), Inches(4.3), Emu(220000),
       "「負けても前進感がある設計」は業界に存在しない",
       8.5, bold=True, color=C_ORG, font=FONT_H)

    # 右：感情山場6段階
    rect_b(s, Inches(4.85), Inches(0.85), Inches(4.95), Inches(3.5),
           C_CARD, C_ORG, 1.8)
    tb(s, Inches(4.97), Inches(0.92), Inches(4.7), Emu(310000),
       "感情の山場設計（6段階）", 10, bold=True, color=C_ORG, font=FONT_H)

    mountains = [
        ("1", "通常時",       "継承ポイントが貯まってきた…",     "積み上げ感", C_STEEL),
        ("2", "CZ",          "今世代で英雄になれるか",           "緊張",       C_GOLD),
        ("3", "継承の儀",     "何を次の世代に残すか",             "能動的選択", C_ORG),
        ("4", "英雄連合",     "3世代繋いだ",                    "達成感",     C_GREEN),
        ("5", "決戦AT発動",   "ついにこの日が来た",              "最大の興奮", C_RED),
        ("6", "宿敵討伐",     "積み上げた意志が実った",           "感動",       C_GOLD2),
    ]
    my = Inches(1.32)
    for num, scene, feel, emotion, col in mountains:
        rect(s, Inches(4.87), my, Emu(340000), Emu(270000), RGBColor(0x1A, 0x0C, 0x2C))
        tb(s, Inches(4.87), my + Emu(22000), Emu(340000), Emu(240000),
           num, 10, bold=True, color=col, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, Inches(5.38), my + Emu(22000), Inches(1.35), Emu(240000),
           scene, 8.5, bold=True, color=col, wrap=False)
        tb(s, Inches(6.78), my + Emu(22000), Inches(1.65), Emu(240000),
           feel, 7.5, color=C_LTGRAY, wrap=False)
        tb(s, Inches(8.48), my + Emu(22000), Inches(1.25), Emu(240000),
           f"[{emotion}]", 7.5, color=col, wrap=False)
        my += Emu(278000)

    # フッター
    rect(s, Inches(0.2), Inches(4.45), Inches(9.6), Emu(580000), RGBColor(0x0E, 0x08, 0x1E))
    rect(s, Inches(0.2), Inches(4.45), Emu(60000), Emu(580000), C_PUR)
    tb(s, Inches(0.45), Inches(4.50), Inches(9.2), Emu(260000),
       "設計思想：「積み上げ型ループ」はRPGの「プレイするほど世界が豊かになる」感覚のパチスロ版",
       9, bold=True, color=C_PUR2, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.80), Inches(9.2), Emu(240000),
       "タイムリープ（Re:ゼロ・まどマギ）が「失敗前提ループ」なら、世代は「積み上げ前提ループ」。この感情設計が業界初。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 世界観（主人公・宿敵・物語の骨格）
# ══════════════════════════════════════════════════════════════
def s_world(prs):
    """世界観：オリジナルIP・不死の魔王・英雄の系譜・物語骨格"""
    s = new_slide(prs)
    hdr(s, "世界観  ──  不死の魔王と英雄の系譜（オリジナルIP）")

    # 3カラム：登場人物
    chars = [
        (Inches(0.2),  "主人公",    "最後の英雄の末裔",
         "代々受け継がれる「英雄の末裔」\n1人では魔王に届かない\n\nプレイヤーは英雄の系譜を\n繋ぐ者",
         C_GOLD, RGBColor(0x1C, 0x14, 0x00)),
        (Inches(3.55), "宿敵",      "不死の魔王「イグナール」",
         "1000年間世界を支配する不死の魔王\n倒されるたびに蘇る\n\n「1世代では倒せない」\nという絶望的な存在",
         C_RED,  RGBColor(0x22, 0x04, 0x04)),
        (Inches(6.9),  "プレイヤー", "英雄の系譜を繋ぐ者",
         "世代から世代へ意志を繋ぐ\n継承の儀で何を渡すかを決める\n\n「自分が選んだ継承が\n物語を変える」",
         C_PUR,  RGBColor(0x18, 0x06, 0x28)),
    ]
    for x, role, name, desc, col, fill in chars:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(2.35), fill, col, 2.0)
        rect(s, x + Emu(80000), Inches(0.90), Inches(2.8), Emu(300000),
             RGBColor(0x08, 0x04, 0x18))
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(260000),
           role, 8, color=col, font=FONT_H, align=PP_ALIGN.CENTER)
        tb(s, x + Emu(80000), Inches(1.25), Inches(2.8), Emu(290000),
           name, 10, bold=True, color=C_WHITE, font=FONT_H, align=PP_ALIGN.CENTER)
        tb(s, x + Emu(80000), Inches(1.62), Inches(2.8), Inches(1.3),
           desc, 8.5, color=C_CREAM)

    # 物語の骨格（フローテキスト）
    rect_b(s, Inches(0.2), Inches(3.30), Inches(9.6), Inches(1.22),
           RGBColor(0x10, 0x08, 0x20), C_ORG, 1.5)
    tb(s, Inches(0.32), Inches(3.36), Inches(9.2), Emu(310000),
       "物語の骨格", 10, bold=True, color=C_ORG, font=FONT_H)

    story_steps = [
        ("第1世代", "魔王に挑む → 力尽きる", C_LTGRAY),
        ("継承",    "「継承の炎」を次の英雄へ", C_ORG),
        ("第2世代", "弱点を突く → また力尽きる", C_LTGRAY),
        ("継承",    "「速さ」と「絆」を次へ",  C_ORG),
        ("英雄連合","かつての英雄たちが集う！",  C_GOLD),
        ("決戦",    "1000年ぶりにイグナールが倒れる", C_RED),
    ]
    sx = Inches(0.3)
    for step, desc, col in story_steps:
        w = Inches(1.57)
        rect(s, sx, Inches(3.75), w, Emu(580000), RGBColor(0x16, 0x0C, 0x28))
        tb(s, sx + Emu(40000), Inches(3.80), w - Emu(80000), Emu(260000),
           step, 7.5, bold=True, color=col, font=FONT_H, align=PP_ALIGN.CENTER)
        tb(s, sx + Emu(40000), Inches(4.08), w - Emu(80000), Emu(260000),
           desc, 7, color=C_CREAM, align=PP_ALIGN.CENTER)
        if sx < Inches(8.8):
            arrow_r(s, sx + w + Emu(20000), Inches(4.23), Emu(100000), C_ORG)
        sx += w + Emu(100000)

    # フッター
    rect(s, Inches(0.2), Inches(4.55), Inches(9.6), Emu(480000), RGBColor(0x12, 0x08, 0x1E))
    rect(s, Inches(0.2), Inches(4.55), Emu(60000), Emu(480000), C_GOLD)
    tb(s, Inches(0.45), Inches(4.60), Inches(9.2), Emu(260000),
       "テーマ：「1人では届かない。でも、意志を継ぐことで、必ず倒せる。」",
       10, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.90), Inches(9.2), Emu(230000),
       "IPなし・完全オリジナル。炎と闇のビジュアルテーマ。英雄の系譜という普遍的な物語構造で幅広い層に響く。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    """ゲームフロー全体図：通常時→CZ→AT→英雄連合→決戦ATの蛇行2段"""
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図  ──  1セッションで世代が積み重なる")

    BW  = Inches(3.0)
    BH  = Inches(1.15)
    GAP = Inches(0.2)
    R1Y = Inches(0.55)
    R2Y = Inches(2.32)
    X1  = Inches(0.2)
    X2  = X1 + BW + GAP
    X3  = X2 + BW + GAP

    # Row1
    row1 = [
        (X1, "通常時 / 前兆",
         "4ステージ（荒野→黄金都市）\n周期100GでCZ到来",
         RGBColor(0x08, 0x06, 0x1C), C_STEEL),
        (X2, "CZ「英雄の試練」",
         "3Gバトル / 成功でAT突入\n失敗でも継承ポイント蓄積",
         RGBColor(0x1C, 0x08, 0x00), C_ORG),
        (X3, "世代AT（1世代50G）",
         "純増4.0枚/G / ストック4種\n赤7揃い → 継承の儀発動！",
         RGBColor(0x18, 0x0C, 0x00), C_GOLD),
    ]
    for x, title, desc, fill, bdr in row1:
        rect_b(s, x, R1Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R1Y + Emu(60000), BW - Emu(160000), Emu(330000),
           title, 10, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R1Y + Emu(410000), BW - Emu(160000), BH - Emu(480000),
           desc, 9, color=C_CREAM)
    for x_l in [X1, X2]:
        arrow_r(s, x_l + BW + Emu(40000), R1Y + BH // 2, GAP - Emu(80000), C_GOLD)

    # 折り返し下向き矢印（X3 row1 → row2）
    AT_CX = X3 + BW // 2
    _aw, _ah = Emu(130000), Emu(360000)
    shp_d = s.shapes.add_shape(13, AT_CX - _aw // 2, R1Y + BH + Emu(70000), _aw, _ah)
    shp_d.rotation = 90
    shp_d.fill.solid()
    shp_d.fill.fore_color.rgb = C_GOLD
    shp_d.line.fill.background()

    # Row2
    row2 = [
        (X3, "世代ループ（継承の儀）",
         "世代が重なるほど台が強化\nストック残 → 次世代へ継続！",
         RGBColor(0x18, 0x0C, 0x00), C_GOLD),
        (X2, "英雄連合（3世代以上）",
         "複数世代の英雄が集結\n継続率75% / 純増6.5枚/G",
         RGBColor(0x04, 0x14, 0x04), C_GREEN),
        (X1, "英雄の残照（引き戻し）",
         "AT終了後5Gのラストチャンス\n絆の継承→確定発動！",
         RGBColor(0x06, 0x08, 0x1C), C_STEEL),
    ]
    for x, title, desc, fill, bdr in row2:
        rect_b(s, x, R2Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R2Y + Emu(60000), BW - Emu(160000), Emu(330000),
           title, 10, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R2Y + Emu(410000), BW - Emu(160000), BH - Emu(480000),
           desc, 9, color=C_CREAM)
    for x_r in [X3, X2]:
        _w = GAP - Emu(80000)
        _h = Emu(150000)
        shp = s.shapes.add_shape(13, x_r - GAP + Emu(40000),
                                  R2Y + BH // 2 - _h // 2, _w, _h)
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = C_ORG
        shp.line.fill.background()

    # ⊓ループバック
    LW   = Emu(55000)
    lx_l = X3 + Emu(210000)
    lx_r = X3 + Emu(620000)
    lp_y = R2Y - Emu(330000)
    rect(s, lx_l - LW // 2, lp_y, LW, R2Y - lp_y, C_GOLD)
    rect(s, lx_l - LW // 2, lp_y - LW // 2, lx_r - lx_l + LW, LW, C_GOLD)
    rect(s, lx_r - LW // 2, lp_y, LW, R2Y - lp_y, C_GOLD)
    tb(s, lx_l + Emu(60000), lp_y + Emu(40000),
       lx_r - lx_l - Emu(60000), Emu(250000),
       "↺ ループ！", 8, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)

    # 決戦AT（下部）
    BOT_Y = Inches(3.70)
    rect_b(s, X1, BOT_Y, Inches(2.9), Emu(870000),
           RGBColor(0x20, 0x04, 0x04), C_RED, 2.0)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(60000), Inches(2.6), Emu(340000),
       "決戦AT  ―世代の決着―", 10, bold=True, color=C_RED, font=FONT_H)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(430000), Inches(2.6), Emu(410000),
       "宿敵弱体化MAX\n＋皇帝の力＋継承コンプリート\n純増8.0枚/G", 8.5, color=C_CREAM)

    rect(s, Inches(3.1), BOT_Y, Inches(6.7), Emu(870000),
         RGBColor(0x10, 0x06, 0x1C))
    tb(s, Inches(3.25), BOT_Y + Emu(80000), Inches(6.3), Emu(360000),
       "歴代全世代の英雄が集結。不死の魔王イグナールを1000年ぶりに討伐する大演出。",
       9, bold=True, color=C_ORG2)
    tb(s, Inches(3.25), BOT_Y + Emu(470000), Inches(6.3), Emu(370000),
       "「自分が今日育てた世代が、ここで結実する」——1セッションの物語が完結する瞬間。MY約3,500枚。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 通常時の設計（何をする・継承ポイント・CZルート）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    """通常時の設計：ステージ4種・周期CZ・継承ポイント・英霊の導き"""
    s = new_slide(prs)
    hdr(s, "通常時の設計  ──  何をする台か・3行で")

    # 3行サマリー
    rect(s, Inches(0.2), Inches(0.82), Inches(9.6), Emu(310000),
         RGBColor(0x16, 0x0A, 0x28))
    tb(s, Inches(0.32), Inches(0.87), Inches(9.3), Emu(270000),
       "100G周期でCZが来る。失敗しても継承ポイントが貯まる。貯まるほど天井が短くなり、300ptで高確モード「英霊の導き」が発動する。",
       9.5, bold=True, color=C_ORG2)

    # 3カラム
    cols = [
        (Inches(0.2),  "ステージ4種（遊技数連動）", C_STEEL,
         [("荒野",     "通常ステージ",         C_LTGRAY),
          ("城下町",   "高確示唆",             C_YELLOW),
          ("王都",     "前兆・CZ高確率",       C_ORANGE),
          ("黄金都市", "AT超高確率状態",       C_GOLD2)]),
        (Inches(3.45), "周期CZ「英雄の試練」", C_ORG,
         [("周期",   "100G毎に規則的に到来",   C_CREAM),
          ("バトル", "3G間の試練演出",         C_LTGRAY),
          ("成功",   "世代AT突入",             C_GOLD2),
          ("失敗",   "継承ポイント蓄積+10pt", C_ORG2)]),
        (Inches(6.7),  "継承ポイントシステム",  C_GOLD,
         [("50pt",  "天井100G短縮（→500G）",   C_YELLOW),
          ("100pt", "天井200G短縮（→400G）",   C_ORANGE),
          ("200pt", "CZ突破率大幅UP",           C_ORG),
          ("300pt", "英霊の導き発動！",         C_GOLD2)]),
    ]
    for x, title, col, items in cols:
        rect_b(s, x, Inches(1.22), Inches(3.1), Inches(2.65), C_CARD, col, 1.5)
        tb(s, x + Emu(80000), Inches(1.27), Inches(2.9), Emu(310000),
           title, 9.5, bold=True, color=col, font=FONT_H)
        iy = Inches(1.65)
        for label, desc, ic in items:
            tb(s, x + Emu(80000), iy, Inches(1.15), Emu(260000),
               label, 8.5, bold=True, color=ic, wrap=False)
            tb(s, x + Emu(1000000), iy, Inches(2.05), Emu(260000),
               desc, 8.5, color=C_CREAM, wrap=False)
            iy += Emu(268000)

    # 英霊の導きボックス
    rect_b(s, Inches(0.2), Inches(4.02), Inches(9.6), Emu(540000),
           RGBColor(0x18, 0x08, 0x30), C_PUR, 2.0)
    tb(s, Inches(0.32), Inches(4.08), Inches(4.0), Emu(310000),
       "高確モード「英霊の導き」（継承300pt到達で発動）",
       10, bold=True, color=C_PUR2, font=FONT_H)
    tb(s, Inches(0.32), Inches(4.40), Inches(4.5), Emu(270000),
       "連続CZ状態。CZ失敗でも次の周期CZが20G後に来る。\n「貯めれば貯めるほど有利になる」具体的な実感。",
       8.5, color=C_CREAM)
    tb(s, Inches(5.0), Inches(4.08), Inches(4.6), Emu(580000),
       "通常時の読み方：ステージが「王都」以上になったら好機。\n失敗しても「育てている」感覚＝ストレスフリー設計。\n★ 周期CZは吉宗の規則的設計を参考に昇華。",
       8.5, color=C_GOLD)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: AT「世代AT」の遊び方（継承の儀3択・効果）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    """世代AT：継承の儀3択の具体的な効果・ストック4種・世代の重なり"""
    s = new_slide(prs)
    hdr(s, "世代AT  ──  継承の儀で「何を次の世代に残すか」を選ぶ")

    # 左：世代AT基本
    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.0), Inches(3.45),
           RGBColor(0x18, 0x0C, 0x00), C_GOLD, 2.0)
    tb(s, Inches(0.3), Inches(0.90), Inches(2.8), Emu(320000),
       "世代AT", 14, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(0.3), Inches(1.28), Inches(2.8), Inches(2.7),
       "1セット50G / 純増4.0枚/G\n\n【ストック4種】\n"
       "  一代の力：ループ率20%\n"
       "  二代の力：ループ率40%\n"
       "  三代の力：ループ率60%\n"
       "  皇帝の力：ループ率80%\n\n"
       "赤7揃いで\n「継承の儀」が発動！",
       9, color=C_CREAM)

    arrow_r(s, Inches(3.3), Inches(2.52), Emu(270000), C_ORG)

    # 中央：継承の儀（3択・効果明確化）
    rect_b(s, Inches(3.75), Inches(0.85), Inches(3.35), Inches(3.45),
           RGBColor(0x18, 0x08, 0x00), C_ORG, 2.5)
    tb(s, Inches(3.85), Inches(0.90), Inches(3.15), Emu(310000),
       "継承の儀（業界初）", 12, bold=True, color=C_ORG, font=FONT_H)
    tb(s, Inches(3.85), Inches(1.25), Inches(3.15), Emu(270000),
       "赤7揃いで発動。何を次世代に残すか選ぶ。", 8.5, color=C_LTGRAY)

    choices = [
        ("速さの継承", "次セット開始から10G高確\nレア役確率1.5倍",         C_GOLD2),
        ("深さの継承", "宿敵弱体化ゲージが\n通常の3倍速で貯まる",          C_RED),
        ("絆の継承",   "「英雄の残照」5Gが\n確定発動（引き戻し保証）",     C_PUR2),
    ]
    cy = Inches(1.62)
    for name, effect, col in choices:
        rect(s, Inches(3.85), cy, Inches(3.05), Emu(540000), RGBColor(0x1A, 0x0C, 0x00))
        tb(s, Inches(3.97), cy + Emu(40000), Inches(2.8), Emu(250000),
           name, 9, bold=True, color=col, font=FONT_H)
        tb(s, Inches(3.97), cy + Emu(290000), Inches(2.8), Emu(270000),
           effect, 8, color=C_CREAM)
        cy += Emu(560000)

    arrow_r(s, Inches(7.2), Inches(2.52), Emu(270000), C_GOLD)

    # 右：世代の重なり＋継承コンプリート
    rect_b(s, Inches(7.65), Inches(0.85), Inches(2.15), Inches(3.45),
           RGBColor(0x16, 0x10, 0x00), C_GOLD2, 2.5)
    tb(s, Inches(7.72), Inches(0.90), Inches(2.0), Emu(320000),
       "世代の重なり", 11, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, Inches(7.72), Inches(1.28), Inches(2.0), Inches(2.2),
       "第1世代\n  ↓ 継承\n第2世代\n  ↓ 継承\n第3世代\n  ↓\n英雄連合へ！\n\n3択すべて選ぶ\n→ 継承コンプリート\n→ 決戦AT条件①",
       8.5, color=C_CREAM)

    # フッター
    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(640000),
         RGBColor(0x10, 0x06, 0x1C))
    rect(s, Inches(0.2), Inches(4.40), Emu(60000), Emu(640000), C_ORG)
    tb(s, Inches(0.45), Inches(4.45), Inches(9.2), Emu(270000),
       "継承の儀の戦略性：「速さ」はループ率UP・「深さ」は決戦AT近道・「絆」は引き戻し保証",
       10, bold=True, color=C_ORG2, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.78), Inches(9.2), Emu(270000),
       "v1の課題「効果が曖昧」を完全解消。プレイヤーが戦略を持って選べる3択設計。ミリゴZ-ZONEを昇華。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 英雄連合（3世代ループ達成の報酬）
# ══════════════════════════════════════════════════════════════
def s_union(prs):
    """英雄連合：3世代ループ達成・継続率75%・英雄の残照"""
    s = new_slide(prs)
    hdr(s, "英雄連合  ──  3世代の意志が1つになる瞬間")

    # 英雄連合ボックス（左）
    rect_b(s, Inches(0.2), Inches(0.85), Inches(5.8), Inches(2.85),
           RGBColor(0x04, 0x14, 0x04), C_GREEN, 2.0)
    tb(s, Inches(0.32), Inches(0.91), Inches(5.6), Emu(330000),
       "英雄連合", 16, bold=True, color=RGBColor(0x60, 0xD0, 0x60), font=FONT_H)

    union_items = [
        ("発動条件", "世代AT 3回以上ループ達成（3世代の英雄が揃う）",   C_CREAM),
        ("継続率",   "75%",                                            C_GOLD2),
        ("純増",     "6.5枚/G",                                        C_GOLD2),
        ("期待枚数", "1,000〜1,500枚",                                  C_ORG2),
        ("演出",     "各世代の継承内容が演出に反映！",                   C_GREEN),
    ]
    uy = Inches(1.38)
    for label, val, col in union_items:
        rect(s, Inches(0.3), uy, Inches(1.6), Emu(265000), RGBColor(0x0A, 0x18, 0x0A))
        tb(s, Inches(0.35), uy + Emu(28000), Inches(1.5), Emu(220000),
           label, 8.5, bold=True, color=C_GRAY, wrap=False)
        tb(s, Inches(1.96), uy + Emu(28000), Inches(3.9), Emu(220000),
           val, 9, bold=True, color=col, wrap=False)
        uy += Emu(272000)

    # 英雄の残照（右）
    rect_b(s, Inches(6.2), Inches(0.85), Inches(3.6), Inches(2.85),
           RGBColor(0x06, 0x08, 0x1C), C_STEEL, 1.8)
    tb(s, Inches(6.32), Inches(0.91), Inches(3.4), Emu(330000),
       "英雄の残照（引き戻し）", 11, bold=True, color=C_STEEL, font=FONT_H)
    tb(s, Inches(6.32), Inches(1.32), Inches(3.4), Inches(2.0),
       "英雄連合終了後に発動する5Gのラストチャンス\n\n"
       "「絆の継承」選択 → 確定発動\n通常時         → 確率発動\n\n"
       "奇数図柄揃いで世代AT再突入！\n★ 吉宗1G連の感動をここで昇華",
       9, color=C_CREAM)

    # 演出詳細
    rect_b(s, Inches(0.2), Inches(3.82), Inches(9.6), Emu(680000),
           RGBColor(0x08, 0x10, 0x08), C_GREEN, 1.5)
    tb(s, Inches(0.32), Inches(3.88), Inches(9.3), Emu(290000),
       "演出設計：継承内容が英雄連合に反映される", 10, bold=True, color=C_GREEN, font=FONT_H)

    演出 = [
        ("速さを継いだ世代", "素早い連続攻撃演出",     C_GOLD2),
        ("深さを継いだ世代", "宿敵の弱点を突く必殺技", C_RED),
        ("絆を継いだ世代",   "他の英雄を援護する演出", C_PUR2),
    ]
    ex = Inches(0.35)
    for title, effect, col in 演出:
        rect(s, ex, Inches(4.26), Inches(3.1), Emu(260000), RGBColor(0x0C, 0x18, 0x0C))
        tb(s, ex + Emu(60000), Inches(4.31), Inches(2.9), Emu(220000),
           f"{title}：{effect}", 8.5, color=col, wrap=False)
        ex += Inches(3.18)

    rect(s, Inches(0.2), Inches(4.58), Inches(9.6), Emu(420000),
         RGBColor(0x0E, 0x06, 0x1C))
    tb(s, Inches(0.35), Inches(4.63), Inches(9.2), Emu(380000),
       "MY設計：世代AT(400枚) × 複数 + 英雄連合(1,200枚) + 決戦AT(1,500枚) ＝ 約3,500枚\n"
       "通常遊技でも英雄連合まで楽しめる。決戦ATは夢ではなく「積み上げれば届く」距離感。",
       8.5, color=C_GOLD)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 決戦AT「世代の決着」
# ══════════════════════════════════════════════════════════════
def s_kessen(prs):
    """決戦AT：発動条件3つ・感動の演出・出玉設計"""
    s = new_slide(prs)
    hdr(s, "決戦AT「世代の決着」  ──  積み上げた意志が報われる瞬間")

    # 発動条件（3条件）
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.5), Inches(2.9),
           RGBColor(0x20, 0x04, 0x04), C_RED, 2.5)
    tb(s, Inches(0.32), Inches(0.91), Inches(4.3), Emu(320000),
       "発動条件（3つすべて必要）", 11, bold=True, color=C_RED, font=FONT_H)

    conditions = [
        ("①", "宿敵弱体化ゲージ MAX",
         "「深さの継承」を重ねるほど早く貯まる",     C_ORG),
        ("②", "皇帝の力ストック獲得",
         "レアストック（ループ率80%）を入手",        C_GOLD),
        ("③", "継承コンプリート",
         "速さ・深さ・絆の3択すべてを選び揃える",    C_PUR2),
    ]
    cy = Inches(1.40)
    for num, cond, detail, col in conditions:
        rect(s, Inches(0.32), cy, Emu(360000), Emu(550000), RGBColor(0x28, 0x06, 0x06))
        tb(s, Inches(0.32), cy + Emu(55000), Emu(360000), Emu(300000),
           num, 14, bold=True, color=col, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, Inches(1.05), cy + Emu(55000), Inches(3.5), Emu(270000),
           cond, 9, bold=True, color=col, font=FONT_H)
        tb(s, Inches(1.05), cy + Emu(330000), Inches(3.5), Emu(260000),
           detail, 8, color=C_LTGRAY)
        cy += Emu(575000)

    # 演出設計（右）
    rect_b(s, Inches(4.85), Inches(0.85), Inches(4.95), Inches(2.9),
           RGBColor(0x1A, 0x04, 0x04), C_RED, 2.0)
    tb(s, Inches(4.97), Inches(0.91), Inches(4.7), Emu(320000),
       "感動の演出設計", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, Inches(4.97), Inches(1.30), Inches(4.7), Inches(2.0),
       "純増：8.0枚/G\n\n"
       "歴代全世代の英雄が集結。\n宿敵イグナールとの最終決戦。\n\n"
       "「積み上げた意志が1つの光になる」\n── クライマックス演出\n\n"
       "1000年ぶりの宿敵討伐\n→ 感動の幕引き",
       9, color=C_CREAM)

    # 出玉設計表
    rect_b(s, Inches(0.2), Inches(3.85), Inches(9.6), Emu(1020000),
           RGBColor(0x0E, 0x06, 0x1C), C_GOLD, 1.5)
    tb(s, Inches(0.32), Inches(3.91), Inches(4.0), Emu(290000),
       "出玉設計", 10, bold=True, color=C_GOLD, font=FONT_H)

    out_data = [
        ("通常AT（1〜2世代）", "300〜400枚",  "最短体験の日"),
        ("英雄連合",           "1,000〜1,500枚", "3世代ループ達成の日"),
        ("決戦AT",            "1,500枚以上", "全条件成立の日"),
        ("MY目安（1日全部）",  "約3,500枚",  "決戦+前後込み"),
    ]
    ox = [Inches(0.28), Inches(3.0), Inches(6.2), Inches(8.5)]
    tb(s, ox[0], Inches(4.26), Inches(2.5), Emu(260000), "状態",     8, bold=True, color=C_GOLD)
    tb(s, ox[1], Inches(4.26), Inches(3.1), Emu(260000), "期待枚数", 8, bold=True, color=C_GOLD)
    tb(s, ox[2], Inches(4.26), Inches(2.2), Emu(260000), "備考",     8, bold=True, color=C_GOLD)
    oy = Inches(4.54)
    for j, (state, amt, note) in enumerate(out_data):
        bg = RGBColor(0x12, 0x08, 0x22) if j % 2 == 0 else RGBColor(0x18, 0x0C, 0x26)
        rect(s, Inches(0.22), oy, Inches(9.56), Emu(270000), bg)
        col_a = C_GOLD2 if state == "MY目安（1日全部）" else C_CREAM
        col_b = C_RED if state == "決戦AT" else (C_GOLD2 if state == "MY目安（1日全部）" else C_ORG2)
        tb(s, ox[0], oy + Emu(28000), Inches(2.5), Emu(230000), state, 8, color=col_a, wrap=False)
        tb(s, ox[1], oy + Emu(28000), Inches(3.1), Emu(230000), amt,   9, bold=True, color=col_b, wrap=False)
        tb(s, ox[2], oy + Emu(28000), Inches(2.2), Emu(230000), note,  7.5, color=C_GRAY, wrap=False)
        oy += Emu(275000)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: 既存機種との差別化
# ══════════════════════════════════════════════════════════════
def s_diff(prs):
    """既存機種との差別化：Re:ゼロ・吉宗・ミリゴとの比較表"""
    s = new_slide(prs)
    hdr(s, "既存機種との差別化  ──  世代が埋める「空席」はどこか")

    # 比較表
    headers = ["比較軸", "Re:ゼロ", "真打吉宗", "ミリゴ", "世代 ―継承の炎―"]
    col_w   = [Inches(1.55), Inches(1.88), Inches(1.88), Inches(1.88), Inches(2.21)]

    # ヘッダー行
    cx = Inches(0.2)
    for i, (h, w) in enumerate(zip(headers, col_w)):
        bg = RGBColor(0x22, 0x12, 0x00) if i == 4 else RGBColor(0x18, 0x14, 0x2C)
        col = C_ORG if i == 4 else C_GOLD
        rect(s, cx, Inches(0.85), w, Emu(350000), bg)
        tb(s, cx + Emu(40000), Inches(0.90), w - Emu(60000), Emu(300000),
           h, 9, bold=True, color=col, font=FONT_H)
        cx += w

    rows = [
        ("ループ設計",    "タイムリープ（リセット）", "周期CZ+1G連",        "A〜D4段階ストック",  "世代交代（積み上げ）"),
        ("失敗の意味",    "やり直し",               "継承なし",            "継承なし",           "次世代への蓄積"),
        ("AT中の選択",    "なし",                   "なし",               "なし",              "継承の儀（3択）"),
        ("長期目標",      "天井のみ",               "皇帝モード",          "PGG 1/16384",       "宿敵弱体化ゲージ"),
        ("感情設計",      "失敗前提ループ",          "瞬間爆発",            "夢を追う",           "積み上げ前提ループ"),
        ("来店継続動機",  "設定・天井期待値",        "設定・天井期待値",    "設定・期待値",       "物語の継続・世代の積み上げ"),
    ]
    ry = Inches(1.28)
    for j, row in enumerate(rows):
        cx = Inches(0.2)
        for i, (val, w) in enumerate(zip(row, col_w)):
            bg = RGBColor(0x10, 0x0E, 0x24) if j % 2 == 0 else RGBColor(0x14, 0x12, 0x28)
            if i == 4:
                bg = RGBColor(0x1C, 0x10, 0x00) if j % 2 == 0 else RGBColor(0x22, 0x14, 0x00)
            rect(s, cx, ry, w, Emu(295000), bg)
            col = C_ORG if i == 4 else (C_LTGRAY if i == 0 else C_GRAY)
            tb(s, cx + Emu(40000), ry + Emu(28000), w - Emu(60000), Emu(248000),
               val, 8, bold=(i == 4), color=col, wrap=False)
            cx += w
        ry += Emu(302000)

    # フッター強調
    rect(s, Inches(0.2), Inches(4.52), Inches(9.6), Emu(560000),
         RGBColor(0x14, 0x08, 0x1E))
    rect(s, Inches(0.2), Inches(4.52), Emu(60000), Emu(560000), C_ORG)
    tb(s, Inches(0.45), Inches(4.57), Inches(9.2), Emu(260000),
       "空席を埋める：「AT中に選択できる台」が業界に存在しない。継承の儀は他機種との最大の差別化点。",
       9, bold=True, color=C_ORG2, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.87), Inches(9.2), Emu(240000),
       "積み上げ型ループは業界初。失敗をポジティブに転換する設計で、RPG世代のリピーター獲得を狙う。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 10: スペック・ターゲット・業界初ポイント
# ══════════════════════════════════════════════════════════════
def s_spec(prs):
    """スペック・ターゲット・業界初3点"""
    s = new_slide(prs)
    hdr(s, "スペック・ターゲット・業界初ポイント")

    # スペック表（左）
    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.7), Inches(3.1),
           C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.32), Inches(0.91), Inches(4.5), Emu(300000),
       "基本スペック（ミドル）", 11, bold=True, color=C_GOLD, font=FONT_H)
    specs = [
        ("タイプ",     "スマスロ（L型）"),
        ("天井",       "600G（継承ポイントで最短400G）"),
        ("純増",       "4.0枚/G / 6.5枚/G / 8.0枚/G"),
        ("機械割",     "設定1：97.5%  /  設定6：109%"),
        ("MY目標",     "約3,500枚（ミドル）"),
        ("コイン単価",  "20円（3枚ベット / 1G＝60円）"),
        ("CV",         "0.22〜0.27"),
    ]
    sy = Inches(1.33)
    for label, val in specs:
        rect(s, Inches(0.3), sy, Inches(1.6), Emu(255000), RGBColor(0x1A, 0x10, 0x00))
        tb(s, Inches(0.35), sy + Emu(28000), Inches(1.5), Emu(200000),
           label, 8.5, bold=True, color=C_GOLD2, wrap=False)
        tb(s, Inches(1.98), sy + Emu(28000), Inches(2.8), Emu(200000),
           val, 8.5, color=C_CREAM, wrap=False)
        sy += Emu(263000)

    # ターゲット（右上）
    rect_b(s, Inches(5.1), Inches(0.85), Inches(4.7), Inches(1.45),
           C_CARD, C_STEEL, 1.3)
    tb(s, Inches(5.22), Inches(0.91), Inches(4.5), Emu(280000),
       "ターゲット", 10, bold=True, color=C_STEEL, font=FONT_H)
    tb(s, Inches(5.22), Inches(1.23), Inches(4.5), Emu(650000),
       "30〜50代男性　ロマサガ2・DQ・FF体験世代\n週1〜2来店リピーター　許容投資3,000〜8,000円/日\n感情ニーズ：「物語に没入」「積み上げを感じたい」",
       8.5, color=C_CREAM)

    # 業界初3点（右下）
    rect_b(s, Inches(5.1), Inches(2.40), Inches(4.7), Inches(1.55),
           RGBColor(0x14, 0x08, 0x00), C_ORG, 1.8)
    tb(s, Inches(5.22), Inches(2.46), Inches(4.5), Emu(280000),
       "業界初ポイント（3点）", 10, bold=True, color=C_ORG, font=FONT_H)
    firsts = [
        ("① 世代交代ループ演出",           "ATループ＝世代交代。負けが物語に変わる"),
        ("② 継承の儀（AT中3択・効果明確）", "戦略的選択。全3種コンプ＝決戦AT解放条件"),
        ("③ 宿敵弱体化ゲージ（1セッション）","長期遊技の目標が具体的に可視化される"),
    ]
    fy = Inches(2.82)
    for title, detail in firsts:
        tb(s, Inches(5.22), fy, Inches(4.5), Emu(240000),
           f"★ {title}",  8.5, bold=True, color=C_ORG2, wrap=False)
        fy += Emu(248000)

    # 下部：学習機種
    rect(s, Inches(0.2), Inches(4.05), Inches(9.6), Emu(30000), C_GOLD)
    rect(s, Inches(0.2), Inches(4.13), Inches(9.6), Emu(550000), RGBColor(0x0C, 0x06, 0x1E))
    tb(s, Inches(0.35), Inches(4.18), Inches(9.2), Emu(260000),
       "学習機種からの昇華：ミリゴZ-ZONE→継承の儀 / 吉宗1G連→英雄の残照 / 祟り神加護→継承積み上げ",
       9, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.35), Inches(4.50), Inches(9.2), Emu(250000),
       "既存メカニクスを「感情設計」として統合。「選ぶ楽しさ」「積み上げる喜び」「決着の感動」を1セッションで体験。",
       8.5, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 11: まとめ・この台が生む体験
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    """まとめ：この台が生む体験・感動の3軸・核心メッセージ"""
    s = new_slide(prs)
    hdr(s, "まとめ  ──  「世代 ―継承の炎―」が生む体験")

    cols_data = [
        (Inches(0.2),
         "積み上げることで感じる",
         "「今日の遊技が\n1つの物語」\n\n失敗しても\n継承ポイントが貯まる\n\n世代が重なるたびに\n演出が豊かになる\n\n前進感のある設計",
         C_ORG,  RGBColor(0x1C, 0x0A, 0x00)),
        (Inches(3.55),
         "選ぶことで感じる",
         "「この台は\n自分の台」\n\n継承の儀で\n選んだ内容が物語を変える\n\n同じ台でも\n2人と同じ体験をしない\n\n能動的没入感",
         C_PUR,  RGBColor(0x1A, 0x06, 0x2C)),
        (Inches(6.9),
         "決着で感じる",
         "「積み上げが\n報われた」\n\n決戦ATは\n「たまたま当たった」\nではなく\n「自分が育てた」感覚\n\nMY3,500枚の爆発が\n物語の完成",
         C_RED,  RGBColor(0x20, 0x04, 0x04)),
    ]
    for x, title, desc, col, fill in cols_data:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.28), fill, col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(330000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.38), Inches(2.8), Inches(2.45),
           desc, 9, color=C_CREAM)

    # 核心メッセージ
    rect(s, Inches(0.2), Inches(4.22), Inches(9.6), Emu(40000), C_GOLD)
    rect(s, Inches(0.2), Inches(4.30), Inches(9.6), Emu(780000), RGBColor(0x12, 0x08, 0x00))
    tb(s, Inches(0.35), Inches(4.35), Inches(9.2), Emu(270000),
       "業界初3点：① 世代交代ループ（積み上げ型）  ② 継承の儀（3択・効果明確）  ③ 宿敵弱体化ゲージ（1セッション完結）",
       9, bold=True, color=C_GOLD2)
    tb(s, Inches(0.35), Inches(4.68), Inches(9.2), Emu(380000),
       "ミリゴのZ-ZONE × 吉宗の1G連 × 祟り神の感情設計 × 世界観（不死の魔王イグナール）を統合。\n"
       "「1人では届かない。でも、意志を継ぐことで、必ず倒せる。」——この感情設計が他の全ての台にない「世代」の核心。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        ("タイトル・3本柱",         s_title),
        ("コアコンセプト",           s_concept),
        ("世界観",                   s_world),
        ("ゲームフロー全体図",       s_flow),
        ("通常時の設計",             s_normal),
        ("世代AT・継承の儀",         s_at),
        ("英雄連合",                 s_union),
        ("決戦AT「世代の決着」",     s_kessen),
        ("既存機種との差別化",       s_diff),
        ("スペック・ターゲット",     s_spec),
        ("まとめ",                   s_matome),
    ]

    print("=" * 60)
    print("  「世代 ―継承の炎―」企画提案書ジェネレーター v2")
    print("=" * 60)
    print()
    for i, (name, func) in enumerate(slides, 1):
        print(f"  {i:2d}/{len(slides)} {name}")
        func(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
