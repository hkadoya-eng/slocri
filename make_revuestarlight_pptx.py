"""
L 少女☆歌劇 レヴュースタァライト -The SLOT-  機種説明＋分析 統合資料 v1
（オーイズミ・2025年3月3日導入）
出力: proposals/機種分析/レヴュースタァライト/revuestarlight_guide_v1.pptx
テーマ: 深紺(C_BG) × 赤(C_RED=#CC2244) × 金(C_GOLD=#C8A840) × ピンク(C_PINK=#CC44AA)
構成: Part A（プレイヤー視点説明）6枚 + Part B（分析）3枚 = 計9枚
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "レヴュースタァライト", "revuestarlight_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深紺×赤×金×ピンク：舞台・華やか）────────────────────────
C_BG    = RGBColor(0x03, 0x03, 0x16)   # 深紺
C_CARD  = RGBColor(0x08, 0x06, 0x20)
C_CARD2 = RGBColor(0x10, 0x08, 0x28)
C_ROW   = RGBColor(0x0C, 0x08, 0x24)
C_RED   = RGBColor(0xCC, 0x22, 0x44)   # 赤（レヴューレッド）
C_RED2  = RGBColor(0xFF, 0x44, 0x66)   # 明るい赤
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)   # 金
C_GOLD2 = RGBColor(0xFF, 0xDD, 0x66)   # 明るい金
C_PINK  = RGBColor(0xCC, 0x44, 0xAA)   # ピンク
C_PINK2 = RGBColor(0xFF, 0x77, 0xDD)   # 明るいピンク
C_CYAN  = RGBColor(0x22, 0xAA, 0xCC)   # シアン（補助色）
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_PURP  = RGBColor(0x88, 0x33, 0xCC)
C_WHITE = RGBColor(0xEC, 0xE8, 0xF4)
C_CREAM = RGBColor(0xD0, 0xC0, 0xA8)
C_GRAY  = RGBColor(0x88, 0x88, 0xAA)
C_LTGRY = RGBColor(0x44, 0x44, 0x66)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景PNG生成（舞台の暗幕×スポットライト）────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (3, 3, 22))
    draw = ImageDraw.Draw(img)
    # 縦の光芒（スポットライト風）
    for i in range(0, w, 120):
        for dy in range(h):
            alpha = max(0, 1.0 - abs(i - w // 2) / (w * 0.6))
            v = int(6 * alpha * (1 - dy / h))
            draw.point((i, dy), fill=(v, 0, v + 4))
    # 上部と下部の微細グラデ
    for y in range(0, 50):
        t = (50 - y) / 50 * 0.4
        draw.line([(0, y), (w, y)], fill=(0, 0, int(12 * t)))
    for y in range(h - 60, h):
        t = (y - (h - 60)) / 60
        draw.line([(0, y), (w, y)], fill=(int(12 * t), 0, int(8 * t)))
    # 舞台床ライン（底部横線）
    for y in range(h - 30, h - 25):
        draw.line([(0, y), (w, y)], fill=(80, 30, 60))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


# ── テキストボックス ─────────────────────────────────────────────────────────
def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


# ── 矩形ヘルパー ─────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


# ── ヘッダー（スライド上部バー）─────────────────────────────────────────────
def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_RED)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_GOLD, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_RED)


# ── ネット解析注記 ────────────────────────────────────────────────────────────
def net_note(slide):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(180000),
       "※ネット解析情報より", 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ── 右矢印 ────────────────────────────────────────────────────────────────────
def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


# ── 下矢印 ────────────────────────────────────────────────────────────────────
def arrow_d(slide, cx, y, col=None):
    shp = slide.shapes.add_shape(20, cx - Emu(90000), y, Emu(180000), Emu(200000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_RED
    shp.line.fill.background()


# ── フッター（設計コメント＋補足説明）──────────────────────────────────────
def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.10)
    fh = Emu(380000)
    rect(slide, 0, fy, SLIDE_W, fh, RGBColor(0x06, 0x03, 0x1C))
    rect(slide, 0, fy, Emu(20000), fh, C_RED)
    tb(slide, Inches(0.22), fy + Emu(40000), Inches(6.0), Emu(160000),
       bold_text, 7.5, bold=True, color=C_RED2)
    if sub_text:
        tb(slide, Inches(0.22), fy + Emu(200000), Inches(9.2), Emu(160000),
           sub_text, 7, color=C_GRAY)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント  [Part A - 1/9]
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    # 左パネル背景（舞台の幕風）
    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x04, 0x02, 0x12))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_RED)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_PINK)

    tb(s, Inches(0.22), Inches(0.38), Inches(5.0), Emu(290000),
       "機種説明＋分析資料　Part A: プレイヤー視点   パチスロアワード2025ノミネート", 9,
       color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.82), Inches(5.1), Emu(800000),
       "L少女☆歌劇\nレヴュースタァライト", 26, bold=True, color=C_PINK2, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.62), Inches(5.1), Emu(290000),
       "── 貫通型AT × 舞台効果強化 × 2種の上位AT", 10, color=C_CREAM, font=FONT_H)

    tb(s, Inches(0.22), Inches(3.18), Inches(4.9), Emu(200000),
       "メーカー：オーイズミ（OIZUMI）　　導入：2025年3月3日", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.43), Inches(4.9), Emu(200000),
       "タイプ：A+AT（スマスロ）　設定：1〜6段階", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.68), Inches(4.9), Emu(200000),
       "AT純増：約2.2枚/G（上位AT時 約3.4枚/G）　天井：最大900G+α", 8.5, color=C_GRAY)
    tb(s, Inches(0.22), Inches(3.93), Inches(4.9), Emu(200000),
       "機械割：設定1→97.8%　設定6→110.0%", 8.5, color=C_GRAY)

    # 右：この台の3ポイント（ボックス間隔を縮めて下端はみ出し防止）
    kws = [
        (C_RED,   "① 業界初「貫通型AT」",
         "ボーナス種類によってATの性能が変化する新設計\nAT中にボーナスを引くほど純増・上乗せが強化される"),
        (C_GOLD,  "② 舞台効果で変わるATの強さ",
         "スポットライト/ロンド/トップスターの3種舞台効果\n「引いたボーナスの種類」がゲーム展開を大きく左右"),
        (C_PINK,  "③ 2種の上位AT・高ループ設計",
         "レヴューデュエット（上乗せ特化）\n星罪のレヴュー（継続率約90%・高ループ型）"),
    ]
    # ボックス高さ1.35inch × 3 + 隙間0.07inch × 2 = 4.19inch → 0.28 + 4.19 = 4.47inch ≤ 4.9 OK
    bx_h = Emu(1234000)   # ボックス高さ約1.35inch
    bx_gap = Emu(1298000)  # 間隔（高さ+隙間0.07inch）
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.28) + i * bx_gap
        rect_b(s, Inches(5.65), y0, Inches(4.1), bx_h, C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), bx_h, ac)
        tb(s, Inches(5.85), y0 + Emu(50000), Inches(3.8), Emu(270000),
           kw, 11, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(320000), Inches(3.8), Emu(850000),
           desc, 8.5, color=C_WHITE)

    footer(s,
           "設計コメント：「ボーナスを引いてATを育てる」貫通型ATは通常のG数上乗せ型とは一線を画す独自設計",
           "補足：原作「少女☆歌劇 レヴュースタァライト」のブシロード×キングレコード作品。舞台少女の演じ合う「レヴュー」をゲーム性に昇華")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図  [Part A - 2/9]
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→CZ→AT→上位ATへの全ルートを蛇行2段で可視化", "2/9")

    # ─── 上段フロー（通常→CZ→AT）───────────────────────────────────────
    tb(s, Inches(0.28), Inches(0.65), Inches(9.0), Emu(230000),
       "▶ 基本ルート（通常時からAT突入）", 8.5, bold=True, color=C_RED2)

    boxes_top = [
        (C_LTGRY, C_WHITE,  "通常時",
         "CZポイント蓄積\nレア役でCZ/BIG/\nAT直撃抽選\n天井:900G"),
        (C_RED,   C_RED2,   "チャレンジ\nレヴュー(CZ)",
         "18G/2パート構成\n成功期待度 約37%\nキャラ参加で変化"),
        (C_CARD2, C_GOLD,   "ボーナス\n(BIG/REG)",
         "BIG:AT抽選\nREG:CZpt獲得\n再生産モードも"),
        (RGBColor(0x14, 0x04, 0x20), C_PINK2, "AT\nレヴュースタァライト",
         "初期30G\n純増約2.2枚/G\nG数/セット管理"),
    ]
    bw, bh = Inches(1.90), Emu(1550000)
    gap = Inches(0.22)
    total = 4 * bw + 3 * gap
    sx = (SLIDE_W - total) / 2
    ty = Inches(0.90)
    cy_top = ty + bh // 2

    for i, (fill, bc, lbl, sub) in enumerate(boxes_top):
        bx0 = sx + i * (bw + gap)
        rect_b(s, bx0, ty, bw, bh, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), ty + Emu(80000),
           bw - Emu(80000), Emu(400000), lbl, 9.5, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), ty + Emu(490000),
           bw - Emu(60000), Emu(920000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw + Emu(10000), cy_top, col=C_RED)

    # ─── 下段フロー（AT内部→上位AT）───────────────────────────────────
    tb(s, Inches(0.28), Inches(2.90), Inches(9.0), Emu(230000),
       "▶ AT内部から上位ATへのルート（貫通型ATの特徴）", 8.5, bold=True, color=C_GOLD)

    by2 = Inches(3.12)
    bh2 = Emu(1600000)

    boxes_bot = [
        (RGBColor(0x14, 0x04, 0x20), C_PINK2, "AT中\nボーナス当選",
         "BIG/REG引くほど\nAT性能が上がる\n「貫通型」の核心"),
        (C_CARD2, C_RED2,   "クライマックス\nレヴュー",
         "G数上乗せ特化ゾーン\n成立役×背景色倍率\nで上乗せ算出"),
        (C_CARD2, C_GOLD,   "舞台効果\n発動",
         "スポットライト/\nロンド/トップスター\n3種でAT強化"),
        (RGBColor(0x10, 0x02, 0x20), C_GOLD2, "上位AT",
         "レヴューデュエット\n(上乗せ特化)\n星罪のレヴュー\n(90%ループ)"),
    ]
    for i, (fill, bc, lbl, sub) in enumerate(boxes_bot):
        bx0 = sx + i * (bw + gap)
        rect_b(s, bx0, by2, bw, bh2, fill, bc, 1.8)
        tb(s, bx0 + Emu(40000), by2 + Emu(80000),
           bw - Emu(80000), Emu(400000), lbl, 9.0, bold=True,
           color=bc, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx0 + Emu(30000), by2 + Emu(490000),
           bw - Emu(60000), Emu(980000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx0 + bw + Emu(10000), by2 + bh2 // 2, col=C_GOLD)

    # 段をつなぐ接続矢印（下向き：右端AT→下段AT中）
    arrow_d(s, sx + 3 * (bw + gap) + bw // 2,
            ty + bh + Emu(30000), col=C_PINK)

    footer(s,
           "設計コメント：通常時は「CZポイント蓄積→CZ→AT」。AT内は「BIG引くほど強化」の貫通型が他機種との決定的な差",
           "補足：上段=AT突入への旅路、下段=AT内部で更に強化されていく過程。天井900G+αは最長ルート")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方  [Part A - 3/9]
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── CZポイント蓄積・天井・AT直撃・3つのルート", "3/9")

    # 3カラム
    cols = [
        (C_RED,  "CZポイント蓄積",
         "レア役成立でCZポイントを獲得。\n10pt到達ごとにCZ\n「チャレンジ・レヴュー」の\n突入抽選が行われる。\n\n"
         "ポイントはモードに応じた\nテーブルで管理され、\n何G目にCZが発動するかの\n見通しが変わる。\n\n"
         "RBボーナスではCZポイントが\n直接加算される特典あり。",
         "pt積み上げが\nCZへの道"),
        (C_GOLD, "天井・AT直撃高確",
         "【天井】最大900G+α到達で\n前兆を経由してAT当選。\nリセット時は600G+αに短縮。\n\n"
         "【AT直撃高確】\n300G付近・600G付近で\n内部的に移行する高確率状態。\n"
         "移行後はレア役成立時に\n前兆を経由してAT直撃。\n\n"
         "高確中はボーナス・AT当選まで\n通常状態に転落しない。",
         "リセット後は\n600G天井"),
        (C_PINK, "再生産モード",
         "赤7 BIG成立後に移行する\n特殊な報酬ゾーン（20G）。\n\n"
         "消化中に成立役を参照して\n「再生産アイコン」を\n獲得できる。\n\n"
         "獲得アイコン数に応じた\nボーナスや上位状態が\n期待できる特殊ルート。\n\n"
         "通常ルート外の\nサプライズ体験として機能。",
         "赤7BIG後の\n隠しルート"),
    ]
    col_w = Inches(2.90)
    col_gap = Inches(0.20)
    col_y = Inches(0.72)
    col_h = Emu(3800000)

    for i, (ac, ch, cb, badge) in enumerate(cols):
        cx0 = Inches(0.28) + i * (col_w + col_gap)
        rect_b(s, cx0, col_y, col_w, col_h, C_CARD, ac, 1.8)
        rect(s, cx0, col_y, Emu(45000), col_h, ac)
        rect(s, cx0 + Emu(45000), col_y, col_w - Emu(45000), Emu(350000),
             RGBColor(0x0C, 0x06, 0x28))
        tb(s, cx0 + Emu(75000), col_y + Emu(60000),
           col_w - Emu(100000), Emu(270000),
           ch, 11, bold=True, color=ac, font=FONT_H)
        tb(s, cx0 + Emu(75000), col_y + Emu(380000),
           col_w - Emu(100000), col_h - Emu(620000),
           cb, 8, color=C_WHITE)
        rect_b(s, cx0 + col_w - Emu(820000), col_y + Emu(60000),
               Emu(790000), Emu(230000), C_CARD2, ac, 1.0)
        tb(s, cx0 + col_w - Emu(810000), col_y + Emu(70000),
           Emu(780000), Emu(210000), badge, 7, bold=True,
           color=ac, align=PP_ALIGN.CENTER)

    # 右端：攻略TIP
    rx = Inches(0.28) + 3 * (col_w + col_gap)
    rw = SLIDE_W - rx - Emu(200000)
    rect_b(s, rx, col_y, rw, col_h, RGBColor(0x06, 0x04, 0x1C), C_RED2, 1.5)
    tb(s, rx + Emu(60000), col_y + Emu(60000), rw - Emu(80000), Emu(260000),
       "攻略TIP", 10, bold=True, color=C_RED2, font=FONT_H)
    tips = [
        (C_RED2,  "リセット後は\n600G天井短縮\n朝一狙い有利"),
        (C_GOLD,  "300G/600G付近は\nAT直撃高確\n期待が高まる"),
        (C_PINK,  "AT終了後の\n引き戻しCZ\n期待度約75%"),
    ]
    tip_h = (col_h - Emu(320000)) // 3
    for j, (tc, tt) in enumerate(tips):
        ty0 = col_y + Emu(320000) + j * tip_h
        rect(s, rx + Emu(30000), ty0 + Emu(30000), Emu(15000),
             tip_h - Emu(60000), tc)
        tb(s, rx + Emu(80000), ty0 + Emu(60000),
           rw - Emu(110000), tip_h - Emu(80000), tt, 7.5, color=C_WHITE)

    footer(s,
           "設計コメント：CZポイント蓄積（基本）+ AT直撃高確（300/600G付近）+ 天井（900G）の3重安全網設計",
           "補足：再生産モード（赤7BIG後）はボーナス当選後の特別ルート。通常より多いアイコン獲得で大量上乗せに繋がる")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み  [Part A - 4/9]
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ「チャレンジ・レヴュー」の仕組み ── 舞台少女たちのレヴューと前兆演出", "4/9")

    # ─── 左半分：CZ基本構成 ───────────────────────────────────────────
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(330000), RGBColor(0x44, 0x10, 0x22))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(240000),
       "CZ基本仕様：18G消化・前半12G＋後半6G・2パート構成　成功期待度：約37%", 9,
       bold=True, color=C_RED2)

    # 2パート視覚表示
    ep_y = ly + Emu(390000)
    ep_h = Emu(440000)
    parts = [
        (C_RED,   "前半パート（12G）",
         "参加するキャラクター（舞台少女）によって\n演出の強さが変わる。\nヒカリ参加で期待度大幅アップ"),
        (C_GOLD,  "後半パート（5G）",
         "最終判定パート。\nここでのボーナス当選でAT濃厚。\n背景・演出でAT当選を告知"),
    ]
    pw = (lw - Emu(30000)) // 2 - Emu(15000)
    for k, (pc, pt, ps) in enumerate(parts):
        px = lx + k * (pw + Emu(30000))
        ph = Emu(1050000)
        rect_b(s, px, ep_y, pw, ph, C_CARD, pc, 1.5)
        rect(s, px, ep_y, pw, Emu(180000), RGBColor(0x18, 0x06, 0x12))
        rect(s, px, ep_y, Emu(18000), ph, pc)
        tb(s, px + Emu(40000), ep_y + Emu(30000), pw - Emu(50000), Emu(150000),
           pt, 8.5, bold=True, color=pc)
        tb(s, px + Emu(30000), ep_y + Emu(210000), pw - Emu(50000), ph - Emu(260000),
           ps, 8, color=C_WHITE)

    # CZポイントテーブル説明
    tb_y = ep_y + Emu(1100000)
    tb_h = Emu(1050000)
    rect_b(s, lx, tb_y, lw, tb_h, C_CARD, C_PINK, 1.5)
    rect(s, lx, tb_y, Emu(45000), tb_h, C_PINK)
    tb(s, lx + Emu(75000), tb_y + Emu(55000), lw - Emu(100000), Emu(260000),
       "CZポイントテーブルとCZ突入タイミング", 10, bold=True, color=C_PINK, font=FONT_H)
    tb(s, lx + Emu(75000), tb_y + Emu(315000), lw - Emu(100000), tb_h - Emu(370000),
       "・通常時はテーブルに沿ったG数でCZポイントを獲得\n"
       "・10pt到達 → CZ突入抽選（テーブルにより発動G数が変わる）\n"
       "・REGボーナス消化中もCZポイント獲得の抽選あり\n"
       "・AT終了後は引き戻しCZが発生（期待度約75%）",
       8, color=C_WHITE)

    # ─── 右半分：前兆演出とCZへのルート ────────────────────────────────
    rx, ry = Inches(5.05), Inches(0.72)
    rw = Inches(4.65)

    # 3ルート説明
    routes = [
        (C_RED,   "① CZポイント経由",
         "10pt蓄積 → CZ突入抽選\nテーブル次第でG数が変動\n最も頻繁に発生するルート"),
        (C_GOLD,  "② BIG/REGボーナス経由",
         "BIG消化中:BAR揃いでAT当選濃厚\nREG消化中:CZpt獲得抽選が行われる\nボーナスが次の起点になる"),
        (C_PINK,  "③ AT直撃高確経由",
         "300G/600G付近で内部移行\nレア役 → 前兆 → AT直撃\n転落なしで確実に引き込む"),
    ]
    rh_single = Emu(1320000)
    for ri, (rc, rt, rb) in enumerate(routes):
        ry0 = ry + ri * (rh_single + Emu(40000))
        rect_b(s, rx, ry0, rw, rh_single, C_CARD, rc, 1.5)
        rect(s, rx, ry0, Emu(45000), rh_single, rc)
        rect(s, rx + Emu(45000), ry0, rw - Emu(45000), Emu(310000),
             RGBColor(0x0C, 0x06, 0x22))
        tb(s, rx + Emu(75000), ry0 + Emu(55000), rw - Emu(100000), Emu(240000),
           rt, 10, bold=True, color=rc, font=FONT_H)
        tb(s, rx + Emu(75000), ry0 + Emu(350000), rw - Emu(100000),
           rh_single - Emu(400000), rb, 8.5, color=C_WHITE)

    footer(s,
           "設計コメント：「CZポイント積み上げ」→「CZ」→「ボーナス」→「AT」という4段ステップが通常時の緊張感を持続させる",
           "補足：前兆演出中はキャラクターの「レヴュー（演じ合い）」が発展し、成功か失敗かが18Gかけて判明する演出設計")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（出玉の伸ばし方）  [Part A - 5/9]
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「レヴュースタァライト」── 何をすれば出玉が伸びるか", "5/9")

    # 上段：AT基本構造
    tb(s, Inches(0.28), Inches(0.68), Inches(9.4), Emu(220000),
       "▶ AT基本仕様と出玉を伸ばす2つのアクション", 9, bold=True, color=C_RED2)

    # AT基本仕様バー
    rect_b(s, Inches(0.28), Inches(0.90), Inches(9.44), Emu(520000),
           C_CARD, C_RED, 1.5)
    rect(s, Inches(0.28), Inches(0.90), Emu(45000), Emu(520000), C_RED)
    tb(s, Inches(0.50), Inches(0.94), Inches(9.0), Emu(250000),
       "AT「レヴュースタァライト」：初期30G以上・純増約2.2枚/G（貫通時 約3.4枚/G）・G数＆セット数管理型", 9,
       bold=True, color=C_RED2, font=FONT_H)
    tb(s, Inches(0.50), Inches(1.22), Inches(9.0), Emu(170000),
       "消化中はレア役でG数上乗せ・セットストック・特化ゾーン（クライマックス・レヴュー）を抽選。"
       "G数カウンター100G到達後は上乗せが全てセットストックに変換される。", 8,
       color=C_WHITE)

    # 下段2カラム：アクション詳細
    cy2 = Inches(1.55)
    ch2 = Emu(2800000)
    cl_w = Inches(4.5)
    cr_w = Inches(4.65)

    # 左: ボーナスを引いてATを強化（貫通型）
    rect_b(s, Inches(0.28), cy2, cl_w, ch2, C_CARD, C_RED, 1.5)
    rect(s, Inches(0.28), cy2, Emu(45000), ch2, C_RED)
    tb(s, Inches(0.50), cy2 + Emu(50000), cl_w - Emu(100000), Emu(270000),
       "① AT中のボーナスでAT性能を強化（貫通型）", 10, bold=True, color=C_RED, font=FONT_H)

    bonus_effects = [
        (C_RED2,  "BIGボーナス当選",
         "AT中にBIGを引くとATの純増が上昇\n「貫通型AT」の核心。\n引けば引くほどAT性能が向上する"),
        (C_GOLD,  "REGボーナス当選",
         "CZポイント追加獲得のチャンス\nセットストック抽選にも期待\nAT継続の下支えとして機能"),
        (C_PINK,  "特殊ボーナス（赤7等）",
         "再生産モードや上位AT突入の起点\n通常よりも強い上乗せ・ストック獲得\n引き時の演出が大きく発展"),
    ]
    be_h = (ch2 - Emu(320000)) // 3
    for bi, (bc, bt, bb) in enumerate(bonus_effects):
        biy = cy2 + Emu(320000) + bi * be_h
        rect_b(s, Inches(0.40), biy + Emu(20000),
               cl_w - Emu(180000), be_h - Emu(40000), C_CARD2, bc, 1.0)
        rect(s, Inches(0.40), biy + Emu(20000), Emu(20000),
             be_h - Emu(40000), bc)
        tb(s, Inches(0.55), biy + Emu(50000),
           cl_w - Emu(240000), Emu(240000), bt, 8.5, bold=True, color=bc)
        tb(s, Inches(0.55), biy + Emu(290000),
           cl_w - Emu(240000), be_h - Emu(330000), bb, 7.5, color=C_WHITE)

    # 右: 特化ゾーン「クライマックス・レヴュー」
    rect_b(s, Inches(5.05), cy2, cr_w, ch2, C_CARD, C_GOLD, 1.5)
    rect(s, Inches(5.05), cy2, Emu(45000), ch2, C_GOLD)
    tb(s, Inches(5.27), cy2 + Emu(50000), cr_w - Emu(100000), Emu(270000),
       "② 特化ゾーン「クライマックス・レヴュー」", 10, bold=True, color=C_GOLD, font=FONT_H)

    cz_items = [
        (C_GOLD2, "発動条件",
         "AT中のレア役一部で突入する\nG数上乗せ特化ゾーン"),
        (C_RED2,  "上乗せ計算方式",
         "「成立役が決める基本G数」×\n「背景色に対応した倍率」\nで上乗せゲーム数が算出"),
        (C_PINK,  "背景色の倍率",
         "背景が赤/金になるほど\n倍率が上昇する演出設計\n「赤背景＝激アツ」の直感設計"),
        (C_CYAN,  "100Gカウンター到達",
         "G数カウンターが100Gに達すると\n以降のG数上乗せがすべて\nセットストックに変換"),
    ]
    ci_h = (ch2 - Emu(320000)) // 4
    for ci, (cc, ct, cb) in enumerate(cz_items):
        ciy = cy2 + Emu(320000) + ci * ci_h
        rect_b(s, Inches(5.17), ciy + Emu(20000),
               cr_w - Emu(180000), ci_h - Emu(40000), C_CARD2, cc, 1.0)
        rect(s, Inches(5.17), ciy + Emu(20000), Emu(20000),
             ci_h - Emu(40000), cc)
        tb(s, Inches(5.32), ciy + Emu(50000),
           cr_w - Emu(250000), Emu(240000), ct, 8.5, bold=True, color=cc)
        tb(s, Inches(5.32), ciy + Emu(290000),
           cr_w - Emu(250000), ci_h - Emu(330000), cb, 7.5, color=C_WHITE)

    footer(s,
           "設計コメント：「ボーナスを引いてATを強化（貫通型）」と「クライマックス・レヴューで上乗せ（特化ゾーン）」の二本柱で出玉を伸ばす",
           "補足：100G到達後のセットストック変換はAT継続セット数の積み増し設計。G数が消化されるほど枚数ではなく「継続」として蓄積される")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方  [Part A - 6/9]
# ══════════════════════════════════════════════════════════════
def s_upper_at(prs):
    s = new_slide(prs)
    hdr(s, "上位ATへの道と遊び方 ── レヴューデュエット・星罪のレヴュー・舞台効果", "6/9")

    # 舞台効果帯（上部）
    rect_b(s, Inches(0.28), Inches(0.68), Inches(9.44), Emu(500000),
           RGBColor(0x14, 0x08, 0x28), C_PINK, 1.5)
    rect(s, Inches(0.28), Inches(0.68), Emu(45000), Emu(500000), C_PINK)
    tb(s, Inches(0.50), Inches(0.72), Inches(9.0), Emu(240000),
       "舞台効果（3種）── AT中に発動しATの性能を段階的に強化する", 10,
       bold=True, color=C_PINK, font=FONT_H)

    stage_effects = [
        (C_RED2,  "スポットライト",
         "レア役で\n100%G数上乗せ\n常時発動"),
        (C_GOLD,  "ロンド",
         "セットストック発生時\n70%で特殊ループ\n継続重視"),
        (C_CYAN,  "トップスター",
         "ボーナスランクが\n改善・強化\n上位ATへの鍵"),
    ]
    sew = Inches(2.6)
    segap = Inches(0.50)
    setotal = 3 * sew + 2 * segap
    sesx = (SLIDE_W - setotal) / 2
    sey = Inches(1.02)
    seh = Emu(930000)
    for i, (sc, st, ss) in enumerate(stage_effects):
        sex = sesx + i * (sew + segap)
        rect_b(s, sex, sey, sew, seh, C_CARD, sc, 2.0)
        rect(s, sex, sey, sew, Emu(180000), RGBColor(0x0C, 0x06, 0x22))
        rect(s, sex, sey, Emu(30000), seh, sc)
        tb(s, sex + Emu(50000), sey + Emu(30000), sew - Emu(70000), Emu(150000),
           st, 10, bold=True, color=sc, align=PP_ALIGN.CENTER)
        tb(s, sex + Emu(30000), sey + Emu(210000), sew - Emu(60000), seh - Emu(260000),
           ss, 9, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 2つの上位AT詳細
    tb(s, Inches(0.28), Inches(2.25), Inches(9.4), Emu(220000),
       "▶ 2種の上位AT ── 目指す頂点は2つある", 9, bold=True, color=C_GOLD)

    cy2 = Inches(2.46)
    ch2 = Emu(2020000)
    cl_w = Inches(4.5)
    cr_w = Inches(4.65)

    # 左: レヴューデュエット
    rect_b(s, Inches(0.28), cy2, cl_w, ch2, C_CARD, C_RED, 2.0)
    rect(s, Inches(0.28), cy2, Emu(45000), ch2, C_RED)
    rect(s, Inches(0.28) + Emu(45000), cy2, cl_w - Emu(45000), Emu(340000),
         RGBColor(0x18, 0x06, 0x16))
    tb(s, Inches(0.50), cy2 + Emu(60000), cl_w - Emu(100000), Emu(270000),
       "上位AT①「レヴューデュエット」（上乗せ特化）", 11, bold=True, color=C_RED, font=FONT_H)
    tb(s, Inches(0.50), cy2 + Emu(380000), cl_w - Emu(100000), ch2 - Emu(420000),
       "■ 上乗せ発生率：約1/13（通常ATより格段にアップ）\n\n"
       "■ 突入契機：AT初当たり時の一部 / セットストック当選後の一部 /\n"
       "　　　　　　CZ成功時の一部 / BIG消化中の一部\n\n"
       "■ 特徴：「キラめき目」が必ず発生するとレア役扱いとなり\n"
       "　　　　さらなる上乗せが連発。出玉が最も速く増える状態\n\n"
       "■ AT中のボーナス→舞台効果→デュエット突入が理想ルート",
       8, color=C_WHITE)

    # 右: 星罪のレヴュー
    rect_b(s, Inches(5.05), cy2, cr_w, ch2, C_CARD, C_GOLD, 2.0)
    rect(s, Inches(5.05), cy2, Emu(45000), ch2, C_GOLD)
    rect(s, Inches(5.05) + Emu(45000), cy2, cr_w - Emu(45000), Emu(340000),
         RGBColor(0x16, 0x10, 0x04))
    tb(s, Inches(5.27), cy2 + Emu(60000), cr_w - Emu(100000), Emu(270000),
       "上位AT②「星罪のレヴュー」（高ループ型・最上位）", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(5.27), cy2 + Emu(380000), cr_w - Emu(100000), ch2 - Emu(420000),
       "■ 継続期待度：約90%（高ループ型の最上位AT）\n\n"
       "■ 突入契機：エピソード「されど舞台はつづく」終了後 /\n"
       "　　　　　　AT規定セット終了時の抽選 /\n"
       "　　　　　　レヴューデュエット中のBBの一部\n\n"
       "■ 特徴：突入前に保持しているストックが全て上位ATのループ\n"
       "　　　　ストックに変換される。ジャッジパートでのストック\n"
       "　　　　獲得率が大幅に上昇し、実質的な長期継続ATとなる",
       8, color=C_WHITE)

    footer(s,
           "設計コメント：舞台効果3種×上位AT2種の組み合わせがAT中の「今どの強さか」を体感させる多層構造",
           "補足：レヴューデュエット=短期爆発型、星罪のレヴュー=長期継続型と、性格の異なる上位ATを2つ用意した設計")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（舞台・演劇とパチスロの融合）  [Part B - 1/3]
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "Part B 分析 ── 面白さの設計（舞台演劇×パチスロの融合・貫通型ATの独自性）", "7/9")

    # 上段横断バー：設計の核心
    rect(s, Inches(0.28), Inches(0.72), Inches(9.44), Emu(280000),
         RGBColor(0x44, 0x10, 0x20))
    tb(s, Inches(0.50), Inches(0.74), Inches(9.0), Emu(250000),
       "核心：「舞台少女たちの演じ合い（レヴュー）」がゲーム進行そのものになっている世界観融合設計", 9,
       bold=True, color=C_RED2)

    col_l_x = Inches(0.28)
    col_l_w = Inches(4.5)
    col_r_x = Inches(5.05)
    col_r_w = Inches(4.65)
    col_y = Inches(1.10)
    col_h = Emu(3200000)

    # 左: 貫通型ATの設計的意義
    rect_b(s, col_l_x, col_y, col_l_w, col_h, C_CARD, C_RED, 1.5)
    rect(s, col_l_x, col_y, Emu(45000), col_h, C_RED)
    tb(s, col_l_x + Emu(75000), col_y + Emu(50000), col_l_w - Emu(100000), Emu(270000),
       "貫通型ATという設計的発明", 10, bold=True, color=C_RED, font=FONT_H)

    comp_items = [
        (C_LTGRY, "一般的なAT",
         "G数が減るだけ → 上乗せ待ちの受動的体験\nボーナスはAT前後にしか意味を持たない"),
        (C_RED,   "貫通型AT（本機）",
         "AT中にボーナスを引くほどATが強化される\nボーナス当選自体が「強化のイベント」になる"),
    ]
    comp_h = Emu(680000)
    comp_y = col_y + Emu(330000)
    for ci, (cc, ct, cx) in enumerate(comp_items):
        ciy = comp_y + ci * (comp_h + Emu(50000))
        rect_b(s, col_l_x + Emu(50000), ciy,
               col_l_w - Emu(100000), comp_h, C_CARD2, cc, 1.2)
        rect(s, col_l_x + Emu(50000), ciy, Emu(30000), comp_h, cc)
        tb(s, col_l_x + Emu(120000), ciy + Emu(60000),
           col_l_w - Emu(200000), Emu(250000), ct, 9, bold=True, color=cc)
        tb(s, col_l_x + Emu(120000), ciy + Emu(310000),
           col_l_w - Emu(200000), Emu(320000), cx, 8, color=C_WHITE)

    psy_y = comp_y + 2 * (comp_h + Emu(50000)) + Emu(80000)
    tb(s, col_l_x + Emu(60000), psy_y, col_l_w - Emu(100000), Emu(220000),
       "生まれる心理効果", 8.5, bold=True, color=C_RED2)
    psys = [
        "① AT中に「次のボーナスを期待する」能動的参加感",
        "② ボーナスが引けるたびに「AT進化」の達成感",
        "③ 「どの舞台効果が発動するか」のドキドキ感",
        "④ 上位ATへ到達した時の「積み上げ報酬」感",
    ]
    for pi, ps in enumerate(psys):
        tb(s, col_l_x + Emu(60000), psy_y + Emu(240000) + pi * Emu(290000),
           col_l_w - Emu(100000), Emu(260000), ps, 8, color=C_WHITE)

    # 右: 舞台演劇とパチスロの融合
    rect_b(s, col_r_x, col_y, col_r_w, col_h, C_CARD, C_PINK, 1.5)
    rect(s, col_r_x, col_y, Emu(45000), col_h, C_PINK)
    tb(s, col_r_x + Emu(75000), col_y + Emu(50000), col_r_w - Emu(100000), Emu(270000),
       "舞台演劇×パチスロの融合設計", 10, bold=True, color=C_PINK, font=FONT_H)

    stage_pts = [
        (C_PINK2,  "「レヴュー」＝戦い＝CZ成否",
         "原作の「レヴュー（演じ合い・戦い）」が\nCZ「チャレンジ・レヴュー」に直結している\nキャラが戦う=プレイヤーの勝負が同期する"),
        (C_GOLD,   "舞台効果＝AT進化の可視化",
         "「スポットライトを浴びる」「ロンド（輪舞）」\n「トップスター」という舞台用語がそのままAT\n強化要素の名前になっている世界観一致"),
        (C_RED2,   "キャラの出演＝CZ期待度変化",
         "どのキャラが「参加する（演じる）」かで\nCZの期待度が変わる。お気に入りキャラの\n登場が確率的な変化要素になっている"),
        (C_CYAN,   "「されど舞台はつづく」エンドレス",
         "原作テーマ「舞台は終わらない」が\n上位AT突入条件に使われている\n原作ファンへのファンサービスと機能設計の融合"),
    ]
    apt_h = (col_h - Emu(320000)) // 4
    for ai, (ac, at, ab) in enumerate(stage_pts):
        aiy = col_y + Emu(320000) + ai * apt_h
        rect(s, col_r_x + Emu(50000), aiy + Emu(30000),
             Emu(25000), apt_h - Emu(60000), ac)
        tb(s, col_r_x + Emu(110000), aiy + Emu(60000),
           col_r_w - Emu(160000), Emu(250000), at, 9, bold=True, color=ac)
        tb(s, col_r_x + Emu(110000), aiy + Emu(310000),
           col_r_w - Emu(160000), apt_h - Emu(360000), ab, 8, color=C_WHITE)

    footer(s,
           "設計コメント：貫通型ATは「AT中も能動的に楽しむ」新設計。舞台用語の機能化は原作ファンとパチスロ初心者の両方を取り込む",
           "補足：「舞台少女が演じ合う＝レヴュー」というコンセプトがゲームフロー全体に統一的に貫かれている意欲的な世界観設計")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題  [Part B - 2/3]
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 独自設計の革新性 vs ユーザー評価の二極化", "8/9")

    # 上段帯
    rect_b(s, Inches(0.28), Inches(0.72), Inches(9.44), Emu(540000),
           C_CARD, C_RED, 1.5)
    rect(s, Inches(0.28), Inches(0.72), Emu(45000), Emu(540000), C_RED)
    tb(s, Inches(0.50), Inches(0.76), Inches(9.0), Emu(260000),
       "評価の現実：ユーザー評価2.27点（低め） vs 継続率90%・多彩な上乗せルートという設計的革新", 9,
       bold=True, color=C_RED2, font=FONT_H)
    tb(s, Inches(0.50), Inches(1.06), Inches(9.0), Emu(220000),
       "「仕組み理解まで大変だが引けると楽しい」という声が示す通り、システム複雑性がリテラシーの壁になっている", 8,
       color=C_WHITE)

    cy2 = Inches(1.42)
    ch2 = Emu(2750000)
    cl_w = Inches(4.5)
    cr_w = Inches(4.65)

    # 左: 良い点
    rect_b(s, Inches(0.28), cy2, cl_w, ch2, C_CARD, C_GREEN, 1.5)
    rect(s, Inches(0.28), cy2, Emu(45000), ch2, C_GREEN)
    tb(s, Inches(0.50), cy2 + Emu(50000), cl_w - Emu(100000), Emu(270000),
       "良い点（革新性・継続性・爆発力）", 11, bold=True, color=C_GREEN, font=FONT_H)

    pros = [
        (C_GREEN, "業界初「貫通型AT」設計",
         "AT中のボーナム当選がATを強化する新発想。\n受動的なG数消化から能動的な強化体験へ。"),
        (C_RED2,  "引き戻しCZ 約75%と高継続性",
         "AT終了後も引き戻しCZが約75%の期待度で発動。\n星罪のレヴューの継続率約90%と合わせて長期稼働を後押し。"),
        (C_GOLD,  "2種の上位ATで2つの爆発ルート",
         "デュエット（高速上乗せ）と星罪（高ループ）という\n性格の異なる上位ATが2種存在し、到達後の体験が豊富。"),
        (C_PINK,  "舞台世界観との高い一体感",
         "舞台効果名・CZ名・上位AT突入条件すべてに\n原作「レヴュースタァライト」の世界観が貫かれている。"),
    ]
    pro_h = (ch2 - Emu(320000)) // 4
    for pi, (pc, pt, pb) in enumerate(pros):
        piy = cy2 + Emu(320000) + pi * pro_h
        rect_b(s, Inches(0.40), piy + Emu(20000),
               cl_w - Emu(180000), pro_h - Emu(40000), C_CARD2, pc, 1.0)
        rect(s, Inches(0.40), piy + Emu(20000), Emu(20000),
             pro_h - Emu(40000), pc)
        tb(s, Inches(0.55), piy + Emu(50000),
           cl_w - Emu(240000), Emu(240000), pt, 8.5, bold=True, color=pc)
        tb(s, Inches(0.55), piy + Emu(290000),
           cl_w - Emu(240000), pro_h - Emu(330000), pb, 7.5, color=C_WHITE)

    # 右: 課題
    rect_b(s, Inches(5.05), cy2, cr_w, ch2, C_CARD, C_RED, 1.5)
    rect(s, Inches(5.05), cy2, Emu(45000), ch2, C_RED)
    tb(s, Inches(5.27), cy2 + Emu(50000), cr_w - Emu(100000), Emu(270000),
       "課題（複雑性・機械割・リテラシー）", 11, bold=True, color=C_RED2, font=FONT_H)

    cons = [
        (C_RED2,  "システム複雑性・リテラシーの壁",
         "貫通型AT・舞台効果・上位AT2種の理解が\n必要で、初心者には仕組みが分かりづらい。"),
        (C_GOLD,  "AT純増が約2.2枚/Gと控えめ",
         "通常AT時の純増2.2枚/Gは中程度。\n貫通型を活かす前に終わると物足りない体験に。"),
        (C_GRAY,  "設定1の機械割 97.8%と低め",
         "低設定での稼働が多い現実では\nプレイヤーが「負ける体験」から離脱しやすい。"),
        (C_LTGRY, "上位AT「星罪のレヴュー」への到達難度",
         "最上位ATへの道筋が複数段階あり\n高設定以外ではなかなか到達できない。"),
    ]
    con_h = (ch2 - Emu(320000)) // 4
    for ci, (cc, ct, cb) in enumerate(cons):
        ciy = cy2 + Emu(320000) + ci * con_h
        rect_b(s, Inches(5.17), ciy + Emu(20000),
               cr_w - Emu(180000), con_h - Emu(40000), C_CARD2, cc, 1.0)
        rect(s, Inches(5.17), ciy + Emu(20000), Emu(20000),
             con_h - Emu(40000), cc)
        tb(s, Inches(5.32), ciy + Emu(50000),
           cr_w - Emu(250000), Emu(240000), ct, 8.5, bold=True, color=cc)
        tb(s, Inches(5.32), ciy + Emu(290000),
           cr_w - Emu(250000), con_h - Emu(330000), cb, 7.5, color=C_WHITE)

    footer(s,
           "設計コメント：「仕組みを理解すると楽しい」台の典型。チュートリアル演出・分かりやすい目標設定の追加で評価が大きく変わる余地あり",
           "補足：低評価の主因は「低設定での体験」と「複雑さによる挫折感」。原作ファン以外へのゲーム性の伝え方が課題")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること  [Part B - 3/3]
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── L少女☆歌劇 レヴュースタァライト 設計から学べること", "9/9")

    # 左: 3つの設計的学び
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(300000), RGBColor(0x44, 0x10, 0x20))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(225000),
       "設計から学べる3つのエッセンス", 10, bold=True, color=C_GOLD, font=FONT_H)

    elems = [
        (C_RED,   "① 貫通型ATという「強化参加型」設計",
         "AT中のボーナム当選がATを育てる設計は\n"
         "「受動的なG数消化」から「能動的な強化体験」への転換。\n"
         "プレイヤーを「育てる主体」にすることで没入感を高める。"),
        (C_GOLD,  "② 舞台効果名・システム名の世界観一致",
         "スポットライト/ロンド/トップスター/レヴューという\n"
         "舞台用語がそのまま機能名として使われる設計は\n"
         "原作IPとゲーム性の融合として模範的な実例。"),
        (C_PINK,  "③ 上位ATの性格分け（爆発型 vs ループ型）",
         "デュエット（速攻上乗せ）と星罪（長期継続）という\n"
         "2種の上位ATはプレイヤーに「2つの夢」を提供する。\n"
         "「短期で爆発か」「長く続くか」を選べない緊張感。"),
    ]
    for i, (ac, t, b) in enumerate(elems):
        ey = ly + Emu(300000) + i * Emu(1250000)
        rect_b(s, lx, ey, lw, Emu(1200000), C_CARD, ac, 1.5)
        rect(s, lx, ey, Emu(45000), Emu(1200000), ac)
        tb(s, lx + Emu(75000), ey + Emu(50000), lw - Emu(95000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), ey + Emu(305000), lw - Emu(95000), Emu(830000),
           b, 8, color=C_WHITE)

    # 右: 設計原則＋総括
    rx, ry = Inches(5.05), Inches(0.72)
    rw = Inches(4.65)

    rect(s, rx, ry, rw, Emu(280000), C_CARD2)
    tb(s, rx + Emu(50000), ry + Emu(45000), rw - Emu(70000), Emu(215000),
       "設計原則（他機種への応用）", 10, bold=True, color=C_GOLD, font=FONT_H)

    principles = [
        (C_RED,   "AT中もボーナスが「育成イベント」になる貫通型設計は継続動機を強化する"),
        (C_PINK,  "IPのコアワード（舞台/レヴュー）を機能名に転用するとIP理解がゲーム理解になる"),
        (C_GOLD,  "舞台効果3種＝「今のATの強さ」の可視化で現状把握を直感的にする"),
        (C_GREEN, "上位AT2種（速攻型/ループ型）は「どちらに行くか」の期待感の分岐を作る"),
        (C_RED2,  "複雑な仕組みこそ「最初の1時間に分かる化」の演出設計が評価を左右する"),
    ]
    for i, (ac, p) in enumerate(principles):
        py0 = ry + Emu(280000) + i * Emu(480000)
        rect(s, rx, py0, Emu(20000), Emu(440000), ac)
        tb(s, rx + Emu(50000), py0 + Emu(55000), rw - Emu(60000), Emu(360000),
           p, 8, bold=(i == 4), color=C_RED2 if i == 4 else C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(2700000), rw, Emu(1050000),
           RGBColor(0x06, 0x03, 0x1C), C_RED, 1.5)
    rect(s, rx, ry + Emu(2700000), Emu(45000), Emu(1050000), C_RED)
    tb(s, rx + Emu(75000), ry + Emu(2750000), rw - Emu(95000), Emu(270000),
       "総括", 10, bold=True, color=C_RED2, font=FONT_H)
    tb(s, rx + Emu(75000), ry + Emu(3020000), rw - Emu(95000), Emu(680000),
       "貫通型ATと舞台演劇世界観の融合は意欲的な革新設計。\n"
       "複雑性の「入口のハードル」を下げることができれば\n"
       "パチスロアワード2025ノミネートの期待に応えうる台。",
       8.5, color=C_WHITE)

    footer(s,
           "設計コメント：「ボーナスを引くほどATが育つ」という分かりやすい一言コンセプトに集約できれば訴求力は大きく変わる",
           "補足：原作「少女☆歌劇 レヴュースタァライト」のファン層＋パチスロファン層の両方を取り込める設計ポテンシャルを持つ")
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Part A: プレイヤー視点説明 6枚
    s_title(prs)      # 1: タイトル・スペック・この台の3ポイント
    s_flow(prs)       # 2: ゲームフロー全体図
    s_normal(prs)     # 3: 通常時の遊び方
    s_cz(prs)         # 4: CZ/前兆の仕組み
    s_at(prs)         # 5: AT/ボーナス（出玉を伸ばす）
    s_upper_at(prs)   # 6: 上位ATへの道と遊び方

    # Part B: 分析パート 3枚
    s_design(prs)     # 7: 面白さの設計（舞台×パチスロの融合）
    s_pros_cons(prs)  # 8: 良い点と課題
    s_matome(prs)     # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
