"""
Lパチスロ 炎炎ノ消防隊2 機種説明＋分析 統合版 PPTXジェネレーター v2
出力: proposals/機種分析/炎炎ノ消防隊2/enenno_guide_v2.pptx
テーマ: 深黒×炎赤×オレンジ×白（炎カラー）
情報源: なな徹・ちょんぼりすた・一撃 各解析ページ（2026年2月〜）
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "炎炎ノ消防隊2", "enenno_guide_v2.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深黒×炎赤×オレンジ×白）──────────────────────────
C_BG    = RGBColor(0x0A, 0x04, 0x04)
C_CARD  = RGBColor(0x14, 0x08, 0x04)
C_CARD2 = RGBColor(0x1C, 0x0C, 0x06)
C_ROW   = RGBColor(0x18, 0x0A, 0x05)
C_FIRE  = RGBColor(0xFF, 0x44, 0x00)   # 炎赤メイン #FF4400
C_FIRE2 = RGBColor(0xFF, 0x66, 0x11)   # 炎オレンジ
C_ORG   = RGBColor(0xFF, 0x88, 0x00)   # オレンジ #FF8800
C_FLAME = RGBColor(0xFF, 0xCC, 0x44)   # 炎先端（金）
C_WHITE = RGBColor(0xF0, 0xEC, 0xE8)
C_CREAM = RGBColor(0xF0, 0xD8, 0xB0)
C_GRAY  = RGBColor(0x99, 0x88, 0x80)
C_LTGRY = RGBColor(0x55, 0x44, 0x40)
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)
C_BLUE  = RGBColor(0x22, 0x77, 0xFF)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)
C_RED   = RGBColor(0xDD, 0x11, 0x11)
C_PINK  = RGBColor(0xFF, 0x44, 0x88)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景・ヘルパー群 ─────────────────────────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (10, 4, 4))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 80):
        draw.line([(i, 0), (0, i)], fill=(16, 6, 4), width=1)
    for y in range(h - 120, h):
        t = (y - (h - 120)) / 120
        r = int(50 * t)
        g = int(12 * t)
        draw.line([(0, y), (w, y)], fill=(r, g, 0))
    for y in range(0, 40):
        t = (40 - y) / 40 * 0.5
        draw.line([(0, y), (w, y)], fill=(int(15 * t), 0, 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_FIRE)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_FIRE)


def net_note(slide):
    tb(slide, Inches(7.8), Inches(5.38), Inches(2.1), Emu(180000),
       "※ネット解析情報より（なな徹・ちょんぼりすた・一撃）", 6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x0E, 0x06, 0x04))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_FIRE)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_ORG)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_FIRE
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_FIRE
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x06, 0x02, 0x02))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_FIRE)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_FIRE)

    tb(s, Inches(0.22), Inches(0.4), Inches(5.0), Emu(330000),
       "機種説明＋分析 統合ガイド  v2（解析情報更新版）", 10, color=C_FIRE2, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.88), Inches(5.1), Emu(900000),
       "炎炎ノ消防隊2", 34, bold=True, color=C_FIRE, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.62), Inches(5.0), Emu(280000),
       "Lパチスロ（スマスロ）── 十字目変換×ストック型ST×高純増5.8枚の設計", 9, color=C_CREAM, font=FONT_H)

    # スペック表（実測値） ※行間縮小・開始Y上げで画面下に収める
    specs = [
        ("メーカー",        "SANKYO　2026年2月2日導入"),
        ("設定",           "1〜6段階"),
        ("AT純増",         "約5.8枚/G（ボーナス中・AT中とも）"),
        ("設定1機械割",     "97.7%"),
        ("設定6機械割",     "114.9%"),
        ("ボーナス初当り①", "設定1：1/272 ／ 設定6：1/227"),
        ("炎炎ループ初当り", "設定1：1/684 ／ 設定6：1/486"),
        ("天井①",         "ボーナス間850G（SPエピソードBONUS保証）"),
        ("天井②",         "炎炎ループ間2,000G（SPエピソードBONUS保証）"),
        ("天井③",         "伝導者の罠5スルーでSPエピ確定"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(2.85) + i * Emu(195000)
        tb(s, Inches(0.22), ry, Inches(1.7), Emu(185000),
           k, 7.0, color=C_GRAY)
        tb(s, Inches(1.92), ry, Inches(3.3), Emu(185000),
           v, 7.0, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_FIRE,  "① 十字目変換フロー（核心）",
         "リプレイ小V→十字目変換→変換演出\n→伝導者決戦でボーナス告知。\n約1/100で発生する緊張感の連続。"),
        (C_ORG,   "② 高純増5.8枚/G×ループ設計",
         "炎炎大戦のループ率約80%で\n短時間に大量獲得が可能。\n紅J大戦なら期待約2,050枚。"),
        (C_CYAN,  "③ 3段階天井＋伝導者の罠スルー",
         "850G天井・2,000G天井・5スルー\n天井の多層セーフティで\n投資計画が立てやすい安心設計。"),
    ]
    # 右パネル3ボックス: y0=0.55+i*1.6, 高さ1.3 → 下端MAX=0.55+2×1.6+1.3=5.05 OK
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.55) + i * Emu(1461600)  # 1.6 inch = 1461600 EMU
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.3), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.3), ac)
        tb(s, Inches(5.85), y0 + Emu(65000), Inches(3.8), Emu(310000),
           kw, 11.5, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(380000), Inches(3.8), Emu(450000),
           desc, 8.5, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「十字目変換×伝導者決戦×ループ型AT」── 毎ゲーム緊張感と多層天井で安心感を両立",
           "純増5.8枚/Gはボーナス中・AT中ともに共通の高速出玉スペック")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（全ルートを蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→AT→上位ATの全ルート", "2/9")

    # 上段：通常→初当り→炎炎激闘 の基本フロー
    top_y = Inches(0.75)
    top_h = Emu(1100000)

    flow1 = [
        (C_CARD2, C_GRAY,  "通常遊技",         "規定G数/レア役\n/十字目変換で\n初当り抽選"),
        (C_CARD,  C_FIRE2, "REGボーナス\n(~85枚)", "消化後→\n伝導者の罠へ"),
        (C_CARD,  C_ORG,   "伝導者の罠",       "エピソードBONUS\n抽選。5スルーで\nSPエピ確定"),
        (C_CARD,  C_FIRE,  "SPエピソード\nBONUS(~200枚)", "直接\n炎炎激闘へ"),
        (C_CARD,  C_FLAME, "炎炎激闘\n(メインST)",  "1セット15G+α\nストック型"),
    ]
    # 5ボックス: margin=0.3, bw=1.65, gap=0.17 → 合計=0.3+5×1.65+4×0.17=9.23inch OK
    bw1 = Inches(1.65)
    gap1 = Inches(0.17)
    sx1 = Inches(0.30)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), top_y + Emu(70000), bw1 - Emu(60000), Emu(360000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), top_y + Emu(480000), bw1 - Emu(45000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), cy1)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1120000), Inches(3.0), Emu(260000),
       "天井①：ボーナス間850G → SPエピ確定", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1360000), Inches(3.5), Emu(5000), C_CYAN)
    tb(s, sx1 + Inches(3.6), top_y + Emu(1120000), Inches(3.5), Emu(260000),
       "天井③：伝導者の罠5スルー → SPエピ確定", 7.5, color=C_ORG)

    # 中段区切り線
    rect(s, 0, Inches(2.1), SLIDE_W, Emu(5000), RGBColor(0x44, 0x18, 0x08))

    # 下段：炎炎激闘内→上位AT昇格ルート（5ボックス）
    # margin=0.3, bw=1.65, gap=0.17 → 合計=0.3+5×1.65+4×0.17=9.23 OK
    bot_y = Inches(2.18)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,  C_FIRE2, "炎炎激闘\n（基本ST）",     "ボーナス期待度\n約57%/15G"),
        (C_CARD,  C_ORG,   "炎炎大戦\n（上位ST）",     "ループ率約80%\n純増5.8枚/G"),
        (C_CARD,  C_PINK,  "紅J大戦\n（特殊上位）",    "紅丸+J参戦\n期待約2,050枚"),
        (RGBColor(0x18,0x08,0x02), C_GOLD,
         "アドラバースト\n(穢レ無キ炎)",  "期待約2,760枚\n森羅万象経由"),
        (C_CARD,  C_CYAN,  "アドラリンク\nCZ(3G)",     "成功率約50%\n上乗せ契機"),
    ]
    bw2 = Inches(1.65)
    gap2 = Inches(0.17)
    sx2 = Inches(0.30)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8)
        tb(s, bx + Emu(35000), bot_y + Emu(70000), bw2 - Emu(60000), Emu(370000),
           lbl, 9.5, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), bot_y + Emu(490000), bw2 - Emu(45000), Emu(470000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw2 + Emu(10000), cy2)

    # 天井②アノテーション
    tb(s, sx2, bot_y + Emu(1100000), Inches(4.0), Emu(260000),
       "天井②：炎炎ループ間2,000G → SPエピ確定", 7.5, color=C_FLAME)
    rect(s, sx2, bot_y + Emu(1340000), Inches(4.5), Emu(5000), C_FLAME)

    net_note(s)
    footer(s, "上段=通常〜炎炎激闘突入ルート（REG経由と直撃の2本）、下段=AT内昇格ルート（炎炎大戦→紅J大戦→アドラバースト）",
           "伝導者の罠5スルー天井という独自の第三天井が投資上限の安心感をさらに強化する")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（天井含む全ルート）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 全AT突入ルート・天井管理・打ち方", "3/9")

    # 左：AT突入ルート図
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(290000), RGBColor(0x55, 0x18, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜炎炎激闘突入ルート（全3系統）", 10, bold=True, color=C_ORG)

    routes = [
        (C_FIRE2, "ルート①  SPエピソードBONUS直撃",
         "初当りがSPエピソードBONUSなら\n炎炎激闘に直行（最短ルート）。\n約200枚消化→そのまま激闘へ突入。"),
        (C_FIRE,  "ルート②  REGボーナス→伝導者の罠→炎炎激闘",
         "初当りがREGの場合は消化後に「伝導者の罠」へ。\n罠でエピソードBONUS当選→炎炎激闘突入。\n5スルーで次SPエピ確定（天井③）。"),
        (C_CYAN,  "ルート③  天井①（ボーナス間850G）",
         "ボーナス間850GでSPエピソードBONUS確定。\n最も頻出の天井到達ルート。\n設定変更後は短縮あり（要確認）。"),
        (C_FLAME, "ルート④  天井②（炎炎ループ間2,000G）",
         "炎炎激闘（ループ）間2,000Gで発動。\nSPエピソードBONUS確定。長期ハマりの出口。\n天井②到達時点は大きな期待値を持つ。"),
    ]
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(290000) + i * Emu(1125000)
        rect_b(s, lx, iy, lw, Emu(1060000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1060000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(320000), lw - Emu(100000), Emu(670000),
           b, 7.5, color=C_WHITE)

    # 右：チャンス役と確率・打ち方
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(290000), RGBColor(0x55, 0x18, 0x00))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "チャンス役の確率・期待度（全設定共通）", 10, bold=True, color=C_ORG)

    chance = [
        (C_GRAY,  "リプレイ",    "約1/8.6",  "小V停止でAT中に十字目変換の起点に"),
        (C_FIRE2, "弱チェリー",  "約1/80",   "十字目変換/前兆移行の抽選対象"),
        (C_FIRE,  "スイカ",      "約1/128",  "十字目変換・ボーナス高確の抽選対象"),
        (C_ORG,   "チャンス目",  "約1/128",  "十字目変換の抽選対象役"),
        (C_FLAME, "十字リプレイ","約1/5000", "アドラリンク当選濃厚・状態問わず"),
        (C_CYAN,  "レア役合算",  "約1/35",   "チェリー・スイカ・チャンス目の合算"),
    ]
    ch_h = Emu(690000)
    for i, (ac, role, prob, desc) in enumerate(chance):
        cy = ry + Emu(290000) + i * ch_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, ch_h, bg)
        rect(s, rx, cy, Emu(35000), ch_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(55000), Inches(1.0), Emu(260000),
           role, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.0), cy + Emu(60000), Inches(0.8), Emu(240000),
           prob, 8, bold=True, color=C_FLAME, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(310000), rw - Inches(0.6), Emu(320000),
           desc, 7.5, color=C_WHITE)

    # 打ち方メモ
    rect_b(s, rx, ry + Emu(4440000), rw, Emu(550000), C_CARD2, C_FIRE, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4490000), rw - Emu(80000), Emu(200000),
       "打ち方：順押しBAR狙い（左リール上段にBAR）", 8.5, bold=True, color=C_FIRE)
    tb(s, rx + Emu(60000), ry + Emu(4700000), rw - Emu(80000), Emu(280000),
       "手順を守らないとペナルティあり。小V停止を見逃さないこと。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "通常時の基本戦略：天井①(850G)を基準に管理し、伝導者の罠スルー回数（5スルー天井）もカウントする",
           "天井②(2,000G)は長期ハマり台の最終出口として機能。十字リプレイは状態問わず激アツ")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: AT「炎炎激闘」の核心：十字目変換フロー
# ══════════════════════════════════════════════════════════════
def s_at_flow(prs):
    s = new_slide(prs)
    hdr(s, "AT「炎炎激闘」の核心 ── 十字目変換フロー詳細（解析値）", "4/9")

    fc_x = Inches(0.28)
    fc_w = Inches(4.4)

    rect(s, fc_x, Inches(0.72), fc_w, Emu(260000), RGBColor(0x55, 0x18, 0x00))
    tb(s, fc_x + Emu(60000), Inches(0.74), fc_w - Emu(80000), Emu(240000),
       "十字目変換フロー（AT1セット15G+α）", 9.5, bold=True, color=C_ORG)

    # STEP 1
    n1_y = Inches(1.06)
    n1_h = Emu(560000)
    rect_b(s, fc_x, n1_y, fc_w, n1_h, C_CARD2, C_GRAY, 1.2)
    tb(s, fc_x + Emu(40000), n1_y + Emu(45000), fc_w - Emu(60000), Emu(200000),
       "STEP 1", 7, bold=True, color=C_GRAY)
    tb(s, fc_x + Emu(40000), n1_y + Emu(240000), fc_w - Emu(60000), Emu(260000),
       "レア役 or リプレイ小V 成立", 10.5, bold=True, color=C_WHITE, font=FONT_H)

    arrow_d(s, fc_x + fc_w // 2, n1_y + n1_h + Emu(15000), C_FIRE)

    # STEP 2
    n2_y = n1_y + n1_h + Emu(215000)
    n2_h = Emu(560000)
    rect_b(s, fc_x, n2_y, fc_w, n2_h, RGBColor(0x20, 0x08, 0x02), C_FIRE, 2.0)
    rect(s, fc_x, n2_y, Emu(35000), n2_h, C_FIRE)
    tb(s, fc_x + Emu(60000), n2_y + Emu(45000), fc_w - Emu(80000), Emu(200000),
       "STEP 2", 7, bold=True, color=C_FIRE)
    tb(s, fc_x + Emu(60000), n2_y + Emu(240000), fc_w - Emu(80000), Emu(260000),
       "十字目変換発生！（約1/100）", 10.5, bold=True, color=C_FIRE2, font=FONT_H)

    arrow_d(s, fc_x + fc_w // 2, n2_y + n2_h + Emu(15000), C_FIRE)

    # STEP 3: 変換演出の種類と期待度
    n3_y = n2_y + n2_h + Emu(215000)
    n3_h = Emu(840000)
    rect_b(s, fc_x, n3_y, fc_w, n3_h, RGBColor(0x22, 0x0A, 0x04), C_ORG, 2.0)
    rect(s, fc_x, n3_y, Emu(35000), n3_h, C_ORG)
    tb(s, fc_x + Emu(60000), n3_y + Emu(50000), fc_w - Emu(80000), Emu(230000),
       "STEP 3  変換演出の種類で期待度確認", 9.5, bold=True, color=C_ORG, font=FONT_H)

    # 変換演出4パターン（実測値）
    conv_infos = [
        (C_GRAY,  "十字マーク\n変換",  "低期待度"),
        (C_CYAN,  "アイリス\n変換",   "約50%+\n(レア役時のみ)"),
        (C_FIRE,  "シンラ\n変換",    "約84%"),
        (C_FLAME, "紅丸\n変換",     "約98%"),
    ]
    cw3 = (fc_w - Emu(80000)) // 4
    for ci, (cc, clbl, cpct) in enumerate(conv_infos):
        cx3 = fc_x + Emu(50000) + ci * cw3
        rect_b(s, cx3 + Emu(8000), n3_y + Emu(290000),
               cw3 - Emu(16000), Emu(490000), C_CARD, cc, 1.5)
        tb(s, cx3 + Emu(15000), n3_y + Emu(320000), cw3 - Emu(25000), Emu(240000),
           clbl, 8, bold=True, color=cc, align=PP_ALIGN.CENTER)
        tb(s, cx3 + Emu(15000), n3_y + Emu(560000), cw3 - Emu(25000), Emu(220000),
           cpct, 7.5, bold=True, color=C_FLAME, align=PP_ALIGN.CENTER)

    arrow_d(s, fc_x + fc_w // 2, n3_y + n3_h + Emu(15000), C_FLAME)

    # STEP 4
    n4_y = n3_y + n3_h + Emu(210000)
    n4_h = Emu(560000)
    rect_b(s, fc_x, n4_y, fc_w, n4_h, RGBColor(0x24, 0x10, 0x00), C_FLAME, 2.0)
    rect(s, fc_x, n4_y, Emu(35000), n4_h, C_FLAME)
    tb(s, fc_x + Emu(60000), n4_y + Emu(45000), fc_w - Emu(80000), Emu(200000),
       "STEP 4  変換成功！", 7, bold=True, color=C_FLAME)
    tb(s, fc_x + Emu(60000), n4_y + Emu(240000), fc_w - Emu(80000), Emu(260000),
       "伝導者決戦 → ボーナス当選の告知演出", 10, bold=True, color=C_GOLD, font=FONT_H)

    # 右パネル：保険設計と発生頻度詳細
    rx, ry = Inches(4.9), Inches(0.72)
    rw = Inches(4.85)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x55, 0x18, 0x00))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "抽選の仕組み・保険設計・発生頻度", 10, bold=True, color=C_ORG)

    insurance = [
        (C_FIRE,  "抽選の3段階フロー（解析値）",
         "① 十字目変換発生の抽選（内部状態参照）\n"
         "② 十字目高確への移行抽選\n"
         "③ 十字目ランクの昇格抽選\n\n"
         "この3段階の抽選が順番に実行され、\n"
         "ランクが高いほど変換演出が豪華になる。\n"
         "トータル発生率：約1/100。"),
        (C_CYAN,  "内部状態と高確の概念",
         "通常時は「ボーナス高確」「十字目高確」の\n"
         "2つの内部状態が存在。どちらが高いかで\n"
         "変換発生率・期待度が変化する。\n\n"
         "伝導者の罠中は十字目高確の概念なし。\n"
         "常に同じ当選率で抽選される。"),
        (C_ORG,   "ストック型継続の仕組み",
         "炎炎激闘はボーナスストックで継続管理。\n"
         "ストックが残る限り15G再セット。\n\n"
         "炎炎（ストック）ボーナス当選→ストック獲得\n"
         "→セット継続のサイクルが出玉の本体。\n"
         "ストック数が\"残弾数\"として機能。"),
    ]
    for i, (ac, t, b) in enumerate(insurance):
        iy = ry + Emu(260000) + i * Emu(1335000)
        rect_b(s, rx, iy, rw, Emu(1265000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(1265000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(830000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "十字目変換の設計核心：3段階内部抽選＋キャラ変換演出で「期待度の可視化」と「毎ゲームの緊張感」を両立",
           "紅丸変換（約98%）は実質告知演出。シンラ変換（約84%）でも十分な激アツ期待度")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 出玉を伸ばす方法（アドラリンク×炎炎大戦×上位演出）
# ══════════════════════════════════════════════════════════════
def s_extend(prs):
    s = new_slide(prs)
    hdr(s, "出玉を伸ばす方法 ── アドラリンク×伝導者決戦×上位昇格", "5/9")

    # 左：アドラリンク（上乗せCZ）
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x18, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "アドラリンク（上乗せ専用CZ・3G）", 10, bold=True, color=C_ORG)

    rect_b(s, lx, ly + Emu(260000), lw, Emu(1400000), C_CARD, C_CYAN, 2.0)
    rect(s, lx, ly + Emu(260000), Emu(40000), Emu(1400000), C_CYAN)
    tb(s, lx + Emu(70000), ly + Emu(310000), lw - Emu(90000), Emu(260000),
       "アドラリンクの仕組み（解析値）", 10, bold=True, color=C_CYAN, font=FONT_H)
    tb(s, lx + Emu(70000), ly + Emu(580000), lw - Emu(90000), Emu(1000000),
       "前兆中に十字目停止で発生する3GのCZ。\n"
       "成功期待度：約50%。\n\n"
       "3G間に小Vリプレイ or レア役成立→ボーナス当選。\n"
       "十字リプレイ経由で発生時→紅J大戦昇格濃厚！\n\n"
       "通常時突入率：約1/450\n"
       "伝導者の罠中突入率：約1/550",
       8.5, color=C_WHITE)

    # アドラリンク発動のリールロック（段数演出）
    rect(s, lx, ly + Emu(1660000), lw, Emu(260000), RGBColor(0x44, 0x18, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(1710000), lw - Emu(80000), Emu(200000),
       "リールロック段数で期待度を視覚的に表示", 9, bold=True, color=C_ORG)

    locks = [
        (C_GRAY,  "1段ロック", "チャンス",      "約30%",  0.30),
        (C_FIRE2, "2段ロック", "期待大",        "約60%",  0.60),
        (C_FLAME, "3段ロック", "当確に近い",    "約90%+", 0.92),
    ]
    lock_y = ly + Emu(1920000)
    lock_h = Emu(760000)
    lock_w = lw / 3

    for i, (ac, lbl, desc, pct_s, pct) in enumerate(locks):
        lkx = lx + i * lock_w
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect_b(s, lkx + Emu(8000), lock_y, lock_w - Emu(16000), lock_h, bg, ac, 1.5)
        tb(s, lkx + Emu(15000), lock_y + Emu(50000), lock_w - Emu(25000), Emu(240000),
           lbl, 8.5, bold=True, color=ac, align=PP_ALIGN.CENTER, wrap=False)
        tb(s, lkx + Emu(15000), lock_y + Emu(290000), lock_w - Emu(25000), Emu(200000),
           desc, 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)
        bar_w = lock_w - Emu(80000)
        rect(s, lkx + Emu(40000), lock_y + Emu(500000), bar_w, Emu(80000), C_LTGRY)
        rect(s, lkx + Emu(40000), lock_y + Emu(500000), int(bar_w * pct), Emu(80000), ac)
        tb(s, lkx + Emu(40000), lock_y + Emu(600000), lock_w - Emu(50000), Emu(180000),
           pct_s, 8, bold=True, color=ac, align=PP_ALIGN.CENTER, wrap=False)

    # 右：伝導者決戦と出玉昇格
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), RGBColor(0x55, 0x18, 0x00))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "伝導者決戦と各ボーナスの出玉役割", 10, bold=True, color=C_ORG)

    densha = [
        (C_FIRE,  "伝導者決戦（ボーナス告知バトル演出）",
         "十字目変換成功後に発生するバトル演出。\n"
         "勝利でボーナス当選が告知される。\n"
         "演出豪華さ・BGM・キャラ登場で期待度が変化。\n"
         "シンラ/炎柱絡み演出は高期待度。"),
        (C_FLAME, "ボーナス種類と出玉・役割",
         "REGボーナス（~85枚）：設定判別チャンス\n"
         "   →消化後に伝導者の罠へ\n"
         "SPエピソードBONUS（~200枚）：炎炎激闘直行\n"
         "炎炎ボーナス（~200枚）：激闘内ストック獲得\n"
         "炎炎BSTボーナス：ループ率90%・期待約550枚\n"
         "紅J大戦直撃：期待約2,050枚"),
        (C_ORG,   "森羅万象CZ → アドラバースト（最強）",
         "ボーナス消化中に突入する隠しCZ。\n"
         "成功でアドラバースト（穢レ無キ炎）確定。\n"
         "アドラチャレンジ成功でも到達可能。\n"
         "期待獲得枚数：約2,760枚（シリーズ最強）。"),
    ]
    for i, (ac, t, b) in enumerate(densha):
        iy = ry + Emu(260000) + i * Emu(1345000)
        rect_b(s, rx, iy, rw, Emu(1275000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(1275000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(260000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(320000), rw - Emu(90000), Emu(850000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "出玉を伸ばす鍵：アドラリンクの3段ロック演出＋十字リプレイ経由での紅J大戦昇格狙いが最大の上振れルート",
           "森羅万象CZ（ボーナス中）→アドラバーストは約2,760枚の最終到達点。終了画面まで要注目")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位モード（到達ルート＋遊び方）
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位モード ── 炎炎大戦・紅J大戦・アドラバーストへの到達と遊び方", "6/9")

    # 上段：昇格ルートフロー
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(260000), RGBColor(0x44, 0x18, 0x00))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(210000),
       "上位モード到達ルート（炎炎激闘からの昇格階層）", 9, bold=True, color=C_ORG)

    route_boxes = [
        (C_FIRE2, "炎炎激闘\n（基本ST）",     "ボーナス期待度\n約57%/15G\nスタート地点"),
        (C_ORG,   "炎炎大戦\n（上位ST）",     "ループ率約80%\n純増5.8枚/G\n主軸出玉ゾーン"),
        (C_PINK,  "紅J大戦\n（特殊上位ST）",  "紅丸+J参戦\n期待約2,050枚\nアドラリンク経由"),
        (C_GOLD,  "アドラバースト\n(穢レ無キ炎)", "森羅万象CZ\n期待約2,760枚\n最強フィナーレ"),
    ]
    bw_r = Inches(2.1)
    gap_r = Inches(0.26)
    sx_r = Inches(0.35)
    cy_r = Inches(1.42)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x18, 0x08, 0x00), ac, 2.0 if i >= 2 else 1.5)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(55000), Emu(370000), lbl, 10, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(30000), cy_r - bh_r // 2 + Emu(460000),
           bw_r - Emu(45000), Emu(480000), sub, 7.5,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw_r + Emu(12000), cy_r)

    # 中区切り
    rect(s, 0, Inches(2.06), SLIDE_W, Emu(5000), RGBColor(0x44, 0x18, 0x08))

    # 下段左：炎炎大戦の遊び方
    lx2, ly2 = Inches(0.28), Inches(2.12)
    lw2 = Inches(4.55)
    lh2 = Emu(2620000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_ORG, 1.8)
    rect(s, lx2, ly2, Emu(45000), lh2, C_ORG)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "炎炎大戦・紅J大戦の遊び方", 11, bold=True, color=C_ORG, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【炎炎大戦（上位ST）】\n"
       "炎炎激闘の一部で突入。純増5.8枚/Gは変わらず\n"
       "ボーナス期待度が約80%に跳ね上がる。\n"
       "15G+α消化後も高確率でループ継続。\n"
       "ボーナス当選時は通常時を経由せず炎炎大戦に復帰。\n\n"
       "【昇格条件】\n"
       "小Vリプレイ・レア役で炎炎大戦への昇格を抽選。\n"
       "昇格後はさらに紅J大戦への昇格も抽選される。\n\n"
       "【紅J大戦（特殊上位ST）】\n"
       "紅丸とジョーカーが参戦する特殊バージョン。\n"
       "期待獲得枚数は約2,050枚。\n"
       "アドラリンクで十字リプレイ経由→昇格濃厚。",
       8, color=C_WHITE)

    # 下段右：アドラバースト
    rx2, ry2 = Inches(5.05), Inches(2.12)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2,
           RGBColor(0x18, 0x08, 0x00), C_GOLD, 2.0)
    rect(s, rx2, ry2, Emu(45000), lh2, C_GOLD)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "アドラバースト（穢レ無キ炎）の遊び方", 11, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【到達ルート】\n"
       "① ボーナス消化中の森羅万象CZ成功\n"
       "② アドラチャレンジ（別途CZ）成功\n\n"
       "【期待枚数約2,760枚の内訳】\n"
       "ボーナス消化中の獲得枚数＋炎炎大戦ループ\n"
       "による出玉の合算。JACゲームとの組み合わせで\n"
       "大量獲得が実現する構造。\n\n"
       "【打ち方・注意点】\n"
       "アドラバースト突入時は演出を最後まで確認。\n"
       "終了画面に設定示唆が表示される場合がある。\n"
       "消化中は離席厳禁。JACゲーム中の役把握も重要。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "上位モードの階層：炎炎激闘→炎炎大戦（ループ率80%）→紅J大戦（2,050枚）→アドラバースト（2,760枚）",
           "アドラリンクで十字リプレイを引けた場合は紅J大戦昇格濃厚。これが最大の出玉上振れルート")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（なぜこの台は面白いのか）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── なぜ炎炎ノ消防隊2は面白いのか", "7/9")

    principles = [
        (C_FIRE,  "① 毎ゲーム「起きるかも」の緊張感",
         "十字目変換は約1/100と頻繁に発生。\n"
         "リプレイを引くたびに「変換するか？」の\n"
         "ドキドキが生まれる。スロットを回す行為\n"
         "そのものに意味を持たせる設計。"),
        (C_FIRE2, "② キャラ変換演出という直感的期待度表示",
         "十字マーク→アイリス→シンラ→紅丸と\n"
         "キャラの格で期待度が直感でわかる。\n"
         "原作ファンには感情移入を促す二重設計。\n"
         "初心者もベテランも楽しめるUX。"),
        (C_ORG,   "③ ストック型継続という「残弾数」の安心感",
         "ストックが残る限りセットが続く設計。\n"
         "「あと何セット確定している」という\n"
         "残弾数感覚が離席を防ぎ継続プレイを促す。\n"
         "15G再セットも含めたストレスフリー設計。"),
        (C_CYAN,  "④ アドラリンクによる自力感の演出",
         "前兆中に割り込む3GのCZ。打ち手が\n"
         "「自分でボーナスを引いた」と感じる仕組み。\n"
         "3段ロックの視覚演出が期待感を段階的に\n"
         "高め、没入感とリピート意欲を生む。"),
        (C_FLAME, "⑤ 常に「上のレベル」が見える多層構造",
         "炎炎激闘→炎炎大戦→紅J大戦→アドラバースト\n"
         "と常に上位状態が存在し続ける。\n"
         "\"次のステージ\"への目標が途切れず\n"
         "プレイヤーの離席抑制につながる設計。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1200000)
    gy = Inches(0.10)

    positions = [
        (Inches(0.28),  Inches(0.72)),
        (Inches(5.17),  Inches(0.72)),
        (Inches(0.28),  Inches(0.72) + bh_p + gy),
        (Inches(5.17),  Inches(0.72) + bh_p + gy),
        (Inches(0.28),  Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1100000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(260000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(330000), pw - Emu(90000), ph - Emu(390000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "面白さの核心：「1/100の十字目変換×キャラ期待度演出×ストック残弾感×自力感×多層目標」の5要素連鎖",
           "炎（原作モチーフ）をキャラ変換演出のランク指標に昇華した点が本機の設計的美しさの真骨頂")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN, "高純増5.8枚/G×ボーナス中も共通",
         "ボーナス消化中・AT中ともに純増5.8枚で\n"
         "出玉速度が統一されている。ループ中の\n"
         "体感出玉スピードが圧倒的。"),
        (C_GREEN, "キャラ変換演出による直感的期待度設計",
         "アイリス・シンラ・紅丸という原作キャラの\n"
         "格が期待度順。原作ファン以外も\n"
         "画面を見ただけで期待感を直感できる。"),
        (C_GREEN, "3段階天井＋スルー天井の多層セーフティ",
         "850G天井・2,000G天井・5スルー天井で\n"
         "ハイエナ立ち回りも計画しやすい。\n"
         "投資上限が明確で安心感が高い。"),
        (C_GREEN, "ストック型×ループ型の複合AT設計",
         "炎炎激闘（ストック管理）と炎炎大戦\n"
         "（ループ型）の2層構造が出玉の\n"
         "予測感と意外性を同時に提供する。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1140000)
        rect_b(s, lx, iy, lw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(300000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,   "天井②（2,000G）の投資負担の大きさ",
         "炎炎ループ間2,000Gは非常に長い設定。\n"
         "天井②を狙う際は大きな投資が必要で\n"
         "ライトユーザーには敷居が高い。"),
        (C_FIRE,  "伝導者の罠スルー天井の複雑さ",
         "5スルー天井という独自カウント概念が\n"
         "初心者には理解しにくい。\n"
         "カウント方法を知らないと損をする場面も。"),
        (C_ORG,   "設定判別難易度の高さ",
         "REGボーナスのシナリオが設定差の主軸だが\n"
         "引けなければ判別不能に近い。\n"
         "高設定確信まで時間とゲーム数が必要。"),
        (C_GRAY,  "炎炎激闘序盤の単調感",
         "ストック数が少ない序盤は15G×数セットの\n"
         "繰り返しで単純作業感が出やすい。\n"
         "炎炎大戦に乗るまでの我慢が必要。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1140000)
        rect_b(s, rx, iy, rw, Emu(1070000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1070000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(250000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(300000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "強みと課題の両面把握が設計学習の本質：強みを他機種に応用し、課題を次世代設計で克服する視点を持つ",
           "伝導者の罠スルー天井の複雑さはゲームセンター的な深みでもあり、ヘビーユーザーには魅力になり得る")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x18, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "炎炎ノ消防隊2 ── 設計的強み総括", 10, bold=True, color=C_ORG)

    strengths = [
        (C_FIRE,  "十字目変換フローという核心設計",
         "レア役→1/100で十字目変換発生→\n"
         "キャラ変換演出で期待度可視化→伝導者決戦\n"
         "という緊張感の連鎖が毎Gを意味ある時間にする。"),
        (C_FLAME, "5.8枚/G×ループ型ATの出玉体験",
         "炎炎大戦ループ率80%×純増5.8枚/Gが\n"
         "他台にない「速くて大きな勝ち体験」を提供。\n"
         "紅J大戦・アドラバーストが夢として機能。"),
        (C_CYAN,  "3段階天井＋スルー天井の安心設計",
         "850G・2,000G・5スルー天井の多層構造で\n"
         "投資計画が立てやすく来店動機になる。\n"
         "ハイエナ立ち回りの明確な指標にもなる。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1280000)
        rect_b(s, lx, iy, lw, Emu(1210000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1210000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(260000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(325000), lw - Emu(90000), Emu(780000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_ORG, font=FONT_H)

    principles = [
        (C_FIRE,  "毎Gの行動に意味を持たせよ",
         "1/100の変換チャンスが打ち手を\n能動的にし、回し続けるモチベーションを生む"),
        (C_FIRE2, "期待度はキャラの格で直感的に伝えよ",
         "アイリス<シンラ<紅丸という\n原作序列＝期待度がUXと物語を融合させる"),
        (C_ORG,   "「残弾数」で安心感を設計せよ",
         "ストック型のセット継続保証が\nストレスを消し長時間稼働を促す"),
        (C_FLAME, "常に上の目標を見せよ",
         "多層の上位ATが\"次のステージ\"を\n常に提示し離席を防ぐ設計の鉄則"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(760000)
        rect_b(s, rx, py0, rw, Emu(710000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(710000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(230000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(285000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3310000), rw, Emu(1060000),
           RGBColor(0x1A, 0x06, 0x02), C_FIRE, 2.0)
    rect(s, rx, ry + Emu(3310000), Emu(40000), Emu(1060000), C_FIRE)
    tb(s, rx + Emu(65000), ry + Emu(3360000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_FIRE)
    tb(s, rx + Emu(65000), ry + Emu(3620000), rw - Emu(85000), Emu(690000),
       "十字目変換×キャラ期待度演出×3段階天井×ループ型AT\n"
       "という4要素の統合設計。\n"
       "原作「炎」の世界観をゲーム性に昇華した\n"
       "2026年導入機の中で設計完成度が高い一台。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "本機の設計思想：「毎G意味付け・キャラ期待度演出・残弾安心感・多層目標」の4原則を次世代機設計に活用せよ",
           "十字目変換フロー＋キャラ変換演出の組み合わせは現代ATの教科書的事例として参照価値が高い")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)     # 1: タイトル・スペック・3ポイント
    s_flow(prs)      # 2: ゲームフロー全体図
    s_normal(prs)    # 3: 通常時の遊び方
    s_at_flow(prs)   # 4: AT炎炎激闘・十字目変換フロー詳細
    s_extend(prs)    # 5: 出玉を伸ばす方法
    s_upper(prs)     # 6: 上位モード
    s_design(prs)    # 7: 面白さの設計
    s_pros_cons(prs) # 8: 良い点と課題
    s_matome(prs)    # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
