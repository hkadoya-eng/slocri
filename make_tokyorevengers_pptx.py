"""
スマスロ 東京リベンジャーズ 機種説明＋分析 統合版 PPTXジェネレーター v1
出力: proposals/機種分析/東京リベンジャーズ/tokyorevengers_guide_v1.pptx
テーマ: 深黒×オレンジ(#FF6600)×青(#2266CC)×金（不良・熱血）
情報源: altema.jp・chonborista.com・1geki.jp・p-gabu.jp・nana-press.com（2025年〜2026年）
パチスロアワード2025ノミネート機種
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__),
           "proposals", "機種分析", "東京リベンジャーズ", "tokyorevengers_guide_v1.pptx")
os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)

# ── カラーパレット（深黒×オレンジ×青×金）──────────────────────────
C_BG    = RGBColor(0x08, 0x06, 0x02)
C_CARD  = RGBColor(0x12, 0x0E, 0x04)
C_CARD2 = RGBColor(0x1A, 0x14, 0x06)
C_ROW   = RGBColor(0x16, 0x10, 0x05)
C_ORG   = RGBColor(0xFF, 0x66, 0x00)   # オレンジ（主カラー）#FF6600
C_ORG2  = RGBColor(0xFF, 0x88, 0x22)   # 薄オレンジ
C_BLUE  = RGBColor(0x22, 0x66, 0xCC)   # 青 #2266CC
C_BLUE2 = RGBColor(0x44, 0x99, 0xFF)   # 薄青
C_GOLD  = RGBColor(0xCC, 0xAA, 0x30)   # 金
C_GOLD2 = RGBColor(0xFF, 0xDD, 0x55)   # 金（明るめ）
C_WHITE = RGBColor(0xF2, 0xEE, 0xE0)
C_CREAM = RGBColor(0xF0, 0xD8, 0xA0)
C_GRAY  = RGBColor(0x99, 0x88, 0x70)
C_LTGRY = RGBColor(0x55, 0x48, 0x38)
C_RED   = RGBColor(0xDD, 0x22, 0x22)
C_GREEN = RGBColor(0x33, 0xBB, 0x66)
C_CYAN  = RGBColor(0x22, 0xCC, 0xDD)
C_PINK  = RGBColor(0xFF, 0x44, 0x88)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景・ヘルパー群（夜の街×オレンジ） ─────────────────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (8, 6, 2))
    draw = ImageDraw.Draw(img)
    # 斜線テクスチャ（都市感）
    for i in range(0, w + h, 70):
        draw.line([(i, 0), (0, i)], fill=(14, 10, 3), width=1)
    # 下部オレンジグロー（ネオン街灯）
    for y in range(h - 110, h):
        t = (y - (h - 110)) / 110
        r = int(60 * t)
        g = int(18 * t)
        b = int(0)
        draw.line([(0, y), (w, y)], fill=(r, g, b))
    # 上部薄暗トップ
    for y in range(0, 35):
        t = (35 - y) / 35 * 0.4
        draw.line([(0, y), (w, y)], fill=(int(8 * t), int(4 * t), 0))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s


def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or C_WHITE
    run.font.name = font or FONT_B


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def rect_b(slide, x, y, w, h, fill, border, bw=1.0):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp


def hdr(slide, title_text, pg=""):
    rect(slide, 0, 0, SLIDE_W, Inches(0.58), C_CARD)
    rect(slide, 0, 0, Emu(45000), Inches(0.58), C_ORG)
    tb(slide, Inches(0.15), Emu(28000), Inches(8.5), Emu(410000),
       title_text, 14, bold=True, color=C_ORG, font=FONT_H)
    if pg:
        tb(slide, Inches(8.8), Emu(45000), Inches(1.0), Emu(340000),
           pg, 8, color=C_GRAY, align=PP_ALIGN.RIGHT)
    rect(slide, 0, Inches(0.58), SLIDE_W, Emu(7000), C_ORG)


def net_note(slide):
    tb(slide, Inches(7.5), Inches(5.38), Inches(2.4), Emu(180000),
       "※ネット解析情報より（altema・ちょんぼりすた・一撃・なな徹）", 6.5, color=C_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, bold_text, sub_text=""):
    fy = Inches(5.08)
    rect(slide, 0, fy, SLIDE_W, Inches(0.545), RGBColor(0x0C, 0x09, 0x02))
    rect(slide, 0, fy, Emu(20000), Inches(0.545), C_ORG)
    tb(slide, Inches(0.18), fy + Emu(40000), Inches(5.5), Emu(340000),
       bold_text, 7.5, bold=True, color=C_ORG)
    if sub_text:
        tb(slide, Inches(5.8), fy + Emu(40000), Inches(4.0), Emu(340000),
           sub_text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


def arrow_r(slide, x, cy, col=None):
    shp = slide.shapes.add_shape(13, x, cy - Emu(90000), Emu(200000), Emu(180000))
    shp.fill.solid()
    shp.fill.fore_color.rgb = col or C_ORG
    shp.line.fill.background()


def arrow_d(slide, cx, y, col=None):
    shp2 = slide.shapes.add_shape(14, cx - Emu(90000), y, Emu(180000), Emu(180000))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = col or C_ORG
    shp2.line.fill.background()
    return shp2


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル・スペック・この台の3ポイント
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, 0, 0, Inches(5.4), SLIDE_H, RGBColor(0x08, 0x05, 0x01))
    rect(s, 0, 0, Emu(55000), SLIDE_H, C_ORG)
    rect(s, Inches(5.4), 0, Emu(8000), SLIDE_H, C_ORG)

    tb(s, Inches(0.22), Inches(0.3), Inches(5.0), Emu(300000),
       "機種説明＋分析 統合ガイド  v1（パチスロアワード2025ノミネート）", 9, color=C_ORG2, font=FONT_H)
    tb(s, Inches(0.22), Inches(0.78), Inches(5.1), Emu(900000),
       "東京リベンジャーズ", 30, bold=True, color=C_ORG, font=FONT_H)
    tb(s, Inches(0.22), Inches(2.5), Inches(5.0), Emu(280000),
       "スマスロ（Lスロット）── タイムリープ×不良バトル×差枚数上乗せ型AT", 9, color=C_CREAM, font=FONT_H)

    # スペック表
    specs = [
        ("メーカー",         "サミー（SAMMY）  2025年9月8日導入"),
        ("設定",            "1〜6段階"),
        ("AT純増（通常）",   "約3.2枚/G（東卍RUSH）"),
        ("AT純増（上位）",   "約8.0枚/G（東卍RUSH BURST）"),
        ("設定1機械割",      "97〜98%前後（投資200Gで105%超）"),
        ("設定6機械割",      "110%前後"),
        ("天井①",          "周期6周期目（最大約500pt×6=3000pt相当）"),
        ("天井②",          "AT間ゲーム数天井（AT当選確定）"),
        ("初当り時AT期待度", "約60%（東卍チャンス経由含む）"),
        ("CZ種類",          "ミッドナイトモード（成功約50%）稀咲陰謀（成功約75%）"),
    ]
    for i, (k, v) in enumerate(specs):
        ry = Inches(3.0) + i * Emu(220000)
        tb(s, Inches(0.22), ry, Inches(1.75), Emu(205000),
           k, 7.5, color=C_GRAY)
        tb(s, Inches(1.97), ry, Inches(3.25), Emu(205000),
           v, 7.5, bold=True, color=C_WHITE)

    # 右パネル：この台の3ポイント
    kws = [
        (C_ORG,   "① タイムリープ演出×不良バトル",
         "原作のタイムリープ要素を「前兆→CZ→バトル」に落とし込み。\n"
         "東卍（東京卍會）の不良バトル世界観が\n通常時から演出を彩る熱血設計。"),
        (C_BLUE,  "② 周期+CZ+天井の多層突入設計",
         "毎Gポイント獲得→規定pt到達で前兆→CZ経由or直接AT。\n"
         "2種CZ（約50%・約75%）と天井で\n初当りルートが多彩。"),
        (C_GOLD,  "③ 上位AT「東卍RUSH BURST」で大爆発",
         "純増8.0枚/Gの上位AT。継続率約80%。\n"
         "特化ゾーン「黒い衝動」では期待3,000枚超。\n"
         "リベンジチャンス成功で到達できる。"),
    ]
    for i, (ac, kw, desc) in enumerate(kws):
        y0 = Inches(0.48) + i * Emu(1530000)
        rect_b(s, Inches(5.65), y0, Inches(4.1), Inches(1.28), C_CARD, ac, 2.0)
        rect(s, Inches(5.65), y0, Emu(60000), Inches(1.28), ac)
        tb(s, Inches(5.85), y0 + Emu(60000), Inches(3.8), Emu(310000),
           kw, 11, bold=True, color=ac, font=FONT_H)
        tb(s, Inches(5.85), y0 + Emu(375000), Inches(3.8), Emu(470000),
           desc, 8, color=C_WHITE)

    net_note(s)
    footer(s, "設計核心：「周期システム×2段CZ×差枚数上乗せ型AT×上位AT爆発」── タイムリープ原作世界観をゲーム性に融合",
           "純増3.2→8.0枚/Gの段階的上位ATが本機最大の出玉爆発力")


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: ゲームフロー全体図（蛇行2段で可視化）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図 ── 通常時→AT→上位ATの全ルート", "2/9")

    # 上段：通常時→初当り→東卍RUSH
    top_y = Inches(0.76)
    top_h = Emu(1060000)

    flow1 = [
        (C_CARD2,  C_GRAY,   "通常遊技",           "毎Gポイント獲得\n規定pt到達で\n前兆・CZへ"),
        (C_CARD,   C_BLUE2,  "東卍アクセル\n(特化5G+α)",  "ポイント獲得\n特化ゾーン\nSTタイプ"),
        (C_CARD,   C_ORG2,   "CZ\nミッドナイト\nモード",   "成功約50%\nフリーズで\n成功濃厚"),
        (C_CARD,   C_ORG,    "CZ\n稀咲陰謀",        "成功約75%\n上位CZ\n突入で激アツ"),
        (C_CARD,   C_GOLD,   "東卍RUSH\n(AT/純増3.2枚)", "差枚数\n管理型\nメインAT"),
    ]
    bw1 = Inches(1.57)
    gap1 = Inches(0.17)
    sx1 = Inches(0.28)
    cy1 = top_y + top_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow1):
        bx = sx1 + i * (bw1 + gap1)
        rect_b(s, bx, top_y, bw1, top_h, fill, ac, 1.8)
        tb(s, bx + Emu(30000), top_y + Emu(65000), bw1 - Emu(55000), Emu(360000),
           lbl, 9, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(20000), top_y + Emu(470000), bw1 - Emu(38000), Emu(520000),
           sub, 7.5, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow_r(s, bx + bw1 + Emu(10000), cy1)

    # 天井アノテーション
    tb(s, sx1, top_y + Emu(1080000), Inches(3.5), Emu(260000),
       "天井①：6周期目で AT当選確定（最大保険）", 7.5, color=C_CYAN)
    rect(s, sx1, top_y + Emu(1320000), Inches(3.8), Emu(5000), C_CYAN)
    tb(s, sx1 + Inches(4.0), top_y + Emu(1080000), Inches(3.8), Emu(260000),
       "天井②：AT間ゲーム数天井でAT確定（長期保険）", 7.5, color=C_BLUE2)

    # 中段区切り線
    rect(s, 0, Inches(2.1), SLIDE_W, Emu(5000), RGBColor(0x44, 0x20, 0x00))

    # 下段：AT内→上位AT昇格ルート
    bot_y = Inches(2.17)
    bot_h = Emu(1080000)

    flow2 = [
        (C_CARD,   C_GOLD,   "東卍RUSH\n（基本AT）",    "BREAK CHANCE\nで初期枚数決定\n純増3.2枚/G"),
        (C_CARD,   C_ORG,    "一触即発\n（AT内バトル）",  "レア役で上乗せ\nLEAP昇格で\n3桁上乗せも"),
        (C_CARD,   C_BLUE,   "黒い衝動\n（特化5G+α）",  "100枚以上\nの上乗せ保証\nSTリセット"),
        (RGBColor(0x16,0x0C,0x02), C_GOLD2,
         "東卍RUSH\nBURST\n（上位AT）",  "純増8.0枚/G\n継続率約80%\n上乗せ50枚以上"),
    ]
    bw2 = Inches(1.98)
    gap2 = Inches(0.22)
    sx2 = Inches(0.28)
    cy2 = bot_y + bot_h // 2

    for i, (fill, ac, lbl, sub) in enumerate(flow2):
        bx = sx2 + i * (bw2 + gap2)
        rect_b(s, bx, bot_y, bw2, bot_h, fill, ac, 1.8 if i == 3 else 1.5)
        tb(s, bx + Emu(35000), bot_y + Emu(65000), bw2 - Emu(60000), Emu(380000),
           lbl, 10, bold=True, color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(25000), bot_y + Emu(490000), bw2 - Emu(45000), Emu(480000),
           sub, 8, color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            arrow_r(s, bx + bw2 + Emu(12000), cy2)

    # 右端：リベンジチャンス説明
    rx_rc = sx2 + 4 * (bw2 + gap2)
    rect_b(s, rx_rc, bot_y, Inches(1.45), bot_h, C_CARD, C_PINK, 1.5)
    rect(s, rx_rc, bot_y, Emu(28000), bot_h, C_PINK)
    tb(s, rx_rc + Emu(48000), bot_y + Emu(70000), Inches(1.2), Emu(310000),
       "リベンジ\nチャンス\n(CZ)", 9, bold=True, color=C_PINK, font=FONT_H, align=PP_ALIGN.CENTER)
    tb(s, rx_rc + Emu(48000), bot_y + Emu(490000), Inches(1.2), Emu(460000),
       "BURST突入\n期待度\n高い", 7.5, color=C_WHITE, align=PP_ALIGN.CENTER)

    net_note(s)
    footer(s, "上段=通常〜東卍RUSH突入ルート（周期+2段CZ+天井）、下段=AT内上乗せ→上位BURST昇格ルート",
           "周期到達で必ずCZ/前兆が発生する設計のため「1周期=1チャンス」の安心感がある")


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: 通常時の遊び方（天井含む全ルート）
# ══════════════════════════════════════════════════════════════
def s_normal(prs):
    s = new_slide(prs)
    hdr(s, "通常時の遊び方 ── 周期システム・全AT突入ルート・天井", "3/9")

    # 左：AT突入ルート図
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(285000), RGBColor(0x55, 0x20, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(220000),
       "通常時〜東卍RUSH突入ルート（全4系統）", 10, bold=True, color=C_ORG)

    routes = [
        (C_ORG,   "ルート①  周期到達→前兆→直接AT（最短）",
         "毎Gポイントが加算（最低1pt）され、規定ptに\n達すると「決戦前夜」前兆に移行→AT抽選。\n1周期最大500pt、6周期目が天井①。"),
        (C_BLUE2, "ルート②  東卍アクセル→ポイント加速→AT",
         "レア役などで東卍アクセル（ST5G+α）が発生。\nポイントを集中的に獲得し周期到達を早める。\nリプレイ・レア役成立で5G再セット（継続）。"),
        (C_ORG2,  "ルート③  CZ（ミッドナイトモード/稀咲陰謀）→AT",
         "規定pt到達 or レア役でCZへ直接突入。\nミッドナイトモード：成功約50%（フリーズで濃厚）\n稀咲陰謀（上位CZ）：成功約75%でAT濃厚。"),
        (C_CYAN,  "ルート④  天井（2種）→AT確定",
         "天井①（6周期目）→AT当選確定で発動。\n天井②（AT間ゲーム数）→AT当選確定。\n設定変更（リセット）後は最大900Gに短縮。"),
    ]
    for i, (ac, t, b) in enumerate(routes):
        iy = ly + Emu(285000) + i * Emu(1118000)
        rect_b(s, lx, iy, lw, Emu(1055000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1055000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(265000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(315000), lw - Emu(100000), Emu(665000),
           b, 7.5, color=C_WHITE)

    # 右：チャンス役確率・打ち方
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(285000), RGBColor(0x55, 0x20, 0x00))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(220000),
       "チャンス役の種類と役割", 10, bold=True, color=C_ORG)

    chance = [
        (C_GRAY,  "リプレイ",      "約1/7〜8",   "東卍アクセル中は継続抽選の起点"),
        (C_ORG2,  "弱チェリー",    "約1/50",    "CZ・ポイント加速の抽選対象"),
        (C_ORG,   "スイカ",        "約1/80〜90","周期短縮・CZ直接突入の抽選対象"),
        (C_BLUE2, "チャンス目",     "約1/100",   "CZ直接突入・AT直当り抽選対象"),
        (C_GOLD,  "卍目（強役）",   "約1/500+",  "AT当選or上位CZ直行の激アツ役"),
        (C_CYAN,  "レア役合算",     "約1/30前後","全体的なAT期待役のまとめ"),
    ]
    ch_h = Emu(680000)
    for i, (ac, role, prob, desc) in enumerate(chance):
        cy = ry + Emu(285000) + i * ch_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, cy, rw, ch_h, bg)
        rect(s, rx, cy, Emu(35000), ch_h, ac)
        tb(s, rx + Emu(55000), cy + Emu(50000), Inches(1.05), Emu(260000),
           role, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(55000) + Inches(1.05), cy + Emu(55000), Inches(0.85), Emu(235000),
           prob, 8, bold=True, color=C_GOLD2, wrap=False)
        tb(s, rx + Emu(55000), cy + Emu(305000), rw - Emu(60000), Emu(320000),
           desc, 7.5, color=C_WHITE)

    # 打ち方メモ
    rect_b(s, rx, ry + Emu(4380000), rw, Emu(540000), C_CARD2, C_ORG, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4430000), rw - Emu(80000), Emu(195000),
       "基本打ち方：通常時は順押しで消化", 8.5, bold=True, color=C_ORG)
    tb(s, rx + Emu(60000), ry + Emu(4635000), rw - Emu(80000), Emu(260000),
       "AT中は押し順ナビに従う。卍目成立時は特に注意。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "通常時の基本戦略：天井①（6周期）到達を軸に管理し、稀咲陰謀（成功75%）突入時は大チャンスと認識する",
           "東卍アクセルはポイント加速の重要エンジン。発生時は必ず継続を狙う5Gプレイを意識")


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: CZ/前兆の仕組み（不良バトル・タイムリープ要素）
# ══════════════════════════════════════════════════════════════
def s_cz(prs):
    s = new_slide(prs)
    hdr(s, "CZ/前兆の仕組み ── 不良バトル演出×タイムリープ要素の絡み", "4/9")

    # 左：CZ詳細
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(280000), RGBColor(0x50, 0x22, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(215000),
       "CZ（チャンスゾーン）の2種類と仕組み", 10, bold=True, color=C_ORG)

    czs = [
        (C_ORG2, "ミッドナイトモード（通常CZ）",
         "成功期待度：約50%\n"
         "フリーズ発生で成功濃厚。\n"
         "消化中は不良バトル演出が展開。\n"
         "負けると通常時に戻るが再チャレンジあり。\n"
         "成功→東卍RUSH（AT）当選ほぼ確定。"),
        (C_ORG,  "稀咲陰謀（上位CZ）",
         "成功期待度：約75%\n"
         "敵サイドの黒幕・稀咲鈴次が登場する\n"
         "上位チャンスゾーン。\n"
         "突入だけで大チャンス状態。\n"
         "成功で東卍RUSH確定級の扱い。"),
        (C_BLUE, "前兆「決戦前夜」",
         "周期到達後に発生する前兆ステージ。\n"
         "東京卍會の結成・決戦エピソードが展開。\n"
         "原作の\"決戦前夜\"シーンを再現した演出で\n"
         "ATへの期待感を煽る。\n"
         "前兆中のレア役でAT期待度UP。"),
    ]
    for i, (ac, t, b) in enumerate(czs):
        iy = ly + Emu(280000) + i * Emu(1460000)
        rect_b(s, lx, iy, lw, Emu(1390000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1390000), ac)
        tb(s, lx + Emu(75000), iy + Emu(50000), lw - Emu(100000), Emu(270000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, lx + Emu(75000), iy + Emu(330000), lw - Emu(100000), Emu(950000),
           b, 8, color=C_WHITE)

    # 右：タイムリープ演出設計
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(280000), RGBColor(0x10, 0x20, 0x50))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(215000),
       "タイムリープ演出設計：原作要素とパチスロの融合", 10, bold=True, color=C_BLUE2)

    tlinfo = [
        (C_BLUE,  "タイムリープ演出の役割",
         "原作主人公・花垣武道が過去に戻って\n"
         "\"最悪の未来\"を変えようとする物語を\n"
         "「AT再挑戦」「CZ再突入」という\n"
         "パチスロ演出に落とし込み。\n"
         "負けても\"もう一度過去に戻る\"演出が入り\n"
         "リトライ演出の納得感を高める。"),
        (C_BLUE2, "不良バトル演出（東卍 vs 芭流覇羅など）",
         "東京卍會（東卍）と敵対勢力とのバトルが\n"
         "CZ・AT中の成否演出に使用される。\n"
         "バトルに勝つ=CZ成功/AT継続の直感演出。\n"
         "原作ファンは登場キャラへの感情移入で\n"
         "盛り上がりが倍増する設計。"),
        (C_GOLD,  "演出と数値の対応（信頼度設計）",
         "バトル演出の激しさ・参戦キャラで\n"
         "CZ成功期待度が変化。\n"
         "マイキー（佐野万次郎）登場→最高峰期待度。\n"
         "演出のキャラ格順＝数値期待度順に揃えた\n"
         "直感的UX設計が本機の強み。"),
    ]
    for i, (ac, t, b) in enumerate(tlinfo):
        iy = ry + Emu(280000) + i * Emu(1445000)
        rect_b(s, rx, iy, rw, Emu(1375000), C_CARD, ac, 1.5)
        rect(s, rx, iy, Emu(40000), Emu(1375000), ac)
        tb(s, rx + Emu(70000), iy + Emu(50000), rw - Emu(90000), Emu(270000),
           t, 9, bold=True, color=ac)
        tb(s, rx + Emu(70000), iy + Emu(330000), rw - Emu(90000), Emu(920000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "CZ設計の核心：2段階CZ（50%・75%）で「チャンス→大チャンス」の段階的期待感を演出。稀咲出現が最大の盛り上がり",
           "タイムリープ演出を「再挑戦の納得感」に転用した点が原作愛・ゲーム性両立の好例")


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: AT/ボーナス（何をすれば出玉が伸びる）
# ══════════════════════════════════════════════════════════════
def s_at(prs):
    s = new_slide(prs)
    hdr(s, "AT「東卍RUSH」中の仕組み ── 出玉が伸びるメカニズム", "5/9")

    # 左：AT内フロー詳細
    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.55)

    rect(s, lx, ly, lw, Emu(280000), RGBColor(0x50, 0x22, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(215000),
       "東卍RUSH（AT）中の上乗せフロー", 10, bold=True, color=C_ORG)

    atflow = [
        (C_GOLD,  "BREAK CHANCE（初期枚数決定）",
         "AT初当り時に「BREAK CHANCE」で\n"
         "初期差枚数（上乗せスタート枚数）を決定。\n"
         "初期枚数次第で序盤の資産が変わる。\n"
         "多いほど以降の展開が余裕になる。"),
        (C_ORG,   "東卍ATTACK（規定ポイント到達で発生）",
         "AT消化中に規定ポイント到達で東卍ATTACK。\n"
         "バトルに勝利→差枚数上乗せ or 特化ゾーン獲得。\n"
         "レア役成立でポイント加速。\n"
         "負けても消えず次の戦いへ継続。"),
        (C_BLUE,  "LEAP上乗せ（レア役連続で昇格）",
         "次Gに再度レア役成立でLEAP上乗せに昇格。\n"
         "通常上乗せより大幅アップ（3桁上乗せも）。\n"
         "「日和ってループ」発生の可能性もあり\n"
         "連続引き時は最大の期待タイミング。"),
        (C_CYAN,  "黒い衝動（特化ゾーン：100枚以上保証）",
         "5G+α継続のST型特化ゾーン。\n"
         "1上乗せあたり100枚以上確定。\n"
         "上乗せ当選時はSTのG数がリセット。\n"
         "通常AT中の最大瞬発力ゾーン。"),
    ]
    for i, (ac, t, b) in enumerate(atflow):
        iy = ly + Emu(280000) + i * Emu(1145000)
        rect_b(s, lx, iy, lw, Emu(1080000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(45000), Emu(1080000), ac)
        tb(s, lx + Emu(75000), iy + Emu(48000), lw - Emu(100000), Emu(265000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(75000), iy + Emu(315000), lw - Emu(100000), Emu(680000),
           b, 7.5, color=C_WHITE)

    # 右：上乗せ分岐表
    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(280000), RGBColor(0x50, 0x22, 0x00))
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(215000),
       "上乗せ種類・期待枚数・発生条件の早見表", 10, bold=True, color=C_ORG)

    upper_tbl = [
        (C_GRAY,   "通常上乗せ",     "10〜99枚",    "東卍ATTACK勝利時"),
        (C_BLUE2,  "LEAP上乗せ",     "100枚以上",   "レア役連続成立で昇格"),
        (C_ORG,    "黒い衝動",       "100枚以上保証", "特化ゾーン突入時"),
        (C_GOLD,   "HEAT UP!!",      "100枚以上濃厚", "BURST中一触即発勝利"),
        (C_GOLD2,  "BURST通常上乗せ","50枚以上確定",  "BURST中レア役成立"),
        (C_PINK,   "BURST LEAP",     "300枚以上濃厚", "BURST中LEAP昇格時"),
    ]
    tbl_h = Emu(710000)
    for i, (ac, nm, pct, cond) in enumerate(upper_tbl):
        ty = ry + Emu(280000) + i * tbl_h
        bg = C_CARD if i % 2 == 0 else C_ROW
        rect(s, rx, ty, rw, tbl_h, bg)
        rect(s, rx, ty, Emu(32000), tbl_h, ac)
        tb(s, rx + Emu(52000), ty + Emu(52000), Inches(1.15), Emu(260000),
           nm, 8.5, bold=True, color=ac, wrap=False)
        tb(s, rx + Emu(52000) + Inches(1.15), ty + Emu(52000), Inches(1.1), Emu(260000),
           pct, 8, bold=True, color=C_GOLD2, wrap=False)
        tb(s, rx + Emu(52000), ty + Emu(320000), rw - Emu(65000), Emu(330000),
           cond, 7.5, color=C_WHITE)

    # 補足コメント
    rect_b(s, rx, ry + Emu(4555000), rw, Emu(475000), C_CARD2, C_GOLD, 1.5)
    tb(s, rx + Emu(60000), ry + Emu(4600000), rw - Emu(80000), Emu(195000),
       "出玉を伸ばすポイント", 8.5, bold=True, color=C_GOLD)
    tb(s, rx + Emu(60000), ry + Emu(4800000), rw - Emu(80000), Emu(205000),
       "LEAP昇格連打 + 黒い衝動複数突入 がセット利益を決定する。", 7.5, color=C_GRAY)

    net_note(s)
    footer(s, "AT中の出玉設計：「東卍ATTACK→上乗せ→特化ゾーン→LEAP」のサイクル積み重ねが基本の利益構造",
           "BURST中はすべての上乗せが強化版に。BURST突入後はコイン単価が劇的に変化する")


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 上位ATへの道と遊び方
# ══════════════════════════════════════════════════════════════
def s_upper(prs):
    s = new_slide(prs)
    hdr(s, "上位AT「東卍RUSH BURST」── 到達ルート・性能・継続設計", "6/9")

    # 上段：昇格ルートフロー
    rect(s, 0, Inches(0.72), SLIDE_W, Emu(265000), RGBColor(0x44, 0x20, 0x00))
    tb(s, Inches(0.35), Inches(0.755), Inches(9.0), Emu(215000),
       "東卍RUSH BURST 到達ルート（通常AT→上位への昇格経路）", 9, bold=True, color=C_ORG)

    route_boxes = [
        (C_GOLD,  "東卍RUSH\n（基本AT）",    "差枚数管理型\n純増3.2枚/G\nスタート地点"),
        (C_ORG,   "リベンジ\nチャンス\n(CZ)", "AT内発生\n上位突入\nの登竜門"),
        (C_GOLD2, "東卍RUSH\nBURST\n（上位AT）", "純増8.0枚/G\n継続率約80%\n上乗せ全強化"),
    ]
    bw_r = Inches(2.7)
    gap_r = Inches(0.42)
    sx_r = Inches(0.55)
    cy_r = Inches(1.43)
    bh_r = Emu(1100000)

    for i, (ac, lbl, sub) in enumerate(route_boxes):
        bx = sx_r + i * (bw_r + gap_r)
        rect_b(s, bx, cy_r - bh_r // 2, bw_r, bh_r,
               C_CARD if i < 2 else RGBColor(0x18, 0x0E, 0x01), ac, 2.0 if i == 2 else 1.5)
        tb(s, bx + Emu(40000), cy_r - bh_r // 2 + Emu(70000),
           bw_r - Emu(60000), Emu(380000), lbl, 11, bold=True,
           color=ac, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, bx + Emu(35000), cy_r - bh_r // 2 + Emu(500000),
           bw_r - Emu(55000), Emu(490000), sub, 8,
           color=C_GRAY, align=PP_ALIGN.CENTER)
        if i < 2:
            arrow_r(s, bx + bw_r + Emu(14000), cy_r, C_GOLD)

    # 中区切り
    rect(s, 0, Inches(2.08), SLIDE_W, Emu(5000), RGBColor(0x44, 0x20, 0x00))

    # 下段左：BURST中の性能詳細
    lx2, ly2 = Inches(0.28), Inches(2.14)
    lw2 = Inches(4.55)
    lh2 = Emu(2620000)

    rect_b(s, lx2, ly2, lw2, lh2, C_CARD, C_GOLD2, 2.0)
    rect(s, lx2, ly2, Emu(45000), lh2, C_GOLD2)
    tb(s, lx2 + Emu(75000), ly2 + Emu(50000), lw2 - Emu(95000), Emu(270000),
       "東卍RUSH BURST 中の性能・遊び方", 11, bold=True, color=C_GOLD2, font=FONT_H)
    tb(s, lx2 + Emu(75000), ly2 + Emu(340000), lw2 - Emu(95000), lh2 - Emu(400000),
       "【基本性能】\n"
       "純増：約8.0枚/G（通常ATの2.5倍）\n"
       "継続率：約80%（天上天下唯我独尊で継続判定）\n\n"
       "【上乗せ強化内容（BURST中）】\n"
       "・上乗せ発生時→50枚以上確定\n"
       "・LEAP時→300枚以上濃厚（通常の3倍）\n"
       "・HEAT UP!!時→100枚以上濃厚\n"
       "・一触即発勝利→HEAT UP!!確定\n\n"
       "【継続判定：天上天下唯我独尊】\n"
       "BURST消化後に継続抽選ステージへ移行。\n"
       "BURST継続期待度：約80%。\n"
       "継続するたびに出玉が蓄積する複利設計。",
       8, color=C_WHITE)

    # 下段右：到達の要件・狙い目
    rx2, ry2 = Inches(5.05), Inches(2.14)
    rw2 = Inches(4.65)

    rect_b(s, rx2, ry2, rw2, lh2,
           RGBColor(0x10, 0x08, 0x01), C_ORG, 2.0)
    rect(s, rx2, ry2, Emu(45000), lh2, C_ORG)
    tb(s, rx2 + Emu(75000), ry2 + Emu(50000), rw2 - Emu(95000), Emu(270000),
       "BURST到達のための考え方と遊び方", 11, bold=True, color=C_ORG, font=FONT_H)
    tb(s, rx2 + Emu(75000), ry2 + Emu(340000), rw2 - Emu(95000), lh2 - Emu(400000),
       "【リベンジチャンスとは】\n"
       "通常AT（東卍RUSH）消化中に発生する\n"
       "BURST昇格のためのCZ。成功でBURST確定。\n\n"
       "【リベンジチャンスへの道】\n"
       "AT中の一触即発勝利・レア役・LEAP上乗せ\n"
       "などを重ねることで内部的に昇格抽選が進む。\n\n"
       "【フリーズによる直接BURST】\n"
       "初当り時やAT中にフリーズが発生した場合は\n"
       "BURST直行または高期待度への昇格となる。\n\n"
       "【BURST中の立ち回り】\n"
       "継続率80%を活かして長く打ち続けることが\n"
       "大量獲得の鍵。通常AT時より単価が高いため\n"
       "BURST中は残差枚数を意識した打ち止め判断を。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "BURST設計の核心：純増8.0枚×継続率80%×上乗せ全強化という3点セットが本機最大の爆発力を生む",
           "天上天下唯我独尊の継続判定演出が原作のマイキーvsキャラ対決と連動。演出の熱さが継続への期待感を倍増")


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: 面白さの設計（タイムリープ×不良バトル×パチスロ）
# ══════════════════════════════════════════════════════════════
def s_design(prs):
    s = new_slide(prs)
    hdr(s, "面白さの設計 ── タイムリープ×不良バトルとパチスロの融合", "7/9")

    principles = [
        (C_ORG,   "① 周期システムによる「毎周期のメリハリ」",
         "毎Gポイントが積み上がり、規定ptで必ず何かが起きる。\n"
         "「あと少しで周期到達」という小目標が生まれ\n"
         "打ち続けるモチベーションが常に維持される。"),
        (C_BLUE,  "② タイムリープ演出＝リトライの納得感",
         "原作の\"何度も過去に戻る\"物語構造をCZ失敗時の\n"
         "再挑戦演出に落とし込んだ秀逸な転用設計。\n"
         "負けが悔しくなく「また戻って再挑戦！」と感じさせる。"),
        (C_ORG2,  "③ キャラの格＝期待度の直感UX",
         "マイキー登場>稀咲登場>他キャラという原作ヒエラルキーが\n"
         "そのままCZ期待度に対応。\n"
         "原作ファン・初心者両方が演出を見ただけで期待度を把握できる。"),
        (C_GOLD,  "④ 2段AT（通常→BURST）による目標の多層化",
         "通常ATでは「黒い衝動」を目指し、\n"
         "BURSTではさらに高い天井を狙う。\n"
         "常に「上のモード」が見えており\n"
         "プレイヤーの離席を自然に抑制する多層目標設計。"),
        (C_CYAN,  "⑤ 差枚数管理型ATによる「残り時間の見える化」",
         "差枚数管理型なので「あと何枚で終わる」が\n"
         "常にわかる。上乗せするたびに終了が遠のく設計が\n"
         "上乗せ演出の達成感を増幅させる仕組み。"),
    ]
    bw_p = Inches(4.55)
    bh_p = Emu(1190000)
    gy = Inches(0.095)

    positions = [
        (Inches(0.28),  Inches(0.72)),
        (Inches(5.17),  Inches(0.72)),
        (Inches(0.28),  Inches(0.72) + bh_p + gy),
        (Inches(5.17),  Inches(0.72) + bh_p + gy),
        (Inches(0.28),  Inches(0.72) + 2 * (bh_p + gy)),
    ]
    for i, (ac, t, b) in enumerate(principles):
        if i == 4:
            px = Inches(0.28)
            pw = Inches(9.44)
            ph = Emu(1080000)
        else:
            px, _ = positions[i]
            pw = bw_p
            ph = bh_p
        _, py = positions[i]

        rect_b(s, px, py, pw, ph, C_CARD, ac, 1.5)
        rect(s, px, py, Emu(40000), ph, ac)
        tb(s, px + Emu(70000), py + Emu(50000), pw - Emu(90000), Emu(255000),
           t, 9.5, bold=True, color=ac, font=FONT_H)
        tb(s, px + Emu(70000), py + Emu(320000), pw - Emu(90000), ph - Emu(385000),
           b, 8, color=C_WHITE)

    net_note(s)
    footer(s, "面白さの核心：「周期メリハリ・リトライ納得感・キャラ期待度直感・2段目標・差枚数見える化」の5要素が相互に機能",
           "タイムリープという原作の核心テーマをパチスロの再チャレンジ設計に転用した点が本機最大の設計的発明")


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 良い点と課題
# ══════════════════════════════════════════════════════════════
def s_pros_cons(prs):
    s = new_slide(prs)
    hdr(s, "良い点と課題 ── 設計の強みと改善余地", "8/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), C_GREEN)
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "良い点 ── 設計的強み", 10, bold=True, color=C_BG)

    pros = [
        (C_GREEN, "周期システムで毎周期にメリハリ",
         "最大6周期の構造で「いつかは当たる」安心感と\n"
         "「今周期こそ」の緊張感が同居。\n"
         "単純なゲーム数天井より体感的に飽きにくい。"),
        (C_GREEN, "2段AT（3.2枚→8.0枚）の爆発力差",
         "BURSTの純増8.0枚はスマスロでも高水準。\n"
         "継続率80%と重なり、BURST継続時の\n"
         "コイン増加速度は圧倒的な体感出玉を生む。"),
        (C_GREEN, "原作世界観を活かした演出設計",
         "タイムリープ・不良バトル・キャラ期待度の\n"
         "3層構造で原作ファンと新規層を両方取り込む。\n"
         "サミー史上最大の演出ボリュームを搭載。"),
        (C_GREEN, "差枚数管理型の透明性・UX",
         "「あと○枚」という残数可視化が\n"
         "上乗せ演出の達成感を直接的に高め\n"
         "「もう少しで終わる→上乗せしたい！」の循環を生む。"),
    ]
    for i, (ac, t, b) in enumerate(pros):
        iy = ly + Emu(260000) + i * Emu(1132000)
        rect_b(s, lx, iy, lw, Emu(1065000), C_CARD, ac, 1.2)
        rect(s, lx, iy, Emu(35000), Emu(1065000), ac)
        tb(s, lx + Emu(60000), iy + Emu(48000), lw - Emu(80000), Emu(248000),
           t, 8.5, bold=True, color=ac)
        tb(s, lx + Emu(60000), iy + Emu(298000), lw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_RED)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "課題 ── 改善余地・注意点", 10, bold=True, color=C_WHITE)

    cons = [
        (C_RED,   "通常ATの出玉スピードが物足りない",
         "純増3.2枚/Gは近年のスマスロとして\n"
         "「遅い」と感じるユーザーも多い。\n"
         "BURST前の通常AT消化が単調になりやすい。"),
        (C_ORG,   "評価が低め（DMMぱちタウン1.84/5点）",
         "ユーザー評価は低評価が多い傾向。\n"
         "「やってること自体はシンプル」との声も。\n"
         "BURSTまでの道のりの長さが不満要因か。"),
        (C_BLUE,  "BURST到達ルートがやや不透明",
         "リベンジチャンスの発生条件が\n"
         "解析情報が少なく内部が不透明。\n"
         "「BURST何回でも引けない」不満につながる場面も。"),
        (C_GRAY,  "設定差・設定判別の難易度",
         "設定差が明確な箇所が限られ\n"
         "高設定確信まで時間がかかる。\n"
         "低設定時の機械割が99%未満になる場合も。"),
    ]
    for i, (ac, t, b) in enumerate(cons):
        iy = ry + Emu(260000) + i * Emu(1132000)
        rect_b(s, rx, iy, rw, Emu(1065000), C_CARD, ac, 1.2)
        rect(s, rx, iy, Emu(35000), Emu(1065000), ac)
        tb(s, rx + Emu(60000), iy + Emu(48000), rw - Emu(80000), Emu(248000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(60000), iy + Emu(298000), rw - Emu(80000), Emu(660000),
           b, 7.5, color=C_WHITE)

    net_note(s)
    footer(s, "課題の本質：通常ATの体感薄さとBURST到達不透明感が改善されれば評価は大きく変わる可能性がある",
           "ユーザー評価1.84点という低評価は「やり取り中の退屈さ」が大きな要因。それだけBURST時の感動は大きい")


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ・設計から学べること
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ ── 設計から学べること", "9/9")

    lx, ly = Inches(0.28), Inches(0.72)
    lw = Inches(4.5)

    rect(s, lx, ly, lw, Emu(260000), RGBColor(0x55, 0x22, 0x00))
    tb(s, lx + Emu(60000), ly + Emu(55000), lw - Emu(80000), Emu(200000),
       "スマスロ東京リベンジャーズ ── 設計的強み総括", 10, bold=True, color=C_ORG)

    strengths = [
        (C_ORG,   "タイムリープ転用という設計的発明",
         "原作の核心テーマ「何度でも過去に戻りやり直す」を\n"
         "CZ失敗後の再挑戦演出に完全転用。\n"
         "負け体験を物語体験に変換した点が唯一無二。"),
        (C_GOLD2, "2段AT構造で「通常とBURSTの格差」を演出",
         "純増3.2枚→8.0枚の大きな段差が\n"
         "\"BURSTに入ったときの特別感\"を最大化。\n"
         "継続率80%と組み合わせで爆発力のギャップが感情を揺さぶる。"),
        (C_BLUE2, "パチスロアワード2025ノミネートの評価軸",
         "演出ボリューム・世界観再現性・サミー新技術\n"
         "という3軸でノミネートに値すると判断された。\n"
         "ユーザー評価と業界評価の乖離が興味深いポイント。"),
    ]
    for i, (ac, t, b) in enumerate(strengths):
        iy = ly + Emu(260000) + i * Emu(1275000)
        rect_b(s, lx, iy, lw, Emu(1205000), C_CARD, ac, 1.5)
        rect(s, lx, iy, Emu(40000), Emu(1205000), ac)
        tb(s, lx + Emu(70000), iy + Emu(50000), lw - Emu(90000), Emu(255000),
           t, 9, bold=True, color=ac)
        tb(s, lx + Emu(70000), iy + Emu(320000), lw - Emu(90000), Emu(775000),
           b, 8, color=C_WHITE)

    rx, ry = Inches(5.0), Inches(0.72)
    rw = Inches(4.7)

    rect(s, rx, ry, rw, Emu(260000), C_CARD2)
    tb(s, rx + Emu(60000), ry + Emu(55000), rw - Emu(80000), Emu(200000),
       "設計から学べる原則", 10, bold=True, color=C_ORG, font=FONT_H)

    principles = [
        (C_ORG,   "原作テーマをゲーム設計に転用せよ",
         "タイムリープ=再挑戦設計という発明を参考に\n物語の核心要素を遊技体験に変換する視点を持つ"),
        (C_BLUE2, "期待度はキャラの格で直感的に伝えよ",
         "マイキー>稀咲>他キャラの原作序列が\n数値期待度に対応。直感UXの手本"),
        (C_GOLD,  "通常とBURSTの格差を大きくせよ",
         "3.2枚→8.0枚という段差設計が\nBURST到達の「特別感」を最大化する"),
        (C_ORG2,  "周期システムは「小目標」を作るエンジン",
         "毎周期に必ず結果が出る設計が\nプレイヤーの打ち続ける理由を生み出す"),
    ]
    for i, (ac, t, b) in enumerate(principles):
        py0 = ry + Emu(260000) + i * Emu(755000)
        rect_b(s, rx, py0, rw, Emu(705000), C_CARD, ac, 1.0)
        rect(s, rx, py0, Emu(30000), Emu(705000), ac)
        tb(s, rx + Emu(55000), py0 + Emu(48000), rw - Emu(75000), Emu(228000),
           t, 8.5, bold=True, color=ac)
        tb(s, rx + Emu(55000), py0 + Emu(280000), rw - Emu(75000), Emu(350000),
           b, 7.5, color=C_WHITE)

    # 総括ボックス
    rect_b(s, rx, ry + Emu(3285000), rw, Emu(1080000),
           RGBColor(0x18, 0x0C, 0x01), C_ORG, 2.0)
    rect(s, rx, ry + Emu(3285000), Emu(40000), Emu(1080000), C_ORG)
    tb(s, rx + Emu(65000), ry + Emu(3335000), rw - Emu(85000), Emu(250000),
       "総括", 9, bold=True, color=C_ORG)
    tb(s, rx + Emu(65000), ry + Emu(3595000), rw - Emu(85000), Emu(720000),
       "周期システム×タイムリープ再挑戦演出×2段AT設計\n"
       "という3要素の統合設計が本機の骨格。\n"
       "ユーザー評価は低めだが、パチスロアワード2025\n"
       "ノミネートは演出設計の完成度が評価された証。",
       8, color=C_WHITE)

    net_note(s)
    footer(s, "本機の設計思想：「原作テーマ転用・キャラ期待度直感・2段AT格差・周期小目標」が次世代機設計への4原則",
           "タイムリープという固有テーマをパチスロ設計に昇華した発明は類似コンテンツの設計参照モデルになり得る")


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s_title(prs)      # 1: タイトル・スペック・この台の3ポイント
    s_flow(prs)       # 2: ゲームフロー全体図
    s_normal(prs)     # 3: 通常時の遊び方
    s_cz(prs)         # 4: CZ/前兆の仕組み
    s_at(prs)         # 5: AT/ボーナス（出玉が伸びる仕組み）
    s_upper(prs)      # 6: 上位ATへの道と遊び方
    s_design(prs)     # 7: 面白さの設計
    s_pros_cons(prs)  # 8: 良い点と課題
    s_matome(prs)     # 9: まとめ・設計から学べること

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
