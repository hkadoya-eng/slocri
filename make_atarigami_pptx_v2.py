"""
「祟り神の章」ゲーム性提案 PowerPoint ジェネレーター v2
v1からの進化：游明朝ヘッダー・純Pillow背景・蛇行フロー・net_note・絵文字なし
出力: proposals/atarigami_proposal_v2.pptx
"""
import io, os, sys
from PIL import Image as PILImage, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

OUT_PATH = os.path.join(os.path.dirname(__file__), "proposals", "atarigami_proposal_v2.pptx")

# ── カラーパレット（祟り神：紺黒・赤・青・紫・金テーマ）──────────
C_BG    = RGBColor(0x08, 0x08, 0x18)
C_CARD  = RGBColor(0x14, 0x14, 0x2C)
C_RED   = RGBColor(0xCC, 0x22, 0x22)
C_DKRED = RGBColor(0x30, 0x00, 0x00)
C_BLUE  = RGBColor(0x33, 0x55, 0xCC)
C_DKBLU = RGBColor(0x00, 0x00, 0x35)
C_PUR   = RGBColor(0x88, 0x22, 0xAA)
C_DKPUR = RGBColor(0x20, 0x08, 0x30)
C_GOLD  = RGBColor(0xC8, 0xA8, 0x40)
C_GOLD2 = RGBColor(0xFF, 0xD7, 0x00)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_CREAM = RGBColor(0xE0, 0xD8, 0xC8)
C_LTGRY = RGBColor(0xBB, 0xBB, 0xBB)
C_GRAY  = RGBColor(0x88, 0x88, 0x88)
C_GREEN = RGBColor(0x22, 0xCC, 0x66)
C_ORANGE= RGBColor(0xFF, 0x99, 0x00)
C_YELLOW= RGBColor(0xFF, 0xEE, 0x44)

FONT_H = "游明朝"
FONT_B = "メイリオ"
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


# ── 背景生成（深夜紺黒×紫の靄×底部の赤い怨念）──────────────────
def make_bg(w=1280, h=720):
    img = PILImage.new("RGB", (w, h), (8, 8, 24))
    draw = ImageDraw.Draw(img)
    for i in range(0, w + h, 55):
        draw.line([(i, 0), (0, i)], fill=(14, 10, 30), width=1)
    for y in range(h - 100, h):
        t = (y - (h - 100)) / 100
        draw.line([(0, y), (w, y)], fill=(int(50 * t), 0, 4))
    for y in range(0, 55):
        t = (55 - y) / 55 * 0.5
        draw.line([(0, y), (w, y)], fill=(int(25 * t), 0, int(40 * t)))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


# ── ヘルパー関数 ──────────────────────────────────────────────
def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = make_bg()
    pic = s.shapes.add_picture(bg, 0, 0, SLIDE_W, SLIDE_H)
    s.shapes._spTree.remove(pic._element)
    s.shapes._spTree.insert(2, pic._element)
    return s

def tb(slide, x, y, w, h, text, size=10, bold=False, color=None,
       align=PP_ALIGN.LEFT, font=None, wrap=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font or FONT_B
    if color:
        run.font.color.rgb = color

def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def rect_b(slide, x, y, w, h, fill, border, bw=1.5):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(bw)
    return shp

def arrow_r(slide, x, y, w, color):
    h = Emu(150000)
    shp = slide.shapes.add_shape(13, x, y - h // 2, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()

def hdr(slide, text):
    rect(slide, Inches(0.15), Inches(0.08), Inches(9.7), Emu(420000),
         RGBColor(0x10, 0x00, 0x20))
    rect(slide, Inches(0.15), Inches(0.08), Emu(60000), Emu(420000), C_RED)
    tb(slide, Inches(0.4), Inches(0.1), Inches(9.2), Emu(380000),
       text, 12, bold=True, color=C_GOLD, font=FONT_H)

def net_note(slide, text="※分析データより"):
    tb(slide, Inches(8.2), Inches(5.38), Inches(1.7), Emu(200000),
       text, 7, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 1: タイトル
# ══════════════════════════════════════════════════════════════
def s_title(prs):
    s = new_slide(prs)

    rect(s, Inches(0), Inches(0), Inches(5.6), SLIDE_H, RGBColor(0x04, 0x04, 0x12))
    rect(s, Inches(0.35), Inches(0.5), Emu(30000), Inches(2.0), C_RED)

    tb(s, Inches(0.5), Inches(0.52), Inches(5), Emu(360000),
       "新機種ゲーム性提案書", 14, color=C_GOLD, font=FONT_H)
    tb(s, Inches(0.5), Inches(1.05), Inches(5.2), Emu(1000000),
       "祟り神の章", 46, bold=True, color=C_WHITE, font=FONT_H)
    tb(s, Inches(0.5), Inches(2.35), Inches(5.0), Emu(300000),
       "── 物語と出玉が交差するスロット ──", 12, color=C_LTGRY, font=FONT_H)

    rect(s, Inches(0.5), Inches(2.95), Inches(4.8), Emu(600000), C_DKRED)
    rect(s, Inches(0.5), Inches(2.95), Emu(55000), Emu(600000), C_RED)
    tb(s, Inches(0.65), Inches(3.03), Inches(4.5), Emu(550000),
       "「 倒した相手を、後から理解する。 」\n\n普通ATで祟り神を戦って倒す。\n特別ATで「なぜそうなったか」が分かる。",
       10.5, color=C_WHITE)

    rect_b(s, Inches(5.8), Inches(0.8), Inches(3.8), Inches(2.6), C_CARD, C_PUR, 1.5)
    tb(s, Inches(5.95), Inches(0.88), Inches(3.5), Emu(320000),
       "基本スペック（概要）", 10, bold=True, color=C_PUR, font=FONT_H)
    quick_specs = [
        ("タイプ",   "スマスロ（L型）"),
        ("純増",     "3.8枚/G"),
        ("天井",     "600G（SP MAX時 300G）"),
        ("特別AT",   "100〜1500枚（加護数で変動）"),
        ("機械割",   "設定1: 97.5% / 設定6: 109%"),
    ]
    qy = Inches(1.35)
    for k, v in quick_specs:
        tb(s, Inches(5.95), qy, Inches(1.4), Emu(255000), k, 8.5, bold=True, color=C_GRAY, wrap=False)
        tb(s, Inches(7.45), qy, Inches(2.05), Emu(255000), v, 8.5, color=C_CREAM, wrap=False)
        qy += Emu(265000)

    rect_b(s, Inches(5.8), Inches(3.55), Inches(3.8), Inches(1.55),
           RGBColor(0x10, 0x08, 0x20), C_GOLD, 1.0)
    tb(s, Inches(5.95), Inches(3.62), Inches(3.5), Emu(280000),
       "3つの設計核", 9, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(5.95), Inches(3.96), Inches(3.5), Emu(600000),
       "I    感情逆転設計\nII   DQ4型章システム\nIII  加護の積み上げ来店設計",
       8.5, color=C_CREAM)

    tb(s, Inches(7.5), Inches(5.2), Inches(2.3), Emu(300000),
       "v2.0  2026.05", 8, color=C_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
#  SLIDE 2: コアコンセプト（感情逆転設計）
# ══════════════════════════════════════════════════════════════
def s_concept(prs):
    s = new_slide(prs)
    hdr(s, "コアコンセプト  ──  倒した相手を、後から理解する")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(3.9), Inches(2.5), C_DKRED, C_RED, 2)
    tb(s, Inches(0.32), Inches(0.92), Inches(3.6), Emu(340000),
       "普通AT（バトル型）", 12, bold=True, color=C_RED, font=FONT_H)
    tb(s, Inches(0.32), Inches(1.40), Inches(3.7), Inches(1.7),
       "祟り神を「敵」として戦い、倒す\n\n"
       "  ザコ戦 → 中ボス → ボス\n"
       "  勝利でセット継続・SP蓄積\n"
       "  敗北で悔しさ → リベンジ欲求",
       9.5, color=C_CREAM)

    rect(s, Inches(4.22), Inches(1.7), Inches(1.56), Emu(380000), RGBColor(0x10, 0x08, 0x20))
    tb(s, Inches(4.1), Inches(1.78), Inches(1.8), Emu(200000),
       "感情が逆転", 11, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    tb(s, Inches(4.05), Inches(2.1), Inches(1.9), Emu(250000),
       "敵  →  共感", 13, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

    rect_b(s, Inches(5.9), Inches(0.85), Inches(3.9), Inches(2.5), C_DKBLU, C_BLUE, 2)
    tb(s, Inches(6.02), Inches(0.92), Inches(3.6), Emu(340000),
       "特別AT（章型）", 12, bold=True, color=C_BLUE, font=FONT_H)
    tb(s, Inches(6.02), Inches(1.40), Inches(3.7), Inches(1.7),
       "祟り神の「なぜそうなったか」を知る\n\n"
       "  村人の視点・僧侶の視点・家族の視点\n"
       "  断片が集まって真実が見える\n"
       "  100〜1500枚獲得チャンス",
       9.5, color=C_CREAM)

    # 比較表
    rect(s, Inches(0.2), Inches(3.42), Inches(9.6), Emu(50000), C_GOLD)
    headers = ["比較軸", "従来機（一般的な戦闘系）", "祟り神の章"]
    hx = [Inches(0.2), Inches(2.1), Inches(5.85)]
    hw = [Inches(1.82), Inches(3.65), Inches(4.05)]
    for i, (h, x, w) in enumerate(zip(headers, hx, hw)):
        bg = RGBColor(0x28, 0x20, 0x00) if i == 2 else RGBColor(0x18, 0x18, 0x30)
        rect(s, x, Inches(3.50), w, Emu(320000), bg)
        col = C_GOLD if i == 2 else C_LTGRY
        tb(s, x + Emu(50000), Inches(3.54), w - Emu(80000), Emu(280000), h, 9, bold=True, color=col)

    rows = [
        ("来店動機",  "設定・期待値",       "物語の続きが見たい"),
        ("リピート",  "1回来店で完結",      "複数来店で章が完成する"),
        ("体験差",    "全員が同じ演出",     "章のランダム順で変わる"),
        ("感情変化",  "プレイ中のみ完結",   "プレイ後も「あの神は…」と残る"),
    ]
    ry = Inches(3.88)
    for j, (axis, old, new) in enumerate(rows):
        bg = RGBColor(0x10, 0x10, 0x22) if j % 2 == 0 else RGBColor(0x14, 0x14, 0x28)
        for x, w in zip(hx, hw):
            rect(s, x, ry, w, Emu(305000), bg)
        tb(s, hx[0] + Emu(50000), ry + Emu(28000), hw[0] - Emu(80000), Emu(260000), axis, 8.5, color=C_LTGRY)
        tb(s, hx[1] + Emu(50000), ry + Emu(28000), hw[1] - Emu(80000), Emu(260000), old,  8.5, color=C_GRAY)
        tb(s, hx[2] + Emu(50000), ry + Emu(28000), hw[2] - Emu(80000), Emu(260000), new,  8.5, bold=True, color=C_GOLD)
        ry += Emu(312000)


# ══════════════════════════════════════════════════════════════
#  SLIDE 3: ゲームフロー全体図（蛇行2段）
# ══════════════════════════════════════════════════════════════
def s_flow(prs):
    s = new_slide(prs)
    hdr(s, "ゲームフロー全体図  ──  普通ATと特別ATが交差する")

    BW  = Inches(2.75)
    BH  = Inches(1.05)
    GAP = Inches(0.22)
    R1Y = Inches(0.62)
    R2Y = Inches(2.15)
    X1  = Inches(0.18)
    X2  = X1 + BW + GAP
    X3  = X2 + BW + GAP
    # X3 + BW ≈ 8.87" → 右余白 1.13" (⊓コネクター用)

    row1 = [
        (X1, "通常時",
         "コイン投入・周期でCZ到来\nSPを積み上げ / ステージ4種",
         RGBColor(0x14, 0x10, 0x28), C_PUR),
        (X2, "普通AT（バトル型）",
         "祟り神を敵として戦い倒す\nザコ→中ボス→ボス / 勝利でセット継続",
         C_DKRED, C_RED),
        (X3, "特別AT抽選",
         "普通AT終了時・SP量に応じて抽選\n当選でどの章が来るか決定！",
         RGBColor(0x10, 0x08, 0x20), C_GOLD),
    ]
    for x, title, desc, fill, bdr in row1:
        rect_b(s, x, R1Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R1Y + Emu(50000), BW - Emu(160000), Emu(310000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R1Y + Emu(370000), BW - Emu(160000), BH - Emu(420000),
           desc, 8, color=C_CREAM)
    for x_l in [X1, X2]:
        arrow_r(s, x_l + BW + Emu(40000), R1Y + BH // 2, GAP - Emu(80000), C_RED)

    # ⊓ コネクター（右端、Row1→Row2の折り返し）
    CON_X = X3 + BW + Emu(80000)
    CON_R = CON_X + Emu(550000)
    LW    = Emu(55000)
    MID_Y = (R1Y + BH + R2Y) // 2
    rect(s, CON_X, R1Y + BH // 2, LW, MID_Y - (R1Y + BH // 2) + Emu(28000), C_GOLD)
    rect(s, CON_X, MID_Y, CON_R - CON_X + LW, LW, C_GOLD)
    rect(s, CON_R, R2Y, LW, MID_Y - R2Y + LW, C_GOLD)
    tb(s, CON_X - Emu(60000), MID_Y - Emu(350000), Emu(800000), Emu(330000),
       "↺", 16, bold=True, color=C_GOLD2, align=PP_ALIGN.CENTER)

    row2 = [
        (X3, "特別AT（章解放）",
         "1〜4章をランダム順で体験\n章クリアで加護を獲得！",
         C_DKBLU, C_BLUE),
        (X2, "加護蓄積",
         "章クリアごとに加護が増える\n4個集まると第5章が解放",
         RGBColor(0x20, 0x10, 0x30), C_PUR),
        (X1, "第5章AT",
         "加護全発動で「超楽しい」\n来るたびに体験が濃くなる",
         RGBColor(0x28, 0x20, 0x00), C_GOLD),
    ]
    for x, title, desc, fill, bdr in row2:
        rect_b(s, x, R2Y, BW, BH, fill, bdr, 1.8)
        tb(s, x + Emu(80000), R2Y + Emu(50000), BW - Emu(160000), Emu(310000),
           title, 9.5, bold=True, color=C_WHITE, font=FONT_H)
        tb(s, x + Emu(80000), R2Y + Emu(370000), BW - Emu(160000), BH - Emu(420000),
           desc, 8, color=C_CREAM)
    for x_r in [X3, X2]:
        _w = GAP - Emu(80000)
        _h = Emu(150000)
        shp = s.shapes.add_shape(13, x_r - GAP + Emu(40000), R2Y + BH // 2 - _h // 2, _w, _h)
        shp.rotation = 180
        shp.fill.solid()
        shp.fill.fore_color.rgb = C_PUR
        shp.line.fill.background()

    # 下部：第0章
    BOT_Y = Inches(3.38)
    rect_b(s, X1, BOT_Y, Inches(2.75), Emu(880000), C_DKBLU, C_BLUE, 2.0)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(55000), Inches(2.5), Emu(310000),
       "第0章（真実と成仏）", 9.5, bold=True, color=C_BLUE, font=FONT_H)
    tb(s, X1 + Emu(80000), BOT_Y + Emu(400000), Inches(2.5), Emu(440000),
       "全章クリア後に解放\n成仏チャレンジ：成功1500枚 / 失敗100枚\n「また来て救おう」という来店動機", 8, color=C_CREAM)

    rect(s, Inches(3.1), BOT_Y, Inches(6.72), Emu(880000), RGBColor(0x14, 0x10, 0x20))
    tb(s, Inches(3.25), BOT_Y + Emu(80000), Inches(6.3), Emu(320000),
       "フロー設計の核心：「倒す」と「理解する」の往復",
       9, bold=True, color=C_GOLD, font=FONT_H)
    tb(s, Inches(3.25), BOT_Y + Emu(440000), Inches(6.3), Emu(400000),
       "普通ATで倒した祟り神が、特別ATで「共感できる存在」になる。\n"
       "この感情の往復が、来店するたびに深まる体験を作る。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 章システム（DQ4型）
# ══════════════════════════════════════════════════════════════
def s_chapters(prs):
    s = new_slide(prs)
    hdr(s, "章システム  ──  DQ4型オムニバス・ランダム章解放設計")

    chapters = [
        ("第1章", "村人の視点",   "なぜ祟られたのか\n分からない",     C_LTGRY, "リプレイの力"),
        ("第2章", "僧侶の視点",   "封じるしかなかった\nあの日の決断", C_LTGRY, "ベルの力"),
        ("第3章", "家族の視点",   "あの人がなぜ\n祟り神に…",         C_LTGRY, "選択の力"),
        ("第4章", "祟り神の記憶", "守りたかった\nだけなのに",         C_RED,   "魂の力"),
        ("第5章", "全視点収束",   "加護4個で\n超楽しいAT",           C_GOLD,  "全加護発動"),
        ("第0章", "真実と成仏",   "1500枚一か八か\n感動のラスト",    C_BLUE,  "全章クリア限定"),
    ]

    cw = Inches(1.5)
    ch = Inches(3.0)
    cy = Inches(0.88)
    for i, (num, viewpoint, story, col, kago) in enumerate(chapters):
        cx = Inches(0.2) + i * (cw + Emu(100000))

        if num == "第5章":
            bg, bdr = RGBColor(0x28, 0x20, 0x00), C_GOLD
        elif num == "第0章":
            bg, bdr = C_DKBLU, C_BLUE
        else:
            bg, bdr = C_CARD, C_PUR

        rect_b(s, cx, cy, cw, ch, bg, bdr, 1.5)
        tb(s, cx + Emu(50000), cy + Emu(60000), cw - Emu(100000), Emu(340000),
           num, 11, bold=True, color=col, align=PP_ALIGN.CENTER, font=FONT_H)
        tb(s, cx + Emu(50000), cy + Emu(380000), cw - Emu(100000), Emu(340000),
           viewpoint, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        tb(s, cx + Emu(50000), cy + Emu(700000), cw - Emu(100000), Emu(480000),
           story, 8.5, color=C_CREAM, align=PP_ALIGN.CENTER)

        # 加護ボックス
        kbox_y = cy + ch - Emu(500000)
        if num not in ("第5章", "第0章"):
            rect(s, cx + Emu(50000), kbox_y, cw - Emu(100000), Emu(420000), RGBColor(0x20, 0x10, 0x30))
            tb(s, cx + Emu(80000), kbox_y + Emu(40000), cw - Emu(160000), Emu(380000),
               f"加護：\n{kago}", 7.5, color=C_PUR, align=PP_ALIGN.CENTER)
        elif num == "第5章":
            rect(s, cx + Emu(50000), kbox_y, cw - Emu(100000), Emu(420000), RGBColor(0x28, 0x18, 0x00))
            tb(s, cx + Emu(80000), kbox_y + Emu(40000), cw - Emu(160000), Emu(380000),
               "加護4個\n全部発動！", 7.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        else:
            rect(s, cx + Emu(50000), kbox_y, cw - Emu(100000), Emu(420000), C_DKBLU)
            tb(s, cx + Emu(80000), kbox_y + Emu(40000), cw - Emu(160000), Emu(380000),
               "成仏\nチャレンジ", 7.5, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # 下段説明
    rect(s, Inches(0.2), Inches(4.02), Inches(4.6), Emu(540000), RGBColor(0x10, 0x18, 0x10))
    tb(s, Inches(0.35), Inches(4.08), Inches(4.3), Emu(490000),
       "章はランダム順で解放。「俺は3章から入った、お前は？」\n"
       "→ ホールで会話が生まれる、体験がプレイヤーごとに変わる", 9, color=C_GREEN)

    rect(s, Inches(5.0), Inches(4.02), Inches(4.8), Emu(540000), RGBColor(0x18, 0x10, 0x00))
    tb(s, Inches(5.15), Inches(4.08), Inches(4.5), Emu(490000),
       "章クリアで加護を獲得 → 第5章のATが強化される\n"
       "→ 来るたびに体験が豊かになる「来店価値の積み上げ」", 9, color=C_GOLD)


# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 加護システム → 第5章
# ══════════════════════════════════════════════════════════════
def s_kago(prs):
    s = new_slide(prs)
    hdr(s, "加護システム  ──  積み上げた力が第5章で一斉発動する")

    kagos = [
        ("第1章", "リプレイの力", "リプレイ成立で\nAT抽選が走る", C_LTGRY),
        ("第2章", "ベルの力",     "打順ベル成立で\n上乗せ抽選",   C_LTGRY),
        ("第3章", "選択の力",     "毎セット\n2択チャレンジ発生", C_LTGRY),
        ("第4章", "魂の力",       "セット継続率が\n底上げされる", C_RED),
    ]

    kw = Inches(1.9)
    kh = Inches(1.65)
    ky = Inches(0.85)
    for i, (chap, name, effect, col) in enumerate(kagos):
        kx = Inches(0.2) + i * (kw + Emu(130000))
        rect_b(s, kx, ky, kw, kh, C_DKPUR, C_PUR, 1.5)
        tb(s, kx, ky + Emu(60000), kw, Emu(300000), chap, 9, color=C_PUR, align=PP_ALIGN.CENTER)
        tb(s, kx, ky + Emu(340000), kw, Emu(320000), name, 10, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, kx, ky + Emu(650000), kw, Emu(400000), effect, 8.5, color=C_CREAM, align=PP_ALIGN.CENTER)

    tb(s, Inches(4.35), Inches(2.6), Inches(1.3), Emu(360000),
       "4個全部\n集まると", 10, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 第5章ボックス
    rect_b(s, Inches(0.2), Inches(3.08), Inches(9.6), Inches(1.38),
           RGBColor(0x28, 0x20, 0x00), C_GOLD, 2.5)
    tb(s, Inches(0.3), Inches(3.14), Inches(9.3), Emu(360000),
       "第5章AT ── 同時多発で何かが起き続ける「超楽しい」",
       13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER, font=FONT_H)

    events = [
        ("リプレイの力", "リプレイ → AT抽選発動！"),
        ("ベルの力",     "ベル → 上乗せ抽選！"),
        ("選択の力",     "毎セット → 2択チャレンジ！"),
        ("魂の力",       "継続率底上げ → なかなか終わらない！"),
    ]
    ex_starts = [Inches(0.4), Inches(2.7), Inches(5.1), Inches(7.3)]
    for (kname, edesc), ex_val in zip(events, ex_starts):
        tb(s, ex_val, Inches(3.58), Inches(2.2), Emu(380000),
           f"{kname}\n→ {edesc}", 8.5, color=C_WHITE)

    rect(s, Inches(0.2), Inches(4.56), Inches(9.6), Emu(480000), RGBColor(0x10, 0x10, 0x10))
    tb(s, Inches(0.35), Inches(4.62), Inches(9.2), Emu(410000),
       "加護が少ないまま第5章に来ると体験が薄い → 「また章を集めに来る」動機が生まれる。"
       "加護4個が揃うと同時多発で何かが起き続ける状態に。これが来店継続の仕掛け。",
       9.5, bold=True, color=C_ORANGE)


# ══════════════════════════════════════════════════════════════
#  SLIDE 6: 第0章
# ══════════════════════════════════════════════════════════════
def s_chapter0(prs):
    s = new_slide(prs)
    hdr(s, "第0章  ──  全章コンプリート者だけが辿り着く「真実と成仏」")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.6), Inches(3.4), C_DKBLU, C_BLUE, 2)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.3), Emu(360000),
       "「真の章が解禁されました」", 11, bold=True, color=C_BLUE, font=FONT_H)
    tb(s, Inches(0.32), Inches(1.45), Inches(4.3), Inches(2.4),
       "全視点が交差して、真実が一枚の絵として完成する\n\n"
       "  なぜ村人は祟られたのか\n"
       "  なぜ僧侶は封じるしかなかったのか\n"
       "  なぜ家族は気づけなかったのか\n"
       "  なぜ祟り神はそうなったのか\n\n"
       "── 全部が繋がる",
       9.5, color=C_CREAM)

    rect_b(s, Inches(5.1), Inches(0.85), Inches(4.7), Inches(3.4), C_DKRED, C_RED, 2)
    tb(s, Inches(5.22), Inches(0.92), Inches(4.4), Emu(360000),
       "成仏チャレンジ  ──  一か八か", 11, bold=True, color=C_RED, font=FONT_H)

    rect(s, Inches(5.22), Inches(1.55), Inches(2.15), Inches(1.1), RGBColor(0x00, 0x22, 0x00))
    tb(s, Inches(5.22), Inches(1.60), Inches(2.15), Emu(350000),
       "成功", 18, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    tb(s, Inches(5.22), Inches(2.08), Inches(2.15), Emu(400000),
       "1500枚\n+「救えた」感動", 9.5, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    rect(s, Inches(7.52), Inches(1.55), Inches(2.15), Inches(1.1), RGBColor(0x30, 0x00, 0x00))
    tb(s, Inches(7.52), Inches(1.60), Inches(2.15), Emu(350000),
       "失敗", 18, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    tb(s, Inches(7.52), Inches(2.08), Inches(2.15), Emu(400000),
       "100枚\n+「また来て救おう」", 9.5, color=C_RED, align=PP_ALIGN.CENTER)

    tb(s, Inches(5.22), Inches(2.85), Inches(4.4), Emu(450000),
       "失敗しても「終わり」ではなく\n「また来たい理由」に変わる。\n物語はまだ終わっていないから。",
       9, color=C_LTGRY)

    rect(s, Inches(0.2), Inches(4.35), Inches(9.6), Emu(560000), RGBColor(0x00, 0x00, 0x30))
    rect(s, Inches(0.2), Inches(4.35), Emu(55000), Emu(560000), C_BLUE)
    tb(s, Inches(0.45), Inches(4.40), Inches(9.1), Emu(240000),
       "第0章は「ゴール」ではなく「また来たい気持ち」の設計",
       10, bold=True, color=C_BLUE, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.72), Inches(9.1), Emu(240000),
       "次の祟り神ストーリーへ続く可能性も。「この台の世界観をもっと知りたい」が来店継続の最終動機。",
       8.5, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 7: スペック・ベンチマーク
# ══════════════════════════════════════════════════════════════
def s_spec(prs):
    s = new_slide(prs)
    hdr(s, "スペック・ベンチマーク  ──  CV0.20〜0.25・後半維持率65%以上")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.1), Inches(3.3), C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(3.8), Emu(320000),
       "推奨スペック（目安）", 11, bold=True, color=C_GOLD, font=FONT_H)
    specs = [
        ("タイプ",     "スマスロ（L型）"),
        ("天井",       "600G（SP MAX時 300G）"),
        ("純増",       "3.8枚/G"),
        ("機械割",     "設定1: 97.5%  /  設定6: 109%"),
        ("特別AT出玉", "100〜1500枚（加護数・章で変動）"),
        ("CV目標",     "0.20〜0.25（設定非依存）"),
    ]
    sy = Inches(1.36)
    for j, (k, v) in enumerate(specs):
        bg = RGBColor(0x12, 0x12, 0x26) if j % 2 == 0 else RGBColor(0x16, 0x16, 0x2C)
        rect(s, Inches(0.25), sy, Inches(3.95), Emu(320000), bg)
        tb(s, Inches(0.32), sy + Emu(28000), Inches(1.3), Emu(270000), k, 8.5, color=C_GRAY)
        tb(s, Inches(1.67), sy + Emu(28000), Inches(2.4), Emu(270000), v, 9, bold=True, color=C_WHITE)
        sy += Emu(330000)

    # ベンチマーク表
    bx = Inches(4.5)
    bW = Inches(5.28)
    rect(s, bx, Inches(0.85), bW, Emu(360000), RGBColor(0x22, 0x18, 0x00))
    bench_headers = ["比較軸", "東京喰種", "モンキーV", "番長4", "本提案"]
    bcolw = [Emu(950000), Emu(840000), Emu(840000), Emu(800000), Emu(980000)]
    bhx_list = []
    cx = bx + Emu(30000)
    for cw in bcolw:
        bhx_list.append(cx)
        cx += cw
    for h, hx_val, cw in zip(bench_headers, bhx_list, bcolw):
        tb(s, hx_val, Inches(0.90), cw - Emu(30000), Emu(300000),
           h, 8.5, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    bench_rows = [
        ("CV値",      "0.18",  "0.23",  "0.68",  "0.20〜0.25"),
        ("来店動機",  "世界観", "SP引継", "期待値", "物語続き"),
        ("一か八か",  "○",    "△",    "◎",    "◎（最大1500枚）"),
        ("20代訴求",  "○",    "△",    "△",    "◎"),
        ("後半維持率","73.7%", "68.2%", "29.6%", "目標65%超"),
    ]
    rby = Inches(1.28)
    for j, row in enumerate(bench_rows):
        rbg = RGBColor(0x10, 0x10, 0x22) if j % 2 == 0 else RGBColor(0x14, 0x14, 0x28)
        rect(s, bx, rby, bW, Emu(330000), rbg)
        for k, (val, hx_val, cw) in enumerate(zip(row, bhx_list, bcolw)):
            col = C_GOLD if k == 4 else (C_RED if val in ("0.68", "29.6%") else C_LTGRY)
            tb(s, hx_val, rby + Emu(28000), cw - Emu(30000), Emu(278000),
               val, 8, color=col, align=PP_ALIGN.CENTER)
        rby += Emu(338000)

    rect(s, Inches(0.2), Inches(4.25), Inches(9.6), Emu(640000), RGBColor(0x10, 0x18, 0x10))
    rect(s, Inches(0.2), Inches(4.25), Emu(55000), Emu(640000), C_GREEN)
    tb(s, Inches(0.45), Inches(4.30), Inches(9.1), Emu(260000),
       "CV0.20〜0.25（東京喰種・モンキーターンV水準）を目標とし、設定1でも「体験の質が変わらない」台を設計。",
       9, bold=True, color=C_GREEN)
    tb(s, Inches(0.45), Inches(4.62), Inches(9.1), Emu(260000),
       "後半維持率65%以上、3ヶ月後も稼働している台を目指す。章システムが来店継続の構造的担保。",
       9, color=C_CREAM)
    net_note(s)


# ══════════════════════════════════════════════════════════════
#  SLIDE 8: ターゲット・市場考察
# ══════════════════════════════════════════════════════════════
def s_market(prs):
    s = new_slide(prs)
    hdr(s, "ターゲット・市場考察  ──  ストーリー消費世代×来店継続設計")

    rect_b(s, Inches(0.2), Inches(0.85), Inches(4.65), Inches(3.45), C_CARD, C_GOLD, 1.5)
    tb(s, Inches(0.32), Inches(0.92), Inches(4.4), Emu(320000),
       "メインターゲット", 11, bold=True, color=C_GOLD, font=FONT_H)
    targets = [
        ("年齢層",         "25〜45代男性",                           C_GOLD2),
        ("世代",           "アニメ・マンガ・ゲーム消費世代",          C_CREAM),
        ("プレイスタイル", "週1〜2来店・「続きが気になる」リピーター", C_CREAM),
        ("感情ニーズ",     "「物語に没入」「積み上げを感じたい」",   C_YELLOW),
        ("許容投資",       "3000〜8000円/日",                        C_ORANGE),
    ]
    ty = Inches(1.38)
    for label, val, col in targets:
        tb(s, Inches(0.32), ty, Inches(1.6), Emu(265000), label, 8.5, bold=True, color=C_GRAY, wrap=False)
        tb(s, Inches(1.97), ty, Inches(2.75), Emu(265000), val, 8.5, color=col, wrap=False)
        ty += Emu(268000)

    rect_b(s, Inches(5.05), Inches(0.85), Inches(4.75), Inches(3.45), C_CARD, C_RED, 1.5)
    tb(s, Inches(5.17), Inches(0.92), Inches(4.5), Emu(320000),
       "競合との差別化", 11, bold=True, color=C_RED, font=FONT_H)
    comps = [
        ("Re:ゼロ",        "タイムリープ・失敗前提",    "感情逆転＝理解が生まれる設計"),
        ("番長4",          "期待値・設定6狙い",         "物語が来店動機の台"),
        ("吉宗",           "1G連の瞬間爆発",            "章を集めることで深まる体験"),
        ("PHOENIX L.",     "炎の継承・積み上げ感情",    "こちらは「他者理解」が核心"),
    ]
    cy2 = Inches(1.38)
    for title, their, ours in comps:
        tb(s, Inches(5.17), cy2, Inches(1.3), Emu(265000), title, 8, bold=True, color=C_LTGRY, wrap=False)
        tb(s, Inches(6.52), cy2, Inches(1.55), Emu(265000), their, 7.5, color=C_GRAY, wrap=False)
        tb(s, Inches(8.12), cy2, Inches(1.6), Emu(265000), f"→ {ours}", 7.5, color=C_CREAM, wrap=False)
        cy2 += Emu(268000)

    rect(s, Inches(0.2), Inches(4.40), Inches(9.6), Emu(700000), RGBColor(0x14, 0x08, 0x20))
    rect(s, Inches(0.2), Inches(4.40), Emu(55000), Emu(700000), C_PUR)
    tb(s, Inches(0.45), Inches(4.46), Inches(9.1), Emu(260000),
       "なぜ今「他者理解」型なのか", 10, bold=True, color=C_PUR, font=FONT_H)
    tb(s, Inches(0.45), Inches(4.78), Inches(9.1), Emu(280000),
       "「強さを積み上げる」台が増える中、「感情を積み上げる」台が空席になっている。"
       "ストーリー消費に慣れた世代には、この設計が刺さると考える。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  SLIDE 9: まとめ
# ══════════════════════════════════════════════════════════════
def s_matome(prs):
    s = new_slide(prs)
    hdr(s, "まとめ  ──  祟り神の章が生み出す体験")

    cols_data = [
        (Inches(0.2),
         "感情が逆転する",
         "「倒した敵が\n共感できる存在になる」\n\n普通ATで倒し、\n特別ATで理解する。\nこの往復が\n来店動機の核心",
         C_RED, C_DKRED),
        (Inches(3.55),
         "来るたびに深まる",
         "「今日どの章が来た？」\n\n章はランダム順\nプレイヤーごとに\n異なる体験\n加護が積み上がり\n第5章が豊かになる",
         C_PUR, C_DKPUR),
        (Inches(6.9),
         "物語が終わらない",
         "「第0章で救えた？」\n\n失敗しても\n「また来たい」\nに変わる設計\n次の祟り神へ\n続く可能性も",
         C_BLUE, C_DKBLU),
    ]
    for x, title, desc, col, fill in cols_data:
        rect_b(s, x, Inches(0.85), Inches(3.0), Inches(3.3), fill, col, 2.0)
        tb(s, x + Emu(80000), Inches(0.92), Inches(2.8), Emu(330000),
           title, 10, bold=True, color=col, font=FONT_H)
        tb(s, x + Emu(80000), Inches(1.38), Inches(2.8), Inches(2.5), desc, 9, color=C_CREAM)

    rect(s, Inches(0.2), Inches(4.25), Inches(9.6), Emu(40000), C_RED)
    rect(s, Inches(0.2), Inches(4.33), Inches(9.6), Emu(760000), RGBColor(0x12, 0x08, 0x20))
    tb(s, Inches(0.35), Inches(4.38), Inches(9.2), Emu(270000),
       "3つの設計核：I 感情逆転設計  II DQ4型章システム  III 加護の積み上げ来店設計",
       9, bold=True, color=C_GOLD)
    tb(s, Inches(0.35), Inches(4.70), Inches(9.2), Emu(350000),
       "祟り神の章（感情逆転・物語没入）→ PHOENIX LEGACY（炎の継承・選択と積み上げ） → 次の提案へ。\n"
       "「倒した相手を、後から理解する。」——この一行が、すべての設計の出発点です。",
       9, color=C_CREAM)


# ══════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    slides = [
        ("タイトル",             s_title),
        ("コアコンセプト",       s_concept),
        ("ゲームフロー全体図",   s_flow),
        ("章システム",           s_chapters),
        ("加護システム",         s_kago),
        ("第0章",               s_chapter0),
        ("スペック・ベンチマーク", s_spec),
        ("ターゲット・市場考察", s_market),
        ("まとめ",               s_matome),
    ]

    print("=" * 55)
    print("  祟り神の章 v2 企画提案書ジェネレーター")
    print("=" * 55)
    print()
    for i, (name, func) in enumerate(slides, 1):
        print(f"  {i:2d}/{len(slides)} {name}")
        func(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\n保存完了: {OUT_PATH}\n")


if __name__ == "__main__":
    main()
